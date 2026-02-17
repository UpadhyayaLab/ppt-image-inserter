"""
Core image insertion functions for PowerPoint presentations.
"""

from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches
import os


def insert_image(
    ppt_path: str,
    slide_index: int,
    image_path: str,
    left: float,
    top: float,
    width: float,
    height: float
) -> None:
    """
    Insert an image into a PowerPoint slide at a custom position.

    Args:
        ppt_path (str): Path to the PowerPoint file
        slide_index (int): Slide number (0-based index, so slide 1 = index 0)
        image_path (str): Path to the image file
        left (float): Left position in inches
        top (float): Top position in inches
        width (float): Image width in inches
        height (float): Image height in inches

    Raises:
        FileNotFoundError: If PPT or image file doesn't exist
        IndexError: If slide_index is out of range
        Exception: For other errors during image insertion

    Example:
        >>> insert_image('presentation.pptx', 0, 'chart.png', 1.0, 2.0, 5.0, 4.0)
    """
    # Validate input files
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    if not os.path.exists(image_path):
        raise FileNotFoundError(f"Image file not found: {image_path}")

    try:
        # Load the presentation
        prs = Presentation(ppt_path)
    except Exception as e:
        raise FileNotFoundError(
            f"Could not open PowerPoint file: {ppt_path}. "
            f"File may be corrupted or in old .ppt format. "
            f"Original error: {e}"
        )

    # Validate slide index
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(
            f"Slide index {slide_index} out of range. "
            f"Presentation has {len(prs.slides)} slides (indices 0-{len(prs.slides)-1})"
        )

    # Get the target slide
    slide = prs.slides[slide_index]

    # Insert the image
    slide.shapes.add_picture(
        image_path,
        Inches(left),
        Inches(top),
        width=Inches(width),
        height=Inches(height)
    )

    # Save the presentation
    try:
        prs.save(ppt_path)
    except PermissionError:
        raise PermissionError(
            f"Permission denied when saving {ppt_path}. "
            f"Make sure the file is not open in PowerPoint."
        )
    except Exception as e:
        raise IOError(f"Error saving PowerPoint file: {e}")

    print(f"Successfully inserted {os.path.basename(image_path)} into slide {slide_index + 1}")


def insert_image_preserve_aspect(
    ppt_path: str,
    slide_index: int,
    image_path: str,
    left: float,
    top: float,
    width: Optional[float] = None,
    height: Optional[float] = None
) -> None:
    """
    Insert an image while preserving its aspect ratio.
    Specify either width or height, and the other dimension will be calculated.

    Args:
        ppt_path (str): Path to the PowerPoint file
        slide_index (int): Slide number (0-based index)
        image_path (str): Path to the image file
        left (float): Left position in inches
        top (float): Top position in inches
        width (float, optional): Image width in inches (height will be calculated)
        height (float, optional): Image height in inches (width will be calculated)

    Raises:
        ValueError: If both or neither width and height are specified
    """
    if (width is None and height is None) or (width is not None and height is not None):
        raise ValueError("Specify exactly one of width or height to preserve aspect ratio")

    # Validate input files
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    if not os.path.exists(image_path):
        raise FileNotFoundError(f"Image file not found: {image_path}")

    try:
        prs = Presentation(ppt_path)
    except Exception as e:
        raise FileNotFoundError(
            f"Could not open PowerPoint file: {ppt_path}. "
            f"File may be corrupted or in old .ppt format. "
            f"Original error: {e}"
        )

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(
            f"Slide index {slide_index} out of range. "
            f"Presentation has {len(prs.slides)} slides"
        )

    slide = prs.slides[slide_index]

    # Insert with only one dimension specified (aspect ratio preserved)
    if width is not None:
        slide.shapes.add_picture(
            image_path,
            Inches(left),
            Inches(top),
            width=Inches(width)
        )
    else:
        slide.shapes.add_picture(
            image_path,
            Inches(left),
            Inches(top),
            height=Inches(height)
        )

    try:
        prs.save(ppt_path)
    except PermissionError:
        raise PermissionError(
            f"Permission denied when saving {ppt_path}. "
            f"Make sure the file is not open in PowerPoint."
        )
    except Exception as e:
        raise IOError(f"Error saving PowerPoint file: {e}")

    print(f"Successfully inserted {os.path.basename(image_path)} with preserved aspect ratio")


def list_slides(ppt_path: str) -> List[str]:
    """
    List all slides in a PowerPoint presentation.

    Args:
        ppt_path (str): Path to the PowerPoint file

    Returns:
        list: List of slide indices and titles (if available)

    Example:
        >>> slides = list_slides('presentation.pptx')
        >>> print(slides)
        [0: Slide 1, 1: Slide 2, ...]
    """
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    prs = Presentation(ppt_path)
    slides_info = []

    for idx, slide in enumerate(prs.slides):
        # Try to get slide title if it exists
        title = "Untitled"
        if slide.shapes.title:
            title = slide.shapes.title.text or "Untitled"

        slides_info.append(f"{idx}: {title}")

    return slides_info
