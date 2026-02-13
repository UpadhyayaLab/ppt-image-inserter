"""
Metadata extraction utilities for PowerPoint presentations.
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
from typing import List, Dict, Any


def extract_image_metadata(ppt_path: str) -> List[Dict[str, Any]]:
    """
    Extract image source metadata from all slides in a presentation.

    This utility function iterates through all slides and extracts metadata
    from picture shapes, including the original image path stored in alt text
    (if it was stored using store_metadata=True in copy_slide_replace_image).

    Args:
        ppt_path (str): Path to the PowerPoint file

    Returns:
        list: List of dictionaries, each containing:
            - slide_index (int): 0-based slide index
            - slide_number (int): 1-based slide number (UI numbering)
            - original_path (str): Original image path from alt text (or None)
            - filename (str): Image filename
            - position (dict): Position/size dict with 'left', 'top', 'width', 'height'

    Raises:
        FileNotFoundError: If PowerPoint file doesn't exist

    Example:
        >>> metadata = extract_image_metadata('presentation.pptx')
        >>> for entry in metadata:
        ...     print(f"Slide {entry['slide_number']}: {entry['filename']}")
        ...     print(f"  Path: {entry['original_path']}")
        ...     print(f"  Position: {entry['position']}")
    """
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    prs = Presentation(ppt_path)
    metadata = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Extract stored metadata from alt text via XML API
                # (python-pptx 1.0.2 doesn't have descr/title properties)
                try:
                    original_path = shape._element._nvXxPr.cNvPr.get("descr", None)
                    filename = shape._element._nvXxPr.cNvPr.get("title", None)
                except Exception:
                    original_path = None
                    filename = None

                # Fallback to image.filename if title not set
                if filename is None and hasattr(shape, 'image'):
                    filename = shape.image.filename

                # Get position/size (convert from EMUs to inches)
                position = {
                    'left': shape.left / 914400,  # EMUs to inches
                    'top': shape.top / 914400,
                    'width': shape.width / 914400,
                    'height': shape.height / 914400
                }

                metadata.append({
                    'slide_index': slide_idx,
                    'slide_number': slide_idx + 1,  # UI numbering
                    'original_path': original_path,
                    'filename': filename,
                    'position': position
                })

    return metadata
