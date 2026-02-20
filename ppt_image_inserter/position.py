"""
Position and unit conversion utilities for PowerPoint images.
"""

from typing import Dict, List
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

# English Metric Units (EMU) conversion
# PowerPoint uses EMUs internally: 914400 EMUs = 1 inch
EMUS_PER_INCH = 914400.0


def cm_to_inches(cm: float) -> float:
    """
    Convert centimeters to inches.

    Args:
        cm (float): Length in centimeters

    Returns:
        float: Length in inches
    """
    return cm / 2.54


def get_image_position(
    ppt_path: str,
    slide_index: int,
    image_index: int = 0
) -> Dict[str, float]:
    """
    Extract position and size information from an image on a slide.

    Useful for getting template image parameters to use when inserting
    replacement images at the same location.

    Args:
        ppt_path (str): Path to the PowerPoint file
        slide_index (int): Slide number (0-based index)
        image_index (int): Which image to inspect if multiple exist (default: 0)

    Returns:
        dict: Position info with keys 'left', 'top', 'width', 'height' (all in inches)

    Raises:
        FileNotFoundError: If PPT file doesn't exist
        IndexError: If slide_index or image_index is out of range
        ValueError: If the slide has no pictures

    Example:
        >>> pos = get_image_position('presentation.pptx', 1, 0)
        >>> print(pos)
        {'left': 0.14, 'top': 0.96, 'width': 5.43, 'height': 2.72}
    """
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    prs = Presentation(ppt_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(
            f"Slide index {slide_index} out of range. "
            f"Presentation has {len(prs.slides)} slides"
        )

    slide = prs.slides[slide_index]

    # Find all pictures on the slide
    pictures = [shape for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]

    if not pictures:
        raise ValueError(f"Slide {slide_index} has no pictures")

    if image_index < 0 or image_index >= len(pictures):
        raise IndexError(
            f"Image index {image_index} out of range. "
            f"Slide has {len(pictures)} pictures (indices 0-{len(pictures)-1})"
        )

    # Get the specified picture
    picture = pictures[image_index]

    # Extract position and size (convert from EMUs to inches)
    position_info = {
        'left': picture.left / EMUS_PER_INCH,
        'top': picture.top / EMUS_PER_INCH,
        'width': picture.width / EMUS_PER_INCH,
        'height': picture.height / EMUS_PER_INCH
    }

    return position_info


def get_all_image_positions(
    ppt_path: str,
    slide_index: int
) -> List[Dict[str, float]]:
    """
    Extract positions and sizes of ALL images on a slide.

    Useful for templates with multiple placeholder images. Returns positions
    in the order they appear in the slide's shape collection.

    Args:
        ppt_path (str): Path to the PowerPoint file
        slide_index (int): Slide number (0-based index)

    Returns:
        list: List of position dicts, each with keys 'left', 'top', 'width', 'height' (all in inches)
              Returns empty list if slide has no pictures

    Raises:
        FileNotFoundError: If PPT file doesn't exist
        IndexError: If slide_index is out of range

    Example:
        >>> positions = get_all_image_positions('presentation.pptx', 1)
        >>> print(positions)
        [
            {'left': 0.5, 'top': 1.0, 'width': 4.0, 'height': 3.0},
            {'left': 5.0, 'top': 1.0, 'width': 4.0, 'height': 3.0}
        ]
    """
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    prs = Presentation(ppt_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(
            f"Slide index {slide_index} out of range. "
            f"Presentation has {len(prs.slides)} slides"
        )

    slide = prs.slides[slide_index]

    # Find all pictures on the slide
    pictures = [shape for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]

    # Convert all picture positions to dicts
    positions = []
    for picture in pictures:
        position_info = {
            'left': picture.left / EMUS_PER_INCH,
            'top': picture.top / EMUS_PER_INCH,
            'width': picture.width / EMUS_PER_INCH,
            'height': picture.height / EMUS_PER_INCH
        }
        positions.append(position_info)

    # Sort by visual reading order: top-to-bottom, then left-to-right within each row.
    # Round top to nearest 0.1" to group images that share a row despite minor pixel offsets.
    positions.sort(key=lambda p: (round(p['top'], 1), round(p['left'], 1)))

    return positions
