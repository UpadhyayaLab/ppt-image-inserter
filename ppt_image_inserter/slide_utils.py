"""
Slide manipulation utilities for PowerPoint presentations.
"""

from pptx import Presentation
from pptx.slide import Slide
from pptx.enum.shapes import MSO_SHAPE_TYPE
import copy
import os
from .backup import backup_presentation


def duplicate_slide(prs: Presentation, slide_index: int) -> Slide:
    """
    Duplicate a slide by creating a new slide and copying all shapes.

    This is an internal helper function. python-pptx doesn't have built-in
    slide duplication, so we manually copy shapes via XML element manipulation.

    Args:
        prs (Presentation): The Presentation object
        slide_index (int): Index of the slide to duplicate (0-based)

    Returns:
        Slide: The newly created slide object

    Raises:
        IndexError: If slide_index is out of range
    """
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(
            f"Slide index {slide_index} out of range. "
            f"Presentation has {len(prs.slides)} slides"
        )

    # Get the source slide
    source_slide = prs.slides[slide_index]

    # Create a new blank slide with the same layout
    blank_slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(blank_slide_layout)

    # Copy all shapes from source to new slide
    for shape in source_slide.shapes:
        # Deep copy the shape's XML element
        el = shape.element
        newel = copy.deepcopy(el)

        # Insert into the new slide's shape tree
        # Insert before 'p:extLst' if it exists, otherwise append
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return new_slide


def remove_pictures_from_slide(slide: Slide) -> int:
    """
    Remove all picture shapes from a slide.

    Uses XML manipulation to delete pictures. python-pptx doesn't have
    a built-in method for removing shapes.

    Args:
        slide (Slide): The slide to remove pictures from

    Returns:
        int: Number of pictures removed

    Note:
        This modifies the slide in-place. Changes are only saved when
        you call prs.save().
    """
    pictures_removed = 0

    # Iterate through shapes in reverse to avoid index issues during removal
    # Convert to list first since we're modifying the collection
    shapes_list = list(slide.shapes)

    for shape in shapes_list:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Remove via XML element manipulation
            parent = shape._element.getparent()
            if parent is not None:
                parent.remove(shape._element)
                pictures_removed += 1
            else:
                # Rare edge case - log but continue
                print(f"Warning: Could not remove picture '{shape.name}' - parent is None")

    return pictures_removed


def remove_all_text_from_slide(slide: Slide) -> int:
    """
    Remove ALL text boxes and placeholders from a slide.

    This removes both regular text boxes AND placeholders (title, body, etc.)
    to create a clean slide for custom content.

    Args:
        slide (Slide): The slide to remove text from

    Returns:
        int: Number of text elements removed
    """
    text_removed = 0
    shapes_list = list(slide.shapes)

    for shape in shapes_list:
        # Remove any shape with a text frame (includes both text boxes and placeholders)
        if shape.has_text_frame and shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            parent = shape._element.getparent()
            if parent is not None:
                parent.remove(shape._element)
                text_removed += 1

    return text_removed


def delete_slide(ppt_path: str, slide_index: int, backup_base: str = 'PPT/backups') -> None:
    """
    Delete a slide from the presentation using XML manipulation.

    Args:
        ppt_path (str): Path to the PowerPoint file
        slide_index (int): Slide index to delete (0-based, so slide 1 = index 0)
        backup_base (str): Base directory for backups (default: 'PPT/backups')

    Raises:
        FileNotFoundError: If PowerPoint file doesn't exist
        IndexError: If slide_index is out of range

    Example:
        >>> delete_slide('presentation.pptx', 3)  # Deletes slide 4 in PowerPoint UI
        >>> delete_slide('presentation.pptx', 3, backup_base='custom/backups')  # Custom backup location
    """
    # Validate input file
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    # Load presentation
    prs = Presentation(ppt_path)

    # Validate slide index
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(
            f"Slide index {slide_index} out of range (presentation has {len(prs.slides)} slides)"
        )

    # Create backup before deletion
    backups = backup_presentation(ppt_path, backup_base=backup_base)

    # Delete the slide using XML manipulation
    # Drop the relationship and remove from slide list
    rId = prs.slides._sldIdLst[slide_index].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[slide_index]

    # Save the presentation
    prs.save(ppt_path)
