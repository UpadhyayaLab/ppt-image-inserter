"""
High-level workflow functions for PowerPoint image replacement.
"""

from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt
import os
from .position import get_image_position, get_all_image_positions
from .slide_utils import duplicate_slide, remove_pictures_from_slide, remove_all_text_from_slide

# Text label defaults (inches)
LABEL_MARGIN = 0.3       # margin from right and bottom edges
LABEL_WIDTH = 5.0
LABEL_FONT_SIZE = 8
LABEL_FONT_NAME = 'Arial'


def _add_text_label(
    slide: Slide,
    image_paths,
    slide_width_inches: float = 13.333,
    slide_height_inches: float = 7.5,
    base_dir: str = None,
) -> None:
    """
    Add a text label with image path(s) in the bottom-right corner of a slide.

    Internal helper function for adding labels to slides.

    Args:
        slide: The slide to add the label to
        image_paths: Full path to the image file (str), or list of paths for multi-image slides
        slide_width_inches: Slide width in inches (used for right-alignment)
        slide_height_inches: Slide height in inches (used for bottom-alignment)
        base_dir: If provided, paths are shown relative to this directory
    """
    if isinstance(image_paths, str):
        image_paths = [image_paths]

    def _rel(p):
        # Only compute relative path for absolute paths — leave placeholders/labels untouched
        if base_dir and os.path.isabs(p):
            try:
                return os.path.relpath(p, base_dir).replace('\\', '/')
            except ValueError:
                return p  # Different drive on Windows — fall back to full path
        return p

    label_text = "\n".join(_rel(p) for p in image_paths)
    n_lines = len(image_paths)

    # Compute height from number of lines (~0.13" per line at 8pt + small buffer)
    height = n_lines * 0.13 + 0.1

    # Position: bottom-right corner
    left = slide_width_inches - LABEL_MARGIN - LABEL_WIDTH
    top = slide_height_inches - LABEL_MARGIN - height

    try:
        textbox = slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(LABEL_WIDTH),
            Inches(height)
        )

        text_frame = textbox.text_frame
        text_frame.word_wrap = False
        text_frame.text = label_text

        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(LABEL_FONT_SIZE)
            paragraph.font.name = LABEL_FONT_NAME

    except Exception as e:
        print(f"[WARNING] Could not add text label: {e}")


def copy_slide_replace_image(ppt_path, source_slide_index, new_image_path, position=None,
                             store_metadata=True, add_label=True, base_dir=None):
    """
    Copy a slide and replace its image with a new one (single-image version).

    This is a convenience wrapper around copy_slide_replace_images() for the common
    case of inserting a single image. It maintains backwards compatibility with
    existing code.

    Args:
        ppt_path (str): Path to the PowerPoint file
        source_slide_index (int): Index of the template slide to copy (0-based)
        new_image_path (str): Path to the new image to insert
        position (dict, optional): Position dict with keys 'left', 'top', 'width', 'height'
                                   If None, auto-detects from first image in source slide
        store_metadata (bool): If True, stores the image path in the picture's alt text. Default: True
        add_label (bool): If True, adds a visible text label. Default: True

    Returns:
        int: Index of the newly created slide

    Raises:
        FileNotFoundError: If PPT or image file doesn't exist
        IndexError: If source_slide_index is out of range
        ValueError: If position is None and source slide has no pictures

    Example:
        >>> # Auto-detect position
        >>> new_idx = copy_slide_replace_image('presentation.pptx', 1, 'new_image.tif')

        >>> # Specify exact position
        >>> pos = {'left': 0.14, 'top': 0.96, 'width': 5.43, 'height': 2.72}
        >>> new_idx = copy_slide_replace_image('presentation.pptx', 1, 'new_image.tif',
        ...                                      position=pos, add_label=False)
    """
    # Convert single position dict to list format for plural version
    positions = [position] if position is not None else None

    # Call the plural version with single-item lists
    return copy_slide_replace_images(
        ppt_path,
        source_slide_index,
        [new_image_path],  # Wrap in list
        positions=positions,
        store_metadata=store_metadata,
        add_label=add_label,
        base_dir=base_dir,
    )


def copy_slide_replace_images(ppt_path, source_slide_index, new_image_paths, positions=None,
                               store_metadata=True, add_label=False, base_dir=None):
    """
    Copy a slide and replace its images with new ones (supports multiple images per slide).

    This is the plural version of copy_slide_replace_image. It:
    1. Duplicates the source slide (creating a new slide at the end)
    2. Removes all pictures from the duplicated slide
    3. Inserts multiple new images at specified positions
    4. Optionally stores metadata in each picture's alt text
    5. Optionally adds visible text labels (disabled by default for multi-image slides)

    Args:
        ppt_path (str): Path to the PowerPoint file
        source_slide_index (int): Index of the template slide to copy (0-based)
        new_image_paths (list): List of image file paths to insert
        positions (list, optional): List of position dicts with keys 'left', 'top', 'width', 'height'
                                    If None, auto-detects from all images in source slide
        store_metadata (bool): If True, stores image paths in pictures' alt text. Default: True
        add_label (bool): If True, adds visible text labels. Default: False (usually too cluttered
                         for multi-image slides)

    Returns:
        int: Index of the newly created slide

    Raises:
        FileNotFoundError: If PPT or any image file doesn't exist
        IndexError: If source_slide_index is out of range
        ValueError: If positions is None and source slide has no pictures, or if
                   len(new_image_paths) != len(positions)

    Example:
        >>> # Auto-detect positions from template with 2 placeholder images
        >>> new_idx = copy_slide_replace_images(
        ...     'presentation.pptx',
        ...     1,  # Template is slide 2
        ...     ['left_image.png', 'right_image.png']
        ... )

        >>> # Specify exact positions for 3 images
        >>> positions = [
        ...     {'left': 0.5, 'top': 1.0, 'width': 3.0, 'height': 2.0},
        ...     {'left': 4.0, 'top': 1.0, 'width': 3.0, 'height': 2.0},
        ...     {'left': 7.5, 'top': 1.0, 'width': 3.0, 'height': 2.0}
        ... ]
        >>> new_idx = copy_slide_replace_images(
        ...     'presentation.pptx', 1,
        ...     ['img1.png', 'img2.png', 'img3.png'],
        ...     positions=positions
        ... )
    """
    # Validate input files
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    for img_path in new_image_paths:
        if not os.path.exists(img_path):
            raise FileNotFoundError(f"Image file not found: {img_path}")

    # Load presentation
    prs = Presentation(ppt_path)

    # Auto-detect positions if not specified
    if positions is None:
        positions = get_all_image_positions(ppt_path, source_slide_index)

        # If no images found in template, raise error
        if not positions:
            raise ValueError(
                f"Source slide {source_slide_index} has no pictures to use as position templates. "
                f"Add placeholder images to the template slide or specify positions manually."
            )

    # Validate that we have the same number of images and positions
    if len(new_image_paths) != len(positions):
        raise ValueError(
            f"Number of images ({len(new_image_paths)}) does not match number of positions "
            f"({len(positions)}). Template slide has {len(positions)} placeholder image(s), "
            f"but config provides {len(new_image_paths)} image(s)."
        )

    # Duplicate the source slide
    new_slide = duplicate_slide(prs, source_slide_index)

    # Get the new slide's index (it's added at the end)
    new_slide_index = len(prs.slides) - 1

    # Remove all pictures from the new slide
    num_removed = remove_pictures_from_slide(new_slide)

    # Remove all text (including placeholders) from the new slide
    num_text_removed = remove_all_text_from_slide(new_slide)

    # Insert all new images at their specified positions
    pictures = []
    for img_path, pos in zip(new_image_paths, positions):
        picture = new_slide.shapes.add_picture(
            img_path,
            Inches(pos['left']),
            Inches(pos['top']),
            width=Inches(pos['width']),
            height=Inches(pos['height'])
        )
        pictures.append(picture)

        # Store metadata in alt text if requested
        if store_metadata:
            try:
                picture._element._nvXxPr.cNvPr.set("descr", img_path)  # Full path
                picture._element._nvXxPr.cNvPr.set("title", os.path.basename(img_path))  # Filename
            except Exception as e:
                print(f"[WARNING] Could not store metadata for {os.path.basename(img_path)}: {e}")

    # Add visible text label in bottom-right corner if requested
    if add_label:
        slide_w = prs.slide_width.inches
        slide_h = prs.slide_height.inches
        _add_text_label(new_slide, new_image_paths, slide_w, slide_h, base_dir=base_dir)

    # Save the presentation
    prs.save(ppt_path)

    return new_slide_index


def replace_image_on_existing_slide(ppt_path: str, slide_index: int, new_image_path: str,
                                    store_metadata: bool = True, add_label: bool = True) -> None:
    """
    Replace the image on an existing slide with a new one.

    This function:
    1. Gets the position of the first image on the slide
    2. Removes all pictures from the slide
    3. Removes any existing text box labels (to avoid duplicates)
    4. Inserts the new image at the original position
    5. Optionally stores metadata in alt text
    6. Optionally adds a visible text label

    Args:
        ppt_path (str): Path to the PowerPoint file
        slide_index (int): Index of the slide to update (0-based)
        new_image_path (str): Path to the new image to insert
        store_metadata (bool): If True, stores image path in alt text. Default: True
        add_label (bool): If True, adds visible text label in bottom left. Default: True

    Raises:
        FileNotFoundError: If PPT or image file doesn't exist
        IndexError: If slide_index is out of range
        ValueError: If the slide has no pictures to determine position

    Example:
        >>> replace_image_on_existing_slide('presentation.pptx', 1, 'updated_image.tif')
    """
    # Validate input files
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    if not os.path.exists(new_image_path):
        raise FileNotFoundError(f"Image file not found: {new_image_path}")

    # Load presentation
    prs = Presentation(ppt_path)

    # Validate slide index
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(
            f"Slide index {slide_index} out of range (presentation has {len(prs.slides)} slides)"
        )

    slide = prs.slides[slide_index]

    # Get position from the first image on the slide
    print(f"Getting image position from slide {slide_index} (slide {slide_index + 1} in UI)...")
    position = get_image_position(ppt_path, slide_index, image_index=0)
    print(f"Detected position: left={position['left']:.2f}\", top={position['top']:.2f}\", "
          f"width={position['width']:.2f}\", height={position['height']:.2f}\"")

    # Remove all pictures from the slide
    print(f"Removing existing pictures...")
    num_removed = remove_pictures_from_slide(slide)
    print(f"Removed {num_removed} picture(s)")

    # Remove any existing text boxes (labels) to avoid duplicates
    print(f"Removing existing text labels...")
    shapes_list = list(slide.shapes)
    textboxes_removed = 0
    for shape in shapes_list:
        if shape.has_text_frame:
            # Check if it's a text box (not a placeholder or title)
            try:
                # Text boxes typically don't have placeholder format
                if not hasattr(shape, 'placeholder_format'):
                    parent = shape._element.getparent()
                    if parent is not None:
                        parent.remove(shape._element)
                        textboxes_removed += 1
            except Exception:
                # Skip if we can't determine - this is a best-effort cleanup
                pass
    print(f"Removed {textboxes_removed} text box(es)")

    # Insert the new image at the original position
    print(f"Inserting {os.path.basename(new_image_path)}...")
    picture = slide.shapes.add_picture(
        new_image_path,
        Inches(position['left']),
        Inches(position['top']),
        width=Inches(position['width']),
        height=Inches(position['height'])
    )

    # Store metadata in alt text if requested
    if store_metadata:
        try:
            picture._element._nvXxPr.cNvPr.set("descr", new_image_path)
            picture._element._nvXxPr.cNvPr.set("title", os.path.basename(new_image_path))
            print(f"Stored metadata: {new_image_path}")
        except Exception as e:
            print(f"[WARNING] Could not store metadata: {e}")

    # Add visible text label if requested
    if add_label:
        slide_w = prs.slide_width.inches
        slide_h = prs.slide_height.inches
        _add_text_label(slide, new_image_path, slide_w, slide_h)
        print(f"Added text label: {os.path.basename(new_image_path)}")

    # Save the presentation
    print(f"Saving presentation...")
    prs.save(ppt_path)

    print(f"[SUCCESS] Updated slide {slide_index} (slide {slide_index + 1} in PowerPoint UI)")


def add_label_to_existing_slide(ppt_path: str, slide_index: int, base_dir: str = None) -> bool:
    """
    Add a visible metadata label to an existing slide by reading image source paths from alt-text.

    Reads the 'descr' alt-text stored by copy_slide_replace_images() and adds a visible
    text label in the bottom-right corner. Falls back to shape name if no alt-text found.

    Args:
        ppt_path (str): Path to the PowerPoint file
        slide_index (int): Index of the slide to label (0-based)
        base_dir (str, optional): Base directory for computing relative paths in the label

    Returns:
        bool: True if a label was added, False if slide has no pictures
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    prs = Presentation(ppt_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]

    if not pictures:
        return False

    # Read descr alt-text from each picture using lxml tag search (robust across shape types)
    found_metadata = []
    for pic in pictures:
        descr = ""
        try:
            for elem in pic._element.iter():
                local = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                if local == 'cNvPr':
                    descr = elem.get("descr", "")
                    if descr:
                        break
        except Exception:
            pass
        found_metadata.append(descr)

    slide_w = prs.slide_width.inches
    slide_h = prs.slide_height.inches

    if any(found_metadata):
        # Show real source paths where available
        img_labels = [d if d else f"[placeholder {i + 1}]" for i, d in enumerate(found_metadata)]
        _add_text_label(slide, img_labels, slide_w, slide_h, base_dir=base_dir)
    else:
        # No metadata stored — just mark this as the template slide
        _add_text_label(slide, ["TEMPLATE SLIDE"], slide_w, slide_h)

    prs.save(ppt_path)
    return True
