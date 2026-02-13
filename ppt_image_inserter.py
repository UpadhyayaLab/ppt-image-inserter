"""
PowerPoint Image Insertion Utility

This module provides functions to insert images into PowerPoint presentations
at custom positions and sizes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import copy
import shutil
import datetime
import glob
from typing import List, Dict, Any, Optional


def insert_image(ppt_path, slide_index, image_path, left, top, width, height):
    """
    Insert an image into a PowerPoint slide at a custom position.

    Args:
        ppt_path (str): Path to the PowerPoint file
        slide_index (int): Slide number (0-based index, so slide 1 = index 0)
        image_path (str): Path to the image file
        left (float): Left position in inches
        top (float): Top position in inches
        width (float): Image width in inches
        height (float): Image height in inches

    Raises:
        FileNotFoundError: If PPT or image file doesn't exist
        IndexError: If slide_index is out of range
        Exception: For other errors during image insertion

    Example:
        >>> insert_image('presentation.pptx', 0, 'chart.png', 1.0, 2.0, 5.0, 4.0)
    """
    # Validate input files
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    if not os.path.exists(image_path):
        raise FileNotFoundError(f"Image file not found: {image_path}")

    try:
        # Load the presentation
        prs = Presentation(ppt_path)

        # Validate slide index
        if slide_index < 0 or slide_index >= len(prs.slides):
            raise IndexError(
                f"Slide index {slide_index} out of range. "
                f"Presentation has {len(prs.slides)} slides (indices 0-{len(prs.slides)-1})"
            )

        # Get the target slide
        slide = prs.slides[slide_index]

        # Insert the image
        slide.shapes.add_picture(
            image_path,
            Inches(left),
            Inches(top),
            width=Inches(width),
            height=Inches(height)
        )

        # Save the presentation
        prs.save(ppt_path)

        print(f"Successfully inserted {os.path.basename(image_path)} into slide {slide_index + 1}")

    except Exception as e:
        raise Exception(f"Error inserting image: {str(e)}")


def cm_to_inches(cm):
    """
    Convert centimeters to inches.

    Args:
        cm (float): Length in centimeters

    Returns:
        float: Length in inches
    """
    return cm / 2.54


def list_slides(ppt_path):
    """
    List all slides in a PowerPoint presentation.

    Args:
        ppt_path (str): Path to the PowerPoint file

    Returns:
        list: List of slide indices and titles (if available)

    Example:
        >>> slides = list_slides('presentation.pptx')
        >>> print(slides)
        [0: Slide 1, 1: Slide 2, ...]
    """
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    prs = Presentation(ppt_path)
    slides_info = []

    for idx, slide in enumerate(prs.slides):
        # Try to get slide title if it exists
        title = "Untitled"
        if slide.shapes.title:
            title = slide.shapes.title.text or "Untitled"

        slides_info.append(f"{idx}: {title}")

    return slides_info


def insert_image_preserve_aspect(ppt_path, slide_index, image_path, left, top, width=None, height=None):
    """
    Insert an image while preserving its aspect ratio.
    Specify either width or height, and the other dimension will be calculated.

    Args:
        ppt_path (str): Path to the PowerPoint file
        slide_index (int): Slide number (0-based index)
        image_path (str): Path to the image file
        left (float): Left position in inches
        top (float): Top position in inches
        width (float, optional): Image width in inches (height will be calculated)
        height (float, optional): Image height in inches (width will be calculated)

    Raises:
        ValueError: If both or neither width and height are specified
    """
    if (width is None and height is None) or (width is not None and height is not None):
        raise ValueError("Specify exactly one of width or height to preserve aspect ratio")

    # Validate input files
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    if not os.path.exists(image_path):
        raise FileNotFoundError(f"Image file not found: {image_path}")

    try:
        prs = Presentation(ppt_path)

        if slide_index < 0 or slide_index >= len(prs.slides):
            raise IndexError(
                f"Slide index {slide_index} out of range. "
                f"Presentation has {len(prs.slides)} slides"
            )

        slide = prs.slides[slide_index]

        # Insert with only one dimension specified (aspect ratio preserved)
        if width is not None:
            slide.shapes.add_picture(
                image_path,
                Inches(left),
                Inches(top),
                width=Inches(width)
            )
        else:
            slide.shapes.add_picture(
                image_path,
                Inches(left),
                Inches(top),
                height=Inches(height)
            )

        prs.save(ppt_path)
        print(f"Successfully inserted {os.path.basename(image_path)} with preserved aspect ratio")

    except Exception as e:
        raise Exception(f"Error inserting image: {str(e)}")


def duplicate_slide(prs, slide_index):
    """
    Duplicate a slide by creating a new slide and copying all shapes.

    This is an internal helper function. python-pptx doesn't have built-in
    slide duplication, so we manually copy shapes via XML element manipulation.

    Args:
        prs (Presentation): The Presentation object
        slide_index (int): Index of the slide to duplicate (0-based)

    Returns:
        Slide: The newly created slide object

    Raises:
        IndexError: If slide_index is out of range
    """
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(
            f"Slide index {slide_index} out of range. "
            f"Presentation has {len(prs.slides)} slides"
        )

    # Get the source slide
    source_slide = prs.slides[slide_index]

    # Create a new blank slide with the same layout
    blank_slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(blank_slide_layout)

    # Copy all shapes from source to new slide
    for shape in source_slide.shapes:
        # Deep copy the shape's XML element
        el = shape.element
        newel = copy.deepcopy(el)

        # Insert into the new slide's shape tree
        # Insert before 'p:extLst' if it exists, otherwise append
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return new_slide


def remove_pictures_from_slide(slide):
    """
    Remove all picture shapes from a slide.

    Uses XML manipulation to delete pictures. python-pptx doesn't have
    a built-in method for removing shapes.

    Args:
        slide (Slide): The slide to remove pictures from

    Returns:
        int: Number of pictures removed

    Note:
        This modifies the slide in-place. Changes are only saved when
        you call prs.save().
    """
    pictures_removed = 0

    # Iterate through shapes in reverse to avoid index issues during removal
    # Convert to list first since we're modifying the collection
    shapes_list = list(slide.shapes)

    for shape in shapes_list:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Remove via XML element manipulation
            parent = shape._element.getparent()
            if parent is not None:
                parent.remove(shape._element)
                pictures_removed += 1
            else:
                # Rare edge case - log but continue
                print(f"Warning: Could not remove picture '{shape.name}' - parent is None")

    return pictures_removed


def remove_all_text_from_slide(slide):
    """
    Remove ALL text boxes and placeholders from a slide.

    This removes both regular text boxes AND placeholders (title, body, etc.)
    to create a clean slide for custom content.

    Args:
        slide (Slide): The slide to remove text from

    Returns:
        int: Number of text elements removed
    """
    text_removed = 0
    shapes_list = list(slide.shapes)

    for shape in shapes_list:
        # Remove any shape with a text frame (includes both text boxes and placeholders)
        if shape.has_text_frame and shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            parent = shape._element.getparent()
            if parent is not None:
                parent.remove(shape._element)
                text_removed += 1

    return text_removed


def get_image_position(ppt_path, slide_index, image_index=0):
    """
    Extract position and size information from an image on a slide.

    Useful for getting template image parameters to use when inserting
    replacement images at the same location.

    Args:
        ppt_path (str): Path to the PowerPoint file
        slide_index (int): Slide number (0-based index)
        image_index (int): Which image to inspect if multiple exist (default: 0)

    Returns:
        dict: Position info with keys 'left', 'top', 'width', 'height' (all in inches)

    Raises:
        FileNotFoundError: If PPT file doesn't exist
        IndexError: If slide_index or image_index is out of range
        ValueError: If the slide has no pictures

    Example:
        >>> pos = get_image_position('presentation.pptx', 1, 0)
        >>> print(pos)
        {'left': 0.14, 'top': 0.96, 'width': 5.43, 'height': 2.72}
    """
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    prs = Presentation(ppt_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(
            f"Slide index {slide_index} out of range. "
            f"Presentation has {len(prs.slides)} slides"
        )

    slide = prs.slides[slide_index]

    # Find all pictures on the slide
    pictures = [shape for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]

    if not pictures:
        raise ValueError(f"Slide {slide_index} has no pictures")

    if image_index < 0 or image_index >= len(pictures):
        raise IndexError(
            f"Image index {image_index} out of range. "
            f"Slide has {len(pictures)} pictures (indices 0-{len(pictures)-1})"
        )

    # Get the specified picture
    picture = pictures[image_index]

    # Extract position and size (convert from EMUs to inches)
    # EMU = English Metric Units (914400 EMUs = 1 inch)
    position_info = {
        'left': picture.left / 914400.0,
        'top': picture.top / 914400.0,
        'width': picture.width / 914400.0,
        'height': picture.height / 914400.0
    }

    return position_info


def copy_slide_replace_image(ppt_path, source_slide_index, new_image_path, position=None,
                             store_metadata=True, add_label=True):
    """
    Copy a slide and replace its image with a new one.

    This is the main high-level function for batch image processing. It:
    1. Duplicates the source slide (creating a new slide at the end)
    2. Removes all pictures from the duplicated slide
    3. Inserts the new image at the specified position
    4. Optionally stores the original image path in the picture's alt text
    5. Optionally adds a visible text label showing filename and folder path

    Args:
        ppt_path (str): Path to the PowerPoint file
        source_slide_index (int): Index of the template slide to copy (0-based)
        new_image_path (str): Path to the new image to insert
        position (dict, optional): Position dict with keys 'left', 'top', 'width', 'height'
                                   If None, auto-detects from first image in source slide
        store_metadata (bool): If True, stores the original image path in the picture's
                              alt text fields (descr = full path, title = filename).
                              This allows tracking image sources after slide reordering.
                              Default: True
        add_label (bool): If True, adds a visible text box in bottom left corner showing
                         the filename and folder path. Useful for visual tracking while
                         viewing slides. Default: True

    Returns:
        int: Index of the newly created slide

    Raises:
        FileNotFoundError: If PPT or image file doesn't exist
        IndexError: If source_slide_index is out of range
        ValueError: If position is None and source slide has no pictures

    Example:
        >>> # Auto-detect position with metadata and label
        >>> new_idx = copy_slide_replace_image(
        ...     'presentation.pptx',
        ...     1,  # Copy slide 2
        ...     'new_image.tif'
        ... )
        >>> print(f"Created slide {new_idx + 1}")

        >>> # Specify exact position without label
        >>> pos = {'left': 0.14, 'top': 0.96, 'width': 5.43, 'height': 2.72}
        >>> new_idx = copy_slide_replace_image(
        ...     'presentation.pptx',
        ...     1,
        ...     'new_image.tif',
        ...     position=pos,
        ...     add_label=False
        ... )
    """
    # Validate input files
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    if not os.path.exists(new_image_path):
        raise FileNotFoundError(f"Image file not found: {new_image_path}")

    # Load presentation
    prs = Presentation(ppt_path)

    # Auto-detect position if not specified
    if position is None:
        print(f"Auto-detecting image position from slide {source_slide_index}...")
        position = get_image_position(ppt_path, source_slide_index, image_index=0)
        print(f"Detected position: left={position['left']:.2f}\", top={position['top']:.2f}\", "
              f"width={position['width']:.2f}\", height={position['height']:.2f}\"")

    # Duplicate the source slide
    print(f"Duplicating slide {source_slide_index}...")
    new_slide = duplicate_slide(prs, source_slide_index)

    # Get the new slide's index (it's added at the end)
    new_slide_index = len(prs.slides) - 1

    # Remove all pictures from the new slide
    print(f"Removing pictures from duplicated slide...")
    num_removed = remove_pictures_from_slide(new_slide)
    print(f"Removed {num_removed} picture(s)")

    # Remove all text (including placeholders) from the new slide
    print(f"Removing all text from duplicated slide...")
    num_text_removed = remove_all_text_from_slide(new_slide)
    print(f"Removed {num_text_removed} text element(s)")

    # Insert the new image at the specified position
    print(f"Inserting {os.path.basename(new_image_path)}...")
    picture = new_slide.shapes.add_picture(
        new_image_path,
        Inches(position['left']),
        Inches(position['top']),
        width=Inches(position['width']),
        height=Inches(position['height'])
    )

    # Store metadata in alt text if requested
    if store_metadata:
        # Use XML API to set alt text (python-pptx 1.0.2 doesn't have descr/title properties)
        try:
            picture._element._nvXxPr.cNvPr.set("descr", new_image_path)  # Full path
            picture._element._nvXxPr.cNvPr.set("title", os.path.basename(new_image_path))  # Filename
            print(f"Stored metadata: {new_image_path}")
        except Exception as e:
            print(f"[WARNING] Could not store metadata: {e}")

    # Add visible text label if requested
    if add_label:
        try:
            # Create text box in bottom left corner
            textbox = new_slide.shapes.add_textbox(
                Inches(0.5),    # Left: 0.5" from edge
                Inches(7.0),    # Top: Near bottom (standard slide is 7.5" tall)
                Inches(5.0),    # Width: 5 inches
                Inches(0.4)     # Height: 0.4 inches
            )

            # Set text content
            text_frame = textbox.text_frame
            folder_path = os.path.dirname(new_image_path)
            filename = os.path.basename(new_image_path)
            text_frame.text = f"File: {filename}\nPath: {folder_path}"

            # Format text - small font
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(8)  # 8pt font
                paragraph.font.name = 'Arial'

            print(f"Added text label: {filename}")
        except Exception as e:
            print(f"[WARNING] Could not add text label: {e}")

    # Save the presentation
    print(f"Saving presentation...")
    prs.save(ppt_path)

    print(f"[SUCCESS] Created slide {new_slide_index} (slide {new_slide_index + 1} in PowerPoint UI)")

    return new_slide_index


def extract_image_metadata(ppt_path: str) -> List[Dict[str, Any]]:
    """
    Extract image source metadata from all slides in a presentation.

    This utility function iterates through all slides and extracts metadata
    from picture shapes, including the original image path stored in alt text
    (if it was stored using store_metadata=True in copy_slide_replace_image).

    Args:
        ppt_path (str): Path to the PowerPoint file

    Returns:
        list: List of dictionaries, each containing:
            - slide_index (int): 0-based slide index
            - slide_number (int): 1-based slide number (UI numbering)
            - original_path (str): Original image path from alt text (or None)
            - filename (str): Image filename
            - position (dict): Position/size dict with 'left', 'top', 'width', 'height'

    Raises:
        FileNotFoundError: If PowerPoint file doesn't exist

    Example:
        >>> metadata = extract_image_metadata('presentation.pptx')
        >>> for entry in metadata:
        ...     print(f"Slide {entry['slide_number']}: {entry['filename']}")
        ...     print(f"  Path: {entry['original_path']}")
        ...     print(f"  Position: {entry['position']}")
    """
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    prs = Presentation(ppt_path)
    metadata = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Extract stored metadata from alt text via XML API
                # (python-pptx 1.0.2 doesn't have descr/title properties)
                try:
                    original_path = shape._element._nvXxPr.cNvPr.get("descr", None)
                    filename = shape._element._nvXxPr.cNvPr.get("title", None)
                except Exception:
                    original_path = None
                    filename = None

                # Fallback to image.filename if title not set
                if filename is None and hasattr(shape, 'image'):
                    filename = shape.image.filename

                # Get position/size (convert from EMUs to inches)
                position = {
                    'left': shape.left / 914400,  # EMUs to inches
                    'top': shape.top / 914400,
                    'width': shape.width / 914400,
                    'height': shape.height / 914400
                }

                metadata.append({
                    'slide_index': slide_idx,
                    'slide_number': slide_idx + 1,  # UI numbering
                    'original_path': original_path,
                    'filename': filename,
                    'position': position
                })

    return metadata


def batch_replace_images(config_file: str) -> List[int]:
    """
    Process multiple image replacements from a YAML configuration file.

    This function reads a YAML config file and processes multiple images in batch,
    creating a new slide for each image by copying a template slide and replacing
    its image.

    Args:
        config_file (str): Path to the YAML configuration file

    Returns:
        list: List of newly created slide indices

    Raises:
        FileNotFoundError: If config file, PPT, or image files don't exist
        ValueError: If configuration is invalid
        yaml.YAMLError: If YAML file is malformed

    Example YAML format:
        ```yaml
        presentation: "path/to/file.pptx"
        template_slide: 1
        auto_position: true
        images:
          - path: "image1.tif"
          - path: "image2.tif"
        ```

    Example usage:
        >>> slide_indices = batch_replace_images('config.yaml')
        >>> print(f"Created {len(slide_indices)} new slides: {slide_indices}")
        Created 3 new slides: [2, 3, 4]
    """
    # Import here to avoid circular dependency and allow standalone use
    from batch_config import load_batch_config

    print("=" * 60)
    print("BATCH IMAGE REPLACEMENT")
    print("=" * 60)
    print()

    # Load and validate configuration
    print(f"Loading configuration from: {config_file}")
    config = load_batch_config(config_file)
    print()

    # Extract configuration
    ppt_path = config['presentation']
    template_slide = config['template_slide']
    images = config['images']
    auto_position = config.get('auto_position', True)
    base_dir = config.get('base_dir', '')

    # Create backup before batch processing
    print("Creating backup before batch processing...")
    backups = backup_presentation(ppt_path)
    if backups:
        print(f"[BACKUP] Created backups in: {', '.join(backups.keys())}")
    print()

    # Determine position
    if auto_position:
        position = None  # Will be auto-detected
        print("Position mode: Auto-detect from template slide")
    else:
        position = config['position']
        print(f"Position mode: Manual (left={position['left']}\", top={position['top']}\", "
              f"width={position['width']}\", height={position['height']}\")")

    if base_dir:
        print(f"Base directory: {base_dir}")

    print()
    print(f"Processing {len(images)} image(s)...")
    print()

    # Process each image
    created_slides = []

    for i, img_entry in enumerate(images, 1):
        # Extract image path (handle both string and dict formats)
        if isinstance(img_entry, dict):
            img_path = img_entry['path']
        else:
            # New format: just filename, combine with base_dir
            if base_dir:
                img_path = os.path.join(base_dir, img_entry)
            else:
                img_path = img_entry

        print(f"[{i}/{len(images)}] Processing: {os.path.basename(img_path)}")
        print("-" * 60)

        try:
            new_slide_index = copy_slide_replace_image(
                ppt_path=ppt_path,
                source_slide_index=template_slide,
                new_image_path=img_path,
                position=position
            )
            created_slides.append(new_slide_index)
            print()

        except Exception as e:
            print(f"[ERROR] Processing image {i}: {e}")
            print()
            # Continue with next image instead of failing completely
            continue

    # Summary
    print("=" * 60)
    print("BATCH PROCESSING COMPLETE")
    print("=" * 60)
    print(f"[SUCCESS] Created {len(created_slides)} slide(s)")
    print(f"  New slide indices: {created_slides}")
    print(f"  PowerPoint slide numbers: {[idx + 1 for idx in created_slides]}")

    if len(created_slides) < len(images):
        failed_count = len(images) - len(created_slides)
        print(f"[WARNING] {failed_count} image(s) failed to process (see errors above)")

    return created_slides


def backup_presentation(ppt_path: str, backup_base: str = 'PPT/backups') -> Dict[str, str]:
    """
    Create backups in multiple time-interval categories (one per category).

    This function implements a smart backup strategy that maintains ONE backup
    per time interval category. When a backup is created, it overwrites any
    existing backup in that category. Backups are only created if the time
    threshold has been exceeded.

    Args:
        ppt_path (str): Path to the PowerPoint file to backup
        backup_base (str): Base directory for backups (default: 'PPT/backups')

    Returns:
        dict: Dictionary mapping category names to backup file paths for backups created
              Example: {'latest': 'PPT/backups/latest/file.pptx'}

    Categories and time intervals:
        - 'latest': Always creates a backup (0 seconds threshold)
        - '5min': Creates backup if >5 minutes since last backup in this category
        - '10min': Creates backup if >10 minutes since last backup in this category
        - '30min': Creates backup if >30 minutes since last backup in this category
        - 'hourly': Creates backup if >1 hour since last backup in this category

    Directory structure:
        PPT/backups/
        ├── latest/
        ├── 5min/
        ├── 10min/
        ├── 30min/
        └── hourly/

    Example:
        >>> backups = backup_presentation('presentation.pptx')
        >>> print(f"Created backups: {list(backups.keys())}")
        Created backups: ['latest', '5min', '10min', '30min', 'hourly']
    """
    # Validate input file
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    # Time intervals in seconds
    intervals = {
        'latest': 0,      # Always backup
        '5min': 300,      # 5 minutes = 300 seconds
        '10min': 600,     # 10 minutes = 600 seconds
        '30min': 1800,    # 30 minutes = 1800 seconds
        'hourly': 3600    # 1 hour = 3600 seconds
    }

    created_backups = {}
    timestamp = datetime.datetime.now()

    for category, threshold_seconds in intervals.items():
        # Create category directory
        category_dir = os.path.join(backup_base, category)
        os.makedirs(category_dir, exist_ok=True)

        # Check most recent backup in this category
        existing_backups = glob.glob(os.path.join(category_dir, '*.pptx'))
        should_backup = True

        if existing_backups and threshold_seconds > 0:
            # Get most recent backup by modification time
            latest_backup = max(existing_backups, key=os.path.getmtime)
            backup_time = datetime.datetime.fromtimestamp(os.path.getmtime(latest_backup))
            time_diff = (timestamp - backup_time).total_seconds()

            # Only backup if threshold exceeded
            should_backup = time_diff >= threshold_seconds

        if should_backup:
            # Use original filename (will overwrite existing backup)
            filename = os.path.basename(ppt_path)
            backup_path = os.path.join(category_dir, filename)

            # Copy the file (preserving metadata)
            shutil.copy2(ppt_path, backup_path)
            created_backups[category] = backup_path

    return created_backups


def replace_image_on_existing_slide(ppt_path: str, slide_index: int, new_image_path: str,
                                    store_metadata: bool = True, add_label: bool = True) -> None:
    """
    Replace the image on an existing slide with a new one.

    This function:
    1. Gets the position of the first image on the slide
    2. Removes all pictures from the slide
    3. Removes any existing text box labels (to avoid duplicates)
    4. Inserts the new image at the original position
    5. Optionally stores metadata in alt text
    6. Optionally adds a visible text label

    Args:
        ppt_path (str): Path to the PowerPoint file
        slide_index (int): Index of the slide to update (0-based)
        new_image_path (str): Path to the new image to insert
        store_metadata (bool): If True, stores image path in alt text. Default: True
        add_label (bool): If True, adds visible text label in bottom left. Default: True

    Raises:
        FileNotFoundError: If PPT or image file doesn't exist
        IndexError: If slide_index is out of range
        ValueError: If the slide has no pictures to determine position

    Example:
        >>> replace_image_on_existing_slide('presentation.pptx', 1, 'updated_image.tif')
    """
    # Validate input files
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    if not os.path.exists(new_image_path):
        raise FileNotFoundError(f"Image file not found: {new_image_path}")

    # Load presentation
    prs = Presentation(ppt_path)

    # Validate slide index
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(
            f"Slide index {slide_index} out of range (presentation has {len(prs.slides)} slides)"
        )

    slide = prs.slides[slide_index]

    # Get position from the first image on the slide
    print(f"Getting image position from slide {slide_index} (slide {slide_index + 1} in UI)...")
    position = get_image_position(ppt_path, slide_index, image_index=0)
    print(f"Detected position: left={position['left']:.2f}\", top={position['top']:.2f}\", "
          f"width={position['width']:.2f}\", height={position['height']:.2f}\"")

    # Remove all pictures from the slide
    print(f"Removing existing pictures...")
    num_removed = remove_pictures_from_slide(slide)
    print(f"Removed {num_removed} picture(s)")

    # Remove any existing text boxes (labels) to avoid duplicates
    print(f"Removing existing text labels...")
    shapes_list = list(slide.shapes)
    textboxes_removed = 0
    for shape in shapes_list:
        if shape.has_text_frame:
            # Check if it's a text box (not a placeholder or title)
            try:
                # Text boxes typically don't have placeholder format
                if not hasattr(shape, 'placeholder_format'):
                    parent = shape._element.getparent()
                    if parent is not None:
                        parent.remove(shape._element)
                        textboxes_removed += 1
            except:
                pass  # Skip if we can't determine
    print(f"Removed {textboxes_removed} text box(es)")

    # Insert the new image at the original position
    print(f"Inserting {os.path.basename(new_image_path)}...")
    picture = slide.shapes.add_picture(
        new_image_path,
        Inches(position['left']),
        Inches(position['top']),
        width=Inches(position['width']),
        height=Inches(position['height'])
    )

    # Store metadata in alt text if requested
    if store_metadata:
        try:
            picture._element._nvXxPr.cNvPr.set("descr", new_image_path)
            picture._element._nvXxPr.cNvPr.set("title", os.path.basename(new_image_path))
            print(f"Stored metadata: {new_image_path}")
        except Exception as e:
            print(f"[WARNING] Could not store metadata: {e}")

    # Add visible text label if requested
    if add_label:
        try:
            textbox = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(7.0),
                Inches(5.0),
                Inches(0.4)
            )

            text_frame = textbox.text_frame
            folder_path = os.path.dirname(new_image_path)
            filename = os.path.basename(new_image_path)
            text_frame.text = f"File: {filename}\nPath: {folder_path}"

            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(8)
                paragraph.font.name = 'Arial'

            print(f"Added text label: {filename}")
        except Exception as e:
            print(f"[WARNING] Could not add text label: {e}")

    # Save the presentation
    print(f"Saving presentation...")
    prs.save(ppt_path)

    print(f"[SUCCESS] Updated slide {slide_index} (slide {slide_index + 1} in PowerPoint UI)")


def delete_slide(ppt_path: str, slide_index: int) -> None:
    """
    Delete a slide from the presentation using XML manipulation.

    Args:
        ppt_path (str): Path to the PowerPoint file
        slide_index (int): Slide index to delete (0-based, so slide 1 = index 0)

    Raises:
        FileNotFoundError: If PowerPoint file doesn't exist
        IndexError: If slide_index is out of range

    Example:
        >>> delete_slide('presentation.pptx', 3)  # Deletes slide 4 in PowerPoint UI
    """
    # Validate input file
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")

    # Load presentation
    prs = Presentation(ppt_path)

    # Validate slide index
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(
            f"Slide index {slide_index} out of range (presentation has {len(prs.slides)} slides)"
        )

    # Create backup before deletion
    print(f"[INFO] Creating backup before deletion...")
    backups = backup_presentation(ppt_path)
    if backups:
        print(f"[BACKUP] Created backups in: {', '.join(backups.keys())}")

    print(f"[INFO] Deleting slide at index {slide_index} (slide {slide_index + 1} in PowerPoint UI)")

    # Delete the slide using XML manipulation
    # Drop the relationship and remove from slide list
    rId = prs.slides._sldIdLst[slide_index].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[slide_index]

    # Save the presentation
    prs.save(ppt_path)
    print(f"[SUCCESS] Deleted slide at index {slide_index}")


if __name__ == "__main__":
    # Example usage
    print("PowerPoint Image Inserter")
    print("Import this module and use the functions:")
    print("  - insert_image(ppt_path, slide_index, image_path, left, top, width, height)")
    print("  - list_slides(ppt_path)")
    print("  - cm_to_inches(cm)")
    print("  - insert_image_preserve_aspect(ppt_path, slide_index, image_path, left, top, width=X)")
    print("  - copy_slide_replace_image(ppt_path, source_slide_index, new_image_path, position=None)")
    print("  - batch_replace_images(config_file)")
