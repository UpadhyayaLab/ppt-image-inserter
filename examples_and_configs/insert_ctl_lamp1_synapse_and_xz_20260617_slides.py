"""
insert_ctl_lamp1_synapse_and_xz_20260617_slides.py

Combined montage deck for the 20260617 fixed CTL experiment
(L:/FF/Nucleus_granules/CTL_fixed/20260617_Fixed_CTLs_glass_centrosome_polarization_granules_nucleus_3min_12min).

LAMP1 plays the role H3K27me3 did in the Ciliobrevin Jurkat decks; MT (β-tubulin)
plays the centrosome-context role (this experiment has no dedicated centrosome
stain — the centrosome is the β-tubulin focus). The two conditions are
TIMEPOINTS — C1 = 3 min (left) and C2 = 12 min (right) — shown side by side, one
FOV (first chunk) per block.

As of 20260701, only the 3 min (C1) montages have landed; every 12 min (C2)
montage folder exists but is empty, so 12 min panels render as "(missing)" and
auto-fill on a rerun once those montages generate.

Block sections (each block = one slide):
  - Broadest-slice XY merge           — physical_scale_images/Lamp1_MT_nuc_bz
    (3-channel combined only; group 'broad')
  - Centrosome-slice XY merge         — physical_scale_images/Lamp1_MT_nuc_com
    (3-channel LAMP1+MT+Nuc at the centrosome's Z; group 'xy_phys')
  - Actin synapse mask                — actin/bottom_slice_seg      (group 'synapse',
    layout-pin only — no embedded scalebar)
  - XZ MIPs (MT/LAMP1/actin + nuc)    — physical_scale_images/*_xz  (group 'xz_phys')
  - 3-panel XZ (MT / LAMP1 / merge)   — physical_scale_images/Lamp1_MT_xz_panel[_nolines]
    (stacked full-width, 3 min over 12 min; group 'xzpanel_phys')
  - LAMP1 + MT synapse-plane XY       — physical_scale_images/Lamp1_MT_syn (group 'syn_phys')
  - Deepest-invagination slice merges — EXTRA_BLOCKS (group 'invag_slice')

  (Per-channel XY panel breakdowns (*_3p) are intentionally omitted — the user
   wants only the 3-channel combined at each flat slice; XZ panels are kept.)

All physical_scale_images/ groups share the 5 μm / 104 px scalebar.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_ctl_lamp1_synapse_and_xz_20260617_slides.py
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

EXPERIMENT_ROOT = (
    "L:/FF/Nucleus_granules/CTL_fixed/"
    "20260617_Fixed_CTLs_glass_centrosome_polarization_granules_nucleus_3min_12min"
)

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/CTL_Glass_Nucleus_Centrosome/CTL_fixed_LAMP1_20260617/"
    "CTL_fixed_LAMP1_combined_20260617.pptx"
)

CONDITIONS = [
    ("C1_3min_aCD3_ICAM1_3SI_660bTub_535Actin_488LAMP1_405Nuc",  "3 min"),
    ("C2_12min_aCD3_ICAM1_3SI_660bTub_535Actin_488LAMP1_405Nuc", "12 min"),
]

# Block list: (subpath_template_with_{cond}, title, scale_group).
# subpath has extra cropped/channels/ vs CilioD experiments.
BLOCKS = [
    # --- Broadest slice: the 3-channel combined view only (no single/2-channel
    # breakdowns — user wants just the three combined at each slice). ---
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_nuc_bz/montages",
        "LAMP1 + MT + Nuc, broadest slice — 3 min vs 12 min",
        "broad",
    ),
    # --- Single-slice XY merge: the 3-channel combined view AT THE CENTROSOME
    # SLICE only (rendered at the centrosome's Z; the centrosome is the
    # β-tubulin/MT focus). The old plain single-slice merges (nucleus, Lamp1_nuc,
    # MT_nuc, Lamp1_MT_nuc) were dropped — their Z-slice was unspecified — and no
    # per-channel panels are shown for the single slice. Landed 12 min-only so
    # far; 3 min renders "(missing)" until it generates. ---
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_nuc_com/montages",
        "LAMP1 + MT + Nuc, centrosome slice — 3 min vs 12 min",
        "xy_phys",
    ),
    # Granules + MT only (no nucleus) at the centrosome slice.
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_com/montages",
        "LAMP1 + MT, centrosome slice — 3 min vs 12 min",
        "xy_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_com_adaptive_merge/montages",
        "LAMP1 + MT, centrosome slice (adaptive) — 3 min vs 12 min",
        "xy_phys",
    ),
    # Panel (per-channel + merge) of granules + MT at the centrosome slice.
    # Only the adaptive-contrast variant exists; wide 3-panel, so it uses the
    # stacked full-width layout (3 min over 12 min) via its own group.
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_com_adaptive_panels/montages",
        "LAMP1 + MT, centrosome slice — channels + merge (adaptive) — 3 min vs 12 min",
        "companel_phys",
    ),
    # --- Actin synapse mask + XZ MIPs ---
    (
        "{cond}/cropped/channels/prog_fixed_cells/actin/bottom_slice_seg/montages",
        "Actin synapse mask (bottom slice) — 3 min vs 12 min",
        "synapse",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/MT_nuc_xz/montages",
        "MT + Nuc XZ MIP — 3 min vs 12 min",
        "xz_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_nuc_xz/montages",
        "LAMP1 + Nuc XZ MIP — 3 min vs 12 min",
        "xz_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_nuc_xz/montages",
        "LAMP1 + MT + Nuc XZ MIP — 3 min vs 12 min",
        "xz_phys",
    ),
    # Actin XZ MIPs (recently landed; same 'xz_phys' scale group). The '_planes'
    # variant overlays dotted cell top/bottom reference lines. The '_nolines'
    # siblings and the empty 'actin_nuc_xz_planes_nolines' are intentionally
    # omitted to avoid duplicate slides.
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_xz/montages",
        "Actin XZ MIP — 3 min vs 12 min",
        "xz_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_nuc_xz/montages",
        "Actin + Nuc XZ MIP — 3 min vs 12 min",
        "xz_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_nuc_xz_planes/montages",
        "Actin + Nuc XZ MIP, cell top/bottom marked — 3 min vs 12 min",
        "xz_phys",
    ),
    # No-lines variants of the actin XZ MIPs (same views, overlay lines removed).
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_xz_nolines/montages",
        "Actin XZ MIP (no lines) — 3 min vs 12 min",
        "xz_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_nuc_xz_planes_nolines/montages",
        "Actin + Nuc XZ MIP, planes (no lines) — 3 min vs 12 min",
        "xz_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_MT_xz_nolines/montages",
        "Actin + MT XZ MIP — 3 min vs 12 min",
        "xz_phys",
    ),
    # No-nucleus XZ MIPs (granules / microtubules alone).
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_xz_nolines/montages",
        "LAMP1 XZ MIP — 3 min vs 12 min",
        "xz_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/MT_xz_nolines/montages",
        "MT XZ MIP — 3 min vs 12 min",
        "xz_phys",
    ),
    # 3-panel XZ MIP: per-cell MT / LAMP1 / merge shown side by side. Much wider
    # aspect than a single XZ MIP, so its own scale group 'xzpanel_phys' (mixing
    # it into 'xz_phys' would shrink the single-panel XZ slides). No-lines variant
    # only. Landed 12 min-only so far — 3 min renders "(missing)" until it generates.
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_xz_panel_nolines/montages",
        "MT / LAMP1 / merge, XZ MIP panel (no lines) — 3 min vs 12 min",
        "xzpanel_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_MT_xz_panel_nolines/montages",
        "Actin / MT / merge, XZ MIP panel (no lines) — 3 min vs 12 min",
        "xzpanel_phys",
    ),
    # XY view of LAMP1 + MT at the synapse plane (no nucleus channel). Lives in
    # physical_scale_images/ so it carries the standard 104-px scalebar; put it
    # in its own scale group rather than the layout-pin-only 'synapse' (actin
    # mask) or the very-tall XZ 'xz_phys' group.
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_syn/montages",
        "LAMP1 + MT, synapse plane — 3 min vs 12 min",
        "syn_phys",
    ),
]

CHUNK_GLOB = "montage_cells_*.png"

# Extra "at-the-synapse" blocks. Source dirs use a deepest_invag_slice path
# under each channel's subdir, with non-standard filenames. The synapse-mask
# concept is "what is at the cell's contact plane (deepest invagination)":
#   - actin_synapse  : already a top BLOCK (actin/bottom_slice_seg)
#   - lamp1_with_MT  : LAMP1 granules merged with MT at the synapse
#   - MT_synapse     : MT alone at the synapse
# Each entry is (relative_subpath_template, filename_glob, title, scale_group).
# These are absent if the run didn't finish — we skip with a warning.
EXTRA_BLOCKS = [
    # Channel signal AT the Z-slice of the deepest nuclear invagination — a
    # nuclear-shape landmark, not the cell-substrate (synapse) interface.
    # Their pixel dims are much smaller than actin/bottom_slice_seg, so they
    # get their own scale group 'invag_slice' — pinning them with the actin
    # synapse mask was shrinking them to ~1/3 their natural fit.
    (
        "{cond}/cropped/channels/prog_fixed_cells/Lamp1/deepest_invag_slice/merges/montages_deepest_invag",
        "montage_cells_*_with_MT.png",
        "LAMP1 + MT, deepest invag slice — 3 min vs 12 min",
        "invag_slice",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/MT/deepest_invag_slice/merges/montages_deepest_invag",
        "montage_cells_*.png",
        "MT, deepest invag slice — 3 min vs 12 min",
        "invag_slice",
    ),
]

# Scalebar invariant for physical_scale_images/ XZ montages. Measured empirically
# on this dataset: every white bar is 104 px wide (same as CilioD => 5 μm bar at
# 20.8 px/μm).
SCALEBAR_PX = 104
SCALEBAR_UM = 5
PPUM_SOURCE = SCALEBAR_PX / SCALEBAR_UM

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

GRID_LEFT = 0.10
GRID_TOP = 0.60
CELL_W = 6.50
CELL_H = SLIDE_H - GRID_TOP - 0.10
LABEL_H = 0.30
IMG_H = CELL_H - LABEL_H
LABEL_FONT_PT = 16
COL_GAP = SLIDE_W - 2 * GRID_LEFT - 2 * CELL_W

CELL_POSITIONS = [
    (GRID_LEFT,                    GRID_TOP),
    (GRID_LEFT + CELL_W + COL_GAP, GRID_TOP),
]

# Stacked full-width layout for the wide 3-panel XZ slides (scale group
# PANEL_GROUP): 3 min on top, 12 min on bottom, each spanning the whole slide
# width. The panels are short-and-wide, so full-width rows render them much
# larger than the side-by-side cell layout.
# Scale groups that use the stacked full-width layout (3 min over 12 min).
PANEL_GROUPS = {"xzpanel_phys", "companel_phys"}
PANEL_IMG_W = SLIDE_W - 2 * GRID_LEFT
PANEL_ROW_H = (SLIDE_H - GRID_TOP - 0.10) / 2
PANEL_ROW_IMG_H = PANEL_ROW_H - LABEL_H
PANEL_ROW_TOPS = [GRID_TOP, GRID_TOP + PANEL_ROW_H]

# ---------------------------------------------------------------------------


def _winlong(p) -> str:
    s = os.path.abspath(str(p))
    if os.name == "nt" and not s.startswith("\\\\?\\"):
        s = "\\\\?\\" + s.replace("/", "\\")
    return s


def _exists_long(p) -> bool:
    return os.path.exists(_winlong(p))


def _png_dims(path: Path) -> Tuple[int, int]:
    with Image.open(_winlong(path)) as im:
        return im.size


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def compute_group_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_image_at_ppi(slide, image_path: Path, ppi: float,
                     area_left: float, area_top: float,
                     area_w: float, area_h: float):
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    left_in = area_left + (area_w - w_in) / 2
    top_in  = area_top  + (area_h - h_in) / 2
    return slide.shapes.add_picture(
        _winlong(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(w_in),
    )


def _chunk_start_index(p: Path) -> int:
    m = re.match(r"montage_cells_(\d+)", p.name)
    return int(m.group(1)) if m else 0


def find_first_chunk(montages_dir: Path, pattern: str = CHUNK_GLOB) -> Optional[Path]:
    # Long-path safe: pathlib.glob() breaks on paths near MAX_PATH even when
    # is_dir() returns True. Use os.listdir over a \\?\ prefixed path and
    # filter with fnmatch, then return a Path joined back to the original.
    import fnmatch
    long_dir = _winlong(montages_dir)
    if not os.path.isdir(long_dir):
        return None
    try:
        names = os.listdir(long_dir)
    except OSError:
        return None
    names = [n for n in names if fnmatch.fnmatch(n, pattern)]
    if not names:
        return None
    names.sort(key=lambda n: int(re.match(r"montage_cells_(\d+)", n).group(1))
               if re.match(r"montage_cells_(\d+)", n) else 0)
    return montages_dir / names[0]


def build_compare_slide(prs, title_text, left_label, left_img,
                        right_label, right_img, slide_ppi):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    cells = [
        (left_label,  left_img,  CELL_POSITIONS[0]),
        (right_label, right_img, CELL_POSITIONS[1]),
    ]
    missing = []
    for label, img_path, (cell_left, cell_top) in cells:
        add_textbox(
            slide, label,
            cell_left, cell_top, CELL_W, LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )
        if img_path is not None and _exists_long(img_path):
            add_image_at_ppi(slide, img_path, slide_ppi,
                             cell_left, cell_top + LABEL_H, CELL_W, IMG_H)
        else:
            add_textbox(
                slide, "(missing)",
                cell_left, cell_top + LABEL_H + IMG_H / 2 - 0.15, CELL_W, 0.3,
                font_pt=14, color=WHITE,
            )
            missing.append(label)
    return slide, missing


def build_stacked_slide(prs, title_text, top_label, top_img,
                        bottom_label, bottom_img, slide_ppi):
    """Full-width stacked layout: top_img above bottom_img, each spanning the
    whole slide width. Used for the wide 3-panel XZ slides so they render large."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    rows = [
        (top_label,    top_img,    PANEL_ROW_TOPS[0]),
        (bottom_label, bottom_img, PANEL_ROW_TOPS[1]),
    ]
    missing = []
    for label, img_path, row_top in rows:
        add_textbox(
            slide, label,
            GRID_LEFT, row_top, PANEL_IMG_W, LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )
        if img_path is not None and _exists_long(img_path):
            add_image_at_ppi(slide, img_path, slide_ppi,
                             GRID_LEFT, row_top + LABEL_H, PANEL_IMG_W, PANEL_ROW_IMG_H)
        else:
            add_textbox(
                slide, "(missing)",
                GRID_LEFT, row_top + LABEL_H + PANEL_ROW_IMG_H / 2 - 0.15,
                PANEL_IMG_W, 0.3, font_pt=14, color=WHITE,
            )
            missing.append(label)
    return slide, missing


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    root = Path(EXPERIMENT_ROOT)
    left_folder,  left_label  = CONDITIONS[0]
    right_folder, right_label = CONDITIONS[1]

    # Pre-pass: walk every block -> collect first-chunk slide specs.
    slide_specs: List[dict] = []
    all_blocks = [(sp, CHUNK_GLOB, t, sg) for (sp, t, sg) in BLOCKS] + list(EXTRA_BLOCKS)
    for subpath_tmpl, glob_pat, title, scale_group in all_blocks:
        left_dir  = root / Path(subpath_tmpl.format(cond=left_folder))
        right_dir = root / Path(subpath_tmpl.format(cond=right_folder))
        l_img = find_first_chunk(left_dir, glob_pat)
        r_img = find_first_chunk(right_dir, glob_pat)
        if l_img is None and r_img is None:
            print(f"WARNING: skipping '{title}' — no montages found "
                  f"(pattern '{glob_pat}' in '{subpath_tmpl}').")
            continue
        slide_specs.append({
            "title": title,
            "left_img": l_img, "right_img": r_img,
            "left_dir": left_dir, "right_dir": right_dir,
            "scale_group": scale_group,
        })

    if not slide_specs:
        print("ERROR: no slides to render — every block was empty.")
        sys.exit(1)

    # One PPI per scale_group, pinned to the largest source image in that group.
    group_ppi: dict = {}
    for spec in slide_specs:
        imgs = [p for p in (spec["left_img"], spec["right_img"])
                if p is not None and _exists_long(p)]
        sg = spec["scale_group"]
        if sg in PANEL_GROUPS:
            own = compute_group_ppi(imgs, PANEL_IMG_W, PANEL_ROW_IMG_H) if imgs else 0.0
        else:
            own = compute_group_ppi(imgs, CELL_W, IMG_H) if imgs else 0.0
        group_ppi[sg] = max(group_ppi.get(sg, 0.0), own)

    # Groups whose source is physical_scale_images/ carry the embedded 104 px
    # scalebar; the rest (actin/bottom_slice_seg, deepest_invag_slice merges)
    # do not, so we report them as layout-pin only.
    PHYS_GROUPS = {"xz_phys", "syn_phys", "broad", "xy_phys", "xzpanel_phys", "companel_phys"}
    print(f"Pinned PPI per scale_group ({len(slide_specs)} slides total):")
    for sg, ppi in sorted(group_ppi.items()):
        if sg in PHYS_GROUPS:
            bar = SCALEBAR_PX / ppi
            note = f"  scalebar = {bar:.3f} in = {bar * 2.54:.3f} cm"
        else:
            note = "  (layout-pin only — no embedded scalebar)"
        print(f"  {sg:>11s}  PPI={ppi:>7.2f}{note}")
    print(f"\nWriting deck to: {OUTPUT_PATH}\n")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    missing_total = []
    for spec in slide_specs:
        slide_ppi = group_ppi[spec["scale_group"]]
        builder = (build_stacked_slide if spec["scale_group"] in PANEL_GROUPS
                   else build_compare_slide)
        _, missing = builder(
            prs, spec["title"],
            left_label,  spec["left_img"],
            right_label, spec["right_img"],
            slide_ppi,
        )
        l = "OK" if spec["left_img"] else "MISSING"
        r = "OK" if spec["right_img"] else "MISSING"
        print(f"[{spec['scale_group']:>7s}]  3min:{l}  12min:{r}  {spec['title']}")
        for cell in missing:
            src = spec["left_dir"] if cell == left_label else spec["right_dir"]
            missing_total.append(f"{spec['title']}/{cell}  ({src})")

    prs.save(str(out_path))
    print(f"\nDone. {len(slide_specs)} slides written to:\n  {out_path}")
    if missing_total:
        print(f"\nMissing ({len(missing_total)}):")
        for m in missing_total:
            print(f"  - {m}")
    else:
        print("\nAll panels found.")


if __name__ == "__main__":
    main()
