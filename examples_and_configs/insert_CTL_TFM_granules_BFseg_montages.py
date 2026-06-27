"""
insert_CTL_TFM_granules_BFseg_montages.py

Assemble paired BF_seg granule-tracking montages (circle overlay + segmentation
overlay) into one slide each, for 13 CTL TFM granule experiments.

Source layout per experiment:
    <base_dir>/granule_TFM_prog/BF_seg/montages_circle/montage_cells_*.png
    <base_dir>/granule_TFM_prog/BF_seg/montages_seg/montage_cells_*.png

Each pair shares a filename. One slide per experiment, two images side-by-side,
aspect ratios preserved (width-constrained insertion).

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_CTL_TFM_granules_BFseg_montages.py
"""

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = "K:/FF/PPT/PPT_autogeneration/CTL_TFM_Granules/CTL_TFM_granules_BFseg_montages.pptx"

MONTAGE_SUBPATH = "granule_TFM_prog/BF_seg"

# (base_dir, condition_label, date_label)
# Order: 1.5 kPa control gels -> CK666 -> SMIFH2 -> DMSO
EXPERIMENTS = [
    # --- 1.5 kPa (old 200nm dark red beads) ---
    ("J:/FF/CTL_TFM_granules/old_TFM_200nm_dark_red_beads/1p5kPa/20230502_Gel1_3p0_3SI_aCD3_ICAM1_F-tractin_EGFP_LysoTracker",
     "1.5 kPa", "Gel1 (2023-05-02)"),
    ("J:/FF/CTL_TFM_granules/old_TFM_200nm_dark_red_beads/1p5kPa/20230502_Gel4_3p0_3SI_aCD3_ICAM1_F-tractin_EGFP_LysoTracker",
     "1.5 kPa", "Gel4 (2023-05-02)"),
    ("J:/FF/CTL_TFM_granules/old_TFM_200nm_dark_red_beads/1p5kPa/20230520_Gel4_3p0_3SI_aCD3_ICAM1_F-Tractin-EGFP_LysoTracker",
     "1.5 kPa", "Gel4 (2023-05-20)"),
    ("J:/FF/CTL_TFM_granules/old_TFM_200nm_dark_red_beads/1p5kPa/20230520_Gel6_3p0_3SI_aCD3_ICAM1_F-Tractin-EGFP_LysoTracker",
     "1.5 kPa", "Gel6 (2023-05-20)"),

    # --- CK666 100 uM ---
    ("J:/FF/CTL_TFM_granules/TFM_100nm_red_beads/20240813_Activated_CTLs_F-tractin-EGFP_LT-DeepRed/3p0_CK666_100uM/cells",
     "CK666 100 μM", "2024-08-13"),
    ("J:/FF/CTL_TFM_granules/TFM_100nm_red_beads/20241123_Activated_CTLs_F-tractin-EGFP_LT-DeepRed/3p0_CK666_100uM/cells",
     "CK666 100 μM", "2024-11-23"),
    ("J:/FF/CTL_TFM_granules/TFM_100nm_red_beads/20250128_Activated_CTLs_F-tractin-EGFP_LT-DeepRed/3p0_CK666_100uM/cells",
     "CK666 100 μM", "2025-01-28"),

    # --- SMIFH2 20 uM ---
    ("J:/FF/CTL_TFM_granules/TFM_100nm_red_beads/20240827_Activated_CTLs_F-tractin-EGFP_LT-DeepRed/3p0_SMIFH2_20uM/cells",
     "SMIFH2 20 μM", "2024-08-27"),
    ("J:/FF/CTL_TFM_granules/TFM_100nm_red_beads/20241210_Activated_CTLs_F-tractin-EGFP_LT-DeepRed/3p0_SMIFH2_20uM/cells",
     "SMIFH2 20 μM", "2024-12-10"),

    # --- DMSO ---
    ("J:/FF/CTL_TFM_granules/TFM_100nm_red_beads/20240827_Activated_CTLs_F-tractin-EGFP_LT-DeepRed/3p0_DMSO/cells",
     "DMSO", "2024-08-27"),
    ("J:/FF/CTL_TFM_granules/TFM_100nm_red_beads/20241123_Activated_CTLs_F-tractin-EGFP_LT-DeepRed/3p0_DMSO/cells",
     "DMSO", "2024-11-23"),
    ("J:/FF/CTL_TFM_granules/TFM_100nm_red_beads/20241210_Activated_CTLs_F-tractin-EGFP_LT-DeepRed/3p0_DMSO/cells",
     "DMSO", "2024-12-10"),
    ("J:/FF/CTL_TFM_granules/TFM_100nm_red_beads/20250128_Activated_CTLs_F-tractin-EGFP_LT-DeepRed/3p0_DMSO/cells",
     "DMSO", "2025-01-28"),
]

# Slide dimensions (widescreen)
SLIDE_WIDTH_IN  = 13.333
SLIDE_HEIGHT_IN = 7.5

# Title
TITLE_LEFT      = 1.167
TITLE_TOP       = 0.0
TITLE_WIDTH     = 11.0
TITLE_HEIGHT    = 0.85
TITLE_FONT_SIZE = Pt(32)

# Side-by-side images: circle (left), seg (right)
IMG_TOP         = 1.0
IMG_WIDTH       = 6.4
LEFT_IMG_LEFT   = 0.2
RIGHT_IMG_LEFT  = 6.73   # 0.2 + 6.4 + 0.13 gutter

# Label (bottom-left)
LABEL_LEFT      = 0.3
LABEL_FONT_SIZE = Pt(8)
LABEL_FONT_NAME = "Arial"

# ---------------------------------------------------------------------------


def find_first_montage(montage_dir: Path):
    """Return the montage with the lowest starting cell number."""
    matches = list(montage_dir.glob("montage_cells_*.png"))
    if not matches:
        return None

    def start_cell(p: Path) -> int:
        try:
            return int(p.stem.split("_")[2])
        except (IndexError, ValueError):
            return 9999

    return sorted(matches, key=start_cell)[0]


def insert_image_width_constrained(slide, image_path: Path, left_in: float, top_in: float, width_in: float) -> None:
    slide.shapes.add_picture(
        str(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(width_in),
    )


def add_title(slide, text: str) -> None:
    txBox = slide.shapes.add_textbox(
        Inches(TITLE_LEFT), Inches(TITLE_TOP),
        Inches(TITLE_WIDTH), Inches(TITLE_HEIGHT),
    )
    tf = txBox.text_frame
    tf.text = text
    para = tf.paragraphs[0]
    para.font.size = TITLE_FONT_SIZE
    para.font.bold = True
    para.alignment = PP_ALIGN.CENTER


def add_path_label(slide, full_path: str) -> None:
    label_height = 0.2
    label_width = SLIDE_WIDTH_IN - LABEL_LEFT - 0.3
    top = SLIDE_HEIGHT_IN - 0.3 - label_height
    txBox = slide.shapes.add_textbox(
        Inches(LABEL_LEFT), Inches(top),
        Inches(label_width), Inches(label_height),
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    tf.text = full_path
    para = tf.paragraphs[0]
    para.font.size = LABEL_FONT_SIZE
    para.font.name = LABEL_FONT_NAME


def main() -> None:
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]

    slides_added = 0
    missing_report = []

    for base_dir, condition_label, date_label in EXPERIMENTS:
        title = "{} — {}".format(condition_label, date_label)
        print("\n[{}]".format(title))

        base = Path(base_dir)
        circle_dir = base / MONTAGE_SUBPATH / "montages_circle"
        seg_dir    = base / MONTAGE_SUBPATH / "montages_seg"

        circle_img = find_first_montage(circle_dir) if circle_dir.exists() else None
        seg_img    = find_first_montage(seg_dir)    if seg_dir.exists()    else None

        print("  circle: {}".format(circle_img.name if circle_img else "NOT FOUND"))
        print("  seg:    {}".format(seg_img.name    if seg_img    else "NOT FOUND"))

        if not circle_dir.exists():
            missing_report.append("{}: montages_circle folder missing".format(base_dir))
        elif not circle_img:
            missing_report.append("{}/montages_circle: no montage found".format(base_dir))

        if not seg_dir.exists():
            missing_report.append("{}: montages_seg folder missing".format(base_dir))
        elif not seg_img:
            missing_report.append("{}/montages_seg: no montage found".format(base_dir))

        slide = prs.slides.add_slide(blank_layout)
        add_title(slide, title)

        if circle_img:
            insert_image_width_constrained(slide, circle_img, LEFT_IMG_LEFT,  IMG_TOP, IMG_WIDTH)
        if seg_img:
            insert_image_width_constrained(slide, seg_img,    RIGHT_IMG_LEFT, IMG_TOP, IMG_WIDTH)

        add_path_label(slide, base_dir)

        slides_added += 1

    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print("\nDone. {} slides saved to:\n  {}".format(slides_added, OUTPUT_PATH))

    if missing_report:
        print("\n[WARNING] Missing images ({}):".format(len(missing_report)))
        for msg in missing_report:
            print("  - {}".format(msg))
        sys.exit(1)
    else:
        print("\nAll images found.")


if __name__ == "__main__":
    main()
