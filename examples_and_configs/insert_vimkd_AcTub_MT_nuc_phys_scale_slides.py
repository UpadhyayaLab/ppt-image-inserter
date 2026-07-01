"""
insert_vimkd_AcTub_MT_nuc_phys_scale_slides.py

siCtrl (left) vs siVim / vimentin-knockdown (right) physical-scale montage deck of
MT + nucleus for the fixed Jurkat VimentinKD acetylated-tubulin (AcTub)
experiments (config/datasets/fixed/Jurkats/VimentinKD). Copied from
insert_vimkd_pericentrin_phys_scale_slides.py; identical machinery.

Same deck-wide PPI pinning pattern: all PNGs are inserted at one shared
pixels-per-inch so the embedded 104 px = 5 um scalebar renders at the same
cm-on-page within each scale group.

Two AcTub experiments with both conditions (siCtrl vs siVim, chronological):
  - 01/17/2024
  - 01/29/2024
(The third AcTub experiment, 20231024, is siVim-only -- no siCtrl -- so it cannot
form a siCtrl-vs-siVim comparison and is omitted.) chan_sub is "channels".
Channels imaged: MT (beta-tubulin) / actin / AcTub (acetylated tubulin) / nucleus.
This deck shows the MT (beta-tubulin) channel with the nucleus; the AcTub channel
is not montaged (no AcTub combos exist). vim is the knockdown target (not imaged),
so the siCtrl-vs-siVim columns carry the comparison.

Only experiments whose physical_scale_images exist render; the rest are skipped
and appear once their montages land (rerun to update).

Slide order is group-major: all of one montage group (across experiments,
date-sorted) before the next group. Every slide is siCtrl (left) | siVim (right),
first chunk only. Groups, in order:
  1. MT + Nuc XZ MIP
  2. Nuc (DNA), broadest slice
  3. MT + Nuc, broadest slice
  4. MT + Nuc, deepest invagination slice

Usage:
    python examples_and_configs/insert_vimkd_AcTub_MT_nuc_phys_scale_slides.py
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/VimentinKD/"
    "VimKD_Jurkats_siCtrl_vs_siVim_AcTub_MT_nuc_phys_scale_montages.pptx"
)

# Each experiment contributes one block of slides (one per combo below).
# tag      — short date tag woven into slide titles (the two 01/23/2024 entries
#            are the Vim- and MT-stained replicates of the same day).
# root     — dataset root dir.
# chan_sub — intermediate path inside each condition folder, before
#            prog_fixed_cells (differs per dataset).
# left/right — (condition_folder, display_label) for the DMSO / Noco columns.
# tp_label — activation string used in slide titles.
EXPERIMENTS = [
    {
        "tag": "01/17/2024",
        "root": (
            "L:/FF/Nucleus_MT/Jurkat_fixed/vimentinKD_tubulin-acetylation_fixed/"
            "20240117_MG_AcTub"
        ),
        "chan_sub": "channels",
        "left":  ("siCtrl_aCD3_E6-1_647BTub_535Actin_488AcTub_Hoechst", "siCtrl"),
        "right": ("siVim_aCD3_E6-1_647BTub_535Actin_488AcTub_Hoechst",  "siVim (KD)"),
        "tp_label": "αCD3",
    },
    {
        "tag": "01/29/2024",
        "root": (
            "L:/FF/Nucleus_MT/Jurkat_fixed/vimentinKD_tubulin-acetylation_fixed/"
            "20240129_MG_AcTub"
        ),
        "chan_sub": "channels",
        "left":  ("siCtrl_aCD3_E6-1_647BTub_535Actin_488AcTub_Hoechst", "siCtrl"),
        "right": ("siVim_aCD3_E6-1_647BTub_535Actin_488AcTub_Hoechst",  "siVim (KD)"),
        "tp_label": "αCD3",
    },
]

# (combo_subfolder, title_template, n_chunks, scale_group, fallback, opts).
# Template gets .format()'d with tp= (e.g. "αCD3") and tag= (e.g. "04/29/2022").
# Layout depends on n_chunks:
#   n_chunks == 1 -> 1 row x 2 cols (DMSO left, Noco right; labels above).
#   n_chunks  > 1 -> 1 row x (2 * n_chunks) cols, all in a single row, with a
#                    banner label spanning each condition's group of columns.
# scale_group: slides with the same group key share one pinned PPI across the
# ENTIRE deck (every experiment), so their embedded scalebars render at the
# same cm.
# opts is a dict with optional flags:
#   "per_exp_scale": bool   — compute this combo's PPI independently per
#                             experiment instead of pinning across all.
#   "compact_layout": bool  — smaller title/label bands so the image area grows.
#   "scale_mult":    float  — render at slide_ppi / scale_mult (>1 enlarges).
#   "solo_layout":   bool   — one condition per slide, image spans the full
#                             slide width (SOLO_CELL_W) so a wide combo (XZ MIP)
#                             can grow ~2x without colliding with its other half.
# Mirrors the CilioD physical-scale deck minus its H3K27me3 slide (no Noco
# analog); the Noco montages have no vim/MT combo.
# fallback is None for every combo: an experiment whose montages aren't
# generated yet contributes no slides at all (instead of placeholder "(missing)"
# slides) and fills in automatically once processed. actin_nuc_xz_nolines is
# present in every ready experiment, so the CilioD cent_nuc_xz fallback is unneeded.
COMBOS = [
    ("MT_nuc_xz",  "MT + Nuc XZ MIP ({tp}, {tag})",                       1, "xz",
        None, {}),
    ("nucleus_bz", "Nuc (DNA), broadest slice ({tp}, {tag})",            1, "broad_1c",
        None, {}),
    ("MT_nuc_bz",  "MT + Nuc, broadest slice ({tp}, {tag})",             1, "broad_1c",
        None, {}),
    ("MT_nuc",     "MT + Nuc, deepest invagination slice ({tp}, {tag})", 1, "broad_1c",
        None, {}),
]

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Slide layout (inches). 13.333 x 7.5 widescreen.
SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

# 1x2 cell grid below the title (label + image per cell)
GRID_LEFT = 0.10
GRID_TOP = 0.60
CELL_W = 6.50
CELL_H = SLIDE_H - GRID_TOP - 0.10   # 6.80"
LABEL_H = 0.30
IMG_H = CELL_H - LABEL_H             # 6.50"
LABEL_FONT_PT = 16
COL_GAP = SLIDE_W - 2 * GRID_LEFT - 2 * CELL_W

CELL_POSITIONS = [
    (GRID_LEFT,                    GRID_TOP),
    (GRID_LEFT + CELL_W + COL_GAP, GRID_TOP),
]

# Compact variant used by combos that opt into compact_layout=True (see COMBOS).
# Shrinks the title band, lifts the grid, and tightens the per-cell label so
# the image area gains ~0.30" of vertical room — ~5% larger images for
# height-bound combos.
COMPACT_TITLE_HEIGHT = 0.35
COMPACT_TITLE_FONT_PT = 22
COMPACT_GRID_TOP = 0.40
COMPACT_LABEL_H = 0.20
COMPACT_LABEL_FONT_PT = 14
COMPACT_CELL_H = SLIDE_H - COMPACT_GRID_TOP - 0.10   # 7.00"
COMPACT_IMG_H = COMPACT_CELL_H - COMPACT_LABEL_H     # 6.80"

# Solo-layout cell: one condition per slide, image spans full slide width.
# Used by combos that opt into solo_layout=True (see COMBOS) so a wide combo
# (e.g. XZ MIP) can render ~1.5-2x larger without overlapping its other panel.
SOLO_CELL_W = SLIDE_W - 2 * GRID_LEFT                # 13.133"

# Scalebar invariant for the Jurkat nucleus/actin fixed-cell physical-scale
# pipeline. Measured empirically with
# examples_and_configs/check_scalebar_pixel_widths.py against the Noco montages
# (04/29/2022 and 02/27/2024): every scalebar is exactly 104 px wide, i.e. the
# rendered PPUM is 104 / 5 = 20.8 px/μm — same as the CilioD nucleus pipeline.
SCALEBAR_PX = 104                            # px (measured)
SCALEBAR_UM = 5                              # μm
PPUM_SOURCE = SCALEBAR_PX / SCALEBAR_UM      # 20.8 px/μm in the rendered PNG

# ---------------------------------------------------------------------------


def _winlong(p) -> str:
    """Return a Win32-safe absolute path string. On Windows, prepends the
    \\\\?\\ extended-length prefix so paths near MAX_PATH (260) still work in
    os.stat / Image.open / python-pptx add_picture. pathlib's glob enumerates
    these files fine, but Path.exists() / os.stat() can silently fail just
    below the limit, so any path that crosses this layer must be wrapped."""
    s = os.path.abspath(str(p))
    if os.name == "nt" and not s.startswith("\\\\?\\"):
        s = "\\\\?\\" + s.replace("/", "\\")
    return s


def _exists_long(p) -> bool:
    """MAX_PATH-safe existence check (pathlib.Path.exists() trips on long paths)."""
    return os.path.exists(_winlong(p))


def montage_dir(root, cond_folder: str, chan_sub: str, combo: str) -> Path:
    """Build the montages dir for one (experiment, condition, combo).

    Noco layout differs from CilioD: each condition folder has its own
    intermediate `chan_sub` (e.g. "cells/channels", "tif/cells/channels",
    "bg_sub_fovs/cells/channels", "Cells/individual channels") before
    `prog_fixed_cells`, so the path is built per experiment rather than from a
    single shared subpath template."""
    return (Path(root) / cond_folder / chan_sub /
            "prog_fixed_cells" / "physical_scale_images" / combo / "montages")


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _png_dims(path: Path) -> Tuple[int, int]:
    """Return (width_px, height_px) of a PNG without fully decoding it."""
    with Image.open(_winlong(path)) as im:
        return im.size


def compute_slide_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    """Smallest ppi such that every image fits in (max_w_in x max_h_in).
    Called per slide (not deck-wide) so each slide's panels share one PPI —
    the embedded scalebars match across panels on a slide, but slides with
    smaller-pixel sources render at a smaller PPI and thus a larger
    scalebar-in-cm than slides with larger-pixel sources."""
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_image_at_ppi(slide, image_path: Path, ppi: float,
                     area_left: float, area_top: float,
                     area_w: float, area_h: float):
    """Center an image at uniform PPI inside an arbitrary (left, top, w, h)
    image area. native_px / ppi gives inches; both dims pinned by ppi so all
    images in the slide render at one physical scale."""
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    left_in = area_left + (area_w - w_in) / 2
    top_in  = area_top  + (area_h - h_in) / 2
    return slide.shapes.add_picture(
        _winlong(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(w_in),
    )


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _chunk_start_index(p: Path) -> int:
    m = re.match(r"montage_cells_(\d+)", p.name)
    return int(m.group(1)) if m else 0


def list_chunks(montages_dir) -> List[Path]:
    """Return the montage chunk PNGs in a dir, sorted by chunk-start index.
    Long-path-safe: pathlib's is_dir()/glob() silently fail on Windows paths
    past MAX_PATH (260) — they return False / empty even when the dir exists
    (verified on the 20240123 MT montages at ~279 chars, where the deep
    `tif/cells/channels/...` prefix tips the path over the limit). os.listdir
    on the \\\\?\\-prefixed path enumerates them correctly. Returns [] if absent."""
    d = _winlong(montages_dir)
    if not os.path.isdir(d):
        return []
    names = [f for f in os.listdir(d)
             if f.startswith("montage_cells_") and f.endswith(".png")]
    return sorted((Path(montages_dir) / n for n in names), key=_chunk_start_index)


def find_first_chunks(montages_dir: Path, n: int) -> List[Optional[Path]]:
    """Return the first n chunks sorted by chunk-start index. Pads with None
    if the folder has fewer chunks (or doesn't exist)."""
    chunks = list_chunks(montages_dir)
    return (chunks + [None] * n)[:n]


def _multichunk_geometry(n_chunks: int):
    """1 row x (2 * n_chunks) columns. DMSO's n_chunks chunks sit on the left,
    Noco's n_chunks chunks sit on the right, all in a single row. A banner
    label spans each condition's group of n_chunks columns above the images.
    Returns (col_w, img_h, col_lefts, img_top, banner_lefts, banner_widths)."""
    h_margin = GRID_LEFT          # slide left/right inset
    gap = 0.10                    # uniform gap between every column
    total_cols = 2 * n_chunks
    col_w = (SLIDE_W - 2 * h_margin - (total_cols - 1) * gap) / total_cols
    col_lefts = [h_margin + i * (col_w + gap) for i in range(total_cols)]

    img_top = GRID_TOP + LABEL_H
    img_h = IMG_H

    # Banner per condition group: spans from the first chunk's left edge to
    # the last chunk's right edge in that group.
    def _group(start_col):
        left = col_lefts[start_col]
        width = (col_lefts[start_col + n_chunks - 1] + col_w) - left
        return left, width
    left_banner_left,  left_banner_w  = _group(0)
    right_banner_left, right_banner_w = _group(n_chunks)
    banner_lefts   = [left_banner_left,  right_banner_left]
    banner_widths  = [left_banner_w,     right_banner_w]
    return col_w, img_h, col_lefts, img_top, banner_lefts, banner_widths


def build_compare_slide(prs, title_text,
                        left_label, left_imgs,
                        right_label, right_imgs,
                        slide_ppi, compact=False, solo=False):
    """Render the comparison slide. With n_chunks=1, lay out DMSO|Noco side
    by side (1-row layout). With n_chunks>1, switch to a single row with all
    DMSO chunks then all Noco chunks. All panels share slide_ppi so embedded
    scalebars match across the slide.
    compact=True shrinks title/label bands so the image area gains ~0.30"
    of vertical room (1-chunk slides only).
    solo=True ignores right_label/right_imgs and renders only the left panel
    at the full slide width (SOLO_CELL_W) — used to enlarge XZ MIPs."""
    n_chunks = len(left_imgs)
    assert len(right_imgs) == n_chunks

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    title_h    = COMPACT_TITLE_HEIGHT    if compact else TITLE_HEIGHT
    title_font = COMPACT_TITLE_FONT_PT   if compact else TITLE_FONT_PT
    grid_top   = COMPACT_GRID_TOP        if compact else GRID_TOP
    label_h    = COMPACT_LABEL_H         if compact else LABEL_H
    label_font = COMPACT_LABEL_FONT_PT   if compact else LABEL_FONT_PT
    img_h      = COMPACT_IMG_H           if compact else IMG_H

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, title_h,
        font_pt=title_font, color=WHITE, bold=True,
    )

    missing = []

    if solo:
        # Single panel filling the slide width. n_chunks > 1 not supported
        # for solo combos right now (XZ MIP is always n_chunks=1 here).
        cell_left = GRID_LEFT
        add_textbox(
            slide, left_label,
            cell_left, grid_top, SOLO_CELL_W, label_h,
            font_pt=label_font, color=WHITE, bold=True,
        )
        img_path = left_imgs[0]
        if img_path is not None and _exists_long(img_path):
            add_image_at_ppi(slide, img_path, slide_ppi,
                             cell_left, grid_top + label_h, SOLO_CELL_W, img_h)
        else:
            add_textbox(
                slide, "(missing)",
                cell_left, grid_top + label_h + img_h / 2 - 0.15,
                SOLO_CELL_W, 0.3,
                font_pt=14, color=WHITE,
            )
            missing.append(left_label)
        return slide, missing

    if n_chunks == 1:
        # 1 x 2 column layout: DMSO left, Noco right, label per column.
        cells = [
            (left_label,  left_imgs[0],  CELL_POSITIONS[0][0]),
            (right_label, right_imgs[0], CELL_POSITIONS[1][0]),
        ]
        for label, img_path, cell_left in cells:
            add_textbox(
                slide, label,
                cell_left, grid_top, CELL_W, label_h,
                font_pt=label_font, color=WHITE, bold=True,
            )
            if img_path is not None and _exists_long(img_path):
                add_image_at_ppi(slide, img_path, slide_ppi,
                                 cell_left, grid_top + label_h, CELL_W, img_h)
            else:
                add_textbox(
                    slide, "(missing)",
                    cell_left, grid_top + label_h + img_h / 2 - 0.15,
                    CELL_W, 0.3,
                    font_pt=14, color=WHITE,
                )
                missing.append(label)
        return slide, missing

    # n_chunks > 1: single row with DMSO chunks then Noco chunks side by side
    # (2 * n_chunks columns total). One banner label per condition group.
    col_w, img_h, col_lefts, img_top, banner_lefts, banner_widths = \
        _multichunk_geometry(n_chunks)

    for cond_idx, (cond_label, cond_imgs) in enumerate(
        [(left_label, left_imgs), (right_label, right_imgs)]
    ):
        add_textbox(
            slide, cond_label,
            banner_lefts[cond_idx], GRID_TOP, banner_widths[cond_idx], LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )
        for chunk_idx, img_path in enumerate(cond_imgs):
            col_idx = cond_idx * n_chunks + chunk_idx
            cell_left = col_lefts[col_idx]
            if img_path is not None and _exists_long(img_path):
                add_image_at_ppi(slide, img_path, slide_ppi,
                                 cell_left, img_top, col_w, img_h)
            else:
                add_textbox(
                    slide, "(missing)",
                    cell_left, img_top + img_h / 2 - 0.15,
                    col_w, 0.3,
                    font_pt=14, color=WHITE,
                )
                missing.append(f"{cond_label} (chunk {chunk_idx + 1})")
    return slide, missing


def _exp_date_key(exp):
    """(year, month, day) from the tag's MM/DD/YYYY prefix, for date-sorting
    experiments within each combo group. Same-date replicates (the 01/23/2024
    Vim and MT entries) keep EXPERIMENTS list order via Python's stable sort."""
    m = re.match(r"(\d{2})/(\d{2})/(\d{4})", exp["tag"])
    return (int(m.group(3)), int(m.group(1)), int(m.group(2))) if m else (9999, 99, 99)


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    # Pre-pass: walk every (experiment, combo) -> collect (left_imgs, right_imgs)
    # where each list has length n_chunks. Also stash scale_group + exp_key so
    # we can pin PPI across all slides in the same (experiment, scale_group).
    slide_specs: List[dict] = []
    # Group-major order: emit all of one combo across experiments before the
    # next combo, so the deck shows each montage group in turn. Experiments are
    # date-sorted within the group (stable, so same-date Vim/MT keep list order).
    for base_combo, base_title, n_chunks, scale_group, fallback, opts in COMBOS:
        for exp in sorted(EXPERIMENTS, key=_exp_date_key):
            root = Path(exp["root"])
            chan_sub = exp["chan_sub"]
            left_folder, left_label = exp["left"]
            right_folder, right_label = exp["right"]
            tag = exp["tag"]
            tp_label = exp["tp_label"]
            exp_key = f"{tag} {tp_label}"
            # Fresh per (combo, experiment) so a fallback swap can't leak to the
            # next experiment in this group.
            combo_folder, title_tmpl = base_combo, base_title
            # If the primary combo has no montages in this experiment:
            #   fallback tuple  -> use the listed (combo_folder, title) instead
            #   fallback is None -> skip this slide for this experiment
            primary_dir = montage_dir(root, left_folder, chan_sub, combo_folder)
            primary_ok = bool(list_chunks(primary_dir))
            if not primary_ok:
                if fallback is None:
                    continue
                combo_folder, title_tmpl = fallback
            title = title_tmpl.format(tp=tp_label, tag=tag)
            left_dir  = montage_dir(root, left_folder,  chan_sub, combo_folder)
            right_dir = montage_dir(root, right_folder, chan_sub, combo_folder)
            left_imgs  = find_first_chunks(left_dir,  n_chunks)
            right_imgs = find_first_chunks(right_dir, n_chunks)
            # Per-experiment scaling makes the PPI key unique per experiment,
            # so this combo's image only sizes against its own experiment's
            # source dims (not the deck-wide max).
            sg_key = (f"{scale_group}@{exp_key}"
                      if opts.get("per_exp_scale") else scale_group)
            common = {
                "log_key": f"{exp_key}/{combo_folder}",
                "n_chunks": n_chunks,
                "scale_group": sg_key,
                "exp_key": exp_key,
                "compact_layout": bool(opts.get("compact_layout")),
                "scale_mult": float(opts.get("scale_mult", 1.0)),
            }
            if opts.get("solo_layout"):
                # Emit one slide per condition, each rendered with the full
                # slide width. The "right" side of the spec is unused; the
                # solo build path reads only left_label / left_imgs / left_dir.
                for side_label, side_imgs, side_dir in (
                    (left_label,  left_imgs,  left_dir),
                    (right_label, right_imgs, right_dir),
                ):
                    slide_specs.append({
                        **common,
                        "title": f"{title} — {side_label}",
                        "left_label": side_label, "left_imgs": side_imgs,
                        "right_label": "",        "right_imgs": [None] * n_chunks,
                        "left_dir": side_dir,     "right_dir": side_dir,
                        "is_solo": True,
                    })
            else:
                slide_specs.append({
                    **common,
                    "title": title,
                    "left_label": left_label,   "left_imgs": left_imgs,
                    "right_label": right_label, "right_imgs": right_imgs,
                    "left_dir": left_dir,       "right_dir": right_dir,
                    "is_solo": False,
                })

    # For each slide compute the smallest PPI that fits its images in its cell
    # box; then take the max across every slide in the same scale_group (across
    # ALL experiments unless the spec opted into per-experiment scaling). Every
    # slide in that group is rendered at that pinned PPI, so all its scalebars
    # match cm-on-page across the deck within the group.
    def _cell_box(n_chunks: int, compact: bool, is_solo: bool = False):
        if is_solo:
            img_h = COMPACT_IMG_H if compact else IMG_H
            return SOLO_CELL_W, img_h
        if n_chunks == 1:
            img_h = COMPACT_IMG_H if compact else IMG_H
            return CELL_W, img_h
        col_w, img_h, *_ = _multichunk_geometry(n_chunks)
        return col_w, img_h

    group_ppi: dict = {}
    for spec in slide_specs:
        max_w, max_h = _cell_box(spec["n_chunks"], spec["compact_layout"],
                                 spec.get("is_solo", False))
        imgs = [p for p in (*spec["left_imgs"], *spec["right_imgs"])
                if p is not None and _exists_long(p)]
        own_ppi = compute_slide_ppi(imgs, max_w, max_h) if imgs else 0.0
        sg = spec["scale_group"]
        group_ppi[sg] = max(group_ppi.get(sg, 0.0), own_ppi)

    print(
        f"Deck-wide PPI pinning across {len(slide_specs)} slides. "
        f"Scalebar invariant: {SCALEBAR_UM} μm = {SCALEBAR_PX} px in source "
        f"(PPUM = {PPUM_SOURCE} px/μm — verify with check_scalebar_pixel_widths.py).\n"
    )
    print("Pinned PPI per scale_group:")
    for sg, ppi in sorted(group_ppi.items()):
        if ppi <= 0:
            print(f"  {sg:>9s}  PPI=   0.00  (no montages found yet)")
            continue
        bar = SCALEBAR_PX / ppi
        print(f"  {sg:>9s}  PPI={ppi:>7.2f}  "
              f"scalebar={bar:.3f} in = {bar * 2.54:.3f} cm")
    print(f"\nWriting deck to: {OUTPUT_PATH}\n")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    missing_total = []
    slides_added = 0
    for spec in slide_specs:
        title       = spec["title"]
        left_label  = spec["left_label"]
        left_imgs   = spec["left_imgs"]
        right_label = spec["right_label"]
        right_imgs  = spec["right_imgs"]
        left_dir    = spec["left_dir"]
        right_dir   = spec["right_dir"]
        log_key     = spec["log_key"]
        n_chunks    = spec["n_chunks"]
        scale_group = spec["scale_group"]
        # group_ppi can be 0 if no experiment in this group has montages yet;
        # add_image_at_ppi is only called for images that exist, so a 0 here is
        # harmless (every panel renders "(missing)").
        slide_ppi   = group_ppi[scale_group] / spec["scale_mult"]
        compact     = spec["compact_layout"]
        solo        = spec.get("is_solo", False)

        _, missing = build_compare_slide(
            prs, title,
            left_label, left_imgs,
            right_label, right_imgs,
            slide_ppi, compact=compact, solo=solo,
        )
        slides_added += 1

        left_ok  = sum(1 for p in left_imgs  if p is not None)
        right_ok = sum(1 for p in right_imgs if p is not None)
        print(
            f"[{log_key}]  group={scale_group:<9s}  "
            f"L:{left_ok}/{n_chunks}  R:{right_ok}/{n_chunks}"
        )

        for cell in missing:
            src = left_dir if cell.startswith(left_label) else right_dir
            missing_total.append(f"{log_key}/{cell}  ({src})")

    if slides_added == 0:
        print("\nNo slides built (no montages found yet) — nothing written. "
              "Rerun once the experiment's physical_scale_images have landed.")
        return
    prs.save(str(out_path))
    print(f"\nDone. {slides_added} slides written to:\n  {out_path}")

    if missing_total:
        print(f"\nMissing ({len(missing_total)}):")
        for m in missing_total:
            print(f"  - {m}")
    else:
        print("\nAll images found - no missing items.")


if __name__ == "__main__":
    main()
