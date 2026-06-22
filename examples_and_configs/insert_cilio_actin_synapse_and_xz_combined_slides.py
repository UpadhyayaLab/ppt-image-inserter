"""
insert_cilio_actin_synapse_and_xz_combined_slides.py

Combined deck for the 06/12/2026 Ciliobrevin Jurkats experiment, mirroring
the CAR T pipeline's combined synapse+XZ deck pattern (see
insert_actin_qc_synapse_mask_xz_mip_20260607_slides.py). Three combos in
one deck, one slide per chunk pair (DMSO left, Cilio right):

  Block 1  Actin synapse mask (bottom-slice actin segmentation)
           Source: actin/bottom_slice_seg/montages/
           Scale group: 'synapse' (layout-pin only — these don't carry a
           5 μm scalebar baked in).

  Block 2  Actin XZ MIP (actin only, dashed cell top/bottom lines)
           Source: physical_scale_images/actin_xz/montages/
           Scale group: 'xz_phys'  (5 μm = 104 px, physical scale).

  Block 3  Actin + Cent + Nuc XZ MIP (3 channels, dashed cell top/bottom)
           Source: physical_scale_images/actin_nuc_xz_planes/montages/
           Scale group: 'xz_phys'  (pinned with Block 2 → same cm-on-page).

Output: K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/CilioD/
        Cilio_Jurkats_actin_synapse_and_xz_06122026.pptx
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/CilioD/"
    "Cilio_Jurkats_actin_synapse_and_xz_06122026_06132026.pptx"
)

# Each experiment contributes one block of slides per non-empty BLOCKS entry.
# Schema mirrors insert_cilio_dmso_nuc_phys_scale_06122026_slides.py.
EXPERIMENTS = [
    {
        "tag": "06/12/2026",
        "root": (
            "L:/FF/Nucleus_H3K27me3/Ciliobrevin_Jurkats/"
            "06122026_firstReplicate_50uM30minCilio_7min_"
            "h3k27me3-640LP45_p561LP45_egfpCentrin2-488LP45_h405LP40_"
        ),
        "left":  ("A1_DMSO_7min_aCD3_",                 "DMSO, 7 min αCD3"),
        "right": ("A2_Ciliobrevin50um30min_7min_aCD3_", "Cilio 50 μM, 7 min αCD3"),
        "tp_label": "7 min αCD3",
    },
    {
        "tag": "06/13/2026",
        "root": (
            "L:/FF/Nucleus_H3K27me3/Ciliobrevin_Jurkats/"
            "06132026_50uM30minCilio_7-12min_"
            "h3k27me3-640LP45_p561LP45_egfpCentrin2-488LP45_h405LP40_"
        ),
        "left":  ("GbA1_DMSO5to1000-30min_7min_aCD3_", "DMSO, 7 min αCD3"),
        "right": ("GbA2_50uMCilio-30min_7min_aCD3_",   "Cilio 50 μM, 7 min αCD3"),
        "tp_label": "7 min αCD3",
    },
    {
        "tag": "06/13/2026",
        "root": (
            "L:/FF/Nucleus_H3K27me3/Ciliobrevin_Jurkats/"
            "06132026_50uM30minCilio_7-12min_"
            "h3k27me3-640LP45_p561LP45_egfpCentrin2-488LP45_h405LP40_"
        ),
        "left":  ("GaA1_5to1000DMSO-30min_12min_aCD3_", "DMSO, 12 min αCD3"),
        "right": ("GaA2_50uMCilio-30min_12min_aCD3_",   "Cilio 50 μM, 12 min αCD3"),
        "tp_label": "12 min αCD3",
    },
]

# Block list. Each block is (subpath_template, title_template, scale_group).
# subpath uses {cond}; title uses {tp} (timepoint) and {tag} (date).
# Blocks with no montages for a given experiment auto-skip (warning printed).
#
# Clean no-lines set: the synapse mask plus one XZ MIP per distinct channel
# combination, all using the *_nolines renderings (no dashed cell-top/bottom
# markers). Only the FIRST montage chunk (one FOV) per block is shown — see
# ONE_FOV below.
BLOCKS = [
    (
        "prog_fixed_cells/{cond}/actin/bottom_slice_seg/montages",
        "Actin synapse mask (bottom slice) ({tp}, {tag})",
        "synapse",
    ),
    (
        "prog_fixed_cells/{cond}/physical_scale_images/actin_nuc_xz_nolines/montages",
        "Actin + Nuc XZ MIP ({tp}, {tag})",
        "xz_phys",
    ),
    (
        "prog_fixed_cells/{cond}/physical_scale_images/actin_cent_nuc_xz_nolines/montages",
        "Actin + Cent + Nuc XZ MIP ({tp}, {tag})",
        "xz_phys",
    ),
    (
        "prog_fixed_cells/{cond}/physical_scale_images/actin_cent_xz_nolines/montages",
        "Actin + Cent XZ MIP ({tp}, {tag})",
        "xz_phys",
    ),
    # Nucleus + centrosome only (no actin). No dashed lines; only the *_xz
    # rendering exists (no separate _nolines variant needed).
    (
        "prog_fixed_cells/{cond}/physical_scale_images/cent_nuc_xz/montages",
        "Cent + Nuc XZ MIP ({tp}, {tag})",
        "xz_phys",
    ),
]

# Show only the first montage chunk (one FOV) per block per experiment.
ONE_FOV = True

CHUNK_GLOB = "montage_cells_*.png"

# Scalebar invariant for the physical_scale_images/ XZ montages.
SCALEBAR_PX = 104           # measured (5 μm == 104 px in source)
SCALEBAR_UM = 5
PPUM_SOURCE = SCALEBAR_PX / SCALEBAR_UM

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

GRID_LEFT = 0.10
GRID_TOP = 0.60
CELL_W = 6.50
CELL_H = SLIDE_H - GRID_TOP - 0.10
LABEL_H = 0.30
IMG_H = CELL_H - LABEL_H
LABEL_FONT_PT = 16
COL_GAP = SLIDE_W - 2 * GRID_LEFT - 2 * CELL_W

CELL_POSITIONS = [
    (GRID_LEFT,                    GRID_TOP),
    (GRID_LEFT + CELL_W + COL_GAP, GRID_TOP),
]

# ---------------------------------------------------------------------------


def _winlong(p) -> str:
    s = os.path.abspath(str(p))
    if os.name == "nt" and not s.startswith("\\\\?\\"):
        s = "\\\\?\\" + s.replace("/", "\\")
    return s


def _exists_long(p) -> bool:
    return os.path.exists(_winlong(p))


def _png_dims(path: Path) -> Tuple[int, int]:
    with Image.open(_winlong(path)) as im:
        return im.size


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def compute_group_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_image_at_ppi(slide, image_path: Path, ppi: float,
                     area_left: float, area_top: float,
                     area_w: float, area_h: float):
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    left_in = area_left + (area_w - w_in) / 2
    top_in  = area_top  + (area_h - h_in) / 2
    return slide.shapes.add_picture(
        _winlong(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(w_in),
    )


def _chunk_range(p: Path) -> Tuple[int, int]:
    m = re.match(r"montage_cells_(\d+)_(\d+)\.png", p.name)
    if not m:
        return (0, 0)
    return (int(m.group(1)), int(m.group(2)))


def list_chunks_sorted(montages_dir: Path) -> List[Path]:
    if not montages_dir.is_dir():
        return []
    return sorted(montages_dir.glob(CHUNK_GLOB), key=lambda p: _chunk_range(p)[0])


def build_compare_slide(prs, title_text, left_label, left_img,
                        right_label, right_img, slide_ppi):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    cells = [
        (left_label,  left_img,  CELL_POSITIONS[0]),
        (right_label, right_img, CELL_POSITIONS[1]),
    ]
    missing = []
    for label, img_path, (cell_left, cell_top) in cells:
        add_textbox(
            slide, label,
            cell_left, cell_top, CELL_W, LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )
        if img_path is not None and _exists_long(img_path):
            add_image_at_ppi(slide, img_path, slide_ppi,
                             cell_left, cell_top + LABEL_H, CELL_W, IMG_H)
        else:
            add_textbox(
                slide, "(missing)",
                cell_left, cell_top + LABEL_H + IMG_H / 2 - 0.15, CELL_W, 0.3,
                font_pt=14, color=WHITE,
            )
            missing.append(label)
    return slide, missing


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    def _pad(lst, n):
        return lst + [None] * (n - len(lst))

    # Pre-pass: walk every (experiment, block) -> collect per-chunk slide specs.
    slide_specs: List[dict] = []
    for exp in EXPERIMENTS:
        root = Path(exp["root"])
        left_folder,  left_label  = exp["left"]
        right_folder, right_label = exp["right"]
        tag = exp["tag"]
        tp_label = exp["tp_label"]
        exp_key = f"{tag} {tp_label}"
        for subpath_tmpl, title_tmpl, scale_group in BLOCKS:
            left_dir  = root / Path(subpath_tmpl.format(cond=left_folder))
            right_dir = root / Path(subpath_tmpl.format(cond=right_folder))
            left_chunks  = list_chunks_sorted(left_dir)
            right_chunks = list_chunks_sorted(right_dir)
            if ONE_FOV:
                # Keep only the first montage chunk (one FOV) per condition.
                left_chunks  = left_chunks[:1]
                right_chunks = right_chunks[:1]
            n_chunks = max(len(left_chunks), len(right_chunks))
            if n_chunks == 0:
                print(f"WARNING: skipping '{exp_key} / {subpath_tmpl}' — no montages.")
                continue

            left_chunks  = _pad(left_chunks,  n_chunks)
            right_chunks = _pad(right_chunks, n_chunks)

            for chunk_idx in range(n_chunks):
                slide_specs.append({
                    "title": title_tmpl.format(tp=tp_label, tag=tag),
                    "left_label":  left_label,
                    "right_label": right_label,
                    "left_img": left_chunks[chunk_idx],
                    "right_img": right_chunks[chunk_idx],
                    "left_dir": left_dir, "right_dir": right_dir,
                    "scale_group": scale_group,
                    "log_key": f"{exp_key}/{scale_group}/chunk {chunk_idx + 1}",
                })

    if not slide_specs:
        print("ERROR: no slides to render — every block was empty.")
        sys.exit(1)

    # One PPI per scale_group, pinned to the largest source image in that group.
    group_ppi: dict = {}
    for spec in slide_specs:
        imgs = [p for p in (spec["left_img"], spec["right_img"])
                if p is not None and _exists_long(p)]
        own = compute_group_ppi(imgs, CELL_W, IMG_H) if imgs else 0.0
        sg = spec["scale_group"]
        group_ppi[sg] = max(group_ppi.get(sg, 0.0), own)

    print(f"Pinned PPI per scale_group ({len(slide_specs)} slides total):")
    for sg, ppi in sorted(group_ppi.items()):
        if sg == "xz_phys":
            bar = SCALEBAR_PX / ppi
            note = f"  scalebar = {bar:.3f} in = {bar * 2.54:.3f} cm"
        else:
            note = "  (layout-pin only — no embedded scalebar)"
        print(f"  {sg:>9s}  PPI={ppi:>7.2f}{note}")
    print(f"\nWriting deck to: {OUTPUT_PATH}\n")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    missing_total = []
    for spec in slide_specs:
        slide_ppi = group_ppi[spec["scale_group"]]
        _, missing = build_compare_slide(
            prs, spec["title"],
            spec["left_label"],  spec["left_img"],
            spec["right_label"], spec["right_img"],
            slide_ppi,
        )
        l_status = "OK" if spec["left_img"]  is not None else "MISSING"
        r_status = "OK" if spec["right_img"] is not None else "MISSING"
        print(f"[{spec['log_key']}]  DMSO:{l_status}  Cilio:{r_status}")
        for cell in missing:
            src = spec["left_dir"] if cell == spec["left_label"] else spec["right_dir"]
            missing_total.append(f"{spec['log_key']}/{cell}  ({src})")

    prs.save(str(out_path))
    print(f"\nDone. {len(slide_specs)} slides written to:\n  {out_path}")
    if missing_total:
        print(f"\nMissing ({len(missing_total)}):")
        for m in missing_total:
            print(f"  - {m}")
    else:
        print("\nAll panels found.")


if __name__ == "__main__":
    main()
