"""
insert_CART_CatB_MT_actin_nuc_montages_slides.py

Combined nucleus/granule montage deck for the 5 CAR-T CTSB-mCherry (CatB)
datasets, modeled on the CTL nucleus/MT/granule montage deck
(insert_ctl_lamp1_synapse_and_xz_20260617_slides.py). CatB (CTSB-mCherry) plays
the granule role LAMP1 plays there; MT (beta-tubulin) plays the centrosome-context
role (these bTub datasets have no dedicated centrosome stain); actin and the
Hoechst/DNA nucleus round out the four channels.

The scientific variable is the CAR construct: CAT (left) vs FMC63 (right), shown
side by side, one FOV (first chunk) per cell. The 5 min and 15 min timepoints are
separate slides. Slides are grouped by combo (kind) — all groups of one combo
consecutively, chronological within each block.

Datasets (marker == "CatB" rows of the compiled manifest):

    Y: drive AP data (actin + CatB + nuc, no MT):
        20231018    5, 15 min
        20240312    5, 15 min
    J: drive (actin + MT + CatB + nuc):
        20231127    15 min only
        20240620    5, 15 min   (day3)
        20240624    5, 15 min   (day5)

SOURCE — the nucleus pipeline writes a *separate* progress folder,
`prog_fixed_cells_nuc/`, NOT the per-channel `prog_fixed_cells/` used by
insert_MT_CatB_actin_synapse_xz_mip_montages_slides.py. Under it,
`physical_scale_images/<combo>/montages/montage_cells_*.png` holds the CTL-style
combined overlays (actin_nuc_xz, CatB_MT_nuc_xz, nucleus_bz, ...), all at
104 px / 5 μm — identical to the CTL LAMP1 deck.

As of 2026-07-03 only the 3 J: (bTub) datasets have the `physical_scale_images/`
tree; the 2 Y: datasets have nucleus segmentation but that render stage has not
been run for them (and they have no MT channel), so their cells render
`(missing)` and auto-fill on a rerun once the stage generates. MT combos likewise
render `(missing)` for the two Y: (no-MT) datasets.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_CART_CatB_MT_actin_nuc_montages_slides.py
    # dry run (print per-combo/per-group CAT/FMC status + pinned PPI, build nothing):
    conda run -n PPT_editing python examples_and_configs/insert_CART_CatB_MT_actin_nuc_montages_slides.py --list
"""

import csv
import os
import re
import sys
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation  # noqa: E402

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

MANIFEST_CSV = (
    "L:/FF/CAR T/actin_compiled_results/all_datasets_actin/"
    "compiled_20260623/dataset_manifest.csv"
)

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/CART/nucleus/"
    "CART_CatB_MT_actin_nuc_montages_20260703.pptx"
)

# Only include manifest rows whose marker matches this set (all 5 CatB datasets).
MARKER_FILTER = {"CatB"}

# The nucleus pipeline's progress folder (sibling of the per-channel
# prog_fixed_cells/). base_dir in the manifest already ends at cells/channels.
PROGRESS_FOLDER = "prog_fixed_cells_nuc"

# Combo table: (subpath under <base_dir>/<PROGRESS_FOLDER>/, title, scale_group).
# Grouped by combo — all groups of one combo consecutively, chronological within
# each block. XY/slice views first, then XZ MIPs, then wide multi-panel views.
# `physical_scale_images/` combos carry the embedded 104 px / 5 μm scalebar; the
# actin synapse mask does not (layout-pin only). CatB = granule; MT = centrosome.
COMBOS = [
    # --- Broadest / centrosome slice XY (3-channel and nucleus-only) ---
    ("physical_scale_images/CatB_MT_nuc_bz/montages",
     "CatB + MT + Nuc, broadest slice", "broad"),
    ("physical_scale_images/nucleus_bz/montages",
     "Nucleus (DNA), broadest slice", "broad_1c"),
    ("physical_scale_images/CatB_MT_nuc_com/montages",
     "CatB + MT + Nuc, centrosome slice", "xy_phys"),
    ("physical_scale_images/CatB_MT_com/montages",
     "CatB + MT, centrosome slice", "xy_phys"),
    ("physical_scale_images/CatB_MT_syn/montages",
     "CatB + MT, synapse plane", "syn_phys"),
    # --- Actin synapse mask (no embedded scalebar; layout-pin only) ---
    ("actin/bottom_slice_seg/montages",
     "Actin synapse mask (bottom slice)", "synapse"),
    # --- XZ MIPs (all share the physical-scale 104 px / 5 μm bar) ---
    ("physical_scale_images/actin_nuc_xz/montages",
     "Actin + Nuc XZ MIP", "xz_phys"),
    ("physical_scale_images/actin_nuc_xz_planes/montages",
     "Actin + Nuc XZ MIP, cell top/bottom marked", "xz_phys"),
    ("physical_scale_images/MT_nuc_xz/montages",
     "MT + Nuc XZ MIP", "xz_phys"),
    ("physical_scale_images/CatB_nuc_xz/montages",
     "CatB + Nuc XZ MIP", "xz_phys"),
    ("physical_scale_images/CatB_MT_nuc_xz/montages",
     "CatB + MT + Nuc XZ MIP", "xz_phys"),
    ("physical_scale_images/actin_MT_xz/montages",
     "Actin + MT XZ MIP", "xz_phys"),
    ("physical_scale_images/actin_xz/montages",
     "Actin XZ MIP", "xz_phys"),
    ("physical_scale_images/CatB_xz_nolines/montages",
     "CatB XZ MIP", "xz_phys"),
    ("physical_scale_images/MT_xz_nolines/montages",
     "MT XZ MIP", "xz_phys"),
    # --- Wide multi-panel views (stacked full-width: CAT over FMC) ---
    ("physical_scale_images/CatB_MT_xz_panel_nolines/montages",
     "CatB / MT / merge, XZ MIP panel", "xzpanel_phys"),
    ("physical_scale_images/actin_MT_xz_panel_nolines/montages",
     "Actin / MT / merge, XZ MIP panel", "xzpanel_phys"),
    ("physical_scale_images/CatB_MT_com_adaptive_panels/montages",
     "CatB + MT, centrosome slice — channels + merge (adaptive)", "companel_phys"),
]

CHUNK_GLOB = "montage_cells_*.png"

# Scalebar invariant for physical_scale_images/ montages: every white bar is
# 104 px wide (5 μm bar at 20.8 px/μm), measured empirically — same as the CTL
# deck. Groups that are not physical_scale_images/ (the actin synapse mask) have
# no embedded bar and are reported as layout-pin only.
SCALEBAR_PX = 104
SCALEBAR_UM = 5
PHYS_GROUPS = {"broad", "broad_1c", "xy_phys", "syn_phys",
               "xz_phys", "xzpanel_phys", "companel_phys"}

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

GRID_LEFT = 0.10
GRID_TOP = 0.60
CELL_W = 6.50
CELL_H = SLIDE_H - GRID_TOP - 0.10
LABEL_H = 0.30
IMG_H = CELL_H - LABEL_H
LABEL_FONT_PT = 16
COL_GAP = SLIDE_W - 2 * GRID_LEFT - 2 * CELL_W

CELL_POSITIONS = [
    (GRID_LEFT,                    GRID_TOP),
    (GRID_LEFT + CELL_W + COL_GAP, GRID_TOP),
]

# Stacked full-width layout for the wide multi-panel combos: CAT on top, FMC on
# bottom, each spanning the whole slide width.
PANEL_GROUPS = {"xzpanel_phys", "companel_phys"}
PANEL_IMG_W = SLIDE_W - 2 * GRID_LEFT
PANEL_ROW_H = (SLIDE_H - GRID_TOP - 0.10) / 2
PANEL_ROW_IMG_H = PANEL_ROW_H - LABEL_H
PANEL_ROW_TOPS = [GRID_TOP, GRID_TOP + PANEL_ROW_H]

_COND_RE = re.compile(r"(CAT|FMC63|FMC(?!63))_?(\d+)(?:min)?", re.IGNORECASE)


# ---------------------------------------------------------------------------
# Long-path-safe filesystem helpers (paths here exceed Windows MAX_PATH)
# ---------------------------------------------------------------------------
def _winlong(p) -> str:
    s = os.path.abspath(str(p))
    if os.name == "nt" and not s.startswith("\\\\?\\"):
        s = "\\\\?\\" + s.replace("/", "\\")
    return s


def _exists_long(p) -> bool:
    return os.path.exists(_winlong(p))


def _png_dims(path: Path) -> Tuple[int, int]:
    with Image.open(_winlong(path)) as im:
        return im.size


def parse_condition(base_dir: str) -> Tuple[Optional[str], Optional[int]]:
    m = _COND_RE.search(base_dir)
    if not m:
        return (None, None)
    cell, mins = m.group(1).upper(), int(m.group(2))
    if cell == "FMC63":
        cell = "FMC"
    return (cell, mins)


def _parse_chunk_range(p: Path) -> Tuple[int, int]:
    m4 = re.match(r"montage_cells_(\d+)_(\d+)_(\d+)_(\d+)\.png$", p.name)
    if m4:
        f_a, c_a, f_b, c_b = (int(x) for x in m4.groups())
        return (f_a * 1000 + c_a, f_b * 1000 + c_b)
    m2 = re.match(r"montage_cells_(\d+)_(\d+)\.png$", p.name)
    if m2:
        return (int(m2.group(1)), int(m2.group(2)))
    return (0, 0)


def _list_chunk_files(montages_dir: Path) -> List[Path]:
    """Long-path-safe list of montage chunk PNGs; [] if dir absent. pathlib
    .is_dir()/.glob() silently fail past MAX_PATH even when the dir exists, so
    enumerate via os.listdir over the \\?\-prefixed path."""
    import fnmatch
    long_dir = _winlong(montages_dir)
    if not os.path.isdir(long_dir):
        return []
    try:
        names = os.listdir(long_dir)
    except OSError:
        return []
    return [montages_dir / n for n in names if fnmatch.fnmatch(n, CHUNK_GLOB)]


def find_first_chunk(montages_dir: Optional[Path]) -> Optional[Path]:
    """First unshadowed chunk (one whose cell range is not fully contained in a
    larger chunk's range), lowest start index — drops partial "smoke" chunks."""
    if montages_dir is None:
        return None
    chunks = _list_chunk_files(montages_dir)
    if not chunks:
        return None
    parsed = [(p, *_parse_chunk_range(p)) for p in chunks]
    keep = []
    for (p, s, e) in parsed:
        shadowed = any(
            s2 <= s and e <= e2 and (s2 < s or e < e2)
            for (p2, s2, e2) in parsed if p2 is not p
        )
        if not shadowed:
            keep.append((p, s, e))
    if not keep:
        return None
    keep.sort(key=lambda x: x[1])
    return keep[0][0]


def resolve_montages_dir(base_dir: str, subpath: str) -> Path:
    return Path(base_dir) / PROGRESS_FOLDER / subpath


# ---------------------------------------------------------------------------
# Slide building
# ---------------------------------------------------------------------------
def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def compute_group_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_image_at_ppi(slide, image_path: Path, ppi: float,
                     area_left: float, area_top: float,
                     area_w: float, area_h: float):
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    left_in = area_left + (area_w - w_in) / 2
    top_in = area_top + (area_h - h_in) / 2
    return slide.shapes.add_picture(
        _winlong(image_path), Inches(left_in), Inches(top_in), width=Inches(w_in))


def build_compare_slide(prs, title_text, left_label, left_img,
                        right_label, right_img, slide_ppi):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BLACK)
    add_textbox(slide, title_text, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH,
                TITLE_HEIGHT, font_pt=TITLE_FONT_PT, color=WHITE, bold=True)

    cells = [
        (left_label,  left_img,  CELL_POSITIONS[0]),
        (right_label, right_img, CELL_POSITIONS[1]),
    ]
    missing = []
    for label, img_path, (cell_left, cell_top) in cells:
        add_textbox(slide, label, cell_left, cell_top, CELL_W, LABEL_H,
                    font_pt=LABEL_FONT_PT, color=WHITE, bold=True)
        if img_path is not None and _exists_long(img_path):
            add_image_at_ppi(slide, img_path, slide_ppi,
                             cell_left, cell_top + LABEL_H, CELL_W, IMG_H)
        else:
            add_textbox(slide, "(missing)",
                        cell_left, cell_top + LABEL_H + IMG_H / 2 - 0.15, CELL_W, 0.3,
                        font_pt=14, color=WHITE)
            missing.append(label)
    return slide, missing


def build_stacked_slide(prs, title_text, top_label, top_img,
                        bottom_label, bottom_img, slide_ppi):
    """Full-width stacked layout for the wide multi-panel combos: top_img above
    bottom_img, each spanning the whole slide width so they render large."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BLACK)
    add_textbox(slide, title_text, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH,
                TITLE_HEIGHT, font_pt=TITLE_FONT_PT, color=WHITE, bold=True)

    rows = [
        (top_label,    top_img,    PANEL_ROW_TOPS[0]),
        (bottom_label, bottom_img, PANEL_ROW_TOPS[1]),
    ]
    missing = []
    for label, img_path, row_top in rows:
        add_textbox(slide, label, GRID_LEFT, row_top, PANEL_IMG_W, LABEL_H,
                    font_pt=LABEL_FONT_PT, color=WHITE, bold=True)
        if img_path is not None and _exists_long(img_path):
            add_image_at_ppi(slide, img_path, slide_ppi,
                             GRID_LEFT, row_top + LABEL_H, PANEL_IMG_W, PANEL_ROW_IMG_H)
        else:
            add_textbox(slide, "(missing)",
                        GRID_LEFT, row_top + LABEL_H + PANEL_ROW_IMG_H / 2 - 0.15,
                        PANEL_IMG_W, 0.3, font_pt=14, color=WHITE)
            missing.append(label)
    return slide, missing


# ---------------------------------------------------------------------------
def main() -> None:
    list_only = "--list" in sys.argv

    manifest_path = Path(MANIFEST_CSV)
    if not manifest_path.exists():
        print(f"ERROR: manifest not found at {manifest_path}")
        sys.exit(1)

    with manifest_path.open("r", newline="") as f:
        rows = list(csv.DictReader(f))

    # Parse the CatB rows into (date, cell, tp, base_dir).
    parsed = []
    skipped_marker = 0
    for idx, row in enumerate(rows):
        if row["marker"].strip() not in MARKER_FILTER:
            skipped_marker += 1
            continue
        base_dir = row["base_dir"].strip().rstrip("\\/")
        cell, tp = parse_condition(base_dir)
        if cell is None or tp is None:
            print(f"WARNING: could not parse condition from manifest row {idx}; skipping.")
            continue
        parsed.append({"date": row["date"].strip(), "cell": cell, "tp": tp,
                       "base_dir": base_dir})

    # Group by (date, tp); each group holds the CAT and FMC base_dirs.
    groups: Dict[Tuple[str, int], Dict[str, Optional[str]]] = {}
    for p in parsed:
        g = groups.setdefault((p["date"], p["tp"]), {"CAT": None, "FMC": None})
        g[p["cell"]] = p["base_dir"]
    sorted_group_keys = sorted(groups.keys())  # chronological: (date, tp)

    # Build slide specs: for each combo (kind), one slide per group.
    slide_specs = []  # (title, scale_group, cat_img, fmc_img, cat_dir, fmc_dir, log_key)
    for subpath, desc, scale_group in COMBOS:
        combo_has_any = False
        combo_specs = []
        for (date, tp) in sorted_group_keys:
            g = groups[(date, tp)]
            cat_dir = resolve_montages_dir(g["CAT"], subpath) if g["CAT"] else None
            fmc_dir = resolve_montages_dir(g["FMC"], subpath) if g["FMC"] else None
            cat_img = find_first_chunk(cat_dir)
            fmc_img = find_first_chunk(fmc_dir)
            if cat_img is not None or fmc_img is not None:
                combo_has_any = True
            title = f"{desc} — {date}, {tp} min"
            log_key = f"{desc}/{date}/{tp}min"
            combo_specs.append((title, scale_group, cat_img, fmc_img,
                                cat_dir, fmc_dir, log_key))
        if not combo_has_any:
            print(f"WARNING: skipping combo '{desc}' — no montages found in any group.")
            continue
        slide_specs.extend(combo_specs)

    if not slide_specs:
        print("ERROR: no slides to render — every combo was empty.")
        sys.exit(1)

    # One PPI per scale_group, pinned to the largest source image in that group,
    # so every physical-scale slide of a group shares one scalebar size.
    group_ppi: Dict[str, float] = {}
    for (_, sg, cat_img, fmc_img, _, _, _) in slide_specs:
        imgs = [p for p in (cat_img, fmc_img) if p is not None and _exists_long(p)]
        if not imgs:
            continue
        if sg in PANEL_GROUPS:
            own = compute_group_ppi(imgs, PANEL_IMG_W, PANEL_ROW_IMG_H)
        else:
            own = compute_group_ppi(imgs, CELL_W, IMG_H)
        group_ppi[sg] = max(group_ppi.get(sg, 0.0), own)

    n_groups = len(sorted_group_keys)
    n_present = sum(1 for (_, _, c, m, _, _, _) in slide_specs
                    for p in (c, m) if p is not None and _exists_long(p))
    print(f"Manifest rows: {len(rows)}  |  filtered out (marker not in "
          f"{sorted(MARKER_FILTER)}): {skipped_marker}")
    print(f"Groups (date, timepoint): {n_groups}  |  slides: {len(slide_specs)}  "
          f"|  present cells: {n_present}/{2 * len(slide_specs)}")
    print("Pinned PPI per scale_group:")
    for sg, ppi in sorted(group_ppi.items()):
        if sg in PHYS_GROUPS:
            bar_in = SCALEBAR_PX / ppi
            note = f"scalebar = {bar_in:.3f} in = {bar_in * 2.54:.3f} cm ({SCALEBAR_UM} μm)"
        else:
            note = "(layout-pin only — no embedded scalebar)"
        print(f"  {sg:>13s}  PPI={ppi:>8.2f}   {note}")
    print(f"\nOutput: {OUTPUT_PATH}\n")

    if list_only:
        for (title, sg, cat_img, fmc_img, _, _, log_key) in slide_specs:
            c = "OK" if cat_img and _exists_long(cat_img) else "MISS"
            m = "OK" if fmc_img and _exists_long(fmc_img) else "MISS"
            print(f"  [CAT:{c:<4s} FMC:{m:<4s}] ({sg}) {log_key}")
        return

    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    missing_total = []
    for (title, sg, cat_img, fmc_img, cat_dir, fmc_dir, log_key) in slide_specs:
        slide_ppi = group_ppi.get(sg, 100.0)
        builder = build_stacked_slide if sg in PANEL_GROUPS else build_compare_slide
        _, missing = builder(prs, title, "CAT", cat_img, "FMC", fmc_img, slide_ppi)
        c = "OK" if cat_img and _exists_long(cat_img) else "MISSING"
        m = "OK" if fmc_img and _exists_long(fmc_img) else "MISSING"
        print(f"[{sg:>13s}]  CAT:{c:<7s} FMC:{m:<7s}  {log_key}")
        for cell in missing:
            src = cat_dir if cell == "CAT" else fmc_dir
            missing_total.append(f"{log_key}/{cell}  ({src})")

    if out_path.exists():
        backup_dir = out_path.parent / "backups"
        created = backup_presentation(str(out_path), backup_base=str(backup_dir))
        if created:
            print(f"\nBacked up previous deck to: {backup_dir}")

    prs.save(str(out_path))
    print(f"\nDone. {len(slide_specs)} slides written to:\n  {out_path}")
    if missing_total:
        print(f"\nMissing ({len(missing_total)}):")
        for mrow in missing_total:
            print(f"  - {mrow}")
    else:
        print("\nAll cells found — no missing items.")


if __name__ == "__main__":
    main()
