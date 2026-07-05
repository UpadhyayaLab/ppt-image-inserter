"""
insert_vimkd_all_summary_slides.py

Metric summary deck for the fixed-Jurkat vimentin-knockdown (VimKD) *all-datasets*
compilation, siCtrl (control) vs siVim (vimentin KD):
    M:/.../VimentinKD_NucleusData_Fixed/results_compilation/
        VimKD_all_siCtrl_vs_siVim_by_day_violins_20260703
Each `<metric>_by_day.png` is a single-axis siCtrl-vs-siVim violin grouped by day.

Companion to insert_vimkd_by_cent_summary_slides.py. That "by_cent" deck restricts
to the 7 datasets that have a centrosome (pericentrin / MT-derived) marker and
shows centrosome-referenced metrics. This "all" deck spans ALL 10 experiments
(chronological: Apr6 ×2 / May4 / May20 / May23 / Jun28 2022, Aug4 2023, Oct24 2023,
Jan17 / Jan29 2024 — the two Apr6 entries are the Vim + Pericentrin runs on that
date), but consequently CANNOT compute any centrosome-referenced metric. So this
deck carries only the centrosome-INDEPENDENT subset (11 metrics): cell/nuclear
spreading, full-nucleus deformation & deepest-invagination, deepest-invag
orientation, and nuclear morphology. Everything `*_by_cent` / `*_around_cent` /
`cent_*` / `nuc_cent_*` / `*_cent_global_ratio` is absent here by construction.

Modeled on insert_vimkd_by_cent_summary_slides.py (blank deck, family dividers,
--list dry-run, backup-before-overwrite); same family order/titles minus the
centrosome families, for cross-deck consistency with the by_cent deck and
insert_ctl_granule_nuc_summary_20260617_slides.py.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_vimkd_all_summary_slides.py
    conda run -n PPT_editing python examples_and_configs/insert_vimkd_all_summary_slides.py --list
"""

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation  # noqa: E402

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(
    "M:/FF/FF_4TB_2_Backup_fullHD/Vimentin_Project_2ndharddrive/"
    "VimentinKD_NucleusData_Fixed/results_compilation/"
    "VimKD_all_siCtrl_vs_siVim_by_day_violins_20260703"
)
# By-day violin panels, flat at the compile root; one PNG per metric, a single-axis
# siCtrl-vs-siVim violin grouped by day, days in CHRONOLOGICAL order (date-only
# x-labels). Written by by_day_violin_panels('Jurkats_VimKD_all_compilation')-style
# rerun on 2026-07-03. Centrosome-referenced panels do not exist in this compile.
GRID_DIR = ROOT
GRID_SUFFIX = "_by_day.png"

OUTPUT_PATH = Path(
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/VimentinKD/"
    "VimKD_Jurkats_siCtrl_vs_siVim_all_summary.pptx"
)

DECK_TITLE = "Vimentin knockdown effects on fixed Jurkat nuclei (all experiments)"
DECK_SUBTITLE = (
    "siCtrl (control) vs siVim (vimentin KD)  ·  fixed Jurkat  ·  all 10 experiments  ·  "
    "centrosome-independent metrics only (no centrosome marker across all runs)  ·  "
    "compiled 2026-07-03"
)

# ---------------------------------------------------------------------------
# Curated metrics, grouped into families (divider slide per family).
# The centrosome-independent subset of the by_cent deck's set — the centrosome
# families (Centrosome ↔ nucleus, Actin around centrosome) and the *_by_cent
# invagination/orientation entries are dropped because they don't exist in the
# all-datasets compile. Family order/titles otherwise match the by_cent deck.
# ---------------------------------------------------------------------------
FAMILIES = [
    ("Cell and nuclear spreading", [
        ("nuc_aspect_ratio",       "Nuclear aspect ratio"),
        ("actin_deform_ratio",     "Cell aspect ratio"),
        ("actin_bottom_mask_area", "Synapse area"),
    ]),
    ("Nuclear deformation and invaginations", [
        ("chull_max_D",                       "Max invag depth over full nucleus"),
        ("deepest_invag_volume",              "Deepest invagination volume"),
        ("deepest_invag_fraction_chull_volume", "Deepest invag: frac of convex hull volume"),
        ("deepest_region_periph_ratio_025um", "DNA levels near invag"),
    ]),
    ("Invagination orientation", [
        ("avg_normal_angle_adaptive_region_growth", "Deepest invag orientation"),
    ]),
    ("Nuclear morphology", [
        ("nuc_solidity",        "Nuclear solidity"),
        ("nuc_mesh_sphericity", "Nuclear sphericity"),
        ("nuc_volume_mesh",     "Nuclear volume"),
        ("nuc_SA_mesh",         "Nuclear surface area"),
    ]),
]

# ---------------------------------------------------------------------------
# Colors / layout (matches the by_cent / bleb / washout summary decks)
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.10

TITLE_LEFT = MARGIN
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.55
TITLE_FONT_PT = 28

IMG_LEFT = MARGIN
IMG_TOP = 0.66
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = SLIDE_H - IMG_TOP - 0.12

# ---------------------------------------------------------------------------


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    """Place an image inside (left, top, w, h), preserving aspect ratio and
    centering on whichever dimension ends up smaller than the box."""
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.6, SLIDE_W - 2 * MARGIN, 1.4,
                font_pt=38, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.1, SLIDE_W - 2 * MARGIN, 1.2,
                font_pt=16, color=BLACK, italic=True)


def build_divider_slide(prs, family_name):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, family_name, MARGIN, 3.1, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=44, color=BLACK, bold=True)


def build_slide(prs, title_text, image_path):
    slide = _new_slide(prs)
    add_textbox(slide, title_text, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title_text), color=BLACK, bold=True)
    missing = not image_path.exists()
    if not missing:
        add_image_in_box(slide, str(image_path), IMG_LEFT, IMG_TOP, IMG_BOX_W, IMG_BOX_H)
    else:
        add_textbox(slide, "(missing)", IMG_LEFT, IMG_TOP + IMG_BOX_H / 2 - 0.2,
                    IMG_BOX_W, 0.4, font_pt=18, color=BLACK)
    return missing


def main():
    list_only = "--list" in sys.argv
    n_metrics = sum(len(items) for _, items in FAMILIES)
    est_slides = 1 + sum(1 + len(items) for _, items in FAMILIES)

    print("Source: {}".format(GRID_DIR))
    print("{} curated metrics across {} families, est. {} slides\n".format(
        n_metrics, len(FAMILIES), est_slides))

    if list_only:
        for fam, items in FAMILIES:
            print("=== {} ({}) ===".format(fam, len(items)))
            for stem, title in items:
                exists = (GRID_DIR / (stem + GRID_SUFFIX)).exists()
                print("  [{}] {:<52s} {}".format(
                    "OK " if exists else "MISS", stem + GRID_SUFFIX, title))
            print("")
        return

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)

    missing = []
    for fam, items in FAMILIES:
        build_divider_slide(prs, fam)
        print("=== {} ===".format(fam))
        for stem, title in items:
            p = GRID_DIR / (stem + GRID_SUFFIX)
            is_missing = build_slide(prs, title, p)
            print("  [{}] {} -> {!r}".format("OK" if not is_missing else "MISSING", stem + GRID_SUFFIX, title))
            if is_missing:
                missing.append(stem + GRID_SUFFIX)
        print("")

    if OUTPUT_PATH.exists():
        backup_dir = OUTPUT_PATH.parent / "backups"
        created = backup_presentation(str(OUTPUT_PATH), backup_base=str(backup_dir))
        if created:
            print("Backed up previous deck to: {}\n".format(backup_dir))

    prs.save(str(OUTPUT_PATH))
    total = len(prs.slides._sldIdLst)
    print("Done. {} metrics, {} slides written to:\n  {}".format(
        n_metrics, total, OUTPUT_PATH))
    if missing:
        print("\nSkipped {} missing panel(s):".format(len(missing)))
        for m in missing:
            print("  - {}".format(m))
    else:
        print("\nAll curated panels found - no missing items.")


if __name__ == "__main__":
    main()
