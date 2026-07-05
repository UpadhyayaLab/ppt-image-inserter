"""
insert_pMLC_actin_synapse_xz_mip_montages_slides.py

CAT (left) vs FMC (right) side-by-side deck for the 3 Kiet pMLC confocal
datasets (20260312, 20260510, 20260607). Grouped by kind — all 9 slides
of one kind consecutively, then next kind:

    Slides  1- 9: Actin at Synapse                    (segmented)
    Slides 10-18: pMLC at Synapse                     (fixed contrast, solo)
    Slides 19-27: pMLC + Actin Composite at Synapse   (composite, fixed contrast)
    Slides 28-36: Actin XZ MIP
    Slides 37-45: pMLC XZ MIP                         (fixed contrast)

Within each kind block, slides are chronological by (date, tp).

Filter: `marker == "pMLC"` rows only. 3 datasets x 3 timepoints x 5 kinds
= 45 slides.

Usage:
    python examples_and_configs/insert_pMLC_actin_synapse_xz_mip_montages_slides.py
"""

import csv
import os
import re
import sys
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

MANIFEST_CSV = (
    "L:/FF/CAR T/actin_compiled_results/all_datasets_actin/"
    "compiled_20260623/dataset_manifest.csv"
)

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/CART/pMLC/"
    "CART_pMLC_actin_synapse_xz_mip_montages_20260625.pptx"
)

# Only include manifest rows whose marker matches this set.
MARKER_FILTER = {"pMLC"}

# (kind label, subpath under <base>/<progress_folder>/).
# Interleaved per (date, tp) group: synapse views first, XZ MIP second.
KINDS = [
    ("Actin at Synapse",                    "actin/synapse/1slice/mask/montages"),
    ("pMLC at Synapse (fixed contrast)",    "pMLC/synapse/1slice/fixed_con/montages"),
    ("pMLC + Actin Composite at Synapse",   "pMLC/synapse/composite/1slice/montages"),
    ("Actin XZ MIP",                        "actin/xz_mip/montages"),
    ("pMLC XZ MIP (fixed contrast)",        "pMLC/xz_mip/fixed_con/montages"),
]

CHUNK_GLOB = "montage_cells_*.png"

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.05
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.05
TITLE_HEIGHT = 0.40
TITLE_FONT_PT = 24

GRID_LEFT = 0.05
GRID_TOP = 0.50
CELL_W = 6.60
CELL_H = SLIDE_H - GRID_TOP - 0.05
LABEL_H = 0.22
IMG_H = CELL_H - LABEL_H
LABEL_FONT_PT = 14
COL_GAP = SLIDE_W - 2 * GRID_LEFT - 2 * CELL_W

CELL_POSITIONS = [
    (GRID_LEFT,                    GRID_TOP),
    (GRID_LEFT + CELL_W + COL_GAP, GRID_TOP),
]

PPUM_SOURCE = 30
SCALEBAR_UM = 5
SCALEBAR_PX = PPUM_SOURCE * SCALEBAR_UM

_COND_RE = re.compile(r"(CAT|FMC63|FMC(?!63))_?(\d+)(?:min)?", re.IGNORECASE)
_DTAG_RE = re.compile(r"_D(\d+)_", re.IGNORECASE)


def _long_path_str(p: Path) -> str:
    s = str(p)
    if os.name == "nt" and len(s) >= 240 and not s.startswith("\\\\?\\"):
        s = s.replace("/", "\\")
        if not s.startswith("\\\\"):
            return "\\\\?\\" + s
    return s


def parse_condition(base_dir: str) -> Tuple[Optional[str], Optional[int]]:
    m = _COND_RE.search(base_dir)
    if not m:
        return (None, None)
    cell, mins = m.group(1).upper(), int(m.group(2))
    if cell == "FMC63":
        cell = "FMC"
    return (cell, mins)


def parse_dtag(base_dir: str) -> str:
    m = _DTAG_RE.search(base_dir)
    return f" D{m.group(1)}" if m else ""


def resolve_kind_dir(base_dir: str, progress_folder: str, kind_subpath: str) -> Optional[Path]:
    cand = Path(base_dir) / progress_folder / kind_subpath
    return cand if os.path.isdir(_long_path_str(cand)) else None


def _parse_chunk_range(p: Path) -> Tuple[int, int]:
    m4 = re.match(r"montage_cells_(\d+)_(\d+)_(\d+)_(\d+)\.png$", p.name)
    if m4:
        f_a, c_a, f_b, c_b = (int(x) for x in m4.groups())
        return (f_a * 1000 + c_a, f_b * 1000 + c_b)
    m2 = re.match(r"montage_cells_(\d+)_(\d+)\.png$", p.name)
    if m2:
        return (int(m2.group(1)), int(m2.group(2)))
    return (0, 0)


def _list_chunk_files(montages_dir):
    """Long-path-safe list of montage chunk PNGs (unsorted); [] if dir absent.
    pathlib .is_dir()/.glob() silently return False/empty past Windows MAX_PATH
    (260) even when the dir exists, so enumerate via os.listdir over the
    \\?\-prefixed path."""
    import fnmatch
    long_dir = _long_path_str(montages_dir)
    if not os.path.isdir(long_dir):
        return []
    return [montages_dir / n for n in os.listdir(long_dir)
            if fnmatch.fnmatch(n, CHUNK_GLOB)]


def find_first_chunk(montages_dir: Optional[Path]) -> Optional[Path]:
    chunks = _list_chunk_files(montages_dir)
    if not chunks:
        return None
    parsed = [(p, *_parse_chunk_range(p)) for p in chunks]
    keep = []
    for (p, s, e) in parsed:
        shadowed = any(
            s2 <= s and e <= e2 and (s2 < s or e < e2)
            for (p2, s2, e2) in parsed
            if p2 is not p
        )
        if not shadowed:
            keep.append((p, s, e))
    if not keep:
        return None
    keep.sort(key=lambda x: x[1])
    return keep[0][0]


def _png_dims(path: Path) -> Tuple[int, int]:
    with Image.open(_long_path_str(path)) as im:
        return im.size


def _path_exists(p: Path) -> bool:
    return os.path.exists(_long_path_str(p))


def compute_deck_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_image_in_cell_at_ppi(slide, image_path: Path, ppi: float,
                             cell_left: float, cell_top: float):
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    img_area_top = cell_top + LABEL_H
    left_in = cell_left + (CELL_W - w_in) / 2
    top_in  = img_area_top + (IMG_H - h_in) / 2
    return slide.shapes.add_picture(
        _long_path_str(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(w_in),
    )


def build_compare_slide(prs, title_text: str,
                        cat_img: Optional[Path], fmc_img: Optional[Path],
                        deck_ppi: float):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    cells = [
        ("CAT", cat_img, CELL_POSITIONS[0]),
        ("FMC", fmc_img, CELL_POSITIONS[1]),
    ]
    missing = []
    for label, img_path, (cell_left, cell_top) in cells:
        add_textbox(
            slide, label,
            cell_left, cell_top, CELL_W, LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )
        if img_path is not None and _path_exists(img_path):
            add_image_in_cell_at_ppi(slide, img_path, deck_ppi, cell_left, cell_top)
        else:
            add_textbox(
                slide, "(missing)",
                cell_left, cell_top + LABEL_H + IMG_H / 2 - 0.15, CELL_W, 0.3,
                font_pt=14, color=WHITE,
            )
            missing.append(label)
    return slide, missing


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    manifest_path = Path(MANIFEST_CSV)
    if not manifest_path.exists():
        print(f"ERROR: manifest not found at {manifest_path}")
        return

    with manifest_path.open("r", newline="") as f:
        rows = list(csv.DictReader(f))

    skipped_marker = 0
    parsed_rows = []
    for idx, row in enumerate(rows):
        marker = row["marker"].strip()
        if marker not in MARKER_FILTER:
            skipped_marker += 1
            continue

        date_tag = row["date"].strip()
        progress_folder = row["progress_folder"].strip()
        base_dir = row["base_dir"].strip().rstrip("\\/")
        modality = "TIRF" if row["single_plane"].strip() == "1" else "confocal"

        cell, tp = parse_condition(base_dir)
        dtag = parse_dtag(base_dir)
        marker_label = f"{marker}{dtag}"

        per_kind: Dict[str, Tuple[Optional[Path], Optional[Path]]] = {}
        for kind_label, kind_subpath in KINDS:
            mdir = resolve_kind_dir(base_dir, progress_folder, kind_subpath)
            chunk = find_first_chunk(mdir)
            per_kind[kind_label] = (mdir, chunk)

        parsed_rows.append({
            "idx": idx,
            "date_tag": date_tag,
            "marker_label": marker_label,
            "modality": modality,
            "cell": cell,
            "tp": tp,
            "per_kind": per_kind,
        })

    # Group by (date, marker_label, tp).
    groups: Dict[Tuple[str, str, int], Dict] = {}
    for p in parsed_rows:
        if p["cell"] is None or p["tp"] is None:
            print(f"WARNING: could not parse condition from manifest row {p['idx']}; skipping.")
            continue
        key = (p["date_tag"], p["marker_label"], p["tp"])
        g = groups.setdefault(key, {
            "per_kind": {
                kind_label: {
                    "cat_chunk": None, "fmc_chunk": None,
                    "cat_dir": None, "fmc_dir": None,
                }
                for kind_label, _ in KINDS
            },
            "first_idx": p["idx"],
            "modality": p["modality"],
        })
        for kind_label, _ in KINDS:
            mdir, chunk = p["per_kind"][kind_label]
            slot = g["per_kind"][kind_label]
            if p["cell"] == "CAT":
                slot["cat_chunk"] = chunk
                slot["cat_dir"] = mdir
            elif p["cell"] == "FMC":
                slot["fmc_chunk"] = chunk
                slot["fmc_dir"] = mdir

    sorted_groups = sorted(
        groups.items(),
        key=lambda kv: (kv[0][0], kv[1]["first_idx"]),
    )

    # Group by kind first, then by date within each kind block.
    slide_specs = []
    for kind_label, _ in KINDS:
        for ((date_tag, marker_label, tp), g) in sorted_groups:
            slot = g["per_kind"][kind_label]
            title = (
                f"{kind_label} — {date_tag} ({marker_label}, "
                f"{g['modality']}): {tp} min"
            )
            log_key = f"{kind_label}/{date_tag}/{marker_label}/{tp}min"
            slide_specs.append((
                title,
                slot["cat_chunk"], slot["fmc_chunk"],
                slot["cat_dir"], slot["fmc_dir"],
                log_key,
            ))

    present_count = sum(
        1 for (_, cat_img, fmc_img, _, _, _) in slide_specs
        for p in (cat_img, fmc_img)
        if p is not None and _path_exists(p)
    )
    n_groups = len(sorted_groups)
    print(
        f"Manifest rows: {len(rows)}  -  filtered out (marker not in "
        f"{sorted(MARKER_FILTER)}): {skipped_marker}\n"
        f"Groups: {n_groups}\n"
        f"Slides emitted: {len(slide_specs)} ({len(KINDS)} kinds x {n_groups} groups)\n"
        f"Present cells: {present_count} / {2 * len(slide_specs)}\n"
        f"Per-slide PPI: each slide pins its own PPI to the larger of its CAT/FMC\n"
        f"  montages so cells fit exactly with no overflow.\n"
        f"Source PPUM = {PPUM_SOURCE} px/μm (locked).\n"
    )
    print(f"Writing deck to: {OUTPUT_PATH}\n")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    missing_total = []
    slides_added = 0
    for (title, cat_img, fmc_img, cat_dir, fmc_dir, log_key) in slide_specs:
        present_on_slide = [
            p for p in (cat_img, fmc_img)
            if p is not None and _path_exists(p)
        ]
        slide_ppi = compute_deck_ppi(present_on_slide, CELL_W, IMG_H) if present_on_slide else 100.0
        bar_cm = (SCALEBAR_PX / slide_ppi) * 2.54

        _, missing = build_compare_slide(prs, title, cat_img, fmc_img, slide_ppi)
        slides_added += 1

        status_parts = [
            "CAT:OK" if cat_img and _path_exists(cat_img) else "CAT:MISSING",
            "FMC:OK" if fmc_img and _path_exists(fmc_img) else "FMC:MISSING",
        ]
        print(
            f"[{log_key}]  PPI={slide_ppi:7.2f}  bar={bar_cm:.3f}cm  "
            + "  ".join(status_parts)
        )

        for cell in missing:
            src = cat_dir if cell == "CAT" else fmc_dir
            missing_total.append(f"{log_key}/{cell}  ({src})")

    prs.save(str(out_path))
    print(f"\nDone. {slides_added} slides written to:\n  {out_path}")

    if missing_total:
        print(f"\nMissing ({len(missing_total)}):")
        for m in missing_total:
            print(f"  - {m}")
    else:
        print("\nAll images found - no missing items.")


if __name__ == "__main__":
    main()
