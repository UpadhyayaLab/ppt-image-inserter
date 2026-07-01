"""
insert_vimkd_vim_nuc_xz_slides.py

siCtrl (left) vs siVim (right) vimentin + nucleus X-Z MIP montage deck for the
fixed-Jurkat VimentinKD experiments. Each slide is one acquisition date; the panel
is the vimentin (channel) + nucleus (DNA) maximum-intensity projection viewed along
the x-z axis, showing how the vimentin cage sits around the nucleus through the cell
depth and how that changes when vimentin is knocked down.

ONLY TWO DATASETS ARE COVERED. The VimentinKD config group
(TCell-3D-Morphodynamics/config/datasets/fixed/Jurkats/VimentinKD) has 10 datasets,
but in 8 of them vimentin is the knockdown *target* and the imaged structural channel
is MT / acetylated-tubulin / pMLC / pericentrin instead — those have no vimentin x-z
MIP to show. Only the two KD-validation experiments actually image a vimentin channel:
  - 04/06/2022  transfection2_48h                  Control / siVIM   (basal)
  - 05/04/2022  Vimentin knockdown validation-48h  siCtrl / siVim    (αCD3)

MONTAGE SOURCE — note the path. These montages are NOT in the physical_scale_images/
tree the sibling decks (noco/bleb/LatA, the vimkd MT/pericentrin decks) read. The
vimentin x-z MIP is emitted under the vimentin output subtree instead:
  <root>/<cond>/cells/individual-channels/prog_fixed_cells/vimentin/XZ_nuc_vim/montages/
These are fixed-canvas montages (every chunk is 1335x1000), not the physical-scale
rendering — so there is no embedded-scalebar invariant to preserve here. PPI is pinned
deck-wide only so the two slides' panels render at the same on-page size (which, given
the uniform canvas, is automatic); it carries no physical-scale claim. This is why the
script name drops the "phys_scale" suffix the sibling scripts carry.

Layout follows the sibling decks: one slide per dataset, siCtrl | siVim two-column,
showing the first montage chunk per condition (n_chunks=1). Long-path-safe enumeration
is load-bearing — the 05/04/2022 montage path runs past Windows MAX_PATH (260), so the
\\?\ prefix is required for os.listdir / Image.open / add_picture, same as the LatA deck.

Usage:
    python examples_and_configs/insert_vimkd_vim_nuc_xz_slides.py
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/VimentinKD/"
    "VimKD_Jurkats_siCtrl_vs_siVim_vim_nuc_xz_montages.pptx"
)

# Shared dataset root (both vim-imaged validation experiments live under it).
ROOT = (
    "M:/FF/FF_4TB_2_Backup_fullHD/Vimentin_Project_2ndharddrive/"
    "VimentinKD_NucleusData_Fixed"
)

# The vimentin x-z MIP combo folder under prog_fixed_cells/vimentin/.
COMBO = "XZ_nuc_vim"

# Each entry is one siCtrl-vs-siVim slide for one acquisition date.
# tag      — date tag woven into slide titles + used for chronological sort.
# root     — dataset root (the date-folder + intermediate condition-group dir).
# chan_sub — intermediate path inside each condition folder before prog_fixed_cells.
# left/right — (condition_folder, display_label) for the siCtrl / siVim columns.
# tp_label — activation string ("αCD3" or "" for the basal 04/06 dataset).
EXPERIMENTS = [
    {
        "tag": "04/06/2022",
        "root": ROOT + "/20220406 - Vimentin siRNA Experiments/transfection2_48h",
        "chan_sub": "cells/individual-channels", "tp_label": "",
        "left":  ("Control", "siCtrl"),
        "right": ("siVIM",   "siVim"),
    },
    {
        "tag": "05/04/2022",
        "root": (ROOT + "/20220504 - Vimentin siRNA Experiments/"
                 "Vimentin knockdown validation - 48h"),
        "chan_sub": "cells/individual-channels", "tp_label": "αCD3",
        "left":  ("siCtrl - aCD3", "siCtrl"),
        "right": ("siVim - aCD3",  "siVim"),
    },
]

# (combo_subfolder, title_template, n_chunks, scale_group, fallback, opts).
# title_template gets " ({sub})" appended in main(), where sub is "tp, tag" (or just
# tag when tp is empty), so the basal title reads "... (04/06/2022)" and the αCD3 title
# reads "... (αCD3, 05/04/2022)". Single combo: the vimentin + nucleus x-z MIP.
COMBOS = [
    (COMBO, "Vim + Nuc XZ MIP", 1, "xz", None, {}),
]

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Slide layout (inches). 13.333 x 7.5 widescreen.
SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

# 1x2 cell grid below the title (label + image per cell)
GRID_LEFT = 0.10
GRID_TOP = 0.60
CELL_W = 6.50
CELL_H = SLIDE_H - GRID_TOP - 0.10   # 6.80"
LABEL_H = 0.30
IMG_H = CELL_H - LABEL_H             # 6.50"
LABEL_FONT_PT = 16
COL_GAP = SLIDE_W - 2 * GRID_LEFT - 2 * CELL_W

CELL_POSITIONS = [
    (GRID_LEFT,                    GRID_TOP),
    (GRID_LEFT + CELL_W + COL_GAP, GRID_TOP),
]

# ---------------------------------------------------------------------------


def _winlong(p) -> str:
    """Return a Win32-safe absolute path string. On Windows, prepends the
    \\\\?\\ extended-length prefix so paths near MAX_PATH (260) still work in
    os.listdir / Image.open / python-pptx add_picture. The 05/04/2022 montage
    path runs past 260 chars, so this is load-bearing, not just defensive."""
    s = os.path.abspath(str(p))
    if os.name == "nt" and not s.startswith("\\\\?\\"):
        s = "\\\\?\\" + s.replace("/", "\\")
    return s


def _exists_long(p) -> bool:
    """MAX_PATH-safe existence check (pathlib.Path.exists() trips on long paths)."""
    return os.path.exists(_winlong(p))


def montage_dir(root, cond_folder: str, chan_sub: str, combo: str) -> Path:
    """Build the montages dir for one (experiment, condition, combo).

    The vimentin x-z MIP lives under the vimentin output subtree, NOT under
    physical_scale_images/ like the sibling decks' combos."""
    return (Path(root) / cond_folder / chan_sub /
            "prog_fixed_cells" / "vimentin" / combo / "montages")


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _png_dims(path: Path) -> Tuple[int, int]:
    """Return (width_px, height_px) of a PNG without fully decoding it."""
    with Image.open(_winlong(path)) as im:
        return im.size


def compute_slide_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    """Smallest ppi such that every image fits in (max_w_in x max_h_in)."""
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_image_at_ppi(slide, image_path: Path, ppi: float,
                     area_left: float, area_top: float,
                     area_w: float, area_h: float):
    """Center an image at uniform PPI inside an arbitrary (left, top, w, h)
    image area. native_px / ppi gives inches; both dims pinned by ppi so all
    images in the slide render at one on-page scale."""
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    left_in = area_left + (area_w - w_in) / 2
    top_in  = area_top  + (area_h - h_in) / 2
    return slide.shapes.add_picture(
        _winlong(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(w_in),
    )


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _chunk_start_index(p: Path) -> int:
    m = re.match(r"montage_cells_(\d+)", p.name)
    return int(m.group(1)) if m else 0


def list_chunks(montages_dir) -> List[Path]:
    """Return the montage chunk PNGs in a dir, sorted by chunk-start index.
    Long-path-safe: pathlib's is_dir()/glob() silently fail on Windows paths
    past MAX_PATH (260) — they return False / empty even when the dir exists.
    os.listdir on the \\\\?\\-prefixed path enumerates them correctly.
    Returns [] if absent (so a dataset without montages is simply skipped)."""
    d = _winlong(montages_dir)
    if not os.path.isdir(d):
        return []
    names = [f for f in os.listdir(d)
             if f.startswith("montage_cells_") and f.endswith(".png")]
    return sorted((Path(montages_dir) / n for n in names), key=_chunk_start_index)


def find_first_chunks(montages_dir: Path, n: int) -> List[Optional[Path]]:
    """Return the first n chunks sorted by chunk-start index. Pads with None
    if the folder has fewer chunks (or doesn't exist)."""
    chunks = list_chunks(montages_dir)
    return (chunks + [None] * n)[:n]


def build_compare_slide(prs, title_text,
                        left_label, left_imgs,
                        right_label, right_imgs,
                        slide_ppi):
    """Render the siCtrl | siVim comparison slide (one montage chunk per side).
    Both panels share slide_ppi so they render at the same on-page size."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    missing = []
    cells = [
        (left_label,  left_imgs[0],  CELL_POSITIONS[0][0]),
        (right_label, right_imgs[0], CELL_POSITIONS[1][0]),
    ]
    for label, img_path, cell_left in cells:
        add_textbox(
            slide, label,
            cell_left, GRID_TOP, CELL_W, LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )
        if img_path is not None and _exists_long(img_path):
            add_image_at_ppi(slide, img_path, slide_ppi,
                             cell_left, GRID_TOP + LABEL_H, CELL_W, IMG_H)
        else:
            add_textbox(
                slide, "(missing)",
                cell_left, GRID_TOP + LABEL_H + IMG_H / 2 - 0.15,
                CELL_W, 0.3,
                font_pt=14, color=WHITE,
            )
            missing.append(label)
    return slide, missing


def _exp_date_key(exp):
    """Sort key from the tag (MM/DD/YYYY). Stable sort keeps EXPERIMENTS order
    for same-key entries."""
    t = exp["tag"]
    m = re.match(r"(\d{2})/(\d{2})/(\d{4})", t)
    if m:
        return (int(m.group(3)), int(m.group(1)), int(m.group(2)))
    return (9999, 99, 99)


def _make_title(base: str, tp_label: str, tag: str) -> str:
    """'<base> (<tp>, <tag>)', dropping the tp part when it's empty (basal
    dataset) -> '<base> (<tag>)'."""
    sub = f"{tp_label}, {tag}" if tp_label else tag
    return f"{base} ({sub})"


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    slide_specs: List[dict] = []
    # Group-major order: emit the combo across experiments (date-sorted, stable).
    for combo_folder, base_title, n_chunks, scale_group, fallback, opts in COMBOS:
        for exp in sorted(EXPERIMENTS, key=_exp_date_key):
            root = Path(exp["root"])
            chan_sub = exp["chan_sub"]
            left_folder, left_label = exp["left"]
            right_folder, right_label = exp["right"]
            tag = exp["tag"]
            tp_label = exp["tp_label"]
            exp_key = f"{tag} {tp_label}".strip()

            primary_dir = montage_dir(root, left_folder, chan_sub, combo_folder)
            if not list_chunks(primary_dir) and fallback is None:
                continue  # no montages for this dataset -> skip the slide

            title = _make_title(base_title, tp_label, tag)
            left_dir  = montage_dir(root, left_folder,  chan_sub, combo_folder)
            right_dir = montage_dir(root, right_folder, chan_sub, combo_folder)
            left_imgs  = find_first_chunks(left_dir,  n_chunks)
            right_imgs = find_first_chunks(right_dir, n_chunks)
            slide_specs.append({
                "log_key": f"{exp_key}/{combo_folder}",
                "n_chunks": n_chunks,
                "scale_group": scale_group,
                "title": title,
                "left_label": left_label,   "left_imgs": left_imgs,
                "right_label": right_label, "right_imgs": right_imgs,
                "left_dir": left_dir,       "right_dir": right_dir,
            })

    # Deck-wide PPI pinning per scale_group, so all panels share one on-page scale.
    group_ppi: dict = {}
    for spec in slide_specs:
        imgs = [p for p in (*spec["left_imgs"], *spec["right_imgs"])
                if p is not None and _exists_long(p)]
        own_ppi = compute_slide_ppi(imgs, CELL_W, IMG_H) if imgs else 0.0
        sg = spec["scale_group"]
        group_ppi[sg] = max(group_ppi.get(sg, 0.0), own_ppi)

    print(f"Pinning PPI per scale_group across {len(slide_specs)} slides "
          f"(uniform on-page panel size; not a physical-scale claim).")
    for sg, ppi in sorted(group_ppi.items()):
        print(f"  {sg:>12s}  PPI={ppi:>7.2f}")
    print(f"\nWriting deck to: {OUTPUT_PATH}\n")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    missing_total = []
    slides_added = 0
    for spec in slide_specs:
        slide_ppi = group_ppi[spec["scale_group"]]
        _, missing = build_compare_slide(
            prs, spec["title"],
            spec["left_label"], spec["left_imgs"],
            spec["right_label"], spec["right_imgs"],
            slide_ppi,
        )
        slides_added += 1

        n_chunks = spec["n_chunks"]
        left_ok  = sum(1 for p in spec["left_imgs"]  if p is not None)
        right_ok = sum(1 for p in spec["right_imgs"] if p is not None)
        print(f"[{spec['log_key']}]  L:{left_ok}/{n_chunks}  R:{right_ok}/{n_chunks}")

        for cell in missing:
            src = spec["left_dir"] if cell == spec["left_label"] else spec["right_dir"]
            missing_total.append(f"{spec['log_key']}/{cell}  ({src})")

    prs.save(str(out_path))
    print(f"\nDone. {slides_added} slides written to:\n  {out_path}")

    if slides_added == 0:
        print(
            "\nNo slides written — no vimentin XZ_nuc_vim montages found. Expected at\n  "
            "<cond>/<chan_sub>/prog_fixed_cells/vimentin/XZ_nuc_vim/montages/ ."
        )
    if missing_total:
        print(f"\nMissing ({len(missing_total)}):")
        for m in missing_total:
            print(f"  - {m}")
    elif slides_added:
        print("\nAll images found - no missing items.")


if __name__ == "__main__":
    main()
