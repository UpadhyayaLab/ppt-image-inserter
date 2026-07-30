"""
insert_vim_restyled_Ctrl_vs_Noco_slides.py

Vimentin geom-density deck from the RESTYLED figures (Ctrl blue vs Noco 1 μM
orange, sig stars). 3 experiments (Apr 29 2022, Jan 23 2024, Feb 27 2024) —
main slides use the _pooled figures; an across-experiments section shows the
per-experiment breakdown. Sourced from geom_density/*.png (the curated
vim_loc_wrto_invag/ copies are .tiff for this dataset, so we use the geom_density
.png files the spec references). Same structure/style as the MT restyled deck,
including the centrosome proximal|distal 1×2 figures and the prox÷distal ratio.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_vim_restyled_Ctrl_vs_Noco_slides.py
    conda run -n PPT_editing python examples_and_configs/insert_vim_restyled_Ctrl_vs_Noco_slides.py --list
"""

import math
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

try:
    from PIL import Image
except ImportError:  # pragma: no cover
    Image = None

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation, safe_path, path_exists  # noqa: E402

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
RESTYLED = Path(
    "J:/FF/fixed_cell/Vimentin/results_across_experiments/"
    "Noco_CD3_20220429_20240123_20240227_20260717/restyled")
PROF = RESTYLED / "geom_density" / "profiles"
ENR = RESTYLED / "geom_density" / "enrichment"
BYCENT = RESTYLED / "geom_density" / "profiles" / "bycent"
NC = RESTYLED / "geom_density" / "near_cent"
GRID = RESTYLED / "grid_panels"

# per-experiment tokens (comma form used by enrichment / near_cent)
EXP_DATES = [("Apr 29 2022", "Apr_29,_2022"),
             ("Jan 23 2024", "Jan_23,_2024"),
             ("Feb 27 2024", "Feb_27,_2024")]

OUTPUT_PATH = Path(
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/"
    "nuc_mesh_struct_outside_nuc/"
    "Vimentin geom-density vs NE geometry, Ctrl vs Noco (restyled, 20260722).pptx")

DECK_TITLE = "Vimentin vs nuclear-envelope geometry — Ctrl vs Noco (restyled)"
DECK_SUBTITLE = (
    "Fixed Jurkats · Ctrl (blue) vs Noco 1 μM (orange) · 3 experiments "
    "(Apr 29 2022, Jan 23 2024, Feb 27 2024) shown SEPARATELY — no pooling · "
    "restyled figures · 2026-07-22")

FOOTER_CMP = "Ctrl (blue) vs Noco 1 μM (orange)"

# ---------------------------------------------------------------------------
# Colors / layout (16:9)
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0x55, 0x55, 0x55)
FIELD_COLOR = RGBColor(0x2E, 0x5A, 0x88)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.08
GAP = 0.10

TITLE_LEFT = MARGIN
TITLE_TOP = 0.04
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 27

IMG_TOP = 0.58
FOOTER_TOP = 7.20
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.28
FOOTER_FONT_PT = 8
MAX_GRID_COLS = 5


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------
def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER, wrap=True):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    image_path = safe_path(image_path)
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.6, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=32, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.0, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=15, color=GREY, italic=True)


def build_divider_slide(prs, title, subtitle=""):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, title, MARGIN, 3.0, SLIDE_W - 2 * MARGIN, 1.2,
                font_pt=36, color=BLACK, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, MARGIN, 4.25, SLIDE_W - 2 * MARGIN, 0.9,
                    font_pt=18, color=GREY, italic=True)


def add_labeled_caption(slide, lines, left, top, width, height):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    for i, (text, pt, color) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.size = Pt(pt)
        run.font.color.rgb = color
    return box


def add_header_box(slide, text, left, top, width, height):
    """Plain bold black experiment-date label placed above a plot."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = BLACK
    return box


def _place_headers(slide, header, left, top, width, height):
    """header may be a single label (one box over the tile) or a list of
    labels spread evenly across the tile (one box per sub-panel)."""
    if isinstance(header, (list, tuple)):
        k = len(header)
        sub_w = width / k
        for j, htext in enumerate(header):
            add_header_box(slide, htext, left + j * sub_w, top, sub_w, height)
    else:
        add_header_box(slide, header, left, top, width, height)


def _tile_caption(slide, cap, path, left, top, width, height, cap_pt):
    lines = []
    if cap:
        lines.append((cap, min(cap_pt, 11), GREY))
    lines.append((Path(str(path)).stem, 8, FIELD_COLOR))
    add_labeled_caption(slide, lines, left, top, width, height)


_AR_CACHE = {}


def _img_ar(path):
    key = str(path)
    if key in _AR_CACHE:
        return _AR_CACHE[key]
    ar = 1.5
    if Image is not None:
        try:
            with Image.open(safe_path(key)) as im:
                w, h = im.size
                if h:
                    ar = w / float(h)
        except Exception:
            pass
    _AR_CACHE[key] = ar
    return ar


def _uniform_img_size(a, cols, rows, area_w, area_h, cap_h):
    cw = (area_w - (cols - 1) * GAP) / cols
    cell_h = (area_h - (rows - 1) * GAP) / rows
    ih_avail = cell_h - cap_h
    if ih_avail <= 0.2 or cw <= 0.2:
        return None
    iw = min(cw, ih_avail * a)
    return iw, iw / a


def _best_cols(ars, n, area_w, area_h, cap_h, cap_cols):
    uniform = (max(ars) / min(ars)) < 1.20
    a = sorted(ars)[len(ars) // 2]
    best = None
    for cols in range(1, min(cap_cols, n) + 1):
        rows = math.ceil(n / cols)
        if uniform:
            sz = _uniform_img_size(a, cols, rows, area_w, area_h, cap_h)
            if sz is None:
                continue
            score = sz[0] * sz[1]
        else:
            cw = (area_w - (cols - 1) * GAP) / cols
            cell_h = (area_h - (rows - 1) * GAP) / rows
            ih_avail = cell_h - cap_h
            if ih_avail <= 0.2 or cw <= 0.2:
                continue
            score = 0.0
            for ai in ars:
                iw = min(cw, ih_avail * ai)
                score += iw * (iw / ai)
        if best is None or score > best[0] + 1e-6:
            best = (score, cols, rows, uniform)
    if best is None:
        c = min(cap_cols, n)
        return c, math.ceil(n / c), uniform
    return best[1], best[2], best[3]


def _norm_entry(e):
    """(path, cap[, header]) -> (path, cap, header) with header default None."""
    path = e[0]
    cap = e[1] if len(e) > 1 else None
    header = e[2] if len(e) > 2 else None
    return path, cap, header


def build_grid_slide(prs, title, entries, footer_text, cap_pt=11, cap_h=0.34,
                     hdr_h=0.30):
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    entries = [_norm_entry(e) for e in entries]
    n = len(entries)
    area_top = IMG_TOP
    area_w = SLIDE_W - 2 * MARGIN
    area_h = FOOTER_TOP - IMG_TOP - 0.02

    # If any tile carries a date-header box, reserve a band above each image.
    hb = hdr_h if any(h for _, _, h in entries) else 0.0
    over = cap_h + hb  # per-cell non-image overhead (header above + caption below)

    ars = [_img_ar(p) for p, _, _ in entries]
    cols, rows, uniform = _best_cols(ars, n, area_w, area_h, over,
                                     min(MAX_GRID_COLS, n))
    if uniform:
        a = sorted(ars)[len(ars) // 2]
        iw, ih = _uniform_img_size(a, cols, rows, area_w, area_h, over)
        block_w = cols * iw + (cols - 1) * GAP
        block_h = rows * (ih + over) + (rows - 1) * GAP
        x0 = MARGIN + (area_w - block_w) / 2
        y0 = area_top + max(0.0, (area_h - block_h) / 2)
        for i, (path, cap, header) in enumerate(entries):
            r, c = divmod(i, cols)
            left = x0 + c * (iw + GAP)
            top = y0 + r * (ih + over + GAP)
            if header:
                _place_headers(slide, header, left, top, iw, hb)
            img_top = top + hb
            add_image_in_box(slide, str(path), left, img_top, iw, ih)
            _tile_caption(slide, cap, path, left, img_top + ih, iw, cap_h, cap_pt)
    else:
        cell_w = (area_w - (cols - 1) * GAP) / cols
        cell_h = (area_h - (rows - 1) * GAP) / rows
        img_h = cell_h - over
        for i, (path, cap, header) in enumerate(entries):
            r, c = divmod(i, cols)
            left = MARGIN + c * (cell_w + GAP)
            top = area_top + r * (cell_h + GAP)
            if header:
                _place_headers(slide, header, left, top, cell_w, hb)
            img_top = top + hb
            add_image_in_box(slide, str(path), left, img_top, cell_w, img_h)
            _tile_caption(slide, cap, path, left, img_top + img_h,
                          cell_w, cap_h, cap_pt)

    add_textbox(slide, footer_text, MARGIN, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


# ---------------------------------------------------------------------------
# Content plan
# ---------------------------------------------------------------------------
def build_plan():
    plan, missing = [], []

    def add_grid(title, pairs, footer=FOOTER_CMP):
        kept, miss = [], []
        for e in pairs:
            (kept if path_exists(e[0]) else miss).append(e)
        missing.extend(e[0] for e in miss)
        if kept:
            plan.append(("grid", title, kept, footer))

    # NOTHING POOLED — every slide shows the 3 experiments (BYEXP tiled figure
    # or the 3 per-experiment figures side by side). per-experiment tiles get a
    # date box (header) above the plot: (path, caption=None, header=date).
    BYEXP_DIR = PROF / "singles"
    EXP_LABELS = [disp for disp, _ in EXP_DATES]  # boxes over the 3 sub-panels

    def per_exp(base, fmt):
        return [(base / fmt.format(d), None, disp) for disp, d in EXP_DATES]

    # --- Sec 1: across-experiment density overview (BYEXP tiled) ---
    plan.append(("divider", "Vimentin density vs NE geometry — across experiments",
                 "one panel per experiment (Apr 29 / Jan 23 / Feb 27) · "
                 "Mean ± SEM · " + FOOTER_CMP))
    add_grid("Vimentin density vs min curvature — by experiment  (headline)",
             [(BYEXP_DIR / "vim_geomdens_mincurv_perinuc05_BYEXP_line.png",
               None, EXP_LABELS)],
             "panels: Apr 29 2022 | Jan 23 2024 | Feb 27 2024 · Mean ± SEM · "
             + FOOTER_CMP)
    add_grid("Vimentin density vs mean curvature — by experiment",
             [(BYEXP_DIR / "vim_geomdens_meancurv_perinuc05_BYEXP_line.png",
               None, EXP_LABELS)],
             "panels: Apr 29 2022 | Jan 23 2024 | Feb 27 2024 · Mean ± SEM · "
             + FOOTER_CMP)
    add_grid("Vimentin density vs hull-boundary distance — by experiment",
             [(BYEXP_DIR / "vim_geomdens_hulldist_perinuc05_BYEXP_line.png",
               None, EXP_LABELS)],
             "panels: Apr 29 2022 | Jan 23 2024 | Feb 27 2024 · Mean ± SEM · "
             + FOOTER_CMP)

    # --- Sec 2: per-experiment enrichment (3 side by side) ---
    plan.append(("divider",
                 "Vimentin per-cell enrichment — per experiment",
                 "region density ÷ cell mean (>1 = enriched) · one dot per cell · "
                 "3 experiments side by side"))
    add_grid("Vimentin enrichment on concave NE (min curv < 0) — per experiment  "
             "(headline)",
             per_exp(ENR, "vim_mincurv_lt0_shells_{}.png"),
             "region density ÷ cell mean · ref line 1 · " + FOOTER_CMP)
    add_grid("Vimentin enrichment on concave NE (mean curv < 0) — per experiment",
             per_exp(ENR, "vim_meancurv_lt0_shells_{}.png"),
             "region density ÷ cell mean · ref line 1 · " + FOOTER_CMP)
    add_grid("Vimentin enrichment in deep invaginations (hull dist > 0.5 μm) — "
             "per experiment",
             per_exp(ENR, "vim_hulldist_gt0.5um_shells_{}.png"),
             "region density ÷ cell mean · ref line 1 · " + FOOTER_CMP)
    add_grid("Vimentin correlation vs min curvature — per experiment",
             per_exp(ENR, "vim_corr_mincurv_shells_{}.png"),
             "per-cell Pearson r · ref line 0 · " + FOOTER_CMP)
    add_grid("Vimentin correlation vs mean curvature — per experiment",
             per_exp(ENR, "vim_corr_meancurv_shells_{}.png"),
             "per-cell Pearson r · ref line 0 · " + FOOTER_CMP)

    # --- Sec 3: centrosome proximal | distal, per experiment (each tile 1×2) ---
    plan.append(("divider", "Centrosome proximal | distal — per experiment",
                 "each tile is one experiment's Proximal | Distal 1×2 · "
                 + FOOTER_CMP))
    add_grid("Vim enrichment on concave NE (min curv < 0), proximal | distal — "
             "per experiment  (headline)",
             per_exp(ENR, "vim_mincurv_lt0_shells_bycent_{}.png"),
             "region density ÷ cell mean · ref line 1 · " + FOOTER_CMP)
    add_grid("Vim enrichment on concave NE (mean curv < 0), proximal | distal — "
             "per experiment",
             per_exp(ENR, "vim_meancurv_lt0_shells_bycent_{}.png"),
             "region density ÷ cell mean · ref line 1 · " + FOOTER_CMP)
    add_grid("Vim density vs min curvature, proximal | distal — per experiment",
             per_exp(BYCENT,
                     "vim_geomdens_mincurv_shells_OVERLAY_line_bycent_{}.png"),
             "Mean ± SEM density vs curvature · " + FOOTER_CMP)
    add_grid("Vim density vs mean curvature, proximal | distal — per experiment",
             per_exp(BYCENT,
                     "vim_geomdens_meancurv_shells_OVERLAY_line_bycent_{}.png"),
             "Mean ± SEM density vs curvature · " + FOOTER_CMP)

    # --- Sec 4: near-centrosome perinuclear vimentin, per experiment ---
    plan.append(("divider", "Vimentin near the centrosome — per experiment",
                 "perinuclear vimentin within 0.5 μm of the nuclear surface AND "
                 "2 μm of the centrosome · " + FOOTER_CMP))
    add_grid("Vimentin centrosomal enrichment (within 0.5 μm nuc & 2 μm cent) — "
             "per experiment  (KEY METRIC)",
             [(GRID / "vim_enrichment_within_half_um_nuc_2_um_cent_grid.png",
               None)],
             "shell-normalized: mean vim in (0.5 μm perinuc shell ∩ ≤2 μm cent) ÷ "
             "mean vim over the whole 0.5 μm shell · ref line 1 · "
             "panels: Apr 29 2022 | Jan 23 2024 | Feb 27 2024 · " + FOOTER_CMP)
    add_grid("Vimentin fraction near the centrosome (within 0.5 μm nuc & 2 μm "
             "cent) — per experiment  (content readout)",
             [(GRID / "vim_frac_perinuc_within_2_um_cent_grid.png", None)],
             "fraction of perinuclear vim signal within 2 μm of the centrosome "
             "(sum ÷ sum) · panels: Apr 29 2022 | Jan 23 2024 | Feb 27 2024 · "
             + FOOTER_CMP)
    add_grid("Vimentin proximal ÷ distal ratio — per experiment",
             per_exp(NC, "vim_near_cent_prox_over_distal_ratio_{}.png"),
             "near-cent ÷ away-cent · ref line 1 · > 1 = vimentin denser near "
             "the centrosome · " + FOOTER_CMP)

    # --- Sec 5: invagination-pocket scalars (grid_panels, rows = experiments) ---
    plan.append(("divider", "Vimentin in nuclear invagination pockets",
                 "voxel convex-hull decomposition · rows = experiments · "
                 + FOOTER_CMP))
    add_grid("Vimentin in nuclear invagination pockets (rows = experiments)",
             [(GRID / "vim_cyto_in_nuc_hull_vs_near_convex_nuc_MFI_ratio_grid.png",
               "grooves ÷ convex surface"),
              (GRID / "vim_cyto_in_nuc_hull_vs_all_perinuc_MFI_ratio_grid.png",
               "grooves ÷ whole perinuc"),
              (GRID / "vim_frac_in_nuc_convex_hull_grid.png",
               "fraction inside nuclear hull")],
             "rows = experiments · ref line 1 · " + FOOTER_CMP)

    # --- Drop orphan dividers ---
    cleaned = []
    for i, it in enumerate(plan):
        if it[0] == "divider":
            if i + 1 >= len(plan) or plan[i + 1][0] == "divider":
                continue
        cleaned.append(it)
    return cleaned, missing


def main():
    list_only = "--list" in sys.argv
    plan, missing = build_plan()

    print("Output: {}".format(OUTPUT_PATH))
    print("{} content slides (+ title) = {} total\n".format(
        len(plan), 1 + len(plan)))
    for it in plan:
        if it[0] == "divider":
            print("\n=== {} ===".format(it[1]))
        else:
            print("  [{}] {}".format(len(it[2]), it[1]))
    if missing:
        print("\nMISSING ({}):".format(len(missing)))
        for m in missing:
            print("  - {}".format(m))
    if list_only:
        return

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)

    for it in plan:
        if it[0] == "divider":
            build_divider_slide(prs, it[1], it[2])
        else:
            _, title, entries, footer = it
            build_grid_slide(prs, title, entries, footer)

    if OUTPUT_PATH.exists():
        created = backup_presentation(str(OUTPUT_PATH),
                                      backup_base=str(OUTPUT_PATH.parent / "backups"))
        if created:
            print("\nBacked up previous deck.")
    prs.save(str(OUTPUT_PATH))
    print("\nDone. {} slides written to:\n  {}".format(
        len(prs.slides._sldIdLst), OUTPUT_PATH))
    if missing:
        print("Skipped {} missing panel(s).".format(len(missing)))


if __name__ == "__main__":
    main()
