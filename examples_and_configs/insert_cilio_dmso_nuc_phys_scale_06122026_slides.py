"""
insert_cilio_dmso_nuc_phys_scale_06122026_slides.py

DMSO (left) vs Ciliobrevin (right) nucleus physical-scale montage deck for the
06/12/2026 Ciliobrevin Jurkats experiment
(L:/FF/Nucleus_H3K27me3/Ciliobrevin_Jurkats/06122026_firstReplicate_*).

Same deck-wide PPI pinning pattern as insert_actin_qc_cat_vs_fmc_20260607_slides.py:
all PNGs are inserted at one shared pixels-per-inch so the embedded scalebar
renders at the same cm-on-page across every panel and slide.

Slides (first chunk of each combo only):
  1. Nucleus (DNA)            -> physical_scale_images/nucleus_bz/montages
  2. Centrin2 + Nucleus XZ    -> physical_scale_images/cent_nuc_xz/montages

Usage:
    python examples_and_configs/insert_cilio_dmso_nuc_phys_scale_06122026_slides.py
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/CilioD/"
    "Cilio_Jurkats_DMSO_vs_Cilio_nuc_phys_scale_06122026.pptx"
)

PHYS_SCALE_SUBPATH = "prog_fixed_cells/{cond}/physical_scale_images/{combo}/montages"

# Each experiment contributes one block of slides (one per combo below).
# tag      — short date tag woven into slide titles (06/12 and 06/13 7min are
#            different replicates of the same nominal timepoint).
# root     — experiment root dir.
# left/right — (condition_folder, display_label) for the DMSO / Cilio columns.
# tp_label — timepoint string used in slide titles.
EXPERIMENTS = [
    {
        "tag": "06/12/2026",
        "root": (
            "L:/FF/Nucleus_H3K27me3/Ciliobrevin_Jurkats/"
            "06122026_firstReplicate_50uM30minCilio_7min_"
            "h3k27me3-640LP45_p561LP45_egfpCentrin2-488LP45_h405LP40_"
        ),
        "left":  ("A1_DMSO_7min_aCD3_",                 "DMSO, 7 min αCD3"),
        "right": ("A2_Ciliobrevin50um30min_7min_aCD3_", "Cilio 50 μM, 7 min αCD3"),
        "tp_label": "7 min αCD3",
    },
    {
        "tag": "06/13/2026",
        "root": (
            "L:/FF/Nucleus_H3K27me3/Ciliobrevin_Jurkats/"
            "06132026_50uM30minCilio_7-12min_"
            "h3k27me3-640LP45_p561LP45_egfpCentrin2-488LP45_h405LP40_"
        ),
        "left":  ("GbA1_DMSO5to1000-30min_7min_aCD3_", "DMSO, 7 min αCD3"),
        "right": ("GbA2_50uMCilio-30min_7min_aCD3_",   "Cilio 50 μM, 7 min αCD3"),
        "tp_label": "7 min αCD3",
    },
    {
        "tag": "06/13/2026",
        "root": (
            "L:/FF/Nucleus_H3K27me3/Ciliobrevin_Jurkats/"
            "06132026_50uM30minCilio_7-12min_"
            "h3k27me3-640LP45_p561LP45_egfpCentrin2-488LP45_h405LP40_"
        ),
        "left":  ("GaA1_5to1000DMSO-30min_12min_aCD3_", "DMSO, 12 min αCD3"),
        "right": ("GaA2_50uMCilio-30min_12min_aCD3_",   "Cilio 50 μM, 12 min αCD3"),
        "tp_label": "12 min αCD3",
    },
]

# (combo_subfolder, title_template, n_chunks). Template gets .format()'d with
# tp= (e.g. "7 min αCD3") and tag= (e.g. "06/13/2026") per experiment.
# Layout depends on n_chunks:
#   n_chunks == 1 -> 1 row x 2 cols (DMSO left, Cilio right; labels above).
#   n_chunks  > 1 -> 1 row x (2 * n_chunks) cols, all in a single row:
#                    DMSO's n_chunks chunks on the left, Cilio's n_chunks
#                    chunks on the right, with a banner label spanning each
#                    condition's group of columns. All panels share one PPI.
# scale_group: slides with the same group key share one pinned PPI across the
# ENTIRE deck (every experiment), so their embedded scalebars render at the
# same cm. Distinct groups are intentionally separate — the H3K27me3 1x4 slide
# has narrower per-panel cells, so it gets its own group and its own (smaller)
# pinned scalebar rather than forcing the wider broadest-slice panels to shrink
# down to match it.
COMBOS = [
    # (combo_subfolder, title_template, n_chunks, scale_group)
    ("cent_nuc_xz",          "Cent + Nuc, XZ MIP ({tp}, {tag})",                         1, "xz"),
    ("nucleus_bz",           "Nuc (DNA), broadest slice ({tp}, {tag})",                  1, "broad_1c"),
    ("cent_nuc_bz",          "Cent + Nuc, broadest slice ({tp}, {tag})",                 1, "broad_1c"),
    ("H3K27me3_nuc_cent_bz", "H3K27me3 + Cent + Nuc, broadest slice ({tp}, {tag})",      2, "broad_4c"),
]

CHUNK_GLOB = "montage_cells_*.png"

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Slide layout (inches). 13.333 x 7.5 widescreen.
SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

# 1x2 cell grid below the title (label + image per cell)
GRID_LEFT = 0.10
GRID_TOP = 0.60
CELL_W = 6.50
CELL_H = SLIDE_H - GRID_TOP - 0.10   # 6.80"
LABEL_H = 0.30
IMG_H = CELL_H - LABEL_H             # 6.50"
LABEL_FONT_PT = 16
COL_GAP = SLIDE_W - 2 * GRID_LEFT - 2 * CELL_W

CELL_POSITIONS = [
    (GRID_LEFT,                    GRID_TOP),
    (GRID_LEFT + CELL_W + COL_GAP, GRID_TOP),
]

# Scalebar invariant for the H3K27me3 / Jurkat nucleus fixed-cell pipeline.
# Measured empirically with examples_and_configs/check_scalebar_pixel_widths.py
# against this dataset's montages: every scalebar is exactly 104 px wide, which
# means the rendered PPUM is 104 / 5 = 20.8 px/μm (different from the actin
# pipeline's 30 px/μm).
SCALEBAR_PX = 104                            # px (measured)
SCALEBAR_UM = 5                              # μm
PPUM_SOURCE = SCALEBAR_PX / SCALEBAR_UM      # 20.8 px/μm in the rendered PNG

# ---------------------------------------------------------------------------


def _winlong(p) -> str:
    """Return a Win32-safe absolute path string. On Windows, prepends the
    \\\\?\\ extended-length prefix so paths near MAX_PATH (260) still work in
    os.stat / Image.open / python-pptx add_picture. pathlib's glob enumerates
    these files fine, but Path.exists() / os.stat() can silently fail just
    below the limit, so any path that crosses this layer must be wrapped."""
    s = os.path.abspath(str(p))
    if os.name == "nt" and not s.startswith("\\\\?\\"):
        s = "\\\\?\\" + s.replace("/", "\\")
    return s


def _exists_long(p) -> bool:
    """MAX_PATH-safe existence check (pathlib.Path.exists() trips on long paths)."""
    return os.path.exists(_winlong(p))


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _png_dims(path: Path) -> Tuple[int, int]:
    """Return (width_px, height_px) of a PNG without fully decoding it."""
    with Image.open(_winlong(path)) as im:
        return im.size


def compute_slide_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    """Smallest ppi such that every image fits in (max_w_in x max_h_in).
    Called per slide (not deck-wide) so each slide's panels share one PPI —
    the embedded scalebars match across panels on a slide, but slides with
    smaller-pixel sources render at a smaller PPI and thus a larger
    scalebar-in-cm than slides with larger-pixel sources."""
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_image_at_ppi(slide, image_path: Path, ppi: float,
                     area_left: float, area_top: float,
                     area_w: float, area_h: float):
    """Center an image at uniform PPI inside an arbitrary (left, top, w, h)
    image area. native_px / ppi gives inches; both dims pinned by ppi so all
    images in the slide render at one physical scale."""
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    left_in = area_left + (area_w - w_in) / 2
    top_in  = area_top  + (area_h - h_in) / 2
    return slide.shapes.add_picture(
        _winlong(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(w_in),
    )


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _chunk_start_index(p: Path) -> int:
    m = re.match(r"montage_cells_(\d+)", p.name)
    return int(m.group(1)) if m else 0


def find_first_chunks(montages_dir: Path, n: int) -> List[Optional[Path]]:
    """Return the first n chunks sorted by chunk-start index. Pads with None
    if the folder has fewer chunks (or doesn't exist)."""
    if not montages_dir.is_dir():
        return [None] * n
    chunks = sorted(montages_dir.glob(CHUNK_GLOB), key=_chunk_start_index)
    return (chunks + [None] * n)[:n]


def _multichunk_geometry(n_chunks: int):
    """1 row x (2 * n_chunks) columns. DMSO's n_chunks chunks sit on the left,
    Cilio's n_chunks chunks sit on the right, all in a single row. A banner
    label spans each condition's group of n_chunks columns above the images.
    Returns (col_w, img_h, col_lefts, img_top, banner_lefts, banner_widths)."""
    h_margin = GRID_LEFT          # slide left/right inset
    gap = 0.10                    # uniform gap between every column
    total_cols = 2 * n_chunks
    col_w = (SLIDE_W - 2 * h_margin - (total_cols - 1) * gap) / total_cols
    col_lefts = [h_margin + i * (col_w + gap) for i in range(total_cols)]

    img_top = GRID_TOP + LABEL_H
    img_h = IMG_H

    # Banner per condition group: spans from the first chunk's left edge to
    # the last chunk's right edge in that group.
    def _group(start_col):
        left = col_lefts[start_col]
        width = (col_lefts[start_col + n_chunks - 1] + col_w) - left
        return left, width
    left_banner_left,  left_banner_w  = _group(0)
    right_banner_left, right_banner_w = _group(n_chunks)
    banner_lefts   = [left_banner_left,  right_banner_left]
    banner_widths  = [left_banner_w,     right_banner_w]
    return col_w, img_h, col_lefts, img_top, banner_lefts, banner_widths


def build_compare_slide(prs, title_text,
                        left_label, left_imgs,
                        right_label, right_imgs,
                        slide_ppi):
    """Render the comparison slide. With n_chunks=1, lay out DMSO|Cilio side
    by side (current 1-row layout). With n_chunks>1, switch to 2-row layout
    (DMSO row on top with all DMSO chunks side by side, Cilio row below).
    All panels share slide_ppi so embedded scalebars match across the slide."""
    n_chunks = len(left_imgs)
    assert len(right_imgs) == n_chunks

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    missing = []

    if n_chunks == 1:
        # Original 1 x 2 column layout: DMSO left, Cilio right, label per column.
        cells = [
            (left_label,  left_imgs[0],  CELL_POSITIONS[0][0]),
            (right_label, right_imgs[0], CELL_POSITIONS[1][0]),
        ]
        for label, img_path, cell_left in cells:
            add_textbox(
                slide, label,
                cell_left, GRID_TOP, CELL_W, LABEL_H,
                font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
            )
            if img_path is not None and _exists_long(img_path):
                add_image_at_ppi(slide, img_path, slide_ppi,
                                 cell_left, GRID_TOP + LABEL_H, CELL_W, IMG_H)
            else:
                add_textbox(
                    slide, "(missing)",
                    cell_left, GRID_TOP + LABEL_H + IMG_H / 2 - 0.15,
                    CELL_W, 0.3,
                    font_pt=14, color=WHITE,
                )
                missing.append(label)
        return slide, missing

    # n_chunks > 1: single row with DMSO chunks then Cilio chunks side by side
    # (2 * n_chunks columns total). One banner label per condition group.
    col_w, img_h, col_lefts, img_top, banner_lefts, banner_widths = \
        _multichunk_geometry(n_chunks)

    for cond_idx, (cond_label, cond_imgs) in enumerate(
        [(left_label, left_imgs), (right_label, right_imgs)]
    ):
        add_textbox(
            slide, cond_label,
            banner_lefts[cond_idx], GRID_TOP, banner_widths[cond_idx], LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )
        for chunk_idx, img_path in enumerate(cond_imgs):
            col_idx = cond_idx * n_chunks + chunk_idx
            cell_left = col_lefts[col_idx]
            if img_path is not None and _exists_long(img_path):
                add_image_at_ppi(slide, img_path, slide_ppi,
                                 cell_left, img_top, col_w, img_h)
            else:
                add_textbox(
                    slide, "(missing)",
                    cell_left, img_top + img_h / 2 - 0.15,
                    col_w, 0.3,
                    font_pt=14, color=WHITE,
                )
                missing.append(f"{cond_label} (chunk {chunk_idx + 1})")
    return slide, missing


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    # Pre-pass: walk every (experiment, combo) -> collect (left_imgs, right_imgs)
    # where each list has length n_chunks. Also stash scale_group + exp_key so
    # we can pin PPI across all slides in the same (experiment, scale_group).
    slide_specs: List[dict] = []
    for exp in EXPERIMENTS:
        root = Path(exp["root"])
        left_folder, left_label = exp["left"]
        right_folder, right_label = exp["right"]
        tag = exp["tag"]
        tp_label = exp["tp_label"]
        exp_key = f"{tag} {tp_label}"
        for combo_folder, title_tmpl, n_chunks, scale_group in COMBOS:
            title = title_tmpl.format(tp=tp_label, tag=tag)
            left_dir  = root / Path(PHYS_SCALE_SUBPATH.format(
                cond=left_folder,  combo=combo_folder))
            right_dir = root / Path(PHYS_SCALE_SUBPATH.format(
                cond=right_folder, combo=combo_folder))
            left_imgs  = find_first_chunks(left_dir,  n_chunks)
            right_imgs = find_first_chunks(right_dir, n_chunks)
            slide_specs.append({
                "title": title,
                "left_label": left_label,   "left_imgs": left_imgs,
                "right_label": right_label, "right_imgs": right_imgs,
                "left_dir": left_dir,       "right_dir": right_dir,
                "log_key": f"{exp_key}/{combo_folder}",
                "n_chunks": n_chunks,
                "scale_group": scale_group,
                "exp_key": exp_key,
            })

    # For each slide compute the smallest PPI that fits its images in its cell
    # box; then take the max across every slide in the same scale_group (across
    # ALL experiments). Every slide in that group is rendered at that pinned
    # PPI, so all its scalebars match cm-on-page across the whole deck.
    def _cell_box(n_chunks: int):
        if n_chunks == 1:
            return CELL_W, IMG_H
        col_w, img_h, *_ = _multichunk_geometry(n_chunks)
        return col_w, img_h

    group_ppi: dict = {}
    for spec in slide_specs:
        max_w, max_h = _cell_box(spec["n_chunks"])
        imgs = [p for p in (*spec["left_imgs"], *spec["right_imgs"])
                if p is not None and _exists_long(p)]
        own_ppi = compute_slide_ppi(imgs, max_w, max_h) if imgs else 0.0
        sg = spec["scale_group"]
        group_ppi[sg] = max(group_ppi.get(sg, 0.0), own_ppi)

    print(
        f"Deck-wide PPI pinning across {len(slide_specs)} slides. "
        f"Scalebar invariant: {SCALEBAR_UM} μm = {SCALEBAR_PX} px in source "
        f"(PPUM = {PPUM_SOURCE} px/μm — verify with check_scalebar_pixel_widths.py).\n"
    )
    print("Pinned PPI per scale_group:")
    for sg, ppi in sorted(group_ppi.items()):
        bar = SCALEBAR_PX / ppi
        print(f"  {sg:>9s}  PPI={ppi:>7.2f}  "
              f"scalebar={bar:.3f} in = {bar * 2.54:.3f} cm")
    print(f"\nWriting deck to: {OUTPUT_PATH}\n")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    missing_total = []
    slides_added = 0
    for spec in slide_specs:
        title       = spec["title"]
        left_label  = spec["left_label"]
        left_imgs   = spec["left_imgs"]
        right_label = spec["right_label"]
        right_imgs  = spec["right_imgs"]
        left_dir    = spec["left_dir"]
        right_dir   = spec["right_dir"]
        log_key     = spec["log_key"]
        n_chunks    = spec["n_chunks"]
        scale_group = spec["scale_group"]
        slide_ppi   = group_ppi[scale_group]

        _, missing = build_compare_slide(
            prs, title,
            left_label, left_imgs,
            right_label, right_imgs,
            slide_ppi,
        )
        slides_added += 1

        left_ok  = sum(1 for p in left_imgs  if p is not None)
        right_ok = sum(1 for p in right_imgs if p is not None)
        print(
            f"[{log_key}]  group={scale_group:<9s}  "
            f"L:{left_ok}/{n_chunks}  R:{right_ok}/{n_chunks}"
        )

        for cell in missing:
            src = left_dir if cell.startswith(left_label) else right_dir
            missing_total.append(f"{log_key}/{cell}  ({src})")

    prs.save(str(out_path))
    print(f"\nDone. {slides_added} slides written to:\n  {out_path}")

    if missing_total:
        print(f"\nMissing ({len(missing_total)}):")
        for m in missing_total:
            print(f"  - {m}")
    else:
        print("\nAll images found - no missing items.")


if __name__ == "__main__":
    main()
