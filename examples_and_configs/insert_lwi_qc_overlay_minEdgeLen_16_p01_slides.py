"""
insert_lwi_qc_overlay_minEdgeLen_16_p01_slides.py

Overlay-only QC deck reading from the minEdgeLen=16 rerun of the
prog_lwi_swi_only pipeline. For each lmnb1 condition subdirectory,
builds one slide with a single large Overlay montage (p01) below a
bold title.

Source dir:
    cropped/channels/prog_lwi_swi_only_minEdgeLen_16/LWI_QC/Montages_filter_pp_load/

Companion to insert_lwi_qc_overlay_p01_slides.py (which reads from the
original `prog/` pipeline). Output filename includes the minEdgeLen_16
tag so both decks coexist.

Usage:
    python examples_and_configs/insert_lwi_qc_overlay_minEdgeLen_16_p01_slides.py
"""

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/"
    "Ctrl, DZNep, bTub, Lamin B, 20260203/"
    "Ctrl_DZNep_upto1hr_LWI_QC_Overlay_minEdgeLen_16_LaminB1_p01.pptx"
)

PARENT_DIR = (
    "J:/FF/Nucleus_Project_up_to_1hr/"
    "020302026_jurkats_dznep_lmnB1orbTub_561LP55_488LP45_405LP40_"
)

QC_SUBPATH = "cropped/channels/prog_lwi_swi_only_minEdgeLen_16/LWI_QC/Montages_filter_pp_load"
IMAGE_FILENAME = "Montage_Overlay_LaminB1_p01.png"

# Only process lmnb1 wells (skip bTub directories)
MARKER_FILTER = "lmnb1"

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Slide layout (inches). 13.333 x 7.5 widescreen.
SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

# Image area below the title — uses essentially the whole remaining slide.
IMG_LEFT = 0.10
IMG_TOP = 0.60
IMG_BOX_W = SLIDE_W - 2 * 0.10           # 13.13"
IMG_BOX_H = SLIDE_H - IMG_TOP - 0.10     # 6.80"

# ---------------------------------------------------------------------------


def format_condition_name(folder_name: str) -> str:
    """G2A3_h2o_8min_lmnb1_ -> 'H2O 8 min: Lamin B1'. Falls back to raw."""
    condition_map = {"h2o": "H2O", "dznep": "DZNep"}
    timepoint_map = {
        "4min": "4 min", "8min": "8 min", "15min": "15 min",
        "30min": "30 min", "1hr": "1 hr",
    }
    marker_map = {"lmnb1": "Lamin B1", "btub": "β-Tub"}

    parts = folder_name.strip("_").split("_")
    if len(parts) != 4:
        return folder_name

    _, condition, timepoint, marker = parts
    return (
        f"{condition_map.get(condition.lower(), condition)} "
        f"{timepoint_map.get(timepoint.lower(), timepoint)}: "
        f"{marker_map.get(marker.lower(), marker)}"
    )


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    """Add a centered textbox with given text and styling."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    """Place an image inside the given bounding box, preserving aspect
    ratio and centering on the dimension that is < box."""
    pic = slide.shapes.add_picture(
        image_path,
        Inches(box_left),
        Inches(box_top),
        width=Inches(box_w),
    )
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path,
            Inches(box_left),
            Inches(box_top),
            height=Inches(box_h),
        )
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def build_condition_slide(prs, cond_dir: Path):
    """Build a single Overlay slide for one condition.
    Returns (slide, missing_flag)."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    title_text = format_condition_name(cond_dir.name)
    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    img_path = cond_dir / QC_SUBPATH / IMAGE_FILENAME
    if img_path.exists():
        add_image_in_box(slide, str(img_path), IMG_LEFT, IMG_TOP, IMG_BOX_W, IMG_BOX_H)
        return slide, False
    else:
        add_textbox(
            slide, "(missing)",
            IMG_LEFT, IMG_TOP + IMG_BOX_H / 2 - 0.2, IMG_BOX_W, 0.4,
            font_pt=18, color=WHITE,
        )
        return slide, True


def main() -> None:
    parent = Path(PARENT_DIR)
    if not parent.exists():
        print(f"ERROR: Parent directory not found: {parent}")
        return

    condition_dirs = sorted([
        d for d in parent.iterdir()
        if d.is_dir() and MARKER_FILTER in d.name.lower()
    ])
    if not condition_dirs:
        print(f"ERROR: No matching subdirectories found in {parent}")
        return

    print(f"Found {len(condition_dirs)} condition directories.")
    print(f"Writing deck to: {OUTPUT_PATH}\n")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    missing = []
    for cond_dir in condition_dirs:
        _, is_missing = build_condition_slide(prs, cond_dir)
        status = "OK" if not is_missing else "MISSING"
        print(f"[{cond_dir.name}]  {status}")
        if is_missing:
            missing.append(cond_dir.name)

    prs.save(OUTPUT_PATH)
    print(f"\nDone. {len(condition_dirs)} slides written to:\n  {OUTPUT_PATH}")

    if missing:
        print(f"\nMissing ({len(missing)}):")
        for c in missing:
            print(f"  - {c}: {IMAGE_FILENAME}")
    else:
        print("\nAll images found - no missing items.")


if __name__ == "__main__":
    main()
