"""Build a Cellpose QC PowerPoint deck from a YAML config.

Usage:
    python build_cellpose_qc_deck.py <config.yaml>

For every FOV the script writes ``{fov}.{ext}`` (actin MIP + colored mask)
into ``cache_dir``. When ``hoechst_channel`` is set, it *also* writes
``{fov}_with_nuc.{ext}`` (actin + Hoechst overlay + colored mask) in the
same TIFF read pass, and emits a sibling ``<output_pptx>_with_nuc.pptx``
deck. Datasets without a Hoechst channel just produce the plain version.

YAML schema (selected keys):
    output_pptx:        path to .pptx (parent dir created if missing)
    cache_dir:          path for the generated composite images
    raw_dir, mask_dir:  source directories
    raw_pattern:        e.g. "{fov}.tif"
    mask_pattern:       e.g. "cellpose_mask_{fov}.tif"
    actin_channel:      int
    pixel_size_um:      float (required; TIFF metadata is unreliable)
    hoechst_channel:    int (optional; enables with-nuc output)
    hoechst_min_lo:     raw-counts floor on Hoechst lower clip (default 120)
    concurrent_fovs:    int (default 1; >1 enables thread-pool FOV reads)
    fov_ids:            list[str] OR "auto"
    grid_rows, grid_cols: int (default 2, 2)
"""

from __future__ import annotations

import argparse
import re
import sys
import threading
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
from mip_mask_compositor import build_fov_composites  # noqa: E402

_print_lock = threading.Lock()


def _log(msg: str) -> None:
    with _print_lock:
        print(msg, flush=True)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.45
TITLE_FONT_PT = 22

GRID_TOP = 0.55
GRID_LEFT = 0.10
GRID_W = SLIDE_W - 2 * GRID_LEFT
GRID_H = SLIDE_H - GRID_TOP - 0.10
CELL_LABEL_H = 0.28
CELL_LABEL_FONT_PT = 12
CELL_GAP = 0.08


def load_config(path: Path) -> Dict[str, Any]:
    with path.open("r") as f:
        cfg = yaml.safe_load(f)
    required = [
        "output_pptx", "cache_dir", "raw_dir", "mask_dir", "fov_ids",
        "pixel_size_um",
    ]
    missing = [k for k in required if k not in cfg]
    if missing:
        raise ValueError(f"Config {path} missing required keys: {missing}")
    cfg.setdefault("deck_title", "Cellpose QC")
    cfg.setdefault("raw_pattern", "{fov}.tif")
    cfg.setdefault("mask_pattern", "cellpose_mask_{fov}.tif")
    cfg.setdefault("actin_channel", 0)
    cfg.setdefault("contrast_low_pct", 1.0)
    cfg.setdefault("contrast_high_pct", 99.5)
    cfg.setdefault("grid_rows", 2)
    cfg.setdefault("grid_cols", 2)
    cfg.setdefault("scalebar_um", 10.0)
    cfg.setdefault("image_format", "png")
    cfg.setdefault("jpeg_quality", 90)
    cfg.setdefault("concurrent_fovs", 1)
    fmt = str(cfg["image_format"]).lower()
    if fmt not in ("png", "jpg", "jpeg"):
        raise ValueError(f"image_format must be 'png' or 'jpg' (got {fmt!r})")

    cfg.setdefault("actin_low_pct", cfg.get("contrast_low_pct", 1.0))
    cfg.setdefault("actin_high_pct", cfg.get("contrast_high_pct", 99.5))
    cfg.setdefault("skip_plain", False)
    # Hoechst overlay enabled when hoechst_channel is present.
    if "hoechst_channel" in cfg:
        cfg.setdefault("hoechst_low_pct", 2.0)
        cfg.setdefault("hoechst_high_pct", 99.5)
        cfg.setdefault("hoechst_min_lo", 120.0)
        cfg.setdefault("hoechst_color", [255, 0, 0])
    if cfg["skip_plain"] and "hoechst_channel" not in cfg:
        raise ValueError("skip_plain: true requires hoechst_channel (otherwise nothing is produced)")
    return cfg


def _with_nuc_path(p: Path) -> Path:
    """Insert ``_with_nuc`` before the suffix of ``p``."""
    return p.with_name(p.stem + "_with_nuc" + p.suffix)


def resolve_fov_ids(cfg: Dict[str, Any]) -> List[str]:
    """If fov_ids == 'auto', glob raw_dir for files matching raw_pattern.

    Optional ``fov_id_filter`` (regex) further restricts auto-discovered
    IDs — e.g. ``"^\\d+$"`` to drop stray non-numeric files in the raw dir.
    """
    fov_ids = cfg["fov_ids"]
    fov_filter = cfg.get("fov_id_filter")
    filt_re = re.compile(fov_filter) if fov_filter else None
    if isinstance(fov_ids, list):
        ids = [str(f) for f in fov_ids]
    elif isinstance(fov_ids, str) and fov_ids.lower() == "auto":
        raw_dir = Path(cfg["raw_dir"])
        pattern = cfg["raw_pattern"]
        if "{fov}" not in pattern:
            raise ValueError(
                f"fov_ids: auto requires '{{fov}}' in raw_pattern (got {pattern!r})"
            )
        regex = re.compile("^" + re.escape(pattern).replace(r"\{fov\}", r"(?P<fov>.+)") + "$")
        ids = []
        for p in sorted(raw_dir.iterdir()):
            if not p.is_file():
                continue
            m = regex.match(p.name)
            if m:
                ids.append(m.group("fov"))
    else:
        raise TypeError(f"fov_ids must be a list or 'auto', got {fov_ids!r}")
    if filt_re:
        ids = [i for i in ids if filt_re.match(i)]
    return ids


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _cell_box(rows: int, cols: int, r: int, c: int):
    """Return (left, top, w, h) of the (r,c) cell within the GRID area."""
    cell_w = (GRID_W - (cols - 1) * CELL_GAP) / cols
    cell_h = (GRID_H - (rows - 1) * CELL_GAP) / rows
    left = GRID_LEFT + c * (cell_w + CELL_GAP)
    top = GRID_TOP + r * (cell_h + CELL_GAP)
    return left, top, cell_w, cell_h


def add_image_in_cell(slide, image_path: Path, cell_left, cell_top, cell_w, cell_h):
    """Add an image fit-to-cell below the cell label, centered."""
    img_top = cell_top + CELL_LABEL_H
    img_h = cell_h - CELL_LABEL_H
    if img_h <= 0:
        return
    pic = slide.shapes.add_picture(
        str(image_path), Inches(cell_left), Inches(img_top), width=Inches(cell_w),
    )
    h_in = pic.height / 914400.0
    if h_in > img_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            str(image_path), Inches(cell_left), Inches(img_top), height=Inches(img_h),
        )
        w_in = pic.width / 914400.0
        pic.left = Inches(cell_left + (cell_w - w_in) / 2)
    else:
        pic.top = Inches(img_top + (img_h - h_in) / 2)


def build_slide(prs, title: str, items, rows: int, cols: int) -> None:
    """items: list of (fov_id, composite_png_path_or_None)."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)
    add_textbox(
        slide, title,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    for idx, (fov, png_path) in enumerate(items):
        r, c = divmod(idx, cols)
        left, top, w, h = _cell_box(rows, cols, r, c)
        add_textbox(
            slide, f"FOV {fov}",
            left, top, w, CELL_LABEL_H,
            font_pt=CELL_LABEL_FONT_PT, color=WHITE, bold=True,
        )
        if png_path is not None and Path(png_path).exists():
            add_image_in_cell(slide, png_path, left, top, w, h)
        else:
            add_textbox(
                slide, "(missing)",
                left, top + CELL_LABEL_H + h / 2 - 0.15, w, 0.3,
                font_pt=12, color=WHITE,
            )


def _process_fov(
    fov: str, cfg: Dict[str, Any], raw_dir: Path, mask_dir: Path,
    cache_dir: Path, ext: str, has_hoechst: bool,
) -> Dict[str, Any]:
    """Generate the missing composite(s) for one FOV. Pure, thread-safe."""
    raw_path = raw_dir / cfg["raw_pattern"].format(fov=fov)
    mask_path = mask_dir / cfg["mask_pattern"].format(fov=fov)
    skip_plain = bool(cfg.get("skip_plain", False)) and has_hoechst
    plain_path = cache_dir / f"{fov}.{ext}"
    nuc_path = (cache_dir / f"{fov}_with_nuc.{ext}") if has_hoechst else None

    if not raw_path.exists():
        return {"fov": fov, "status": "missing_raw", "detail": str(raw_path),
                "plain": None, "with_nuc": None}
    if not mask_path.exists():
        return {"fov": fov, "status": "missing_mask", "detail": str(mask_path),
                "plain": None, "with_nuc": None}

    need_plain = (not skip_plain) and (not plain_path.exists())
    need_nuc = has_hoechst and not nuc_path.exists()

    if not need_plain and not need_nuc:
        return {"fov": fov, "status": "cached",
                "plain": plain_path, "with_nuc": nuc_path}

    build_fov_composites(
        raw_tif_path=raw_path,
        mask_tif_path=mask_path,
        actin_channel=int(cfg["actin_channel"]),
        out_plain=plain_path if need_plain else None,
        out_with_nuc=nuc_path if need_nuc else None,
        pixel_size_um=float(cfg["pixel_size_um"]),
        scalebar_um=float(cfg["scalebar_um"]),
        actin_low_pct=float(cfg["actin_low_pct"]),
        actin_high_pct=float(cfg["actin_high_pct"]),
        hoechst_channel=int(cfg["hoechst_channel"]) if has_hoechst else None,
        hoechst_low_pct=float(cfg.get("hoechst_low_pct", 2.0)),
        hoechst_high_pct=float(cfg.get("hoechst_high_pct", 99.5)),
        hoechst_min_lo=float(cfg.get("hoechst_min_lo", 120.0)),
        hoechst_color=tuple(int(x) for x in cfg.get("hoechst_color", [255, 0, 0])),
        image_format=str(cfg["image_format"]),
        jpeg_quality=int(cfg["jpeg_quality"]),
    )
    return {"fov": fov, "status": "composed",
            "plain": plain_path, "with_nuc": nuc_path}


def _assemble_deck(
    out_pptx: Path, deck_title: str, items: List[Tuple[str, Optional[Path]]],
    rows: int, cols: int,
) -> int:
    """Build and save one deck. Returns slide count."""
    per_slide = rows * cols
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide_idx = 0
    for start in range(0, len(items), per_slide):
        slide_idx += 1
        chunk = items[start:start + per_slide]
        first_fov = chunk[0][0]
        last_fov = chunk[-1][0]
        title = f"{deck_title} - FOVs {first_fov}-{last_fov}"
        build_slide(prs, title, chunk, rows, cols)
    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))
    return slide_idx


def build_deck(cfg: Dict[str, Any]) -> Dict[str, List[str]]:
    raw_dir = Path(cfg["raw_dir"])
    mask_dir = Path(cfg["mask_dir"])
    cache_dir = Path(cfg["cache_dir"])
    out_pptx = Path(cfg["output_pptx"])
    cache_dir.mkdir(parents=True, exist_ok=True)
    out_pptx.parent.mkdir(parents=True, exist_ok=True)

    rows = int(cfg["grid_rows"])
    cols = int(cfg["grid_cols"])
    has_hoechst = "hoechst_channel" in cfg
    fmt = str(cfg["image_format"]).lower()
    ext = "jpg" if fmt in ("jpg", "jpeg") else "png"
    n_workers = max(1, int(cfg["concurrent_fovs"]))

    fov_ids = resolve_fov_ids(cfg)
    _log(f"Found {len(fov_ids)} FOV ids; grid {rows}x{cols} per slide; "
         f"hoechst={'on' if has_hoechst else 'off'}; workers={n_workers}")

    results: Dict[str, Dict[str, Any]] = {}

    def _record(r: Dict[str, Any]) -> None:
        results[r["fov"]] = r
        if r["status"] == "missing_raw":
            _log(f"[{r['fov']}] MISSING raw: {r['detail']}")
        elif r["status"] == "missing_mask":
            _log(f"[{r['fov']}] MISSING mask: {r['detail']}")
        elif r["status"] == "cached":
            _log(f"[{r['fov']}] cached")
        else:
            _log(f"[{r['fov']}] composed")

    if n_workers == 1:
        for fov in fov_ids:
            _record(_process_fov(fov, cfg, raw_dir, mask_dir, cache_dir, ext, has_hoechst))
    else:
        with ThreadPoolExecutor(max_workers=n_workers) as ex:
            futures = {
                ex.submit(_process_fov, fov, cfg, raw_dir, mask_dir, cache_dir,
                          ext, has_hoechst): fov
                for fov in fov_ids
            }
            for fut in as_completed(futures):
                fov = futures[fut]
                try:
                    _record(fut.result())
                except Exception as e:  # noqa: BLE001
                    _log(f"[{fov}] ERROR: {e!r}")
                    results[fov] = {"fov": fov, "status": "error",
                                    "plain": None, "with_nuc": None}

    # Preserve original fov_ids order for slides.
    skip_plain = bool(cfg.get("skip_plain", False)) and has_hoechst
    if not skip_plain:
        plain_items: List[Tuple[str, Optional[Path]]] = [
            (fov, results[fov]["plain"]) for fov in fov_ids
        ]
        plain_slides = _assemble_deck(out_pptx, cfg["deck_title"], plain_items, rows, cols)
        _log(f"\nDeck written: {out_pptx}  ({plain_slides} slides)")

    if has_hoechst:
        nuc_items: List[Tuple[str, Optional[Path]]] = [
            (fov, results[fov]["with_nuc"]) for fov in fov_ids
        ]
        nuc_pptx = _with_nuc_path(out_pptx)
        nuc_slides = _assemble_deck(
            nuc_pptx, f"{cfg['deck_title']} (with Hoechst)", nuc_items, rows, cols,
        )
        _log(f"Deck written: {nuc_pptx}  ({nuc_slides} slides)")

    done = [f for f, r in results.items() if r["status"] in ("cached", "composed")]
    missing = [
        f"{f}: {r['status']} ({r.get('detail','')})"
        for f, r in results.items()
        if r["status"].startswith("missing") or r["status"] == "error"
    ]
    if missing:
        _log(f"\nMissing/errored ({len(missing)}):")
        for m in missing:
            _log(f"  - {m}")
    return {"done": done, "missing": missing}


def main(argv: List[str]) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("config", help="Path to YAML config")
    args = parser.parse_args(argv)
    cfg = load_config(Path(args.config))
    result = build_deck(cfg)
    return 1 if result["missing"] else 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
