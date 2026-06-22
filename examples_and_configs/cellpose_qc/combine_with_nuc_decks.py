"""Build one combined "with_nuc" deck from per-condition cache directories.

Each cache dir is one condition. For every condition the script emits a
black section-title slide followed by NxM grid slides of the FOV
``{fov}_with_nuc.jpg`` composites (same layout as the per-condition
decks).

Usage:
    python combine_with_nuc_decks.py --out master.pptx \\
        cache/dir_A cache/dir_B cache/dir_C ...
"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent))
from build_cellpose_qc_deck import (  # noqa: E402
    SLIDE_W, SLIDE_H, BLACK, WHITE,
    add_textbox, set_slide_background, build_slide,
)


def _fov_sort_key(p: Path, suffix: str = "_with_nuc") -> Tuple[int, str]:
    stem = p.stem
    if suffix and stem.endswith(suffix):
        stem = stem[: -len(suffix)]
    m = re.match(r"(\d+)", stem)
    if m:
        return (int(m.group(1)), stem)
    return (10**9, stem)


def _section_label_from_cache_name(cache_name: str) -> str:
    """Strip common prefixes for a tidier section title."""
    s = cache_name
    for prefix in ("cart_20260607_rerun_", "cart_20260607_", "cilio_06132026_"):
        if s.startswith(prefix):
            s = s[len(prefix):]
            return f"{prefix.rstrip('_')}: {s}"
    return s


def add_section_slide(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BLACK)
    add_textbox(
        slide, title,
        0.5, SLIDE_H / 2 - 0.6, SLIDE_W - 1.0, 1.2,
        font_pt=36, color=WHITE, bold=True,
    )


def build_combined(
    caches: List[Path], out_path: Path, deck_title: str,
    rows: int, cols: int, suffix: str = "_with_nuc",
) -> int:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # Lead-in title slide.
    add_section_slide(prs, deck_title)

    per_slide = rows * cols
    total_added = 1

    for cache_dir in caches:
        cd = Path(cache_dir)
        if not cd.is_dir():
            print(f"[{cd.name}] SKIP (not a directory)", flush=True)
            continue
        if suffix:
            jpg_pat, png_pat = f"*{suffix}.jpg", f"*{suffix}.png"
        else:
            # Plain mode: only files that are NOT _with_nuc.
            jpg_pat, png_pat = "*.jpg", "*.png"
        nuc_files = [
            p for p in sorted(cd.glob(jpg_pat),
                              key=lambda x: _fov_sort_key(x, suffix))
            if suffix or "_with_nuc" not in p.stem
        ]
        if not nuc_files:
            nuc_files = [
                p for p in sorted(cd.glob(png_pat),
                                  key=lambda x: _fov_sort_key(x, suffix))
                if suffix or "_with_nuc" not in p.stem
            ]
        if not nuc_files:
            print(f"[{cd.name}] SKIP (no matching composites in cache)", flush=True)
            continue

        section_title = _section_label_from_cache_name(cd.name)
        add_section_slide(prs, section_title)
        total_added += 1

        items: List[Tuple[str, Optional[Path]]] = [
            (p.stem[: -len(suffix)] if suffix and p.stem.endswith(suffix) else p.stem, p)
            for p in nuc_files
        ]
        for start in range(0, len(items), per_slide):
            chunk = items[start:start + per_slide]
            sub = f"{section_title}  FOVs {chunk[0][0]}-{chunk[-1][0]}"
            build_slide(prs, sub, chunk, rows, cols)
            total_added += 1

        print(f"[{cd.name}] {len(nuc_files)} FOVs -> "
              f"{(len(nuc_files) + per_slide - 1)//per_slide} slides", flush=True)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return total_added


def main(argv: List[str]) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--out", required=True, help="Output .pptx path")
    parser.add_argument("--title", default="Cellpose QC: all conditions (with Hoechst)")
    parser.add_argument("--rows", type=int, default=3)
    parser.add_argument("--cols", type=int, default=3)
    parser.add_argument(
        "--composite-suffix", default="_with_nuc",
        help="Filename suffix before the extension. Default '_with_nuc'. "
             "Pass '' (empty) to pick plain {fov}.jpg/png composites.",
    )
    parser.add_argument("caches", nargs="+", help="One or more cache directories")
    args = parser.parse_args(argv)

    caches = [Path(c) for c in args.caches]
    slides = build_combined(
        caches, Path(args.out), args.title, args.rows, args.cols,
        suffix=args.composite_suffix,
    )
    print(f"\nDeck written: {args.out}  ({slides} slides)", flush=True)
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
