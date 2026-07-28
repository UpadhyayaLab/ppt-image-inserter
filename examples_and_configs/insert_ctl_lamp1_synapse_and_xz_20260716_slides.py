"""
insert_ctl_lamp1_synapse_and_xz_20260716_slides.py

Combined montage deck for the 20260716 fixed CTL experiment
(L:/FF/Nucleus_granules/CTL_fixed/20260716_Fixed_CTLs_glass_centrosome_polarization_granules_nucleus).

Modeled on insert_ctl_lamp1_synapse_and_xz_20260617_slides.py — same block set,
same channel mapping (LAMP1 = granules, MT = β-tubulin/centrosome-context, no
dedicated centrosome stain). Difference: this run has THREE timepoints
(3 min, 5 min, 12 min) instead of two, shown side by side per slide in
chronological order (C3=3min | C2=5min | C1=12min), one FOV (first chunk) each.

Same scale-group scheme:
  - broad / xy_phys / xz_phys / syn_phys / invag_slice / synapse           — 3-panel row layout
  - xzpanel_phys / companel_phys (wide 3-panel-image sources)               — full-width, 3 stacked rows

All physical_scale_images/ groups share the 5 µm / 104 px scalebar.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_ctl_lamp1_synapse_and_xz_20260716_slides.py
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

EXPERIMENT_ROOT = (
    "L:/FF/Nucleus_granules/CTL_fixed/"
    "20260716_Fixed_CTLs_glass_centrosome_polarization_granules_nucleus"
)

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/CTL_Glass_Nucleus_Centrosome/CTL_fixed_LAMP1/"
    "CTL_fixed_LAMP1_combined_20260716.pptx"
)

# Chronological order (3 min → 5 min → 12 min).
CONDITIONS = [
    ("C3_3min_aCD3_ICAM1_3SI_660bTub_535Actin_488LAMP1_405Nuc",  "3 min"),
    ("C2_5min_aCD3_ICAM1_3SI_660bTub_535Actin_488LAMP1_405Nuc",  "5 min"),
    ("C1_12min_aCD3_ICAM1_3SI_660bTub_535Actin_488LAMP1_405Nuc", "12 min"),
]

TP_TITLE = " — ".join(c[1] for c in CONDITIONS)  # "3 min — 5 min — 12 min"

# Block list: (subpath_template_with_{cond}, title, scale_group).
BLOCKS = [
    # --- Broadest slice: the 3-channel combined view only. ---
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_nuc_bz/montages",
        f"LAMP1 + MT + Nuc, broadest slice — {TP_TITLE}",
        "broad",
    ),
    # --- Single-slice XY merge at the centrosome (β-tubulin) Z-slice. ---
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_nuc_com/montages",
        f"LAMP1 + MT + Nuc, centrosome slice — {TP_TITLE}",
        "xy_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_com/montages",
        f"LAMP1 + MT, centrosome slice — {TP_TITLE}",
        "xy_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_com_adaptive_merge/montages",
        f"LAMP1 + MT, centrosome slice (adaptive) — {TP_TITLE}",
        "xy_phys",
    ),
    # Wide per-channel + merge panel — uses stacked full-width layout.
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_com_adaptive_panels/montages",
        f"LAMP1 + MT, centrosome slice — channels + merge (adaptive) — {TP_TITLE}",
        "companel_phys",
    ),
    # --- Actin synapse mask + XZ MIPs ---
    (
        "{cond}/cropped/channels/prog_fixed_cells/actin/bottom_slice_seg/montages",
        f"Actin synapse mask (bottom slice) — {TP_TITLE}",
        "synapse",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/MT_nuc_xz/montages",
        f"MT + Nuc XZ MIP — {TP_TITLE}",
        "xz_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_nuc_xz/montages",
        f"LAMP1 + Nuc XZ MIP — {TP_TITLE}",
        "xz_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_nuc_xz/montages",
        f"LAMP1 + MT + Nuc XZ MIP — {TP_TITLE}",
        "xz_phys",
    ),
    # Actin XZ MIPs.
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_xz/montages",
        f"Actin XZ MIP — {TP_TITLE}",
        "xz_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_nuc_xz/montages",
        f"Actin + Nuc XZ MIP — {TP_TITLE}",
        "xz_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_nuc_xz_planes/montages",
        f"Actin + Nuc XZ MIP, cell top/bottom marked — {TP_TITLE}",
        "xz_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_xz_nolines/montages",
        f"Actin XZ MIP (no lines) — {TP_TITLE}",
        "xz_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_nuc_xz_planes_nolines/montages",
        f"Actin + Nuc XZ MIP, planes (no lines) — {TP_TITLE}",
        "xz_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_MT_xz_nolines/montages",
        f"Actin + MT XZ MIP — {TP_TITLE}",
        "xz_phys",
    ),
    # No-nucleus XZ MIPs.
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_xz_nolines/montages",
        f"LAMP1 XZ MIP — {TP_TITLE}",
        "xz_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/MT_xz_nolines/montages",
        f"MT XZ MIP — {TP_TITLE}",
        "xz_phys",
    ),
    # 3-panel XZ MIP — stacked full-width layout.
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_xz_panel_nolines/montages",
        f"MT / LAMP1 / merge, XZ MIP panel (no lines) — {TP_TITLE}",
        "xzpanel_phys",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_MT_xz_panel_nolines/montages",
        f"Actin / MT / merge, XZ MIP panel (no lines) — {TP_TITLE}",
        "xzpanel_phys",
    ),
    # LAMP1 + MT at the synapse plane (no nucleus).
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_syn/montages",
        f"LAMP1 + MT, synapse plane — {TP_TITLE}",
        "syn_phys",
    ),
    # LAMP1 + actin at the synapse (3-slice MIP, panel view: channels + merge).
    (
        "{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_actin_syn3mip_fixed/montages/panels",
        f"LAMP1 + Actin, synapse 3-slice MIP (panels) — {TP_TITLE}",
        "syn3mip_phys",
    ),
]

CHUNK_GLOB = "montage_cells_*.png"

# Deepest-invag-slice blocks (non-standard filenames).
EXTRA_BLOCKS = [
    (
        "{cond}/cropped/channels/prog_fixed_cells/Lamp1/deepest_invag_slice/merges/montages_deepest_invag",
        "montage_cells_*_with_MT.png",
        f"LAMP1 + MT, deepest invag slice — {TP_TITLE}",
        "invag_slice",
    ),
    (
        "{cond}/cropped/channels/prog_fixed_cells/MT/deepest_invag_slice/merges/montages_deepest_invag",
        "montage_cells_*.png",
        f"MT, deepest invag slice — {TP_TITLE}",
        "invag_slice",
    ),
]

# Scalebar invariant (5 µm = 104 px in every physical_scale_images/ montage).
SCALEBAR_PX = 104
SCALEBAR_UM = 5
PPUM_SOURCE = SCALEBAR_PX / SCALEBAR_UM

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = 13.333
SLIDE_H = 7.5

# MAXIMALLY AGGRESSIVE — zero reserved margins / gaps / title band. Image
# cells claim the entire slide. Title and per-cell labels are drawn as
# overlays AFTER the images (so they win z-order) at the slide/cell top,
# where they sit on the black dead-space of width-bound panels or lightly
# overlap the black edge (above cell membrane) of height-bound MIPs.
TITLE_LEFT = 0.02
TITLE_TOP = 0.01
TITLE_WIDTH = SLIDE_W - 2 * 0.02
TITLE_HEIGHT = 0.24
TITLE_FONT_PT = 18

GRID_LEFT = 0.0
GRID_TOP = 0.0
GRID_BOT_MARGIN = 0.0
GRID_H = SLIDE_H - GRID_TOP - GRID_BOT_MARGIN     # 7.50
LABEL_H = 0.22
LABEL_FONT_PT = 12

# Per-cell labels overlay the top of each image cell. For the top row the
# label would collide with the title box (both anchored near y=0), so we
# push top-row labels down to just below the title band.
TITLE_BOTTOM = TITLE_TOP + TITLE_HEIGHT           # 0.25
def label_y_for(cell_top: float) -> float:
    return max(cell_top, TITLE_BOTTOM)

# N-column side-by-side layout (N = len(CONDITIONS)).
N_COND = len(CONDITIONS)
COL_GAP = 0.0
CELL_W = (SLIDE_W - 2 * GRID_LEFT - (N_COND - 1) * COL_GAP) / N_COND
IMG_H = GRID_H                       # full cell — labels overlay, not reserve
CELL_LEFTS = [GRID_LEFT + i * (CELL_W + COL_GAP) for i in range(N_COND)]

# 2-top-1-bottom layout — for wide-ish groups (xz_phys), each cell is
# ~half-slide wide × ~half-height tall so images can be wider than in 3-column.
# Row 1 has cells 0 and 1; row 2 has cell 2 centered. All cells same size so
# every panel shares the same physical scale.
TWO_TOP_ONE_BOT_GROUPS = {"xz_phys", "syn3mip_phys"}
TT_CELL_W = (SLIDE_W - 2 * GRID_LEFT - COL_GAP) / 2
TT_ROW_H = GRID_H / 2
TT_IMG_H = TT_ROW_H                  # full row — label overlays image top edge
TT_ROW1_TOP = GRID_TOP
TT_ROW2_TOP = GRID_TOP + TT_ROW_H
TT_POSITIONS = [
    (GRID_LEFT,                        TT_ROW1_TOP),   # 3 min
    (GRID_LEFT + TT_CELL_W + COL_GAP,  TT_ROW1_TOP),   # 5 min
    ((SLIDE_W - TT_CELL_W) / 2,        TT_ROW2_TOP),   # 12 min centered
]

# N-row stacked full-width layout — for wide 3-panel groups
# (companel_phys, xzpanel_phys). Each row spans the full slide width so
# panels are at max size under group-pinned physical scale. Chronological
# order (3 min top, 5 middle, 12 bottom) means 3+5 sit above 12.
PANEL_GROUPS = {"xzpanel_phys", "companel_phys"}
PANEL_IMG_W = SLIDE_W - 2 * GRID_LEFT
PANEL_ROW_H = GRID_H / N_COND
PANEL_ROW_IMG_H = PANEL_ROW_H
PANEL_ROW_TOPS = [GRID_TOP + i * PANEL_ROW_H for i in range(N_COND)]

# WIDE_2TOP layout — REMOVED. It gave the 12 min panel a huge full-width
# bottom cell, but at the cost of per-cell PPI (different physical scales
# across the three panels on one slide). Kept the empty set so the render
# dispatch still compiles; nothing routes to it.
WIDE_2TOP_GROUPS: set = set()
W2T_TOP_CELL_W = (SLIDE_W - 2 * GRID_LEFT - COL_GAP) / 2
W2T_TOP_H = GRID_H / 2
W2T_BOT_CELL_W = SLIDE_W - 2 * GRID_LEFT
W2T_BOT_H = GRID_H / 2
W2T_TOP_LEFTS = [GRID_LEFT, GRID_LEFT + W2T_TOP_CELL_W + COL_GAP]
W2T_BOT_LEFT = GRID_LEFT
W2T_BOT_TOP = GRID_TOP + W2T_TOP_H
WIDE_2TOP_SHRINK: dict = {}

# ---------------------------------------------------------------------------


def _winlong(p) -> str:
    s = os.path.abspath(str(p))
    if os.name == "nt" and not s.startswith("\\\\?\\"):
        s = "\\\\?\\" + s.replace("/", "\\")
    return s


def _exists_long(p) -> bool:
    return os.path.exists(_winlong(p))


def _png_dims(path: Path) -> Tuple[int, int]:
    with Image.open(_winlong(path)) as im:
        return im.size


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    # Zero internal padding so a narrow text box (e.g. 0.14" label) actually
    # fits the intended font size instead of getting eaten by default margins.
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def compute_group_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_image_at_ppi(slide, image_path: Path, ppi: float,
                     area_left: float, area_top: float,
                     area_w: float, area_h: float):
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    left_in = area_left + (area_w - w_in) / 2
    top_in  = area_top  + (area_h - h_in) / 2
    return slide.shapes.add_picture(
        _winlong(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(w_in),
    )


def _chunk_start_index(p: Path) -> int:
    m = re.match(r"montage_cells_(\d+)", p.name)
    return int(m.group(1)) if m else 0


def find_first_chunk(montages_dir: Path, pattern: str = CHUNK_GLOB) -> Optional[Path]:
    # Long-path safe: pathlib.glob() breaks on paths near MAX_PATH even when
    # is_dir() returns True. Use os.listdir over a \\?\ prefixed path.
    import fnmatch
    long_dir = _winlong(montages_dir)
    if not os.path.isdir(long_dir):
        return None
    try:
        names = os.listdir(long_dir)
    except OSError:
        return None
    names = [n for n in names if fnmatch.fnmatch(n, pattern)]
    if not names:
        return None
    names.sort(key=lambda n: int(re.match(r"montage_cells_(\d+)", n).group(1))
               if re.match(r"montage_cells_(\d+)", n) else 0)
    return montages_dir / names[0]


def build_multi_compare_slide(prs, title_text, panels, slide_ppi):
    """N-column side-by-side layout. panels = list of (label, img_path) tuples
    in the same order as CELL_LEFTS."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    missing = []
    # Image first, label overlays on top so it wins the z-order.
    for i, (label, img_path) in enumerate(panels):
        cell_left = CELL_LEFTS[i]
        if img_path is not None and _exists_long(img_path):
            add_image_at_ppi(slide, img_path, slide_ppi,
                             cell_left, GRID_TOP, CELL_W, IMG_H)
        else:
            add_textbox(
                slide, "(missing)",
                cell_left, GRID_TOP + IMG_H / 2 - 0.15, CELL_W, 0.3,
                font_pt=14, color=WHITE,
            )
            missing.append(label)
    for i, (label, _img) in enumerate(panels):
        cell_left = CELL_LEFTS[i]
        add_textbox(
            slide, label,
            cell_left, label_y_for(GRID_TOP), CELL_W, LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )
    return slide, missing


def build_multi_2top1bot_slide(prs, title_text, panels, slide_ppi):
    """2-top-1-bottom layout for wide-ish groups (xz_phys). All cells share
    the same size, so panels stay at the same physical scale."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    missing = []
    for i, (label, img_path) in enumerate(panels):
        cell_left, cell_top = TT_POSITIONS[i]
        if img_path is not None and _exists_long(img_path):
            add_image_at_ppi(slide, img_path, slide_ppi,
                             cell_left, cell_top, TT_CELL_W, TT_IMG_H)
        else:
            add_textbox(
                slide, "(missing)",
                cell_left, cell_top + TT_IMG_H / 2 - 0.15,
                TT_CELL_W, 0.3, font_pt=14, color=WHITE,
            )
            missing.append(label)
    for i, (label, _img) in enumerate(panels):
        cell_left, cell_top = TT_POSITIONS[i]
        add_textbox(
            slide, label,
            cell_left, label_y_for(cell_top), TT_CELL_W, LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )
    return slide, missing


def build_wide_2top_slide(prs, title_text, panels, shrink: float = 1.0):
    """3+5 min on top row (half-width side-by-side), 12 min on bottom row
    (FULL width). Each panel renders at its own max-fit PPI (physical
    scale is NOT pinned across cells here) so the bottom cell's extra
    horizontal room actually gets used. `shrink` (0<s<=1) scales the
    fit-target area — 1.0 fills the cell, 0.93 renders ~13% smaller."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    # Cells: [top-left, top-right, bottom-full-width]
    cells = [
        (W2T_TOP_LEFTS[0], GRID_TOP,     W2T_TOP_CELL_W, W2T_TOP_H),
        (W2T_TOP_LEFTS[1], GRID_TOP,     W2T_TOP_CELL_W, W2T_TOP_H),
        (W2T_BOT_LEFT,     W2T_BOT_TOP,  W2T_BOT_CELL_W, W2T_BOT_H),
    ]

    missing = []
    for i, (label, img_path) in enumerate(panels):
        cl, ct, cw, ch = cells[i]
        if img_path is not None and _exists_long(img_path):
            # Per-cell PPI — unpinned. Each image renders at its natural
            # best fit inside its cell (scaled by `shrink`), aspect preserved.
            w_px, h_px = _png_dims(img_path)
            fit_w = cw * shrink
            fit_h = ch * shrink
            cell_ppi = max(w_px / fit_w, h_px / fit_h)
            add_image_at_ppi(slide, img_path, cell_ppi, cl, ct, cw, ch)
        else:
            add_textbox(
                slide, "(missing)",
                cl, ct + ch / 2 - 0.15, cw, 0.3, font_pt=14, color=WHITE,
            )
            missing.append(label)
    for i, (label, _img) in enumerate(panels):
        cl, ct, cw, ch = cells[i]
        add_textbox(
            slide, label,
            cl, label_y_for(ct), cw, LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )
    return slide, missing


def build_multi_stacked_slide(prs, title_text, panels, slide_ppi):
    """N-row stacked full-width layout (each row spans the whole slide width).
    panels = list of (label, img_path) tuples in top→bottom order."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    missing = []
    for i, (label, img_path) in enumerate(panels):
        row_top = PANEL_ROW_TOPS[i]
        if img_path is not None and _exists_long(img_path):
            add_image_at_ppi(slide, img_path, slide_ppi,
                             GRID_LEFT, row_top, PANEL_IMG_W, PANEL_ROW_IMG_H)
        else:
            add_textbox(
                slide, "(missing)",
                GRID_LEFT, row_top + PANEL_ROW_IMG_H / 2 - 0.15,
                PANEL_IMG_W, 0.3, font_pt=14, color=WHITE,
            )
            missing.append(label)
    for i, (label, _img) in enumerate(panels):
        row_top = PANEL_ROW_TOPS[i]
        add_textbox(
            slide, label,
            GRID_LEFT, label_y_for(row_top), PANEL_IMG_W, LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )
    return slide, missing


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    root = Path(EXPERIMENT_ROOT)

    # Pre-pass: walk every block × every CONDITION -> collect first-chunk panels.
    slide_specs: List[dict] = []
    all_blocks = [(sp, CHUNK_GLOB, t, sg) for (sp, t, sg) in BLOCKS] + list(EXTRA_BLOCKS)
    for subpath_tmpl, glob_pat, title, scale_group in all_blocks:
        dirs = [root / Path(subpath_tmpl.format(cond=cond)) for cond, _ in CONDITIONS]
        imgs = [find_first_chunk(d, glob_pat) for d in dirs]
        if all(img is None for img in imgs):
            print(f"WARNING: skipping '{title}' — no montages found "
                  f"(pattern '{glob_pat}' in '{subpath_tmpl}').")
            continue
        slide_specs.append({
            "title": title,
            "dirs": dirs,
            "imgs": imgs,
            "scale_group": scale_group,
        })

    if not slide_specs:
        print("ERROR: no slides to render — every block was empty.")
        sys.exit(1)

    # One PPI per scale_group, pinned to the largest source image in that group.
    # PPI-binding cell dimensions depend on the layout each group uses:
    #   PANEL_GROUPS         → stacked full-width rows (PANEL_IMG_W × PANEL_ROW_IMG_H)
    #   TWO_TOP_ONE_BOT_GROUPS → 2-top-1-bot half-width cells (TT_CELL_W × TT_IMG_H)
    #   everything else       → 3-column side-by-side (CELL_W × IMG_H)
    group_ppi: dict = {}
    for spec in slide_specs:
        imgs = [p for p in spec["imgs"] if p is not None and _exists_long(p)]
        sg = spec["scale_group"]
        if sg in WIDE_2TOP_GROUPS:
            # Per-cell PPI inside the builder — no group-level pin.
            group_ppi.setdefault(sg, 0.0)
            continue
        if sg in PANEL_GROUPS:
            own = compute_group_ppi(imgs, PANEL_IMG_W, PANEL_ROW_IMG_H) if imgs else 0.0
        elif sg in TWO_TOP_ONE_BOT_GROUPS:
            own = compute_group_ppi(imgs, TT_CELL_W, TT_IMG_H) if imgs else 0.0
        else:
            own = compute_group_ppi(imgs, CELL_W, IMG_H) if imgs else 0.0
        group_ppi[sg] = max(group_ppi.get(sg, 0.0), own)

    PHYS_GROUPS = {"xz_phys", "syn_phys", "syn3mip_phys", "broad", "xy_phys", "xzpanel_phys", "companel_phys"}
    print(f"Pinned PPI per scale_group ({len(slide_specs)} slides total):")
    for sg, ppi in sorted(group_ppi.items()):
        if sg in WIDE_2TOP_GROUPS:
            note = "  (per-cell PPI — physical scale is per-panel, not group-pinned)"
            print(f"  {sg:>13s}  PPI=  <per-cell>{note}")
            continue
        if sg in PHYS_GROUPS:
            bar = SCALEBAR_PX / ppi
            note = f"  scalebar = {bar:.3f} in = {bar * 2.54:.3f} cm"
        else:
            note = "  (layout-pin only — no embedded scalebar)"
        print(f"  {sg:>13s}  PPI={ppi:>7.2f}{note}")
    print(f"\nWriting deck to: {OUTPUT_PATH}\n")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    tp_labels = [tp for _, tp in CONDITIONS]
    missing_total = []
    for spec in slide_specs:
        slide_ppi = group_ppi[spec["scale_group"]]
        panels = list(zip(tp_labels, spec["imgs"]))
        sg = spec["scale_group"]
        if sg in WIDE_2TOP_GROUPS:
            shrink = WIDE_2TOP_SHRINK.get(sg, 1.0)
            _, missing = build_wide_2top_slide(prs, spec["title"], panels, shrink=shrink)
        else:
            if sg in PANEL_GROUPS:
                builder = build_multi_stacked_slide
            elif sg in TWO_TOP_ONE_BOT_GROUPS:
                builder = build_multi_2top1bot_slide
            else:
                builder = build_multi_compare_slide
            _, missing = builder(prs, spec["title"], panels, slide_ppi)
        status = "  ".join(
            f"{tp}:{'OK' if img else 'MISSING'}" for tp, img in panels
        )
        print(f"[{sg:>13s}]  {status}  {spec['title']}")
        for label in missing:
            idx = tp_labels.index(label) if label in tp_labels else 0
            missing_total.append(f"{spec['title']}/{label}  ({spec['dirs'][idx]})")

    prs.save(str(out_path))
    n_slides = sum(1 for _ in prs.slides)
    print(f"\nDone. {n_slides} slides written to:\n  {out_path}")
    if missing_total:
        print(f"\nMissing ({len(missing_total)}):")
        for m in missing_total:
            print(f"  - {m}")
    else:
        print("\nAll panels found.")


if __name__ == "__main__":
    main()
