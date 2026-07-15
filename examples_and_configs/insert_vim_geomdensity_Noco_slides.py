"""
insert_vim_geomdensity_DMSO_slides.py

Build the Vimentin geometry-density summary deck for the Jan-23-2024
Nocodazole experiment, DMSO arm only (single condition, n≈109).

Vimentin is a cytoplasmic (struct_out_nuc) channel: the question is whether Vimentin density
enriches in / tracks nuclear invaginations and concave NE, sampled in the
perinuc 0.5 µm shell (0.5 µm OUTSIDE the NE — the shell to feature). Because Vimentin
is cytoplasmic its shells are shell_max, shell_mean, perinuc05 (perinuc-1
dropped); there is no NE-boundary shell and no ≥2-shell averaged overlay (so no
02a-c), so this deck leads with the enrichment / correlation violins plus the
geom_density/profiles density-vs-geometry figures.

Sections (single condition, one violin column):
  1. Vimentin vs NE geometry -- enrichment in deep invaginations (03a-b) and on the
     concave surface (03c-e); per-cell Vimentin-vs-geometry correlation (04a-c); and
     Vimentin-vs-curvature correlation within deep invaginations (deepcorr).
  2. Vimentin density profiles -- density vs hull distance / min / mean curvature
     (per-cell lines, Mean±SEM, bars).

Self-contained blank 16:9 deck. Missing panels are skipped; a previous deck at
the output path is backed up first.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_vim_geomdensity_DMSO_slides.py
    conda run -n PPT_editing python examples_and_configs/insert_vim_geomdensity_DMSO_slides.py --list
"""

import math
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation, safe_path, path_exists  # noqa: E402

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
# Channel: filenames use lowercase "vim"; titles display "Vimentin". (Ignore the
# duplicate "Vimentin_"-prefixed files in the compile per user.)
CH = "vim"
DISP = "Vimentin"

ROOT = Path(
    "M:/FF/FF_4TB_2_Backup_fullHD/Nucleus Project_2ndharddrive/"
    "Nucleus deformations analysis/Nucleus - Fixed Cell Data/"
    "20240123_E6-1_Nocodazole_Vimentin/results_compilation_Vim_Noco_geomdensity_20260713"
)
CURATED = ROOT / (CH + "_loc_wrto_invag")            # 03a-e enrichment, 04a-c correlation (single-shell)
ENRICH = ROOT / "geom_density" / "enrichment"        # + vim_deepcorr_*_shells (not in curated)
SINGLES = ENRICH / "singles"
PROFILES = ROOT / "geom_density" / "profiles"        # vim_geomdens_* density-vs-geometry

OUTPUT_PATH = Path(
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/"
    "nuc_mesh_struct_outside_nuc/Vimentin geom-density vs NE geometry, DMSO vs Noco (20240123).pptx"
)

DECK_TITLE = "Vimentin density vs nuclear-envelope geometry — DMSO vs Noco"
DECK_SUBTITLE = (
    "Fixed Jurkats · Nocodazole exp (Jan 23 2024) · DMSO vs Noco 1µM · Vimentin = cytoplasmic (struct_out_nuc), "
    "perinuc 0.5 µm outside NE · density profiles + per-cell scalar comparison violins · compiled 2026-07-13"
)

# Single shell (perinuc 0.5 µm, 0.5 µm outside the NE). Vimentin is cytoplasmic so
# has no NE-boundary shell; every enrichment/correlation violin is single-panel and
# there is no ≥2-shell averaged overlay (no 02a-c).
SHELLS_ENR = "perinuc 0.5 μm shell (0.5 μm outside NE)  ·  one dot per cell  ·  ref line 1 (enriched > 1)"
SHELLS_CORR = "perinuc 0.5 μm shell (0.5 μm outside NE)  ·  one dot per cell  ·  ref line 0"

# ---------------------------------------------------------------------------
# Colors / layout (16:9)  -- shared with the SUN1 geom-density deck.
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0x55, 0x55, 0x55)
FIELD_COLOR = RGBColor(0x2E, 0x5A, 0x88)   # full metric/field name (tracking ID)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.10
GAP = 0.16

TITLE_LEFT = MARGIN
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.55
TITLE_FONT_PT = 28

IMG_TOP = 0.66
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = 6.36

FOOTER_TOP = 7.06
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.40
FOOTER_FONT_PT = 9

CAPTION_HEIGHT = 0.55


# ---------------------------------------------------------------------------
# Slide helpers  (same layout engine as insert_sun1_dna_geomdensity_N1vsN2_slides.py)
# ---------------------------------------------------------------------------
def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    """Add an image scaled to fit within (box_w, box_h), centered in the box."""
    image_path = safe_path(image_path)
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.7, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=38, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.1, SLIDE_W - 2 * MARGIN, 1.1,
                font_pt=16, color=GREY, italic=True)


def build_divider_slide(prs, title, subtitle=""):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, title, MARGIN, 3.0, SLIDE_W - 2 * MARGIN, 1.2,
                font_pt=40, color=BLACK, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, MARGIN, 4.25, SLIDE_W - 2 * MARGIN, 0.9,
                    font_pt=18, color=GREY, italic=True)


def build_single_slide(prs, title, image_path, caption, footer=None):
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    cap_h = CAPTION_HEIGHT if caption else 0.0
    img_h = IMG_BOX_H - cap_h
    add_image_in_box(slide, str(image_path), MARGIN, IMG_TOP, IMG_BOX_W, img_h)
    if caption:
        add_textbox(slide, caption, MARGIN, IMG_TOP + img_h, IMG_BOX_W, cap_h,
                    font_pt=12, color=GREY)
    add_textbox(slide, footer or rel_footer(image_path), MARGIN, FOOTER_TOP,
                FOOTER_WIDTH, FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


def build_hero_slide(prs, title, entries, footer_text, cap_h=0.0,
                     cap_pt=11, top_frac=0.55):
    """One "hero" image on top spanning the full width, the rest in a row below."""
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    area_top = IMG_TOP
    area_h = FOOTER_TOP - IMG_TOP - 0.02
    top_h = (area_h - GAP) * top_frac
    bot_h = (area_h - GAP) * (1 - top_frac)

    def place(path, cap, left, top, w, h):
        add_image_in_box(slide, str(path), left, top, w, h - cap_h)
        if cap and cap_h:
            add_textbox(slide, cap, left, top + h - cap_h, w, cap_h,
                        font_pt=cap_pt, color=GREY)

    place(entries[0][0], entries[0][1], MARGIN, area_top, IMG_BOX_W, top_h)
    rest = entries[1:]
    if rest:
        nb = len(rest)
        cw = (IMG_BOX_W - (nb - 1) * GAP) / nb
        bot_top = area_top + top_h + GAP
        for i, (path, cap) in enumerate(rest):
            place(path, cap, MARGIN + i * (cw + GAP), bot_top, cw, bot_h)
    add_textbox(slide, footer_text, MARGIN, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


def add_labeled_caption(slide, lines, left, top, width, height):
    """Multi-line caption; each line is (text, font_pt, color). Used to show a
    human hint plus the FULL filename stem (the tracking ID) under each tile."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    for i, (text, pt, color) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.size = Pt(pt)
        run.font.color.rgb = color
    return box


def _tile_caption(slide, cap, path, left, top, width, height, cap_pt):
    lines = []
    if cap:
        lines.append((cap, min(cap_pt, 10), GREY))
    lines.append((Path(str(path)).stem, 8, FIELD_COLOR))   # full field name for tracking
    add_labeled_caption(slide, lines, left, top, width, height)


def build_grid_slide(prs, title, entries, footer_text, max_cols,
                     cap_pt=11, cap_h=0.28):
    """Tile entries in a grid (max_cols=1 => vertical stack for wide images)."""
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    n = len(entries)
    cols = min(max_cols, n)
    rows = math.ceil(n / cols)
    area_top = IMG_TOP
    area_h = FOOTER_TOP - IMG_TOP - 0.02
    cell_w = (SLIDE_W - 2 * MARGIN - (cols - 1) * GAP) / cols
    cell_h = (area_h - (rows - 1) * GAP) / rows
    img_h = cell_h - cap_h
    for i, (path, cap) in enumerate(entries):
        r, c = divmod(i, cols)
        left = MARGIN + c * (cell_w + GAP)
        top = area_top + r * (cell_h + GAP)
        add_image_in_box(slide, str(path), left, top, cell_w, img_h)
        _tile_caption(slide, cap, path, left, top + img_h, cell_w, cap_h, cap_pt)
    add_textbox(slide, footer_text, MARGIN, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


def rel_footer(path):
    try:
        return Path(path).relative_to(ROOT).as_posix()
    except ValueError:
        return Path(path).name


# ---------------------------------------------------------------------------
# Content
# ---------------------------------------------------------------------------
def _tiles(pairs):
    """(path, caption) pairs -> kept present ones; record missing."""
    kept, miss = [], []
    for path, cap in pairs:
        (kept if path_exists(path) else miss).append((path, cap))
    return kept, [p for p, _ in miss]


def build_plan():
    """DMSO vs Noco-1µM comparison (perinuc 0.5 µm). Density profiles + the
    per-cell scalar comparison violins (grid_panels/). NOTE: the curated
    mesh-face enrichment/correlation violins (vim_loc_wrto_invag 03a-e/04a-c)
    are DEGENERATE in this compile (every cell collapses to one value → flat
    lines), so they are deliberately skipped; the density profiles use the same
    data correctly and are fine."""
    plan, missing = [], []

    # Density profiles (DMSO vs Noco-1µM), one big figure per geometry.
    plan.append(("divider", "Vimentin density vs NE geometry — DMSO vs Noco 1µM",
                 "density vs geometry · per-cell lines, Mean±SEM, bars · perinuc 0.5 µm"))
    prof_cap = ("perinuc 0.5 µm (0.5 µm outside NE)  ·  DMSO vs Noco 1µM  ·  "
                "rows: per-cell lines · Mean±SEM · bars")
    for geom, nice in (("hulldist", "hull-boundary distance"),
                       ("mincurv", "min curvature"), ("meancurv", "mean curvature")):
        p = PROFILES / "vim_geomdens_{}_Jan_23_2024.png".format(geom)
        if path_exists(p):
            plan.append(("single", "Vimentin density vs {}".format(nice), p, prof_cap))
        else:
            missing.append(p.name)

    # --- Vimentin invagination & centrosome scalars, DMSO-vs-Noco COMPARISON
    # violins from grid_panels/. ---
    GRID = ROOT / "grid_panels"

    def g(stem, cap):
        return (GRID / (stem + "_grid.png"), cap)

    plan.append(("divider", "Vimentin invagination & centrosome scalars",
                 "DMSO vs Noco 1µM comparison violins · one dot per cell · ref line 1"))

    # Featured: Vimentin in the nuclear invagination pockets (voxel convex-hull)
    kept, miss = _tiles([
        g("vim_cyto_in_nuc_hull_vs_near_convex_nuc_MFI_ratio", "grooves ÷ convex surface"),
        g("vim_cyto_in_nuc_hull_vs_all_perinuc_MFI_ratio", "grooves ÷ whole perinuc shell"),
    ])
    missing += miss
    if kept:
        plan.append(("grid", "Vimentin in the nuclear invagination pockets — DMSO vs Noco", kept,
                     "voxel convex-hull decomposition · >1 = Vimentin enriched in the grooves (inside hull, outside nucleus)",
                     2, 13, 0.30))

    # Convex-hull supporting metrics
    kept, miss = _tiles([
        g("vim_cyto_in_nuc_hull_MFI", "grooves MFI (raw)"),
        g("vim_cyto_in_nuc_hull_sig_fraction", "fraction of Vimentin signal in grooves"),
        g("vim_frac_in_nuc_convex_hull", "fraction of Vimentin inside the hull"),
        g("vim_invag_within_chull_vs_all_chull_MFI_ratio", "invag interior ÷ all within-hull"),
        g("vim_invag_within_chull_vs_convex_within_chull_MFI_ratio", "invag interior ÷ convex rim"),
    ])
    missing += miss
    if kept:
        plan.append(("grid", "Vimentin in the nuclear convex hull — supporting metrics — DMSO vs Noco", kept,
                     "DMSO vs Noco 1µM · ref line 1", 3, 11, 0.30))

    # Deepest invagination
    kept, miss = _tiles([
        g("vim_ratio_by_deepest_invag_all_0_5um", "all faces · 0.5 μm"),
        g("vim_ratio_by_deepest_invag_all_1um", "all faces · 1 μm"),
        g("vim_ratio_by_deepest_invag_away_0_5um", "away faces · 0.5 μm"),
        g("vim_ratio_by_deepest_invag_away_1um", "away faces · 1 μm"),
    ])
    missing += miss
    if kept:
        plan.append(("grid", "Vimentin at the deepest invagination — DMSO vs Noco", kept,
                     "Vimentin at the single deepest invagination ÷ reference · ref line 1",
                     4, 11, 0.30))

    # Centrosome -- NE-facing (cent-side vs away-side; near_cent_level absent here)
    kept, miss = _tiles([
        g("vim_ratio_by_cent_away_0_5um", "cent-side ÷ away-side · 0.5 µm shell"),
        g("vim_ratio_by_cent_away_1um", "cent-side ÷ away-side · 1 µm shell"),
    ])
    missing += miss
    if kept:
        plan.append(("grid", "Vimentin near the centrosome (NE-facing) — DMSO vs Noco", kept,
                     "cent-side vs away-side of the NE · ref line 1", 2, 12, 0.30))

    # Centrosome -- cytoplasmic pool near the centrosome
    kept, miss = _tiles([
        g("vim_frac_around_cent_2um", "fraction of Vimentin within 2 µm of centrosome"),
        g("vim_MFI_around_cent_2um", "Vimentin MFI within 2 µm of centrosome"),
    ])
    missing += miss
    if kept:
        plan.append(("grid", "Vimentin clustered near the centrosome (cytoplasmic pool) — DMSO vs Noco", kept,
                     "3D ball around the centrosome, not NE-restricted", 2, 12, 0.30))

    # Drop orphan dividers (sections whose content is all absent in this compile).
    cleaned = []
    for i, it in enumerate(plan):
        if it[0] == "divider" and (i + 1 >= len(plan) or plan[i + 1][0] == "divider"):
            continue
        cleaned.append(it)
    return cleaned, missing


def main():
    list_only = "--list" in sys.argv
    plan, missing = build_plan()

    n_slides = 1 + len(plan)
    print("Output: {}".format(OUTPUT_PATH))
    print("{} content slides (+ title) = {} total\n".format(len(plan), n_slides))
    for it in plan:
        if it[0] == "divider":
            print("\n=== {} ===".format(it[1]))
        elif it[0] == "single":
            print("  [1] {}".format(it[1]))
        else:  # grid / hero
            tag = "▲" if it[0] == "hero" else ""
            print("  [{}{}] {}".format(len(it[2]), tag, it[1]))
    if missing:
        print("\nMISSING ({}): {}".format(len(missing), missing))
    if list_only:
        return

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)
    for it in plan:
        if it[0] == "divider":
            build_divider_slide(prs, it[1], it[2])
        elif it[0] == "single":
            build_single_slide(prs, it[1], it[2], it[3])
        elif it[0] == "hero":
            _, title, entries, footer, cap_h = it
            build_hero_slide(prs, title, entries, footer, cap_h)
        elif it[0] == "grid":
            _, title, entries, footer, max_cols, cap_pt, cap_h = it
            build_grid_slide(prs, title, entries, footer, max_cols, cap_pt, cap_h)

    if OUTPUT_PATH.exists():
        backup_dir = OUTPUT_PATH.parent / "backups"
        created = backup_presentation(str(OUTPUT_PATH), backup_base=str(backup_dir))
        if created:
            print("\nBacked up previous deck under: {}".format(backup_dir))

    prs.save(str(OUTPUT_PATH))
    total = len(prs.slides._sldIdLst)
    print("\nDone. {} slides written to:\n  {}".format(total, OUTPUT_PATH))
    if missing:
        print("Skipped {} missing panel(s).".format(len(missing)))


if __name__ == "__main__":
    main()
