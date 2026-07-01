"""
insert_ctl_granule_nuc_summary_20260617_slides.py

Condition-comparison summary deck for the 20260617 fixed activated-CTL experiment
(20260617_Fixed_CTLs_glass_centrosome_polarization_granules_nucleus_3min_12min),
compiled into compiled_results/CTL_Glass_nuc_MT_granules_20260617_ncdist_neg0p5_20260701.
The two conditions are TIMEPOINTS — 3 min (early activation) vs 12 min (established
polarization) — αCD3/ICAM1/3SI on glass. Channels: LAMP1 (lytic granules), MT
(β-tubulin; also the centrosome-context stain — no dedicated centrosome marker),
actin, and Hoechst/DNA. Each grid panel is a 3 min vs 12 min comparison plot.
Companion to the montage deck (insert_ctl_lamp1_synapse_and_xz_20260617_slides.py).

Modeled on the blebbistatin summary deck (insert_bleb_summary_slides.py) for the
per-slide layout, title slide, family dividers, and --list dry-run, plus the
noco-washout deck (insert_noco_washout_summary_slides.py) for the side-by-side
multi-panel layout. FAMILIES is built around the granule (LAMP1) metrics as the
focus, carries over the nucleus/centrosome/invagination families from the LatA
summary decks, and adds MT and actin channel families (MT stands in for the
vimentin family used in the LatA/bleb decks, since vimentin was not imaged here).

A FAMILIES entry is (stem, title) for a single panel, or ([(stem, sublabel), ...],
title) to place related panels (e.g. 1 μm / 2 μm thresholds) side by side on one
slide with a sublabel above each. The compile has 288 grid panels total; FAMILIES
is this curated subset — add a row to grow the deck.

Self-contained: builds a blank deck (no template .pptx). Missing panels render
"(missing)" rather than failing. A previous deck is backed up before overwrite.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_ctl_granule_nuc_summary_20260617_slides.py
    # dry run (print planned families/titles, build nothing):
    conda run -n PPT_editing python examples_and_configs/insert_ctl_granule_nuc_summary_20260617_slides.py --list
"""

import os
import sys
import uuid
from pathlib import Path
from xml.sax.saxutils import escape

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation  # noqa: E402

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(
    "L:/FF/Nucleus_granules/CTL_fixed/"
    "20260617_Fixed_CTLs_glass_centrosome_polarization_granules_nucleus_3min_12min/"
    "compiled_results/CTL_Glass_nuc_MT_granules_20260617_ncdist_neg0p5_20260701"
)
GRID_DIR = ROOT / "grid_panels"
CELL_COUNTS_PNG = ROOT / "cell_counts_barplot.png"   # context slide (optional)

OUTPUT_PATH = Path(
    "K:/FF/PPT/PPT_autogeneration/CTL_Glass_Nucleus_Centrosome/"
    "CTL_fixed_LAMP1_20260617/"                       # co-located with the LAMP1 montage deck
    "CTL_fixed_granule_nuc_summary_20260617.pptx"
)

GRID_SUFFIX = "_grid.png"

DECK_TITLE = "Granule polarization and nuclear morphology in activated CTLs"
# n's from cell_counts.csv (3 min 109, 12 min 100). αCD3/ICAM1/3SI, glass.
DECK_SUBTITLE = (
    "3 min (n = 109) vs 12 min (n = 100)  ·  αCD3/ICAM1/3SI, glass  ·  "
    "LAMP1 / MT / actin / DNA  ·  fixed 06/17/2026  ·  compiled 2026-07-01"
)

# ---------------------------------------------------------------------------
# Curated metrics, grouped into families (divider slide per family). Each entry
# is (grid-panel stem, slide title) for one panel, or ([(stem, sublabel), ...],
# slide title) to show related panels side by side on one slide. The stem +
# GRID_SUFFIX is the PNG under grid_panels/. Titles are for navigation; each plot
# carries its own authoritative y-axis label and 3 min / 12 min x-labels.
# ---------------------------------------------------------------------------
FAMILIES = [
    ("Cell and nuclear spreading", [
        ([("nuc_aspect_ratio", "nucleus"),
          ("actin_deform_ratio", "cell")], "Nuclear and cell aspect ratio"),
        ([("actin_bottom_mask_area", "synapse area"),
          ("nuc_broadest_slice_area", "nuclear broadest slice")],
         "Synapse and nuclear broadest-slice area"),
    ]),
    ("Granules — polarization and synapse delivery", [
        ([("centrosome_center_z_rel_bottom_actin_plane", "centrosome"),
          ("Lamp1_zCOF_cell_bottom_distance", "granules")],
         "Centrosome and granule distance to synapse"),
        ([("Lamp1_z50_cell_bottom_distance", "z₅₀"),
          ("Lamp1_z75_cell_bottom_distance", "z₇₅")], "Granule z₅₀ / z₇₅ distance to synapse"),
        ("Lamp1_synapse_g_ave",              "Granule clustering at synapse (g)"),
        ("Lamp1_synapse_inner_outer_ratio",  "Granule synapse inner/outer ratio"),
    ]),
    ("Granules — dispersion", [
        ([("Lamp1_FDD_3D", "3D"),
          ("Lamp1_z_FDD", "axial (z)")], "Granule dispersion (FDD): 3D vs axial"),
        ([("Lamp1_FDD_3D_rel_cent", "3D"),
          ("Lamp1_z_FDD_rel_cent", "axial (z)")], "Granule dispersion rel. centrosome: 3D vs axial"),
    ]),
    ("Granules — signal and centrosome localization", [
        ("Lamp1_total_sig",   "Granule total signal (whole cell)"),
        ("Lamp1_peak_sig",    "Granule peak signal"),
        ("Lamp1_synapse_MFI", "Granule MFI at synapse"),
        ([("Lamp1_synapse_total_sig", "single slice"),
          ("Lamp1_synapse_total_sig_3mip", "3-slice MIP")],
         "Total granule signal at synapse"),
        ([("Lamp1_frac_around_cent_1um", "1 μm"),
          ("Lamp1_frac_around_cent_2um", "2 μm")], "Granule fraction around centrosome"),
        ([("Lamp1_MFI_around_cent_1um", "1 μm"),
          ("Lamp1_MFI_around_cent_2um", "2 μm")], "Granule MFI around centrosome"),
    ]),
    ("Granules — perinuclear and centrosome enrichment", [
        # Lamp1_perinuc_sig_fraction is not computed yet, so its panel renders
        # "no data yet" — the slide appears now (paired with perinuclear MFI) so
        # it is not lost, and auto-fills once the cells are reprocessed.
        ([("Lamp1_all_perinuc_MFI", "MFI"),
          ("Lamp1_perinuc_sig_fraction", "signal fraction")],
         "Granule perinuclear MFI and signal fraction"),
        ([("Lamp1_frac_perinuc_within_1_um_cent", "1 μm cent"),
          ("Lamp1_frac_perinuc_within_2_um_cent", "2 μm cent")],
         "Granule perinuclear fraction near centrosome"),
        ([("Lamp1_cyto_in_nuc_hull_MFI", "MFI"),
          ("Lamp1_cyto_in_nuc_hull_sig_fraction", "signal fraction")],
         "Granule MFI and fraction in cytoplasm within nuclear hull"),
        ("Lamp1_frac_in_nuc_convex_hull", "Granule fraction in nuclear convex hull"),
        ([("Lamp1_enrichment_within_half_um_nuc_2_um_cent", "granule"),
          ("MT_enrichment_within_half_um_nuc_2_um_cent", "MT")],
         "Enrichment near centrosome (0.5 μm of nucleus, 2 μm of cent)"),
    ]),
    ("Centrosome ↔ nucleus", [
        ("nuc_cent_closest_dist",            "Nucleus-centrosome closest distance"),
        ("cent_nuc_norm_dist_sphere_rad",    "Centrosome-to-nuclear-centroid distance (norm. to nuclear sphere radius)"),
        ("centrosome_dist_deepest_real_avg_periphery_ratio", "Centrosome distance to deepest invag vs avg periphery ratio"),
        # Centrosome radial position (r_norm) is defined/labeled in the pipeline
        # (config/get_default_display_labels.m: centrosome_r_norm_MIP,
        # centrosome_r_norm_bottom_plane) but is only computed in the dedicated
        # centrosome-channel path (process_centrosome_channel.m). This experiment
        # derives the centrosome from MT (no centrosome stain), so no panel exists.
        # Enable after adding centrosome_radial_pos to the MT-derived path
        # (process_MT_channel.m) and reprocessing + recompiling:
        # ([("centrosome_r_norm_bottom_plane", "synapse plane"),
        #   ("centrosome_r_norm_MIP", "MIP")], "Centrosome radial position (r_norm)"),
    ]),
    ("Nuclear deformation and invaginations", [
        ("chull_max_D",                       "Max invag depth over full nucleus"),
        ("chull_max_D_by_cent",               "Invagination depth near centrosome"),
        ("chull_mean_D_cent_global_ratio",    "Centrosomal Invagination Index (global)"),
        ([("C_min_F_mean_by_cent", "min principal"),
          ("C_mean_F_mean_by_cent", "mean")],
         "Nuclear surface curvature near centrosome"),
        ("deepest_invag_fraction_chull_volume", "Deepest invag: frac of convex hull volume"),
        ("deepest_region_periph_ratio_025um", "DNA levels near invag"),
        ([("invag_by_cent_centroid_z_cell_bottom_distance_from_MT", "centroid"),
          ("invag_by_cent_tip_z_cell_bottom_distance_from_MT", "tip")],
         "Invagination region (near centrosome): height above synapse"),
    ]),
    ("Invagination orientation", [
        ("avg_normal_angle_adaptive_region_growth",         "Deepest invag orientation"),
        ("avg_normal_angle_adaptive_region_growth_by_cent", "Invag orientation (adaptive) near centrosome"),
        ("avg_normal_angle_by_cent",                        "Invag orientation near centrosome"),
    ]),
    ("Nuclear morphology", [
        ("nuc_solidity",    "Nuclear solidity"),
        ("nuc_volume_mesh", "Nuclear volume"),
        ("nuc_SA_mesh",     "Nuclear surface area"),
    ]),
    ("Microtubules (β-tubulin)", [
        ("MT_frac_in_nuc_convex_hull",                  "MT fraction in nuclear convex hull"),
        ("MT_frac_around_cent_2um",  "MT fraction around centrosome (2 μm)"),
        ("MT_MFI_around_cent_2um",   "MT MFI around centrosome (2 μm)"),
    ]),
    ("Actin — levels and localization", [
        ("actin_total_sig", "Total actin signal (whole cell)"),
        ([("actin_bottom_MFI", "single slice"),
          ("actin_bottom_MFI_3mip", "3-slice MIP")], "Actin MFI at synapse"),
        ([("actin_bottom_total_sig", "single slice"),
          ("actin_bottom_total_sig_3mip", "3-slice MIP")], "Total actin signal at synapse"),
        ("actin_bottom_inner_outer_ratio", "Actin synapse inner/outer ratio"),
        ([("actin_MFI_around_cent_1um", "1 μm"),
          ("actin_MFI_around_cent_2um", "2 μm")], "Actin MFI around centrosome"),
        ([("actin_frac_around_cent_1um", "1 μm"),
          ("actin_frac_around_cent_2um", "2 μm")], "Actin fraction around centrosome"),
    ]),
    # --- Appendix: additional context metrics (added 2026-07-01) ---------------
    ("Cell and nuclear flattening (context)", [
        ("actin_height",        "Cell height"),
        ("nuc_height",          "Nuclear height"),
        ("nuc_centroid_z",      "Nuclear centroid height above synapse"),
        ("nuc_mesh_sphericity", "Nuclear sphericity"),
        ("actin_MIP_circularity", "Cell footprint circularity"),
    ]),
    ("Chromatin / DNA distribution", [
        ("nuc_all_CV",          "DNA intensity CV (heterogeneity)"),
        ("nuc_all_prop_gr_2med", "DNA fraction > 2× median (bright foci)"),
        ("nuc_all_skewness",    "DNA intensity skewness"),
        ("nuc_all_norm_entropy", "DNA distribution normalized entropy"),
    ]),
]

# Pairwise (scatter) section drawn from the compilation's pairwise_plots/ output
# (suite 6, GranuleDelivery_vs_InvagZ): the two centrosome-region invagination
# height metrics, each vs the granule distance-to-synapse metrics, one X-metric
# per slide with its four scatter plots side by side.
_PW_INVAGZ = "pairwise_plots/GranuleDelivery_vs_InvagZ"
_GRANULE_DIST_YS = [
    ("Lamp1_zCOF_cell_bottom_distance", "zCOF"),
    ("Lamp1_z50_cell_bottom_distance",  "z₅₀"),
    ("Lamp1_z75_cell_bottom_distance",  "z₇₅"),
    ("Lamp1_z90_cell_bottom_distance",  "z₉₀"),
]
FAMILIES.append((
    "Pairwise — invagination-region height vs granule delivery",
    [([("{}/{}_VS_{}.png".format(_PW_INVAGZ, _x, _ystem), _ylab)
       for _ystem, _ylab in _GRANULE_DIST_YS], _title)
     for _x, _title in [
        ("invag_by_cent_centroid_z_cell_bottom_distance_from_MT",
         "Invag-region centroid height vs granule distance to synapse"),
     ]],
))

# Pairwise: granule distance to the centrosome (Lamp1_*_rel_cent = intensity-
# weighted mean granule-to-centrosome distance) vs invagination depth near the
# centrosome (chull_max_D_by_cent). From pairwise suite 7 (GranuleCent_vs_
# InvagDepth) — populates after the next recompile; "no data yet" until then.
_PW_CENTDEPTH = "pairwise_plots/GranuleCent_vs_InvagDepth"


def _pw_centdepth(ys):
    """Pairwise panels: chull_max_D_by_cent (X) vs each granule y-metric (Y)."""
    return [("{}/chull_max_D_by_cent_from_MT_VS_{}.png".format(_PW_CENTDEPTH, _y), _lab)
            for _y, _lab in ys]


FAMILIES.append((
    "Pairwise — granule clustering at centrosome vs invagination depth",
    [
        (_pw_centdepth([("Lamp1_FDD_3D_rel_cent", "avg dist to cent"),
                        ("Lamp1_z_FDD_rel_cent",  "axial dist to cent")]),
         "Invag depth near centrosome vs granule distance to centrosome"),
        (_pw_centdepth([("Lamp1_MFI_around_cent_1um", "1 μm"),
                        ("Lamp1_MFI_around_cent_2um", "2 μm")]),
         "Invag depth near centrosome vs granule MFI around centrosome"),
        (_pw_centdepth([("Lamp1_frac_around_cent_1um", "1 μm"),
                        ("Lamp1_frac_around_cent_2um", "2 μm")]),
         "Invag depth near centrosome vs granule fraction around centrosome"),
    ],
))

# ---------------------------------------------------------------------------
# Colors / layout (matches the bleb/noco/vimkd summary decks)
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.10

TITLE_LEFT = MARGIN
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.55
TITLE_FONT_PT = 28

IMG_LEFT = MARGIN
IMG_TOP = 0.66
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = 6.36

# Side-by-side (multi-panel) slides: a sublabel band above the figures, and a
# gap between columns.
COL_GAP = 0.15
SUBLABEL_H = 0.30
SUBLABEL_GAP = 0.04
SUBLABEL_FONT_PT = 16

FOOTER_LEFT = MARGIN
FOOTER_TOP = 7.06
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.40
FOOTER_FONT_PT = 9

# ---------------------------------------------------------------------------


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    """Place an image inside (left, top, w, h), preserving aspect ratio and
    centering on whichever dimension ends up smaller than the box. The grid
    panels are near-square, so they fit to height and center horizontally."""
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def rel_footer(path):
    try:
        return path.relative_to(ROOT).as_posix()
    except ValueError:
        return path.as_posix()


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.7, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=40, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.1, SLIDE_W - 2 * MARGIN, 1.0,
                font_pt=18, color=BLACK, italic=True)


def build_divider_slide(prs, family_name):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, family_name, MARGIN, 3.1, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=44, color=BLACK, bold=True)


def _ext(path):
    """Windows extended-length path (\\\\?\\...) so files whose absolute path
    exceeds the 260-char MAX_PATH are still found and opened. Long metric names
    (and pairwise <X>_VS_<Y> filenames) under the deep compiled_results tree bust
    MAX_PATH; without this prefix Path.exists()/PIL.open silently fail. No-op off
    Windows."""
    p = os.path.abspath(str(path))
    if os.name == "nt" and not p.startswith("\\\\?\\"):
        p = "\\\\?\\" + p.replace("/", "\\")
    return p


def _exists(path):
    return os.path.exists(_ext(path))


def _place_missing(slide, left, top, width):
    add_textbox(slide, "no data yet", left, top + IMG_BOX_H / 2 - 0.2,
                width, 0.4, font_pt=18, color=BLACK)


def build_slide(prs, title_text, panels, footer_text):
    """Title + one or more panels (aspect preserved) + source-path footer.
    `panels` is a list of (path, sublabel); a single entry fills the image box,
    multiple entries are laid out in side-by-side columns with a sublabel above
    each. Returns the list of missing panel paths (not on disk)."""
    slide = _new_slide(prs)
    add_textbox(slide, title_text, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title_text), color=BLACK, bold=True)

    missing = []
    if len(panels) == 1:
        path, _ = panels[0]
        if _exists(path):
            add_image_in_box(slide, _ext(path), IMG_LEFT, IMG_TOP, IMG_BOX_W, IMG_BOX_H)
        else:
            _place_missing(slide, IMG_LEFT, IMG_TOP, IMG_BOX_W)
            missing.append(path)
    else:
        n = len(panels)
        col_w = (IMG_BOX_W - (n - 1) * COL_GAP) / n
        band_top = IMG_TOP + SUBLABEL_H + SUBLABEL_GAP
        band_h = IMG_BOX_H - SUBLABEL_H - SUBLABEL_GAP
        for i, (path, sublabel) in enumerate(panels):
            left = IMG_LEFT + i * (col_w + COL_GAP)
            sub_box = add_textbox(slide, sublabel or "", left, IMG_TOP, col_w, SUBLABEL_H,
                                  font_pt=SUBLABEL_FONT_PT, color=BLACK, bold=True)
            # Bottom-align so the label sits just above its plot, not floating high.
            sub_box.text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
            if _exists(path):
                add_image_in_box(slide, _ext(path), left, band_top, col_w, band_h)
            else:
                _place_missing(slide, left, band_top, col_w)
                missing.append(path)

    add_textbox(slide, footer_text, FOOTER_LEFT, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=BLACK)
    return missing


def _panel_path(stem):
    """Resolve a panel reference to a file path. A plain metric stem maps to
    GRID_DIR/<stem>_grid.png; a reference ending in '.png' is taken relative to
    the compilation ROOT (used for pairwise_plots/<suite>/<X>_VS_<Y>.png)."""
    if stem.endswith(".png"):
        return ROOT / stem
    return GRID_DIR / (stem + GRID_SUFFIX)


def entry_panels(entry_stem):
    """Normalize a FAMILIES entry's stem field into a list of (path, sublabel)."""
    if isinstance(entry_stem, str):
        return [(_panel_path(entry_stem), None)]
    return [(_panel_path(stem), sublabel) for stem, sublabel in entry_stem]


def add_sections(prs, section_spec):
    """Add native PowerPoint sections (the collapsible, named groups shown as
    tabs in the slide navigator / slide sorter). python-pptx has no API for
    this, so inject the p14:sectionLst extension into the presentation part.

    `section_spec` is an ordered list of (section_name, n_slides); the counts
    must sum to the total slide count in build order, and the first section
    must contain the first slide (a PowerPoint requirement)."""
    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    slide_ids = [sldId.get("id") for sldId in prs.slides._sldIdLst]

    parts = ['<p:extLst xmlns:p="{}">'.format(P),
             '<p:ext uri="{{521415D9-36F7-43E2-AB2F-B90AF26B5E84}}">',
             '<p14:sectionLst xmlns:p14="{}">'.format(P14)]
    i = 0
    for name, count in section_spec:
        parts.append('<p14:section name="{}" id="{{{}}}">'.format(
            escape(name), str(uuid.uuid4()).upper()))
        parts.append('<p14:sldIdLst>')
        for sid in slide_ids[i:i + count]:
            parts.append('<p14:sldId id="{}"/>'.format(sid))
        parts.append('</p14:sldIdLst></p14:section>')
        i += count
    parts.append('</p14:sectionLst></p:ext></p:extLst>')

    # The sldIdLst's parent is the <p:presentation> element; extLst goes last.
    prs.slides._sldIdLst.getparent().append(parse_xml("".join(parts)))


def main():
    list_only = "--list" in sys.argv

    n_metrics = sum(len(items) for _, items in FAMILIES)
    # title + (cell-counts if present) + per-family (divider + metric slides)
    est_slides = 1 + (1 if CELL_COUNTS_PNG.exists() else 0) + \
        sum(1 + len(items) for _, items in FAMILIES)

    print("Source: {}".format(GRID_DIR))
    print("{} metric slides across {} families, est. {} slides\n".format(
        n_metrics, len(FAMILIES), est_slides))

    if list_only:
        for fam, items in FAMILIES:
            print("=== {} ({}) ===".format(fam, len(items)))
            for entry_stem, title in items:
                panels = entry_panels(entry_stem)
                flags = " ".join(
                    "{}:{}".format(sub or "-", "OK" if _exists(p) else "MISS")
                    for p, sub in panels)
                print("  [{}] {:<50s} {}".format(
                    "OK " if all(_exists(p) for p, _ in panels) else "MISS",
                    title, flags))
            print("")
        return

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)

    if CELL_COUNTS_PNG.exists():
        build_slide(prs, "Cell counts (3 min vs 12 min)",
                    [(CELL_COUNTS_PNG, None)], rel_footer(CELL_COUNTS_PNG))
    else:
        print("Note: {} not found - skipping cell-counts slide.\n".format(
            CELL_COUNTS_PNG.name))

    # Native PowerPoint sections: an intro section (title + cell counts) then
    # one section per family (its divider slide + metric slides).
    section_spec = [("Overview", len(prs.slides._sldIdLst))]

    missing = []
    for fam, items in FAMILIES:
        build_divider_slide(prs, fam)
        print("=== {} ===".format(fam))
        for entry_stem, title in items:
            panels = entry_panels(entry_stem)
            footer = " | ".join(rel_footer(p) for p, _ in panels)
            miss = build_slide(prs, title, panels, footer)
            status = "OK" if not miss else "MISSING"
            print("  [{}] {} -> {!r}".format(
                status, ", ".join(p.name for p, _ in panels), title))
            missing.extend(p.name for p in miss)
        section_spec.append((fam, 1 + len(items)))
        print("")

    add_sections(prs, section_spec)

    if OUTPUT_PATH.exists():
        backup_dir = OUTPUT_PATH.parent / "backups"
        created = backup_presentation(str(OUTPUT_PATH), backup_base=str(backup_dir))
        if created:
            print("Backed up previous deck to: {}\n".format(backup_dir))

    prs.save(str(OUTPUT_PATH))
    total = len(prs.slides._sldIdLst)
    print("Done. {} metric slides, {} slides written to:\n  {}".format(
        n_metrics, total, OUTPUT_PATH))
    if missing:
        print("\nSkipped {} missing panel(s) (not on disk):".format(len(missing)))
        for m in missing:
            print("  - {}".format(m))
    else:
        print("\nAll curated panels found - no missing items.")


if __name__ == "__main__":
    main()
