"""
insert_chromark_h3k27me3_summary_slides.py

Build a condition-comparison summary deck for the fixed activated-CTL nucleus
"chromark" analysis (tifsCTLsFixed101010aCD3aCD28ICAM_H3K27me3_01292024,
compiled 2026-06-22). Every metric in the curated grid_panels_curated/ folder
is included; each is shown across three comparison views:

  1. all conditions          - 8-condition overview (wide, no stat brackets)
  2. stiffness comparison     - stiffness_27hr + stiffness_51hr side by side
  3. timecourses              - timepoint_1p5kPa + _12kPa + _glass

Metrics are discovered from disk and auto-classified into families
(intensity, GLCM texture, GLCM mid-slice, chromatin organization, 3D/2D
morphology). Ordering: 3D before 2D; DNA (Hoechst) before H3K27me3; where both
channels have the SAME metric they are a "pair" whose slides interleave per view
(DNA all-conditions, H3K27me3 all-conditions, DNA stiffness, ...). Titles use
"DNA" for Hoechst, tag (2D)/(3D) unless obvious (e.g. volume), and give GLCM
distance in pixels. Unrecognized stems land in an "Other metrics" family.

Self-contained: builds a blank deck (no template .pptx). Missing/excluded panels
are skipped (no placeholder, no failure). Previous decks are backed up.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_chromark_h3k27me3_summary_slides.py
    # dry run (print the planned families/titles, build nothing):
    conda run -n PPT_editing python examples_and_configs/insert_chromark_h3k27me3_summary_slides.py --list
"""

import os
import re
import sys
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation  # noqa: E402

# ---------------------------------------------------------------------------
# Datasets. One deck per experiment (or the cross-experiment pool). Select with
# --dataset <key>. Per-dataset: source root, grid subdir, stiffness-view
# timepoints, output path and titles. Everything else (classifier, families,
# all_conditions + timecourse views) is shared across datasets.
# ---------------------------------------------------------------------------
EXP_BASE = "J:/FF/fixed_cell/CTL_nucleus/tifsFixed3SIactivatedCTLs_nucleus"
CROSS_SEARCH = EXP_BASE + "/chromark_cross_experiment"
OUT_DIR = "K:/FF/PPT/PPT_autogeneration/Naive_CTL/chromark"

_STIFF_EARLY = [("stiffness_10min", "10 min"), ("stiffness_3hr", "3 h")]
_STIFF_LATE = [("stiffness_27hr", "27 h"), ("stiffness_51hr", "51 h")]

# Each dataset auto-resolves the NEWEST dated subfolder under `search` matching
# `pat` (so a re-compile to a new date is picked up automatically — no path edit).
# `mark` is the display label for the h3k27me3 channel token (panels use that token
# generically even when the actual mark differs).
DATASETS = {
    "cross": dict(
        search=CROSS_SEARCH, pat="CTL_3marker_*", grid="grid_panels_curated",
        stiffness=_STIFF_EARLY, mark="chromatin mark",
        out=OUT_DIR + "/CTL_chromark_crossexp_3exp_summary.pptx",
        title="CTL chromark — cross-experiment (3 marks as rows)",
        subtitle="Chromatin mark / DNA nuclear chromark (curated) - rows: "
                 "H3K27me3 (0321) / H3K9me3 (0424) / H3K27ac (0531)"),
    "cross_full": dict(
        search=CROSS_SEARCH, pat="CTL_3marker_*", grid="grid_panels",
        stiffness=_STIFF_EARLY, mark="chromatin mark",
        out=OUT_DIR + "/CTL_chromark_crossexp_3exp_FULL_summary.pptx",
        title="CTL chromark — cross-experiment (3 marks as rows, full set)",
        subtitle="Chromatin mark / DNA nuclear chromark (FULL) - rows: "
                 "H3K27me3 / H3K9me3 / H3K27ac"),
    "0129": dict(
        search=EXP_BASE + "/tifsCTLsFixed101010aCD3aCD28ICAM_H3K27me3_01292024/chromark",
        pat="compiled_all_*", grid="grid_panels_curated",
        stiffness=_STIFF_LATE, mark="H3K27me3",
        out=OUT_DIR + "/CTL_chromark_H3K27me3_01292024_summary.pptx",
        title="Activated CTLs on stiffness substrates — H3K27me3 (01/29/2024)",
        subtitle="H3K27me3 / DNA nuclear chromark (curated) - aCD3/aCD28/ICAM"),
    "0321": dict(
        search=EXP_BASE + "/tifs3SICTLsFixed1010aCD3ICAM_H3K27me3_03212024/chromark",
        pat="compiled_all_*", grid="grid_panels_curated",
        stiffness=_STIFF_EARLY, mark="H3K27me3",
        out=OUT_DIR + "/CTL_chromark_H3K27me3_03212024_summary.pptx",
        title="Activated CTLs on stiffness substrates — H3K27me3 (03/21/2024)",
        subtitle="H3K27me3 / DNA nuclear chromark (curated) - aCD3/ICAM - early"),
    # Pending the per-mark chromark compile (auto-picks up the newest when it lands):
    "0424": dict(
        search=EXP_BASE + "/tifs3SICTLsFixed1010aCD3ICAM_H3K9me3_04242024/chromark",
        pat="compiled_all_*", grid="grid_panels_curated",
        stiffness=_STIFF_EARLY, mark="H3K9me3",
        out=OUT_DIR + "/CTL_chromark_H3K9me3_04242024_summary.pptx",
        title="Activated CTLs on stiffness substrates — H3K9me3 (04/24/2024)",
        subtitle="H3K9me3 / DNA nuclear chromark (curated) - aCD3/ICAM - early"),
    "0531": dict(
        search=EXP_BASE + "/tifs3SICTLsFixed101010aCD3aCD28ICAM_H3K27ac_05312024/chromark",
        pat="compiled_all_*", grid="grid_panels_curated",
        stiffness=_STIFF_EARLY, mark="H3K27ac",
        out=OUT_DIR + "/CTL_chromark_H3K27ac_05312024_summary.pptx",
        title="Activated CTLs on stiffness substrates — H3K27ac (05/31/2024)",
        subtitle="H3K27ac / DNA nuclear chromark (curated) - aCD3/aCD28/ICAM - early"),
}
DEFAULT_DATASET = "cross"

GRID_SUFFIX = "_grid.png"

# Comparison views shared across datasets (folder under grid dir -> caption).
ALL_COND_VIEW = "all_conditions"
TIMECOURSE_VIEWS = [
    ("timepoint_1p5kPa", "1.5 kPa"),
    ("timepoint_12kPa", "12 kPa"),
    ("timepoint_glass", "glass"),
]
VIEW_ORDER = ("all", "stiffness", "timecourse")


MARK_LABEL = "H3K27me3"  # display label for the h3k27me3 channel token; per-dataset


def apply_dataset(key):
    """Point the module globals at one dataset, auto-resolving the NEWEST dated
    compile under its search/pat (so new date folders are picked up automatically)."""
    global ROOT, GRID_DIR, STIFFNESS_VIEWS, OUTPUT_PATH, DECK_TITLE, DECK_SUBTITLE, MARK_LABEL
    ds = DATASETS[key]
    search = Path(ds["search"])
    dated = sorted(d for d in search.glob(ds["pat"]) if d.is_dir()) if search.is_dir() else []
    ROOT = dated[-1] if dated else (search / "__no_compile_yet__")
    GRID_DIR = ROOT / ds["grid"]
    STIFFNESS_VIEWS = ds["stiffness"]
    MARK_LABEL = ds.get("mark", "H3K27me3")
    OUTPUT_PATH = Path(ds["out"])
    DECK_TITLE = ds["title"]
    DECK_SUBTITLE = ds["subtitle"]


apply_dataset(DEFAULT_DATASET)

EXCLUDE_PANELS_RAW = set()
EXCLUDE_PANELS = {
    p.split("grid_panels_curated/", 1)[-1].split("grid_panels/", 1)[-1].lstrip("/")
    for p in EXCLUDE_PANELS_RAW
}
# Whole-channel mean/SD are redundant with (and less clearly nuclear-mask than)
# the explicit `*_3d_nuclear_mean/std_int`; drop them so there is a single
# nuclear-mask "mean intensity" / "SD intensity".
DROP_METRICS = {
    "chan_mean_hoechst_3d_int", "chan_mean_h3k27me3_3d_int",
    "chan_std_hoechst_3d_int", "chan_std_h3k27me3_3d_int",
}

# Family order (index -> divider title).
FAMILIES_ORDER = [
    "Nuclear intensity",
    "Texture (GLCM)",
    "Texture (GLCM, max area slice)",
    "Mark–DNA co-distribution",
    "Chromatin organization",
    "Nuclear morphology (3D)",
    "Nuclear morphology (2D)",
    "Nuclear morphology (maximal area slice)",
    "Other metrics",
]
FAM = {name: i for i, name in enumerate(FAMILIES_ORDER)}

# ---------------------------------------------------------------------------
# Auto-classification / titling
# ---------------------------------------------------------------------------
CH_TITLE = {"hoechst": "DNA", "h3k27me3": "H3K27me3"}

INT_STAT_RANK = {"mean": 0, "std": 1, "sd": 1, "skewness": 2, "kurtosis": 3,
                 "median": 4, "q25": 5, "q75": 6, "min": 7, "max": 8,
                 "mode": 9, "rel": 10, "d25": 11, "d75": 12}
INT_STAT_TITLE = {"mean": "mean", "std": "SD", "sd": "SD", "skewness": "skewness",
                  "kurtosis": "excess kurtosis", "median": "median", "q25": "Q25",
                  "q75": "Q75", "min": "min", "max": "max", "mode": "mode",
                  "rel": "relative", "d25": "d25", "d75": "d75"}

# Suffixes marking a DNA-normalised twin of another metric. Whichever spelling
# the upstream pipeline emits, the metric is titled "<base>, per DNA" and sorted
# directly after its unnormalised counterpart.
NORM_SUFFIXES = ("_per_dna", "_dnanorm", "_dna_norm", "_norm_dna", "_perdna")

GLCM_RANK = {"contrast": 0, "correlation": 1, "dissimilarity": 2, "energy": 3,
             "homogeneity": 4, "asm": 5}
GLCM_TITLE = {"asm": "ASM", "contrast": "contrast", "correlation": "correlation",
              "dissimilarity": "dissimilarity", "energy": "energy",
              "homogeneity": "homogeneity"}

# Scope token -> display label. "2dmid" is the MAXIMUM-AREA slice (a real optical
# section), not a midplane -- the upstream field name is historical.
DIM_LABEL = {"2d": "2D", "2dmid": "max area slice", "3d": "3D"}

# Chromatin names -> (sub-rank, title-without-channel-or-dimension).
# Channel is prepended and the scope label appended, so the same entry serves the
# 2D (MIP) and max-area-slice variants.
CHROM2D = {
    "i80_i20": (0, "I80/I20 ratio"),
    "hc_area_ec_area": (1, "heterochromatin/euchromatin area"),
    "hc_area_nuc_area": (2, "heterochromatin/nucleus area"),
    "hc_content_ec_content": (3, "HC/EC content"),
    "hc_content_dna_content": (4, "HC/DNA content"),
    "nhigh_nlow": (9, "N_high / N_low"),
}

# Marker <-> DNA co-distribution (new 2026-07-27). No _2d_ (MIP) variant by design.
CODIST = {
    "dna_corr": (0, "× DNA correlation"),
    "dna_dense_reg_rel_level": (1, "relative level in DNA-dense regions"),
    "dna_sparse_reg_rel_level": (2, "relative level in DNA-sparse regions"),
}

# Morphology titles (channel-independent). Insertion order = display order.
MORPH3D_TITLES = {
    # Shape/flatness first -- aspect ratio is the featured mechano readout.
    "morph3d_aspect_ratio": "Lateral-axial aspect ratio",
    "morph3d_z_height": "Nuclear height",
    "morph3d_sphericity": "Sphericity",
    "morph3d_nuclear_volume": "Nuclear volume",
    "morph3d_surface_area": "Surface area",
    "morph3d_convex_hull_vol": "Convex hull volume",
    "morph3d_equivalent_diameter": "Equivalent diameter (3D)",
    "morph3d_extent": "Nuclear extent (vol/bbox)",
    "morph3d_solidity": "Solidity (3D)",
    "morph3d_concavity_3d": "Concavity (3D)",
    "morph3d_major_axis_length": "Major axis length (3D)",
    "morph3d_minor_axis_length": "Minor axis length (3D)",
}
MORPH2D_TITLES = {
    # size / global shape
    "morph2d_area": "Area (2D)",
    "morph2d_perimeter": "Perimeter (2D)",
    "morph2d_convex_area": "Convex area (2D)",
    "morph2d_bbox_area": "Bounding-box area (2D)",
    "morph2d_area_bbarea": "Area / bbox area (2D)",
    "morph2d_equivalent_diameter": "Equivalent diameter (2D)",
    "morph2d_eccentricity": "Eccentricity (2D)",
    "morph2d_solidity": "Solidity (2D)",
    "morph2d_shape_factor": "Shape factor (2D)",
    "morph2d_a_r": "Aspect ratio (2D)",
    "morph2d_orientation": "Orientation (2D)",
    "morph2d_major_axis_length": "Major axis length (2D)",
    "morph2d_minor_axis_length": "Minor axis length (2D)",
    "morph2d_feret_max": "Max Feret (2D)",
    "morph2d_max_calliper": "Max calliper (2D)",
    "morph2d_min_calliper": "Min calliper (2D)",
    "morph2d_smallest_largest_calliper": "Smallest/largest calliper (2D)",
    # radius
    "morph2d_avg_radius": "Avg radius (2D)",
    "morph2d_med_radius": "Median radius (2D)",
    "morph2d_std_radius": "Radius SD (2D)",
    "morph2d_min_radius": "Min radius (2D)",
    "morph2d_max_radius": "Max radius (2D)",
    "morph2d_mode_radius": "Mode radius (2D)",
    "morph2d_d25_radius": "Radius d25 (2D)",
    "morph2d_d75_radius": "Radius d75 (2D)",
    # curvature
    "morph2d_avg_curvature": "Avg curvature (2D)",
    "morph2d_std_curvature": "Curvature SD (2D)",
    "morph2d_avg_posi_curv": "Avg positive curvature (2D)",
    "morph2d_avg_neg_curv": "Avg negative curvature (2D)",
    "morph2d_std_posi_curv": "Positive-curvature SD (2D)",
    "morph2d_std_neg_curv": "Negative-curvature SD (2D)",
    "morph2d_max_posi_curv": "Max positive curvature (2D)",
    "morph2d_max_neg_curv": "Max negative curvature (2D)",
    "morph2d_med_posi_curv": "Median positive curvature (2D)",
    "morph2d_med_neg_curv": "Median negative curvature (2D)",
    "morph2d_sum_posi_curv": "Sum positive curvature (2D)",
    "morph2d_sum_neg_curv": "Sum negative curvature (2D)",
    "morph2d_len_posi_curv": "Length of positive curvature (2D)",
    "morph2d_len_neg_curv": "Length of negative curvature (2D)",
    "morph2d_frac_peri_w_posi_curvature": "Fraction of perimeter w/ positive curvature (2D)",
    "morph2d_frac_peri_w_neg_curvature": "Fraction of perimeter w/ negative curvature (2D)",
    "morph2d_concavity": "Concavity (2D)",
    # prominence / polarity
    "morph2d_prominant_pos_curv": "Prominent positive curvature (2D)",
    "morph2d_prominant_neg_curv": "Prominent negative curvature (2D)",
    "morph2d_num_prominant_pos_curv": "Number of prominent positive-curvature points (2D)",
    "morph2d_num_prominant_neg_curv": "Number of prominent negative-curvature points (2D)",
    "morph2d_prominance_prominant_pos_curv": "Prominence of prominent positive curvature (2D)",
    "morph2d_prominance_prominant_neg_curv": "Prominence of prominent negative curvature (2D)",
    "morph2d_width_prominant_pos_curv": "Width of prominent positive curvature (2D)",
    "morph2d_width_prominant_neg_curv": "Width of prominent negative curvature (2D)",
    "morph2d_npolarity_changes": "Number of curvature sign changes (2D)",
    "morph2d_frac_peri_w_polarity_changes": "Fraction of perimeter w/ curvature sign changes (2D)",
}
MORPH_ORDER = {stem: i for i, stem in enumerate(
    list(MORPH3D_TITLES) + list(MORPH2D_TITLES))}

# Featured metrics: the `violin_selected/` set uses bare stems (no morph*/hoechst*
# prefix). Order here = display order; aspect ratio leads (the headline mechano
# readout), volume is the constant-volume control.
FEATURED_TITLES = {
    "aspect_ratio": "Lateral-axial aspect ratio",
    "volume": "Nuclear volume",
    "dna_skewness": "DNA intensity skewness",
    "dna_kurtosis": "DNA intensity excess kurtosis",
    "dna_hc_area_fraction": "DNA heterochromatin area fraction",
    "dna_hc_volume_fraction": "DNA heterochromatin volume fraction",
}
FEATURED_ORDER = {stem: i for i, stem in enumerate(FEATURED_TITLES)}
FEATURED_FAM = {
    "aspect_ratio": "Nuclear morphology (3D)",
    "volume": "Nuclear morphology (3D)",
    "dna_skewness": "Nuclear intensity",
    "dna_kurtosis": "Nuclear intensity",
    "dna_hc_area_fraction": "Chromatin organization",
    "dna_hc_volume_fraction": "Chromatin organization",
}


# Distribution-SHAPE moments read wrong with "intensity" appended ("excess
# kurtosis intensity"), so they are titled bare; level statistics keep it
# ("mean intensity").
SHAPE_MOMENTS = {"skewness", "kurtosis"}


def _int_title(ch_label, stat, dim_label):
    """'<channel> mean intensity (3D)' but '<channel> excess kurtosis (3D)'."""
    word = INT_STAT_TITLE.get(stat, stat)
    noun = "" if stat in SHAPE_MOMENTS else " intensity"
    return "{} {}{} ({})".format(ch_label, word, noun, dim_label)


# For a chromatin MARK channel, "heterochromatin/euchromatin" is wrong -- these
# thresholds are just high/low signal of that antibody (H3K27ac is an ACTIVE
# mark). Only the DNA channel gets the chromatin-state wording.
_HC_SWAPS = [
    ("relative heterochromatin volume", "high-signal volume fraction"),
    ("relative euchromatin volume", "low-signal volume fraction"),
    ("heterochromatin/euchromatin", "high-signal/low-signal"),
    ("heterochromatin/nucleus", "high-signal/nucleus"),
    ("HC/EC volume ratio", "high-signal/low-signal volume ratio"),
    ("HC/EC content", "high-signal/low-signal content"),
    ("HC/DNA content", "high-signal content / DNA content"),
]


def _hc_terms(ch_tok, text):
    if ch_tok == "hoechst":
        return text
    for a, b in _HC_SWAPS:
        text = text.replace(a, b)
    return text


def _ch(tok):
    # "hoechst" -> DNA; the mark tokens ("h3k27me3", or the generic "mark" used by
    # the cross-experiment violin_marks set) -> the dataset's MARK_LABEL (H3K27me3,
    # H3K9me3, H3K27ac, or "Chromatin mark" when the columns are 3 different marks).
    return "DNA" if tok == "hoechst" else MARK_LABEL


def _periph_dist(tok):
    """Parse a peripheral chromatin/enrichment distance token -> (display, sortkey).
    Handles pixel (``10``), micron (``0p5um``, ``1um``, ``2um``) and radial-percent
    (``r10pct``) forms. sortkey groups px < um < pct, then by value."""
    m = re.match(r"^(\d+)$", tok)
    if m:
        return ("{} px".format(tok), (0, int(tok)))
    m = re.match(r"^(\d+)(?:p(\d+))?um$", tok)
    if m:
        whole, frac = m.group(1), m.group(2)
        disp = ("{}.{}".format(whole, frac) if frac else whole) + " μm"
        num = float("{}.{}".format(whole, frac)) if frac else float(whole)
        return (disp, (1, num))
    m = re.match(r"^r(\d+)pct$", tok)
    if m:
        return ("r{}%".format(m.group(1)), (2, int(m.group(1))))
    return (tok, (3, 0))


def classify(stem):
    """Return a record dict or None. Record keys: fam, sort, channel, pair, title.
    `fam` is a FAMILIES_ORDER name; `sort` orders within the family; `pair` is a
    channel-agnostic key (records sharing it across DNA/H3K27me3 form a pair);
    `channel` in {'DNA','H3K27me3', None}."""

    # --- DNA-normalised readouts (checked FIRST) ---------------------------
    # The mark channel's absolute level is batch-sensitive, so the interpretable
    # version is normalised to DNA. Two spellings are supported:
    #   <ch>_dna_ratio             -> the overall mark/DNA ratio
    #   <base-metric><NORM_SUFFIX> -> any metric's DNA-normalised twin
    # Both sort IMMEDIATELY AFTER their unnormalised counterpart. This must run
    # before the specific rules below, or e.g. the peripheral-distance regex
    # would swallow the suffix as part of the shell token.
    m = re.match(r"^(hoechst|h3k27me3|mark)_dna_ratio$", stem)
    if m:
        ch = m.group(1)
        return dict(fam="Nuclear intensity",
                    sort=(0, INT_STAT_RANK["mean"], 1, stem),
                    channel=_ch(ch), pair=("dna_ratio",),
                    title="{} / DNA ratio".format(_ch(ch)))

    for suf in NORM_SUFFIXES:
        if stem.endswith(suf):
            base = classify(stem[: -len(suf)])
            if base:
                rec = dict(base)
                rec["title"] = base["title"] + ", per DNA"
                # same sort key as the base, nudged one place later
                rec["sort"] = tuple(base["sort"][:-1]) + (str(base["sort"][-1]) + "~",)
                rec["pair"] = ("norm",) + tuple(base["pair"])
                return rec

    # --- intensity: chan_<stat>_<ch>_3d_int ---
    m = re.match(r"^chan_(.+)_(hoechst|h3k27me3|mark)_3d_int$", stem)
    if m:
        stat, ch = m.group(1), m.group(2)
        return dict(fam="Nuclear intensity",
                    sort=(0, INT_STAT_RANK.get(stat, 50), 1, stem),
                    channel=_ch(ch), pair=("int_chan", stat),
                    title=_int_title(_ch(ch), stat, "3D"))

    # --- intensity: <ch>_3d_nuclear_<stat>_int ---
    m = re.match(r"^(hoechst|h3k27me3|mark)_3d_nuclear_(.+)_int$", stem)
    if m:
        ch, stat = m.group(1), m.group(2)
        return dict(fam="Nuclear intensity",
                    sort=(0, INT_STAT_RANK.get(stat, 50), 0, stem),
                    channel=_ch(ch), pair=("int_nuclear", stat),
                    title=_int_title(_ch(ch), stat, "3D"))

    # --- marker <-> DNA co-distribution: <ch>_(3d|2dmid)_dna_<measure> ---
    m = re.match(r"^(hoechst|h3k27me3|mark)_(3d|2dmid)_(dna_corr|dna_dense_reg_rel_level"
                 r"|dna_sparse_reg_rel_level)$", stem)
    if m:
        ch, dim, key = m.groups()
        sub, ttl = CODIST[key]
        return dict(fam="Mark–DNA co-distribution",
                    sort=(0 if dim == "3d" else 1, sub, stem),
                    channel=_ch(ch), pair=("codist", dim, key),
                    title="{} {} ({})".format(_ch(ch), ttl, DIM_LABEL[dim]))

    # --- intensity: <ch>_(2d|2dmid)_int_<stat> ---
    m = re.match(r"^(hoechst|h3k27me3|mark)_(2d|2dmid)_int_(.+)$", stem)
    if m:
        ch, dim, stat = m.groups()
        return dict(fam="Nuclear intensity",
                    sort=(1, INT_STAT_RANK.get(stat, 50), 0, stem),
                    channel=_ch(ch), pair=("int_" + dim, stat),
                    title=_int_title(_ch(ch), stat, DIM_LABEL[dim]))

    # --- intensity shape moments: <ch>_(2d|2dmid)_skewness | _kurtosis ---
    m = re.match(r"^(hoechst|h3k27me3|mark)_(2d|2dmid)_(skewness|kurtosis)$", stem)
    if m:
        ch, dim, stat = m.groups()
        return dict(fam="Nuclear intensity",
                    sort=(1, INT_STAT_RANK[stat], 0, stem),
                    channel=_ch(ch), pair=("int_" + dim, stat),
                    title=_int_title(_ch(ch), stat, DIM_LABEL[dim]))

    # --- GLCM: <ch>_(2d|2dmid)_<type>_<dist> ---
    m = re.match(r"^(hoechst|h3k27me3|mark)_(2d|2dmid)_"
                 r"(asm|contrast|correlation|dissimilarity|energy|homogeneity)_(\d+)$", stem)
    if m:
        ch, variant, gt, dist = m.groups()
        fam = "Texture (GLCM)" if variant == "2d" else "Texture (GLCM, max area slice)"
        suffix = "2D" if variant == "2d" else "max area slice"
        return dict(fam=fam,
                    sort=(GLCM_RANK[gt], int(dist), stem),
                    channel=_ch(ch), pair=("glcm", variant, gt, dist),
                    title="{} GLCM {} ({} px, {})".format(
                        _ch(ch), GLCM_TITLE[gt], dist, suffix))

    # --- texture: <ch>_(2d|2dmid)_entropy ---
    m = re.match(r"^(hoechst|h3k27me3|mark)_(2d|2dmid)_entropy$", stem)
    if m:
        ch, dim = m.group(1), m.group(2)
        fam = "Texture (GLCM)" if dim == "2d" else "Texture (GLCM, max area slice)"
        return dict(fam=fam, sort=(90, 0, stem),
                    channel=_ch(ch), pair=("entropy", dim),
                    title="{} entropy ({})".format(_ch(ch), DIM_LABEL[dim]))

    # --- chromatin: <ch>_3d_rdp_<n> ---
    m = re.match(r"^(hoechst|h3k27me3|mark)_3d_rdp_(\d+)$", stem)
    if m:
        ch, n = m.group(1), m.group(2)
        return dict(fam="Chromatin organization",
                    sort=(0, 0, int(n), stem), channel=_ch(ch), pair=("rdp", n),
                    title="{} radial density profile, shell {} (3D)".format(_ch(ch), n))

    # --- chromatin: <ch>_3d_rel_(hc|ec)_volume ---
    m = re.match(r"^(hoechst|h3k27me3|mark)_3d_rel_(hc|ec)_volume$", stem)
    if m:
        ch, which = m.group(1), m.group(2)
        full = "heterochromatin" if which == "hc" else "euchromatin"
        return dict(fam="Chromatin organization",
                    sort=(0, 1, 0 if which == "hc" else 1, stem),
                    channel=_ch(ch), pair=("rel_vol", which),
                    title=_hc_terms(ch, "{} relative {} volume".format(_ch(ch), full)))

    # --- chromatin: <ch>_3d_hc_ec_ratio_3d ---
    m = re.match(r"^(hoechst|h3k27me3|mark)_3d_hc_ec_ratio_3d$", stem)
    if m:
        ch = m.group(1)
        return dict(fam="Chromatin organization", sort=(0, 2, 0, stem),
                    channel=_ch(ch), pair=("hc_ec_ratio",),
                    title=_hc_terms(ch, "{} HC/EC volume ratio (3D)".format(_ch(ch))))

    # --- chromatin: <ch>_(2d|3d)_peripheral_(chromatin|enrichment)_<dist> ---
    # dist may be pixels (10), microns (0p5um/1um/2um) or radial-percent (r10pct).
    m = re.match(r"^(hoechst|h3k27me3|mark)_(2d|2dmid|3d)_peripheral_"
                 r"(chromatin|enrichment)_(.+)$", stem)
    if m:
        ch, dim, kind, tok = m.groups()
        disp, dsort = _periph_dist(tok)
        sub = 6 if kind == "chromatin" else 7
        return dict(fam="Chromatin organization",
                    sort=(0 if dim == "3d" else 1, sub, dsort, stem), channel=_ch(ch),
                    pair=("peripheral", dim, kind, tok),
                    title="{} peripheral {}, {} ({})".format(
                        _ch(ch), kind, disp, DIM_LABEL[dim]))

    # --- chromatin: named 2D ratios/contents ---
    m = re.match(r"^(hoechst|h3k27me3|mark)_(2d|2dmid)_(.+)$", stem)
    if m and m.group(3) in CHROM2D:
        ch, dim, key = m.groups()
        sub, ttl = CHROM2D[key]
        return dict(fam="Chromatin organization", sort=(1, sub, 0, stem),
                    channel=_ch(ch), pair=("chrom2d", dim, key),
                    title=_hc_terms(ch, "{} {} ({})".format(
                        _ch(ch), ttl, DIM_LABEL[dim])))

    # --- featured / pre-named metrics (violin_selected uses bare stems) ---
    if stem in FEATURED_TITLES:
        return dict(fam=FEATURED_FAM[stem], sort=(FEATURED_ORDER[stem], stem),
                    channel=None, pair=("featured", stem), title=FEATURED_TITLES[stem])

    # --- morphology (channel-independent) ---
    if stem in MORPH3D_TITLES:
        return dict(fam="Nuclear morphology (3D)", sort=(MORPH_ORDER[stem], stem),
                    channel=None, pair=("morph", stem), title=MORPH3D_TITLES[stem])
    if stem in MORPH2D_TITLES:
        return dict(fam="Nuclear morphology (2D)", sort=(MORPH_ORDER[stem], stem),
                    channel=None, pair=("morph", stem), title=MORPH2D_TITLES[stem])

    # --- equatorial-slice morphology (morph2dmid_*): same quantities as morph2d_*,
    # measured on the max-area optical slice rather than the MIP silhouette. Reuse
    # the 2D titles with a "(slice)" suffix so the pair is obvious.
    if stem.startswith("morph2dmid_"):
        base = "morph2d_" + stem[len("morph2dmid_"):]
        if base in MORPH2D_TITLES:
            ttl = MORPH2D_TITLES[base].replace("(2D)", "(max area slice)")
            if "(max area slice)" not in ttl:
                ttl += " (max area slice)"
            order = MORPH_ORDER.get(base, 999)
        else:
            ttl = (stem[len("morph2dmid_"):].replace("_", " ").capitalize()
                   + " (max area slice)")
            order = 999
        return dict(fam="Nuclear morphology (maximal area slice)", sort=(order, stem),
                    channel=None, pair=("morphmid", stem), title=ttl)

    if stem.startswith(("morph3d_", "morph2d_")):
        fam = "Nuclear morphology (3D)" if stem.startswith("morph3d_") else "Nuclear morphology (2D)"
        return dict(fam=fam, sort=(999, stem), channel=None, pair=("morph", stem),
                    title=stem.replace("morph2d_", "").replace("morph3d_", "").replace("_", " ").capitalize())

    return None


def fallback_title(stem):
    t = stem.replace("hoechst", "DNA").replace("h3k27me3", "H3K27me3")
    return t.replace("_", " ")


# ---------------------------------------------------------------------------
# Colors / layout
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.08  # tight margins to maximize plot area

TITLE_LEFT = MARGIN
TITLE_TOP = 0.03
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

# Single (wide) image: fills almost the whole slide below the title.
IMG_LEFT = MARGIN
IMG_TOP = 0.56
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = 6.76          # 0.56 -> 7.32

FOOTER_LEFT = MARGIN
FOOTER_TOP = 7.33
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.15
FOOTER_FONT_PT = 7

# Multi-image (N-up): thin caption row, images fill the rest.
COL_GAP = 0.12
MULTI_LABEL_TOP = 0.55
MULTI_LABEL_HEIGHT = 0.34
MULTI_IMG_TOP = 0.92
MULTI_IMG_H = 6.40        # 0.92 -> 7.32
MULTI_FOOTER_FONT_PT = 7


# ---------------------------------------------------------------------------
# Generic slide helpers
# ---------------------------------------------------------------------------
def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def build_slide(prs, title_text, image_path, footer_text):
    slide = _new_slide(prs)
    add_textbox(slide, title_text, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title_text), color=BLACK, bold=True)
    add_image_in_box(slide, str(image_path), IMG_LEFT, IMG_TOP, IMG_BOX_W, IMG_BOX_H)
    add_textbox(slide, footer_text, FOOTER_LEFT, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=BLACK)
    return slide


def build_multi_slide(prs, title, image_paths, labels, footers):
    n = len(image_paths)
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    col_w = (SLIDE_W - 2 * MARGIN - (n - 1) * COL_GAP) / n
    label_font = 22 if n <= 2 else 18
    for i in range(n):
        left = MARGIN + i * (col_w + COL_GAP)
        add_textbox(slide, labels[i], left, MULTI_LABEL_TOP, col_w, MULTI_LABEL_HEIGHT,
                    font_pt=label_font, color=BLACK, bold=True)
        add_image_in_box(slide, str(image_paths[i]), left, MULTI_IMG_TOP, col_w, MULTI_IMG_H)
    box = slide.shapes.add_textbox(
        Inches(FOOTER_LEFT), Inches(FOOTER_TOP), Inches(FOOTER_WIDTH), Inches(FOOTER_HEIGHT))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = footers[0]
    for idx, footer in enumerate(footers):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        if idx != 0:
            para.text = footer
        para.alignment = PP_ALIGN.CENTER
        para.runs[0].font.size = Pt(MULTI_FOOTER_FONT_PT)
        para.runs[0].font.color.rgb = BLACK
    return slide


def build_stack_slide(prs, title, image_paths, labels, footers):
    """Stack N images vertically (full width), each under a caption. Used for the
    all-conditions combined view: nucleus (DNA) on top, histone mark on bottom."""
    n = len(image_paths)
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    top_y = TITLE_TOP + TITLE_HEIGHT + 0.03
    bottom_y = SLIDE_H - 0.06
    row_h = (bottom_y - top_y) / n
    cap_h = 0.30
    for i in range(n):
        rt = top_y + i * row_h
        add_textbox(slide, labels[i], MARGIN, rt, SLIDE_W - 2 * MARGIN, cap_h,
                    font_pt=18, color=BLACK, bold=True)
        add_image_in_box(slide, str(image_paths[i]), MARGIN, rt + cap_h + 0.02,
                         SLIDE_W - 2 * MARGIN, row_h - cap_h - 0.08)
    return slide


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.7, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=40, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.1, SLIDE_W - 2 * MARGIN, 1.0,
                font_pt=18, color=BLACK, italic=True)


def build_divider_slide(prs, family_name):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, family_name, MARGIN, 3.1, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=44, color=BLACK, bold=True)


# ---------------------------------------------------------------------------
# Deck assembly
# ---------------------------------------------------------------------------
def rel_footer(path):
    try:
        return path.relative_to(ROOT).as_posix()
    except ValueError:
        return path.as_posix()


def _excluded(view, fname):
    return "{}/{}".format(view, fname) in EXCLUDE_PANELS


def _emit_multi(prs, pretty, kind, views, fname, missing, omitted, sep):
    kept = []
    for v, lab in views:
        if _excluded(v, fname):
            omitted.append("{}/{}".format(v, fname))
        elif not (GRID_DIR / v / fname).exists():
            missing.append("{}/{}".format(v, fname))
        else:
            kept.append((v, lab))
    if not kept:
        return
    paths = [GRID_DIR / v / fname for v, _ in kept]
    labels = [lab for _, lab in kept]
    title = "{} — {} ({})".format(pretty, kind, sep.join(labels))
    build_multi_slide(prs, title, paths, labels, [rel_footer(pp) for pp in paths])


def emit_view(prs, view_kind, stem, pretty, missing, omitted):
    fname = stem + GRID_SUFFIX
    if view_kind == "all":
        if _excluded(ALL_COND_VIEW, fname):
            omitted.append("{}/{}".format(ALL_COND_VIEW, fname))
            return
        p = GRID_DIR / ALL_COND_VIEW / fname
        if p.exists():
            build_slide(prs, "{} — all conditions".format(pretty), p, rel_footer(p))
        else:
            missing.append("{}/{}".format(ALL_COND_VIEW, fname))
    elif view_kind == "stiffness":
        _emit_multi(prs, pretty, "stiffness comparison", STIFFNESS_VIEWS, fname,
                    missing, omitted, " / ")
    elif view_kind == "timecourse":
        _emit_multi(prs, pretty, "timecourses", TIMECOURSE_VIEWS, fname,
                    missing, omitted, ", ")


def build_item(prs, item, missing, omitted):
    if item[0] == "pair":
        _, dna_stem, dna_title, h3_stem, h3_title = item
        for view in VIEW_ORDER:
            emit_view(prs, view, dna_stem, dna_title, missing, omitted)
            emit_view(prs, view, h3_stem, h3_title, missing, omitted)
    else:  # solo
        _, stem, title = item
        for view in VIEW_ORDER:
            emit_view(prs, view, stem, title, missing, omitted)


def build_item_allcond(prs, item, missing):
    """All-conditions-only variant: one slide per item. A pair STACKS the two
    all-conditions panels vertically -- nucleus (DNA) on top, histone mark on
    bottom -- under a shared metric title; a solo shows its single panel."""
    if item[0] == "pair":
        _, dna_stem, dna_title, h3_stem, h3_title = item
        base = dna_title[len("DNA "):] if dna_title.startswith("DNA ") else dna_title
        rows = []
        for stem, lab in ((dna_stem, "Nucleus (DNA)"), (h3_stem, MARK_LABEL)):
            fn = stem + GRID_SUFFIX
            p = GRID_DIR / ALL_COND_VIEW / fn
            if p.exists() and not _excluded(ALL_COND_VIEW, fn):
                rows.append((p, lab))
            else:
                missing.append("{}/{}".format(ALL_COND_VIEW, fn))
        if not rows:
            return
        title = "{} — all conditions".format(base)
        if len(rows) == 1:
            build_slide(prs, title, rows[0][0], rel_footer(rows[0][0]))
        else:
            build_stack_slide(prs, title, [r[0] for r in rows],
                              [r[1] for r in rows], [rel_footer(r[0]) for r in rows])
    else:  # solo
        _, stem, title = item
        fn = stem + GRID_SUFFIX
        p = GRID_DIR / ALL_COND_VIEW / fn
        if p.exists() and not _excluded(ALL_COND_VIEW, fn):
            build_slide(prs, "{} — all conditions".format(title), p, rel_footer(p))
        else:
            missing.append("{}/{}".format(ALL_COND_VIEW, fn))


def build_item_allcond_channel(prs, item, channel, missing):
    """Emit ONE full-size all-conditions slide for a single channel
    ('nucleus' or 'mark'). One plot per slide, filling the slide (build_slide).
    Channel-independent nuclear-morphology solos go to the nucleus deck only."""
    if item[0] == "pair":
        _, dna_stem, dna_title, h3_stem, h3_title = item
        stem, title = (dna_stem, dna_title) if channel == "nucleus" else (h3_stem, h3_title)
    else:  # solo -> nuclear morphology, nucleus deck only
        if channel != "nucleus":
            return
        _, stem, title = item
    fn = stem + GRID_SUFFIX
    p = GRID_DIR / ALL_COND_VIEW / fn
    if p.exists() and not _excluded(ALL_COND_VIEW, fn):
        build_slide(prs, "{} — all conditions".format(title), p, rel_footer(p))
    else:
        missing.append("{}/{}".format(ALL_COND_VIEW, fn))


def _item_stems(item):
    return [item[1], item[3]] if item[0] == "pair" else [item[1]]


def _item_log(item):
    return "{} | {}".format(item[2], item[4]) if item[0] == "pair" else item[2]


def discover_stems():
    """Union of metric stems across all views (so view-specific metrics are kept)."""
    stems = set()
    if not GRID_DIR.is_dir():
        return []
    for d in GRID_DIR.iterdir():
        if d.is_dir():
            for p in d.glob("*" + GRID_SUFFIX):
                stems.add(p.name[: -len(GRID_SUFFIX)])
    return sorted(stems)


def build_families():
    """Discover + classify all metrics into ordered (family_name, [items])."""
    stems = [s for s in discover_stems() if s not in DROP_METRICS]
    records, unknown = [], []
    for s in stems:
        rec = classify(s)
        if rec is None:
            unknown.append(s)
        else:
            rec["stem"] = s
            records.append(rec)

    # Group records into pair/solo items by (family, pair-key).
    groups = defaultdict(list)
    for r in records:
        groups[(r["fam"], r["pair"])].append(r)

    fam_items = defaultdict(list)  # fam -> list of (sort_key, item)
    for (fam, _pair), recs in groups.items():
        chans = {r["channel"] for r in recs}
        sort_key = min(r["sort"] for r in recs)
        if len(recs) == 2 and chans == {"DNA", MARK_LABEL}:
            dna = next(r for r in recs if r["channel"] == "DNA")
            h3 = next(r for r in recs if r["channel"] == MARK_LABEL)
            fam_items[fam].append((sort_key, ("pair", dna["stem"], dna["title"],
                                              h3["stem"], h3["title"])))
        else:
            for r in recs:  # solos (incl. any unpaired channel metric)
                fam_items[fam].append((r["sort"], ("solo", r["stem"], r["title"])))

    if unknown:
        fam_items["Other metrics"].extend(
            ((s,), ("solo", s, fallback_title(s))) for s in sorted(unknown))

    families = []
    for fam in FAMILIES_ORDER:
        if fam in fam_items:
            ordered = [it for _, it in sorted(fam_items[fam], key=lambda x: x[0])]
            families.append((fam, ordered))
    return families, unknown


def main():
    args = sys.argv[1:]
    list_only = "--list" in args
    allcond = "--allcond" in args  # all-conditions-only deck (DNA vs H3K side by side)
    ds_key = DEFAULT_DATASET
    if "--dataset" in args:
        i = args.index("--dataset")
        ds_key = args[i + 1] if i + 1 < len(args) else ds_key
    if ds_key not in DATASETS:
        print("Unknown dataset '{}'. Choose from: {}".format(ds_key, list(DATASETS)))
        sys.exit(2)
    apply_dataset(ds_key)
    print("Dataset: {}".format(ds_key))

    families, unknown = build_families()

    n_metrics = sum(len(_item_stems(it)) for _, items in families for it in items)
    n_items = sum(len(items) for _, items in families)
    n_pairs = sum(1 for _, items in families for it in items if it[0] == "pair")
    est_slides = 1 + len(families) + (n_items if allcond else n_metrics * 3)

    print("Source: {}".format(GRID_DIR))
    print("{} metrics ({} pairs), est. {} slides across {} families{}\n".format(
        n_metrics, n_pairs, est_slides, len(families),
        "  [all-conditions only]" if allcond else ""))

    if list_only:
        for fam, items in families:
            print("=== {} ({} items) ===".format(fam, len(items)))
            for it in items:
                print("  {}".format(_item_log(it)))
            print("")
        if unknown:
            print("UNKNOWN (-> Other): {}".format(unknown))
        return

    if n_metrics == 0:
        print("No metrics found under {}\n  -> data not compiled yet? Skipping build."
              .format(GRID_DIR))
        return

    if allcond:
        # Split the all-conditions comparison into SEPARATE full-size decks:
        # one nucleus (DNA) deck, one chromatin-mark deck. One plot per slide,
        # filling the slide. Nuclear-morphology solos go to the nucleus deck.
        for channel, clabel in (("nucleus", "Nucleus (DNA)"), ("mark", MARK_LABEL)):
            out_path = OUTPUT_PATH.with_name(
                OUTPUT_PATH.stem + "_all_conditions_" + channel + ".pptx")
            out_path.parent.mkdir(parents=True, exist_ok=True)
            prs = Presentation()
            prs.slide_width = Inches(SLIDE_W)
            prs.slide_height = Inches(SLIDE_H)
            build_title_slide(prs, DECK_TITLE,
                              DECK_SUBTITLE + "  -  all conditions, {} only".format(clabel))
            missing = []
            for fam, items in families:
                build_divider_slide(prs, fam)
                for it in items:
                    build_item_allcond_channel(prs, it, channel, missing)
            if out_path.exists():
                backup_presentation(str(out_path), backup_base=str(out_path.parent / "backups"))
            prs.save(str(out_path))
            print("[{} only] {} slides -> {}".format(
                clabel, len(prs.slides._sldIdLst), out_path.name))
        return

    # Full deck: both channels interleaved across the three views.
    out_path = OUTPUT_PATH
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)

    missing, omitted = [], []
    for fam, items in families:
        build_divider_slide(prs, fam)
        print("=== {} ===".format(fam))
        for it in items:
            build_item(prs, it, missing, omitted)
            print("  {}".format(_item_log(it)))
        print("")

    if out_path.exists():
        backup_presentation(str(out_path), backup_base=str(out_path.parent / "backups"))

    prs.save(str(out_path))
    total = len(prs.slides._sldIdLst)
    print("Done. {} metrics, {} slides written to:\n  {}".format(n_metrics, total, out_path))
    if unknown:
        print("\n{} unrecognized stem(s) under 'Other metrics': {}".format(len(unknown), unknown))
    if missing:
        print("\nSkipped {} missing panel(s) (not on disk).".format(len(missing)))


if __name__ == "__main__":
    main()
