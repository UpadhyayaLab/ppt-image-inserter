"""
insert_chromark_consistency_violins_slides.py

Build a deck from the hand-curated cross-experiment CONSISTENCY violin set
(J:/.../chromark_cross_experiment/ppt_violins_consistent_20260719/): 6 metrics
that change consistently with substrate stiffness across all three acquisitions
(H3K27me3 03/21, H3K9me3 04/24, H3K27ac 05/31 — the 3 columns in every panel).

Layout: title slide -> story slide -> per metric: a divider (metric + effect
direction + Cliff's delta) then its 6 view PNGs, one full-size plot per slide.
Reuses the maximized single-image layout from
insert_chromark_h3k27me3_summary_slides.py.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_chromark_consistency_violins_slides.py
"""

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import insert_chromark_h3k27me3_summary_slides as M  # noqa: E402

ROOT = Path(
    "J:/FF/fixed_cell/CTL_nucleus/tifsFixed3SIactivatedCTLs_nucleus/"
    "chromark_cross_experiment/ppt_violins_consistent_20260719"
)
OUTPUT = Path(
    "K:/FF/PPT/PPT_autogeneration/Naive_CTL/chromark/"
    "CTL_chromark_consistency_violins_summary.pptx"
)

DECK_TITLE = "CTL chromark — cross-experiment consistency set"
DECK_SUBTITLE = ("6 metrics that track substrate stiffness in ALL 3 datasets  -  "
                 "every panel: 3 columns = H3K27me3 (0321) | H3K9me3 (0424) | H3K27ac (0531)  -  "
                 "compiled 2026-07-19")

# (folder, slide title, effect note for the divider)
METRICS = [
    ("01_nuclear_area_2D", "Nuclear area (2D)",
     "morph2d_area — 2D footprint. LARGER on stiff glass (nucleus spreads).  "
     "Cliff's d = -0.84 (10 min) / -0.54 (3 h)"),
    ("02_nuclear_minoraxis_3D", "Nuclear minor axis (3D, ~ height)",
     "morph3d_minor_axis_length — shortest 3D axis. LARGER on soft gel "
     "(glass nuclei are flatter).  d = +0.87 / +0.69"),
    ("03_peripheral_chromatin", "Peripheral chromatin (DNA rim, outer 20%)",
     "hoechst_3d_peripheral_enrichment_r20pct — MORE rim-enriched on soft gel.  "
     "d = +0.88 / +0.61"),
    ("04_radial_DNA_rdp7", "Radial DNA profile (RDP, outer shell 7)",
     "hoechst_3d_rdp_7 — higher on stiff glass.  d = -0.88 / -0.76"),
    ("05_texture_contrast", "Chromatin texture (GLCM contrast, 1 px)",
     "hoechst_2d_contrast_1 — coarser / higher-contrast on soft gel.  d = +0.75 / +0.61"),
    ("06_envelope_curvature", "Nuclear envelope curvature (avg)",
     "morph2d_avg_curvature — more crenelated on soft gel.  d = +0.83 / +0.45"),
]

# (filename, view label)
VIEWS = [
    ("all_conditions.png", "all 8 conditions (overview)"),
    ("stiffness_10min.png", "stiffness @ 10 min (1.5 vs 12 kPa vs glass)"),
    ("stiffness_3hr.png", "stiffness @ 3 h (1.5 vs 12 kPa vs glass)"),
    ("timepoint_1p5kPa.png", "1.5 kPa gel: 10 min vs 3 h"),
    ("timepoint_12kPa.png", "12 kPa gel: 10 min vs 3 h"),
    ("timepoint_glass.png", "glass: 10 min vs 3 h"),
]

STORY = [
    "The reproducible story — holds in all three antibody datasets (H3K27me3 / H3K9me3 / H3K27ac):",
    "",
    "•  On STIFF glass the nucleus SPREADS and FLATTENS — larger 2D area, smaller 3D minor axis.",
    "•  Chromatin on stiff glass is LESS rim-enriched, FINER-textured, with a SMOOTHER envelope;",
    "    on SOFT gel the opposite (more peripheral DNA, coarser texture, more crenelated boundary).",
    "•  Effects are consistently STRONGER at 10 min than at 3 h (they relax over time).",
    "",
    "Deliberately excluded: 3D nuclear VOLUME — not reproducible (voxel-count volume is noisy across",
    "acquisitions). The size/shape story uses 2D area + 3D minor axis instead.",
]


def divider(prs, title, note):
    slide = M._new_slide(prs, bg=M.DIVIDER_BG)
    M.add_textbox(slide, title, M.MARGIN, 2.6, M.SLIDE_W - 2 * M.MARGIN, 1.1,
                  font_pt=40, color=M.BLACK, bold=True)
    M.add_textbox(slide, note, M.MARGIN, 3.85, M.SLIDE_W - 2 * M.MARGIN, 1.6,
                  font_pt=18, color=M.BLACK, italic=True)


def story_slide(prs, title, lines):
    slide = M._new_slide(prs)
    M.add_textbox(slide, title, M.MARGIN, 0.35, M.SLIDE_W - 2 * M.MARGIN, 0.7,
                  font_pt=30, color=M.BLACK, bold=True)
    box = slide.shapes.add_textbox(
        Inches(M.MARGIN + 0.2), Inches(1.4),
        Inches(M.SLIDE_W - 2 * M.MARGIN - 0.4), Inches(5.6))
    tf = box.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = ln
        para.alignment = PP_ALIGN.LEFT
        if para.runs:
            para.runs[0].font.size = Pt(18)
            para.runs[0].font.color.rgb = M.BLACK
            para.runs[0].font.bold = i == 0


def main():
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(M.SLIDE_W)
    prs.slide_height = Inches(M.SLIDE_H)

    M.build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)
    story_slide(prs, "Reading the consistency set", STORY)

    missing = []
    for folder, title, note in METRICS:
        divider(prs, title, note)
        print("=== {} ===".format(title))
        for vfile, vlabel in VIEWS:
            p = ROOT / folder / vfile
            if p.exists():
                M.build_slide(prs, "{} — {}".format(title, vlabel), p,
                              "{}/{}".format(folder, vfile))
                print("  {}".format(vlabel))
            else:
                missing.append("{}/{}".format(folder, vfile))
                print("  MISSING: {}".format(vfile))

    if OUTPUT.exists():
        M.backup_presentation(str(OUTPUT), backup_base=str(OUTPUT.parent / "backups"))
    prs.save(str(OUTPUT))
    total = len(prs.slides._sldIdLst)
    print("\nDone. {} metrics, {} slides -> {}".format(len(METRICS), total, OUTPUT))
    if missing:
        print("\nMissing {} PNG(s):".format(len(missing)))
        for m in missing:
            print("  - {}".format(m))


if __name__ == "__main__":
    main()
