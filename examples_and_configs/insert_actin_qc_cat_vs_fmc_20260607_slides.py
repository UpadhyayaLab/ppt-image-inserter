"""
insert_actin_qc_cat_vs_fmc_20260607_slides.py

CAT (left) vs FMC (right) side-by-side actin QC compare deck SCOPED TO
the 20260607 Kiet CART dataset only
(Y:/User_data/Kiet/20260607_pMLC_CART_actin_hoescht). Same layout as
insert_actin_qc_cat_vs_fmc_slides.py but trimmed to one dataset.

Result: 1 dataset x 3 timepoints x (2 synapse interleaved + 1 XZ MIP) =
9 slides.

Until pipeline results exist, every cell will render as a `(missing)`
placeholder.

Usage:
    python examples_and_configs/insert_actin_qc_cat_vs_fmc_20260607_slides.py
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/CART_actin_only/"
    "CART_actin_QC_CAT_vs_FMC_20260607.pptx"
)

KIET_ROOT = "Y:/User_data/Kiet"

# Single dataset scope.
DATASETS = [
    ("20260607", "20260607_pMLC_CART_actin_hoescht"),
]

TIMEPOINTS = ["5min", "10min", "15min"]

CONDITION_SUBPATH = "converted/cropped/split_channels"
PROG_SUBPATH = "prog_fixed_cells_actin_only/actin"

# Kind blocks. Block 1 (interleaved synapse pair: no-rings then with-rings)
# is followed by Block 2 (XZ MIP alone).
# Slide order: 6 synapse interleaved + 3 XZ MIP = 9 slides.
# Stage AG path mappings:
#   synapse/mask/                      -> synapse/1slice/mask/
#   synapse/inner_outer/bot_combined/  -> synapse/inner_outer/1slice_combined/
KIND_BLOCKS = [
    [
        ("Actin at Synapse",             "synapse/1slice/mask/montages"),
        ("Inner-Outer Ratio Definition", "synapse/inner_outer/1slice_combined/montages"),
    ],
    [
        ("Actin XZ MIP", "xz_mip/montages"),
    ],
]

CHUNK_GLOB = "montage_cells_*.png"

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Slide layout (inches). 13.333 x 7.5 widescreen.
SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

# 1x2 cell grid below the title (label + image per cell)
GRID_LEFT = 0.10
GRID_TOP = 0.60
CELL_W = 6.50
CELL_H = SLIDE_H - GRID_TOP - 0.10   # 6.80"
LABEL_H = 0.30
IMG_H = CELL_H - LABEL_H             # 6.50"
LABEL_FONT_PT = 16
COL_GAP = SLIDE_W - 2 * GRID_LEFT - 2 * CELL_W

CELL_POSITIONS = [
    (GRID_LEFT,                    GRID_TOP),
    (GRID_LEFT + CELL_W + COL_GAP, GRID_TOP),
]

# Scalebar invariant from the MATLAB CART_fixed_cell_analysis pipeline (Stage AF).
# Every per-cell tile is rendered at PPUM_SOURCE px/μm in its data area; the
# 5 μm scalebar is therefore SCALEBAR_PX pixels in every montage tile.
PPUM_SOURCE = 30          # px/μm in source PNGs
SCALEBAR_UM = 5           # μm
SCALEBAR_PX = PPUM_SOURCE * SCALEBAR_UM  # 150 px

# ---------------------------------------------------------------------------


def format_timepoint(tp: str) -> str:
    tp_map = {"5min": "5 min", "10min": "10 min", "15min": "15 min"}
    return tp_map.get(tp.lower(), tp)


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _png_dims(path: Path) -> Tuple[int, int]:
    """Return (width_px, height_px) of a PNG without fully decoding it."""
    with Image.open(str(path)) as im:
        return im.size


def compute_deck_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    """Smallest ppi such that every image fits in (max_w_in x max_h_in).
    Used to pin px/inch across the deck so the 5 μm scalebar lands at the
    same cm on every slide."""
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_image_in_cell_at_ppi(slide, image_path: Path, ppi: float,
                             cell_left: float, cell_top: float):
    """Place an image inside a labelled cell using a uniform deck px/inch.
    Image is centered in the (CELL_W x IMG_H) image area below the label;
    both dims = native_px / ppi inches. Pinning ppi across the deck keeps
    the embedded 5 μm scalebar at a constant cm on every slide."""
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    img_area_top = cell_top + LABEL_H
    left_in = cell_left + (CELL_W - w_in) / 2
    top_in  = img_area_top + (IMG_H - h_in) / 2
    return slide.shapes.add_picture(
        str(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(w_in),
    )


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _chunk_start_index(p: Path) -> int:
    m = re.match(r"montage_cells_(\d+)", p.name)
    return int(m.group(1)) if m else 0


def find_first_chunk(montages_dir: Path) -> Optional[Path]:
    if not montages_dir.is_dir():
        return None
    chunks = list(montages_dir.glob(CHUNK_GLOB))
    if not chunks:
        return None
    chunks.sort(key=_chunk_start_index)
    return chunks[0]


def build_compare_slide(prs, title_text, cat_img, fmc_img, deck_ppi):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    cells = [
        ("CAT", cat_img, CELL_POSITIONS[0]),
        ("FMC", fmc_img, CELL_POSITIONS[1]),
    ]
    missing = []
    for label, img_path, (cell_left, cell_top) in cells:
        add_textbox(
            slide, label,
            cell_left, cell_top, CELL_W, LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )
        if img_path is not None and img_path.exists():
            add_image_in_cell_at_ppi(slide, img_path, deck_ppi, cell_left, cell_top)
        else:
            add_textbox(
                slide, "(missing)",
                cell_left, cell_top + LABEL_H + IMG_H / 2 - 0.15, CELL_W, 0.3,
                font_pt=14, color=WHITE,
            )
            missing.append(label)
    return slide, missing


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    kiet_root = Path(KIET_ROOT)

    # Pre-pass: walk every (block, date, tp, kind) → collect (cat_img, fmc_img)
    # plus the title so we can pin deck-wide px/inch before slide build.
    slide_specs: List[Tuple[str, Optional[Path], Optional[Path], Path, Path, str]] = []
    for block in KIND_BLOCKS:
        for date_tag, dataset_name in DATASETS:
            for tp in TIMEPOINTS:
                tp_pretty = format_timepoint(tp)
                for kind_label, kind_subpath in block:
                    cat_dir = (
                        kiet_root / dataset_name / f"CAT_{tp}"
                        / CONDITION_SUBPATH / PROG_SUBPATH / kind_subpath
                    )
                    fmc_dir = (
                        kiet_root / dataset_name / f"FMC_{tp}"
                        / CONDITION_SUBPATH / PROG_SUBPATH / kind_subpath
                    )
                    cat_img = find_first_chunk(cat_dir)
                    fmc_img = find_first_chunk(fmc_dir)
                    title = f"{kind_label}: {tp_pretty} ({date_tag})"
                    log_key = f"{kind_label}/{date_tag}/{tp}"
                    slide_specs.append((title, cat_img, fmc_img, cat_dir, fmc_dir, log_key))

    present: List[Path] = []
    for (_, cat_img, fmc_img, _, _, _) in slide_specs:
        for p in (cat_img, fmc_img):
            if p is not None and p.exists():
                present.append(p)

    if not present:
        print("WARNING: no real images found — using fallback PPI=100.")
        deck_ppi = 100.0
    else:
        deck_ppi = compute_deck_ppi(present, CELL_W, IMG_H)

    bar_in = SCALEBAR_PX / deck_ppi
    print(
        f"Deck-wide PPI = {deck_ppi:.2f} (pinned across all {len(present)} present cells "
        f"in {len(slide_specs)} slides).\n"
        f"  Scalebar invariant: {SCALEBAR_UM} μm = {SCALEBAR_PX} px in source "
        f"=> {bar_in:.3f} in = {bar_in * 2.54:.3f} cm on every cell.\n"
        f"  Source PPUM = {PPUM_SOURCE} px/μm (locked).\n"
    )
    print(f"Writing deck to: {OUTPUT_PATH}\n")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    missing_total = []
    slides_added = 0
    for (title, cat_img, fmc_img, cat_dir, fmc_dir, log_key) in slide_specs:
        _, missing = build_compare_slide(prs, title, cat_img, fmc_img, deck_ppi)
        slides_added += 1

        status_parts = [
            "CAT:OK" if cat_img else "CAT:MISSING",
            "FMC:OK" if fmc_img else "FMC:MISSING",
        ]
        print(f"[{log_key}]  " + "  ".join(status_parts))

        for cell in missing:
            src = cat_dir if cell == "CAT" else fmc_dir
            missing_total.append(f"{log_key}/{cell}  ({src})")

    prs.save(str(out_path))
    print(f"\nDone. {slides_added} slides written to:\n  {out_path}")

    if missing_total:
        print(f"\nMissing ({len(missing_total)}):")
        for m in missing_total:
            print(f"  - {m}")
    else:
        print("\nAll images found - no missing items.")


if __name__ == "__main__":
    main()
