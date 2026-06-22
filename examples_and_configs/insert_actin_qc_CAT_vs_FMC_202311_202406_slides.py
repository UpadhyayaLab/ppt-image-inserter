"""
insert_actin_qc_CAT_vs_FMC_202311_202406_slides.py

CAT (left) vs FMC (right) side-by-side actin QC compare deck for the
J: drive CART experiments using the new Stage AG pipeline output
(`prog_fixed_cells_foci/actin/` with `synapse/1slice/...` layout).

Three experiments:
    20231127 — 1 timepoint (15 min)              -> 2 + 1 = 3 slides
    20240620 — 2 timepoints (5 min, 15 min)      -> 4 + 2 = 6 slides
    20240624 — 2 timepoints (5 min, 15 min)      -> 4 + 2 = 6 slides
Total: 15 slides (10 interleaved synapse + 5 XZ MIP).

Folder naming differs per dataset:
    20231127:  CAR=CAT  -> W3_NA_bCD19_CAT_{tp}_CTSB-mCh_647bTub_488Actin_Hoechst
               CAR=FMC63-> W1_NA_bCD19_FMC63_{tp}_CTSB-mCh_647bTub_488Actin_Hoechst
    20240620:  CAR=CAT  -> CAT{tp}      CAR=FMC63-> FMC{tp}
    20240624:  CAR=CAT  -> CAT_{tp}     CAR=FMC63-> FMC63_{tp}

Deck-wide PPI is pinned across all inserted montages so the 5 μm
scalebar lands at the same cm on every slide (Stage AF/AG closes the
within-tile invariant; this script closes the cross-slide one).

Supersedes insert_actin_qc_CAT_vs_FMC_202406_slides.py (which used the
older `prog_fixed_cells_actin_only/` outputs with old `synapse/mask/`
paths and constant-slide-width insertion).

Usage:
    python examples_and_configs/insert_actin_qc_CAT_vs_FMC_202311_202406_slides.py
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/CART/actin_only/"
    "CART_actin_QC_CAT_vs_FMC_202311_202406.pptx"
)

CTSB_ROOT = "J:/FF/fixed_cell/CAR_TCell"

# (date_tag, dataset_folder, cell_templates, timepoints)
# Each cell_templates entry: CAR_label -> (CAR_subdir, condition_format)
DATASETS = [
    (
        "20231127",
        "20231127_Fixed_CAR-Tcells_CTSB-mCherry_bTub",
        {
            "CAT": ("CAT",   "W3_NA_bCD19_CAT_{tp}_CTSB-mCh_647bTub_488Actin_Hoechst"),
            "FMC": ("FMC63", "W1_NA_bCD19_FMC63_{tp}_CTSB-mCh_647bTub_488Actin_Hoechst"),
        },
        ["15min"],
    ),
    (
        "20240620",
        "20240620_day3_Fixed_CAR-Tcells_CTSB-mCherry_bTub",
        {
            "CAT": ("CAT",   "CAT{tp}"),
            "FMC": ("FMC63", "FMC{tp}"),
        },
        ["5min", "15min"],
    ),
    (
        "20240624",
        "20240624_day5_Fixed_CAR-Tcells_CTSB-mCherry_bTub",
        {
            "CAT": ("CAT",   "CAT_{tp}"),
            "FMC": ("FMC63", "FMC63_{tp}"),
        },
        ["5min", "15min"],
    ),
]

CONDITION_SUBPATH = "cells/channels"
# Stage AG: outputs moved from prog_fixed_cells_actin_only/ -> prog_fixed_cells_foci/
PROG_SUBPATH = "prog_fixed_cells_foci/actin"

# Stage AG kind subpaths (synapse/mask -> synapse/1slice/mask, etc.)
KIND_BLOCKS = [
    [
        ("Actin at Synapse",             "synapse/1slice/mask/montages"),
        ("Inner-Outer Ratio Definition", "synapse/inner_outer/1slice_combined/montages"),
    ],
    [
        ("Actin XZ MIP", "xz_mip/montages"),
    ],
]

CHUNK_GLOB = "montage_cells_*.png"

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Slide layout (inches). 13.333 x 7.5 widescreen.
SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

# 1x2 cell grid below the title (label + image per cell)
GRID_LEFT = 0.10
GRID_TOP = 0.60
CELL_W = 6.50
CELL_H = SLIDE_H - GRID_TOP - 0.10   # 6.80"
LABEL_H = 0.30
IMG_H = CELL_H - LABEL_H             # 6.50"
LABEL_FONT_PT = 16
COL_GAP = SLIDE_W - 2 * GRID_LEFT - 2 * CELL_W

CELL_POSITIONS = [
    (GRID_LEFT,                    GRID_TOP),
    (GRID_LEFT + CELL_W + COL_GAP, GRID_TOP),
]

# Scalebar invariant from the MATLAB CART_fixed_cell_analysis pipeline
# (Stage AF + AG). Every per-cell tile is rendered at PPUM_SOURCE px/μm
# in its data area; the 5 μm scalebar is therefore SCALEBAR_PX pixels.
PPUM_SOURCE = 30          # px/μm in source PNGs
SCALEBAR_UM = 5           # μm
SCALEBAR_PX = PPUM_SOURCE * SCALEBAR_UM  # 150 px (nominal; measured ~156 due to LineWidth=4)

# ---------------------------------------------------------------------------


def format_timepoint(tp: str) -> str:
    tp_map = {"5min": "5 min", "15min": "15 min"}
    return tp_map.get(tp.lower(), tp)


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _parse_chunk_range(p: Path) -> Tuple[int, int]:
    """Return (start, end) cell-id range for a montage_cells_*.png filename.
    Handles both the 4-int FOV-padded pattern and the 2-int J:-drive pattern."""
    m4 = re.match(r"montage_cells_(\d+)_(\d+)_(\d+)_(\d+)\.png$", p.name)
    if m4:
        f_a, c_a, f_b, c_b = (int(x) for x in m4.groups())
        return (f_a * 1000 + c_a, f_b * 1000 + c_b)
    m2 = re.match(r"montage_cells_(\d+)_(\d+)\.png$", p.name)
    if m2:
        return (int(m2.group(1)), int(m2.group(2)))
    return (0, 0)


def find_first_chunk(montages_dir: Path) -> Optional[Path]:
    """Pick the lowest-start chunk, BUT first drop any chunk whose [start, end]
    range is strictly contained in another chunk's range. Catches leftover
    smoke chunks (e.g. `montage_cells_1_12.png` next to `montage_cells_1_26.png`)."""
    if not montages_dir.is_dir():
        return None
    chunks = list(montages_dir.glob(CHUNK_GLOB))
    if not chunks:
        return None
    parsed = [(p, *_parse_chunk_range(p)) for p in chunks]
    keep = []
    for (p, s, e) in parsed:
        shadowed = any(
            s2 <= s and e <= e2 and (s2 < s or e < e2)
            for (p2, s2, e2) in parsed
            if p2 is not p
        )
        if not shadowed:
            keep.append((p, s, e))
    if not keep:
        return None
    keep.sort(key=lambda x: x[1])
    return keep[0][0]


def _png_dims(path: Path) -> Tuple[int, int]:
    with Image.open(str(path)) as im:
        return im.size


def compute_deck_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    """Smallest ppi such that every image fits in (max_w_in x max_h_in).
    Pinning ppi across the deck keeps the embedded scalebar at a constant
    cm on every slide."""
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_image_in_cell_at_ppi(slide, image_path: Path, ppi: float,
                             cell_left: float, cell_top: float):
    """Place an image inside a labelled cell using a uniform deck px/inch.
    Both dims = native_px / ppi inches; image centered in (CELL_W x IMG_H)
    image area below the label."""
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    img_area_top = cell_top + LABEL_H
    left_in = cell_left + (CELL_W - w_in) / 2
    top_in  = img_area_top + (IMG_H - h_in) / 2
    return slide.shapes.add_picture(
        str(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(w_in),
    )


def build_compare_slide(prs, title_text, cat_img, fmc_img, deck_ppi):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    cells = [
        ("CAT", cat_img, CELL_POSITIONS[0]),
        ("FMC", fmc_img, CELL_POSITIONS[1]),
    ]
    missing = []
    for label, img_path, (cell_left, cell_top) in cells:
        add_textbox(
            slide, label,
            cell_left, cell_top, CELL_W, LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )
        if img_path is not None and img_path.exists():
            add_image_in_cell_at_ppi(slide, img_path, deck_ppi, cell_left, cell_top)
        else:
            add_textbox(
                slide, "(missing)",
                cell_left, cell_top + LABEL_H + IMG_H / 2 - 0.15, CELL_W, 0.3,
                font_pt=14, color=WHITE,
            )
            missing.append(label)
    return slide, missing


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    root = Path(CTSB_ROOT)

    # Pre-pass: walk every (block, dataset, tp, kind) -> resolve CAT/FMC
    # montage paths and pre-load their pixel dimensions to compute deck-wide
    # ppi BEFORE slide build.
    slide_specs: List[Tuple[str, Optional[Path], Optional[Path], Path, Path, str]] = []
    for block in KIND_BLOCKS:
        for (date_tag, dataset_folder, cell_templates, timepoints) in DATASETS:
            for tp in timepoints:
                tp_pretty = format_timepoint(tp)
                for kind_label, kind_subpath in block:
                    def resolve(cell):
                        car_sub, cond_template = cell_templates[cell]
                        cond_folder = cond_template.format(tp=tp)
                        return (
                            root / dataset_folder / car_sub / cond_folder
                            / CONDITION_SUBPATH / PROG_SUBPATH / kind_subpath
                        )
                    cat_dir = resolve("CAT")
                    fmc_dir = resolve("FMC")
                    cat_img = find_first_chunk(cat_dir)
                    fmc_img = find_first_chunk(fmc_dir)
                    title = f"{kind_label}: {tp_pretty} ({date_tag})"
                    log_key = f"{kind_label}/{date_tag}/{tp}"
                    slide_specs.append((title, cat_img, fmc_img, cat_dir, fmc_dir, log_key))

    present: List[Path] = []
    for (_, cat_img, fmc_img, _, _, _) in slide_specs:
        for p in (cat_img, fmc_img):
            if p is not None and p.exists():
                present.append(p)

    if not present:
        print("WARNING: no real images found - using fallback PPI=100.")
        deck_ppi = 100.0
    else:
        deck_ppi = compute_deck_ppi(present, CELL_W, IMG_H)

    bar_in = SCALEBAR_PX / deck_ppi
    print(
        f"Deck-wide PPI = {deck_ppi:.2f} (pinned across all {len(present)} present cells "
        f"in {len(slide_specs)} slides).\n"
        f"  Scalebar invariant: {SCALEBAR_UM} μm = {SCALEBAR_PX} px in source "
        f"=> {bar_in:.3f} in = {bar_in * 2.54:.3f} cm on every cell.\n"
        f"  Source PPUM = {PPUM_SOURCE} px/μm (locked).\n"
    )
    print(f"Writing deck to: {OUTPUT_PATH}\n")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    missing_total = []
    slides_added = 0
    for (title, cat_img, fmc_img, cat_dir, fmc_dir, log_key) in slide_specs:
        _, missing = build_compare_slide(prs, title, cat_img, fmc_img, deck_ppi)
        slides_added += 1

        status_parts = [
            "CAT:OK" if cat_img else "CAT:MISSING",
            "FMC:OK" if fmc_img else "FMC:MISSING",
        ]
        print(f"[{log_key}]  " + "  ".join(status_parts))

        for cell in missing:
            src = cat_dir if cell == "CAT" else fmc_dir
            missing_total.append(f"{log_key}/{cell}  ({src})")

    prs.save(str(out_path))
    print(f"\nDone. {slides_added} slides written to:\n  {out_path}")

    if missing_total:
        print(f"\nMissing ({len(missing_total)}):")
        for m in missing_total:
            print(f"  - {m}")
    else:
        print("\nAll images found - no missing items.")


if __name__ == "__main__":
    main()
