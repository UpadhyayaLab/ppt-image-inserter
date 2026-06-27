"""
insert_noco_MT_nuc_slides.py

Focused physical-scale deck of the nucleus + microtubule (β-tubulin) montages for
the 01/23/2024 nocodazole MT experiment (20240123_E6-1_Nocodazole_Vimentin, W2
DMSO vs W1 1 μM noco, β-tubulin / MT stain). Companion to
insert_noco_dmso_vs_noco_phys_scale_slides.py — that deck covers actin/cent/nucleus
across all experiments; the MT marker channel only exists in the MT experiment's
montages, so it gets its own deck here. Nocodazole depolymerizes microtubules, so
these are the direct DMSO-vs-noco readout.

Two montage groups, group-major (all of one, then the next):
  1. MT + Nuc (DNA), XZ MIP                       (MT_nuc_xz)
  2. MT + Nuc (DNA), deepest invagination slice   (MT_nuc)
DMSO (left) vs noco (right) montages are paired by chunk index — one slide per
pair, up to the smaller condition's chunk count; extra chunks on the longer side
(DMSO) are dropped. Each combo gets its own pinned pixels-per-inch (the XZ MIP and
the XY slice differ in size), so the embedded 5 μm scalebar (104 px) matches within
a group.

Usage:
    python examples_and_configs/insert_noco_MT_nuc_slides.py
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/Noco/"
    "Noco_Jurkats_MT_nuc_montages_20240123.pptx"
)

ROOT = (
    "M:/FF/FF_4TB_2_Backup_fullHD/Nucleus Project_2ndharddrive/"
    "Nucleus deformations analysis/Nucleus - Fixed Cell Data/"
    "20240123_E6-1_Nocodazole_Vimentin"
)
CHAN_SUB = "tif/cells/channels"
DATE_TAG = "01/23/2024"

# (combo_folder, title). Rendered group-major: all of one combo's paired slides,
# then the next.
COMBOS = [
    ("MT_nuc_xz", "MT + Nuc (DNA), XZ MIP"),
    ("MT_nuc",    "MT + Nuc (DNA), deepest invagination slice"),
]

# (condition_folder, display_label). DMSO (control) first, then noco.
CONDITIONS = [
    ("W2_aCD3_E6-1_EGFP-Cen2_DMSO_AF647bTub_535Actin_Hoechst",    "DMSO, αCD3"),
    ("W1_aCD3_E6-1_EGFP-Cen2_1uMNoco_AF647bTub_535Actin_Hoechst", "Noco 1 μM, αCD3"),
]

# Set to an int to cap chunks per condition (e.g. 1 = one representative
# montage each); None renders every chunk.
MAX_CHUNKS_PER_COND = None

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Slide layout (inches). 13.333 x 7.5 widescreen. Two columns: DMSO left,
# noco right (mirrors the main Noco phys-scale deck's side-by-side slides).
SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

GRID_LEFT = 0.10
GRID_TOP = 0.60
LABEL_H = 0.30
LABEL_FONT_PT = 16
CELL_W = 6.50
IMG_H = SLIDE_H - GRID_TOP - LABEL_H - 0.10      # 6.50"
COL_GAP = SLIDE_W - 2 * GRID_LEFT - 2 * CELL_W   # 0.133"
CELL_LEFTS = [GRID_LEFT, GRID_LEFT + CELL_W + COL_GAP]   # [DMSO, noco]

# Scalebar invariant for the Jurkat nucleus/actin fixed-cell physical-scale
# pipeline: 104 px = 5 μm (20.8 px/μm). Verify with
# examples_and_configs/check_scalebar_pixel_widths.py.
SCALEBAR_PX = 104
SCALEBAR_UM = 5
PPUM_SOURCE = SCALEBAR_PX / SCALEBAR_UM

# ---------------------------------------------------------------------------


def _winlong(p) -> str:
    """Win32-safe absolute path: prepend \\\\?\\ so paths past MAX_PATH (260)
    work in os.stat / os.listdir / Image.open / add_picture. The MT montages
    sit ~279 chars deep, past the limit."""
    s = os.path.abspath(str(p))
    if os.name == "nt" and not s.startswith("\\\\?\\"):
        s = "\\\\?\\" + s.replace("/", "\\")
    return s


def _exists_long(p) -> bool:
    return os.path.exists(_winlong(p))


def _chunk_range(p: Path) -> Tuple[int, int]:
    """(start, end) cell indices parsed from montage_cells_<a>_<b>.png."""
    m = re.match(r"montage_cells_(\d+)_(\d+)", p.name)
    return (int(m.group(1)), int(m.group(2))) if m else (0, 0)


def list_chunks(montages_dir) -> List[Path]:
    """Montage chunk PNGs in a dir, sorted by start index. Long-path-safe:
    pathlib's is_dir()/glob() silently return False/empty past MAX_PATH, so
    enumerate via os.listdir on the \\\\?\\-prefixed path. [] if absent."""
    d = _winlong(montages_dir)
    if not os.path.isdir(d):
        return []
    names = [f for f in os.listdir(d)
             if f.startswith("montage_cells_") and f.endswith(".png")]
    return sorted((Path(montages_dir) / n for n in names),
                  key=lambda p: _chunk_range(p)[0])


def montages_dir(cond_folder: str, combo: str) -> Path:
    return (Path(ROOT) / cond_folder / CHAN_SUB /
            "prog_fixed_cells" / "physical_scale_images" / combo / "montages")


def _png_dims(path: Path) -> Tuple[int, int]:
    with Image.open(_winlong(path)) as im:
        return im.size


def compute_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    """Smallest ppi such that every image fits in (max_w_in x max_h_in)."""
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_image_at_ppi(slide, image_path: Path, ppi: float,
                     area_left: float, area_top: float,
                     area_w: float, area_h: float):
    """Center an image at uniform PPI inside (left, top, w, h)."""
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    left_in = area_left + (area_w - w_in) / 2
    top_in  = area_top  + (area_h - h_in) / 2
    return slide.shapes.add_picture(
        _winlong(image_path), Inches(left_in), Inches(top_in), width=Inches(w_in),
    )


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    (dmso_folder, dmso_label), (noco_folder, noco_label) = CONDITIONS

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank_layout = prs.slide_layouts[6]
    print(f"Writing deck to: {OUTPUT_PATH}\n")

    total = 0
    for combo, combo_title in COMBOS:
        dmso_chunks = list_chunks(montages_dir(dmso_folder, combo))
        noco_chunks = list_chunks(montages_dir(noco_folder, combo))
        if MAX_CHUNKS_PER_COND is not None:
            dmso_chunks = dmso_chunks[:MAX_CHUNKS_PER_COND]
            noco_chunks = noco_chunks[:MAX_CHUNKS_PER_COND]

        # Pair by chunk index: DMSO chunk i (left) vs noco chunk i (right). The
        # chunks are different cells, so we pair positionally up to the smaller
        # count; extra chunks on the longer side (DMSO) are dropped.
        n_pairs = min(len(dmso_chunks), len(noco_chunks))
        if n_pairs == 0:
            print(f"[{combo}] no paired montages (DMSO={len(dmso_chunks)}, "
                  f"noco={len(noco_chunks)}) — skipped.")
            continue
        pairs = [(dmso_chunks[i], noco_chunks[i]) for i in range(n_pairs)]

        # Per-combo pinned PPI: the XZ MIP and the XY slice differ in size, so
        # each combo group shares its own scale (scalebar matches within group).
        all_imgs = [p for pair in pairs for p in pair]
        ppi = compute_ppi(all_imgs, CELL_W, IMG_H)
        bar_in = SCALEBAR_PX / ppi
        print(f"[{combo}] {n_pairs} paired slides (DMSO {len(dmso_chunks)} / noco "
              f"{len(noco_chunks)} chunks; dropped DMSO={len(dmso_chunks) - n_pairs}). "
              f"PPI={ppi:.2f}  scalebar={bar_in * 2.54:.3f} cm.")

        for dchunk, nchunk in pairs:
            slide = prs.slides.add_slide(blank_layout)
            set_slide_background(slide, BLACK)
            add_textbox(
                slide, f"{combo_title} ({DATE_TAG} αCD3)",
                TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
            )
            for label, ch, cell_left in (
                (dmso_label, dchunk, CELL_LEFTS[0]),
                (noco_label, nchunk, CELL_LEFTS[1]),
            ):
                a, b = _chunk_range(ch)
                add_textbox(
                    slide, f"{label} · cells {a}–{b}",
                    cell_left, GRID_TOP, CELL_W, LABEL_H,
                    font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
                )
                if _exists_long(ch):
                    add_image_at_ppi(slide, ch, ppi,
                                     cell_left, GRID_TOP + LABEL_H, CELL_W, IMG_H)
            total += 1

    if total == 0:
        print("No montages found — nothing written. Has the MT experiment "
              "finished processing into physical_scale_images?")
        return
    prs.save(str(out_path))
    print(f"\nDone. {total} slides written to:\n  {out_path}")


if __name__ == "__main__":
    main()
