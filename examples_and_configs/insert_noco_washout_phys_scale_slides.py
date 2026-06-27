"""
insert_noco_washout_phys_scale_slides.py

DMSO (left) vs 1 µM nocodazole (right) physical-scale montage deck for the
20240203 Jurkat nocodazole-WASHOUT time-course experiment
(Jurkats_NocoWashout_Vim_MG_20240203 — L:/FF/Nucleus_centrosome/
vimentin_nocodazole-washout_fixed/20240203_MG_NocoWashout). Companion to
insert_noco_dmso_vs_noco_phys_scale_slides.py; same physical-scale PPI pinning
(104 px = 5 µm) and long-path-safe enumeration.

Design: 2 conditions × 3 washout timepoints (minutes after noco washout):
  noco 1 µM:  W1 (0 min)  W2 (4 min)  W3 (15 min)
  DMSO:       W4 (0 min)  W5 (4 min)  W6 (15 min)
Channels: Vimentin / Actin / Centrin / Hoechst (nucleus). The physical_scale
montages are actin/cent/nucleus only (no vim combo), so the deck is marker-
agnostic and the DMSO-vs-noco columns carry the comparison.

Layout (Option A): one slide per (combo group, washout timepoint), DMSO left |
noco right (first montage chunk each). Group-major: all of one combo group
across timepoints (0 → 4 → 15 min), then the next group. A (combo, timepoint)
slide is emitted only when BOTH its DMSO and noco montages exist, so the deck
fills in as the pipeline finishes; pending pairs are logged. Per scale group
(xz / broad) one pinned PPI is shared so scalebars match within a group.

Usage:
    python examples_and_configs/insert_noco_washout_phys_scale_slides.py
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/Noco/"
    "Noco_Jurkats_washout_phys_scale_montages_20240203.pptx"
)

ROOT = ("L:/FF/Nucleus_centrosome/vimentin_nocodazole-washout_fixed/"
        "20240203_MG_NocoWashout")
CHAN_SUB = "channels"
DATE_TAG = "02/03/2024"

DMSO_LABEL = "DMSO, αCD3"
NOCO_LABEL = "Noco 1 μM, αCD3"

# (minutes, dmso_folder, noco_folder, label). One slide per timepoint per combo
# group, ordered by minutes (0 → 4 → 15) within each group. `minutes` is the
# sort key; `label` is what's shown on the slide (0 min = "pre-washout").
TIMEPOINTS = [
    (0,  "W4_E6-1_EGFPCen2_aCD3_DMSO_0m_647Vim_535Actin_Hoechst",
         "W1_E6-1_EGFPCen2_aCD3_1uM-noco_0m_647Vim_535Actin_Hoechst",  "pre-washout"),
    (4,  "W5_E6-1_EGFPCen2_aCD3_DMSO_4m_647Vim_535Actin_Hoechst",
         "W2_E6-1_EGFPCen2_aCD3_1uM-noco_4m_647Vim_535Actin_Hoechst",  "4 min after washout"),
    (15, "W6_E6-1_EGFPCen2_aCD3_DMSO_15m_647Vim_535Actin_Hoechst",
         "W3_E6-1_EGFPCen2_aCD3_1uM-noco_15m_647Vim_535Actin_Hoechst", "15 min after washout"),
]

# (combo_folder, title, scale_group). Group-major, same combo set / scale groups
# as the main Noco phys-scale deck. XZ uses the *_nolines variant (no dashed
# cell top/bottom markers).
COMBOS = [
    ("actin_nuc_xz_nolines",  "Actin + Nuc XZ MIP",               "xz"),
    ("actin_cent_xz_nolines", "Actin + Cent XZ MIP",              "xz"),
    ("nucleus_bz",            "Nuc (DNA), broadest slice",        "broad"),
    ("cent_nuc_bz",           "Cent + Nuc, broadest slice",       "broad"),
    ("cent_nuc",              "Cent + Nuc, deepest invag. slice", "broad"),
]

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Slide layout (inches). 13.333 x 7.5 widescreen. Two columns: DMSO left, noco
# right (mirrors the main Noco phys-scale deck's side-by-side slides).
SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

GRID_LEFT = 0.10
GRID_TOP = 0.60
LABEL_H = 0.30
LABEL_FONT_PT = 16
CELL_W = 6.50
IMG_H = SLIDE_H - GRID_TOP - LABEL_H - 0.10      # 6.50"
COL_GAP = SLIDE_W - 2 * GRID_LEFT - 2 * CELL_W   # 0.133"
CELL_LEFTS = [GRID_LEFT, GRID_LEFT + CELL_W + COL_GAP]   # [DMSO, noco]

# Scalebar invariant for the Jurkat nucleus/actin fixed-cell physical-scale
# pipeline: 104 px = 5 μm (20.8 px/μm). Verify with
# examples_and_configs/check_scalebar_pixel_widths.py.
SCALEBAR_PX = 104
SCALEBAR_UM = 5
PPUM_SOURCE = SCALEBAR_PX / SCALEBAR_UM

# ---------------------------------------------------------------------------


def _winlong(p) -> str:
    """Win32-safe absolute path: prepend \\\\?\\ so paths past MAX_PATH (260)
    work in os.stat / os.listdir / Image.open / add_picture."""
    s = os.path.abspath(str(p))
    if os.name == "nt" and not s.startswith("\\\\?\\"):
        s = "\\\\?\\" + s.replace("/", "\\")
    return s


def _exists_long(p) -> bool:
    return os.path.exists(_winlong(p))


def _chunk_start_index(p: Path) -> int:
    m = re.match(r"montage_cells_(\d+)", p.name)
    return int(m.group(1)) if m else 0


def list_chunks(montages_dir) -> List[Path]:
    """Montage chunk PNGs in a dir, sorted by start index. Long-path-safe:
    pathlib's is_dir()/glob() silently return False/empty past MAX_PATH, so
    enumerate via os.listdir on the \\\\?\\-prefixed path. [] if absent."""
    d = _winlong(montages_dir)
    if not os.path.isdir(d):
        return []
    names = [f for f in os.listdir(d)
             if f.startswith("montage_cells_") and f.endswith(".png")]
    return sorted((Path(montages_dir) / n for n in names), key=_chunk_start_index)


def montages_dir(cond_folder: str, combo: str) -> Path:
    return (Path(ROOT) / cond_folder / CHAN_SUB /
            "prog_fixed_cells" / "physical_scale_images" / combo / "montages")


def first_chunk(cond_folder: str, combo: str) -> Optional[Path]:
    chunks = list_chunks(montages_dir(cond_folder, combo))
    return chunks[0] if chunks else None


def _png_dims(path: Path) -> Tuple[int, int]:
    with Image.open(_winlong(path)) as im:
        return im.size


def compute_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    """Smallest ppi such that every image fits in (max_w_in x max_h_in)."""
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_image_at_ppi(slide, image_path: Path, ppi: float,
                     area_left: float, area_top: float,
                     area_w: float, area_h: float):
    """Center an image at uniform PPI inside (left, top, w, h)."""
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    left_in = area_left + (area_w - w_in) / 2
    top_in  = area_top  + (area_h - h_in) / 2
    return slide.shapes.add_picture(
        _winlong(image_path), Inches(left_in), Inches(top_in), width=Inches(w_in),
    )


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    # Build specs group-major (combo outer, timepoint inner, sorted by minutes).
    # A (combo, timepoint) is emitted only when BOTH DMSO and noco montages
    # exist; otherwise it's logged as pending so the deck fills in over time.
    specs: List[dict] = []
    pending: List[str] = []
    for combo, ctitle, sgroup in COMBOS:
        for mins, dmso_f, noco_f, tplabel in sorted(TIMEPOINTS, key=lambda t: t[0]):
            d_img = first_chunk(dmso_f, combo)
            n_img = first_chunk(noco_f, combo)
            if d_img is not None and n_img is not None:
                specs.append({"combo": combo, "title": ctitle, "scale_group": sgroup,
                              "label": tplabel, "d_img": d_img, "n_img": n_img})
            else:
                miss = [s for s, img in (("DMSO", d_img), ("noco", n_img)) if img is None]
                pending.append(f"{combo} @ {tplabel} — missing {', '.join(miss)}")

    # One pinned PPI per scale group (max over all its slides) so scalebars match
    # within a group across timepoints and combos.
    group_ppi: dict = {}
    for s in specs:
        ppi = compute_ppi([s["d_img"], s["n_img"]], CELL_W, IMG_H)
        group_ppi[s["scale_group"]] = max(group_ppi.get(s["scale_group"], 0.0), ppi)

    print(f"{len(specs)} complete DMSO|noco slides, {len(pending)} pending. "
          f"Scalebar invariant: {SCALEBAR_UM} μm = {SCALEBAR_PX} px.")
    for sg, ppi in sorted(group_ppi.items()):
        bar = SCALEBAR_PX / ppi
        print(f"  group {sg:<6s} PPI={ppi:>7.2f}  scalebar={bar * 2.54:.3f} cm")
    if pending:
        print("Pending (montages not generated yet):")
        for p in pending:
            print(f"  - {p}")

    if not specs:
        print("\nNo complete DMSO|noco pairs yet — nothing written. Rerun once the "
              "remaining conditions finish processing into physical_scale_images.")
        return

    print(f"\nWriting deck to: {OUTPUT_PATH}\n")
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank_layout = prs.slide_layouts[6]

    for s in specs:
        ppi = group_ppi[s["scale_group"]]
        slide = prs.slides.add_slide(blank_layout)
        set_slide_background(slide, BLACK)
        add_textbox(
            slide, f"{s['title']} · {s['label']} ({DATE_TAG})",
            TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
            font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
        )
        for label, img, cell_left in (
            (DMSO_LABEL, s["d_img"], CELL_LEFTS[0]),
            (NOCO_LABEL, s["n_img"], CELL_LEFTS[1]),
        ):
            add_textbox(
                slide, label,
                cell_left, GRID_TOP, CELL_W, LABEL_H,
                font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
            )
            add_image_at_ppi(slide, img, ppi, cell_left, GRID_TOP + LABEL_H, CELL_W, IMG_H)
        print(f"[{s['combo']} @ {s['label']}]  group={s['scale_group']}")

    prs.save(str(out_path))
    print(f"\nDone. {len(specs)} slides written to:\n  {out_path}")


if __name__ == "__main__":
    main()
