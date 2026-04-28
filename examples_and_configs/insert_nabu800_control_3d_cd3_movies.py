#!/usr/bin/env python
"""
Build a PowerPoint deck with two movie slides per cell.

For each cell number N present in all four source folders, this script creates:

1. A "panel_1x3" slide with xz on top and yz on bottom
2. A "panel_raw_cent" slide with xz on top and yz on bottom

The script creates a new presentation from scratch rather than modifying an
existing template. This keeps the core image-inserter package unchanged while
providing a repeatable workflow for this specific movie-layout task.
"""

from __future__ import annotations

import argparse
import re
import shutil
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional, Sequence, Set

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image


REPO_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO_ROOT))

from ppt_image_inserter import (
    backup_presentation,
    force_autoplay_in_pptx,
    remove_all_text_from_slide,
)

try:
    import imageio.v3 as iio
except ImportError:
    iio = None


DEFAULT_PANEL_1X3_XZ = Path(
    "F:/FF/nucleus_live_cell/jurkat_nucleus_centrosome/NaBu800 Experiments/"
    "Control_3D_CD3/all_cells_together/prog_live_cells/raw_with_meshes_movies/"
    "panel_1x3/xz_white_grey_blue"
)
DEFAULT_PANEL_1X3_YZ = Path(
    "F:/FF/nucleus_live_cell/jurkat_nucleus_centrosome/NaBu800 Experiments/"
    "Control_3D_CD3/all_cells_together/prog_live_cells/raw_with_meshes_movies/"
    "panel_1x3/yz_white_grey_blue"
)
DEFAULT_RAW_CENT_XZ = Path(
    "F:/FF/nucleus_live_cell/jurkat_nucleus_centrosome/NaBu800 Experiments/"
    "Control_3D_CD3/all_cells_together/prog_live_cells/raw_with_meshes_movies/"
    "panel_raw_cent/xz"
)
DEFAULT_RAW_CENT_YZ = Path(
    "F:/FF/nucleus_live_cell/jurkat_nucleus_centrosome/NaBu800 Experiments/"
    "Control_3D_CD3/all_cells_together/prog_live_cells/raw_with_meshes_movies/"
    "panel_raw_cent/yz"
)
SECOND_PANEL_1X3_XZ = Path(
    "F:/FF/nucleus_live_cell/jurkat_nucleus_centrosome/GFP-Centrin_SiR-DNA/"
    "Control/cells/all_cells_together/prog_live_cells/raw_with_meshes_movies/"
    "panel_1x3/xz_white_grey_blue"
)
SECOND_PANEL_1X3_YZ = Path(
    "F:/FF/nucleus_live_cell/jurkat_nucleus_centrosome/GFP-Centrin_SiR-DNA/"
    "Control/cells/all_cells_together/prog_live_cells/raw_with_meshes_movies/"
    "panel_1x3/yz_white_grey_blue"
)
SECOND_RAW_CENT_XZ = Path(
    "F:/FF/nucleus_live_cell/jurkat_nucleus_centrosome/GFP-Centrin_SiR-DNA/"
    "Control/cells/all_cells_together/prog_live_cells/raw_with_meshes_movies/"
    "panel_raw_cent/xz"
)
SECOND_RAW_CENT_YZ = Path(
    "F:/FF/nucleus_live_cell/jurkat_nucleus_centrosome/GFP-Centrin_SiR-DNA/"
    "Control/cells/all_cells_together/prog_live_cells/raw_with_meshes_movies/"
    "panel_raw_cent/yz"
)
DEFAULT_OUTPUT = Path("NaBu800_Control_3D_CD3_movies.pptx")
DEFAULT_POSTER_DIR = Path(__file__).resolve().parent / "_movie_posters_tmp"

CELL_PATTERN = re.compile(r"^Cell(?P<cell>\d+)_.*\.mp4$", re.IGNORECASE)

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
MOVIE_LEFT_IN = 0.7
MOVIE_BOX_WIDTH_IN = 11.9
TOP_MOVIE_TOP_IN = 0.78
BOTTOM_MOVIE_TOP_IN = 4.02
MOVIE_BOX_HEIGHT_IN = 2.62


@dataclass(frozen=True)
class CellMovieSet:
    """All four movie paths needed to build the two slides for one cell."""

    experiment_label: str
    cell_number: int
    panel_1x3_xz: Path
    panel_1x3_yz: Path
    raw_cent_xz: Path
    raw_cent_yz: Path


@dataclass(frozen=True)
class ExperimentSource:
    """Folder configuration for one experiment to add into the deck."""

    experiment_label: str
    panel_1x3_xz: Path
    panel_1x3_yz: Path
    raw_cent_xz: Path
    raw_cent_yz: Path


def parse_args() -> argparse.Namespace:
    """Parse command-line arguments."""
    parser = argparse.ArgumentParser(
        description=(
            "Create a PowerPoint deck with two movie slides per cell: "
            "panel_1x3 (xz over yz) and panel_raw_cent (xz over yz)."
        )
    )
    parser.add_argument(
        "--panel-1x3-xz",
        type=Path,
        default=DEFAULT_PANEL_1X3_XZ,
        help=f"Folder containing CellN panel_1x3 xz movies. Default: {DEFAULT_PANEL_1X3_XZ}",
    )
    parser.add_argument(
        "--panel-1x3-yz",
        type=Path,
        default=DEFAULT_PANEL_1X3_YZ,
        help=f"Folder containing CellN panel_1x3 yz movies. Default: {DEFAULT_PANEL_1X3_YZ}",
    )
    parser.add_argument(
        "--raw-cent-xz",
        type=Path,
        default=DEFAULT_RAW_CENT_XZ,
        help=f"Folder containing CellN panel_raw_cent xz movies. Default: {DEFAULT_RAW_CENT_XZ}",
    )
    parser.add_argument(
        "--raw-cent-yz",
        type=Path,
        default=DEFAULT_RAW_CENT_YZ,
        help=f"Folder containing CellN panel_raw_cent yz movies. Default: {DEFAULT_RAW_CENT_YZ}",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=DEFAULT_OUTPUT,
        help=f"Output PowerPoint path. Default: {DEFAULT_OUTPUT}",
    )
    parser.add_argument(
        "--cells",
        type=str,
        default="",
        help="Optional comma-separated list of cell numbers to include, e.g. 1,2,5,10",
    )
    return parser.parse_args()


def parse_requested_cells(cells_arg: str) -> Optional[Set[int]]:
    """Parse optional comma-separated cell numbers."""
    if not cells_arg.strip():
        return None

    requested: Set[int] = set()
    for token in cells_arg.split(","):
        token = token.strip()
        if not token:
            continue
        try:
            requested.add(int(token))
        except ValueError as exc:
            raise ValueError(f"Invalid cell number in --cells: {token}") from exc

    if not requested:
        raise ValueError("--cells was provided but no valid cell numbers were found")

    return requested


def validate_folder(folder: Path, label: str) -> None:
    """Ensure a required movie folder exists."""
    if not folder.exists():
        raise FileNotFoundError(f"{label} folder not found: {folder}")
    if not folder.is_dir():
        raise NotADirectoryError(f"{label} is not a folder: {folder}")


def collect_movies(folder: Path, label: str) -> Dict[int, Path]:
    """Collect CellN movie files from a folder."""
    validate_folder(folder, label)

    movies: Dict[int, Path] = {}
    for movie_path in sorted(folder.glob("*.mp4")):
        match = CELL_PATTERN.match(movie_path.name)
        if not match:
            print(f"[WARNING] Skipping unmatched filename in {label}: {movie_path.name}")
            continue

        cell_number = int(match.group("cell"))
        if cell_number in movies:
            raise ValueError(
                f"Duplicate Cell{cell_number} movie found in {label}: "
                f"{movies[cell_number].name} and {movie_path.name}"
            )
        movies[cell_number] = movie_path

    if not movies:
        raise ValueError(f"No CellN .mp4 files found in {label}: {folder}")

    return movies


def build_complete_cell_sets(
    experiment_label: str,
    panel_1x3_xz: Dict[int, Path],
    panel_1x3_yz: Dict[int, Path],
    raw_cent_xz: Dict[int, Path],
    raw_cent_yz: Dict[int, Path],
    requested_cells: Optional[Set[int]] = None,
) -> List[CellMovieSet]:
    """Build complete four-movie records for cells present in all required folders."""
    available_in_all = (
        set(panel_1x3_xz)
        & set(panel_1x3_yz)
        & set(raw_cent_xz)
        & set(raw_cent_yz)
    )

    if requested_cells is not None:
        available_in_all &= requested_cells

    cell_sets = [
        CellMovieSet(
            experiment_label=experiment_label,
            cell_number=cell_number,
            panel_1x3_xz=panel_1x3_xz[cell_number],
            panel_1x3_yz=panel_1x3_yz[cell_number],
            raw_cent_xz=raw_cent_xz[cell_number],
            raw_cent_yz=raw_cent_yz[cell_number],
        )
        for cell_number in sorted(available_in_all)
    ]

    return cell_sets


def report_missing_cells(
    experiment_label: str,
    label_to_movies: Dict[str, Dict[int, Path]],
    requested_cells: Optional[Set[int]] = None,
) -> None:
    """Print warnings for cells missing one or more required movie files."""
    all_cells = set()
    for movies in label_to_movies.values():
        all_cells.update(movies)

    if requested_cells is not None:
        all_cells &= requested_cells

    for cell_number in sorted(all_cells):
        missing_from = [
            label for label, movies in label_to_movies.items() if cell_number not in movies
        ]
        if missing_from:
            missing_text = ", ".join(missing_from)
            print(
                f"[WARNING] {experiment_label}: Cell{cell_number} skipped; "
                f"missing movie in: {missing_text}"
            )

    if requested_cells is not None:
        undiscovered = requested_cells - all_cells
        for cell_number in sorted(undiscovered):
            print(
                f"[WARNING] {experiment_label}: Cell{cell_number} skipped; "
                "not found in any source folder"
            )


def collect_cell_sets_for_experiment(
    source: ExperimentSource,
    requested_cells: Optional[Set[int]] = None,
) -> List[CellMovieSet]:
    """Collect complete cell movie sets for one experiment."""
    folder_maps = {
        "panel_1x3_xz": collect_movies(source.panel_1x3_xz, f"{source.experiment_label} panel_1x3_xz"),
        "panel_1x3_yz": collect_movies(source.panel_1x3_yz, f"{source.experiment_label} panel_1x3_yz"),
        "raw_cent_xz": collect_movies(source.raw_cent_xz, f"{source.experiment_label} raw_cent_xz"),
        "raw_cent_yz": collect_movies(source.raw_cent_yz, f"{source.experiment_label} raw_cent_yz"),
    }

    report_missing_cells(
        source.experiment_label,
        folder_maps,
        requested_cells=requested_cells,
    )

    return build_complete_cell_sets(
        source.experiment_label,
        folder_maps["panel_1x3_xz"],
        folder_maps["panel_1x3_yz"],
        folder_maps["raw_cent_xz"],
        folder_maps["raw_cent_yz"],
        requested_cells=requested_cells,
    )


def get_blank_layout(prs: Presentation):
    """Return a blank slide layout when available, otherwise fall back to the last layout."""
    for layout in prs.slide_layouts:
        if len(layout.placeholders) == 0:
            return layout
    return prs.slide_layouts[len(prs.slide_layouts) - 1]


def add_title(slide, text: str) -> None:
    """Add a slide title."""
    title_box = slide.shapes.add_textbox(
        Inches(0.45), Inches(0.12), Inches(12.4), Inches(0.35)
    )
    text_frame = title_box.text_frame
    text_frame.text = text
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(22)
    paragraph.font.bold = True
    paragraph.font.name = "Arial"


def add_region_label(slide, text: str, top_inches: float) -> None:
    """Add a small label above a movie."""
    label_box = slide.shapes.add_textbox(
        Inches(MOVIE_LEFT_IN),
        Inches(top_inches),
        Inches(2.0),
        Inches(0.18),
    )
    text_frame = label_box.text_frame
    text_frame.text = text
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(11)
    paragraph.font.bold = True
    paragraph.font.name = "Arial"


def extract_first_frame(movie_path: Path, poster_path: Path) -> tuple[int, int]:
    """Extract the first frame from a movie, save it, and return width/height."""
    if iio is None:
        raise RuntimeError(
            "imageio is required to extract poster frames. "
            "Install imageio and imageio-ffmpeg in the active environment."
        )

    frame = iio.imread(movie_path, index=0)
    poster = Image.fromarray(frame)
    poster.save(poster_path)
    return frame.shape[1], frame.shape[0]


def fit_within_box(
    content_width: int,
    content_height: int,
    box_left: float,
    box_top: float,
    box_width: float,
    box_height: float,
) -> tuple[float, float, float, float]:
    """Fit content inside a box while preserving aspect ratio."""
    content_aspect = content_width / content_height
    box_aspect = box_width / box_height

    if content_aspect >= box_aspect:
        final_width = box_width
        final_height = box_width / content_aspect
    else:
        final_height = box_height
        final_width = box_height * content_aspect

    final_left = box_left + (box_width - final_width) / 2
    final_top = box_top + (box_height - final_height) / 2
    return final_left, final_top, final_width, final_height


def add_movie(
    slide,
    movie_path: Path,
    poster_path: Path,
    frame_width: int,
    frame_height: int,
    box_top_inches: float,
) -> None:
    """Embed one movie onto the slide."""
    if not hasattr(slide.shapes, "add_movie"):
        raise RuntimeError(
            "This python-pptx installation does not support add_movie(). "
            "Upgrade python-pptx to a version that includes movie embedding."
        )

    left_inches, top_inches, width_inches, height_inches = fit_within_box(
        frame_width,
        frame_height,
        MOVIE_LEFT_IN,
        box_top_inches,
        MOVIE_BOX_WIDTH_IN,
        MOVIE_BOX_HEIGHT_IN,
    )

    slide.shapes.add_movie(
        str(movie_path),
        Inches(left_inches),
        Inches(top_inches),
        Inches(width_inches),
        Inches(height_inches),
        poster_frame_image=str(poster_path),
        mime_type="video/mp4",
    )


def add_movie_slide(
    prs: Presentation,
    title: str,
    top_movie: Path,
    bottom_movie: Path,
    poster_dir: Path,
) -> None:
    """Create one slide with xz on top and yz on bottom."""
    slide = prs.slides.add_slide(get_blank_layout(prs))
    remove_all_text_from_slide(slide)
    top_poster = poster_dir / f"{top_movie.stem}_poster.png"
    bottom_poster = poster_dir / f"{bottom_movie.stem}_poster.png"

    top_frame_width, top_frame_height = extract_first_frame(top_movie, top_poster)
    bottom_frame_width, bottom_frame_height = extract_first_frame(bottom_movie, bottom_poster)

    add_title(slide, title)
    add_region_label(slide, "xz", 0.55)
    add_movie(
        slide,
        top_movie,
        top_poster,
        top_frame_width,
        top_frame_height,
        TOP_MOVIE_TOP_IN,
    )
    add_region_label(slide, "yz", 3.79)
    add_movie(
        slide,
        bottom_movie,
        bottom_poster,
        bottom_frame_width,
        bottom_frame_height,
        BOTTOM_MOVIE_TOP_IN,
    )


def format_panel_1x3_title(cell_number: int, experiment_label: str) -> str:
    """Return the title for the panel_1x3 slide."""
    return f"Raw, invag depth, min curvature: Cell {cell_number} ({experiment_label})"


def format_raw_cent_title(cell_number: int, experiment_label: str) -> str:
    """Return the title for the panel_raw_cent slide."""
    return f"Raw, region near centrosome: Cell {cell_number} ({experiment_label})"


def create_presentation(cell_sets: Sequence[CellMovieSet], poster_dir: Path) -> Presentation:
    """Build the presentation in memory."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    for cell_set in cell_sets:
        add_movie_slide(
            prs,
            format_panel_1x3_title(cell_set.cell_number, cell_set.experiment_label),
            cell_set.panel_1x3_xz,
            cell_set.panel_1x3_yz,
            poster_dir,
        )
        add_movie_slide(
            prs,
            format_raw_cent_title(cell_set.cell_number, cell_set.experiment_label),
            cell_set.raw_cent_xz,
            cell_set.raw_cent_yz,
            poster_dir,
        )

    return prs


def ensure_parent_dir(output_path: Path) -> None:
    """Create the output directory if needed."""
    if output_path.parent and not output_path.parent.exists():
        output_path.parent.mkdir(parents=True, exist_ok=True)


def prepare_poster_dir(poster_dir: Path) -> Path:
    """Create a clean poster-cache directory."""
    if poster_dir.exists():
        shutil.rmtree(poster_dir, ignore_errors=True)
    poster_dir.mkdir(parents=True, exist_ok=True)
    return poster_dir


def main() -> int:
    """Run the movie deck generation workflow."""
    args = parse_args()
    experiment_sources = [
        ExperimentSource(
            experiment_label="032022 expt",
            panel_1x3_xz=args.panel_1x3_xz,
            panel_1x3_yz=args.panel_1x3_yz,
            raw_cent_xz=args.raw_cent_xz,
            raw_cent_yz=args.raw_cent_yz,
        ),
        ExperimentSource(
            experiment_label="04142022 expt",
            panel_1x3_xz=SECOND_PANEL_1X3_XZ,
            panel_1x3_yz=SECOND_PANEL_1X3_YZ,
            raw_cent_xz=SECOND_RAW_CENT_XZ,
            raw_cent_yz=SECOND_RAW_CENT_YZ,
        ),
    ]

    try:
        requested_cells = parse_requested_cells(args.cells)
    except ValueError as exc:
        print(f"[ERROR] {exc}")
        return 1

    try:
        cell_sets: List[CellMovieSet] = []
        for source in experiment_sources:
            experiment_cell_sets = collect_cell_sets_for_experiment(
                source,
                requested_cells=requested_cells,
            )
            print(f"{source.experiment_label}: found {len(experiment_cell_sets)} complete cell set(s)")
            cell_sets.extend(experiment_cell_sets)
    except Exception as exc:
        print(f"[ERROR] {exc}")
        return 1

    if not cell_sets:
        print("[ERROR] No cells have a complete set of all four required movies")
        return 1

    print(f"Found {len(cell_sets)} complete cell set(s) across all experiments")
    print(f"Will create {len(cell_sets) * 2} slide(s)")

    output_path = args.output.resolve()
    ensure_parent_dir(output_path)

    if output_path.exists():
        backup_dir = output_path.parent / "backups"
        try:
            backups_created = backup_presentation(str(output_path), backup_base=str(backup_dir))
            print(f"Backed up existing output to: {', '.join(sorted(backups_created))}")
        except Exception as exc:
            print(f"[WARNING] Could not back up existing output: {exc}")

    poster_dir = prepare_poster_dir(DEFAULT_POSTER_DIR)

    try:
        prs = create_presentation(cell_sets, poster_dir)
    except Exception as exc:
        print(f"[ERROR] Failed while building slides: {exc}")
        return 1

    try:
        prs.save(str(output_path))
    except PermissionError:
        print(
            f"[ERROR] Permission denied when saving {output_path}. "
            "Make sure the PowerPoint file is closed."
        )
        return 1
    except Exception as exc:
        print(f"[ERROR] Failed to save presentation: {exc}")
        return 1

    try:
        rewritten = force_autoplay_in_pptx(str(output_path))
        print(f"Set {rewritten} slide(s) to autoplay")
    except Exception as exc:
        print(f"[WARNING] Could not set autoplay timing: {exc}")

    print(f"[SUCCESS] Saved presentation: {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
