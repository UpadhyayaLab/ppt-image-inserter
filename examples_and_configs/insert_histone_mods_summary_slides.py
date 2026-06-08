"""
insert_histone_mods_summary_slides.py

Build summary decks for the Jurkat histone-modification / nuclear-correlation
analyses. Each slide shows one TIFF/TIF plot with:
  - bold title (explicit per slide)
  - the plot, full-width, aspect ratio preserved
  - the full source filepath in a small footer (for traceability)

Each DATASET produces one .pptx. To add or remove slides, edit the relevant
SLIDES_* list (entries are 3-tuples: (subfolder, filename, title)). Each run
regenerates every deck (idempotent); previous decks are backed up first.

Modeled after insert_actin_summary_slides.py.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_histone_mods_summary_slides.py
"""

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation  # noqa: E402

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

# Slide entries: (subfolder_relative_to_compiled_root, filename, title).
# Joined with each dataset's compiled_root at deck-build time. Missing files
# render "(missing)" so the deck still builds.
#
# Sentinel subfolder SUB_CROSSCORR_PAIR builds a 2-image slide with the same
# filename pulled from BOTH crosscorr_profiles/XY_MIP/ and
# crosscorr_profiles/nuc_broadest_slice/ side-by-side.

SUB_CROSSCORR_PAIR = "crosscorr_pair"

# Deck 1 — Jurkats H3K27me3 + PLL/aCD3 7min
SLIDES_H3K27ME3 = [
    ("grid_panels",
        "H3K27me3_Hoechst_corr_grid.tif",
        "Pearson Correlation: DNA and H3K27me3"),

    (SUB_CROSSCORR_PAIR,
        "H3K27me3_autocorr_pooled.png",
        "H3K27me3 autocorrelation (pooled)"),
    (SUB_CROSSCORR_PAIR,
        "H3K27me3_Hoechst_crosscorr_pooled.png",
        "H3K27me3 × DNA cross-correlation (pooled)"),
]

# Deck 2 — Jurkats CD3-vs-PLL HistoneMods (H3K27ac, H3K9me3)
SLIDES_HISTONEMODS = [
    ("grid_panels/H3K27ac",
        "struct_in_nuc_Hoechst_corr.tiff",
        "Pearson Correlation: DNA and H3K27ac"),
    (SUB_CROSSCORR_PAIR,
        "H3K27ac_autocorr.png",
        "H3K27ac autocorrelation"),
    (SUB_CROSSCORR_PAIR,
        "H3K27ac_Hoechst_crosscorr.png",
        "H3K27ac × DNA cross-correlation"),

    ("grid_panels/H3K9me3",
        "struct_in_nuc_Hoechst_corr.tiff",
        "Pearson Correlation: DNA and H3K9me3"),
    (SUB_CROSSCORR_PAIR,
        "H3K9me3_autocorr.png",
        "H3K9me3 autocorrelation"),
    (SUB_CROSSCORR_PAIR,
        "H3K9me3_Hoechst_crosscorr.png",
        "H3K9me3 × DNA cross-correlation"),
]

# Deck 3 — Jurkats RNAPII
SLIDES_RNAPII = [
    ("grid_panels/aCD3",
        "RNAPII_Hoechst_corr.tiff",
        "Pearson Correlation: DNA and RNAP II"),

    (SUB_CROSSCORR_PAIR,
        "RNAPII_autocorr.png",
        "RNAP II autocorrelation"),
    (SUB_CROSSCORR_PAIR,
        "RNAPII_Hoechst_crosscorr.png",
        "RNAP II × DNA cross-correlation"),
]

# Each dataset -> one deck. compiled_root points at the experiment folder
# (parent of crosscorr_profiles/ and grid_panels/).
OUTPUT_DIR = Path(
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/Histone Mods_RNAPII"
)

DATASETS = [
    {
        "label": "Jurkats H3K27me3 PLL/aCD3 7min",
        "compiled_root": Path(
            "H:/FF/Nucleus_Data/3D_Nucleus/Fixed/results_compilation/"
            "Jurkats_H3K27me3_PLL_aCD3_7min_20260606"
        ),
        "output_path": OUTPUT_DIR / "Jurkats_H3K27me3_PLL_aCD3_7min_corr_summary.pptx",
        "slides": SLIDES_H3K27ME3,
    },
    {
        "label": "Jurkats HistoneMods CD3vsPLL",
        "compiled_root": Path(
            "J:/FF/fixed_cell/Jurkat_nucleus/histone_mods/results_compilation/"
            "Jurkats_CD3vsPLL_HistoneMods_20240617_20260606"
        ),
        "output_path": OUTPUT_DIR / "Jurkats_HistoneMods_CD3vsPLL_corr_summary.pptx",
        "slides": SLIDES_HISTONEMODS,
    },
    {
        "label": "Jurkats RNAPII",
        "compiled_root": Path(
            "J:/FF/fixed_cell/RNAP_II/compiled_results/"
            "Jurkats_RNAPII_20251117_20260606"
        ),
        "output_path": OUTPUT_DIR / "Jurkats_RNAPII_corr_summary.pptx",
        "slides": SLIDES_RNAPII,
    },
]

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Slide layout (inches). 13.333 x 7.5 widescreen.
SLIDE_W = 13.333
SLIDE_H = 7.5

MARGIN = 0.10

TITLE_LEFT = MARGIN
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

SUBTITLE_LEFT = MARGIN
SUBTITLE_TOP = 0.55
SUBTITLE_WIDTH = SLIDE_W - 2 * MARGIN
SUBTITLE_HEIGHT = 0.32
SUBTITLE_FONT_PT = 14

# Image placement: bigger when there's no subtitle, smaller when there is.
IMG_LEFT = MARGIN
IMG_TOP = 0.60
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = 6.40

IMG_TOP_WITH_SUBTITLE = 0.90
IMG_BOX_H_WITH_SUBTITLE = 6.10

FOOTER_LEFT = MARGIN
FOOTER_TOP = 7.05
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.40
FOOTER_FONT_PT = 9

# Paired-image slide layout: two images side-by-side, captioned above each.
PAIR_GAP = 0.20
PAIR_IMG_W = (SLIDE_W - 2 * MARGIN - PAIR_GAP) / 2
PAIR_IMG_LEFT_L = MARGIN
PAIR_IMG_LEFT_R = MARGIN + PAIR_IMG_W + PAIR_GAP
PAIR_LABEL_TOP = 0.55
PAIR_LABEL_HEIGHT = 0.50
PAIR_LABEL_FONT_PT = 28
PAIR_IMG_TOP = 1.10
PAIR_IMG_H = 5.90
PAIR_FOOTER_FONT_PT = 8

# ---------------------------------------------------------------------------


def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False):
    """Add a centered textbox with given text and styling."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    """Place an image inside the given bounding box, preserving aspect ratio
    and centering on the dimension that is < box."""
    pic = slide.shapes.add_picture(
        image_path,
        Inches(box_left),
        Inches(box_top),
        width=Inches(box_w),
    )
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path,
            Inches(box_left),
            Inches(box_top),
            height=Inches(box_h),
        )
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def build_slide(prs, title_text: str, image_path: Path, footer_text: str,
                subtitle_text=None):
    """Build one slide: title + optional subtitle + full image + filepath footer.
    Returns (slide, missing_flag)."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, WHITE)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=BLACK, bold=True,
    )

    if subtitle_text:
        add_textbox(
            slide, subtitle_text,
            SUBTITLE_LEFT, SUBTITLE_TOP, SUBTITLE_WIDTH, SUBTITLE_HEIGHT,
            font_pt=SUBTITLE_FONT_PT, color=BLACK, italic=True,
        )
        img_top, img_h = IMG_TOP_WITH_SUBTITLE, IMG_BOX_H_WITH_SUBTITLE
    else:
        img_top, img_h = IMG_TOP, IMG_BOX_H

    missing = not image_path.exists()
    if not missing:
        add_image_in_box(
            slide, str(image_path), IMG_LEFT, img_top, IMG_BOX_W, img_h
        )
    else:
        add_textbox(
            slide, "(missing)",
            IMG_LEFT, img_top + img_h / 2 - 0.2, IMG_BOX_W, 0.4,
            font_pt=18, color=BLACK,
        )

    add_textbox(
        slide, footer_text,
        FOOTER_LEFT, FOOTER_TOP, FOOTER_WIDTH, FOOTER_HEIGHT,
        font_pt=FOOTER_FONT_PT, color=BLACK, bold=False,
    )
    return slide, missing


def build_pair_slide(prs, title: str, image_paths, labels, footers):
    """Build one slide with two images side-by-side under per-image captions.
    Returns (slide, missing_left, missing_right)."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, WHITE)

    add_textbox(
        slide, title,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=BLACK, bold=True,
    )

    lefts = (PAIR_IMG_LEFT_L, PAIR_IMG_LEFT_R)
    missing_flags = [False, False]
    for i in (0, 1):
        add_textbox(
            slide, labels[i],
            lefts[i], PAIR_LABEL_TOP, PAIR_IMG_W, PAIR_LABEL_HEIGHT,
            font_pt=PAIR_LABEL_FONT_PT, color=BLACK, bold=True,
        )
        if image_paths[i].exists():
            add_image_in_box(
                slide, str(image_paths[i]),
                lefts[i], PAIR_IMG_TOP, PAIR_IMG_W, PAIR_IMG_H,
            )
        else:
            missing_flags[i] = True
            add_textbox(
                slide, "(missing)",
                lefts[i], PAIR_IMG_TOP + PAIR_IMG_H / 2 - 0.2,
                PAIR_IMG_W, 0.4, font_pt=14, color=BLACK,
            )

    # Two-line footer: one filepath per line.
    box = slide.shapes.add_textbox(
        Inches(FOOTER_LEFT), Inches(FOOTER_TOP),
        Inches(FOOTER_WIDTH), Inches(FOOTER_HEIGHT),
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = footers[0]
    para1 = tf.paragraphs[0]
    para1.alignment = PP_ALIGN.CENTER
    para1.runs[0].font.size = Pt(PAIR_FOOTER_FONT_PT)
    para1.runs[0].font.color.rgb = BLACK
    para2 = tf.add_paragraph()
    para2.text = footers[1]
    para2.alignment = PP_ALIGN.CENTER
    para2.runs[0].font.size = Pt(PAIR_FOOTER_FONT_PT)
    para2.runs[0].font.color.rgb = BLACK

    return slide, missing_flags[0], missing_flags[1]


def build_deck(dataset) -> int:
    """Build one deck for the given dataset entry. Returns missing-file count."""
    label = dataset["label"]
    compiled_root = dataset["compiled_root"]
    output_path = dataset["output_path"]
    slides = dataset["slides"]

    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    print(f"\n=== Dataset: {label} ===")
    print(f"Source root: {compiled_root}")
    print(f"Writing deck to: {output_path}\n")

    missing = []
    for sub, filename, title in slides:
        if sub == SUB_CROSSCORR_PAIR:
            path_xy = compiled_root / "crosscorr_profiles" / "XY_MIP" / filename
            path_bs = (compiled_root / "crosscorr_profiles"
                       / "nuc_broadest_slice" / filename)
            _, miss_l, miss_r = build_pair_slide(
                prs, title,
                (path_xy, path_bs),
                ("XY MIP", "Broadest slice"),
                (path_xy.as_posix(), path_bs.as_posix()),
            )
            status_l = "OK" if not miss_l else "MISSING"
            status_r = "OK" if not miss_r else "MISSING"
            print(f"[{filename}]  XY:{status_l} BS:{status_r}  -> {title!r}")
            if miss_l:
                missing.append(f"{filename} (XY_MIP)  ({path_xy})")
            if miss_r:
                missing.append(f"{filename} (broadest slice)  ({path_bs})")
            continue

        image_path = compiled_root / sub / filename
        footer = image_path.as_posix()
        _, is_missing = build_slide(prs, title, image_path, footer)
        status = "OK" if not is_missing else "MISSING"
        print(f"[{filename}]  {status}  -> {title!r}")
        if is_missing:
            missing.append(f"{filename}  ({image_path})")

    # Snapshot the previous deck (if any) into backups/ before overwriting.
    if output_path.exists():
        backup_dir = output_path.parent / "backups"
        created = backup_presentation(str(output_path), backup_base=str(backup_dir))
        if created:
            print(f"\nBacked up previous deck to: {backup_dir}")
            for cat, path in created.items():
                print(f"  [{cat:6}] {path}")

    prs.save(str(output_path))
    print(f"\nDone. {len(slides)} slides written to:\n  {output_path}")

    if missing:
        print(f"\nMissing ({len(missing)}):")
        for m in missing:
            print(f"  - {m}")
    else:
        print("\nAll images found - no missing items.")

    return len(missing)


def main() -> None:
    total_missing = 0
    for dataset in DATASETS:
        total_missing += build_deck(dataset)
    if total_missing:
        sys.exit(1)


if __name__ == "__main__":
    main()
