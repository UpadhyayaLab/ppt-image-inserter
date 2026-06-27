"""
insert_noco_washout_summary_slides.py

Metric summary deck for the fixed Jurkat nocodazole-WASHOUT experiment
(20240203_MG_NocoWashout), compiled into results/compiled_20260626. This is the
washout time-course companion to the standard noco DMSO-vs-noco decks: cells were
treated with 1 uM nocodazole (or DMSO control) and fixed at three washout
timepoints (0, 4, 15 min after washout).

The compile already groups each metric into per-comparison subfolders. Here we
use the three DMSO-vs-noco pairwise folders:
    grid_panels/DMSO_vs_Noco_0m
    grid_panels/DMSO_vs_Noco_4m
    grid_panels/DMSO_vs_Noco_15m
Each panel (`<stem>.png`) is a DMSO (left) vs noco (right) violin/swarm comparison
with significance brackets at that one washout timepoint. One slide therefore
shows a single metric as a washout time course: the 0 / 4 / 15 min panels side by
side, so the DMSO-vs-noco effect and its recovery are read left-to-right.

Modeled on insert_bleb_summary_slides.py (title slide, family dividers, curated
FAMILIES, --list dry-run, backup-before-overwrite); the only structural change is
the 3-column-per-slide timepoint layout instead of one panel per slide. PNG (not
the sibling .tiff) is used because python-pptx adds PNG reliably.

Metric selection: the curated fixed-Jurkat drug-vs-DMSO set shared by the noco /
bleb decks (vimentin included, since this experiment is vim-stained). The
centrosome family leads, opening on Centrosome-synapse distance.

Self-contained: builds a blank deck (no template .pptx). Missing panels render
"(missing)" rather than failing. A previous deck is backed up before overwrite.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_noco_washout_summary_slides.py
    # dry run (print planned families/titles + which panels are on disk):
    conda run -n PPT_editing python examples_and_configs/insert_noco_washout_summary_slides.py --list
"""

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation  # noqa: E402

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(
    "L:/FF/Nucleus_centrosome/vimentin_nocodazole-washout_fixed/"
    "20240203_MG_NocoWashout/results/compiled_20260626"
)
GRID_DIR = ROOT / "grid_panels"

OUTPUT_PATH = Path(
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/Noco/"
    "Noco_Jurkats_washout_summary_20240203.pptx"
)

# (subfolder, column label) per washout timepoint, left -> right.
# DMSO_vs_Noco_* are the post-rename folders (DMSO left, noco right). The earlier
# Noco_vs_DMSO_* folders are stale leftovers from a pre-change compile and were
# NOT overwritten — do not point at them.
TIMEPOINTS = [
    ("DMSO_vs_Noco_0m",  "Pre-washout"),
    ("DMSO_vs_Noco_4m",  "4 min post-washout"),
    ("DMSO_vs_Noco_15m", "15 min post-washout"),
]
PANEL_SUFFIX = ".png"

DECK_TITLE = "Effects of nocodazole washout on fixed Jurkat nuclei"
DECK_SUBTITLE = (
    "DMSO vs nocodazole 1 μM, washout time course (0 / 4 / 15 min)  ·  "
    "each panel is DMSO vs noco at one timepoint  ·  αCD3, E6-1  ·  "
    "fixed 02/03/2024  ·  compiled 2026-06-26"
)

# ---------------------------------------------------------------------------
# Curated metrics, grouped into families (divider slide per family).
# Each entry is (panel stem, slide title). Titles match the noco / CilioD decks;
# each plot carries its own authoritative y-axis label and DMSO/noco x-labels.
# Leads on the centrosome family (Centrosome-synapse distance first), as in the
# noco metric-grid decks.
# ---------------------------------------------------------------------------
FAMILIES = [
    ("Centrosome ↔ nucleus", [
        ("centrosome_center_z_rel_bottom_actin_plane", "Centrosome-synapse distance"),
        ("nuc_cent_closest_dist",                "Nucleus-centrosome closest distance"),
        ("cent_nuc_norm_dist_sphere_rad",        "Centrosome-nucleus distance (norm. to nuclear sphere radius)"),
        ("centrosome_dist_deepest_real_avg_periphery_ratio", "Centrosome distance to deepest invag vs avg periphery ratio"),
    ]),
    ("Nuclear morphology", [
        ("nuc_aspect_ratio",  "Nuclear aspect ratio"),
        ("nuc_solidity",      "Nuclear solidity"),
        ("nuc_volume_mesh",   "Nuclear volume"),
        ("nuc_SA_mesh",       "Nuclear surface area"),
    ]),
    ("Nuclear deformation & invaginations", [
        ("chull_max_D",                       "Max invag depth over full nucleus"),
        ("chull_max_D_by_cent",               "Invagination depth near centrosome"),
        ("chull_mean_D_cent_global_ratio",    "Centrosomal Invagination Index (global)"),
        ("concavity_index_around_cent",       "Concavity index around centrosome"),
        ("deepest_invag_fraction_chull_volume", "Deepest invag: frac of convex hull volume"),
        ("deepest_region_periph_ratio_025um", "DNA levels near invag"),
        ("avg_normal_angle_adaptive_region_growth", "Deepest Invag Orientation"),
    ]),
    ("Actin", [
        ("actin_deform_ratio",        "Cell Aspect Ratio"),
        ("actin_bottom_mask_area",    "Synapse Area"),
        ("actin_MFI_around_cent_2um", "Actin MFI around centrosome (2 μm)"),
        ("actin_frac_around_cent_2um", "Actin fraction around centrosome (2 μm)"),
    ]),
    ("Vimentin", [
        ("vim_frac_in_nuc_convex_hull",                  "Vimentin fraction in nuclear convex hull"),
        ("vim_enrichment_within_half_um_nuc_2_um_cent",  "Vimentin enrichment (0.5 μm of nuc, 2 μm of cent)"),
        ("vim_cyto_in_nuc_hull_vs_near_convex_nuc_MFI_ratio", "Vimentin MFI: cytoplasmic (in nuc hull) vs near-convex regions of nuc"),
        ("VCRAI_0_2um",               "VCRAI (0–2 μm)"),
        ("vim_frac_around_cent_2um",  "Vimentin fraction around centrosome (2 μm)"),
        ("vim_MFI_around_cent_2um",   "Vimentin MFI around centrosome (2 μm)"),
        ("vim_ratio_above_below_0_5um", "Vimentin ratio above/below 0.5 μm"),
    ]),
]

# ---------------------------------------------------------------------------
# Colors / layout
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.10

TITLE_LEFT = MARGIN
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.55
TITLE_FONT_PT = 28

# Three timepoint columns.
COL_GAP = 0.12
N_COLS = len(TIMEPOINTS)
COL_W = (SLIDE_W - 2 * MARGIN - (N_COLS - 1) * COL_GAP) / N_COLS
COL_LEFTS = [MARGIN + i * (COL_W + COL_GAP) for i in range(N_COLS)]

COL_LABEL_H = 0.32
COL_LABEL_FONT_PT = 16
COL_LABEL_GAP = 0.03          # gap between the label and the figure beneath it

# Images are vertically centered in this band (title bottom -> near slide bottom,
# no footer). Each column's label is then placed right above wherever its
# centered image actually lands, so the labels hug the figures.
IMG_TOP = 0.66
IMG_BOX_H = SLIDE_H - IMG_TOP - 0.15             # ~6.69"

# ---------------------------------------------------------------------------


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    """Place an image inside (left, top, w, h), preserving aspect ratio and
    centering. Panels are near-square, so they usually fit to width; if that
    overflows the box height, refit to height and center horizontally."""
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def panel_path(timepoint_sub, stem):
    return GRID_DIR / timepoint_sub / (stem + PANEL_SUFFIX)


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.6, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=40, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.0, SLIDE_W - 2 * MARGIN, 1.4,
                font_pt=16, color=BLACK, italic=True)


def build_divider_slide(prs, family_name):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, family_name, MARGIN, 3.1, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=44, color=BLACK, bold=True)


def build_metric_slide(prs, title_text, stem):
    """Title + the 0/4/15 min panels in a row (each a DMSO-vs-noco comparison).
    The timepoint label for each column is placed just above where its centered
    image lands, so labels sit right on top of the figures. Returns the list of
    timepoint labels whose panel was missing on disk."""
    slide = _new_slide(prs)
    add_textbox(slide, title_text, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title_text), color=BLACK, bold=True)

    missing = []
    for col, (sub, tp_label) in enumerate(TIMEPOINTS):
        left = COL_LEFTS[col]
        p = panel_path(sub, stem)
        if p.exists():
            pic = add_image_in_box(slide, str(p), left, IMG_TOP, COL_W, IMG_BOX_H)
            img_top_in = pic.top / 914400.0
            label_top = img_top_in - COL_LABEL_H - COL_LABEL_GAP
        else:
            add_textbox(slide, "(missing)", left, IMG_TOP + IMG_BOX_H / 2 - 0.2,
                        COL_W, 0.4, font_pt=14, color=BLACK)
            missing.append(tp_label)
            label_top = IMG_TOP + IMG_BOX_H / 2 - 0.7
        add_textbox(slide, tp_label, left, label_top, COL_W, COL_LABEL_H,
                    font_pt=COL_LABEL_FONT_PT, color=BLACK, bold=True)
    return missing


def main():
    list_only = "--list" in sys.argv

    n_metrics = sum(len(items) for _, items in FAMILIES)
    est_slides = 1 + sum(1 + len(items) for _, items in FAMILIES)

    print("Source: {}".format(GRID_DIR))
    print("Timepoints: {}".format(", ".join(lbl for _, lbl in TIMEPOINTS)))
    print("{} curated metrics across {} families, est. {} slides\n".format(
        n_metrics, len(FAMILIES), est_slides))

    if list_only:
        for fam, items in FAMILIES:
            print("=== {} ({}) ===".format(fam, len(items)))
            for stem, title in items:
                flags = "".join(
                    "{}:{} ".format(lbl, "OK" if panel_path(sub, stem).exists() else "MISS")
                    for sub, lbl in TIMEPOINTS)
                print("  {:<50s} {:<60s} {}".format(stem, title, flags))
            print("")
        return

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)

    missing = []
    for fam, items in FAMILIES:
        build_divider_slide(prs, fam)
        print("=== {} ===".format(fam))
        for stem, title in items:
            miss_tps = build_metric_slide(prs, title, stem)
            status = "OK" if not miss_tps else "MISSING " + ",".join(miss_tps)
            print("  [{}] {} -> {!r}".format(status, stem, title))
            for tp in miss_tps:
                missing.append("{} ({})".format(stem, tp))
        print("")

    if OUTPUT_PATH.exists():
        backup_dir = OUTPUT_PATH.parent / "backups"
        created = backup_presentation(str(OUTPUT_PATH), backup_base=str(backup_dir))
        if created:
            print("Backed up previous deck to: {}\n".format(backup_dir))

    prs.save(str(OUTPUT_PATH))
    total = len(prs.slides._sldIdLst)
    print("Done. {} metrics, {} slides written to:\n  {}".format(
        n_metrics, total, OUTPUT_PATH))
    if missing:
        print("\nSkipped {} missing panel(s) (not on disk):".format(len(missing)))
        for m in missing:
            print("  - {}".format(m))
    else:
        print("\nAll curated panels found - no missing items.")


if __name__ == "__main__":
    main()
