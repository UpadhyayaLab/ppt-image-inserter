"""
insert_pMLC_qc_CAT_vs_FMC_Kiet_slides.py

CAT (left) vs FMC (right) side-by-side pMLC QC compare deck for the 3
Kiet Y:-drive datasets using the new unified `prog_fixed_cells/` layout.

Three datasets (Y:/User_data/Kiet/...), each with 6 conditions
(CAT_5/10/15min, FMC_5/10/15min):

    20260312  03122026_pMLC_actin_CAR_T
    20260510  20260510_pMLC_Actin561_nucleus_CAR_Tcell_
    20260607  20260607_pMLC_CART_actin_hoescht

Per condition, the pMLC montage tree is:

    <dataset>/<COND>/converted/cropped/split_channels/prog_fixed_cells/pMLC/
        synapse/1slice/raw/montages/
        synapse/3slice/raw/montages/
        xz_mip/montages/

Result: 2 synapse kinds (interleaved) + 1 XZ MIP kind, × 3 datasets × 3
timepoints = 27 slides.

Per-kind PPI: each kind has its own PPI tuned to its widest montage so
the synapse kinds (more square aspect) don't drag the XZ MIP kind (wide
aspect) down to a tighter scale.

Modeled on insert_actin_qc_cat_vs_fmc_20260607_slides.py (path layout)
and insert_MT_qc_CAT_vs_FMC_202311_202406_slides.py (per-kind PPI).

Usage:
    python examples_and_configs/insert_pMLC_qc_CAT_vs_FMC_Kiet_slides.py
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/CART/pMLC/"
    "CART_pMLC_QC_CAT_vs_FMC_Kiet.pptx"
)

KIET_ROOT = "Y:/User_data/Kiet"

# (YYYYMMDD acquisition date, dataset folder name).
# 20260312 folder is MMDDYYYY (03122026); the other two are already YYYYMMDD.
DATASETS = [
    ("20260312", "03122026_pMLC_actin_CAR_T"),
    ("20260510", "20260510_pMLC_Actin561_nucleus_CAR_Tcell_"),
    ("20260607", "20260607_pMLC_CART_actin_hoescht"),
]

TIMEPOINTS = ["5min", "10min", "15min"]

CONDITION_SUBPATH = "converted/cropped/split_channels"
PROG_SUBPATH = "prog_fixed_cells/pMLC"

# pMLC has no segmentation mask. Block 1 pairs the two synapse depth views
# (single z-slice and 3-slice slab MIP); block 2 is the XZ MIP.
KIND_BLOCKS = [
    [
        ("pMLC at Synapse",         "synapse/1slice/raw/montages"),
        ("pMLC 3-Slice at Synapse", "synapse/3slice/raw/montages"),
    ],
    [
        ("pMLC XZ MIP", "xz_mip/montages"),
    ],
]

CHUNK_GLOB = "montage_cells_*.png"

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Slide layout (inches). 13.333 x 7.5 widescreen.
SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

# 1x2 cell grid below the title (label + image per cell).
GRID_LEFT = 0.10
GRID_TOP = 0.60
CELL_W = 6.50
CELL_H = SLIDE_H - GRID_TOP - 0.10    # 6.80"
LABEL_H = 0.30
IMG_H = CELL_H - LABEL_H              # 6.50"
LABEL_FONT_PT = 16
COL_GAP = SLIDE_W - 2 * GRID_LEFT - 2 * CELL_W

CELL_POSITIONS = [
    (GRID_LEFT,                    GRID_TOP),  # left = CAT
    (GRID_LEFT + CELL_W + COL_GAP, GRID_TOP),  # right = FMC
]

# Scalebar invariant (Stage AF/AG): per-tile 5 μm bar = 150 px nominal.
PPUM_SOURCE = 30
SCALEBAR_UM = 5
SCALEBAR_PX = PPUM_SOURCE * SCALEBAR_UM

# ---------------------------------------------------------------------------


def format_timepoint(tp: str) -> str:
    tp_map = {"5min": "5 min", "10min": "10 min", "15min": "15 min"}
    return tp_map.get(tp.lower(), tp)


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _parse_chunk_range(p: Path) -> Tuple[int, int]:
    """Return (start, end) cell-id range for a montage_cells_*.png filename.
    Handles both the 4-int FOV-padded pattern and the 2-int pattern."""
    m4 = re.match(r"montage_cells_(\d+)_(\d+)_(\d+)_(\d+)\.png$", p.name)
    if m4:
        f_a, c_a, f_b, c_b = (int(x) for x in m4.groups())
        return (f_a * 1000 + c_a, f_b * 1000 + c_b)
    m2 = re.match(r"montage_cells_(\d+)_(\d+)\.png$", p.name)
    if m2:
        return (int(m2.group(1)), int(m2.group(2)))
    return (0, 0)


def find_first_chunk(montages_dir: Path) -> Optional[Path]:
    """Pick the lowest-start chunk, dropping any whose [start, end] range is
    strictly contained in another chunk's range (smoke-shadow filter)."""
    if not montages_dir.is_dir():
        return None
    chunks = list(montages_dir.glob(CHUNK_GLOB))
    if not chunks:
        return None
    parsed = [(p, *_parse_chunk_range(p)) for p in chunks]
    keep = []
    for (p, s, e) in parsed:
        shadowed = any(
            s2 <= s and e <= e2 and (s2 < s or e < e2)
            for (p2, s2, e2) in parsed
            if p2 is not p
        )
        if not shadowed:
            keep.append((p, s, e))
    if not keep:
        return None
    keep.sort(key=lambda x: x[1])
    return keep[0][0]


def _png_dims(path: Path) -> Tuple[int, int]:
    with Image.open(str(path)) as im:
        return im.size


def compute_deck_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_image_in_cell_at_ppi(slide, image_path: Path, ppi: float,
                             cell_left: float, cell_top: float):
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    img_area_top = cell_top + LABEL_H
    left_in = cell_left + (CELL_W - w_in) / 2
    top_in  = img_area_top + (IMG_H - h_in) / 2
    return slide.shapes.add_picture(
        str(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(w_in),
    )


def build_compare_slide(prs, title_text, cat_img, fmc_img, deck_ppi):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    cells = [
        ("CAT", cat_img, CELL_POSITIONS[0]),
        ("FMC", fmc_img, CELL_POSITIONS[1]),
    ]
    missing = []
    for label, img_path, (cell_left, cell_top) in cells:
        add_textbox(
            slide, label,
            cell_left, cell_top, CELL_W, LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )
        if img_path is not None and img_path.exists():
            add_image_in_cell_at_ppi(slide, img_path, deck_ppi, cell_left, cell_top)
        else:
            add_textbox(
                slide, "(missing)",
                cell_left, cell_top + LABEL_H + IMG_H / 2 - 0.15, CELL_W, 0.3,
                font_pt=14, color=WHITE,
            )
            missing.append(label)
    return slide, missing


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    kiet_root = Path(KIET_ROOT)

    # Pre-pass: walk every (block, dataset, tp, kind) → resolve CAT/FMC paths
    # and pre-load montage pixel dims to compute per-kind PPI BEFORE slide build.
    slide_specs: List[Tuple[str, Optional[Path], Optional[Path], Path, Path, str, str]] = []
    for block in KIND_BLOCKS:
        for date_tag, dataset_name in DATASETS:
            for tp in TIMEPOINTS:
                tp_pretty = format_timepoint(tp)
                for kind_label, kind_subpath in block:
                    cat_dir = (
                        kiet_root / dataset_name / f"CAT_{tp}"
                        / CONDITION_SUBPATH / PROG_SUBPATH / kind_subpath
                    )
                    fmc_dir = (
                        kiet_root / dataset_name / f"FMC_{tp}"
                        / CONDITION_SUBPATH / PROG_SUBPATH / kind_subpath
                    )
                    cat_img = find_first_chunk(cat_dir)
                    fmc_img = find_first_chunk(fmc_dir)
                    title = f"{kind_label}: {tp_pretty} ({date_tag})"
                    log_key = f"{kind_label}/{date_tag}/{tp}"
                    slide_specs.append((title, cat_img, fmc_img, cat_dir, fmc_dir, log_key, kind_label))

    # Group present images by kind and compute one PPI per kind.
    present_by_kind: dict = {}
    for (_, cat_img, fmc_img, _, _, _, kind_label) in slide_specs:
        bucket = present_by_kind.setdefault(kind_label, [])
        for p in (cat_img, fmc_img):
            if p is not None and p.exists():
                bucket.append(p)

    ppi_by_kind = {}
    for kind_label, paths in present_by_kind.items():
        if not paths:
            ppi_by_kind[kind_label] = 100.0
        else:
            ppi_by_kind[kind_label] = compute_deck_ppi(paths, CELL_W, IMG_H)

    print("Per-kind PPI (each kind internally consistent, cross-kind may differ):")
    for kind_label in [k for block in KIND_BLOCKS for (k, _) in block]:
        ppi = ppi_by_kind.get(kind_label)
        if ppi is None:
            continue
        bar_in = SCALEBAR_PX / ppi
        n = len(present_by_kind.get(kind_label, []))
        print(
            f"  {kind_label:<26} -> PPI {ppi:7.2f}  "
            f"(5 μm = {bar_in:.3f} in = {bar_in * 2.54:.3f} cm)  [{n} cells]"
        )
    print(f"\nSource PPUM = {PPUM_SOURCE} px/μm (locked).\n")
    print(f"Writing deck to: {OUTPUT_PATH}\n")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    missing_total = []
    slides_added = 0
    for (title, cat_img, fmc_img, cat_dir, fmc_dir, log_key, kind_label) in slide_specs:
        kind_ppi = ppi_by_kind[kind_label]
        _, missing = build_compare_slide(prs, title, cat_img, fmc_img, kind_ppi)
        slides_added += 1

        status_parts = [
            "CAT:OK" if cat_img else "CAT:MISSING",
            "FMC:OK" if fmc_img else "FMC:MISSING",
        ]
        print(f"[{log_key}]  " + "  ".join(status_parts))

        for cell in missing:
            src = cat_dir if cell == "CAT" else fmc_dir
            missing_total.append(f"{log_key}/{cell}  ({src})")

    prs.save(str(out_path))
    print(f"\nDone. {slides_added} slides written to:\n  {out_path}")

    if missing_total:
        print(f"\nMissing ({len(missing_total)}):")
        for m in missing_total:
            print(f"  - {m}")
    else:
        print("\nAll images found - no missing items.")


if __name__ == "__main__":
    main()
