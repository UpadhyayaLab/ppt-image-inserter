"""
insert_actin_qc_cat_vs_fmc_slides.py

CAT (left) vs FMC (right) side-by-side comparison deck for the actin-only
pipeline QC montages, paired by timepoint.

For each (kind, dataset, timepoint) triple, builds one slide with two
labelled cells:

    Slide title — "<Kind>: <timepoint> (<YYYYMMDD>)"

       CAT                              FMC
    +---------+                    +---------+
    | montage |                    | montage |
    +---------+                    +---------+

Result: 2 kinds x 3 datasets x 3 timepoints = 18 slides.

Sources (per condition, same as the per-condition deck):
    <base>/prog_fixed_cells_actin_only/actin/synapse/inner_outer/bot_combined/montages/
    <base>/prog_fixed_cells_actin_only/actin/xz_mip/montages/
where <base> = <dataset>/CAT_<tp>/converted/cropped/split_channels/
          or  <dataset>/FMC_<tp>/converted/cropped/split_channels/

Within each montages/ directory we glob `montage_cells_*.png` and pick
the first sorted file (first chunk, ~25 cells).

Companion to insert_actin_qc_synapse_mask_xz_mip_slides.py (per-condition,
36 slides). Aspect-preserving cell-fit helper copied from
insert_lwi_qc_part1_minEdgeLen_16_p01_slides.py.

Usage:
    python examples_and_configs/insert_actin_qc_cat_vs_fmc_slides.py
"""

import os
import re
import sys
from pathlib import Path
from typing import Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/CART/actin_only/"
    "CART_actin_QC_CAT_vs_FMC.pptx"
)

KIET_ROOT = "Y:/User_data/Kiet"

# (YYYYMMDD acquisition date, dataset folder name)
# D1 folder name is MMDDYYYY (03122026 = 12 Mar 2026); D2/D3 are already YYYYMMDD.
DATASETS = [
    ("20260312", "03122026_pMLC_actin_CAR_T"),
    ("20260510", "20260510_pMLC_Actin561_nucleus_CAR_Tcell_"),
    ("20260414", "20260414_p_PKC_theta_Phalloidin561_"),
]

TIMEPOINTS = ["5min", "10min", "15min"]

CONDITION_SUBPATH = "converted/cropped/split_channels"
PROG_SUBPATH = "prog_fixed_cells_actin_only/actin"

# Kind blocks. Each block is a list of (kind label, subpath under PROG_SUBPATH/)
# tuples that get INTERLEAVED at each (dataset, timepoint) — i.e. the inner loop
# walks the kinds within the block. Between blocks we restart the
# (dataset, timepoint) walk.
#
# Slide order for the current config (3 datasets x 3 timepoints):
#   Block 1 (18 slides): for each (date, tp): Synapse Mask -> Actin at Synapse
#   Block 2 ( 9 slides): for each (date, tp): Actin XZ MIP
# Total: 27 slides.
KIND_BLOCKS = [
    [   # Block 1: interleaved synapse pair (no-rings then with-rings)
        ("Actin at Synapse",            "synapse/mask/montages"),
        ("Inner-Outer Ratio Definition", "synapse/inner_outer/bot_combined/montages"),
    ],
    [   # Block 2: XZ MIP on its own
        ("Actin XZ MIP", "xz_mip/montages"),
    ],
]

CHUNK_GLOB = "montage_cells_*.png"

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Slide layout (inches). 13.333 x 7.5 widescreen.
SLIDE_W = 13.333
SLIDE_H = 7.5

# Title strip at top
TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

# 1x2 cell grid below the title (label + image per cell)
GRID_LEFT = 0.10
GRID_TOP = 0.60
CELL_W = 6.50          # each cell ~half the slide width
CELL_H = SLIDE_H - GRID_TOP - 0.10   # 6.80"
LABEL_H = 0.30
IMG_H = CELL_H - LABEL_H             # 6.50"
LABEL_FONT_PT = 16
COL_GAP = SLIDE_W - 2 * GRID_LEFT - 2 * CELL_W   # gap between left & right cells

CELL_POSITIONS = [
    (GRID_LEFT,                    GRID_TOP),   # Left  (CAT)
    (GRID_LEFT + CELL_W + COL_GAP, GRID_TOP),   # Right (FMC)
]

# ---------------------------------------------------------------------------


def format_timepoint(tp: str) -> str:
    """5min -> '5 min'."""
    tp_map = {"5min": "5 min", "10min": "10 min", "15min": "15 min"}
    return tp_map.get(tp.lower(), tp)


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    """Add a centered textbox with given text and styling."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_image_in_cell(slide, image_path, cell_left, cell_top):
    """Place an image inside a labelled cell (below the label area),
    preserving aspect ratio. Width-bind first, fall back to height-bind
    if the image would overflow vertically. Centers on the non-binding axis.
    """
    pic = slide.shapes.add_picture(
        image_path,
        Inches(cell_left),
        Inches(cell_top + LABEL_H),
        width=Inches(CELL_W),
    )
    actual_h_in = pic.height / 914400.0
    if actual_h_in > IMG_H:
        # Too tall — refit by height and re-center horizontally.
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path,
            Inches(cell_left),
            Inches(cell_top + LABEL_H),
            height=Inches(IMG_H),
        )
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(cell_left + (CELL_W - actual_w_in) / 2)
    else:
        # Fits width-bound — center vertically.
        pic.top = Inches(cell_top + LABEL_H + (IMG_H - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _parse_chunk_range(p: Path) -> Tuple[int, int]:
    """Return (start, end) cell-id range for a montage_cells_*.png filename.
    Handles both the 4-int FOV-padded pattern and the 2-int pattern."""
    m4 = re.match(r"montage_cells_(\d+)_(\d+)_(\d+)_(\d+)\.png$", p.name)
    if m4:
        f_a, c_a, f_b, c_b = (int(x) for x in m4.groups())
        return (f_a * 1000 + c_a, f_b * 1000 + c_b)
    m2 = re.match(r"montage_cells_(\d+)_(\d+)\.png$", p.name)
    if m2:
        return (int(m2.group(1)), int(m2.group(2)))
    return (0, 0)


def find_first_chunk(montages_dir: Path) -> Optional[Path]:
    """Pick the lowest-start chunk, BUT first drop any chunk whose [start, end]
    range is strictly contained in another chunk's range. Catches leftover
    smoke chunks (e.g. `montage_cells_1_12.png` shadowed by `montage_cells_1_26.png`)."""
    if not montages_dir.is_dir():
        return None
    chunks = list(montages_dir.glob(CHUNK_GLOB))
    if not chunks:
        return None
    parsed = [(p, *_parse_chunk_range(p)) for p in chunks]
    keep = []
    for (p, s, e) in parsed:
        shadowed = any(
            s2 <= s and e <= e2 and (s2 < s or e < e2)
            for (p2, s2, e2) in parsed
            if p2 is not p
        )
        if not shadowed:
            keep.append((p, s, e))
    if not keep:
        return None
    keep.sort(key=lambda x: x[1])
    return keep[0][0]


def build_compare_slide(prs, title_text, cat_img, fmc_img):
    """Build one CAT-vs-FMC slide. Returns (slide, missing_list)."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    cells = [
        ("CAT", cat_img, CELL_POSITIONS[0]),
        ("FMC", fmc_img, CELL_POSITIONS[1]),
    ]
    missing = []
    for label, img_path, (cell_left, cell_top) in cells:
        add_textbox(
            slide, label,
            cell_left, cell_top, CELL_W, LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )
        if img_path is not None and img_path.exists():
            add_image_in_cell(slide, str(img_path), cell_left, cell_top)
        else:
            add_textbox(
                slide, "(missing)",
                cell_left, cell_top + LABEL_H + IMG_H / 2 - 0.15, CELL_W, 0.3,
                font_pt=14, color=WHITE,
            )
            missing.append(label)
    return slide, missing


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    kiet_root = Path(KIET_ROOT)
    missing_total = []
    slides_added = 0

    print(f"Writing deck to: {OUTPUT_PATH}\n")

    # Loop order: block -> dataset -> timepoint -> kind (within block).
    # The two synapse kinds in block 1 interleave at each (date, tp):
    #   slide 1: Synapse Mask (no rings), slide 2: Actin at Synapse (rings)
    # Block 2 is XZ MIP on its own — one slide per (date, tp).
    for block in KIND_BLOCKS:
        for date_tag, dataset_name in DATASETS:
            for tp in TIMEPOINTS:
                tp_pretty = format_timepoint(tp)
                for kind_label, kind_subpath in block:
                    cat_dir = (
                        kiet_root / dataset_name / f"CAT_{tp}"
                        / CONDITION_SUBPATH / PROG_SUBPATH / kind_subpath
                    )
                    fmc_dir = (
                        kiet_root / dataset_name / f"FMC_{tp}"
                        / CONDITION_SUBPATH / PROG_SUBPATH / kind_subpath
                    )
                    cat_img = find_first_chunk(cat_dir)
                    fmc_img = find_first_chunk(fmc_dir)

                    title = f"{kind_label}: {tp_pretty} ({date_tag})"
                    _, missing = build_compare_slide(prs, title, cat_img, fmc_img)
                    slides_added += 1

                    status_parts = [
                        "CAT:OK" if cat_img else "CAT:MISSING",
                        "FMC:OK" if fmc_img else "FMC:MISSING",
                    ]
                    print(f"[{kind_label}/{date_tag}/{tp}]  " + "  ".join(status_parts))

                    for cell in missing:
                        src = cat_dir if cell == "CAT" else fmc_dir
                        missing_total.append(
                            f"{kind_label}/{date_tag}/{tp}/{cell}  ({src})"
                        )

    prs.save(str(out_path))
    print(f"\nDone. {slides_added} slides written to:\n  {out_path}")

    if missing_total:
        print(f"\nMissing ({len(missing_total)}):")
        for m in missing_total:
            print(f"  - {m}")
    else:
        print("\nAll images found - no missing items.")


if __name__ == "__main__":
    main()
