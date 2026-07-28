"""
insert_chromark_nucleus_morphology_3expt_slides.py

Build the cross-experiment NUCLEUS MORPHOLOGY deck from the curated violin set:

  chromatin-analysis-figures/cross_experiment/nucleus_morphology_3expt/
      violin_selected/     4 featured metrics (aspect ratio, volume, DNA skew/kurtosis)
      violin_all_curated/  137 metrics, the full curated body

Every PNG is a 3-column grid (one column per acquisition: H3K27me3 2024-03-21 |
H3K9me3 2024-04-24 | H3K27ac 2024-05-31), so cross-experiment consistency reads
left-to-right. Each metric has three views:

    grid_<stem>.png             all conditions (wide banner)
    grid_<stem>_stiffness.png   substrate contrast, with significance stars
    grid_<stem>_timepoint.png   10 min vs 3 h, with stars

Layout: ONE full-size plot per slide (plots fill the slide), a divider per metric
family, featured metrics first. Metric titles/families/ordering come from the
shared classifier in insert_chromark_h3k27me3_summary_slides.py, so conventions
(DNA not Hoechst, 3D before 2D, dimension tags) match the other chromark decks.

NOTE ON SCOPE: every curated metric ships. The cross-experiment report argues that
texture / peripheral / curvature partly track the constant-volume flattening rather
than being independent chromatin reorganization; that caveat is stated once on the
"How to read" slide and nothing is dropped or re-ordered because of it.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_chromark_nucleus_morphology_3expt_slides.py
    ... --list        # dry run: families/titles/counts, build nothing
    ... --selected    # small featured-only deck
"""

import os
import sys
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import insert_chromark_h3k27me3_summary_slides as M  # noqa: E402

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
BASE = Path(
    "J:/FF/fixed_cell/CTL_nucleus/tifsFixed3SIactivatedCTLs_nucleus/"
    "chromatin-analysis-figures/cross_experiment/nucleus_morphology_3expt"
)
DIR_ALL = BASE / "violin_all_curated"
DIR_SEL = BASE / "violin_selected"
DIR_MARKS = BASE / "violin_marks"

OUT_DIR = Path("K:/FF/PPT/PPT_autogeneration/Naive_CTL/chromark")
OUT_FULL = OUT_DIR / "CTL_chromark_nucleus_morphology_3expt_summary.pptx"
OUT_SEL = OUT_DIR / "CTL_chromark_nucleus_morphology_3expt_featured.pptx"
OUT_MARKS = OUT_DIR / "CTL_chromark_marks_3expt_summary.pptx"

COLUMNS = ("every panel: 3 columns = H3K27me3 (2024-03-21) | "
           "H3K9me3 (2024-04-24) | H3K27ac (2024-05-31)")

DECK_TITLE = "Nuclear morphology — cross-experiment (3 acquisitions)"
DECK_SUBTITLE = "Curated violin set  -  " + COLUMNS

# --marks: same metrics/trims/ordering, measured on the chromatin-mark channel.
# Each column is a DIFFERENT mark, so the channel is labelled generically.
MARKS_TITLE = "Chromatin mark metrics — cross-experiment (3 acquisitions)"
MARKS_SUBTITLE = (
    "Mark channel  -  " + COLUMNS + "  -  NB: one experiment per mark (no replicates); "
    "raw mark intensity is batch-sensitive"
)
MARKS_LABEL = "Chromatin mark"

PREFIX = "grid_"
# (filename suffix, view label used in the slide title)
VIEWS = [
    ("", "all conditions"),
    ("_stiffness", "stiffness comparison"),
    ("_timepoint", "timepoint comparison"),
]

# Body section order for THIS deck. Deliberately morphology-first: this source is
# nucleus_morphology_3expt (96 of 135 metrics are morphology) and the headline is the
# constant-volume flattening, so the shape families lead and intensity/texture follow
# as supporting material. Deck-local on purpose -- the shared M.FAMILIES_ORDER still
# governs the other chromark decks, which keep their intensity-first order.
BODY_ORDER = [
    "Nuclear morphology (maximal area slice)",
    "Mark–DNA co-distribution",
    "Chromatin organization",
    "Texture (GLCM)",
    "Texture (GLCM, max area slice)",
    "Other metrics",
]

# Highlights open the deck: the 3D morphology metrics, the 3D intensity metrics,
# and the maximal-area-slice cross-section area. These are MOVED here (not copied),
# so nothing repeats later in the deck.
HIGHLIGHT_FAMS = ["Nuclear morphology (3D)", "Nuclear intensity"]
HIGHLIGHT_EXTRA = ["morph2dmid_area"]
HIGHLIGHTS_TITLE = "Highlights"

# ---------------------------------------------------------------------------
# Deck-local trims (user-directed). These shape THIS deck only; the shared
# module's metric set is untouched, so the other chromark decks are unaffected.
# ---------------------------------------------------------------------------
#  - 2D morphology is taken from the maximal-area slice (a real optical section),
#    so the XY MIP silhouette versions (morph2d_*) are dropped entirely.
#  - Nuclear intensity: 3D only (drop the 2D intensity moments).
#  - Peripheral enrichment: skip the 2 um shell and the r33% shell.
DROP_PREFIXES = ("morph2d_",)
DROP_SUBSTRINGS = ("_peripheral_enrichment_2um", "_peripheral_enrichment_r33pct")

# Channel tokens, normalised to "@" so one drop-list covers the DNA deck and the
# chromatin-mark deck (violin_marks uses a generic "mark" token).
CHANNEL_TOKENS = ("hoechst", "h3k27me3", "mark")

#  - Orientation (angle of the fitted ellipse in the imaging plane) carries no
#    biological meaning here -- cells land in arbitrary rotations.
#  - chan_mean/chan_std duplicate the explicit nuclear-mask *_3d_nuclear_*_int
#    (identical titles), so keep the nuclear-mask version only.
#  - chan_rel_*_3d_int ("relative intensity") is a monotonic rescale of the
#    nuclear mean -- identical Cliff's delta in all three marks -- so it carries
#    no information the mean doesn't already show.
#  - _2dmid_int_* are absolute slice levels; upstream drops them from the
#    analysis bank (is_dna_intensity) and we keep intensity 3D-only anyway.
#    The slice SHAPE moments (skewness / kurtosis / entropy) are kept.
DROP_EXACT_NORM = {
    "@_2d_int_mean", "@_2d_int_sd", "@_2d_int_median", "@_2d_int_max",
    "@_2d_skewness", "@_2d_kurtosis",
    "@_2dmid_int_mean", "@_2dmid_int_median", "@_2dmid_int_sd", "@_2dmid_int_max",
    "@_2dmid_int_min", "@_2dmid_int_mode", "@_2dmid_int_d25", "@_2dmid_int_d75",
    "chan_mean_@_3d_int", "chan_std_@_3d_int", "chan_rel_@_3d_int",
    "morph2d_orientation", "morph2dmid_orientation",
}


def _norm_channel(stem):
    for tok in CHANNEL_TOKENS:
        stem = stem.replace(tok, "@")
    return stem


def dropped(stem):
    """True if this metric is trimmed out of this deck (channel-agnostic)."""
    if stem in M.DROP_METRICS:
        return True
    if _norm_channel(stem) in DROP_EXACT_NORM:
        return True
    if stem.startswith(DROP_PREFIXES):
        return True
    return any(sub in stem for sub in DROP_SUBSTRINGS)

HOW_TO_READ = [
    "How to read these slides",
    "",
    "•  Every panel is a 3-column grid — one column per acquisition "
    "(H3K27me3 0321 | H3K9me3 0424 | H3K27ac 0531).",
    "    A result that reproduces shows the same pattern left-to-right.",
    "•  Each metric gets three slides: all 8 conditions (overview), the stiffness "
    "contrast, and the 10 min vs 3 h",
    "    timepoint contrast. The contrast views carry FDR significance stars; the "
    "overview does not.",
    "•  White dot = median, bar = IQR, points = individual nuclei.",
    "",
    "Headline from the cross-experiment report: the reproducible response to substrate "
    "stiffness is a",
    "constant-volume flattening — on stiff glass the nucleus spreads and flattens "
    "(aspect ratio δ ≈ −0.96),",
    "on soft gel it stays rounder and taller, while total volume does not change.",
    "",
    "Caveat carried from that report: texture, peripheral-enrichment and curvature "
    "metrics partly track the",
    "flattening (they weaken or reverse when nuclei are matched on area). They are all "
    "included here — read",
    "them alongside the shape metrics rather than as independent chromatin reorganization.",
]


# ---------------------------------------------------------------------------
def discover(directory):
    """stem -> set of available view suffixes, from grid_<stem>[_view].png."""
    found = defaultdict(set)
    if not directory.is_dir():
        return found
    known = [s for s, _ in VIEWS if s]
    for p in sorted(directory.glob(PREFIX + "*.png")):
        stem = p.stem[len(PREFIX):]
        # Apply the trims (shared drop-list + this deck's own, see dropped()).
        base = stem
        for suf in known:
            if base.endswith(suf):
                base = base[: -len(suf)]
                break
        if dropped(base):
            continue
        for suf in known:
            if stem.endswith(suf):
                found[stem[: -len(suf)]].add(suf)
                break
        else:
            found[stem].add("")
    return found


def group_families(stems):
    """Classify stems into ordered (family, [(stem, title)]) using the shared
    classifier, so titles/ordering match the other chromark decks."""
    fam_items = defaultdict(list)
    unknown = []
    for s in sorted(stems):
        rec = M.classify(s)
        if rec is None:
            unknown.append(s)
            fam_items["Other metrics"].append(((999, s), s, M.fallback_title(s)))
        else:
            fam_items[rec["fam"]].append((rec["sort"], s, rec["title"]))
    ordered = {fam: [(s, t) for _, s, t in sorted(v, key=lambda x: x[0])]
               for fam, v in fam_items.items()}

    # Highlights open the deck: 3D morphology, then 3D intensity, then the
    # maximal-area-slice cross-section area. Moved, not copied.
    highlights = []
    for fam in HIGHLIGHT_FAMS:
        highlights.extend(ordered.pop(fam, []))
    for stem in HIGHLIGHT_EXTRA:
        for fam, items in ordered.items():
            hit = [(s, t) for s, t in items if s == stem]
            if hit:
                highlights.extend(hit)
                ordered[fam] = [(s, t) for s, t in items if s != stem]
                break

    out = []
    if highlights:
        out.append((HIGHLIGHTS_TITLE, highlights))
    for fam in BODY_ORDER:
        if ordered.get(fam):
            out.append((fam, ordered[fam]))
    return out, unknown


def emit_metric(prs, directory, stem, title, available, missing):
    """One full-size slide per available view of this metric."""
    n = 0
    for suf, vlabel in VIEWS:
        if suf not in available:
            continue
        p = directory / "{}{}{}.png".format(PREFIX, stem, suf)
        if not p.exists():
            missing.append(p.name)
            continue
        M.build_slide(prs, "{} — {}".format(title, vlabel), p,
                      "{}/{}".format(directory.name, p.name))
        n += 1
    return n


def text_slide(prs, lines):
    slide = M._new_slide(prs)
    M.add_textbox(slide, lines[0], M.MARGIN, 0.30, M.SLIDE_W - 2 * M.MARGIN, 0.7,
                  font_pt=30, color=M.BLACK, bold=True)
    box = slide.shapes.add_textbox(
        Inches(M.MARGIN + 0.25), Inches(1.15),
        Inches(M.SLIDE_W - 2 * M.MARGIN - 0.5), Inches(6.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines[1:]):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = ln
        para.alignment = PP_ALIGN.LEFT
        if para.runs:
            para.runs[0].font.size = Pt(15)
            para.runs[0].font.color.rgb = M.BLACK


def new_deck(subtitle, title=DECK_TITLE):
    prs = Presentation()
    prs.slide_width = Inches(M.SLIDE_W)
    prs.slide_height = Inches(M.SLIDE_H)
    M.build_title_slide(prs, title, subtitle)
    return prs


def save(prs, out_path, label):
    out_path.parent.mkdir(parents=True, exist_ok=True)
    if out_path.exists():
        M.backup_presentation(str(out_path), backup_base=str(out_path.parent / "backups"))
    prs.save(str(out_path))
    print("[{}] {} slides -> {}".format(label, len(prs.slides._sldIdLst), out_path.name))


def main():
    args = sys.argv[1:]
    list_only = "--list" in args
    selected_only = "--selected" in args
    marks = "--marks" in args

    if marks:
        # Same pipeline, mark channel. MARK_LABEL drives the titles ("Chromatin
        # mark ..." rather than a specific mark) because the 3 columns are 3
        # different marks.
        M.MARK_LABEL = MARKS_LABEL
        allc = discover(DIR_MARKS)
        fams, unknown = group_families(allc)
        print("marks : {} metrics ({})".format(len(allc), DIR_MARKS.name))
        print("unclassified: {}".format(len(unknown)))
        if list_only:
            for fam, items in fams:
                print("=== {} ({}) ===".format(fam, len(items)))
                for stem, title in items:
                    print("   {:<44} {}".format(stem, title))
            if unknown:
                print("\nUNKNOWN:", unknown)
            return
        missing = []
        prs = new_deck(MARKS_SUBTITLE, MARKS_TITLE)
        text_slide(prs, HOW_TO_READ)
        for fam, items in fams:
            M.build_divider_slide(prs, fam)
            print("=== {} ({}) ===".format(fam, len(items)))
            for stem, title in items:
                emit_metric(prs, DIR_MARKS, stem, title, allc[stem], missing)
        save(prs, OUT_MARKS, "marks")
        if missing:
            print("\nMissing {} panel(s): {}".format(len(missing), missing[:8]))
        else:
            print("\nAll panels found.")
        return

    sel = discover(DIR_SEL)
    allc = discover(DIR_ALL)
    all_fams, all_unknown = group_families(allc)

    # The featured set is small and deliberately ordered (aspect ratio leads --
    # it is the headline mechano readout; volume is the constant-volume control).
    # Keep FEATURED_TITLES order rather than family-grouping it.
    sel_items, sel_unknown = [], []
    for stem in M.FEATURED_TITLES:
        if stem in sel:
            sel_items.append((stem, M.FEATURED_TITLES[stem]))
    for stem in sorted(sel):
        if stem not in M.FEATURED_TITLES:
            rec = M.classify(stem)
            sel_items.append((stem, rec["title"] if rec else M.fallback_title(stem)))
            if rec is None:
                sel_unknown.append(stem)
    sel_fams = [("Featured metrics", sel_items)]

    print("featured : {} metrics ({})".format(len(sel), DIR_SEL.name))
    print("curated  : {} metrics ({})".format(len(allc), DIR_ALL.name))
    print("unclassified: featured={} curated={}".format(len(sel_unknown), len(all_unknown)))

    if list_only:
        for tag, fams in (("FEATURED", sel_fams), ("CURATED", all_fams)):
            print("\n########## {} ##########".format(tag))
            for fam, items in fams:
                print("=== {} ({}) ===".format(fam, len(items)))
                for stem, title in items:
                    print("   {:<48} {}".format(stem, title))
        if sel_unknown or all_unknown:
            print("\nUNKNOWN:", sel_unknown + all_unknown)
        # Two stems sharing a title make near-identical slides -- flag them.
        seen = defaultdict(list)
        for fam, items in all_fams:
            for stem, title in items:
                seen[title].append(stem)
        dupes = {t: s for t, s in seen.items() if len(s) > 1}
        print("\nduplicate titles: {}".format(len(dupes)))
        for t, s in dupes.items():
            print("   {!r} <- {}".format(t, s))
        return

    missing = []

    # ---- featured-only deck ----
    if selected_only:
        prs = new_deck(DECK_SUBTITLE + "  -  featured metrics")
        text_slide(prs, HOW_TO_READ)
        for fam, items in sel_fams:
            M.build_divider_slide(prs, fam)
            for stem, title in items:
                emit_metric(prs, DIR_SEL, stem, title, sel[stem], missing)
                print("  {}".format(title))
        save(prs, OUT_SEL, "featured")
        if missing:
            print("\nMissing {} panel(s): {}".format(len(missing), missing[:8]))
        return

    # ---- full deck: the curated body only ----
    # No separate "Featured" section: the 2026-07-27 re-render added aspect_ratio,
    # nuclear_volume and the 3D skew/kurtosis to violin_all_curated, so a featured
    # block would just duplicate them. Morphology leads, and aspect ratio (the
    # headline flatness readout) is first within it, so the deck still opens on it.
    # The featured-only deck is still available via --selected.
    prs = new_deck(DECK_SUBTITLE)
    text_slide(prs, HOW_TO_READ)

    for fam, items in all_fams:
        M.build_divider_slide(prs, fam)
        print("=== {} ({}) ===".format(fam, len(items)))
        for stem, title in items:
            emit_metric(prs, DIR_ALL, stem, title, allc[stem], missing)

    save(prs, OUT_FULL, "full")
    if missing:
        print("\nMissing {} panel(s): {}".format(len(missing), missing[:10]))
    else:
        print("\nAll panels found.")


if __name__ == "__main__":
    main()
