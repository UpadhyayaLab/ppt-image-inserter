"""
insert_cart_nuc_only_summary_20260706_slides.py

Nucleus-ONLY construct-comparison summary deck for Kiet's fixed CAR-T dataset,
compiled into
    Y:/User_data/Kiet/CART_compiled_results/CART_nucleus_across_dates_processed/
        _by_day_violins_20260706
The two conditions are the CAR constructs — CAT (blue) vs FMC63 (red) — shown as
adjacent violins per experiment date/timepoint (12 date×timepoint groups,
Oct 2023 -> Mar 2026), chronological, with a per-pair significance bracket and a
CAT/FMC63 legend.

Companion to insert_cart_catb_nuc_summary_20260705_slides.py — that deck covers
CatB granules + MT + actin + nucleus for the 20260705 compile. This deck is
NUCLEUS-FOCUSED: it carries the VimKD/CTL curated nuclear set (spreading,
deformation & invaginations, orientation, morphology, chromatin/DNA) plus the two
standard cell-spreading context metrics (cell aspect ratio, synapse area). The
Kiet compile has NO centrosome marker pooled across all dates, so — like the VimKD
"all" deck — no centrosome-referenced metric exists; those are omitted.
(The compile emits ~120 by_day panels; FAMILIES is this curated subset.)

Same self-contained layout as the sibling: blank deck (no template), title slide,
family dividers, native PowerPoint sections, footer, --list dry-run, and
backup-before-overwrite. Panels are flat at the compile root (BYDAY_DIR = ROOT),
not in a by_day_panels/ subdir. Missing panels render "no data yet".

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_cart_nuc_only_summary_20260706_slides.py
    conda run -n PPT_editing python examples_and_configs/insert_cart_nuc_only_summary_20260706_slides.py --list
"""

import os
import sys
import uuid
from pathlib import Path
from xml.sax.saxutils import escape

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation  # noqa: E402

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(
    "Y:/User_data/Kiet/CART_compiled_results/"
    "CART_nucleus_across_dates_processed/_by_day_violins_20260706"
)
# By-day violin panels are flat at this dir (one <stem>_by_day.png per metric),
# NOT in a by_day_panels/ subdir like the 20260705 CatB compile.
BYDAY_DIR = ROOT
# cell_counts barplot lives one level up in the processed dir (optional context).
CELL_COUNTS_PNG = ROOT.parent / "cell_counts_barplot.png"

OUTPUT_PATH = Path(
    "K:/FF/PPT/PPT_autogeneration/CART/nucleus/"
    "CART_nuc_only_summary_across_dates_20260706.pptx"
)

PANEL_SUFFIX = "_by_day.png"

# Compile date parsed from the dated ROOT folder (…_YYYYMMDD), shown in the footer.
_d = ROOT.name.rsplit("_", 1)[-1]          # e.g. "20260706"
COMPILE_DATE = "{}-{}-{}".format(_d[:4], _d[4:6], _d[6:8])

DECK_TITLE = "Nuclear morphology in CAR T cells (CAT vs FMC63)"
DECK_SUBTITLE = (
    "CAT vs FMC63 by date/timepoint  ·  nuclear morphology, deformation & chromatin "
    "(+ cell-spreading context)  ·  12 date×timepoint groups, Oct 2023 – Mar 2026  ·  "
    "compiled 2026-07-06"
)

# ---------------------------------------------------------------------------
# Curated metrics, grouped into families (divider slide per family). Each entry
# is (metric stem, slide title) for ONE by_day panel per slide. These are the
# 16 nucleus panels the Kiet compile emits; family names/titles mirror the
# sibling CART / VimKD / CTL summary decks.
# ---------------------------------------------------------------------------
FAMILIES = [
    ("Cell and nuclear spreading", [
        ("nuc_aspect_ratio",        "Nuclear aspect ratio"),
        ("actin_deform_ratio",      "Cell aspect ratio"),
        ("actin_bottom_mask_area",  "Synapse area"),
        ("nuc_broadest_slice_area", "Nuclear broadest-slice area"),
        ("nuc_MIP_area",            "Nuclear MIP area"),
        ("nuc_MIP_MajorAxisLength", "Nuclear MIP major axis length"),
        ("nuc_MIP_MinorAxisLength", "Nuclear MIP minor axis length"),
    ]),
    ("Nuclear deformation and invaginations", [
        ("chull_max_D",                         "Max invag depth over full nucleus"),
        ("chull_mean_D",                        "Mean invag depth over full nucleus"),
        ("deepest_invag_volume",                "Deepest invagination volume"),
        ("deepest_invag_fraction_chull_volume", "Deepest invag: frac of convex hull volume"),
        ("deepest_region_periph_ratio_025um",   "DNA levels near invag"),
    ]),
    ("Invagination orientation", [
        ("avg_normal_angle_adaptive_region_growth", "Deepest invag orientation"),
    ]),
    ("Nuclear morphology", [
        ("nuc_solidity",        "Nuclear solidity"),
        ("nuc_mesh_sphericity", "Nuclear sphericity"),
        ("nuc_volume_mesh",     "Nuclear volume"),
        ("nuc_SA_mesh",         "Nuclear surface area"),
        ("nuc_height",          "Nuclear height"),
    ]),
    ("Chromatin / DNA distribution", [
        ("nuc_MFI",              "DNA mean intensity (MFI)"),
        ("nuc_all_CV",           "DNA intensity CV (heterogeneity)"),
        ("nuc_all_skewness",     "DNA intensity skewness"),
        ("nuc_all_norm_entropy", "DNA distribution normalized entropy"),
    ]),
]

# ---------------------------------------------------------------------------
# Colors / layout (matches the CatB/bleb/noco/vimkd/CTL summary decks)
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.10

TITLE_LEFT = MARGIN
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.55
TITLE_FONT_PT = 28

IMG_LEFT = MARGIN
IMG_TOP = 0.66
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = 6.36

FOOTER_LEFT = MARGIN
FOOTER_TOP = 7.06
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.40
FOOTER_FONT_PT = 9

# ---------------------------------------------------------------------------


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    """Place an image inside (left, top, w, h), preserving aspect ratio and
    centering. The by_day panels are wide/landscape, so they fit to width and
    center vertically."""
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def _ext(path):
    """Windows extended-length path (\\\\?\\...) so files whose absolute path
    exceeds MAX_PATH (260) are still found/opened (the Y: compiled tree is deep).
    No-op off Windows."""
    p = os.path.abspath(str(path))
    if os.name == "nt" and not p.startswith("\\\\?\\"):
        p = "\\\\?\\" + p.replace("/", "\\")
    return p


def _exists(path):
    return os.path.exists(_ext(path))


def rel_footer(path):
    try:
        rel = path.relative_to(ROOT).as_posix()
    except ValueError:
        rel = path.as_posix()
    return "{} / {}".format(COMPILE_DATE, rel)


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.7, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=40, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.1, SLIDE_W - 2 * MARGIN, 1.2,
                font_pt=16, color=BLACK, italic=True)


def build_divider_slide(prs, family_name):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, family_name, MARGIN, 3.1, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=44, color=BLACK, bold=True)


def build_slide(prs, title_text, image_path, footer_text):
    """Title + one by_day panel (aspect preserved) + source-path footer.
    Returns True if the panel was missing on disk."""
    slide = _new_slide(prs)
    add_textbox(slide, title_text, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title_text), color=BLACK, bold=True)
    missing = not _exists(image_path)
    if not missing:
        add_image_in_box(slide, _ext(image_path), IMG_LEFT, IMG_TOP, IMG_BOX_W, IMG_BOX_H)
    else:
        add_textbox(slide, "no data yet", IMG_LEFT, IMG_TOP + IMG_BOX_H / 2 - 0.2,
                    IMG_BOX_W, 0.4, font_pt=18, color=BLACK)
    add_textbox(slide, footer_text, FOOTER_LEFT, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=BLACK)
    return missing


def _panel_path(stem):
    return BYDAY_DIR / (stem + PANEL_SUFFIX)


def add_sections(prs, section_spec):
    """Native PowerPoint sections (named groups in the slide navigator).
    python-pptx has no API for this, so inject the p14:sectionLst extension.
    `section_spec` is an ordered list of (name, n_slides) summing to the slide
    count in build order; the first section must contain the first slide."""
    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    slide_ids = [sldId.get("id") for sldId in prs.slides._sldIdLst]

    parts = ['<p:extLst xmlns:p="{}">'.format(P),
             '<p:ext uri="{{521415D9-36F7-43E2-AB2F-B90AF26B5E84}}">',
             '<p14:sectionLst xmlns:p14="{}">'.format(P14)]
    i = 0
    for name, count in section_spec:
        parts.append('<p14:section name="{}" id="{{{}}}">'.format(
            escape(name), str(uuid.uuid4()).upper()))
        parts.append('<p14:sldIdLst>')
        for sid in slide_ids[i:i + count]:
            parts.append('<p14:sldId id="{}"/>'.format(sid))
        parts.append('</p14:sldIdLst></p14:section>')
        i += count
    parts.append('</p14:sectionLst></p:ext></p:extLst>')
    prs.slides._sldIdLst.getparent().append(parse_xml("".join(parts)))


def main():
    list_only = "--list" in sys.argv

    n_metrics = sum(len(items) for _, items in FAMILIES)
    est_slides = 1 + (1 if _exists(CELL_COUNTS_PNG) else 0) + \
        sum(1 + len(items) for _, items in FAMILIES)

    print("Source: {}".format(BYDAY_DIR))
    print("{} metric slides across {} families, est. {} slides\n".format(
        n_metrics, len(FAMILIES), est_slides))

    if list_only:
        for fam, items in FAMILIES:
            print("=== {} ({}) ===".format(fam, len(items)))
            for stem, title in items:
                ok = _exists(_panel_path(stem))
                print("  [{}] {:<50s} {}".format(
                    "OK " if ok else "MISS", title, stem + PANEL_SUFFIX))
            print("")
        return

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)

    if _exists(CELL_COUNTS_PNG):
        build_slide(prs, "Cell counts (CAT vs FMC63)", CELL_COUNTS_PNG,
                    rel_footer(CELL_COUNTS_PNG))
    else:
        print("Note: {} not found - skipping cell-counts slide.\n".format(
            CELL_COUNTS_PNG.name))

    section_spec = [("Overview", len(prs.slides._sldIdLst))]

    missing = []
    for fam, items in FAMILIES:
        build_divider_slide(prs, fam)
        print("=== {} ===".format(fam))
        for stem, title in items:
            p = _panel_path(stem)
            is_missing = build_slide(prs, title, p, rel_footer(p))
            print("  [{}] {} -> {!r}".format(
                "OK" if not is_missing else "MISSING", p.name, title))
            if is_missing:
                missing.append(p.name)
        section_spec.append((fam, 1 + len(items)))
        print("")

    add_sections(prs, section_spec)

    if OUTPUT_PATH.exists():
        backup_dir = OUTPUT_PATH.parent / "backups"
        created = backup_presentation(str(OUTPUT_PATH), backup_base=str(backup_dir))
        if created:
            print("Backed up previous deck to: {}\n".format(backup_dir))

    prs.save(str(OUTPUT_PATH))
    total = len(prs.slides._sldIdLst)
    print("Done. {} metric slides, {} slides written to:\n  {}".format(
        n_metrics, total, OUTPUT_PATH))
    if missing:
        print("\nSkipped {} missing panel(s):".format(len(missing)))
        for m in missing:
            print("  - {}".format(m))
    else:
        print("\nAll curated panels found - no missing items.")


if __name__ == "__main__":
    main()
