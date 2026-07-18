"""
insert_master_vim_activation_vs_drug_slides.py

Master vimentin deck: for each metric, ONE slide showing BOTH comparisons of
vimentin vs nuclear-envelope geometry side by side —
  activation : aCD3 12 min vs PLL 12 min (MG compile, 2 exps: Dec 12 2024,
               Mar 18 2024)
  drug       : DMSO vs Noco 1 μM (Noco_CD3 compile, 3 exps: Apr 29 2022,
               Jan 23 2024, Feb 27 2024)

Per experiment, nothing pooled — each tile is labelled by comparison + date.
Same layout engine and topic set as insert_master_struct_vs_NE_geometry_slides.py.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_master_vim_activation_vs_drug_slides.py
    conda run -n PPT_editing python examples_and_configs/insert_master_vim_activation_vs_drug_slides.py --list
"""

import math
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

try:
    from PIL import Image
except ImportError:  # pragma: no cover
    Image = None

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation, safe_path, path_exists  # noqa: E402

# ---------------------------------------------------------------------------
# Channel source roots
# ---------------------------------------------------------------------------
# Two vimentin comparisons, one metric per slide:
#   ACT  = activation  (aCD3 12 min vs PLL 12 min), MG compile, 2 experiments
#   DRUG = drug        (DMSO vs Noco 1 μM), Noco_CD3 compile, 3 experiments
ACT_ROOT = Path(
    "J:/FF/vim_data/20240318_MG_E6-1_Vimentin/"
    "results_compilation_Vim_MG_geomdensity_combined_20260716")
DRUG_ROOT = Path(
    "J:/FF/fixed_cell/Vimentin/results_across_experiments/"
    "Noco_CD3_20220429_20240123_20240227_20260717")

OUTPUT_PATH = Path(
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/"
    "nuc_mesh_struct_outside_nuc/"
    "Master, vimentin vs NE geometry — activation vs drug "
    "(aCD3-PLL & DMSO-Noco).pptx")

DECK_TITLE = "Vimentin vs nuclear-envelope geometry — activation vs drug"
DECK_SUBTITLE = (
    "Fixed Jurkats · vim = perinuc 0.5 μm · each slide shows both comparisons: "
    "aCD3 12 min vs PLL 12 min (activation) and DMSO vs Noco 1 μM (drug) · "
    "per experiment, nothing pooled")

ACT_DATES = [("Dec 12 2024", "Dec_12,_2024"),
             ("Mar 18 2024", "Mar_18,_2024")]
DRUG_DATES = [("Apr 29 2022", "Apr_29,_2022"),
              ("Jan 23 2024", "Jan_23,_2024"),
              ("Feb 27 2024", "Feb_27,_2024")]

# ---------------------------------------------------------------------------
# Colors / layout (16:9)
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0x55, 0x55, 0x55)
FIELD_COLOR = RGBColor(0x2E, 0x5A, 0x88)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.08
GAP = 0.10

TITLE_LEFT = MARGIN
TITLE_TOP = 0.04
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 27

IMG_TOP = 0.58
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = 6.60

FOOTER_TOP = 7.20
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.28
FOOTER_FONT_PT = 8

CAPTION_HEIGHT = 0.50
MAX_GRID_COLS = 5


# ---------------------------------------------------------------------------
# Slide helpers (shared with the per-channel decks)
# ---------------------------------------------------------------------------
def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER, wrap=True):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    image_path = safe_path(image_path)
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.6, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=34, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.0, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=15, color=GREY, italic=True)


def build_divider_slide(prs, title, subtitle=""):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, title, MARGIN, 3.0, SLIDE_W - 2 * MARGIN, 1.2,
                font_pt=38, color=BLACK, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, MARGIN, 4.25, SLIDE_W - 2 * MARGIN, 0.9,
                    font_pt=18, color=GREY, italic=True)


def add_labeled_caption(slide, lines, left, top, width, height):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    for i, (text, pt, color) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.size = Pt(pt)
        run.font.color.rgb = color
    return box


def _tile_caption(slide, cap, path, left, top, width, height, cap_pt):
    lines = []
    if cap:
        lines.append((cap, min(cap_pt, 10), GREY))
    lines.append((Path(str(path)).stem, 8, FIELD_COLOR))
    add_labeled_caption(slide, lines, left, top, width, height)


_AR_CACHE = {}


def _img_ar(path):
    key = str(path)
    if key in _AR_CACHE:
        return _AR_CACHE[key]
    ar = 1.5
    if Image is not None:
        try:
            with Image.open(safe_path(key)) as im:
                w, h = im.size
                if h:
                    ar = w / float(h)
        except Exception:
            pass
    _AR_CACHE[key] = ar
    return ar


def _uniform_img_size(a, cols, rows, area_w, area_h, cap_h):
    cw = (area_w - (cols - 1) * GAP) / cols
    cell_h = (area_h - (rows - 1) * GAP) / rows
    ih_avail = cell_h - cap_h
    if ih_avail <= 0.2 or cw <= 0.2:
        return None
    iw = min(cw, ih_avail * a)
    return iw, iw / a


def _best_cols(ars, n, area_w, area_h, cap_h, cap_cols):
    uniform = (max(ars) / min(ars)) < 1.20
    a = sorted(ars)[len(ars) // 2]
    best = None
    for cols in range(1, min(cap_cols, n) + 1):
        rows = math.ceil(n / cols)
        if uniform:
            sz = _uniform_img_size(a, cols, rows, area_w, area_h, cap_h)
            if sz is None:
                continue
            score = sz[0] * sz[1]
        else:
            cw = (area_w - (cols - 1) * GAP) / cols
            cell_h = (area_h - (rows - 1) * GAP) / rows
            ih_avail = cell_h - cap_h
            if ih_avail <= 0.2 or cw <= 0.2:
                continue
            score = 0.0
            for ai in ars:
                iw = min(cw, ih_avail * ai)
                score += iw * (iw / ai)
        if best is None or score > best[0] + 1e-6:
            best = (score, cols, rows, uniform)
    if best is None:
        c = min(cap_cols, n)
        return c, math.ceil(n / c), uniform
    return best[1], best[2], best[3]


def build_grid_slide(prs, title, entries, footer_text, cap_pt=11, cap_h=0.30):
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    n = len(entries)
    area_top = IMG_TOP
    area_w = SLIDE_W - 2 * MARGIN
    area_h = FOOTER_TOP - IMG_TOP - 0.02

    ars = [_img_ar(p) for p, _ in entries]
    cap_cols = min(MAX_GRID_COLS, n)
    cols, rows, uniform = _best_cols(ars, n, area_w, area_h, cap_h, cap_cols)

    if uniform:
        a = sorted(ars)[len(ars) // 2]
        iw, ih = _uniform_img_size(a, cols, rows, area_w, area_h, cap_h)
        block_w = cols * iw + (cols - 1) * GAP
        block_h = rows * (ih + cap_h) + (rows - 1) * GAP
        x0 = MARGIN + (area_w - block_w) / 2
        y0 = area_top + max(0.0, (area_h - block_h) / 2)
        for i, (path, cap) in enumerate(entries):
            r, c = divmod(i, cols)
            left = x0 + c * (iw + GAP)
            top = y0 + r * (ih + cap_h + GAP)
            add_image_in_box(slide, str(path), left, top, iw, ih)
            _tile_caption(slide, cap, path, left, top + ih, iw, cap_h, cap_pt)
    else:
        cell_w = (area_w - (cols - 1) * GAP) / cols
        cell_h = (area_h - (rows - 1) * GAP) / rows
        img_h = cell_h - cap_h
        for i, (path, cap) in enumerate(entries):
            r, c = divmod(i, cols)
            left = MARGIN + c * (cell_w + GAP)
            top = area_top + r * (cell_h + GAP)
            add_image_in_box(slide, str(path), left, top, cell_w, img_h)
            _tile_caption(slide, cap, path, left, top + img_h,
                          cell_w, cap_h, cap_pt)

    add_textbox(slide, footer_text, MARGIN, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


# ---------------------------------------------------------------------------
# Per-channel file maps.  Each accessor returns a list of (path, caption).
# ---------------------------------------------------------------------------
def _vim_comparison(root, name, sub, dates):
    """One vimentin comparison (activation OR drug). `dates` = list of
    (display, comma_token) per experiment; BYEXP images combine experiments."""
    S = root / "geom_density" / "profiles" / "singles"
    E = root / "geom_density" / "enrichment"
    NC = root / "geom_density" / "near_cent"
    G = root / "grid_panels"
    D = root / "invag_depth_profiles" / "singles"
    exp_note = " · ".join(disp for disp, _ in dates)
    return {
        "name": name,
        "sub": sub,
        "density": lambda g: [
            (S / "vim_geomdens_{}_perinuc05_BYEXP_line.png".format(g), exp_note)],
        "enrich": lambda: [
            (E / "vim_hulldist_gt0.5um_shells_{}.png".format(d), disp)
            for disp, d in dates],
        "corr": lambda: [
            (E / "vim_corr_hulldist_shells_{}.png".format(d), disp)
            for disp, d in dates],
        "pocket": lambda: [
            (G / "vim_cyto_in_nuc_hull_vs_near_convex_nuc_MFI_ratio_grid.png",
             sub)],
        "nearcent": lambda: [
            (NC / "vim_near_cent_level_{}.png".format(d), disp)
            for disp, d in dates],
        "depth": lambda: [
            (D / "vim_invag_depth_0_5um_BYEXP_line.png", exp_note)],
    }


CHANNELS = [
    _vim_comparison(ACT_ROOT, "aCD3 vs PLL",
                    "activation · 12 min", ACT_DATES),
    _vim_comparison(DRUG_ROOT, "DMSO vs Noco",
                    "drug · Noco 1 μM", DRUG_DATES),
]

# Topic: (accessor key, geom-or-None, divider title, divider subtitle)
TOPICS = [
    ("density", "hulldist", "Density vs hull-boundary distance",
     "how far into the perinuclear region each protein sits vs distance "
     "from the nuclear surface"),
    ("density", "mincurv", "Density vs minimum curvature",
     "signal vs concave (groove) / convex nuclear-envelope curvature"),
    ("density", "meancurv", "Density vs mean curvature",
     "signal vs mean nuclear-envelope curvature"),
    ("enrich", None, "Enrichment near invaginations (dist > 0.5 μm)",
     "signal fraction ÷ expected in shells around deep invaginations"),
    ("corr", None, "Per-cell correlation vs hull distance",
     "per-cell correlation of signal with distance into invaginations"),
    ("pocket", None, "Signal in nuclear invagination pockets",
     "grooves ÷ convex nuclear surface (MFI ratio)"),
    ("nearcent", None, "Level near the centrosome",
     "signal at the nuclear envelope facing the centrosome"),
    ("depth", None, "Invagination depth profiles",
     "signal vs depth within nuclear invaginations"),
]

CHAN_ORDER_NOTE = "aCD3-vs-PLL (activation) then DMSO-vs-Noco (drug)"

# Both comparisons use clean BYEXP strips / violins (no dense monster figures),
# so every metric condenses onto one slide.
HUGE_TOPICS = set()


# ---------------------------------------------------------------------------
# Plan
# ---------------------------------------------------------------------------
def build_plan():
    plan, missing = [], []
    for key, geom, dtitle, dsub in TOPICS:
        short = dtitle.split(" (")[0]
        chan_tiles = []
        for ch in CHANNELS:
            tiles = ch[key](geom) if key == "density" else ch[key]()
            kept = [(p, c) for p, c in tiles if path_exists(p)]
            missing.extend(p for p, _ in tiles if not path_exists(p))
            if kept:
                chan_tiles.append((ch, kept))
        if not chan_tiles:
            continue

        if key in HUGE_TOPICS:
            # One slide per channel (dense figures need the whole slide).
            plan.append(("divider", dtitle,
                         dsub + " · one slide per channel · " + CHAN_ORDER_NOTE))
            for ch, kept in chan_tiles:
                plan.append(("grid", "{} — {}".format(short, ch["name"]), kept,
                             "{} · {} · perinuc 0.5 μm".format(
                                 ch["name"], ch["sub"])))
        else:
            # All channels on one slide, each tile labelled by channel.
            combined = []
            for ch, kept in chan_tiles:
                for p, c in kept:
                    combined.append((p, "{} · {}".format(ch["name"], c)))
            plan.append(("grid", "{}  ·  activation vs drug".format(short), combined,
                         dsub + " · " + CHAN_ORDER_NOTE + " · perinuc 0.5 μm"))
    return plan, missing


def main():
    list_only = "--list" in sys.argv
    plan, missing = build_plan()

    print("Output: {}".format(OUTPUT_PATH))
    print("{} content slides (+ title) = {} total\n".format(
        len(plan), 1 + len(plan)))
    for it in plan:
        if it[0] == "divider":
            print("\n=== {} ===".format(it[1]))
        else:
            print("  [{}] {}".format(len(it[2]), it[1]))
    if missing:
        print("\nMISSING ({}):".format(len(missing)))
        for m in missing:
            print("  - {}".format(m))
    if list_only:
        return

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)

    for it in plan:
        if it[0] == "divider":
            build_divider_slide(prs, it[1], it[2])
        else:
            _, title, entries, footer = it
            build_grid_slide(prs, title, entries, footer)

    if OUTPUT_PATH.exists():
        backup_dir = OUTPUT_PATH.parent / "backups"
        created = backup_presentation(str(OUTPUT_PATH),
                                      backup_base=str(backup_dir))
        if created:
            print("\nBacked up previous deck under: {}".format(backup_dir))

    prs.save(str(OUTPUT_PATH))
    total = len(prs.slides._sldIdLst)
    print("\nDone. {} slides written to:\n  {}".format(total, OUTPUT_PATH))
    if missing:
        print("Skipped {} missing panel(s).".format(len(missing)))


if __name__ == "__main__":
    main()
