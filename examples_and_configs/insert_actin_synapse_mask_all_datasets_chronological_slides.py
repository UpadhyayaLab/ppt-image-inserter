"""
insert_actin_synapse_mask_all_datasets_chronological_slides.py

CAT (left) vs FMC (right) side-by-side actin-synapse-mask deck across
every (date, marker, timepoint) group in the compiled manifest at

    L:/FF/CAR T/actin_compiled_results/all_datasets_actin/compiled_20260623/
      dataset_manifest.csv

Manifest rows (67) get grouped by (date, marker_label, timepoint_minutes).
For each group: CAT row's first chunk -> left cell, FMC row's first chunk
-> right cell, both at the same deck-wide PPI so the 5 μm scalebar is the
same cm on every slide.

For each row, resolves the actin synapse mask montages dir by probing:

    <base_dir>/<progress_folder>/actin/synapse/1slice/mask/montages/   (Stage AG)
    <base_dir>/<progress_folder>/actin/synapse/mask/montages/          (pre-AG)

Includes a Windows long-path (\\?\) shim — several manifest paths exceed
260 chars and would otherwise silently fail `.exists()` / `open()` while
`.glob()` still works.

Usage:
    python examples_and_configs/insert_actin_synapse_mask_all_datasets_chronological_slides.py
"""

import csv
import os
import re
import sys
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

MANIFEST_CSV = (
    "L:/FF/CAR T/actin_compiled_results/all_datasets_actin/"
    "compiled_20260623/dataset_manifest.csv"
)

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/CART/actin_only/"
    "CART_actin_synapse_mask_all_datasets_20260623.pptx"
)

# Stage AG layout first, pre-AG fallback. Both joined under
# `<base_dir>/<progress_folder>/actin/synapse/`.
KIND_SUBPATHS = [
    "1slice/mask/montages",
    "mask/montages",
]

CHUNK_GLOB = "montage_cells_*.png"

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Slide layout (inches). 13.333 x 7.5 widescreen.
# Tighter than the original MT/CatB compare layout: smaller title + label
# fonts and smaller margins so each cell gets a bigger image area.
SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.05
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.05
TITLE_HEIGHT = 0.40
TITLE_FONT_PT = 24

# 1x2 cell grid below the title (label + image per cell).
GRID_LEFT = 0.05
GRID_TOP = 0.50
CELL_W = 6.60
CELL_H = SLIDE_H - GRID_TOP - 0.05    # 6.95"
LABEL_H = 0.22
IMG_H = CELL_H - LABEL_H              # 6.73"
LABEL_FONT_PT = 14
COL_GAP = SLIDE_W - 2 * GRID_LEFT - 2 * CELL_W

CELL_POSITIONS = [
    (GRID_LEFT,                    GRID_TOP),  # left = CAT
    (GRID_LEFT + CELL_W + COL_GAP, GRID_TOP),  # right = FMC
]

# Scalebar invariant (Stage AF/AG).
PPUM_SOURCE = 30
SCALEBAR_UM = 5
SCALEBAR_PX = PPUM_SOURCE * SCALEBAR_UM   # 150 px

# Condition / experiment parsing.
# `FMC(?!63)` prevents `FMC` from matching at the start of `FMC63` and
# consuming `63` as the minutes digits.
_COND_RE = re.compile(r"(CAT|FMC63|FMC(?!63))_?(\d+)(?:min)?", re.IGNORECASE)
_DTAG_RE = re.compile(r"_D(\d+)_", re.IGNORECASE)


def _long_path_str(p: Path) -> str:
    """Return a path string that bypasses Windows' 260-char MAX_PATH limit.
    Several manifest entries land 260+ chars deep; `.glob()` works but
    `.exists()` / `open()` / `Image.open()` silently fail without the
    `\\?\` prefix."""
    s = str(p)
    if os.name == "nt" and len(s) >= 240 and not s.startswith("\\\\?\\"):
        s = s.replace("/", "\\")
        if not s.startswith("\\\\"):
            return "\\\\?\\" + s
    return s


def parse_condition(base_dir: str) -> Tuple[Optional[str], Optional[int]]:
    """Find the first CAT|FMC63|FMC + digits match in the path. Returns
    (cell_type_normalized, minutes) or (None, None). FMC63 is normalized
    to "FMC" — the dataset identity already encodes which CAR variant."""
    m = _COND_RE.search(base_dir)
    if not m:
        return (None, None)
    cell, mins = m.group(1).upper(), int(m.group(2))
    if cell == "FMC63":
        cell = "FMC"
    return (cell, mins)


def parse_dtag(base_dir: str) -> str:
    """Return ' D3' / ' D5' suffix if `_D<n>_` appears in the base_dir."""
    m = _DTAG_RE.search(base_dir)
    return f" D{m.group(1)}" if m else ""


def resolve_montages_dir(base_dir: str, progress_folder: str) -> Optional[Path]:
    """Try Stage-AG layout first, fall back to pre-AG. Returns the first
    existing dir, or None."""
    parent = Path(base_dir) / progress_folder / "actin" / "synapse"
    for kind in KIND_SUBPATHS:
        cand = parent / kind
        if cand.is_dir():
            return cand
    return None


def _parse_chunk_range(p: Path) -> Tuple[int, int]:
    """Return (start, end) cell-id range for a montage_cells_*.png filename.
    Handles both 4-int FOV-padded and 2-int patterns."""
    m4 = re.match(r"montage_cells_(\d+)_(\d+)_(\d+)_(\d+)\.png$", p.name)
    if m4:
        f_a, c_a, f_b, c_b = (int(x) for x in m4.groups())
        return (f_a * 1000 + c_a, f_b * 1000 + c_b)
    m2 = re.match(r"montage_cells_(\d+)_(\d+)\.png$", p.name)
    if m2:
        return (int(m2.group(1)), int(m2.group(2)))
    return (0, 0)


def find_first_chunk(montages_dir: Path) -> Optional[Path]:
    """Pick the lowest-start chunk, dropping any whose [start, end] range is
    strictly contained in another chunk's range (smoke-shadow filter)."""
    if not montages_dir.is_dir():
        return None
    chunks = list(montages_dir.glob(CHUNK_GLOB))
    if not chunks:
        return None
    parsed = [(p, *_parse_chunk_range(p)) for p in chunks]
    keep = []
    for (p, s, e) in parsed:
        shadowed = any(
            s2 <= s and e <= e2 and (s2 < s or e < e2)
            for (p2, s2, e2) in parsed
            if p2 is not p
        )
        if not shadowed:
            keep.append((p, s, e))
    if not keep:
        return None
    keep.sort(key=lambda x: x[1])
    return keep[0][0]


def _png_dims(path: Path) -> Tuple[int, int]:
    with Image.open(_long_path_str(path)) as im:
        return im.size


def _path_exists(p: Path) -> bool:
    return os.path.exists(_long_path_str(p))


def compute_deck_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_image_in_cell_at_ppi(slide, image_path: Path, ppi: float,
                             cell_left: float, cell_top: float):
    """Place image inside a labelled cell at uniform deck px/inch. Both
    dims = native_px / ppi; image centered in (CELL_W x IMG_H) area below
    the label."""
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    img_area_top = cell_top + LABEL_H
    left_in = cell_left + (CELL_W - w_in) / 2
    top_in  = img_area_top + (IMG_H - h_in) / 2
    return slide.shapes.add_picture(
        _long_path_str(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(w_in),
    )


def build_compare_slide(prs, title_text: str,
                        cat_img: Optional[Path], fmc_img: Optional[Path],
                        deck_ppi: float):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    cells = [
        ("CAT", cat_img, CELL_POSITIONS[0]),
        ("FMC", fmc_img, CELL_POSITIONS[1]),
    ]
    missing = []
    for label, img_path, (cell_left, cell_top) in cells:
        add_textbox(
            slide, label,
            cell_left, cell_top, CELL_W, LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )
        if img_path is not None and _path_exists(img_path):
            add_image_in_cell_at_ppi(slide, img_path, deck_ppi, cell_left, cell_top)
        else:
            add_textbox(
                slide, "(missing)",
                cell_left, cell_top + LABEL_H + IMG_H / 2 - 0.15, CELL_W, 0.3,
                font_pt=14, color=WHITE,
            )
            missing.append(label)
    return slide, missing


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    manifest_path = Path(MANIFEST_CSV)
    if not manifest_path.exists():
        print(f"ERROR: manifest not found at {manifest_path}")
        return

    with manifest_path.open("r", newline="") as f:
        rows = list(csv.DictReader(f))

    # Per-row parse: derive cell / tp / marker_label / montages_dir / chunk.
    parsed_rows = []
    for idx, row in enumerate(rows):
        date_tag = row["date"].strip()
        marker = row["marker"].strip()
        progress_folder = row["progress_folder"].strip()
        base_dir = row["base_dir"].strip().rstrip("\\/")

        cell, tp = parse_condition(base_dir)
        dtag = parse_dtag(base_dir)
        marker_label = f"{marker}{dtag}"

        montages_dir = resolve_montages_dir(base_dir, progress_folder)
        chunk = find_first_chunk(montages_dir) if montages_dir is not None else None

        parsed_rows.append({
            "idx": idx,
            "date_tag": date_tag,
            "marker_label": marker_label,
            "cell": cell,
            "tp": tp,
            "montages_dir": montages_dir,
            "chunk": chunk,
        })

    # Group by (date_tag, marker_label, tp). One slide per group.
    groups: Dict[Tuple[str, str, int], Dict] = {}
    for p in parsed_rows:
        if p["cell"] is None or p["tp"] is None:
            print(f"WARNING: could not parse condition from manifest row {p['idx']}; skipping.")
            continue
        key = (p["date_tag"], p["marker_label"], p["tp"])
        g = groups.setdefault(key, {
            "cat_chunk": None, "fmc_chunk": None,
            "cat_dir": None, "fmc_dir": None,
            "first_idx": p["idx"],
        })
        if p["cell"] == "CAT":
            g["cat_chunk"] = p["chunk"]
            g["cat_dir"] = p["montages_dir"]
        elif p["cell"] == "FMC":
            g["fmc_chunk"] = p["chunk"]
            g["fmc_dir"] = p["montages_dir"]

    # Sort groups chronologically: by date, then by manifest order of the
    # first row that landed in the group. Preserves the experimenter's
    # ordering of conditions within a date.
    sorted_groups = sorted(
        groups.items(),
        key=lambda kv: (kv[0][0], kv[1]["first_idx"]),
    )

    # Build slide specs.
    slide_specs = []
    for ((date_tag, marker_label, tp), g) in sorted_groups:
        title = (
            f"Actin at Synapse — {date_tag} ({marker_label}): {tp} min"
        )
        log_key = f"{date_tag}/{marker_label}/{tp}min"
        slide_specs.append((
            title, g["cat_chunk"], g["fmc_chunk"],
            g["cat_dir"], g["fmc_dir"], log_key,
        ))

    # Deck-wide PPI.
    present: List[Path] = []
    for (_, cat_img, fmc_img, _, _, _) in slide_specs:
        for p in (cat_img, fmc_img):
            if p is not None and _path_exists(p):
                present.append(p)

    if not present:
        print("WARNING: no real images found - using fallback PPI=100.")
        deck_ppi = 100.0
    else:
        deck_ppi = compute_deck_ppi(present, CELL_W, IMG_H)

    bar_in = SCALEBAR_PX / deck_ppi
    print(
        f"Manifest rows: {len(rows)}  -  groups (slides): {len(slide_specs)}\n"
        f"Present cells: {len(present)} / {2 * len(slide_specs)}\n"
        f"Deck-wide PPI = {deck_ppi:.2f} -> "
        f"5 μm = {bar_in:.3f} in = {bar_in * 2.54:.3f} cm on every cell.\n"
        f"Source PPUM = {PPUM_SOURCE} px/μm (locked).\n"
    )
    print(f"Writing deck to: {OUTPUT_PATH}\n")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    missing_total = []
    slides_added = 0
    for (title, cat_img, fmc_img, cat_dir, fmc_dir, log_key) in slide_specs:
        _, missing = build_compare_slide(prs, title, cat_img, fmc_img, deck_ppi)
        slides_added += 1

        status_parts = [
            "CAT:OK" if cat_img and _path_exists(cat_img) else "CAT:MISSING",
            "FMC:OK" if fmc_img and _path_exists(fmc_img) else "FMC:MISSING",
        ]
        print(f"[{log_key}]  " + "  ".join(status_parts))

        for cell in missing:
            src = cat_dir if cell == "CAT" else fmc_dir
            missing_total.append(f"{log_key}/{cell}  ({src})")

    prs.save(str(out_path))
    print(f"\nDone. {slides_added} slides written to:\n  {out_path}")

    if missing_total:
        print(f"\nMissing ({len(missing_total)}):")
        for m in missing_total:
            print(f"  - {m}")
    else:
        print("\nAll images found - no missing items.")


if __name__ == "__main__":
    main()
