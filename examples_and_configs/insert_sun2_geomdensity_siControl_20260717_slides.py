"""
insert_sun2_geomdensity_siControl_20260717_slides.py

Build the SUN2 geometry-density summary deck from the 20260717 compile.

Source:
  Jurkats_SUN2_20220704_siControl_geomdensity_20260717  (single condition,
  siControl only, one experiment)

This is the single-condition analogue of the SUN1 N1-vs-N2 deck
(insert_sun1_geomdensity_N1vsN2_20260715_slides.py).  Because SUN2 is a
single replicate there is no N1-vs-N2 comparison, no grid_panels/key_figures,
and morphology_scatter/near_cent are empty; the scalar violins in
violin_plots/ take the place of the SUN1 grid_panels/N1_vs_N2 scalar panels.

Sections:
  1. SUN2 density profiles (density vs geometry)
  2. Curated enrichment & correlation + shell-resolved correlation
  3. Invagination depth profiles
  4. SUN2 scalars (violin_plots — one violin per metric)
  5. QC raw images (top / low ranked cells)

Self-contained blank 16:9 deck.  Missing panels are skipped; a previous deck
at the output path is backed up first.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_sun2_geomdensity_siControl_20260717_slides.py
    conda run -n PPT_editing python examples_and_configs/insert_sun2_geomdensity_siControl_20260717_slides.py --list
"""

import math
import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation, safe_path, path_exists  # noqa: E402

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(
    "O:/FF_backup/Jurkat_nucleus/from_Ivan_HD/0704-2022 - Sun2 AND Ycomp/"
    "results_compilation/Jurkats_SUN2_20220704_siControl_geomdensity_20260717"
)
CURATED = ROOT / "SUN2_loc_wrto_invag"
ENRICH = ROOT / "geom_density" / "enrichment"
PROFILES = ROOT / "geom_density" / "profiles"
POOLED = ROOT / "geom_density" / "pooled"
DEPTH = ROOT / "invag_depth_profiles"
VIOLINS = ROOT / "violin_plots"

OUTPUT_PATH = Path(
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/SUN/"
    "SUN2 vs NE geometry, siControl (20260717).pptx"
)

DECK_TITLE = "SUN2 density vs nuclear-envelope geometry"
DECK_SUBTITLE = (
    "Fixed Jurkats · siControl · 07/04/2022 · "
    "SUN2 = perinuc 0.5 μm outside NE · compiled 2026-07-17"
)

# ---------------------------------------------------------------------------
# Colors / layout (16:9)
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0x55, 0x55, 0x55)
FIELD_COLOR = RGBColor(0x2E, 0x5A, 0x88)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.10
GAP = 0.16

TITLE_LEFT = MARGIN
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.55
TITLE_FONT_PT = 28

IMG_TOP = 0.66
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = 6.36

FOOTER_TOP = 7.06
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.40
FOOTER_FONT_PT = 9

CAPTION_HEIGHT = 0.55

QC_MAX_COLS = 5


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------
def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    image_path = safe_path(image_path)
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.7, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=38, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.1, SLIDE_W - 2 * MARGIN, 1.1,
                font_pt=16, color=GREY, italic=True)


def build_divider_slide(prs, title, subtitle=""):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, title, MARGIN, 3.0, SLIDE_W - 2 * MARGIN, 1.2,
                font_pt=40, color=BLACK, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, MARGIN, 4.25, SLIDE_W - 2 * MARGIN, 0.9,
                    font_pt=18, color=GREY, italic=True)


def build_single_slide(prs, title, image_path, caption, footer=None):
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    cap_h = CAPTION_HEIGHT if caption else 0.0
    img_h = IMG_BOX_H - cap_h
    add_image_in_box(slide, str(image_path), MARGIN, IMG_TOP, IMG_BOX_W, img_h)
    if caption:
        add_textbox(slide, caption, MARGIN, IMG_TOP + img_h, IMG_BOX_W, cap_h,
                    font_pt=12, color=GREY)
    add_textbox(slide, footer or rel_footer(image_path), MARGIN, FOOTER_TOP,
                FOOTER_WIDTH, FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


def add_labeled_caption(slide, lines, left, top, width, height):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    for i, (text, pt, color) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.size = Pt(pt)
        run.font.color.rgb = color
    return box


def _tile_caption(slide, cap, path, left, top, width, height, cap_pt):
    lines = []
    if cap:
        lines.append((cap, min(cap_pt, 10), GREY))
    lines.append((Path(str(path)).stem, 8, FIELD_COLOR))
    add_labeled_caption(slide, lines, left, top, width, height)


def build_grid_slide(prs, title, entries, footer_text, max_cols,
                     cap_pt=11, cap_h=0.28):
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    n = len(entries)
    cols = min(max_cols, n)
    rows = math.ceil(n / cols)
    area_top = IMG_TOP
    area_h = FOOTER_TOP - IMG_TOP - 0.02
    cell_w = (SLIDE_W - 2 * MARGIN - (cols - 1) * GAP) / cols
    cell_h = (area_h - (rows - 1) * GAP) / rows
    img_h = cell_h - cap_h
    for i, (path, cap) in enumerate(entries):
        r, c = divmod(i, cols)
        left = MARGIN + c * (cell_w + GAP)
        top = area_top + r * (cell_h + GAP)
        add_image_in_box(slide, str(path), left, top, cell_w, img_h)
        _tile_caption(slide, cap, path, left, top + img_h, cell_w, cap_h, cap_pt)
    add_textbox(slide, footer_text, MARGIN, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


def rel_footer(path):
    try:
        return Path(path).relative_to(ROOT).as_posix()
    except ValueError:
        return Path(path).name


# ---------------------------------------------------------------------------
# QC helpers (single-condition: siControl_top1_cell57.png)
# ---------------------------------------------------------------------------
def parse_qc(fname):
    """siControl_top1_cell57.png -> (cond, cond_disp, sortkey, caption)."""
    stem = fname[:-4] if fname.lower().endswith(".png") else fname
    m = re.match(r"^(.+?)_(top|low)(\d+)_cell(\d+)$", stem)
    if not m:
        return None
    cond, rank, rnum, cell = m.groups()
    cond_disp = cond.replace("_", " ")
    sortkey = (0 if rank == "top" else 1, int(rnum))
    caption = "{}{} · cell {}".format(rank, rnum, cell)
    return cond, cond_disp, sortkey, caption


def qc_slides_for(folder):
    """Return list of ("grid", title, entries, footer, cols, cap_pt, cap_h)."""
    qc_dir = folder / "qc_raw"
    if not qc_dir.is_dir():
        return []
    by_cond = {}
    for p in sorted(qc_dir.glob("*.png")):
        parsed = parse_qc(p.name)
        if parsed is None:
            continue
        cond, cond_disp, sortkey, caption = parsed
        by_cond.setdefault((cond, cond_disp), []).append((sortkey, p, caption))
    slides = []
    for (cond, cond_disp), items in sorted(by_cond.items()):
        items.sort(key=lambda t: t[0])
        entries = [(p, cap) for _, p, cap in items]
        title = "SUN2 QC — raw signal at deepest invagination"
        footer = ("top / low = cells with highest / lowest SUN2 in deep NE "
                  "invaginations (÷ cell mean)  ·  SUN2 orange, nucleus cyan")
        slides.append(("grid", title, entries, footer, QC_MAX_COLS, 10, 0.30))
    return slides


# ---------------------------------------------------------------------------
# Content plan
# ---------------------------------------------------------------------------
def _tiles(pairs):
    """Filter to tiles that exist on disk; return (kept, missing_paths)."""
    kept, miss = [], []
    for path, cap in pairs:
        (kept if path_exists(path) else miss).append((path, cap))
    return kept, [p for p, _ in miss]


def build_plan():
    plan, missing = [], []

    # =======================================================================
    # Sec 1: SUN2 density profiles
    # =======================================================================
    plan.append(("divider",
                 "SUN2 density profiles",
                 "relative density (÷ cell mean) vs NE geometry · perinuc shells"))

    kept, miss = _tiles([
        (PROFILES / "SUN2_geomdens_hulldist.png", "vs hull-boundary distance"),
        (PROFILES / "SUN2_geomdens_mincurv.png", "vs min curvature"),
        (PROFILES / "SUN2_geomdens_meancurv.png", "vs mean curvature"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN2 density vs nuclear-envelope geometry",
                     kept,
                     "relative density (÷ cell mean) vs hull-boundary distance "
                     "/ curvature · perinuc shells",
                     3, 10, 0.28))

    # =======================================================================
    # Sec 2: Curated enrichment & correlation
    # =======================================================================
    plan.append(("divider",
                 "SUN2 enrichment & correlation",
                 "curated evidence · single condition (siControl)"))

    # Average level (02a + 02b + 02c): 3-up
    kept, miss = _tiles([
        (CURATED / "02a_SUN2_avg_hulldist.png", "avg vs hull dist"),
        (CURATED / "02b_SUN2_avg_minCurvature_concave.png",
         "avg vs min curv (concave)"),
        (CURATED / "02c_SUN2_avg_meanCurvature_concave.png",
         "avg vs mean curv (concave)"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN2 average level vs NE geometry",
                     kept,
                     "relative density (÷ cell mean) · perinuc 0.5 μm",
                     3, 10, 0.28))

    # Depth enrichment (03a + 03b): 2-up
    kept, miss = _tiles([
        (CURATED / "03a_SUN2_enrichment_hulldist_gt0.5um.png",
         "hull dist > 0.5 μm"),
        (CURATED / "03b_SUN2_enrichment_hulldist_gt1.0um.png",
         "hull dist > 1.0 μm"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN2 levels near invaginations",
                     kept,
                     "NE + perinuc shell · ref line 1",
                     2, 10, 0.30))

    # Curvature enrichment (03c + 03d + 03e): 3-up
    kept, miss = _tiles([
        (CURATED / "03c_SUN2_enrichment_minCurvature_lt0.png", "min curv < 0"),
        (CURATED / "03d_SUN2_enrichment_minCurvature_ltm0.25.png",
         "min curv < −0.25"),
        (CURATED / "03e_SUN2_enrichment_meanCurvature.png", "mean curv < 0"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN2 levels on concave NE surfaces",
                     kept,
                     "NE + perinuc shell · ref line 1",
                     3, 10, 0.28))

    # Correlation (04a + 04b + 04c): 3-up
    kept, miss = _tiles([
        (CURATED / "04a_SUN2_correlation_hulldist.png", "vs hull-boundary dist"),
        (CURATED / "04b_SUN2_correlation_minCurvature.png", "vs min curvature"),
        (CURATED / "04c_SUN2_correlation_meanCurvature.png", "vs mean curvature"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN2 per-cell correlation vs NE geometry",
                     kept,
                     "ref line 0",
                     3, 10, 0.28))

    # Shell-resolved correlation (enrichment/*_corr_*_shells): 3-up
    kept, miss = _tiles([
        (ENRICH / "SUN2_corr_hulldist_shells.png", "vs hull-boundary dist"),
        (ENRICH / "SUN2_corr_mincurv_shells.png", "vs min curvature"),
        (ENRICH / "SUN2_corr_meancurv_shells.png", "vs mean curvature"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN2 correlation vs NE geometry — by shell",
                     kept,
                     "correlation resolved across perinuc shells · ref line 0",
                     3, 10, 0.28))

    # Deep-invagination correlation (deepcorr): 2-up
    kept, miss = _tiles([
        (ENRICH / "SUN2_deepcorr_mincurv_shells.png",
         "vs min curvature (deep)"),
        (ENRICH / "SUN2_deepcorr_meancurv_shells.png",
         "vs mean curvature (deep)"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN2 correlation vs curvature in deep invaginations",
                     kept,
                     "deep-invagination shells · ref line 0",
                     2, 10, 0.28))

    # =======================================================================
    # Sec 3: Invagination depth profiles
    # =======================================================================
    plan.append(("divider",
                 "SUN2 in invaginations — depth profiles",
                 "centrosome-stratified (near vs away)"))

    for fname, cap in [
        ("SUN2_invag_depth_profiles.png",
         "depth profile"),
        ("SUN2_invag_depth_profiles_0_5um_cent_stratified.png",
         "depth profile, 0.5 μm perinuc"),
        ("SUN2_invag_depth_profiles_1um_cent_stratified.png",
         "depth profile, 1 μm perinuc"),
        ("SUN2_invag_enrichment_vs_centdist.png",
         "enrichment vs centrosome distance"),
    ]:
        p = DEPTH / fname
        if path_exists(p):
            plan.append(("single",
                         "SUN2 in invaginations — {}".format(cap),
                         p,
                         "near vs away from centrosome"))
        else:
            missing.append(p)

    # =======================================================================
    # Sec 4: Scalar violins (violin_plots)
    # =======================================================================
    plan.append(("divider",
                 "SUN2 scalars",
                 "one violin per metric · siControl"))

    def vp(stem, cap):
        return (VIOLINS / (stem + ".png"), cap)

    # Featured invagination pockets
    kept, miss = _tiles([
        vp("SUN2_cyto_in_nuc_hull_vs_near_convex_nuc_MFI_ratio",
           "grooves ÷ convex surface"),
        vp("SUN2_cyto_in_nuc_hull_vs_all_perinuc_MFI_ratio",
           "grooves ÷ whole perinuc shell"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN2 in the nuclear invagination pockets",
                     kept,
                     "voxel convex-hull decomposition · "
                     ">1 = enriched in grooves",
                     2, 13, 0.30))

    # SUN2 × DNA NE correlation
    kept, miss = _tiles([
        vp("SUN2_NE_Hoechst_corr", "SUN2 × DNA NE correlation"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN2 × DNA correlation at the nuclear envelope",
                     kept,
                     "NE boundary shell",
                     2, 13, 0.30))

    # Convex-hull supporting
    kept, miss = _tiles([
        vp("SUN2_cyto_in_nuc_hull_MFI", "grooves MFI (raw)"),
        vp("SUN2_cyto_in_nuc_hull_sig_fraction",
           "fraction of SUN2 signal in grooves"),
        vp("SUN2_frac_in_nuc_convex_hull",
           "fraction of SUN2 inside the hull"),
        vp("SUN2_invag_within_chull_vs_all_chull_MFI_ratio",
           "invag interior ÷ all within-hull"),
        vp("SUN2_invag_within_chull_vs_convex_within_chull_MFI_ratio",
           "invag interior ÷ convex rim"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN2 in the nuclear convex hull — supporting metrics",
                     kept,
                     "single condition (siControl)",
                     3, 11, 0.30))

    # Deepest invagination
    kept, miss = _tiles([
        vp("SUN2_deepest_invag_ratio_edge", "deepest invag ÷ edge"),
        vp("SUN2_deepest_invag_ratio_outer_shell",
           "deepest invag ÷ outer shell"),
        vp("SUN2_deepest_invag_Hoechst_corr_025um",
           "SUN2×DNA corr 0.25 μm"),
        vp("SUN2_deepest_invag_Hoechst_corr_05um",
           "SUN2×DNA corr 0.5 μm"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN2 at the deepest invagination",
                     kept,
                     "single condition (siControl)",
                     4, 11, 0.30))

    # Invag-vs-other ratios
    kept, miss = _tiles([
        vp("SUN2_invag_other_ratio_edge_025um",
           "invag÷other edge 0.25 μm"),
        vp("SUN2_invag_other_ratio_edge_05um",
           "invag÷other edge 0.5 μm"),
        vp("SUN2_invag_other_ratio_edge_1um",
           "invag÷other edge 1 μm"),
        vp("SUN2_invag_other_ratio_outer_shell_025um",
           "invag÷other outer shell 0.25 μm"),
        vp("SUN2_invag_other_ratio_outer_shell_05um",
           "invag÷other outer shell 0.5 μm"),
        vp("SUN2_invag_other_ratio_outer_shell_1um",
           "invag÷other outer shell 1 μm"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN2 invagination vs other NE ratios",
                     kept,
                     "ref line 1",
                     3, 10, 0.28))

    # Ratio above/below by centrosome side
    kept, miss = _tiles([
        vp("SUN2_ratio_above_below_0_5um", "above÷below 0.5 μm"),
        vp("SUN2_ratio_above_below_1um", "above÷below 1 μm"),
        vp("SUN2_ratio_above_by_side_0_5um", "above by side 0.5 μm"),
        vp("SUN2_ratio_above_by_side_1um", "above by side 1 μm"),
        vp("SUN2_ratio_below_by_side_0_5um", "below by side 0.5 μm"),
        vp("SUN2_ratio_below_by_side_1um", "below by side 1 μm"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN2 above/below ratios by centrosome side",
                     kept,
                     "ref line 1",
                     3, 10, 0.28))

    # By deepest invagination
    kept, miss = _tiles([
        vp("SUN2_ratio_by_deepest_invag_all_0_5um", "all faces · 0.5 μm"),
        vp("SUN2_ratio_by_deepest_invag_all_1um", "all faces · 1 μm"),
        vp("SUN2_ratio_by_deepest_invag_away_0_5um", "away faces · 0.5 μm"),
        vp("SUN2_ratio_by_deepest_invag_away_1um", "away faces · 1 μm"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN2 at the deepest invagination — ratio by side",
                     kept,
                     "SUN2 at deepest invag ÷ reference · ref line 1",
                     4, 11, 0.30))

    # Signal distribution / peak / boundary / outer shell metrics
    kept, miss = _tiles([
        vp("SUN2_MFI_at_nuc_boundary", "MFI at boundary"),
        vp("SUN2_MFI_outer_shell", "MFI outer shell"),
        vp("SUN2_peak_sig", "peak signal"),
        vp("SUN2_all_perinuc_MFI", "all perinuc MFI"),
        vp("SUN2_perinuc_sig_fraction", "perinuc signal fraction"),
        vp("SUN2_total_sig", "total signal"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN2 signal intensity & distribution",
                     kept,
                     "single condition (siControl)",
                     3, 10, 0.28))

    # Heterogeneity metrics (CV, SD, skewness, entropy)
    kept, miss = _tiles([
        vp("SUN2_CV_at_nuc_boundary", "CV at boundary"),
        vp("SUN2_CV_outer_shell", "CV outer shell"),
        vp("SUN2_SD_at_nuc_boundary", "SD at boundary"),
        vp("SUN2_SD_outer_shell", "SD outer shell"),
        vp("SUN2_norm_entropy_at_nuc_boundary", "norm entropy at boundary"),
        vp("SUN2_norm_entropy_outer_shell", "norm entropy outer shell"),
        vp("SUN2_skewness_at_nuc_boundary", "skewness at boundary"),
        vp("SUN2_skewness_outer_shell", "skewness outer shell"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN2 heterogeneity at the nuclear envelope",
                     kept,
                     "single condition (siControl)",
                     4, 10, 0.28))

    # =======================================================================
    # Sec 5: QC raw images
    # =======================================================================
    plan.append(("divider",
                 "QC — raw signal images",
                 "SUN2 signal at deepest invagination · top / low ranked cells"))

    for qc_it in qc_slides_for(CURATED):
        plan.append(qc_it)

    # --- Drop orphan dividers ---
    cleaned = []
    for i, it in enumerate(plan):
        if it[0] == "divider":
            if i + 1 >= len(plan) or plan[i + 1][0] == "divider":
                continue
        cleaned.append(it)
    return cleaned, missing


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    list_only = "--list" in sys.argv
    plan, missing = build_plan()

    n_slides = 1 + len(plan)
    print("Output: {}".format(OUTPUT_PATH))
    print("{} content slides (+ title) = {} total\n".format(len(plan), n_slides))
    for it in plan:
        if it[0] == "divider":
            print("\n=== {} ===".format(it[1]))
        elif it[0] == "single":
            print("  [1] {}".format(it[1]))
        else:
            print("  [{}] {}".format(len(it[2]), it[1]))
    if missing:
        print("\nMISSING ({}):".format(len(missing)))
        for m in missing:
            print("  - {}".format(Path(m).name if hasattr(m, "name") else m))
    if list_only:
        return

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)

    for it in plan:
        if it[0] == "divider":
            build_divider_slide(prs, it[1], it[2])
        elif it[0] == "single":
            build_single_slide(prs, it[1], it[2], it[3])
        elif it[0] == "grid":
            _, title, entries, footer, max_cols, cap_pt, cap_h = it
            build_grid_slide(prs, title, entries, footer, max_cols, cap_pt, cap_h)

    if OUTPUT_PATH.exists():
        backup_dir = OUTPUT_PATH.parent / "backups"
        created = backup_presentation(str(OUTPUT_PATH),
                                      backup_base=str(backup_dir))
        if created:
            print("\nBacked up previous deck under: {}".format(backup_dir))

    prs.save(str(OUTPUT_PATH))
    total = len(prs.slides._sldIdLst)
    print("\nDone. {} slides written to:\n  {}".format(total, OUTPUT_PATH))
    if missing:
        print("Skipped {} missing panel(s).".format(len(missing)))


if __name__ == "__main__":
    main()
