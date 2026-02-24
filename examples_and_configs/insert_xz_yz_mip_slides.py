"""
insert_xz_yz_mip_slides.py

Iterates through condition subdirectories, creates one blank slide per condition
in an existing PowerPoint, and inserts XZ MIP (left) and YZ MIP (right) montage
images where they exist. Slide title is set to the folder name.

Positions are rough defaults — adjust manually in PowerPoint afterwards.

Usage:
    python examples_and_configs/insert_xz_yz_mip_slides.py
"""

import glob
import os
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Configuration — edit these paths as needed
# ---------------------------------------------------------------------------

PPT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/"
    "Ctrl, DZNep, bTub, Lamin B, 20260203/"
    "Ctrl_DZNep_upto1hr_XZMIP_montages_nuc_deform_ratio_template.pptx"
)

# Output path: same folder as PPT_PATH but with "_template" removed from filename
OUTPUT_PATH = PPT_PATH.replace("_template.pptx", ".pptx")

PARENT_DIR = (
    "J:/FF/Nucleus_Project_up_to_1hr/"
    "020302026_jurkats_dznep_lmnB1orbTub_561LP55_488LP45_405LP40_"
)

# Subfolder within each condition directory that contains the montage images
MONTAGE_SUBPATH = "cropped/channels/montages_deformed_nuc"

# Glob patterns for the two image types
XZ_GLOB = "xz_mip_montage_nuc_gray_deformed_*.png"
YZ_GLOB = "yz_mip*.png"

# Slide dimensions (inches) — default PowerPoint widescreen is 13.33 x 7.5
SLIDE_WIDTH_IN = 13.33
SLIDE_HEIGHT_IN = 7.5

# Title textbox position/size (inches)
TITLE_LEFT = 0.3
TITLE_TOP = 0.1
TITLE_WIDTH = 12.7
TITLE_HEIGHT = 0.7
TITLE_FONT_SIZE = Pt(24)

# Image positions (inches) — height is auto-calculated to preserve aspect ratio
XZ_LEFT = 0.3
XZ_TOP = 0.9
XZ_WIDTH = 6.0

YZ_LEFT = 6.8
YZ_TOP = 0.9
YZ_WIDTH = 6.0

# ---------------------------------------------------------------------------


def format_condition_name(folder_name: str) -> str:
    """Convert a raw folder name to a readable slide title.

    Expected folder format: [replicate]_[condition]_[timepoint]_[marker]_
    Example: G2A3_h2o_8min_lmnb1_ → 'H2O 8 min: Lamin B1'

    Falls back to the raw folder name if the pattern isn't recognised.
    """
    condition_map = {
        "h2o": "H2O",
        "dznep": "DZNep",
    }
    timepoint_map = {
        "4min": "4 min",
        "8min": "8 min",
        "15min": "15 min",
        "30min": "30 min",
        "1hr": "1 hr",
    }
    marker_map = {
        "lmnb1": "Lamin B1",
        "btub": "\u03b2-Tub",  # β-Tub
    }

    # Strip trailing underscore(s) and split
    parts = folder_name.strip("_").split("_")
    # Expect exactly 4 tokens: replicate, condition, timepoint, marker
    if len(parts) != 4:
        return folder_name

    _, condition, timepoint, marker = parts
    condition_label = condition_map.get(condition.lower(), condition)
    timepoint_label = timepoint_map.get(timepoint.lower(), timepoint)
    marker_label = marker_map.get(marker.lower(), marker)

    return f"{condition_label} {timepoint_label}: {marker_label}"


def find_image(folder: Path, pattern: str) -> Path | None:
    """Return first match for glob pattern in folder, or None."""
    matches = sorted(folder.glob(pattern))
    return matches[0] if matches else None


def add_title(slide, text: str) -> None:
    """Add a title textbox to the slide."""
    txBox = slide.shapes.add_textbox(
        Inches(TITLE_LEFT), Inches(TITLE_TOP),
        Inches(TITLE_WIDTH), Inches(TITLE_HEIGHT)
    )
    tf = txBox.text_frame
    tf.text = text
    para = tf.paragraphs[0]
    para.font.size = TITLE_FONT_SIZE
    para.font.bold = True
    para.alignment = PP_ALIGN.CENTER


def insert_image_preserve_aspect(slide, image_path: Path, left_in: float, top_in: float, width_in: float) -> None:
    """Insert image at given position; height auto-calculated to preserve aspect ratio."""
    slide.shapes.add_picture(
        str(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(width_in)
        # height omitted → python-pptx preserves aspect ratio
    )


def main() -> None:
    parent = Path(PARENT_DIR)
    if not parent.exists():
        print(f"ERROR: Parent directory not found: {parent}")
        return

    # Collect subdirectories (condition folders), sorted
    condition_dirs = sorted([d for d in parent.iterdir() if d.is_dir()])
    if not condition_dirs:
        print(f"ERROR: No subdirectories found in {parent}")
        return

    print(f"Found {len(condition_dirs)} condition directories.")
    print(f"Opening PPT: {PPT_PATH}\n")

    prs = Presentation(PPT_PATH)
    blank_layout = prs.slide_layouts[6]  # Blank layout

    slides_added = 0
    missing_report = []

    for cond_dir in condition_dirs:
        cond_name = cond_dir.name
        montage_dir = cond_dir / MONTAGE_SUBPATH

        xz_img = find_image(montage_dir, XZ_GLOB) if montage_dir.exists() else None
        yz_img = find_image(montage_dir, YZ_GLOB) if montage_dir.exists() else None

        # Report status
        xz_status = xz_img.name if xz_img else "NOT FOUND"
        yz_status = yz_img.name if yz_img else "NOT FOUND"
        print(f"[{cond_name}]")
        print(f"  XZ: {xz_status}")
        print(f"  YZ: {yz_status}")

        if not montage_dir.exists():
            missing_report.append(f"{cond_name}: montage folder missing ({MONTAGE_SUBPATH})")
        if not xz_img:
            missing_report.append(f"{cond_name}: XZ MIP not found")
        if not yz_img:
            missing_report.append(f"{cond_name}: YZ MIP not found")

        # Add blank slide and title regardless of whether images exist
        slide = prs.slides.add_slide(blank_layout)
        add_title(slide, format_condition_name(cond_name))

        if xz_img:
            insert_image_preserve_aspect(slide, xz_img, XZ_LEFT, XZ_TOP, XZ_WIDTH)
        if yz_img:
            insert_image_preserve_aspect(slide, yz_img, YZ_LEFT, YZ_TOP, YZ_WIDTH)

        slides_added += 1

    prs.save(OUTPUT_PATH)
    print(f"\nDone. {slides_added} slides added. Saved to:\n  {OUTPUT_PATH}")

    if missing_report:
        print(f"\nMissing images / folders ({len(missing_report)}):")
        for msg in missing_report:
            print(f"  - {msg}")
    else:
        print("\nAll images found — no missing items.")


if __name__ == "__main__":
    main()
