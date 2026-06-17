"""
insert_actin_qc_synapse_mask_xz_mip_20260607_slides.py

Per-condition single-image actin QC deck SCOPED TO the 20260607 Kiet
CART dataset only (Y:/User_data/Kiet/20260607_pMLC_CART_actin_hoescht).
Same layout as insert_actin_qc_synapse_mask_xz_mip_slides.py but
trimmed to one dataset so it can be run as soon as that dataset's
MATLAB pipeline outputs land.

Result: 1 dataset x 6 conditions x 2 kinds = 12 slides.

Until pipeline results exist, every slide will render as a
`(missing)` placeholder.

Usage:
    python examples_and_configs/insert_actin_qc_synapse_mask_xz_mip_20260607_slides.py
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/CART_actin_only/"
    "CART_actin_QC_synapse_mask_xz_mip_20260607.pptx"
)

KIET_ROOT = "Y:/User_data/Kiet"

# Single dataset scope. YYYYMMDD acquisition date shown in slide titles.
DATASETS = [
    ("20260607", "20260607_pMLC_CART_actin_hoescht"),
]

CONDITIONS = [
    "CAT_5min",  "FMC_5min",
    "CAT_10min", "FMC_10min",
    "CAT_15min", "FMC_15min",
]

CONDITION_SUBPATH = "converted/cropped/split_channels"
PROG_SUBPATH = "prog_fixed_cells_actin_only/actin"

# (kind label, subpath under PROG_SUBPATH/)
# Stage AG path mapping: synapse/inner_outer/bot_combined/ -> 1slice_combined/
KINDS = [
    ("Actin at Synapse", "synapse/inner_outer/1slice_combined/montages"),
    ("Actin XZ MIP",     "xz_mip/montages"),
]

CHUNK_GLOB = "montage_cells_*.png"

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Slide layout (inches). 13.333 x 7.5 widescreen.
SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

IMG_LEFT = 0.10
IMG_TOP = 0.60
IMG_BOX_W = SLIDE_W - 2 * 0.10           # 13.13"
IMG_BOX_H = SLIDE_H - IMG_TOP - 0.10     # 6.80"

# Scalebar invariant from the MATLAB CART_fixed_cell_analysis pipeline (Stage AF).
# Every per-cell tile is rendered at PPUM_SOURCE px/μm in its data area; the
# 5 μm scalebar is therefore SCALEBAR_PX pixels in every montage tile.
# If the MATLAB constants change, update these and rerun all decks.
PPUM_SOURCE = 30          # px/μm in source PNGs
SCALEBAR_UM = 5           # μm
SCALEBAR_PX = PPUM_SOURCE * SCALEBAR_UM  # 150 px

# ---------------------------------------------------------------------------


def format_condition(folder_name: str) -> str:
    """CAT_5min -> 'CAT 5 min', FMC_10min -> 'FMC 10 min'."""
    time_map = {"5min": "5 min", "10min": "10 min", "15min": "15 min"}
    parts = folder_name.split("_", 1)
    if len(parts) != 2:
        return folder_name
    cell, time = parts
    return f"{cell} {time_map.get(time.lower(), time)}"


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    """Add a centered textbox with given text and styling."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _png_dims(path: Path) -> Tuple[int, int]:
    """Return (width_px, height_px) of a PNG without fully decoding it."""
    with Image.open(str(path)) as im:
        return im.size


def compute_deck_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    """Smallest ppi such that every image fits in (max_w_in x max_h_in).

    Used to pin px/inch across an entire deck so the 5 μm scalebar in every
    inserted montage lands at the same cm on every slide.
    """
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_image_at_ppi(slide, image_path: Path, ppi: float,
                     box_left: float, box_top: float,
                     box_w: float, box_h: float):
    """Insert image at width_in = w_px / ppi, height_in = h_px / ppi, centered
    in the (box_left, box_top, box_w, box_h) bounding box. Every slide that
    shares `ppi` therefore has the same slide px/inch, which keeps the
    embedded 5 μm scalebar at the same cm across the deck."""
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    left_in = box_left + (box_w - w_in) / 2
    top_in  = box_top  + (box_h - h_in) / 2
    return slide.shapes.add_picture(
        str(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(w_in),
    )


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _chunk_start_index(p: Path) -> int:
    """Extract the first integer after 'montage_cells_' for natural sort."""
    m = re.match(r"montage_cells_(\d+)", p.name)
    return int(m.group(1)) if m else 0


def find_first_chunk(montages_dir: Path) -> Optional[Path]:
    """Return the lowest-numbered montage_cells_*.png in the dir, or None."""
    if not montages_dir.is_dir():
        return None
    chunks = list(montages_dir.glob(CHUNK_GLOB))
    if not chunks:
        return None
    chunks.sort(key=_chunk_start_index)
    return chunks[0]


def build_slide(prs, title_text: str, image_path: Optional[Path], deck_ppi: float):
    """Build one full-image slide with title. Returns (slide, missing_flag)."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    if image_path is not None and image_path.exists():
        add_image_at_ppi(slide, image_path, deck_ppi,
                         IMG_LEFT, IMG_TOP, IMG_BOX_W, IMG_BOX_H)
        return slide, False

    add_textbox(
        slide, "(missing)",
        IMG_LEFT, IMG_TOP + IMG_BOX_H / 2 - 0.2, IMG_BOX_W, 0.4,
        font_pt=18, color=WHITE,
    )
    return slide, True


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    kiet_root = Path(KIET_ROOT)

    # Pre-pass: resolve every (kind, dataset, condition) montage and collect
    # the present images so we can pin deck-wide px/inch BEFORE slide build.
    slide_specs: List[Tuple[str, Optional[Path], str]] = []  # (title, image_path, log_key)
    for kind_label, kind_subpath in KINDS:
        for date_tag, dataset_name in DATASETS:
            for cond in CONDITIONS:
                cond_pretty = format_condition(cond)
                montages_dir = (
                    kiet_root / dataset_name / cond
                    / CONDITION_SUBPATH / PROG_SUBPATH / kind_subpath
                )
                chunk = find_first_chunk(montages_dir)
                title = f"{kind_label}: {cond_pretty} ({date_tag})"
                log_key = f"{kind_label}/{date_tag}/{cond}"
                slide_specs.append((title, chunk, log_key))

    present = [p for (_, p, _) in slide_specs if p is not None and p.exists()]
    if not present:
        print("WARNING: no real images found — using fallback PPI=100.")
        deck_ppi = 100.0
    else:
        deck_ppi = compute_deck_ppi(present, IMG_BOX_W, IMG_BOX_H)

    bar_in = SCALEBAR_PX / deck_ppi
    print(
        f"Deck-wide PPI = {deck_ppi:.2f} (pinned across all {len(present)}/"
        f"{len(slide_specs)} present slides).\n"
        f"  Scalebar invariant: {SCALEBAR_UM} μm = {SCALEBAR_PX} px in source "
        f"=> {bar_in:.3f} in = {bar_in * 2.54:.3f} cm on every slide.\n"
        f"  Source PPUM = {PPUM_SOURCE} px/μm (locked; update if MATLAB pipeline changes).\n"
    )
    print(f"Writing deck to: {OUTPUT_PATH}\n")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    missing = []
    slides_added = 0
    for (title, chunk, log_key) in slide_specs:
        _, is_missing = build_slide(prs, title, chunk, deck_ppi)
        slides_added += 1
        if is_missing:
            print(f"[{log_key}]  MISSING")
            missing.append(log_key)
        else:
            w_px, h_px = _png_dims(chunk)
            print(
                f"[{log_key}]  OK ({chunk.name}, {w_px}x{h_px} px "
                f"-> {w_px/deck_ppi:.2f}x{h_px/deck_ppi:.2f} in)"
            )

    prs.save(str(out_path))
    print(f"\nDone. {slides_added} slides written to:\n  {out_path}")

    if missing:
        print(f"\nMissing ({len(missing)}):")
        for m in missing:
            print(f"  - {m}")
    else:
        print("\nAll images found - no missing items.")


if __name__ == "__main__":
    main()
