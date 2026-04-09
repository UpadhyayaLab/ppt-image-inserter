"""
insert_VimMG_deepest_invag_montages.py

Creates one slide per condition (5 slides total) comparing both Vimentin MG experiments.
Each slide shows:
  - Experiment 1 (20240318) _with_cent montage (top)
  - Experiment 2 (20241212) _with_cent montage (bottom)
Source: merges/montages_deepest_invag in each experiment.

Conditions: PLL 3min, PLL 12min, αCD3 3min, αCD3 7min, αCD3 12min

Usage:
    python examples_and_configs/insert_VimMG_deepest_invag_montages.py
"""

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/Vimentin_MG_deepest_invag_montages.pptx"

EXPERIMENTS = [
    {
        "base_dir":        "J:/FF/vim_data/20240318_MG_E6-1_Vimentin",
        "montage_subpath": "channels/prog_fixed_cells/vimentin/deepest_invag_slice/merges/montages_deepest_invag",
        "tp_suffix":       "m",     # folder names like aCD3_3m
    },
    {
        "base_dir":        "J:/FF/vim_data/20241212_MG_E6-1_Vimentin",
        "montage_subpath": "individual_channels/prog_fixed_cells/vimentin/deepest_invag_slice/merges/montages_deepest_invag",
        "tp_suffix":       "min",   # folder names like aCD3_3min
    },
]

# Conditions in display order: (treatment_folder, timepoint_num, display_label)
CONDITIONS = [
    ("PLL",  "3",  "PLL, 3 min"),
    ("PLL",  "12", "PLL, 12 min"),
    ("aCD3", "3",  "\u03b1CD3, 3 min"),
    ("aCD3", "7",  "\u03b1CD3, 7 min"),
    ("aCD3", "12", "\u03b1CD3, 12 min"),
]

# Slide dimensions (widescreen)
SLIDE_WIDTH_IN  = 13.333
SLIDE_HEIGHT_IN = 7.5

# Title
TITLE_LEFT   = 1.167
TITLE_TOP    = 0.0
TITLE_WIDTH  = 11.0
TITLE_HEIGHT = 0.85
TITLE_FONT_SIZE = Pt(32)

# Images: exp1 top, exp2 bottom
IMG_LEFT      = 1.167
IMG_MAX_WIDTH = 11.0
TOP_IMG_TOP   = 0.9
BOT_IMG_TOP   = 4.3

# Label (bottom-left)
LABEL_LEFT      = 0.3
LABEL_FONT_SIZE = Pt(8)
LABEL_FONT_NAME = "Arial"

# ---------------------------------------------------------------------------


def find_first_with_cent(montage_dir: Path) -> Path | None:
    """Return the _with_cent montage with the lowest starting cell number."""
    matches = list(montage_dir.glob("montage_cells_*_with_cent.png"))
    if not matches:
        return None

    def start_cell(p: Path) -> int:
        try:
            return int(p.stem.split("_")[2])
        except (IndexError, ValueError):
            return 9999

    return sorted(matches, key=start_cell)[0]


def insert_image_width_constrained(slide, image_path: Path, left_in: float, top_in: float, width_in: float) -> None:
    slide.shapes.add_picture(
        str(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(width_in),
    )


def add_title(slide, text: str) -> None:
    txBox = slide.shapes.add_textbox(
        Inches(TITLE_LEFT), Inches(TITLE_TOP),
        Inches(TITLE_WIDTH), Inches(TITLE_HEIGHT),
    )
    tf = txBox.text_frame
    tf.text = text
    para = tf.paragraphs[0]
    para.font.size = TITLE_FONT_SIZE
    para.font.bold = True
    para.alignment = PP_ALIGN.CENTER


def add_path_label(slide, full_path: str) -> None:
    label_height = 0.2
    label_width = SLIDE_WIDTH_IN - LABEL_LEFT - 0.3
    top = SLIDE_HEIGHT_IN - 0.3 - label_height
    txBox = slide.shapes.add_textbox(
        Inches(LABEL_LEFT), Inches(top),
        Inches(label_width), Inches(label_height),
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    tf.text = full_path
    para = tf.paragraphs[0]
    para.font.size = LABEL_FONT_SIZE
    para.font.name = LABEL_FONT_NAME


def main() -> None:
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]

    slides_added = 0
    missing_report = []

    for treatment, tp_num, display_label in CONDITIONS:
        print(f"\n[{display_label}]")

        images = []
        for exp in EXPERIMENTS:
            base = Path(exp["base_dir"])
            folder_name = f"{treatment}_{tp_num}{exp['tp_suffix']}"
            montage_dir = base / folder_name / exp["montage_subpath"]

            img = find_first_with_cent(montage_dir) if montage_dir.exists() else None
            status = img.name if img else "NOT FOUND"
            print(f"  {base.name}/{folder_name}: {status}")

            if not montage_dir.exists():
                missing_report.append(f"{base.name}/{folder_name}: montage folder missing")
            elif not img:
                missing_report.append(f"{base.name}/{folder_name}: no _with_cent montage found")

            images.append(img)

        slide = prs.slides.add_slide(blank_layout)
        add_title(slide, f"Vimentin deepest invag: {display_label}")

        top_positions = [TOP_IMG_TOP, BOT_IMG_TOP]
        for img, top in zip(images, top_positions):
            if img:
                insert_image_width_constrained(slide, img, IMG_LEFT, top, IMG_MAX_WIDTH)

        label_path = str(next((img for img in images if img), montage_dir))
        add_path_label(slide, label_path)

        slides_added += 1

    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"\nDone. {slides_added} slides saved to:\n  {OUTPUT_PATH}")

    if missing_report:
        print(f"\n[WARNING] Missing images ({len(missing_report)}):")
        for msg in missing_report:
            print(f"  - {msg}")
        sys.exit(1)
    else:
        print("\nAll images found.")


if __name__ == "__main__":
    main()
