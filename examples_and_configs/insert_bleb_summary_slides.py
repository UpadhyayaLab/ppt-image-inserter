"""
insert_bleb_summary_slides.py

Condition-comparison summary deck for the fixed Jurkat blebbistatin experiment
(20240123_Fixed_E6-1_blebbistatin_Vimentin), compiled into
results_compilation/Jurkats_Blebb_Vim_20240123_20260626. Each grid panel is a
DMSO (blue) vs Blebbistatin (gold) violin/swarm comparison with significance
brackets, so one slide = one metric (title + the comparison plot + a source
footer). Companion to the montage deck (insert_bleb_dmso_vs_bleb_phys_scale_slides.py).

Modeled on the CART actin summary deck (insert_actin_summary_slides.py) for the
per-slide layout, plus the chromark deck's title slide / family dividers /
--list dry-run (insert_chromark_h3k27me3_summary_slides.py).

Metric selection + titles: kept identical to the noco-washout summary deck
(insert_noco_washout_summary_slides.py) — same metrics, same family order
(centrosome family leads), and the same per-metric titles — so the bleb and
noco/washout summary decks line up one-to-one (only the deck title/subtitle
differ, by experiment + date). The compile has 247 grid panels total; FAMILIES
is this curated subset (26 metrics) — add a (stem, title) row to grow the deck.

Self-contained: builds a blank deck (no template .pptx). Missing panels render
"(missing)" rather than failing. A previous deck is backed up before overwrite.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_bleb_summary_slides.py
    # dry run (print planned families/titles, build nothing):
    conda run -n PPT_editing python examples_and_configs/insert_bleb_summary_slides.py --list
"""

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation  # noqa: E402

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(
    "H:/FF/Nucleus_Data/3D_Nucleus/Fixed/Blebbistatin/"
    "20240123_Fixed_E6-1_blebbistatin_Vimentin/results_compilation/"
    "Jurkats_Blebb_Vim_20240123_20260626"
)
GRID_DIR = ROOT / "grid_panels"
CELL_COUNTS_PNG = ROOT / "cell_counts_barplot.png"   # context slide (optional)

OUTPUT_PATH = Path(
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/Blebbistatin/"
    "Bleb_Jurkats_DMSO_vs_Bleb_summary.pptx"
)

GRID_SUFFIX = "_grid.png"

DECK_TITLE = "Effects of blebbistatin on fixed Jurkat nuclei"
# n's from cell_counts.csv (DMSO 109, Blebbistatin 100). αCD3-activated, E6-1.
DECK_SUBTITLE = (
    "DMSO (n = 109) vs blebbistatin (n = 100)  ·  αCD3, E6-1  ·  "
    "fixed 01/23/2024  ·  compiled 2026-06-26"
)

# ---------------------------------------------------------------------------
# Curated metrics, grouped into families (divider slide per family).
# Each entry is (grid-panel stem, slide title). The stem + GRID_SUFFIX is the
# PNG under grid_panels/. Titles are for navigation; each plot carries its own
# authoritative y-axis label and DMSO/Blebbistatin x-labels.
# ---------------------------------------------------------------------------
FAMILIES = [
    ("Centrosome ↔ nucleus", [
        ("centrosome_center_z_rel_bottom_actin_plane", "Centrosome-synapse distance"),
        ("nuc_cent_closest_dist",                "Nucleus-centrosome closest distance"),
        ("cent_nuc_norm_dist_sphere_rad",        "Centrosome-to-nuclear-centroid distance (norm. to equiv sphere radius)"),
        ("centrosome_dist_deepest_real_avg_periphery_ratio", "Centrosome distance to deepest invag vs avg periphery ratio"),
    ]),
    ("Nuclear morphology", [
        ("nuc_aspect_ratio",  "Nuclear aspect ratio"),
        ("nuc_solidity",      "Nuclear solidity"),
        ("nuc_volume_mesh",   "Nuclear volume"),
        ("nuc_SA_mesh",       "Nuclear surface area"),
    ]),
    ("Nuclear deformation & invaginations", [
        ("chull_max_D",                       "Max invag depth over full nucleus"),
        ("chull_max_D_by_cent",               "Invagination depth near centrosome"),
        ("chull_mean_D_cent_global_ratio",    "Centrosomal Invagination Index (global)"),
        ("concavity_index_around_cent",       "Concavity index around centrosome"),
        ("deepest_invag_fraction_chull_volume", "Deepest invag: frac of convex hull volume"),
        ("deepest_region_periph_ratio_025um", "DNA levels near invag"),
        ("avg_normal_angle_adaptive_region_growth", "Deepest Invag Orientation"),
    ]),
    ("Actin", [
        ("actin_deform_ratio",        "Cell Aspect Ratio"),
        ("actin_bottom_mask_area",    "Synapse Area"),
        ("actin_MFI_around_cent_2um", "Actin MFI around centrosome (2 μm)"),
        ("actin_frac_around_cent_2um", "Actin fraction around centrosome (2 μm)"),
    ]),
    ("Vimentin", [
        ("vim_frac_in_nuc_convex_hull",                  "Vimentin fraction in nuclear convex hull"),
        ("vim_enrichment_within_half_um_nuc_2_um_cent",  "Vimentin enrichment (0.5 μm of nuc, 2 μm of cent)"),
        ("vim_cyto_in_nuc_hull_vs_near_convex_nuc_MFI_ratio", "Vimentin MFI: cytoplasmic (in nuc hull) vs near-convex regions of nuc"),
        ("VCRAI_0_2um",               "VCRAI (0–2 μm)"),
        ("vim_frac_around_cent_2um",  "Vimentin fraction around centrosome (2 μm)"),
        ("vim_MFI_around_cent_2um",   "Vimentin MFI around centrosome (2 μm)"),
        ("vim_ratio_above_below_0_5um", "Vimentin ratio above/below 0.5 μm"),
    ]),
]

# ---------------------------------------------------------------------------
# Colors / layout (matches the actin/chromark summary decks)
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.10

TITLE_LEFT = MARGIN
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.55
TITLE_FONT_PT = 28

IMG_LEFT = MARGIN
IMG_TOP = 0.66
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = 6.36

FOOTER_LEFT = MARGIN
FOOTER_TOP = 7.06
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.40
FOOTER_FONT_PT = 9

# ---------------------------------------------------------------------------


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    """Place an image inside (left, top, w, h), preserving aspect ratio and
    centering on whichever dimension ends up smaller than the box. The bleb
    grid panels are near-square, so they fit to height and center horizontally."""
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def rel_footer(path):
    try:
        return path.relative_to(ROOT).as_posix()
    except ValueError:
        return path.as_posix()


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.7, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=40, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.1, SLIDE_W - 2 * MARGIN, 1.0,
                font_pt=18, color=BLACK, italic=True)


def build_divider_slide(prs, family_name):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, family_name, MARGIN, 3.1, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=44, color=BLACK, bold=True)


def build_slide(prs, title_text, image_path, footer_text):
    """Title + full image (aspect preserved) + source-path footer.
    Returns missing flag (image not on disk)."""
    slide = _new_slide(prs)
    add_textbox(slide, title_text, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title_text), color=BLACK, bold=True)
    missing = not image_path.exists()
    if not missing:
        add_image_in_box(slide, str(image_path), IMG_LEFT, IMG_TOP, IMG_BOX_W, IMG_BOX_H)
    else:
        add_textbox(slide, "(missing)", IMG_LEFT, IMG_TOP + IMG_BOX_H / 2 - 0.2,
                    IMG_BOX_W, 0.4, font_pt=18, color=BLACK)
    add_textbox(slide, footer_text, FOOTER_LEFT, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=BLACK)
    return missing


def main():
    list_only = "--list" in sys.argv

    n_metrics = sum(len(items) for _, items in FAMILIES)
    # title + (cell-counts if present) + per-family (divider + metrics)
    est_slides = 1 + (1 if CELL_COUNTS_PNG.exists() else 0) + \
        sum(1 + len(items) for _, items in FAMILIES)

    print("Source: {}".format(GRID_DIR))
    print("{} curated metrics across {} families, est. {} slides\n".format(
        n_metrics, len(FAMILIES), est_slides))

    if list_only:
        for fam, items in FAMILIES:
            print("=== {} ({}) ===".format(fam, len(items)))
            for stem, title in items:
                exists = (GRID_DIR / (stem + GRID_SUFFIX)).exists()
                print("  [{}] {:<52s} {}".format(
                    "OK " if exists else "MISS", stem + GRID_SUFFIX, title))
            print("")
        return

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)

    if CELL_COUNTS_PNG.exists():
        build_slide(prs, "Cell counts (DMSO vs blebbistatin)",
                    CELL_COUNTS_PNG, rel_footer(CELL_COUNTS_PNG))
    else:
        print("Note: {} not found - skipping cell-counts slide.\n".format(
            CELL_COUNTS_PNG.name))

    missing = []
    for fam, items in FAMILIES:
        build_divider_slide(prs, fam)
        print("=== {} ===".format(fam))
        for stem, title in items:
            p = GRID_DIR / (stem + GRID_SUFFIX)
            is_missing = build_slide(prs, title, p, rel_footer(p))
            status = "OK" if not is_missing else "MISSING"
            print("  [{}] {} -> {!r}".format(status, stem + GRID_SUFFIX, title))
            if is_missing:
                missing.append(stem + GRID_SUFFIX)
        print("")

    if OUTPUT_PATH.exists():
        backup_dir = OUTPUT_PATH.parent / "backups"
        created = backup_presentation(str(OUTPUT_PATH), backup_base=str(backup_dir))
        if created:
            print("Backed up previous deck to: {}\n".format(backup_dir))

    prs.save(str(OUTPUT_PATH))
    total = len(prs.slides._sldIdLst)
    print("Done. {} metrics, {} slides written to:\n  {}".format(
        n_metrics, total, OUTPUT_PATH))
    if missing:
        print("\nSkipped {} missing panel(s) (not on disk):".format(len(missing)))
        for m in missing:
            print("  - {}".format(m))
    else:
        print("\nAll curated panels found - no missing items.")


if __name__ == "__main__":
    main()
