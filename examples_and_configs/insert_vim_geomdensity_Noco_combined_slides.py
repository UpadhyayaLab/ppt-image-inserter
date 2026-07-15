"""
insert_vim_geomdensity_Noco_combined_slides.py

Build the Vimentin geometry-density summary deck for the COMBINED Nocodazole
experiment (Apr 29 2022 + Jan 23 2024), DMSO vs Noco 1μM.

This is the two-experiment combined compile:
  results_compilation_Vim_Noco_geomdensity_combined_20260714

Focus: PER-EXPERIMENT comparisons.  Each enrichment / correlation violin shows
DMSO vs Noco within ONE experiment; the two experiments (Apr | Jan) are tiled
side-by-side for direct comparison.  Curated evidence (vim_loc_wrto_invag 03/04)
is pooled — it is the headline result and only exists pooled.

Sections:
  1. Curated key evidence (pooled) — 03a-e enrichment, 04a-c correlation
  2. Per-experiment enrichment & correlation violins
  3. Invagination depth profiles (per-expt, centrosome-stratified)
  4. Morphology scatter & centrosome NE metrics (per-expt)
  5. Scalar violins (by_day_panels — 4 violins: Apr DMSO, Apr Noco, Jan DMSO, Jan Noco)

Self-contained blank 16:9 deck.  Missing panels are skipped; a previous deck at
the output path is backed up first.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_vim_geomdensity_Noco_combined_slides.py
    conda run -n PPT_editing python examples_and_configs/insert_vim_geomdensity_Noco_combined_slides.py --list
"""

import math
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation, safe_path, path_exists  # noqa: E402

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(
    "M:/FF/FF_4TB_2_Backup_fullHD/Nucleus Project_2ndharddrive/"
    "Nucleus deformations analysis/Nucleus - Fixed Cell Data/"
    "20240123_E6-1_Nocodazole_Vimentin/"
    "results_compilation_Vim_Noco_geomdensity_combined_20260714"
)
ENRICH = ROOT / "geom_density" / "enrichment"
SINGLES = ROOT / "geom_density" / "profiles" / "singles"
NC = ROOT / "geom_density" / "near_cent"
MORPH = ROOT / "geom_density" / "morphology_scatter"
DEPTH = ROOT / "invag_depth_profiles"
BYDAY = ROOT / "by_day_panels"

OUTPUT_PATH = Path(
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/"
    "nuc_mesh_struct_outside_nuc/"
    "Vimentin geom-density vs NE geometry, DMSO vs Noco combined"
    " (20220429 + 20240123).pptx"
)

DECK_TITLE = "Vimentin density vs nuclear-envelope geometry — DMSO vs Noco"
DECK_SUBTITLE = (
    "Fixed Jurkats · Nocodazole (Apr 29 2022 + Jan 23 2024) · "
    "DMSO vs Noco 1μM · "
    "Vimentin = cytoplasmic (struct_out_nuc), perinuc 0.5 μm outside NE · "
    "per-experiment where available · compiled 2026-07-14"
)

# Experiment suffixes in per-experiment filenames
APR = "Apr_29,_2022"
JAN = "Jan_23,_2024"
APR_DISP = "Apr 29 2022"
JAN_DISP = "Jan 23 2024"
# invag_depth_profiles uses double underscore (no comma)
APR_DEPTH = "Apr_29__2022"
JAN_DEPTH = "Jan_23__2024"

# ---------------------------------------------------------------------------
# Colors / layout (16:9)
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0x55, 0x55, 0x55)
FIELD_COLOR = RGBColor(0x2E, 0x5A, 0x88)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.10
GAP = 0.16

TITLE_LEFT = MARGIN
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.55
TITLE_FONT_PT = 28

IMG_TOP = 0.66
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = 6.36

FOOTER_TOP = 7.06
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.40
FOOTER_FONT_PT = 9

CAPTION_HEIGHT = 0.55


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------
def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    image_path = safe_path(image_path)
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.7, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=38, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.1, SLIDE_W - 2 * MARGIN, 1.1,
                font_pt=16, color=GREY, italic=True)


def build_divider_slide(prs, title, subtitle=""):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, title, MARGIN, 3.0, SLIDE_W - 2 * MARGIN, 1.2,
                font_pt=40, color=BLACK, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, MARGIN, 4.25, SLIDE_W - 2 * MARGIN, 0.9,
                    font_pt=18, color=GREY, italic=True)


def build_single_slide(prs, title, image_path, caption, footer=None):
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    cap_h = CAPTION_HEIGHT if caption else 0.0
    img_h = IMG_BOX_H - cap_h
    add_image_in_box(slide, str(image_path), MARGIN, IMG_TOP, IMG_BOX_W, img_h)
    if caption:
        add_textbox(slide, caption, MARGIN, IMG_TOP + img_h, IMG_BOX_W, cap_h,
                    font_pt=12, color=GREY)
    add_textbox(slide, footer or rel_footer(image_path), MARGIN, FOOTER_TOP,
                FOOTER_WIDTH, FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


def add_labeled_caption(slide, lines, left, top, width, height):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    for i, (text, pt, color) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.size = Pt(pt)
        run.font.color.rgb = color
    return box


def _tile_caption(slide, cap, path, left, top, width, height, cap_pt):
    lines = []
    if cap:
        lines.append((cap, min(cap_pt, 10), GREY))
    lines.append((Path(str(path)).stem, 8, FIELD_COLOR))
    add_labeled_caption(slide, lines, left, top, width, height)


def build_grid_slide(prs, title, entries, footer_text, max_cols,
                     cap_pt=11, cap_h=0.28):
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    n = len(entries)
    cols = min(max_cols, n)
    rows = math.ceil(n / cols)
    area_top = IMG_TOP
    area_h = FOOTER_TOP - IMG_TOP - 0.02
    cell_w = (SLIDE_W - 2 * MARGIN - (cols - 1) * GAP) / cols
    cell_h = (area_h - (rows - 1) * GAP) / rows
    img_h = cell_h - cap_h
    for i, (path, cap) in enumerate(entries):
        r, c = divmod(i, cols)
        left = MARGIN + c * (cell_w + GAP)
        top = area_top + r * (cell_h + GAP)
        add_image_in_box(slide, str(path), left, top, cell_w, img_h)
        _tile_caption(slide, cap, path, left, top + img_h, cell_w, cap_h, cap_pt)
    add_textbox(slide, footer_text, MARGIN, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


def rel_footer(path):
    try:
        return Path(path).relative_to(ROOT).as_posix()
    except ValueError:
        return Path(path).name


# ---------------------------------------------------------------------------
# Content plan
# ---------------------------------------------------------------------------
def _tiles(pairs):
    """Filter to tiles that exist on disk; return (kept, missing_paths)."""
    kept, miss = [], []
    for path, cap in pairs:
        (kept if path_exists(path) else miss).append((path, cap))
    return kept, [p for p, _ in miss]


def _enr_pair(stem, cap_hint):
    """Return [(apr_path, cap), (jan_path, cap)] for a per-experiment enrichment."""
    return [
        (ENRICH / "{}_shells_{}.png".format(stem, APR),
         "{} — {}".format(cap_hint, APR_DISP)),
        (ENRICH / "{}_shells_{}.png".format(stem, JAN),
         "{} — {}".format(cap_hint, JAN_DISP)),
    ]


def _depth_pair(stem, cap_hint):
    """Return [(apr_path, cap), (jan_path, cap)] for invag_depth_profiles."""
    return [
        (DEPTH / "{}_{}.png".format(stem, APR_DEPTH),
         "{} — {}".format(cap_hint, APR_DISP)),
        (DEPTH / "{}_{}.png".format(stem, JAN_DEPTH),
         "{} — {}".format(cap_hint, JAN_DISP)),
    ]


def build_plan():
    plan, missing = [], []

    # =======================================================================
    # Sec 0: Density profiles — per experiment (DMSO and Noco separately)
    # =======================================================================
    plan.append(("divider",
                 "Vimentin density profiles — per experiment",
                 "Mean±SEM · perinuc 0.5 μm shell · "
                 "DMSO vs Noco per experiment"))

    def _profile_quad(geom, suffix=""):
        """4 per-condition per-experiment profile singletons for one geometry."""
        base = "vim_geomdens_{}_perinuc05_CD3".format(geom)
        return [
            (SINGLES / "{}_{}_line{}.png".format(base, "DMSO_" + APR_DEPTH, suffix),
             "DMSO — {}".format(APR_DISP)),
            (SINGLES / "{}_{}_line{}.png".format(base, "Noco_" + APR_DEPTH, suffix),
             "Noco — {}".format(APR_DISP)),
            (SINGLES / "{}_{}_line{}.png".format(base, "DMSO_" + JAN_DEPTH, suffix),
             "DMSO — {}".format(JAN_DISP)),
            (SINGLES / "{}_{}_line{}.png".format(base, "Noco_" + JAN_DEPTH, suffix),
             "Noco — {}".format(JAN_DISP)),
        ]

    # Hull-boundary distance (= invagination depth)
    kept, miss = _tiles(_profile_quad("hulldist"))
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vimentin density vs invagination depth — per experiment",
                     kept,
                     "Mean±SEM relative density (÷ cell mean) vs "
                     "hull-boundary distance · perinuc 0.5 μm",
                     2, 10, 0.28))

    # Min curvature (concave half)
    kept, miss = _tiles(_profile_quad("mincurv", "_concave"))
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vimentin density vs min curvature — per experiment",
                     kept,
                     "Mean±SEM relative density · concave half · "
                     "perinuc 0.5 μm",
                     2, 10, 0.28))

    # Mean curvature (concave half)
    kept, miss = _tiles(_profile_quad("meancurv", "_concave"))
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vimentin density vs mean curvature — per experiment",
                     kept,
                     "Mean±SEM relative density · concave half · "
                     "perinuc 0.5 μm",
                     2, 10, 0.28))

    # =======================================================================
    # Sec 1: Per-experiment enrichment & correlation
    # =======================================================================
    plan.append(("divider",
                 "Vimentin enrichment & correlation — per experiment",
                 "each violin = DMSO vs Noco within one experiment · "
                 "paired Apr 2022 vs Jan 2024"))

    # Depth enrichment per expt (4 tiles: gt0.5 Apr|Jan, gt1.0 Apr|Jan)
    tiles = (_enr_pair("vim_hulldist_gt0.5um", "hull dist > 0.5 μm")
             + _enr_pair("vim_hulldist_gt1.0um", "hull dist > 1.0 μm"))
    kept, miss = _tiles(tiles)
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vimentin levels near invaginations — per experiment",
                     kept,
                     "DMSO vs Noco per experiment · ref line 1",
                     4, 10, 0.30))

    # Curvature enrichment per expt (6 tiles, max_cols=2: Apr|Jan per threshold)
    tiles = (_enr_pair("vim_mincurv_lt0", "min curv < 0")
             + _enr_pair("vim_mincurv_ltm0.25", "min curv < −0.25")
             + _enr_pair("vim_meancurv_lt0", "mean curv < 0"))
    kept, miss = _tiles(tiles)
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vimentin levels on concave NE surfaces — per experiment",
                     kept,
                     "DMSO vs Noco per experiment · ref line 1",
                     2, 10, 0.28))

    # Correlation per expt (6 tiles, max_cols=2)
    tiles = (_enr_pair("vim_corr_hulldist", "vs hull-boundary dist")
             + _enr_pair("vim_corr_mincurv", "vs min curvature")
             + _enr_pair("vim_corr_meancurv", "vs mean curvature"))
    kept, miss = _tiles(tiles)
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vimentin per-cell correlation vs NE geometry — per experiment",
                     kept,
                     "DMSO vs Noco per experiment · ref line 0",
                     2, 10, 0.28))

    # Deepcorr per expt (4 tiles, max_cols=2)
    tiles = (_enr_pair("vim_deepcorr_mincurv", "vs min curvature (deep)")
             + _enr_pair("vim_deepcorr_meancurv", "vs mean curvature (deep)"))
    kept, miss = _tiles(tiles)
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vimentin correlation vs curvature in deep invaginations "
                     "— per experiment",
                     kept,
                     "DMSO vs Noco per experiment · ref line 0",
                     2, 10, 0.28))

    # =======================================================================
    # Sec 3: Invagination depth profiles (per-expt, centrosome-stratified)
    # =======================================================================
    plan.append(("divider",
                 "Vimentin in invaginations — depth profiles",
                 "centrosome-stratified (near vs away) · per experiment"))

    # 0.5 μm cent-stratified: Apr | Jan side by side
    kept, miss = _tiles(_depth_pair(
        "vim_invag_depth_profiles_0_5um_cent_stratified",
        "0.5 μm perinuc"))
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vimentin in invaginations — depth profile, "
                     "0.5 μm perinuc",
                     kept,
                     "near vs away from centrosome · per experiment",
                     2, 10, 0.28))

    # 1 μm cent-stratified: Apr | Jan side by side
    kept, miss = _tiles(_depth_pair(
        "vim_invag_depth_profiles_1um_cent_stratified",
        "1 μm perinuc"))
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vimentin in invaginations — depth profile, "
                     "1 μm perinuc",
                     kept,
                     "near vs away from centrosome · per experiment",
                     2, 10, 0.28))

    # Enrichment vs centrosome distance: Apr | Jan side by side
    kept, miss = _tiles(_depth_pair(
        "vim_invag_enrichment_vs_centdist",
        "enrichment vs centrosome dist"))
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vimentin invag enrichment vs centrosome distance "
                     "— per experiment",
                     kept,
                     "per experiment",
                     2, 10, 0.28))

    # =======================================================================
    # Sec 4: Morphology scatter & centrosome NE — per experiment
    # =======================================================================
    plan.append(("divider",
                 "Morphology scatter & centrosome metrics — per experiment",
                 "enrichment / correlation vs deepest invagination depth, "
                 "by centrosome side"))

    # Morphology scatter (6-up: top row = Apr, bottom = Jan)
    scat_stems = [
        ("vim_hull_enrichment_2um_cent_vs_chull_max_D_by_cent",
         "hull enrich 2μm cent vs max D"),
        ("vim_r_hulldist_perinuc05_bycent_vs_chull_max_D_by_cent",
         "r(hulldist) by cent vs max D"),
        ("vim_ratio_invag_away_1um_vs_chull_max_D_by_cent",
         "ratio invag/away 1μm vs max D"),
    ]
    scat_tiles = []
    for stem, cap in scat_stems:
        scat_tiles.append((MORPH / "{}_{}.png".format(stem, APR),
                           "{} — {}".format(cap, APR_DISP)))
    for stem, cap in scat_stems:
        scat_tiles.append((MORPH / "{}_{}.png".format(stem, JAN),
                           "{} — {}".format(cap, JAN_DISP)))
    kept, miss = _tiles(scat_tiles)
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vimentin vs invagination depth — morphology scatter",
                     kept,
                     "per experiment · color = centrosome side",
                     3, 9, 0.30))

    # Near centrosome level per expt (2-up)
    kept, miss = _tiles([
        (NC / "vim_near_cent_level_{}.png".format(APR),
         "Vim near cent — {}".format(APR_DISP)),
        (NC / "vim_near_cent_level_{}.png".format(JAN),
         "Vim near cent — {}".format(JAN_DISP)),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vimentin near the centrosome (NE-facing) — per experiment",
                     kept,
                     "Vim ÷ cell mean near centrosome NE · "
                     "DMSO vs Noco per experiment · ref line 1",
                     2, 12, 0.30))

    # =======================================================================
    # Sec 5: Scalar violins (by_day_panels — per-expt split)
    # =======================================================================
    plan.append(("divider",
                 "Vimentin scalars — by experiment (by_day panels)",
                 "4 violins per panel: Apr DMSO, Apr Noco, Jan DMSO, Jan Noco "
                 "· ref line 1"))

    def bd(stem, cap):
        return (BYDAY / (stem + "_by_day.png"), cap)

    # Featured invagination pockets
    kept, miss = _tiles([
        bd("vim_cyto_in_nuc_hull_vs_near_convex_nuc_MFI_ratio",
           "grooves ÷ convex surface"),
        bd("vim_cyto_in_nuc_hull_vs_all_perinuc_MFI_ratio",
           "grooves ÷ whole perinuc shell"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vimentin in the nuclear invagination pockets — by experiment",
                     kept,
                     "voxel convex-hull decomposition · "
                     ">1 = enriched in grooves · per-expt DMSO vs Noco",
                     2, 13, 0.30))

    # Convex-hull supporting (exclude vs_nuc_MFI_ratio — meaningless for cyto)
    kept, miss = _tiles([
        bd("vim_cyto_in_nuc_hull_MFI",
           "grooves MFI (raw)"),
        bd("vim_cyto_in_nuc_hull_sig_fraction",
           "fraction of Vim signal in grooves"),
        bd("vim_frac_in_nuc_convex_hull",
           "fraction of Vim inside the hull"),
        bd("vim_invag_within_chull_vs_all_chull_MFI_ratio",
           "invag interior ÷ all within-hull"),
        bd("vim_invag_within_chull_vs_convex_within_chull_MFI_ratio",
           "invag interior ÷ convex rim"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vimentin in the nuclear convex hull — supporting metrics",
                     kept,
                     "per-expt DMSO vs Noco · ref line 1",
                     3, 11, 0.30))

    # Deepest invagination
    kept, miss = _tiles([
        bd("vim_ratio_by_deepest_invag_all_0_5um",
           "all faces · 0.5 μm"),
        bd("vim_ratio_by_deepest_invag_all_1um",
           "all faces · 1 μm"),
        bd("vim_ratio_by_deepest_invag_away_0_5um",
           "away faces · 0.5 μm"),
        bd("vim_ratio_by_deepest_invag_away_1um",
           "away faces · 1 μm"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vimentin at the deepest invagination — by experiment",
                     kept,
                     "Vim at the single deepest invag ÷ reference · "
                     "per-expt DMSO vs Noco · ref line 1",
                     4, 11, 0.30))

    # Centrosome NE-facing by_day
    kept, miss = _tiles([
        bd("vim_ratio_by_cent_away_0_5um",
           "cent-side ÷ away-side · 0.5 μm"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vimentin NE cent-side vs away-side — by experiment",
                     kept,
                     "per-expt DMSO vs Noco · ref line 1",
                     2, 12, 0.30))

    # Centoplasmic pool around centrosome
    kept, miss = _tiles([
        bd("vim_frac_around_cent_2um",
           "fraction of Vim within 2 μm of centrosome"),
        bd("vim_MFI_around_cent_2um",
           "Vim MFI within 2 μm of centrosome"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vimentin clustered near the centrosome (cytoplasmic pool)",
                     kept,
                     "3D ball around centrosome, not NE-restricted · "
                     "per-expt DMSO vs Noco",
                     2, 12, 0.30))

    # --- Drop orphan dividers ---
    cleaned = []
    for i, it in enumerate(plan):
        if it[0] == "divider":
            if i + 1 >= len(plan) or plan[i + 1][0] == "divider":
                continue
        cleaned.append(it)
    return cleaned, missing


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    list_only = "--list" in sys.argv
    plan, missing = build_plan()

    n_slides = 1 + len(plan)
    print("Output: {}".format(OUTPUT_PATH))
    print("{} content slides (+ title) = {} total\n".format(len(plan), n_slides))
    for it in plan:
        if it[0] == "divider":
            print("\n=== {} ===".format(it[1]))
        elif it[0] == "single":
            print("  [1] {}".format(it[1]))
        else:
            print("  [{}] {}".format(len(it[2]), it[1]))
    if missing:
        print("\nMISSING ({}):".format(len(missing)))
        for m in missing:
            print("  - {}".format(Path(m).name if hasattr(m, "name") else m))
    if list_only:
        return

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)

    for it in plan:
        if it[0] == "divider":
            build_divider_slide(prs, it[1], it[2])
        elif it[0] == "single":
            build_single_slide(prs, it[1], it[2], it[3])
        elif it[0] == "grid":
            _, title, entries, footer, max_cols, cap_pt, cap_h = it
            build_grid_slide(prs, title, entries, footer, max_cols, cap_pt, cap_h)

    if OUTPUT_PATH.exists():
        backup_dir = OUTPUT_PATH.parent / "backups"
        created = backup_presentation(str(OUTPUT_PATH),
                                      backup_base=str(backup_dir))
        if created:
            print("\nBacked up previous deck under: {}".format(backup_dir))

    prs.save(str(OUTPUT_PATH))
    total = len(prs.slides._sldIdLst)
    print("\nDone. {} slides written to:\n  {}".format(total, OUTPUT_PATH))
    if missing:
        print("Skipped {} missing panel(s).".format(len(missing)))


if __name__ == "__main__":
    main()
