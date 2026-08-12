"""
insert_actin_summary_slides.py

Build one CART actin summary deck per dataset listed in DATASETS.
Each slide shows one PNG plot with:
  - bold title derived from the filename (with overrides)
  - the plot, full-width, aspect ratio preserved
  - the full source filepath in a small footer (for traceability)

To add a new dataset, append an entry to DATASETS with its grid_panels/
root, output path, and PNG list. To add slides, extend the relevant
*_METRICS list. Each run regenerates every deck (idempotent).

Modeled after insert_actin_qc_synapse_mask_xz_mip_slides.py.

Usage:
    python examples_and_configs/insert_actin_summary_slides.py
"""

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation  # noqa: E402

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

# Subfolder names — identical across every dataset's grid_panels/ tree.
SUB_SCATTER    = "timecourse_scatter_plots"
SUB_CAT_VS_FMC = "CAT_vs_FMC_by_timepoint"

# Slide entries: (subfolder, png_filename). Joined with each dataset's
# compiled_root at deck-build time. Missing files render "(missing)".

SLICE_METRICS = [
    # area + spatial concentration
    (SUB_SCATTER, "actin_bottom_mask_area_grid.png"),
    (SUB_SCATTER, "actin_bottom_slice_inner_outer_ratio_50pct_grid.png"),
    (SUB_SCATTER, "actin_bottom_slice_inner_outer_ratio_70pct_grid.png"),
    # intensity
    (SUB_SCATTER, "actin_bottom_MFI_grid.png"),
    (SUB_SCATTER, "actin_bottom_total_sig_grid.png"),
    # shape
    (SUB_SCATTER, "actin_bottom_slice_perimeter_grid.png"),
    (SUB_SCATTER, "actin_bottom_slice_circularity_grid.png"),
    (SUB_SCATTER, "actin_bottom_slice_solidity_grid.png"),
    (SUB_SCATTER, "actin_bottom_slice_eccentricity_grid.png"),
]
RAD_PROFILE_SLICE = (
    SUB_CAT_VS_FMC,
    "actin_bottom_slice_rad_profile_auc1_all_cells_with_average_grid.png",
)
# Slice block with rad_profile spliced in after the two IOR slides
# (slot 4 in 1-based) per the agreed order.
SLICE_METRICS_WITH_RAD = (
    SLICE_METRICS[:3] + [RAD_PROFILE_SLICE] + SLICE_METRICS[3:]
)

THREE_SLICE_METRICS = [
    # area + spatial concentration
    (SUB_SCATTER, "actin_bottom_3slice_mask_area_grid.png"),
    (SUB_SCATTER, "actin_bottom_3slice_mip_inner_outer_ratio_50pct_grid.png"),
    (SUB_SCATTER, "actin_bottom_3slice_mip_inner_outer_ratio_70pct_grid.png"),
    # intensity
    (SUB_SCATTER, "actin_bottom_3slice_MFI_grid.png"),
    (SUB_SCATTER, "actin_bottom_3slice_total_sig_grid.png"),
    # shape
    (SUB_SCATTER, "actin_bottom_3slice_perimeter_grid.png"),
    (SUB_SCATTER, "actin_bottom_3slice_circularity_grid.png"),
    (SUB_SCATTER, "actin_bottom_3slice_solidity_grid.png"),
    (SUB_SCATTER, "actin_bottom_3slice_eccentricity_grid.png"),
]
RAD_PROFILE_3SLICE = (
    SUB_CAT_VS_FMC,
    "actin_bottom_3slice_mip_rad_profile_auc1_all_cells_with_average_grid.png",
)
THREE_SLICE_METRICS_WITH_RAD = (
    THREE_SLICE_METRICS[:3] + [RAD_PROFILE_3SLICE] + THREE_SLICE_METRICS[3:]
)

# IRM-derived synapse area + actin median radius, present in compiles
# that ran the newer footprint code. Splice into a per-dataset png_files:
# IRM after mask_area, R_median after the two inner/outer ratios.
IRM_SYNAPSE_AREA = (SUB_SCATTER, "IRM_mask_area_grid.png")
R_MEDIAN_SLICE = (SUB_SCATTER, "actin_bottom_slice_r_median_grid.png")
R_MEDIAN_3SLICE = (SUB_SCATTER, "actin_bottom_3slice_mip_r_median_grid.png")

# Optional extra blocks appended to actin-only decks whose compiles
# ran the fuller feature set (foci scoring, whole-cell MIP metrics,
# Z-distribution). Older compiles may lack some of these PNGs — check
# the source before splicing into an existing DATASETS entry.

# Actin foci (8): 4 canonical + 4 scoring/threshold extras from the
# newer foci-detection code.
ACTIN_FOCI_METRICS = [
    (SUB_SCATTER, "actin_foci_count_grid.png"),
    (SUB_SCATTER, "actin_foci_area_um2_grid.png"),
    (SUB_SCATTER, "actin_foci_area_fraction_grid.png"),
    (SUB_SCATTER, "actin_foci_mean_intensity_grid.png"),
    (SUB_SCATTER, "actin_foci_total_intensity_grid.png"),
    (SUB_SCATTER, "actin_foci_max_norm_score_grid.png"),
    (SUB_SCATTER, "actin_foci_mean_norm_score_grid.png"),
    (SUB_SCATTER, "actin_foci_threshold_used_grid.png"),
]

# Actin whole-cell metrics (6): MIP footprint (4) + integrated intensity (2).
ACTIN_WHOLE_CELL_METRICS = [
    (SUB_SCATTER, "actin_MIP_area_grid.png"),
    (SUB_SCATTER, "actin_MIP_major_axis_length_grid.png"),
    (SUB_SCATTER, "actin_MIP_minor_axis_length_grid.png"),
    (SUB_SCATTER, "actin_MIP_mask_nonzero_grid.png"),
    (SUB_SCATTER, "actin_total_sig_grid.png"),
    (SUB_SCATTER, "actin_peak_sig_grid.png"),
]

# Actin Z-distribution (10): position of signal in cell height + cell extent.
ACTIN_Z_METRICS = [
    (SUB_SCATTER, "actin_z50_rel_cell_bottom_grid.png"),
    (SUB_SCATTER, "actin_z75_rel_cell_bottom_grid.png"),
    (SUB_SCATTER, "actin_z90_rel_cell_bottom_grid.png"),
    (SUB_SCATTER, "actin_zprofile_peak_idx_grid.png"),
    (SUB_SCATTER, "actin_zprofile_peak_rel_bot_grid.png"),
    (SUB_SCATTER, "actin_bottom_slice_num_grid.png"),
    (SUB_SCATTER, "actin_top_slice_num_grid.png"),
    (SUB_SCATTER, "actin_height_grid.png"),
    (SUB_SCATTER, "actin_above_FOV_grid.png"),
    (SUB_SCATTER, "actin_below_FOV_grid.png"),
]

# by_timepoint compiles are flat: every PNG lives under
# <by_timepoint_root>/all_<tp>_<date>/all_conditions/, with no _grid suffix.
BY_TIMEPOINT_ROOT = Path(
    "J:/FF/fixed_cell/CAR_TCell/results_compiled/actin/by_timepoint"
)


def _by_timepoint(tp_short: str, date: str = "20260618") -> Path:
    """Path to one timepoint's flat all_conditions/ dir."""
    return BY_TIMEPOINT_ROOT / f"all_{tp_short}_{date}" / "all_conditions"


# Flat png_files list usable with strip_grid_suffix=True datasets (the
# subfolder slot is "" because the PNGs sit directly under compiled_root).
ALL_CONDITIONS_PNG_FILES = [
    ("", fn) for _, fn in (SLICE_METRICS_WITH_RAD + THREE_SLICE_METRICS_WITH_RAD)
]

# Same shape but slice-only (for compiles with no 3-slice MIP variants).
SLICE_ONLY_PNG_FILES = [("", fn) for _, fn in SLICE_METRICS_WITH_RAD]

# L: drive `all_datasets_actin` by-timepoint compile (different naming
# than the prior J: drive one: separate `compiled_<date>` segment).
ALL_DATASETS_BY_TIMEPOINT_ROOT = Path(
    "L:/FF/CAR T/actin_compiled_results/all_datasets_actin/by_timepoint"
)


PLAT_TIRF_ROOT = Path(
    "L:/FF/CAR T/CAR T data for actin area trouble shoot/"
    "CAR T Actin Fixed using PLAT data 5 15min/"
    "results_compiled_tirf/compiled_plat_tirf_20260701/grid_panels"
)
OUTPUT_DIR_PLAT = Path("K:/FF/PPT/PPT_autogeneration/CART/pLAT")

# The 20260701 pLAT compile uses different subfolder names than the
# older timecourse_20260630 compile — timecourse_scatter_plots was
# renamed to scatter_timecourse_5_vs_15, CAT_vs_FMC_by_timepoint to
# CAT_vs_FMC, and the binned distribution is a single top-level PNG
# (no per-date breakdown).
SUB_PLAT_SCATTER = "scatter_timecourse_5_vs_15"
SUB_PLAT_CAT     = "CAT_vs_FMC"

# 3 correlation curves — parity with pZap70. g_ave + autocorr live in
# the pLAT-specific subfolder for this compile.
PLAT_CORR_METRICS = [
    (SUB_PLAT_SCATTER,
     "pLAT_synapse_g_ave_grid.png"),
    (SUB_PLAT_CAT,
     "pLAT_synapse_autocorr_c_all_cells_with_average_grid.png"),
    (SUB_PLAT_CAT,
     "foci_pLAT_cross_corr_profile_all_cells_with_average_grid.png"),
]

PZAP70_TIRF_ROOT = Path(
    "L:/FF/CAR T/CAR T data for actin area trouble shoot/"
    "results_compiled_tirf/compiled_pzap70_tirf_timecourse_20260701/grid_panels"
)
OUTPUT_DIR_PZAP70 = Path("K:/FF/PPT/PPT_autogeneration/CART/pZap70")

PMLC_JOINT_ROOT = Path(
    "Y:/User_data/Kiet/results_compiled/pmlc_joint/"
    "compiled_pmlc_v2_20260702/grid_panels"
)
OUTPUT_DIR_PMLC = Path("K:/FF/PPT/PPT_autogeneration/CART/pMLC")

PKC_ROOT = Path(
    "Y:/User_data/Kiet/01312026_fixed_CART_PKC_Phalloidin_/"
    "results_compiled/compiled_pkc_20260702/grid_panels"
)
OUTPUT_DIR_PKC = Path("K:/FF/PPT/PPT_autogeneration/CART/PKC")

PPKC_ROOT = Path(
    "Y:/User_data/Kiet/20260414_p_PKC_theta_Phalloidin561_/"
    "results_compiled/compiled_ppkc_20260702/grid_panels"
)
OUTPUT_DIR_PPKC = Path("K:/FF/PPT/PPT_autogeneration/CART/pPKC")

CTSB_ALL_ROOT = Path(
    "J:/FF/fixed_cell/CAR_TCell/results_compiled/ctsb_all/"
    "compiled_2023_2024_20260702/grid_panels"
)
OUTPUT_DIR_CTSB = Path("K:/FF/PPT/PPT_autogeneration/CART/MT_CatB")

# 13 CatB metrics — 12 in scatter + 1 rad_profile in CAT_vs_FMC.
# CatB rad_profile in this compile is `_all_cells_with_average` only
# (no `_auc1_` variant), so use that.
CTSB_METRICS = [
    (SUB_SCATTER,    "CathepsinB_synapse_MFI_grid.png"),
    (SUB_SCATTER,    "CathepsinB_synapse_total_sig_grid.png"),
    (SUB_SCATTER,    "CathepsinB_total_sig_grid.png"),
    (SUB_SCATTER,    "CathepsinB_peak_sig_grid.png"),
    (SUB_SCATTER,    "CathepsinB_synapse_inner_outer_ratio_50pct_grid.png"),
    (SUB_SCATTER,    "CathepsinB_synapse_inner_outer_ratio_70pct_grid.png"),
    (SUB_SCATTER,    "CathepsinB_synapse_inner_mask_MFI_grid.png"),
    (SUB_SCATTER,    "CathepsinB_synapse_outer_mask_MFI_grid.png"),
    (SUB_SCATTER,    "CathepsinB_synapse_r_eff_grid.png"),
    (SUB_CAT_VS_FMC, "CathepsinB_synapse_rad_profile_all_cells_with_average_grid.png"),
    (SUB_SCATTER,    "CathepsinB_z50_rel_cell_bottom_grid.png"),
    (SUB_SCATTER,    "CathepsinB_z50_norm_cell_height_grid.png"),
    (SUB_SCATTER,    "CathepsinB_z75_rel_cell_bottom_grid.png"),
    (SUB_SCATTER,    "CathepsinB_z75_norm_cell_height_grid.png"),
    (SUB_SCATTER,    "CathepsinB_z90_rel_cell_bottom_grid.png"),
    (SUB_SCATTER,    "CathepsinB_z90_norm_cell_height_grid.png"),
]

# 4 CatB Z-center-of-fluorescence variants.
CTSB_ZCOF = [
    (SUB_SCATTER, "CathepsinB_zCOF_grid.png"),
    (SUB_SCATTER, "CathepsinB_zCOF_actin_scale_grid.png"),
    (SUB_SCATTER, "CathepsinB_zCOF_cell_bottom_distance_grid.png"),
    (SUB_SCATTER, "CathepsinB_zCOF_cell_bottom_distance_norm_cell_height_grid.png"),
]

# 6 CatB FDD/structural metrics: 3 z_FDD (1D fluorescence distribution
# derivative in z) + 3 FDD_3D variants (base, RMS, rel_cent).
CTSB_FDD = [
    (SUB_SCATTER, "CathepsinB_z_FDD_grid.png"),
    (SUB_SCATTER, "CathepsinB_z_FDD_RMS_grid.png"),
    (SUB_SCATTER, "CathepsinB_z_FDD_rel_cent_grid.png"),
    (SUB_SCATTER, "CathepsinB_FDD_3D_grid.png"),
    (SUB_SCATTER, "CathepsinB_FDD_3D_RMS_grid.png"),
    (SUB_SCATTER, "CathepsinB_FDD_3D_rel_cent_grid.png"),
]

CTSB_CORR_METRICS = [
    (SUB_SCATTER,    "CathepsinB_synapse_g_ave_grid.png"),
    (SUB_CAT_VS_FMC, "CathepsinB_synapse_autocorr_c_all_cells_with_average_grid.png"),
]

FOCI_METRICS_CTSB = [
    (SUB_SCATTER, "actin_foci_count_grid.png"),
    (SUB_SCATTER, "actin_foci_area_um2_grid.png"),
    (SUB_SCATTER, "actin_foci_area_fraction_grid.png"),
    (SUB_SCATTER, "actin_foci_mean_intensity_grid.png"),
]

# Centrosome-Synapse Distance (raw + polarized). Polarization grid lives
# directly at grid_panels/ top level (subfolder "").
CTSB_CENTROSOME = [
    (SUB_SCATTER, "centrosome_center_z_cell_bottom_distance_grid.png"),
    ("",          "centrosome_center_z_cell_bottom_distance_polarization_grid.png"),
]

# 3 CatB centrosome-proximity slides (fraction of CatB signal within
# 1/2/3 μm of the centrosome).
CTSB_FRAC_AROUND_CENT = [
    (SUB_CAT_VS_FMC, "CathepsinB_frac_around_cent_1um_grid.png"),
    (SUB_CAT_VS_FMC, "CathepsinB_frac_around_cent_2um_grid.png"),
    (SUB_CAT_VS_FMC, "CathepsinB_frac_around_cent_3um_grid.png"),
]

# CatB Z50 / Z75 polarization grids (top-level grid_panels/).
CTSB_Z_POLARIZATION = [
    ("", "CathepsinB_z50_rel_cell_bottom_polarization_grid.png"),
    ("", "CathepsinB_z75_rel_cell_bottom_polarization_grid.png"),
]

# Actin single-slice with the non-auc1 rad_profile variant — this
# compile lacks the `_auc1_all_cells_with_average` variant, only the
# plain `_all_cells_with_average` exists. 3-slice MIP has NO
# rad_profile at all in this compile.
CTSB_ACTIN_METRICS_WITH_RAD = (
    SLICE_METRICS[:3]
    + [(SUB_CAT_VS_FMC, "actin_bottom_slice_rad_profile_all_cells_with_average_grid.png")]
    + SLICE_METRICS[3:]
)

# 6 Foci x CatB colocalization slides. This compile uses `foci_CatB_*`
# (short "CatB"), not `foci_CathepsinB_*` like the earlier MT_CatB compile.
FOCI_CTSB_COLOC = [
    (SUB_SCATTER, "foci_CatB_enrichment_ratio_grid.png"),
    (SUB_SCATTER, "foci_CatB_mean_intensity_synapse_grid.png"),
    (SUB_SCATTER, "foci_CatB_mean_intensity_in_foci_grid.png"),
    (SUB_SCATTER, "foci_CatB_mean_intensity_out_foci_grid.png"),
    (SUB_SCATTER, "foci_CatB_m1_intensity_in_foci_grid.png"),
    (SUB_SCATTER, "foci_CatB_pearsons_coeff_grid.png"),
]

# 13 pPKC (phospho-PKC theta) metrics — mirrors PMLC_METRICS/PKC_METRICS.
PPKC_METRICS = [
    (SUB_SCATTER,    "pPKC_synapse_MFI_grid.png"),
    (SUB_SCATTER,    "pPKC_synapse_total_sig_grid.png"),
    (SUB_SCATTER,    "pPKC_total_sig_grid.png"),
    (SUB_SCATTER,    "pPKC_peak_sig_grid.png"),
    (SUB_SCATTER,    "pPKC_synapse_inner_outer_ratio_50pct_grid.png"),
    (SUB_SCATTER,    "pPKC_synapse_inner_outer_ratio_70pct_grid.png"),
    (SUB_SCATTER,    "pPKC_synapse_inner_mask_MFI_grid.png"),
    (SUB_SCATTER,    "pPKC_synapse_outer_mask_MFI_grid.png"),
    (SUB_SCATTER,    "pPKC_synapse_r_eff_grid.png"),
    (SUB_CAT_VS_FMC, "pPKC_synapse_rad_profile_auc1_all_cells_with_average_grid.png"),
    (SUB_SCATTER,    "pPKC_z50_rel_cell_bottom_grid.png"),
    (SUB_SCATTER,    "pPKC_z75_rel_cell_bottom_grid.png"),
    (SUB_SCATTER,    "pPKC_z90_rel_cell_bottom_grid.png"),
]

PPKC_CORR_METRICS = [
    (SUB_SCATTER,    "pPKC_synapse_g_ave_grid.png"),
    (SUB_CAT_VS_FMC, "pPKC_synapse_autocorr_c_all_cells_with_average_grid.png"),
]

FOCI_METRICS_PPKC = [
    (SUB_SCATTER, "actin_foci_count_grid.png"),
    (SUB_SCATTER, "actin_foci_area_um2_grid.png"),
    (SUB_SCATTER, "actin_foci_area_fraction_grid.png"),
    (SUB_SCATTER, "actin_foci_mean_intensity_grid.png"),
]

# 13 PKC metrics — 12 in scatter, 1 (rad_profile) in CAT_vs_FMC.
# Same shape as PMLC_METRICS; ends with Z50/Z75/Z90 rel_cell_bottom.
PKC_METRICS = [
    (SUB_SCATTER,    "PKC_synapse_MFI_grid.png"),
    (SUB_SCATTER,    "PKC_synapse_total_sig_grid.png"),
    (SUB_SCATTER,    "PKC_total_sig_grid.png"),
    (SUB_SCATTER,    "PKC_peak_sig_grid.png"),
    (SUB_SCATTER,    "PKC_synapse_inner_outer_ratio_50pct_grid.png"),
    (SUB_SCATTER,    "PKC_synapse_inner_outer_ratio_70pct_grid.png"),
    (SUB_SCATTER,    "PKC_synapse_inner_mask_MFI_grid.png"),
    (SUB_SCATTER,    "PKC_synapse_outer_mask_MFI_grid.png"),
    (SUB_SCATTER,    "PKC_synapse_r_eff_grid.png"),
    (SUB_CAT_VS_FMC, "PKC_synapse_rad_profile_auc1_all_cells_with_average_grid.png"),
    (SUB_SCATTER,    "PKC_z50_rel_cell_bottom_grid.png"),
    (SUB_SCATTER,    "PKC_z75_rel_cell_bottom_grid.png"),
    (SUB_SCATTER,    "PKC_z90_rel_cell_bottom_grid.png"),
]

# 2 PKC correlation curves — cross-corr and foci_PKC not yet generated.
PKC_CORR_METRICS = [
    (SUB_SCATTER,    "PKC_synapse_g_ave_grid.png"),
    (SUB_CAT_VS_FMC, "PKC_synapse_autocorr_c_all_cells_with_average_grid.png"),
]

FOCI_METRICS_PKC = [
    (SUB_SCATTER, "actin_foci_count_grid.png"),
    (SUB_SCATTER, "actin_foci_area_um2_grid.png"),
    (SUB_SCATTER, "actin_foci_area_fraction_grid.png"),
    (SUB_SCATTER, "actin_foci_mean_intensity_grid.png"),
]

# 13 pMLC metrics — 12 in scatter, 1 (rad_profile) in CAT_vs_FMC.
# Ends with the 3 Z-distribution slides (z50/z75/z90 rel_cell_bottom).
PMLC_METRICS = [
    (SUB_SCATTER,    "pMLC_synapse_MFI_grid.png"),
    (SUB_SCATTER,    "pMLC_synapse_total_sig_grid.png"),
    (SUB_SCATTER,    "pMLC_total_sig_grid.png"),
    (SUB_SCATTER,    "pMLC_peak_sig_grid.png"),
    (SUB_SCATTER,    "pMLC_synapse_inner_outer_ratio_50pct_grid.png"),
    (SUB_SCATTER,    "pMLC_synapse_inner_outer_ratio_70pct_grid.png"),
    (SUB_SCATTER,    "pMLC_synapse_inner_mask_MFI_grid.png"),
    (SUB_SCATTER,    "pMLC_synapse_outer_mask_MFI_grid.png"),
    (SUB_SCATTER,    "pMLC_synapse_r_eff_grid.png"),
    (SUB_CAT_VS_FMC, "pMLC_synapse_rad_profile_auc1_all_cells_with_average_grid.png"),
    (SUB_SCATTER,    "pMLC_z50_rel_cell_bottom_grid.png"),
    (SUB_SCATTER,    "pMLC_z75_rel_cell_bottom_grid.png"),
    (SUB_SCATTER,    "pMLC_z90_rel_cell_bottom_grid.png"),
]

# 2 pMLC correlation curves — cross-corr and foci_pMLC family not yet
# generated in the 20260702 compile.
PMLC_CORR_METRICS = [
    (SUB_SCATTER,    "pMLC_synapse_g_ave_grid.png"),
    (SUB_CAT_VS_FMC, "pMLC_synapse_autocorr_c_all_cells_with_average_grid.png"),
]

# 4 actin foci — no `foci_pMLC_*` colocalization PNGs yet, so this
# block is smaller than pZap70/pLAT (which have 5 more slides here).
FOCI_METRICS_PMLC = [
    (SUB_SCATTER, "actin_foci_count_grid.png"),
    (SUB_SCATTER, "actin_foci_area_um2_grid.png"),
    (SUB_SCATTER, "actin_foci_area_fraction_grid.png"),
    (SUB_SCATTER, "actin_foci_mean_intensity_grid.png"),
]

# Autocorrelation / cross-correlation curves — 3 slides that show the
# actual g(r) and its scalar summary. `_all_cells_with_average` variants
# match the rad_profile convention (per-cell overlay + population mean).
PZAP70_CORR_METRICS = [
    (SUB_SCATTER,    "pZap70_synapse_g_ave_grid.png"),
    (SUB_CAT_VS_FMC, "pZap70_synapse_autocorr_c_all_cells_with_average_grid.png"),
    (SUB_CAT_VS_FMC, "foci_pZap70_cross_corr_profile_all_cells_with_average_grid.png"),
]

# 10 pZap70 metrics (9 from timecourse_scatter_plots, 1 from
# CAT_vs_FMC for the rad_profile AUC1).
PZAP70_METRICS = [
    (SUB_SCATTER,    "pZap70_synapse_MFI_grid.png"),
    (SUB_SCATTER,    "pZap70_synapse_total_sig_grid.png"),
    (SUB_SCATTER,    "pZap70_total_sig_grid.png"),
    (SUB_SCATTER,    "pZap70_peak_sig_grid.png"),
    (SUB_SCATTER,    "pZap70_synapse_inner_outer_ratio_50pct_grid.png"),
    (SUB_SCATTER,    "pZap70_synapse_inner_outer_ratio_70pct_grid.png"),
    (SUB_SCATTER,    "pZap70_synapse_inner_mask_MFI_grid.png"),
    (SUB_SCATTER,    "pZap70_synapse_outer_mask_MFI_grid.png"),
    (SUB_SCATTER,    "pZap70_synapse_r_eff_grid.png"),
    (SUB_CAT_VS_FMC, "pZap70_synapse_rad_profile_auc1_all_cells_with_average_grid.png"),
]

# 9 foci metrics: 4 actin foci + 5 foci x pZap70 colocalization.
FOCI_METRICS_PZAP70 = [
    (SUB_SCATTER, "actin_foci_count_grid.png"),
    (SUB_SCATTER, "actin_foci_area_um2_grid.png"),
    (SUB_SCATTER, "actin_foci_area_fraction_grid.png"),
    (SUB_SCATTER, "actin_foci_mean_intensity_grid.png"),
    (SUB_SCATTER, "foci_pZap70_enrichment_ratio_grid.png"),
    (SUB_SCATTER, "foci_pZap70_mean_intensity_in_foci_grid.png"),
    (SUB_SCATTER, "foci_pZap70_mean_intensity_out_foci_grid.png"),
    (SUB_SCATTER, "foci_pZap70_m1_intensity_in_foci_grid.png"),
    (SUB_SCATTER, "foci_pZap70_pearsons_coeff_grid.png"),
]

# pZap70 MFI binned by distance to nearest actin focus, per acquisition
# date. Same shape as PLAT_DIST_BIN_PER_DATE (3-tuple title-suffix entry);
# this compile has 4 date subdirs (Sep 21 2023, Oct 3 2023, Feb 17 2025
# D3, Feb 17 2025 D5 — chronological in the deck).
PZAP70_DIST_BIN_PER_DATE = [
    ("binned_distribution_plots/Sep_21,_2023_(6_min)",
     "foci_pZap70_dist_bin_mfi_norm_grid.png", "(Sep 21, 2023 — 6 min)"),
    ("binned_distribution_plots/Oct_03,_2023_(6_min)",
     "foci_pZap70_dist_bin_mfi_norm_grid.png", "(Oct 3, 2023 — 6 min)"),
    ("binned_distribution_plots/Feb_17,_2025_(D3)",
     "foci_pZap70_dist_bin_mfi_norm_grid.png", "(Feb 17, 2025 — D3)"),
    ("binned_distribution_plots/Feb_17,_2025_(D5)",
     "foci_pZap70_dist_bin_mfi_norm_grid.png", "(Feb 17, 2025 — D5)"),
]

# Actin block in the pLAT-20260701 layout (relocated from
# SLICE_METRICS_WITH_RAD to use SUB_PLAT_SCATTER / SUB_PLAT_CAT).
PLAT_ACTIN_METRICS_WITH_RAD = (
    [(SUB_PLAT_SCATTER, fn) for _, fn in SLICE_METRICS[:3]]
    + [(SUB_PLAT_CAT, RAD_PROFILE_SLICE[1])]
    + [(SUB_PLAT_SCATTER, fn) for _, fn in SLICE_METRICS[3:]]
)

# 10 pLAT-specific metrics — 9 from scatter_timecourse_5_vs_15,
# 1 from CAT_vs_FMC (rad_profile AUC1).
PLAT_METRICS = [
    (SUB_PLAT_SCATTER, "pLAT_synapse_MFI_grid.png"),
    (SUB_PLAT_SCATTER, "pLAT_synapse_total_sig_grid.png"),
    (SUB_PLAT_SCATTER, "pLAT_total_sig_grid.png"),
    (SUB_PLAT_SCATTER, "pLAT_peak_sig_grid.png"),
    (SUB_PLAT_SCATTER, "pLAT_synapse_inner_outer_ratio_50pct_grid.png"),
    (SUB_PLAT_SCATTER, "pLAT_synapse_inner_outer_ratio_70pct_grid.png"),
    (SUB_PLAT_SCATTER, "pLAT_synapse_inner_mask_MFI_grid.png"),
    (SUB_PLAT_SCATTER, "pLAT_synapse_outer_mask_MFI_grid.png"),
    (SUB_PLAT_SCATTER, "pLAT_synapse_r_eff_grid.png"),
    (SUB_PLAT_CAT,     "pLAT_synapse_rad_profile_auc1_all_cells_with_average_grid.png"),
]

# 9 foci metrics: 4 actin foci + 5 foci x pLAT colocalization.
FOCI_METRICS_PLAT = [
    (SUB_PLAT_SCATTER, "actin_foci_count_grid.png"),
    (SUB_PLAT_SCATTER, "actin_foci_area_um2_grid.png"),
    (SUB_PLAT_SCATTER, "actin_foci_area_fraction_grid.png"),
    (SUB_PLAT_SCATTER, "actin_foci_mean_intensity_grid.png"),
    (SUB_PLAT_SCATTER, "foci_pLAT_enrichment_ratio_grid.png"),
    (SUB_PLAT_SCATTER, "foci_pLAT_mean_intensity_in_foci_grid.png"),
    (SUB_PLAT_SCATTER, "foci_pLAT_mean_intensity_out_foci_grid.png"),
    (SUB_PLAT_SCATTER, "foci_pLAT_m1_intensity_in_foci_grid.png"),
    (SUB_PLAT_SCATTER, "foci_pLAT_pearsons_coeff_grid.png"),
]

# Single combined dist-bin plot at top level of binned_distribution_plots
# (no per-date breakdown in the 20260701 compile).
PLAT_DIST_BIN = [
    ("binned_distribution_plots", "foci_pLAT_dist_bin_mfi_norm_grid.png"),
]

CONFOCAL_3SLICE_ROOT = Path(
    "L:/FF/CAR T/actin_compiled_results/all_datasets_actin/"
    "confocal_3slice/compiled_20260625/all_conditions"
)

# All 20 PNGs from the confocal_3slice/compiled_20260625/all_conditions dir.
# Subfolder is "" because PNGs sit directly under compiled_root. Canonical
# `_grid.png` names so TITLE_OVERRIDES + prettify resolve as elsewhere;
# strip_grid_suffix=True on the dataset entry strips it for disk lookup.
CONFOCAL_3SLICE_METRICS = [
    # area + spatial
    ("", "actin_bottom_3slice_mask_area_grid.png"),
    ("", "actin_bottom_3slice_mip_inner_outer_ratio_50pct_grid.png"),
    ("", "actin_bottom_3slice_mip_inner_outer_ratio_50pct_ylim_0_3_grid.png"),
    ("", "actin_bottom_3slice_mip_inner_outer_ratio_70pct_grid.png"),
    ("", "actin_bottom_3slice_mip_inner_outer_ratio_70pct_ylim_0_3_grid.png"),
    # radial profile variants
    ("", "actin_bottom_3slice_mip_rad_profile_grid.png"),
    ("", "actin_bottom_3slice_mip_rad_profile_all_cells_grid.png"),
    ("", "actin_bottom_3slice_mip_rad_profile_all_cells_with_average_grid.png"),
    ("", "actin_bottom_3slice_mip_rad_profile_auc1_grid.png"),
    ("", "actin_bottom_3slice_mip_rad_profile_auc1_all_cells_grid.png"),
    ("", "actin_bottom_3slice_mip_rad_profile_auc1_all_cells_with_average_grid.png"),
    # intensity
    ("", "actin_bottom_3slice_MFI_grid.png"),
    ("", "actin_bottom_3slice_total_sig_grid.png"),
    # shape
    ("", "actin_bottom_3slice_perimeter_grid.png"),
    ("", "actin_bottom_3slice_circularity_grid.png"),
    ("", "actin_bottom_3slice_solidity_grid.png"),
    ("", "actin_bottom_3slice_eccentricity_grid.png"),
    # whole-cell MIP
    ("", "actin_MIP_area_grid.png"),
    ("", "actin_MIP_major_axis_length_grid.png"),
    ("", "actin_MIP_minor_axis_length_grid.png"),
]


def _by_timepoint_all_datasets(tp_short: str, date: str = "20260625") -> Path:
    """Path to one timepoint's all_conditions/ dir in the L: drive
    `all_datasets_actin` compile."""
    return (
        ALL_DATASETS_BY_TIMEPOINT_ROOT
        / f"all_{tp_short}"
        / f"compiled_{date}"
        / "all_conditions"
    )


def _by_timepoint_confocal_3slice(tp_short: str, date: str = "20260625") -> Path:
    """Path to one timepoint's all_conditions/ dir under the confocal_3slice
    by_timepoint tree (TIRF datasets are excluded — no 3-slice MIP for them)."""
    return (
        Path("L:/FF/CAR T/actin_compiled_results/all_datasets_actin/confocal_3slice/by_timepoint")
        / f"all_{tp_short}"
        / f"compiled_{date}"
        / "all_conditions"
    )


# Flat png_files for the 3-slice-only stacked deck. 10 metrics — same
# shape and order as SLICE_ONLY_PNG_FILES but for the 3-slice MIP block.
THREE_SLICE_ONLY_PNG_FILES = [
    ("", fn) for _, fn in THREE_SLICE_METRICS_WITH_RAD
]

# ---------------------------------------------------------------------------
# MT/CatB compile (CART_summary_MT_CatB deck)

MT_CATB_ROOT = Path(
    "J:/FF/fixed_cell/CAR_TCell/results_compiled/"
    "MT_CATB_20231127_20240620_20240624/compiled_20260702/grid_panels"
)
OUTPUT_DIR_MT_CATB = Path("K:/FF/PPT/PPT_autogeneration/CART/MT_CatB")

# 1 centrosome slide (user picked the z-distance variant; skipped the
# center_slice_from_MT variant).
CENTROSOME_METRICS = [
    (SUB_SCATTER, "centrosome_center_z_cell_bottom_distance_grid.png"),
]

# All 48 CathepsinB PNGs, grouped by analysis category.
CATB_METRICS = [
    # --- synapse intensity & spatial (16) ----------------------------------
    (SUB_SCATTER, "CathepsinB_synapse_MFI_grid.png"),
    (SUB_SCATTER, "CathepsinB_synapse_MFI_3mip_grid.png"),
    (SUB_SCATTER, "CathepsinB_synapse_total_sig_grid.png"),
    (SUB_SCATTER, "CathepsinB_synapse_total_sig_3mip_grid.png"),
    (SUB_SCATTER, "CathepsinB_synapse_inner_outer_ratio_grid.png"),
    (SUB_SCATTER, "CathepsinB_synapse_inner_outer_ratio_3mip_grid.png"),
    (SUB_SCATTER, "CathepsinB_synapse_inner_outer_ratio_50pct_grid.png"),
    (SUB_SCATTER, "CathepsinB_synapse_inner_outer_ratio_70pct_grid.png"),
    (SUB_SCATTER, "CathepsinB_synapse_inner_mask_MFI_grid.png"),
    (SUB_SCATTER, "CathepsinB_synapse_inner_mask_MFI_3mip_grid.png"),
    (SUB_SCATTER, "CathepsinB_synapse_outer_mask_MFI_grid.png"),
    (SUB_SCATTER, "CathepsinB_synapse_outer_mask_MFI_3mip_grid.png"),
    (SUB_SCATTER, "CathepsinB_synapse_g_ave_grid.png"),
    (SUB_SCATTER, "CathepsinB_synapse_g_ave_3mip_grid.png"),
    (SUB_SCATTER, "CathepsinB_synapse_r_eff_grid.png"),
    (SUB_SCATTER, "CathepsinB_synapse_r_eff_3mip_grid.png"),
    # --- synapse autocorr (2) ----------------------------------------------
    (SUB_SCATTER, "CathepsinB_synapse_autocorr_rmax_um_grid.png"),
    (SUB_SCATTER, "CathepsinB_synapse_autocorr_smoothSig_grid.png"),
    # --- total + peak signal (2) -------------------------------------------
    (SUB_SCATTER, "CathepsinB_total_sig_grid.png"),
    (SUB_SCATTER, "CathepsinB_peak_sig_grid.png"),
    # --- granule morphology (6) --------------------------------------------
    (SUB_SCATTER, "CathepsinB_granule_area_um2_grid.png"),
    (SUB_SCATTER, "CathepsinB_granule_area_fraction_grid.png"),
    (SUB_SCATTER, "CathepsinB_granule_mean_intensity_segmented_grid.png"),
    (SUB_SCATTER, "CathepsinB_granule_mean_intensity_nonsegmented_grid.png"),
    (SUB_SCATTER, "CathepsinB_granule_total_intensity_segmented_grid.png"),
    (SUB_SCATTER, "CathepsinB_granule_total_intensity_nonsegmented_grid.png"),
    # --- centrosome proximity (6) ------------------------------------------
    (SUB_SCATTER, "CathepsinB_MFI_around_cent_1um_grid.png"),
    (SUB_SCATTER, "CathepsinB_MFI_around_cent_2um_grid.png"),
    (SUB_SCATTER, "CathepsinB_MFI_around_cent_3um_grid.png"),
    (SUB_SCATTER, "CathepsinB_frac_around_cent_1um_grid.png"),
    (SUB_SCATTER, "CathepsinB_frac_around_cent_2um_grid.png"),
    (SUB_SCATTER, "CathepsinB_frac_around_cent_3um_grid.png"),
    # --- Z-distribution (10) -----------------------------------------------
    (SUB_SCATTER, "CathepsinB_z50_rel_cell_bottom_grid.png"),
    (SUB_SCATTER, "CathepsinB_z50_norm_cell_height_grid.png"),
    (SUB_SCATTER, "CathepsinB_z75_rel_cell_bottom_grid.png"),
    (SUB_SCATTER, "CathepsinB_z75_norm_cell_height_grid.png"),
    (SUB_SCATTER, "CathepsinB_z90_rel_cell_bottom_grid.png"),
    (SUB_SCATTER, "CathepsinB_z90_norm_cell_height_grid.png"),
    (SUB_SCATTER, "CathepsinB_zCOF_grid.png"),
    (SUB_SCATTER, "CathepsinB_zCOF_actin_scale_grid.png"),
    (SUB_SCATTER, "CathepsinB_zCOF_cell_bottom_distance_grid.png"),
    (SUB_SCATTER, "CathepsinB_zCOF_cell_bottom_distance_norm_cell_height_grid.png"),
    # --- FDD / structural (6) ----------------------------------------------
    (SUB_SCATTER, "CathepsinB_FDD_3D_grid.png"),
    (SUB_SCATTER, "CathepsinB_FDD_3D_RMS_grid.png"),
    (SUB_SCATTER, "CathepsinB_FDD_3D_rel_cent_grid.png"),
    (SUB_SCATTER, "CathepsinB_z_FDD_grid.png"),
    (SUB_SCATTER, "CathepsinB_z_FDD_RMS_grid.png"),
    (SUB_SCATTER, "CathepsinB_z_FDD_rel_cent_grid.png"),
]

CATB_Z50_METRICS = [
    (SUB_SCATTER, "CathepsinB_z50_rel_cell_bottom_grid.png"),
    (SUB_SCATTER, "CathepsinB_z50_norm_cell_height_grid.png"),
]
CATB_Z75_METRICS = [
    (SUB_SCATTER, "CathepsinB_z75_rel_cell_bottom_grid.png"),
    (SUB_SCATTER, "CathepsinB_z75_norm_cell_height_grid.png"),
]

CATB_METRICS_CONCISE = [
    # spatial concentration (granule_area_um2 dropped per user)
    (SUB_SCATTER, "CathepsinB_synapse_inner_outer_ratio_50pct_grid.png"),
    (SUB_SCATTER, "CathepsinB_synapse_inner_outer_ratio_70pct_grid.png"),
    # intensity
    (SUB_SCATTER, "CathepsinB_synapse_MFI_grid.png"),
    (SUB_SCATTER, "CathepsinB_synapse_total_sig_grid.png"),
    (SUB_SCATTER, "CathepsinB_total_sig_grid.png"),
    # centrosome-proximity (fraction of CatB signal within N μm of centrosome)
    (SUB_SCATTER, "CathepsinB_frac_around_cent_1um_grid.png"),
    (SUB_SCATTER, "CathepsinB_frac_around_cent_2um_grid.png"),
    (SUB_SCATTER, "CathepsinB_frac_around_cent_3um_grid.png"),
]

# Polarization grid PNGs: one per target field; each PNG embeds all
# configured thresholds in a panel grid. Output by
# assemble_polarization_grids.m at `<compiled_root>/<target>_polarization_grid.png`
# (directly under grid_panels/, NOT in a sub-folder). Filenames will be
# MISSING until the MATLAB compile is rerun with the polarization flag
# turned on for MT_CatB (config: d20231127_20240620_20240624_compilation_config.m).
POLARIZATION_METRICS = [
    ("", "centrosome_center_z_cell_bottom_distance_polarization_grid.png"),
    ("", "CathepsinB_z50_rel_cell_bottom_polarization_grid.png"),
    ("", "CathepsinB_z75_rel_cell_bottom_polarization_grid.png"),
]

# 10 foci PNGs: 4 actin foci + 6 foci x CatB coloc. The 20260702
# recompile uses `foci_CatB_*` prefix and dropped several older
# `foci_CathepsinB_granule_*` / `_overlap_fraction` variants.
FOCI_METRICS = [
    # --- actin foci (4) ----------------------------------------------------
    (SUB_SCATTER, "actin_foci_count_grid.png"),
    (SUB_SCATTER, "actin_foci_area_um2_grid.png"),
    (SUB_SCATTER, "actin_foci_area_fraction_grid.png"),
    (SUB_SCATTER, "actin_foci_mean_intensity_grid.png"),
    # --- foci x CatB colocalization (6) ------------------------------------
    (SUB_SCATTER, "foci_CatB_enrichment_ratio_grid.png"),
    (SUB_SCATTER, "foci_CatB_mean_intensity_synapse_grid.png"),
    (SUB_SCATTER, "foci_CatB_mean_intensity_in_foci_grid.png"),
    (SUB_SCATTER, "foci_CatB_mean_intensity_out_foci_grid.png"),
    (SUB_SCATTER, "foci_CatB_m1_intensity_in_foci_grid.png"),
    (SUB_SCATTER, "foci_CatB_pearsons_coeff_grid.png"),
]

# 3 per-date binned distribution plots: CatB MFI by distance to nearest
# actin focus, per acquisition date. Same 3-tuple (sub, fn, title_suffix)
# pattern as PLAT_DIST_BIN_PER_DATE / PZAP70_DIST_BIN_PER_DATE.
MT_CATB_DIST_BIN_PER_DATE = [
    ("binned_distribution_plots/Nov_27,_2023",
     "foci_CatB_dist_bin_mfi_norm_grid.png", "(Nov 27, 2023)"),
    ("binned_distribution_plots/Jun_20,_2024",
     "foci_CatB_dist_bin_mfi_norm_grid.png", "(Jun 20, 2024)"),
    ("binned_distribution_plots/Jun_24,_2024",
     "foci_CatB_dist_bin_mfi_norm_grid.png", "(Jun 24, 2024)"),
]

# Each dataset -> one deck. compiled_root points at the grid_panels/ folder.
OUTPUT_DIR = Path("K:/FF/PPT/PPT_autogeneration/CART/actin_only")

DATASETS = [
    {
        "label": "Kiet 2026",
        "compiled_root": Path(
            "Y:/User_data/Kiet/results_compiled/actin_only/"
            "compiled_20260312_20260414_20260510_20260604/grid_panels"
        ),
        "output_path": OUTPUT_DIR / "CART_actin_summary.pptx",
        # Kiet's CAT_vs_FMC dir has no 3slice rad_profile PNG, so we
        # leave both rad_profile entries off this deck (18 slides).
        "png_files": SLICE_METRICS + THREE_SLICE_METRICS,
    },
    {
        "label": "2024-06 data",
        "compiled_root": Path(
            "J:/FF/fixed_cell/CAR_TCell/results_compiled/actin_only/"
            "compiled_20240620_20240624_20260605/grid_panels"
        ),
        "output_path": OUTPUT_DIR / "CART_actin_summary_202406_data.pptx",
        "png_files": SLICE_METRICS_WITH_RAD + THREE_SLICE_METRICS_WITH_RAD,
    },
    {
        "label": "Kiet 20260607",
        # Compile name: <dataset_folder>_<compile_date> (no `_test_` segment).
        "compiled_root": Path(
            "Y:/User_data/Kiet/results_compiled/actin_only/"
            "compiled_20260607_pMLC_CART_actin_hoescht_20260616/grid_panels"
        ),
        "output_path": OUTPUT_DIR / "CART_actin_summary_20260607.pptx",
        "png_files": SLICE_METRICS_WITH_RAD + THREE_SLICE_METRICS_WITH_RAD,
    },
    {
        "label": "20260717 + 20260723 (pooled)",
        # Two experiment dates (Jul 17 + Jul 23, 2026) pooled; compile-run
        # 2026-08-11. Same grid_panels/{timecourse_scatter_plots,
        # CAT_vs_FMC_by_timepoint} layout as the "2024-06 data" entry.
        "compiled_root": Path(
            "H:/results_compiled/actin_only/"
            "compiled_20260717_20260723_20260811/grid_panels"
        ),
        "output_path": OUTPUT_DIR / "CART_actin_summary_20260717_20260723.pptx",
        # Splice extras into the standard blocks:
        #   slice  : mask_area | IRM | IOR50 | IOR70 | R_median | rad_prof | ...
        #   3-slice: mask_area | IOR50 | IOR70 | R_median | rad_prof | ...
        # Axial-overlay slides come from a peer compile dir (absolute Path
        # in the "sub" slot bypasses compiled_root via pathlib semantics).
        "png_files": (
            SLICE_METRICS_WITH_RAD[:1]                # mask_area
            + [IRM_SYNAPSE_AREA]                      # + IRM
            + SLICE_METRICS_WITH_RAD[1:3]             # IOR50, IOR70
            + [R_MEDIAN_SLICE]                        # + R_median
            + SLICE_METRICS_WITH_RAD[3:]              # rad_prof, MFI, ...
            + THREE_SLICE_METRICS_WITH_RAD[:3]        # 3-slice mask_area, IOR50/70
            + [R_MEDIAN_3SLICE]                       # + R_median (3-slice)
            + THREE_SLICE_METRICS_WITH_RAD[3:]        # rad_prof, MFI, ...
            + ACTIN_FOCI_METRICS
            + [
                (Path("H:/results_compiled/actin_only/"
                      "compiled_20260717_20260723/axial_profile_overlays"),
                 "axial_overlay_perslice_Jul_17_2026.png"),
                (Path("H:/results_compiled/actin_only/"
                      "compiled_20260717_20260723/axial_profile_overlays"),
                 "axial_overlay_perslice_Jul_23_2026.png"),
            ]
        ),
    },
    {
        "label": "5min by_timepoint",
        # Flat layout: PNGs live directly under compiled_root with no
        # _grid suffix. strip_grid_suffix=True tells build_deck to drop
        # `_grid` when resolving disk paths while keeping the canonical
        # names for prettify / TITLE_OVERRIDES lookup.
        "compiled_root": _by_timepoint("5min"),
        "output_path": OUTPUT_DIR / "CART_actin_summary_5min.pptx",
        "png_files": ALL_CONDITIONS_PNG_FILES,
        "strip_grid_suffix": True,
    },
    {
        "label": "All timepoints",
        # 30 slides: 10 metrics x 3 timepoints, interleaved per-metric.
        # Each metric gets 3 consecutive slides (5/10/15 min). Sourced
        # from the L: drive `all_datasets_actin` compile. Single-slice
        # metrics only — the new compile has no 3-slice MIP PNGs.
        "timepoints": [
            ("5 min",  _by_timepoint_all_datasets("5min")),
            ("10 min", _by_timepoint_all_datasets("10min")),
            ("15 min", _by_timepoint_all_datasets("15min")),
        ],
        "output_path": OUTPUT_DIR / "CART_actin_summary_all_timepoints.pptx",
        "png_files": SLICE_ONLY_PNG_FILES,
        "strip_grid_suffix": True,
    },
    {
        "label": "CatB (CTSB all datasets)",
        # 64 slides: 2 centrosome + 16 CatB (Z rel_cell_bottom + norm_
        # cell_height interleaved) + 4 CatB Z-COF + 6 CatB FDD + 3 CatB
        # frac_around_cent + 2 CatB Z-polarization + 10 actin + 9 actin
        # 3-slice MIP (no rad_profile) + 2 correlation + 4 actin foci
        # + 6 foci x CatB coloc.
        "compiled_root": CTSB_ALL_ROOT,
        "output_path": OUTPUT_DIR_CTSB / "CART_CatB_summary.pptx",
        "png_files": (
            CTSB_CENTROSOME                # 2 centrosome dist + polarized
            + CTSB_METRICS                 # 16 CatB (incl. Z rel + norm)
            + CTSB_ZCOF                    # 4 Z-COF variants
            + CTSB_FDD                     # 6 FDD/structural
            + CTSB_FRAC_AROUND_CENT        # 3 frac around cent
            + CTSB_Z_POLARIZATION          # 2 Z50/Z75 polarization
            + CTSB_ACTIN_METRICS_WITH_RAD  # 10 actin (non-auc1 rad)
            + THREE_SLICE_METRICS          # 9 actin 3-slice MIP (no rad)
            + CTSB_CORR_METRICS            # 2 correlation
            + FOCI_METRICS_CTSB            # 4 actin foci
            + FOCI_CTSB_COLOC              # 6 foci x CatB coloc
        ),
    },
    {
        "label": "pPKC",
        # 39 slides: 13 pPKC (leads with Synapse MFI, closes with Z50/
        # Z75/Z90) + 10 actin + 10 actin 3-slice MIP + 2 correlation +
        # 4 actin foci. Confocal compile — mirrors pMLC/PKC deck shape.
        "compiled_root": PPKC_ROOT,
        "output_path": OUTPUT_DIR_PPKC / "CART_pPKC_summary.pptx",
        "png_files": (
            PPKC_METRICS                     # 13 pPKC
            + SLICE_METRICS_WITH_RAD         # 10 actin
            + THREE_SLICE_METRICS_WITH_RAD   # 10 actin 3-slice MIP
            + PPKC_CORR_METRICS              # 2 correlation
            + FOCI_METRICS_PPKC              # 4 foci
        ),
    },
    {
        "label": "PKC",
        # 39 slides: 13 PKC (leads with Synapse MFI, closes with Z50/
        # Z75/Z90) + 10 actin + 10 actin 3-slice MIP + 2 correlation
        # (g_ave + autocorr; no cross-corr) + 4 actin foci (no
        # foci_PKC coloc yet). Confocal compile — mirrors pMLC deck.
        "compiled_root": PKC_ROOT,
        "output_path": OUTPUT_DIR_PKC / "CART_PKC_summary.pptx",
        "png_files": (
            PKC_METRICS                      # 13 PKC
            + SLICE_METRICS_WITH_RAD         # 10 actin
            + THREE_SLICE_METRICS_WITH_RAD   # 10 actin 3-slice MIP
            + PKC_CORR_METRICS               # 2 correlation
            + FOCI_METRICS_PKC               # 4 foci
        ),
    },
    {
        "label": "pMLC joint",
        # 39 slides: 13 pMLC (leads with Synapse MFI, closes with
        # z50/z75/z90) + 10 actin + 10 actin 3-slice MIP + 2 correlation
        # (g_ave + autocorr; cross-corr not yet generated) + 4 actin
        # foci (no foci_pMLC coloc yet). Confocal compile — includes
        # the 3-slice MIP block absent from TIRF pZap70/pLAT decks.
        "compiled_root": PMLC_JOINT_ROOT,
        "output_path": OUTPUT_DIR_PMLC / "CART_pMLC_summary.pptx",
        "png_files": (
            PMLC_METRICS                     # 13 pMLC
            + SLICE_METRICS_WITH_RAD         # 10 actin
            + THREE_SLICE_METRICS_WITH_RAD   # 10 actin 3-slice MIP
            + PMLC_CORR_METRICS              # 2 correlation
            + FOCI_METRICS_PMLC              # 4 foci
        ),
    },
    {
        "label": "pZap70 TIRF",
        # 36 slides: 10 pZap70 (leads with Synapse MFI) + 10 actin +
        # 3 correlation + 9 foci + 4 per-date dist-bin. Sources from
        # the 20260701 recompile which added g(r) autocorrelation and
        # foci x pZap70 cross-correlation.
        "compiled_root": PZAP70_TIRF_ROOT,
        "output_path": OUTPUT_DIR_PZAP70 / "CART_pZap70_summary.pptx",
        "png_files": (
            PZAP70_METRICS               # 10 pZap70 (leads with Synapse MFI)
            + SLICE_METRICS_WITH_RAD     # 10 actin
            + PZAP70_CORR_METRICS        # 3 g(r) autocorr + cross-corr
            + FOCI_METRICS_PZAP70        # 9 foci
            + PZAP70_DIST_BIN_PER_DATE   # 4 per-date dist-bin
        ),
    },
    {
        "label": "pLAT TIRF",
        # 33 slides: 10 pLAT (leads with Synapse MFI) + 10 actin +
        # 3 correlation (g_ave + autocorr + cross-corr) + 9 foci +
        # 1 combined dist-bin. Sourced from the 20260701 compile which
        # renamed subfolders (scatter_timecourse_5_vs_15, CAT_vs_FMC)
        # and consolidated the per-date dist-bin plots into one.
        "compiled_root": PLAT_TIRF_ROOT,
        "output_path": OUTPUT_DIR_PLAT / "CART_pLAT_summary.pptx",
        "png_files": (
            PLAT_METRICS                  # 10 pLAT (leads with Synapse MFI)
            + PLAT_ACTIN_METRICS_WITH_RAD # 10 actin
            + PLAT_CORR_METRICS           # 3 g_ave + autocorr + cross-corr
            + FOCI_METRICS_PLAT           # 9 foci
            + PLAT_DIST_BIN               # 1 combined dist-bin
        ),
    },
    {
        "label": "Confocal 3-slice",
        # 20 slides: every PNG from the confocal_3slice compile, all
        # condition-collapsed (no timepoint split). Includes the
        # canonical 10 3-slice metrics + IOR ylim variants + extra
        # rad_profile variants + 3 whole-cell MIP metrics.
        "compiled_root": CONFOCAL_3SLICE_ROOT,
        "output_path": OUTPUT_DIR / "CART_actin_summary_3slice_confocal.pptx",
        "png_files": CONFOCAL_3SLICE_METRICS,
        "strip_grid_suffix": True,
    },
    {
        "label": "Confocal 3-slice (stacked)",
        # 10 slides: one per 3-slice metric, with 3 timepoints stacked
        # vertically (5 min top, 10 min middle, 15 min bottom). Sourced
        # from the confocal_3slice/by_timepoint compile (TIRF datasets
        # are not represented — no 3-slice MIP for them).
        "timepoints": [
            ("5 min",  _by_timepoint_confocal_3slice("5min")),
            ("10 min", _by_timepoint_confocal_3slice("10min")),
            ("15 min", _by_timepoint_confocal_3slice("15min")),
        ],
        "output_path": OUTPUT_DIR / "CART_actin_summary_3slice_confocal_stacked.pptx",
        "png_files": THREE_SLICE_ONLY_PNG_FILES,
        "strip_grid_suffix": True,
        "stack_timepoints": True,
    },
    {
        "label": "All timepoints (stacked)",
        # 10 slides: one per metric, with 3 timepoints stacked vertically
        # (5 min top, 10 min middle, 15 min bottom). 1/3 the slide count
        # of the per-timepoint version above; same source data.
        "timepoints": [
            ("5 min",  _by_timepoint_all_datasets("5min")),
            ("10 min", _by_timepoint_all_datasets("10min")),
            ("15 min", _by_timepoint_all_datasets("15min")),
        ],
        "output_path": OUTPUT_DIR / "CART_actin_summary_all_timepoints_stacked.pptx",
        "png_files": SLICE_ONLY_PNG_FILES,
        "strip_grid_suffix": True,
        "stack_timepoints": True,
    },
    {
        "label": "MT_CatB",
        # 24 slides: 1 centrosome + 1 cent-polarization + 8 CatB +
        # 2 Z50 (raw + polarized) + 2 Z75 (raw + polarized) + 10 foci
        # (4 actin foci + 6 foci_CatB coloc) + 3 per-date dist_bin.
        # Sourced from the 20260702 recompile which added the per-date
        # binned distribution plots and renamed foci_CathepsinB_ →
        # foci_CatB_.
        "compiled_root": MT_CATB_ROOT,
        "output_path": OUTPUT_DIR_MT_CATB / "CART_summary_MT_CatB.pptx",
        "png_files": (
            CENTROSOME_METRICS                              # 1: distance
            + [POLARIZATION_METRICS[0]]                     # 1: centrosome polarized
            + CATB_METRICS_CONCISE                          # 8: CatB analogs
            + CATB_Z50_METRICS + [POLARIZATION_METRICS[1]]  # 2: Z50 raw + polarized
            + CATB_Z75_METRICS + [POLARIZATION_METRICS[2]]  # 2: Z75 raw + polarized
            + FOCI_METRICS                                  # 10 foci
            + MT_CATB_DIST_BIN_PER_DATE                     # 3 per-date dist_bin
        ),
    },
]

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Slide layout (inches). 13.333 x 7.5 widescreen.
SLIDE_W = 13.333
SLIDE_H = 7.5

MARGIN = 0.10

TITLE_LEFT = MARGIN
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

SUBTITLE_LEFT = MARGIN
SUBTITLE_TOP = 0.55
SUBTITLE_WIDTH = SLIDE_W - 2 * MARGIN
SUBTITLE_HEIGHT = 0.32
SUBTITLE_FONT_PT = 14

# Image placement: bigger when there's no subtitle, smaller when there is.
IMG_LEFT = MARGIN
IMG_TOP = 0.60
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = 6.40

IMG_TOP_WITH_SUBTITLE = 0.90
IMG_BOX_H_WITH_SUBTITLE = 6.10

FOOTER_LEFT = MARGIN
FOOTER_TOP = 7.05
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.40
FOOTER_FONT_PT = 9

# Stacked-timepoint layout: 1 metric per slide, 3 timepoint images
# stacked vertically (5 min on top, 10 in middle, 15 on bottom).
# Title smaller + footer smaller so each row gets more vertical space.
STACK_TITLE_HEIGHT = 0.40
STACK_TITLE_FONT_PT = 22
STACK_FOOTER_TOP = 7.22
STACK_FOOTER_HEIGHT = 0.25
STACK_FOOTER_FONT_PT = 8
STACK_ROW_AREA_TOP = TITLE_TOP + STACK_TITLE_HEIGHT + 0.05   # 0.50
STACK_ROW_AREA_BOTTOM = STACK_FOOTER_TOP - 0.02              # 7.20
STACK_ROW_H = (STACK_ROW_AREA_BOTTOM - STACK_ROW_AREA_TOP) / 3  # ~2.23
STACK_ROW_LABEL_H = 0.20
STACK_ROW_LABEL_FONT_PT = 13
STACK_ROW_IMG_H = STACK_ROW_H - STACK_ROW_LABEL_H            # ~2.03

# Per-token substitutions for prettify_metric_name (keys are lowercase).
TOKEN_MAP = {
    "mfi": "MFI",
    "mip": "MIP",
    "sig": "Signal",
    "num": "Number",
    "idx": "Index",
    "rel": "Rel.",
    "3slice": "(3-Slice MIP)",
    "bottom": "Synapse",
    # MT_CatB additions
    "cathepsinb": "CatB",
    "3mip": "(3-Slice MIP)",
    "um": "μm",
    "um2": "μm²",
    "cent": "Centrosome",
    "frac": "Fraction",
    "fdd": "FDD",
    "3d": "3D",
    "rms": "RMS",
    "zcof": "Z-COF",
    "rad": "Radial",
    "auc1": "AUC₁",
    "plat": "pLAT",
    "pzap70": "pZap70",
    "pmlc": "pMLC",
    "pkc": "PKC",
    "ppkc": "pPKC",
    "catb": "CatB",
}

# Tokens dropped entirely from the title.
SKIP_TOKENS = {"slice"}

# Per-filename overrides — applied before any other title logic.
TITLE_OVERRIDES = {
    "actin_bottom_mask_area_grid.png": "Synapse Area",
    "actin_bottom_3slice_mask_area_grid.png": "Synapse (3-Slice MIP) Area",
    "actin_bottom_slice_rad_profile_auc1_all_cells_with_average_grid.png":
        "Actin Synapse Radial Profile AUC₁ (all cells + avg)",
    "actin_bottom_3slice_mip_rad_profile_auc1_all_cells_with_average_grid.png":
        "Actin Synapse (3-Slice MIP) Radial Profile AUC₁ (all cells + avg)",
    "actin_bottom_3slice_mip_inner_outer_ratio_50pct_ylim_0_3_grid.png":
        "Actin Synapse (3-Slice MIP) Inner/Outer Ratio (50% Eff Rad Thresh, y-limit 0–3)",
    "actin_bottom_3slice_mip_inner_outer_ratio_70pct_ylim_0_3_grid.png":
        "Actin Synapse (3-Slice MIP) Inner/Outer Ratio (70% Eff Rad Thresh, y-limit 0–3)",
    # --- pLAT TIRF deck overrides ---------------------------------------
    "pLAT_synapse_r_eff_grid.png": "pLAT Synapse Effective Radius",
    "pLAT_synapse_rad_profile_auc1_all_cells_with_average_grid.png":
        "pLAT Synapse Radial Profile AUC₁ (all cells + avg)",
    "pLAT_synapse_g_ave_grid.png":
        "pLAT Synapse g(r) Average (scalar)",
    "pLAT_synapse_autocorr_c_all_cells_with_average_grid.png":
        "pLAT Synapse Autocorrelation g(r) (all cells + avg)",
    "foci_pLAT_cross_corr_profile_all_cells_with_average_grid.png":
        "Foci × pLAT Cross-Correlation Profile (all cells + avg)",
    "foci_pLAT_dist_bin_mfi_norm_grid.png":
        "pLAT MFI by Distance Bin to Actin Foci (norm. to synapse mean)",
    # --- pZap70 TIRF deck overrides -------------------------------------
    "pZap70_synapse_r_eff_grid.png": "pZap70 Synapse Effective Radius",
    "pZap70_synapse_rad_profile_auc1_all_cells_with_average_grid.png":
        "pZap70 Synapse Radial Profile AUC₁ (all cells + avg)",
    # --- pMLC joint deck overrides --------------------------------------
    "pMLC_synapse_r_eff_grid.png": "pMLC Synapse Effective Radius",
    "pMLC_synapse_rad_profile_auc1_all_cells_with_average_grid.png":
        "pMLC Synapse Radial Profile AUC₁ (all cells + avg)",
    "pMLC_synapse_g_ave_grid.png":
        "pMLC Synapse g(r) Average (scalar)",
    "pMLC_synapse_autocorr_c_all_cells_with_average_grid.png":
        "pMLC Synapse Autocorrelation g(r) (all cells + avg)",
    "pMLC_z50_rel_cell_bottom_grid.png": "pMLC Z₅₀ (rel. cell bottom)",
    "pMLC_z75_rel_cell_bottom_grid.png": "pMLC Z₇₅ (rel. cell bottom)",
    "pMLC_z90_rel_cell_bottom_grid.png": "pMLC Z₉₀ (rel. cell bottom)",
    # --- PKC deck overrides ---------------------------------------------
    "PKC_synapse_r_eff_grid.png": "PKC Synapse Effective Radius",
    "PKC_synapse_rad_profile_auc1_all_cells_with_average_grid.png":
        "PKC Synapse Radial Profile AUC₁ (all cells + avg)",
    "PKC_synapse_g_ave_grid.png":
        "PKC Synapse g(r) Average (scalar)",
    "PKC_synapse_autocorr_c_all_cells_with_average_grid.png":
        "PKC Synapse Autocorrelation g(r) (all cells + avg)",
    "PKC_z50_rel_cell_bottom_grid.png": "PKC Z₅₀ (rel. cell bottom)",
    "PKC_z75_rel_cell_bottom_grid.png": "PKC Z₇₅ (rel. cell bottom)",
    "PKC_z90_rel_cell_bottom_grid.png": "PKC Z₉₀ (rel. cell bottom)",
    # --- pPKC (phospho-PKC θ) deck overrides ----------------------------
    "pPKC_synapse_r_eff_grid.png": "pPKC Synapse Effective Radius",
    "pPKC_synapse_rad_profile_auc1_all_cells_with_average_grid.png":
        "pPKC Synapse Radial Profile AUC₁ (all cells + avg)",
    "pPKC_synapse_g_ave_grid.png":
        "pPKC Synapse g(r) Average (scalar)",
    "pPKC_synapse_autocorr_c_all_cells_with_average_grid.png":
        "pPKC Synapse Autocorrelation g(r) (all cells + avg)",
    "pPKC_z50_rel_cell_bottom_grid.png": "pPKC Z₅₀ (rel. cell bottom)",
    "pPKC_z75_rel_cell_bottom_grid.png": "pPKC Z₇₅ (rel. cell bottom)",
    "pPKC_z90_rel_cell_bottom_grid.png": "pPKC Z₉₀ (rel. cell bottom)",
    # --- CTSB (CatB) deck overrides -------------------------------------
    "CathepsinB_synapse_rad_profile_all_cells_with_average_grid.png":
        "CatB Synapse Radial Profile (all cells + avg)",
    "CathepsinB_synapse_autocorr_c_all_cells_with_average_grid.png":
        "CatB Synapse Autocorrelation g(r) (all cells + avg)",
    # (Z50/Z75/Z90 rel_cell_bottom + norm_cell_height overrides live in
    # the MT_CatB block below — consolidated to avoid dict-shadow bugs.)
    # Non-auc1 actin rad_profile variants (for CTSB compile).
    "actin_bottom_slice_rad_profile_all_cells_with_average_grid.png":
        "Actin Synapse Radial Profile (all cells + avg)",
    "actin_bottom_3slice_mip_rad_profile_all_cells_with_average_grid.png":
        "Actin Synapse (3-Slice MIP) Radial Profile (all cells + avg)",
    "pZap70_synapse_g_ave_grid.png":
        "pZap70 Synapse g(r) Average (scalar)",
    "pZap70_synapse_autocorr_c_all_cells_with_average_grid.png":
        "pZap70 Synapse Autocorrelation g(r) (all cells + avg)",
    "foci_pZap70_cross_corr_profile_all_cells_with_average_grid.png":
        "Foci × pZap70 Cross-Correlation Profile (all cells + avg)",
    "foci_pZap70_dist_bin_mfi_norm_grid.png":
        "pZap70 MFI by Distance Bin to Actin Foci (norm. to synapse mean)",
    "foci_CatB_dist_bin_mfi_norm_grid.png":
        "CatB MFI by Distance Bin to Actin Foci (norm. to synapse mean)",
    # --- MT_CatB compile overrides ----------------------------------------
    "centrosome_center_z_cell_bottom_distance_grid.png":
        "Centrosome-Synapse Distance",
    "centrosome_center_z_cell_bottom_distance_polarization_grid.png":
        "Centrosome-Synapse Distance: Fraction Polarized",
    "CathepsinB_z50_rel_cell_bottom_grid.png":
        "CatB Z₅₀ (μm from cell bottom)",
    "CathepsinB_z75_rel_cell_bottom_grid.png":
        "CatB Z₇₅ (μm from cell bottom)",
    "CathepsinB_z90_rel_cell_bottom_grid.png":
        "CatB Z₉₀ (μm from cell bottom)",
    "CathepsinB_z50_norm_cell_height_grid.png":
        "CatB Z₅₀ (normalized by cell height)",
    "CathepsinB_z75_norm_cell_height_grid.png":
        "CatB Z₇₅ (normalized by cell height)",
    "CathepsinB_z90_norm_cell_height_grid.png":
        "CatB Z₉₀ (normalized by cell height)",
    "CathepsinB_z50_rel_cell_bottom_polarization_grid.png":
        "CatB Z₅₀: Fraction Polarized",
    "CathepsinB_z75_rel_cell_bottom_polarization_grid.png":
        "CatB Z₇₅: Fraction Polarized",
    "CathepsinB_granule_area_um2_grid.png": "CatB Synapse Area",
    "CathepsinB_synapse_g_ave_grid.png": "CatB Synapse g(r) Average",
    "CathepsinB_synapse_g_ave_3mip_grid.png":
        "CatB Synapse g(r) Average (3-Slice MIP)",
    "CathepsinB_synapse_r_eff_grid.png": "CatB Synapse Effective Radius",
    "CathepsinB_synapse_r_eff_3mip_grid.png":
        "CatB Synapse Effective Radius (3-Slice MIP)",
    "CathepsinB_synapse_autocorr_rmax_um_grid.png":
        "CatB Synapse Autocorrelation: r_max (μm)",
    "CathepsinB_synapse_autocorr_smoothSig_grid.png":
        "CatB Synapse Autocorrelation: smoothed σ",
    "CathepsinB_zCOF_grid.png": "CatB Z-Center of Fluorescence",
    "CathepsinB_zCOF_actin_scale_grid.png":
        "CatB Z-COF (actin scale)",
    "CathepsinB_zCOF_cell_bottom_distance_grid.png":
        "CatB Z-COF: Distance from Cell Bottom",
    "CathepsinB_zCOF_cell_bottom_distance_norm_cell_height_grid.png":
        "CatB Z-COF: Distance from Cell Bottom (normalized by cell height)",
    # --- Actin foci (canonical + scoring/threshold extras) --------------
    "actin_foci_count_grid.png": "Actin Foci Count",
    "actin_foci_area_um2_grid.png": "Actin Foci Area (μm²)",
    "actin_foci_area_fraction_grid.png":
        "Actin Foci Area Fraction (of synapse)",
    "actin_foci_mean_intensity_grid.png": "Actin Foci Mean Intensity",
    "actin_foci_total_intensity_grid.png": "Actin Foci Total Intensity",
    "actin_foci_max_norm_score_grid.png":
        "Actin Foci Max Normalized Score",
    "actin_foci_mean_norm_score_grid.png":
        "Actin Foci Mean Normalized Score",
    "actin_foci_threshold_used_grid.png":
        "Actin Foci Detection Threshold Used",
    # --- Actin whole-cell (MIP footprint + integrated intensity) --------
    "actin_MIP_area_grid.png": "Actin Whole-Cell MIP Area",
    "actin_MIP_major_axis_length_grid.png":
        "Actin Whole-Cell MIP Major Axis Length",
    "actin_MIP_minor_axis_length_grid.png":
        "Actin Whole-Cell MIP Minor Axis Length",
    "actin_MIP_mask_nonzero_grid.png":
        "Actin Whole-Cell MIP Nonzero Pixel Count",
    "actin_total_sig_grid.png": "Actin Whole-Cell Total Signal",
    "actin_peak_sig_grid.png": "Actin Whole-Cell Peak Signal",
    # --- Actin Z-distribution (signal position + cell extent) -----------
    "actin_z50_rel_cell_bottom_grid.png":
        "Actin Z₅₀ (μm from cell bottom)",
    "actin_z75_rel_cell_bottom_grid.png":
        "Actin Z₇₅ (μm from cell bottom)",
    "actin_z90_rel_cell_bottom_grid.png":
        "Actin Z₉₀ (μm from cell bottom)",
    "actin_zprofile_peak_idx_grid.png":
        "Actin Z-Profile Peak Slice Index",
    "actin_zprofile_peak_rel_bot_grid.png":
        "Actin Z-Profile Peak (μm from cell bottom)",
    "actin_bottom_slice_num_grid.png": "Actin Cell Bottom Slice #",
    "actin_top_slice_num_grid.png": "Actin Cell Top Slice #",
    "actin_height_grid.png": "Cell Height (slices)",
    "actin_above_FOV_grid.png": "Cell Extends Above FOV (flag)",
    "actin_below_FOV_grid.png": "Cell Extends Below FOV (flag)",
    # --- IRM synapse area + actin median radius ------------------------
    "IRM_mask_area_grid.png": "Synapse Area (IRM)",
    "actin_bottom_slice_r_median_grid.png": "Actin Synapse R_median",
    "actin_bottom_3slice_mip_r_median_grid.png":
        "Actin Synapse (3-Slice MIP) R_median",
    # --- Axial-profile overlays (per-experiment) -----------------------
    "axial_overlay_perslice_Jul_17_2026.png":
        "Actin Axial Profile Overlay — per slice (Jul 17, 2026)",
    "axial_overlay_perslice_Jul_23_2026.png":
        "Actin Axial Profile Overlay — per slice (Jul 23, 2026)",
}

# Optional per-filename subtitle shown under the title (formula / units /
# interpretation). Missing entries -> no subtitle line. Currently empty.
SUBTITLE_BY_FILENAME = {}

# ---------------------------------------------------------------------------


def prettify_metric_name(filename: str) -> str:
    """Convert a PNG filename to a readable slide title.

    Checks TITLE_OVERRIDES first. Otherwise strips `.png` and `_grid`,
    splits on `_`, applies per-token substitutions (dropping tokens in
    SKIP_TOKENS and collapsing the redundant `3slice_mip` pair into a
    single `(3-Slice MIP)`), and rewrites `Inner Outer` as `Inner/Outer`.
    """
    if filename in TITLE_OVERRIDES:
        return TITLE_OVERRIDES[filename]

    stem = filename
    if stem.lower().endswith(".png"):
        stem = stem[:-4]
    if stem.lower().endswith("_grid"):
        stem = stem[:-5]

    tokens = stem.split("_")
    out = []
    i = 0
    while i < len(tokens):
        tok = tokens[i]
        low = tok.lower()
        nxt = tokens[i + 1].lower() if i + 1 < len(tokens) else None

        if low == "3slice" and nxt == "mip":
            out.append("(3-Slice MIP)")
            i += 2
            continue
        if low in SKIP_TOKENS:
            i += 1
            continue
        if low in TOKEN_MAP:
            out.append(TOKEN_MAP[low])
        elif low.endswith("pct") and low[:-3].isdigit():
            out.append(f"({low[:-3]}% Eff Rad Thresh)")
        elif low.endswith("um2") and low[:-3].isdigit():
            out.append(f"{low[:-3]} μm²")
        elif low.endswith("um") and low[:-2].isdigit():
            out.append(f"{low[:-2]} μm")
        else:
            out.append(tok.capitalize())
        i += 1

    pretty = " ".join(out)
    pretty = pretty.replace("Inner Outer", "Inner/Outer")
    return pretty


def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False):
    """Add a centered textbox with given text and styling."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    """Place an image inside the given bounding box, preserving aspect
    ratio and centering on the dimension that is < box."""
    pic = slide.shapes.add_picture(
        image_path,
        Inches(box_left),
        Inches(box_top),
        width=Inches(box_w),
    )
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path,
            Inches(box_left),
            Inches(box_top),
            height=Inches(box_h),
        )
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def build_slide(prs, title_text: str, image_path: Path, footer_text: str,
                subtitle_text=None):
    """Build one slide: title + optional subtitle + full image + filepath footer.
    Returns (slide, missing_flag)."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, WHITE)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=BLACK, bold=True,
    )

    if subtitle_text:
        add_textbox(
            slide, subtitle_text,
            SUBTITLE_LEFT, SUBTITLE_TOP, SUBTITLE_WIDTH, SUBTITLE_HEIGHT,
            font_pt=SUBTITLE_FONT_PT, color=BLACK, italic=True,
        )
        img_top, img_h = IMG_TOP_WITH_SUBTITLE, IMG_BOX_H_WITH_SUBTITLE
    else:
        img_top, img_h = IMG_TOP, IMG_BOX_H

    missing = not image_path.exists()
    if not missing:
        add_image_in_box(
            slide, str(image_path), IMG_LEFT, img_top, IMG_BOX_W, img_h
        )
    else:
        add_textbox(
            slide, "(missing)",
            IMG_LEFT, img_top + img_h / 2 - 0.2, IMG_BOX_W, 0.4,
            font_pt=18, color=BLACK,
        )

    add_textbox(
        slide, footer_text,
        FOOTER_LEFT, FOOTER_TOP, FOOTER_WIDTH, FOOTER_HEIGHT,
        font_pt=FOOTER_FONT_PT, color=BLACK, bold=False,
    )
    return slide, missing


def build_stacked_slide(prs, title_text: str,
                        rows: list, footer_text: str):
    """Build one slide with N timepoint images stacked vertically.
    `rows` is a list of (label, image_path) tuples (typically 3).
    Returns (slide, list_of_missing_paths)."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, WHITE)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, STACK_TITLE_HEIGHT,
        font_pt=STACK_TITLE_FONT_PT, color=BLACK, bold=True,
    )

    missing_paths = []
    row_w = SLIDE_W - 2 * MARGIN
    for i, (label, image_path) in enumerate(rows):
        row_top = STACK_ROW_AREA_TOP + i * STACK_ROW_H
        add_textbox(
            slide, label,
            MARGIN, row_top, row_w, STACK_ROW_LABEL_H,
            font_pt=STACK_ROW_LABEL_FONT_PT, color=BLACK, bold=True,
        )
        img_top = row_top + STACK_ROW_LABEL_H
        if image_path.exists():
            add_image_in_box(slide, str(image_path),
                             MARGIN, img_top, row_w, STACK_ROW_IMG_H)
        else:
            add_textbox(
                slide, "(missing)",
                MARGIN, img_top + STACK_ROW_IMG_H / 2 - 0.15,
                row_w, 0.3, font_pt=14, color=BLACK,
            )
            missing_paths.append(image_path)

    add_textbox(
        slide, footer_text,
        FOOTER_LEFT, STACK_FOOTER_TOP, FOOTER_WIDTH, STACK_FOOTER_HEIGHT,
        font_pt=STACK_FOOTER_FONT_PT, color=BLACK, bold=False,
    )
    return slide, missing_paths


def build_deck(dataset) -> int:
    """Build one deck for the given dataset entry. Returns missing-file count."""
    label = dataset["label"]
    output_path = dataset["output_path"]
    png_files = dataset["png_files"]
    strip_grid = dataset.get("strip_grid_suffix", False)
    stack_timepoints = dataset.get("stack_timepoints", False)
    # Either a list of (timepoint_label, root) pairs (combined deck),
    # or fall back to the single compiled_root field. The single-root
    # case is modeled as a 1-element list with empty timepoint label.
    timepoints = dataset.get("timepoints") or [("", dataset["compiled_root"])]

    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    print(f"\n=== Dataset: {label} ===")
    if len(timepoints) == 1 and not timepoints[0][0]:
        print(f"Source root: {timepoints[0][1]}")
    else:
        print(f"Timepoints: {', '.join(tp for tp, _ in timepoints)}{' (stacked)' if stack_timepoints else ''}")
    print(f"Writing deck to: {output_path}\n")

    missing = []
    slides_written = 0
    for entry in png_files:
        # Entries are either (sub, png_name) or (sub, png_name, title_suffix).
        # title_suffix is appended to the prettified title (useful when the
        # same filename appears in multiple subfolders, e.g. per-date variants).
        sub, png_name = entry[0], entry[1]
        title_suffix = entry[2] if len(entry) >= 3 else None
        if stack_timepoints:
            # One slide per metric, all timepoints stacked vertically.
            disk_name = png_name.replace("_grid.png", ".png") if strip_grid else png_name
            rows = [
                (tp_label, tp_root / sub / disk_name)
                for tp_label, tp_root in timepoints
            ]
            title = prettify_metric_name(png_name)
            if title_suffix:
                title = f"{title} {title_suffix}"
            # Footer uses {tp1,tp2,...} brace notation to denote the timepoint dirs.
            tp_brace = "{" + ",".join(
                str(tp_root).split("\\all_")[-1].split("/all_")[-1].split("\\")[0].split("/")[0]
                for _, tp_root in timepoints
            ) + "}"
            first_path_parts = rows[0][1].as_posix()
            # Replace the first timepoint dir name with the brace expansion.
            footer = first_path_parts.replace(
                f"/all_{timepoints[0][0].replace(' ', '')}/", f"/all_{tp_brace}/"
            )
            _, missing_paths = build_stacked_slide(prs, title, rows, footer)
            slides_written += 1
            status = "OK" if not missing_paths else f"MISSING ({len(missing_paths)}/{len(rows)})"
            print(f"[{png_name}]  {status}  -> {title!r}")
            for p in missing_paths:
                missing.append(f"{png_name}  ({p})")
        else:
            for tp_label, tp_root in timepoints:
                disk_name = png_name.replace("_grid.png", ".png") if strip_grid else png_name
                image_path = tp_root / sub / disk_name
                base_title = prettify_metric_name(png_name)
                if title_suffix:
                    base_title = f"{base_title} {title_suffix}"
                title = f"{base_title} ({tp_label})" if tp_label else base_title
                footer = image_path.as_posix()
                subtitle = SUBTITLE_BY_FILENAME.get(png_name)
                _, is_missing = build_slide(
                    prs, title, image_path, footer, subtitle_text=subtitle
                )
                slides_written += 1
                status = "OK" if not is_missing else "MISSING"
                tag = f"{png_name} @ {tp_label}" if tp_label else png_name
                print(f"[{tag}]  {status}  -> {title!r}")
                if is_missing:
                    missing.append(f"{png_name}  ({image_path})")

    # Snapshot the previous deck (if any) into backups/ before overwriting.
    if output_path.exists():
        backup_dir = output_path.parent / "backups"
        created = backup_presentation(str(output_path), backup_base=str(backup_dir))
        if created:
            print(f"\nBacked up previous deck to: {backup_dir}")
            for cat, path in created.items():
                print(f"  [{cat:6}] {path}")

    prs.save(str(output_path))
    print(f"\nDone. {slides_written} slides written to:\n  {output_path}")

    if missing:
        print(f"\nMissing ({len(missing)}):")
        for m in missing:
            print(f"  - {m}")
    else:
        print("\nAll images found - no missing items.")

    return len(missing)


def main() -> None:
    # Optional CLI filters: any args are treated as case-insensitive
    # substrings matched against dataset["label"]. No args = build all.
    filters = [a.lower() for a in sys.argv[1:]]
    if filters:
        selected = [d for d in DATASETS
                    if any(f in d["label"].lower() for f in filters)]
        if not selected:
            print(f"No datasets match filters: {filters}")
            print("Available labels:")
            for d in DATASETS:
                print(f"  - {d['label']}")
            sys.exit(1)
    else:
        selected = DATASETS

    total_missing = 0
    for dataset in selected:
        total_missing += build_deck(dataset)
    if total_missing:
        sys.exit(1)


if __name__ == "__main__":
    main()
