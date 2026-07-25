"""
insert_ctl_lamp1_master_interleaved_slides.py

Master interleaved deck combining the 20260617 (2-timepoint) and 20260716
(3-timepoint) fixed CTL LAMP1 experiments into one PowerPoint. For each block
(broadest slice, XZ MIPs, synapse mask, ...), the 20260617 slide comes first,
followed immediately by the 20260716 slide. Total: 22 blocks × 2 experiments
= 44 slides.

Scale groups are pinned ACROSS both experiments so a 5 µm scalebar renders at
the same cm-on-page whether the slide is from 20260617 or 20260716. Each
slide's layout adapts to its experiment's # of conditions:
  - 20260617 (2 conds: 3 min, 12 min)  → 2-column side-by-side (or 2-row stacked for panel groups)
  - 20260716 (3 conds: 3 min, 5 min, 12 min) → 3-column side-by-side (or 3-row stacked for panel groups)

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_ctl_lamp1_master_interleaved_slides.py
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/CTL_Glass_Nucleus_Centrosome/CTL_fixed_LAMP1/"
    "CTL_fixed_LAMP1_master_interleaved_20260617_20260716.pptx"
)

EXPERIMENTS = [
    {
        "tag": "20260617",
        "root": (
            "L:/FF/Nucleus_granules/CTL_fixed/"
            "20260617_Fixed_CTLs_glass_centrosome_polarization_granules_nucleus_3min_12min"
        ),
        "conditions": [
            ("C1_3min_aCD3_ICAM1_3SI_660bTub_535Actin_488LAMP1_405Nuc",  "3 min"),
            ("C2_12min_aCD3_ICAM1_3SI_660bTub_535Actin_488LAMP1_405Nuc", "12 min"),
        ],
    },
    {
        "tag": "20260716",
        "root": (
            "L:/FF/Nucleus_granules/CTL_fixed/"
            "20260716_Fixed_CTLs_glass_centrosome_polarization_granules_nucleus"
        ),
        "conditions": [
            ("C3_3min_aCD3_ICAM1_3SI_660bTub_535Actin_488LAMP1_405Nuc",  "3 min"),
            ("C2_5min_aCD3_ICAM1_3SI_660bTub_535Actin_488LAMP1_405Nuc",  "5 min"),
            ("C1_12min_aCD3_ICAM1_3SI_660bTub_535Actin_488LAMP1_405Nuc", "12 min"),
        ],
    },
]

# Block list: shared between both experiments. Title template gets .format(tag=)
# per experiment. subpath uses {cond}.
BLOCKS = [
    ("{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_nuc_bz/montages",
        "LAMP1 + MT + Nuc, broadest slice ({tag})", "broad"),
    ("{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_nuc_com/montages",
        "LAMP1 + MT + Nuc, centrosome slice ({tag})", "xy_phys"),
    ("{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_com/montages",
        "LAMP1 + MT, centrosome slice ({tag})", "xy_phys"),
    ("{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_com_adaptive_merge/montages",
        "LAMP1 + MT, centrosome slice (adaptive) ({tag})", "xy_phys"),
    ("{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_com_adaptive_panels/montages",
        "LAMP1 + MT, centrosome slice — channels + merge (adaptive) ({tag})", "companel_phys"),
    ("{cond}/cropped/channels/prog_fixed_cells/actin/bottom_slice_seg/montages",
        "Actin synapse mask (bottom slice) ({tag})", "synapse"),
    ("{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/MT_nuc_xz/montages",
        "MT + Nuc XZ MIP ({tag})", "xz_phys"),
    ("{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_nuc_xz/montages",
        "LAMP1 + Nuc XZ MIP ({tag})", "xz_phys"),
    ("{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_nuc_xz/montages",
        "LAMP1 + MT + Nuc XZ MIP ({tag})", "xz_phys"),
    ("{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_xz/montages",
        "Actin XZ MIP ({tag})", "xz_phys"),
    ("{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_nuc_xz/montages",
        "Actin + Nuc XZ MIP ({tag})", "xz_phys"),
    ("{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_nuc_xz_planes/montages",
        "Actin + Nuc XZ MIP, cell top/bottom marked ({tag})", "xz_phys"),
    ("{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_xz_nolines/montages",
        "Actin XZ MIP (no lines) ({tag})", "xz_phys"),
    ("{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_nuc_xz_planes_nolines/montages",
        "Actin + Nuc XZ MIP, planes (no lines) ({tag})", "xz_phys"),
    ("{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_MT_xz_nolines/montages",
        "Actin + MT XZ MIP ({tag})", "xz_phys"),
    ("{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_xz_nolines/montages",
        "LAMP1 XZ MIP ({tag})", "xz_phys"),
    ("{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/MT_xz_nolines/montages",
        "MT XZ MIP ({tag})", "xz_phys"),
    ("{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_xz_panel_nolines/montages",
        "MT / LAMP1 / merge, XZ MIP panel (no lines) ({tag})", "xzpanel_phys"),
    ("{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/actin_MT_xz_panel_nolines/montages",
        "Actin / MT / merge, XZ MIP panel (no lines) ({tag})", "xzpanel_phys"),
    ("{cond}/cropped/channels/prog_fixed_cells/physical_scale_images/Lamp1_MT_syn/montages",
        "LAMP1 + MT, synapse plane ({tag})", "syn_phys"),
]

CHUNK_GLOB = "montage_cells_*.png"

EXTRA_BLOCKS = [
    ("{cond}/cropped/channels/prog_fixed_cells/Lamp1/deepest_invag_slice/merges/montages_deepest_invag",
        "montage_cells_*_with_MT.png",
        "LAMP1 + MT, deepest invag slice ({tag})", "invag_slice"),
    ("{cond}/cropped/channels/prog_fixed_cells/MT/deepest_invag_slice/merges/montages_deepest_invag",
        "montage_cells_*.png",
        "MT, deepest invag slice ({tag})", "invag_slice"),
]

# Scalebar invariant (5 µm = 104 px in every physical_scale_images/ montage).
SCALEBAR_PX = 104
SCALEBAR_UM = 5
PPUM_SOURCE = SCALEBAR_PX / SCALEBAR_UM

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 24

GRID_LEFT = 0.10
GRID_TOP = 0.60
GRID_BOT_MARGIN = 0.10
GRID_H = SLIDE_H - GRID_TOP - GRID_BOT_MARGIN
LABEL_H = 0.30
LABEL_FONT_PT = 16
COL_GAP = 0.10

PANEL_GROUPS = {"xzpanel_phys", "companel_phys"}
PANEL_IMG_W = SLIDE_W - 2 * GRID_LEFT

# For 3-condition slides, some scale_groups are wide-short (XZ MIPs) and get
# more area with a 2-on-top / 1-below layout (each panel ~6.5" wide × ~3.1"
# tall). Others are XY-tall and get more area from the classic 3-column
# layout (each panel ~4.3" wide × 6.5" tall).
TWO_TOP_ONE_BOT_GROUPS = {"xz_phys"}

# ---------------------------------------------------------------------------


def _winlong(p) -> str:
    s = os.path.abspath(str(p))
    if os.name == "nt" and not s.startswith("\\\\?\\"):
        s = "\\\\?\\" + s.replace("/", "\\")
    return s


def _exists_long(p) -> bool:
    return os.path.exists(_winlong(p))


def _png_dims(path: Path) -> Tuple[int, int]:
    with Image.open(_winlong(path)) as im:
        return im.size


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def compute_group_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_image_at_ppi(slide, image_path: Path, ppi: float,
                     area_left: float, area_top: float,
                     area_w: float, area_h: float):
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    left_in = area_left + (area_w - w_in) / 2
    top_in  = area_top  + (area_h - h_in) / 2
    return slide.shapes.add_picture(
        _winlong(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(w_in),
    )


def find_first_chunk(montages_dir: Path, pattern: str = CHUNK_GLOB) -> Optional[Path]:
    import fnmatch
    long_dir = _winlong(montages_dir)
    if not os.path.isdir(long_dir):
        return None
    try:
        names = os.listdir(long_dir)
    except OSError:
        return None
    names = [n for n in names if fnmatch.fnmatch(n, pattern)]
    if not names:
        return None
    names.sort(key=lambda n: int(re.match(r"montage_cells_(\d+)", n).group(1))
               if re.match(r"montage_cells_(\d+)", n) else 0)
    return montages_dir / names[0]


def side_by_side_cells(n_cond: int) -> Tuple[float, List[float], float]:
    """(cell_w, cell_lefts, img_h) for an N-column side-by-side layout."""
    cell_w = (SLIDE_W - 2 * GRID_LEFT - (n_cond - 1) * COL_GAP) / n_cond
    cell_lefts = [GRID_LEFT + i * (cell_w + COL_GAP) for i in range(n_cond)]
    img_h = GRID_H - LABEL_H
    return cell_w, cell_lefts, img_h


def stacked_rows(n_cond: int) -> Tuple[List[float], float]:
    """(row_tops, row_img_h) for an N-row full-width stacked layout."""
    row_h = GRID_H / n_cond
    row_img_h = row_h - LABEL_H
    row_tops = [GRID_TOP + i * row_h for i in range(n_cond)]
    return row_tops, row_img_h


def two_top_one_bottom_cells() -> Tuple[float, float, List[Tuple[float, float]]]:
    """For 3-condition slides: (cell_w, img_h, cell_positions[(left, top)]).
    Row 1: 2 half-slide-width cells (top).
    Row 2: 1 half-slide-width cell centered (bottom). Same cell size as row 1
    so all three panels share a single physical scale."""
    cell_w = (SLIDE_W - 2 * GRID_LEFT - COL_GAP) / 2
    row_h = GRID_H / 2
    img_h = row_h - LABEL_H
    row1_top = GRID_TOP
    row2_top = GRID_TOP + row_h
    positions = [
        (GRID_LEFT,                     row1_top),  # 3 min
        (GRID_LEFT + cell_w + COL_GAP,  row1_top),  # 5 min
        ((SLIDE_W - cell_w) / 2,        row2_top),  # 12 min centered
    ]
    return cell_w, img_h, positions


def build_multi_compare_slide(prs, title_text, panels, slide_ppi, layout="ncol"):
    """layout controls the arrangement:
      'ncol'   -> single-row N-column side-by-side (default; works for any N)
      '2t1b'   -> 2 panels on top row + 1 centered on bottom row (n_cond=3 only)
    """
    n_cond = len(panels)

    if layout == "2t1b" and n_cond == 3:
        cell_w, img_h, positions = two_top_one_bottom_cells()
    else:
        cw, cell_lefts, ih = side_by_side_cells(n_cond)
        cell_w, img_h = cw, ih
        positions = [(cl, GRID_TOP) for cl in cell_lefts]

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True)

    missing = []
    for (label, img_path), (cell_left, cell_top) in zip(panels, positions):
        add_textbox(slide, label,
            cell_left, cell_top, cell_w, LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True)
        if img_path is not None and _exists_long(img_path):
            add_image_at_ppi(slide, img_path, slide_ppi,
                             cell_left, cell_top + LABEL_H, cell_w, img_h)
        else:
            add_textbox(slide, "(missing)",
                cell_left, cell_top + LABEL_H + img_h / 2 - 0.15, cell_w, 0.3,
                font_pt=14, color=WHITE)
            missing.append(label)
    return slide, missing


def build_multi_stacked_slide(prs, title_text, panels, slide_ppi):
    """N-row stacked full-width. panels = list of (label, img_path) top→bottom."""
    n_cond = len(panels)
    row_tops, row_img_h = stacked_rows(n_cond)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True)

    missing = []
    for i, (label, img_path) in enumerate(panels):
        row_top = row_tops[i]
        add_textbox(slide, label,
            GRID_LEFT, row_top, PANEL_IMG_W, LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True)
        if img_path is not None and _exists_long(img_path):
            add_image_at_ppi(slide, img_path, slide_ppi,
                             GRID_LEFT, row_top + LABEL_H, PANEL_IMG_W, row_img_h)
        else:
            add_textbox(slide, "(missing)",
                GRID_LEFT, row_top + LABEL_H + row_img_h / 2 - 0.15,
                PANEL_IMG_W, 0.3, font_pt=14, color=WHITE)
            missing.append(label)
    return slide, missing


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    # Pre-pass: for each (block, experiment) build a slide spec.
    # Interleave order: block1-exp1, block1-exp2, block2-exp1, block2-exp2, ...
    all_blocks = [(sp, CHUNK_GLOB, t, sg) for (sp, t, sg) in BLOCKS] + list(EXTRA_BLOCKS)
    slide_specs: List[dict] = []
    for subpath_tmpl, glob_pat, title_tmpl, scale_group in all_blocks:
        for exp in EXPERIMENTS:
            root = Path(exp["root"])
            dirs = [root / Path(subpath_tmpl.format(cond=cond))
                    for cond, _ in exp["conditions"]]
            imgs = [find_first_chunk(d, glob_pat) for d in dirs]
            if all(img is None for img in imgs):
                print(f"WARNING: skipping '{title_tmpl.format(tag=exp['tag'])}' "
                      f"— no montages found in '{subpath_tmpl}' for exp {exp['tag']}.")
                continue
            slide_specs.append({
                "title": title_tmpl.format(tag=exp["tag"]),
                "dirs": dirs,
                "imgs": imgs,
                "scale_group": scale_group,
                "tag": exp["tag"],
                "tp_labels": [tp for _, tp in exp["conditions"]],
            })

    if not slide_specs:
        print("ERROR: no slides to render — every (block, experiment) was empty.")
        sys.exit(1)

    # One PPI per (scale_group, experiment_tag) — NOT unified across
    # experiments. This keeps each experiment's slides at the same size they'd
    # have in its solo deck (so a 20260617 slide in this master ≡ the same
    # slide in the 20260617 solo deck). Scalebars are consistent WITHIN an
    # experiment but may differ between experiments.
    def _use_2t1b(spec) -> bool:
        return (len(spec["imgs"]) == 3
                and spec["scale_group"] in TWO_TOP_ONE_BOT_GROUPS)

    def _cell_box(spec) -> Tuple[float, float]:
        n_cond = len(spec["imgs"])
        if spec["scale_group"] in PANEL_GROUPS:
            _, row_img_h = stacked_rows(n_cond)
            return PANEL_IMG_W, row_img_h
        if _use_2t1b(spec):
            cell_w, img_h, _ = two_top_one_bottom_cells()
            return cell_w, img_h
        cell_w, _, img_h = side_by_side_cells(n_cond)
        return cell_w, img_h

    group_ppi: dict = {}
    for spec in slide_specs:
        max_w, max_h = _cell_box(spec)
        imgs = [p for p in spec["imgs"] if p is not None and _exists_long(p)]
        own = compute_group_ppi(imgs, max_w, max_h) if imgs else 0.0
        key = (spec["scale_group"], spec["tag"])
        group_ppi[key] = max(group_ppi.get(key, 0.0), own)

    PHYS_GROUPS = {"xz_phys", "syn_phys", "broad", "xy_phys", "xzpanel_phys", "companel_phys"}
    print(f"Pinned PPI per (scale_group, experiment) ({len(slide_specs)} slides total):")
    for (sg, tag), ppi in sorted(group_ppi.items()):
        if sg in PHYS_GROUPS:
            bar = SCALEBAR_PX / ppi
            note = f"  scalebar = {bar:.3f} in = {bar * 2.54:.3f} cm"
        else:
            note = "  (layout-pin only — no embedded scalebar)"
        print(f"  {sg:>13s} @ {tag}  PPI={ppi:>7.2f}{note}")
    print(f"\nWriting deck to: {OUTPUT_PATH}\n")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    missing_total = []
    for spec in slide_specs:
        slide_ppi = group_ppi[(spec["scale_group"], spec["tag"])]
        panels = list(zip(spec["tp_labels"], spec["imgs"]))
        if spec["scale_group"] in PANEL_GROUPS:
            _, missing = build_multi_stacked_slide(prs, spec["title"], panels, slide_ppi)
        else:
            layout = "2t1b" if _use_2t1b(spec) else "ncol"
            _, missing = build_multi_compare_slide(prs, spec["title"], panels, slide_ppi, layout=layout)
        status = " ".join(
            f"{tp}:{'OK' if img else 'MISS'}" for tp, img in panels
        )
        print(f"[{spec['tag']} {spec['scale_group']:>13s}]  {status}  {spec['title']}")
        for label in missing:
            idx = spec["tp_labels"].index(label) if label in spec["tp_labels"] else 0
            missing_total.append(f"{spec['title']}/{label}  ({spec['dirs'][idx]})")

    prs.save(str(out_path))
    print(f"\nDone. {len(slide_specs)} slides written to:\n  {out_path}")
    if missing_total:
        print(f"\nMissing ({len(missing_total)}):")
        for m in missing_total:
            print(f"  - {m}")
    else:
        print("\nAll panels found.")


if __name__ == "__main__":
    main()
