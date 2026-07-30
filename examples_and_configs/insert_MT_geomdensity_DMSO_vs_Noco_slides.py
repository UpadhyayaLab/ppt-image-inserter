"""
insert_MT_geomdensity_DMSO_vs_Noco_slides.py

Build the microtubule (MT) geometry-density summary deck from the 20260716
compile of the Nocodazole 20240123 experiment.

Single experiment (Jan 23 2024), two conditions: DMSO vs Noco 1 μM. Every
panel is one DMSO-vs-Noco comparison (blue vs yellow), so there are no
by-experiment / matrix layouts — each metric is a single tile. Modeled on
insert_vim_geomdensity_12min_combined_slides.py (same aspect-aware layout
engine).

Uses the 20260718 compile (DMSO n=77, Noco n=43 after the Noco arm was
reprocessed — the earlier 20260716 compile had only n=9 Noco).

Sections:
  1. MT density profiles (perinuc 0.5 μm)
  2. DNA density profiles (one slide per geometry, 3 shells)
  3. MT enrichment & correlation
  4. Invag depth profiles
  5. Curvature profiles
  6. Morphology scatter & centrosome
  7. Scalar violins (grid_panels)
  8. Histograms & QC

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_MT_geomdensity_DMSO_vs_Noco_slides.py
    conda run -n PPT_editing python examples_and_configs/insert_MT_geomdensity_DMSO_vs_Noco_slides.py --list
"""

import math
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

try:
    from PIL import Image
except ImportError:  # pragma: no cover
    Image = None

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation, safe_path, path_exists  # noqa: E402

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(
    "M:/FF/FF_4TB_2_Backup_fullHD/Nucleus Project_2ndharddrive/"
    "prog_fixed_Noco_MT_20240123/results_compilation_MT_geomdensity_20260718"
)
SINGLES = ROOT / "geom_density" / "profiles" / "singles"
ENRICH = ROOT / "geom_density" / "enrichment"
NC = ROOT / "geom_density" / "near_cent"
MORPH = ROOT / "geom_density" / "morphology_scatter"
DEPTH = ROOT / "invag_depth_profiles"
CURV = ROOT / "curvature_profiles"
GRID = ROOT / "grid_panels"
HIST = ROOT / "histograms"
SIMP = ROOT / "geom_density" / "simpson"

OUTPUT_PATH = Path(
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/"
    "nuc_mesh_struct_outside_nuc/"
    "MT geom-density vs NE geometry, DMSO vs Noco (20240123).pptx"
)

DECK_TITLE = "Microtubule density vs nuclear-envelope geometry — DMSO vs Noco"
DECK_SUBTITLE = (
    "Fixed Jurkats · Nocodazole experiment (Jan 23 2024) · "
    "DMSO (n=77) vs Noco 1 μM (n=43) · MT = perinuc 0.5 μm · "
    "single experiment · compiled 2026-07-18"
)

# Date tokens vary by directory.
D_DBL = "Jan_23__2024"     # profiles/singles, depth, histograms
D_COMMA = "Jan_23,_2024"   # enrichment, near_cent, morphology, simpson
D_SGL = "Jan_23_2024"      # curvature_profiles

DNA_SHELLS = ["boundary", "inperinuc025", "inperinuc05"]
DNA_SHELL_DISP = {
    "boundary": "boundary",
    "inperinuc025": "perinuc 0.25 μm",
    "inperinuc05": "perinuc 0.5 μm",
}
GEOM_LABELS_SHORT = {
    "hulldist": "hull dist",
    "mincurv": "min curv",
    "meancurv": "mean curv",
}

# ---------------------------------------------------------------------------
# Colors / layout (16:9)
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0x55, 0x55, 0x55)
FIELD_COLOR = RGBColor(0x2E, 0x5A, 0x88)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.08
GAP = 0.10

TITLE_LEFT = MARGIN
TITLE_TOP = 0.04
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 27

IMG_TOP = 0.58
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = 6.60

FOOTER_TOP = 7.20
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.28
FOOTER_FONT_PT = 8

CAPTION_HEIGHT = 0.50

MAX_GRID_COLS = 5


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------
def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER, wrap=True):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    image_path = safe_path(image_path)
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.7, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=36, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.1, SLIDE_W - 2 * MARGIN, 1.1,
                font_pt=16, color=GREY, italic=True)


def build_divider_slide(prs, title, subtitle=""):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, title, MARGIN, 3.0, SLIDE_W - 2 * MARGIN, 1.2,
                font_pt=40, color=BLACK, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, MARGIN, 4.25, SLIDE_W - 2 * MARGIN, 0.9,
                    font_pt=18, color=GREY, italic=True)


def add_labeled_caption(slide, lines, left, top, width, height):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    for i, (text, pt, color) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.size = Pt(pt)
        run.font.color.rgb = color
    return box


def _tile_caption(slide, cap, path, left, top, width, height, cap_pt):
    lines = []
    if cap:
        lines.append((cap, min(cap_pt, 10), GREY))
    lines.append((Path(str(path)).stem, 8, FIELD_COLOR))
    add_labeled_caption(slide, lines, left, top, width, height)


_AR_CACHE = {}


def _img_ar(path):
    key = str(path)
    if key in _AR_CACHE:
        return _AR_CACHE[key]
    ar = 1.5
    if Image is not None:
        try:
            with Image.open(safe_path(key)) as im:
                w, h = im.size
                if h:
                    ar = w / float(h)
        except Exception:
            pass
    _AR_CACHE[key] = ar
    return ar


def _uniform_img_size(a, cols, rows, area_w, area_h, cap_h):
    cw = (area_w - (cols - 1) * GAP) / cols
    cell_h = (area_h - (rows - 1) * GAP) / rows
    ih_avail = cell_h - cap_h
    if ih_avail <= 0.2 or cw <= 0.2:
        return None
    iw = min(cw, ih_avail * a)
    return iw, iw / a


def _best_cols(ars, n, area_w, area_h, cap_h, cap_cols):
    uniform = (max(ars) / min(ars)) < 1.20
    a = sorted(ars)[len(ars) // 2]
    best = None
    for cols in range(1, min(cap_cols, n) + 1):
        rows = math.ceil(n / cols)
        if uniform:
            sz = _uniform_img_size(a, cols, rows, area_w, area_h, cap_h)
            if sz is None:
                continue
            score = sz[0] * sz[1]
        else:
            cw = (area_w - (cols - 1) * GAP) / cols
            cell_h = (area_h - (rows - 1) * GAP) / rows
            ih_avail = cell_h - cap_h
            if ih_avail <= 0.2 or cw <= 0.2:
                continue
            score = 0.0
            for ai in ars:
                iw = min(cw, ih_avail * ai)
                score += iw * (iw / ai)
        if best is None or score > best[0] + 1e-6:
            best = (score, cols, rows, uniform)
    if best is None:
        c = min(cap_cols, n)
        return c, math.ceil(n / c), uniform
    return best[1], best[2], best[3]


def build_grid_slide(prs, title, entries, footer_text, max_cols,
                     cap_pt=11, cap_h=0.28):
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    n = len(entries)
    area_top = IMG_TOP
    area_w = SLIDE_W - 2 * MARGIN
    area_h = FOOTER_TOP - IMG_TOP - 0.02

    ars = [_img_ar(p) for p, _ in entries]
    cap_cols = min(MAX_GRID_COLS, n)
    cols, rows, uniform = _best_cols(ars, n, area_w, area_h, cap_h, cap_cols)

    if uniform:
        a = sorted(ars)[len(ars) // 2]
        iw, ih = _uniform_img_size(a, cols, rows, area_w, area_h, cap_h)
        block_w = cols * iw + (cols - 1) * GAP
        block_h = rows * (ih + cap_h) + (rows - 1) * GAP
        x0 = MARGIN + (area_w - block_w) / 2
        y0 = area_top + max(0.0, (area_h - block_h) / 2)
        for i, (path, cap) in enumerate(entries):
            r, c = divmod(i, cols)
            left = x0 + c * (iw + GAP)
            top = y0 + r * (ih + cap_h + GAP)
            add_image_in_box(slide, str(path), left, top, iw, ih)
            _tile_caption(slide, cap, path, left, top + ih, iw, cap_h, cap_pt)
    else:
        cell_w = (area_w - (cols - 1) * GAP) / cols
        cell_h = (area_h - (rows - 1) * GAP) / rows
        img_h = cell_h - cap_h
        for i, (path, cap) in enumerate(entries):
            r, c = divmod(i, cols)
            left = MARGIN + c * (cell_w + GAP)
            top = area_top + r * (cell_h + GAP)
            add_image_in_box(slide, str(path), left, top, cell_w, img_h)
            _tile_caption(slide, cap, path, left, top + img_h,
                          cell_w, cap_h, cap_pt)

    add_textbox(slide, footer_text, MARGIN, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


# ---------------------------------------------------------------------------
# Content plan
# ---------------------------------------------------------------------------
def _tiles(pairs):
    kept, miss = [], []
    for path, cap in pairs:
        (kept if path_exists(path) else miss).append((path, cap))
    return kept, [p for p, _ in miss]


def build_plan():
    plan, missing = [], []

    def add_grid(title, pairs, footer, cap_pt=10, cap_h=0.28):
        kept, miss = _tiles(pairs)
        missing.extend(miss)
        if kept:
            plan.append(("grid", title, kept, footer, MAX_GRID_COLS,
                         cap_pt, cap_h))

    # ===================================================================
    # Sec 1: MT density profiles (perinuc 0.5 μm) — DMSO vs Noco overlay
    # ===================================================================
    plan.append(("divider",
                 "Microtubule density profiles",
                 "perinuc 0.5 μm shell · DMSO vs Noco 1 μM"))

    add_grid(
        "MT density profiles vs NE geometry",
        [(SINGLES / "MT_geomdens_{}_perinuc05_OVERLAY_line_{}.png".format(
            geom, D_DBL), GEOM_LABELS_SHORT[geom])
         for geom in ["hulldist", "mincurv", "meancurv"]],
        "DMSO (blue) vs Noco 1 μM (yellow) · perinuc 0.5 μm")

    # ===================================================================
    # Sec 2: DNA density profiles — one slide per geometry (3 shells)
    # ===================================================================
    plan.append(("divider",
                 "DNA density profiles",
                 "boundary / perinuc 0.25 μm / perinuc 0.5 μm · "
                 "DMSO vs Noco"))

    for geom in ["hulldist", "mincurv", "meancurv"]:
        add_grid(
            "DNA density vs {} — DMSO vs Noco".format(
                GEOM_LABELS_SHORT[geom]),
            [(SINGLES / "DNA_geomdens_{}_{}_OVERLAY_line_{}.png".format(
                geom, shell, D_DBL), DNA_SHELL_DISP[shell])
             for shell in DNA_SHELLS],
            "DMSO (blue) vs Noco 1 μM (yellow) · rows = shell")

    # ===================================================================
    # Sec 3: MT enrichment & correlation
    # ===================================================================
    plan.append(("divider",
                 "MT enrichment & correlation",
                 "NE + perinuc shells · DMSO vs Noco"))

    def en(stem):
        return ENRICH / "MT_{}_shells_{}.png".format(stem, D_COMMA)

    add_grid(
        "MT enrichment near invaginations & concave NE surfaces",
        [(en("hulldist_gt0.5um"), "dist > 0.5 μm"),
         (en("hulldist_gt1.0um"), "dist > 1.0 μm"),
         (en("mincurv_lt0"), "min curv < 0"),
         (en("mincurv_ltm0.25"), "min curv < −0.25"),
         (en("meancurv_lt0"), "mean curv < 0")],
        "enrichment = fraction ÷ expected · ref line 1 · DMSO vs Noco")

    add_grid(
        "MT per-cell correlation vs NE geometry",
        [(en("corr_hulldist"), "corr vs hull dist"),
         (en("corr_mincurv"), "corr vs min curv"),
         (en("corr_meancurv"), "corr vs mean curv"),
         (en("deepcorr_mincurv"), "deep corr vs min curv"),
         (en("deepcorr_meancurv"), "deep corr vs mean curv")],
        "incl. deep invaginations · ref line 0 · DMSO vs Noco")

    # ===================================================================
    # Sec 4: Invag depth profiles
    # ===================================================================
    plan.append(("divider",
                 "MT in invaginations — depth profiles",
                 "centrosome-stratified · DMSO vs Noco"))

    add_grid(
        "MT invagination depth profiles",
        [(DEPTH / "MT_invag_depth_profiles_{}.png".format(D_DBL),
          "unstratified"),
         (DEPTH / "MT_invag_depth_profiles_0_5um_cent_stratified_{}.png".format(
             D_DBL), "0.5 μm cent-stratified"),
         (DEPTH / "MT_invag_depth_profiles_1um_cent_stratified_{}.png".format(
             D_DBL), "1 μm cent-stratified"),
         (DEPTH / "MT_invag_enrichment_vs_centdist_{}.png".format(D_DBL),
          "enrichment vs centrosome dist")],
        "DMSO vs Noco 1 μM")

    # ===================================================================
    # Sec 5: Curvature profiles
    # ===================================================================
    plan.append(("divider",
                 "Curvature-stratified MT profiles",
                 "min & mean curvature · centrosome-stratified"))

    for curv, clbl in [("min", "min-curvature"), ("mean", "mean-curvature")]:
        add_grid(
            "MT {} profiles".format(clbl),
            [(CURV / "MT_{}_curv_profiles_{}.png".format(curv, D_SGL),
              "unstratified"),
             (CURV / "MT_{}_curv_profiles_0_5um_cent_stratified_{}.png".format(
                 curv, D_SGL), "0.5 μm cent-stratified"),
             (CURV / "MT_{}_curv_profiles_1um_cent_stratified_{}.png".format(
                 curv, D_SGL), "1 μm cent-stratified")],
            "per stratification · DMSO vs Noco")

    # ===================================================================
    # Sec 6: Morphology scatter & centrosome level
    # ===================================================================
    plan.append(("divider",
                 "Morphology scatter & centrosome metrics",
                 "DMSO vs Noco"))

    add_grid(
        "MT morphology scatter & centrosome level",
        [(MORPH / "MT_r_hulldist_perinuc05_bycent_vs_chull_max_D_by_cent_"
                  "{}.png".format(D_COMMA), "r(hulldist) vs max D"),
         (MORPH / "MT_hull_enrichment_2um_cent_vs_chull_max_D_by_cent_"
                  "{}.png".format(D_COMMA), "hull enrich 2 μm cent vs max D"),
         (MORPH / "MT_ratio_invag_away_1um_vs_chull_max_D_by_cent_"
                  "{}.png".format(D_COMMA), "ratio invag/away 1 μm vs max D"),
         (NC / "MT_near_cent_level_{}.png".format(D_COMMA),
          "level near centrosome")],
        "color = centrosome side · DMSO vs Noco")

    # ===================================================================
    # Sec 7: Scalar violins (grid_panels) — DMSO vs Noco
    # ===================================================================
    plan.append(("divider",
                 "MT scalars — DMSO vs Noco",
                 "grid panels · DMSO vs Noco 1 μM"))

    def gp(stem, cap):
        return (GRID / (stem + "_grid.png"), cap)

    add_grid(
        "Metric overview & invagination pockets",
        [(GRID / "top_differences_barplot.png", "top differences"),
         gp("MT_cyto_in_nuc_hull_vs_near_convex_nuc_MFI_ratio",
            "grooves ÷ convex surface"),
         gp("MT_cyto_in_nuc_hull_vs_all_perinuc_MFI_ratio",
            "grooves ÷ whole perinuc shell")],
        "grid_panels overview + voxel convex-hull decomposition · "
        "DMSO vs Noco")

    add_grid(
        "MT in the nuclear convex hull — supporting metrics",
        [gp("MT_cyto_in_nuc_hull_MFI", "grooves MFI (raw)"),
         gp("MT_cyto_in_nuc_hull_sig_fraction",
            "fraction of MT signal in grooves"),
         gp("MT_frac_in_nuc_convex_hull",
            "fraction of MT inside the hull"),
         gp("MT_invag_within_chull_vs_all_chull_MFI_ratio",
            "invag interior ÷ all within-hull"),
         gp("MT_invag_within_chull_vs_convex_within_chull_MFI_ratio",
            "invag interior ÷ convex rim")],
        "DMSO vs Noco")

    add_grid(
        "MT at the deepest invagination — ratio by side",
        [gp("MT_ratio_by_deepest_invag_all_0_5um", "all faces · 0.5 μm"),
         gp("MT_ratio_by_deepest_invag_all_1um", "all faces · 1 μm"),
         gp("MT_ratio_by_deepest_invag_away_0_5um", "away faces · 0.5 μm"),
         gp("MT_ratio_by_deepest_invag_away_1um", "away faces · 1 μm")],
        "MT at deepest invag ÷ reference · ref line 1 · DMSO vs Noco")

    add_grid(
        "MT above/below ratios by centrosome side",
        [gp("MT_ratio_above_below_0_5um", "above÷below 0.5 μm"),
         gp("MT_ratio_above_below_1um", "above÷below 1 μm"),
         gp("MT_ratio_above_by_side_0_5um", "above by side 0.5 μm"),
         gp("MT_ratio_above_by_side_1um", "above by side 1 μm"),
         gp("MT_ratio_below_by_side_0_5um", "below by side 0.5 μm"),
         gp("MT_ratio_below_by_side_1um", "below by side 1 μm")],
        "ref line 1 · DMSO vs Noco")

    add_grid(
        "MT around the centrosome — fraction & intensity",
        [gp("MT_frac_around_cent_1um", "frac around cent 1 μm"),
         gp("MT_frac_around_cent_2um", "frac around cent 2 μm"),
         gp("MT_frac_around_cent_3um", "frac around cent 3 μm"),
         gp("MT_MFI_around_cent_1um", "MFI around cent 1 μm"),
         gp("MT_MFI_around_cent_2um", "MFI around cent 2 μm"),
         gp("MT_MFI_around_cent_3um", "MFI around cent 3 μm")],
        "DMSO vs Noco")

    add_grid(
        "MT enrichment & perinuclear fraction near the centrosome",
        [gp("MT_enrichment_within_half_um_nuc_1_um_cent",
            "enrich 0.5 μm nuc, 1 μm cent"),
         gp("MT_enrichment_within_half_um_nuc_2_um_cent",
            "enrich 0.5 μm nuc, 2 μm cent"),
         gp("MT_frac_perinuc_within_1_um_cent",
            "perinuc frac within 1 μm cent"),
         gp("MT_frac_perinuc_within_2_um_cent",
            "perinuc frac within 2 μm cent")],
        "DMSO vs Noco")

    add_grid(
        "MT signal intensity & cent/away ratios",
        [gp("MT_total_sig", "total signal"),
         gp("MT_all_perinuc_MFI", "all perinuc MFI"),
         gp("MT_perinuc_sig_fraction", "perinuc signal fraction"),
         gp("MT_ratio_by_cent_away_0_5um", "cent÷away 0.5 μm"),
         gp("MT_ratio_by_cent_away_1um", "cent÷away 1 μm")],
        "DMSO vs Noco")

    # ===================================================================
    # Sec 8: Histograms & QC
    # ===================================================================
    plan.append(("divider",
                 "Histograms & QC",
                 "DMSO vs Noco"))

    add_grid(
        "MT histograms & Simpson's paradox check",
        [(HIST / "MT_cyto_near_vs_away_deepest_invag_ratio_histogram_"
                 "{}.png".format(D_DBL), "near vs away ratio"),
         (SIMP / "MT_simpson_check_{}.png".format(D_COMMA), "Simpson check")],
        "DMSO vs Noco 1 μM")

    # --- Drop orphan dividers ---
    cleaned = []
    for i, it in enumerate(plan):
        if it[0] == "divider":
            if i + 1 >= len(plan) or plan[i + 1][0] == "divider":
                continue
        cleaned.append(it)
    return cleaned, missing


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    list_only = "--list" in sys.argv
    plan, missing = build_plan()

    n_slides = 1 + len(plan)
    print("Output: {}".format(OUTPUT_PATH))
    print("{} content slides (+ title) = {} total\n".format(len(plan), n_slides))
    for it in plan:
        if it[0] == "divider":
            print("\n=== {} ===".format(it[1]))
        else:
            print("  [{}] {}".format(len(it[2]), it[1]))
    if missing:
        print("\nMISSING ({}):".format(len(missing)))
        for m in missing:
            print("  - {}".format(Path(m).name if hasattr(m, "name") else m))
    if list_only:
        return

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)

    for it in plan:
        if it[0] == "divider":
            build_divider_slide(prs, it[1], it[2])
        elif it[0] == "grid":
            _, title, entries, footer, max_cols, cap_pt, cap_h = it
            build_grid_slide(prs, title, entries, footer, max_cols, cap_pt, cap_h)

    if OUTPUT_PATH.exists():
        backup_dir = OUTPUT_PATH.parent / "backups"
        created = backup_presentation(str(OUTPUT_PATH),
                                      backup_base=str(backup_dir))
        if created:
            print("\nBacked up previous deck under: {}".format(backup_dir))

    prs.save(str(OUTPUT_PATH))
    total = len(prs.slides._sldIdLst)
    print("\nDone. {} slides written to:\n  {}".format(total, OUTPUT_PATH))
    if missing:
        print("Skipped {} missing panel(s).".format(len(missing)))


if __name__ == "__main__":
    main()
