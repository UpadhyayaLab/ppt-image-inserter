"""
insert_lata_dmso_vs_lata_phys_scale_slides.py

DMSO (left) vs Latrunculin A (right) physical-scale montage deck for the fixed
Jurkat LatA experiments. Copied from insert_bleb_dmso_vs_bleb_phys_scale_slides.py
(itself from the Noco deck) and repointed at the LatA datasets defined in
TCell-3D-Morphodynamics/config/datasets/fixed/Jurkats/LatA.

Same deck pattern as the noco/bleb decks: each slide is DMSO | LatA at one
condition, panels share a pinned pixels-per-inch so the embedded scalebar renders
at the same cm-on-page within the slide; long-path-safe enumeration; group-major
ordering; experiments whose montages aren't generated yet are skipped and fill in
automatically once processed.

KEY DIFFERENCE FROM NOCO/BLEB — no actin combos:
LatA depolymerizes F-actin, so actin was NOT analyzed in these experiments. The
physical_scale_images pipeline therefore emits only centrosome+nucleus combos
(cent_nuc*, nucleus*) — there are no actin_* combos. So this deck drops the two
actin XZ combos the bleb/noco decks led with and uses cent_nuc_xz for the XZ MIP.

Multiple experiments × concentrations × activation. LatA was tested at several
concentrations (50/100/250 nM) and, for the 2024 datasets, under both αCD3
activation and PLL (non-activated control). Each concentration is its own 2-column
DMSO-vs-LatA[conc] block (the DMSO control is the shared reference and appears in
each of its concentration comparisons). Blocks, chronological:
  - 03/2023     LatA_032023      (basal)  DMSO vs 100 / 250 nM
  - 11/07/2023  LatA_11072023    (basal)  DMSO vs 100 / 250 nM   [montages ready]
  - 05/09/2024  20240509 αCD3             DMSO vs 50 / 100 nM    [montages ready]
  - 05/09/2024  20240509 PLL              DMSO vs 50 / 100 nM    [montages ready]
  - 06/14/2024  20240614 αCD3             DMSO vs 50 nM
  - 06/14/2024  20240614 PLL              DMSO vs 50 nM
chan_sub differs: the 2023 (H:) datasets keep prog_fixed_cells under "combined";
the 2024 (J:) datasets keep it under "cells/channels". The 2024 montage paths run
~300 chars — past Windows MAX_PATH (260) — so the \\?\ long-path handling is load-
bearing here, not just defensive.

PPI is pinned DECK-WIDE per scale_group (xz, broad_1c), like the noco/bleb decks.
check_scalebar_pixel_widths.py confirms every LatA montage renders the 5 µm
scalebar at exactly 104 px (20.8 px/µm) across all 4 experiments, both activations,
and all combos, so one PPI per scale_group yields a uniform physical scale
(consistent scalebar-cm) across the whole deck. (Earlier drafts used
per_exp_scale=True as a hedge against differing µm/px; the measurement showed they
match, so it was dropped.)

Usage:
    python examples_and_configs/insert_lata_dmso_vs_lata_phys_scale_slides.py
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/LatA/"
    "LatA_Jurkats_DMSO_vs_LatA_phys_scale_montages.pptx"
)

# Dataset roots (one per acquisition; concentrations are condition folders within).
ROOT_MAR2023 = "H:/FF/Nucleus_Data/3D_Nucleus_Centrosome/Fixed/LatA_032023"
ROOT_NOV2023 = "H:/FF/Nucleus_Data/3D_Nucleus_Centrosome/Fixed/LatA_11072023"
ROOT_MAY2024 = (
    "J:/FF/fixed_cell/Vimentin/20240509_Jurkat-LatA/"
    "05092024_jurkats_LatA_50nMor100nM_1minPreIncubation_7minSpreading_"
    "antiVim647LP55_PhalloidinLP60_EGFPCent2LP55_hoeschstLP40"
)
ROOT_JUN2024 = (
    "J:/FF/fixed_cell/Vimentin/20240614_Jurkat-LatA/"
    "06142024_jurkats_LatA_50nM_1minPreIncubation_7minSpreading_"
    "antiVim647LP55_PhalloidinLP60_EGFPCent2LP55_hoeschstLP40"
)

# Mar-2023 condition folders share a long laser-line suffix; alias it once.
_M23 = "_405LP30_488LP40_561LP30_640LP40_100ms_"

# Each entry is one DMSO-vs-(one LatA condition) 2-column block.
# tag      — date tag woven into slide titles + used for chronological sort.
#            "MM/YYYY" (day unknown) or "MM/DD/YYYY"; see _exp_date_key.
# root     — dataset root dir.
# chan_sub — intermediate path inside each condition folder before
#            prog_fixed_cells ("combined" for 2023 H:, "cells/channels" for 2024 J:).
# left/right — (condition_folder, display_label) for the DMSO / LatA columns.
# tp_label — activation string ("αCD3", "PLL", or "" for the basal 2023 datasets).
EXPERIMENTS = [
    # --- 03/2023 (basal) ---------------------------------------------------
    {
        "tag": "03/2023", "root": ROOT_MAR2023, "chan_sub": "combined", "tp_label": "",
        "left":  ("Cropped_W1_DMSO" + _M23,     "DMSO"),
        "right": ("Cropped_W2_LatA100nM" + _M23, "LatA 100 nM"),
    },
    {
        "tag": "03/2023", "root": ROOT_MAR2023, "chan_sub": "combined", "tp_label": "",
        "left":  ("Cropped_W1_DMSO" + _M23,      "DMSO"),
        "right": ("Cropped_W3_LatA250nM" + _M23, "LatA 250 nM"),
    },
    # --- 11/07/2023 (basal) ------------------------------------------------
    {
        "tag": "11/07/2023", "root": ROOT_NOV2023, "chan_sub": "combined", "tp_label": "",
        "left":  ("f-Control-DMSO-1107",        "DMSO"),
        "right": ("f-LatrunculinA-100nM-1107",  "LatA 100 nM"),
    },
    {
        "tag": "11/07/2023", "root": ROOT_NOV2023, "chan_sub": "combined", "tp_label": "",
        "left":  ("f-Control-DMSO-1107",        "DMSO"),
        "right": ("f-LatrunculinA-250nM-1107",  "LatA 250 nM"),
    },
    # --- 05/09/2024 αCD3 ---------------------------------------------------
    {
        "tag": "05/09/2024", "root": ROOT_MAY2024, "chan_sub": "cells/channels", "tp_label": "αCD3",
        "left":  ("GcA1_CD3_DMSO_1to2000_1-20", "DMSO, αCD3"),
        "right": ("GcA2_CD3_50nM_1-20",         "LatA 50 nM, αCD3"),
    },
    {
        "tag": "05/09/2024", "root": ROOT_MAY2024, "chan_sub": "cells/channels", "tp_label": "αCD3",
        "left":  ("GcA1_CD3_DMSO_1to2000_1-20", "DMSO, αCD3"),
        "right": ("GcA3_CD3_100nM_1-20",        "LatA 100 nM, αCD3"),
    },
    # --- 05/09/2024 PLL ----------------------------------------------------
    {
        "tag": "05/09/2024", "root": ROOT_MAY2024, "chan_sub": "cells/channels", "tp_label": "PLL",
        "left":  ("GpB1_PLL_DMSO_1to2000_1-20", "DMSO, PLL"),
        "right": ("GpB2_PLL_50nM_1-20",         "LatA 50 nM, PLL"),
    },
    {
        "tag": "05/09/2024", "root": ROOT_MAY2024, "chan_sub": "cells/channels", "tp_label": "PLL",
        "left":  ("GpB1_PLL_DMSO_1to2000_1-20", "DMSO, PLL"),
        "right": ("GpB3_PLL_100nM_1-20",        "LatA 100 nM, PLL"),
    },
    # --- 06/14/2024 αCD3 / PLL (50 nM only) --------------------------------
    {
        "tag": "06/14/2024", "root": ROOT_JUN2024, "chan_sub": "cells/channels", "tp_label": "αCD3",
        "left":  ("CD3_DMSO_W1", "DMSO, αCD3"),
        "right": ("CD3_LatA_W1", "LatA 50 nM, αCD3"),
    },
    {
        "tag": "06/14/2024", "root": ROOT_JUN2024, "chan_sub": "cells/channels", "tp_label": "PLL",
        "left":  ("PLL wells/PLL_DMSO_W1", "DMSO, PLL"),
        "right": ("PLL wells/PLL_LatA_W1", "LatA 50 nM, PLL"),
    },
]

# (combo_subfolder, title_template, n_chunks, scale_group, fallback, opts).
# title_template gets " ({sub})" appended in main(), where sub is "tp, tag" (or
# just tag when tp is empty), so basal 2023 titles read "... (11/07/2023)" and
# 2024 titles read "... (αCD3, 05/09/2024)".
# Actin-free combo set (see module docstring): cent_nuc_xz replaces the bleb deck's
# two actin XZ combos; the broad/deepest slices are the cent/nucleus ones that the
# bleb deck also used. opts={} -> deck-wide PPI pinning per scale_group (like the
# noco deck): verified safe because every LatA montage renders the 5 µm scalebar at
# 104 px across all 4 experiments + combos, so one PPI per group == one physical
# scale deck-wide. Two groups (xz vs broad_1c) since their aspect ratios differ.
COMBOS = [
    ("cent_nuc_xz",      "Cent + Nuc XZ MIP",                       1, "xz",       None, {}),
    ("nucleus_bz",       "Nuc (DNA), broadest slice",               1, "broad_1c", None, {}),
    ("cent_nuc_bz",      "Cent + Nuc, broadest slice",              1, "broad_1c", None, {}),
    ("cent_nuc",         "Cent + Nuc, deepest invagination slice",  1, "broad_1c", None, {}),
    # Centrosome-plane combos (render at the centrosome's Z). cent_nuc_com is the
    # 2-channel nucleus+centrosome; vim_cent_nuc_com adds vimentin (orange) for the
    # 3-channel view. Only datasets reprocessed with the centrosome-plane pipeline
    # have these (11072023, 20240509); 032023 (no vimentin / not reprocessed) and
    # 20240614 (not processed yet) auto-skip. Same broad_1c physical scale.
    ("cent_nuc_com",     "Cent + Nuc, centrosome plane",            1, "broad_1c", None, {}),
    ("vim_cent_nuc_com", "Vim + Cent + Nuc, centrosome plane",      1, "broad_1c", None, {}),
]

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Slide layout (inches). 13.333 x 7.5 widescreen.
SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

# 1x2 cell grid below the title (label + image per cell)
GRID_LEFT = 0.10
GRID_TOP = 0.60
CELL_W = 6.50
CELL_H = SLIDE_H - GRID_TOP - 0.10   # 6.80"
LABEL_H = 0.30
IMG_H = CELL_H - LABEL_H             # 6.50"
LABEL_FONT_PT = 16
COL_GAP = SLIDE_W - 2 * GRID_LEFT - 2 * CELL_W

CELL_POSITIONS = [
    (GRID_LEFT,                    GRID_TOP),
    (GRID_LEFT + CELL_W + COL_GAP, GRID_TOP),
]

# Compact variant used by combos that opt into compact_layout=True (see COMBOS).
COMPACT_TITLE_HEIGHT = 0.35
COMPACT_TITLE_FONT_PT = 22
COMPACT_GRID_TOP = 0.40
COMPACT_LABEL_H = 0.20
COMPACT_LABEL_FONT_PT = 14
COMPACT_CELL_H = SLIDE_H - COMPACT_GRID_TOP - 0.10   # 7.00"
COMPACT_IMG_H = COMPACT_CELL_H - COMPACT_LABEL_H     # 6.80"

# Solo-layout cell: one condition per slide, image spans full slide width.
SOLO_CELL_W = SLIDE_W - 2 * GRID_LEFT                # 13.133"

# Scalebar invariant for the Jurkat nucleus fixed-cell physical-scale pipeline.
# Measured 104 px = 5 µm for EVERY LatA montage (all 4 experiments, both
# activations, all combos — verified with check_scalebar_pixel_widths.py), same as
# noco/CilioD. Because it's uniform, deck-wide PPI pinning gives a uniform physical
# scale. The constant itself is diagnostic-only (the printed scalebar-cm); the
# layout uses measured PNG dims + PPI pinning, so the deck is correct regardless.
SCALEBAR_PX = 104                            # px (measured across all LatA montages)
SCALEBAR_UM = 5                              # μm
PPUM_SOURCE = SCALEBAR_PX / SCALEBAR_UM      # 20.8 px/μm in the rendered PNG

# ---------------------------------------------------------------------------


def _winlong(p) -> str:
    """Return a Win32-safe absolute path string. On Windows, prepends the
    \\\\?\\ extended-length prefix so paths near MAX_PATH (260) still work in
    os.stat / Image.open / python-pptx add_picture. The 2024 LatA montage paths
    run ~300 chars, so this is load-bearing, not just defensive."""
    s = os.path.abspath(str(p))
    if os.name == "nt" and not s.startswith("\\\\?\\"):
        s = "\\\\?\\" + s.replace("/", "\\")
    return s


def _exists_long(p) -> bool:
    """MAX_PATH-safe existence check (pathlib.Path.exists() trips on long paths)."""
    return os.path.exists(_winlong(p))


def montage_dir(root, cond_folder: str, chan_sub: str, combo: str) -> Path:
    """Build the montages dir for one (experiment, condition, combo).

    Each condition folder has its own intermediate `chan_sub` before
    `prog_fixed_cells` ("combined" for the 2023 H: datasets, "cells/channels" for
    the 2024 J: datasets), so the path is built per experiment."""
    return (Path(root) / cond_folder / chan_sub /
            "prog_fixed_cells" / "physical_scale_images" / combo / "montages")


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _png_dims(path: Path) -> Tuple[int, int]:
    """Return (width_px, height_px) of a PNG without fully decoding it."""
    with Image.open(_winlong(path)) as im:
        return im.size


def compute_slide_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    """Smallest ppi such that every image fits in (max_w_in x max_h_in)."""
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_image_at_ppi(slide, image_path: Path, ppi: float,
                     area_left: float, area_top: float,
                     area_w: float, area_h: float):
    """Center an image at uniform PPI inside an arbitrary (left, top, w, h)
    image area. native_px / ppi gives inches; both dims pinned by ppi so all
    images in the slide render at one physical scale."""
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    left_in = area_left + (area_w - w_in) / 2
    top_in  = area_top  + (area_h - h_in) / 2
    return slide.shapes.add_picture(
        _winlong(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(w_in),
    )


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _chunk_start_index(p: Path) -> int:
    m = re.match(r"montage_cells_(\d+)", p.name)
    return int(m.group(1)) if m else 0


def list_chunks(montages_dir) -> List[Path]:
    """Return the montage chunk PNGs in a dir, sorted by chunk-start index.
    Long-path-safe: pathlib's is_dir()/glob() silently fail on Windows paths
    past MAX_PATH (260) — they return False / empty even when the dir exists.
    os.listdir on the \\\\?\\-prefixed path enumerates them correctly.
    Returns [] if absent (so a not-yet-generated experiment is simply skipped)."""
    d = _winlong(montages_dir)
    if not os.path.isdir(d):
        return []
    names = [f for f in os.listdir(d)
             if f.startswith("montage_cells_") and f.endswith(".png")]
    return sorted((Path(montages_dir) / n for n in names), key=_chunk_start_index)


def find_first_chunks(montages_dir: Path, n: int) -> List[Optional[Path]]:
    """Return the first n chunks sorted by chunk-start index. Pads with None
    if the folder has fewer chunks (or doesn't exist)."""
    chunks = list_chunks(montages_dir)
    return (chunks + [None] * n)[:n]


def _multichunk_geometry(n_chunks: int):
    """1 row x (2 * n_chunks) columns. DMSO's chunks left, LatA's chunks right,
    all in a single row, with a banner label spanning each condition's group of
    columns. Returns (col_w, img_h, col_lefts, img_top, banner_lefts, banner_widths)."""
    h_margin = GRID_LEFT          # slide left/right inset
    gap = 0.10                    # uniform gap between every column
    total_cols = 2 * n_chunks
    col_w = (SLIDE_W - 2 * h_margin - (total_cols - 1) * gap) / total_cols
    col_lefts = [h_margin + i * (col_w + gap) for i in range(total_cols)]

    img_top = GRID_TOP + LABEL_H
    img_h = IMG_H

    def _group(start_col):
        left = col_lefts[start_col]
        width = (col_lefts[start_col + n_chunks - 1] + col_w) - left
        return left, width
    left_banner_left,  left_banner_w  = _group(0)
    right_banner_left, right_banner_w = _group(n_chunks)
    banner_lefts   = [left_banner_left,  right_banner_left]
    banner_widths  = [left_banner_w,     right_banner_w]
    return col_w, img_h, col_lefts, img_top, banner_lefts, banner_widths


def build_compare_slide(prs, title_text,
                        left_label, left_imgs,
                        right_label, right_imgs,
                        slide_ppi, compact=False, solo=False):
    """Render the comparison slide. With n_chunks=1, lay out DMSO|LatA side
    by side (1-row layout). With n_chunks>1, switch to a single row with all
    DMSO chunks then all LatA chunks. All panels share slide_ppi so embedded
    scalebars match across the slide."""
    n_chunks = len(left_imgs)
    assert len(right_imgs) == n_chunks

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    title_h    = COMPACT_TITLE_HEIGHT    if compact else TITLE_HEIGHT
    title_font = COMPACT_TITLE_FONT_PT   if compact else TITLE_FONT_PT
    grid_top   = COMPACT_GRID_TOP        if compact else GRID_TOP
    label_h    = COMPACT_LABEL_H         if compact else LABEL_H
    label_font = COMPACT_LABEL_FONT_PT   if compact else LABEL_FONT_PT
    img_h      = COMPACT_IMG_H           if compact else IMG_H

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, title_h,
        font_pt=title_font, color=WHITE, bold=True,
    )

    missing = []

    if solo:
        cell_left = GRID_LEFT
        add_textbox(
            slide, left_label,
            cell_left, grid_top, SOLO_CELL_W, label_h,
            font_pt=label_font, color=WHITE, bold=True,
        )
        img_path = left_imgs[0]
        if img_path is not None and _exists_long(img_path):
            add_image_at_ppi(slide, img_path, slide_ppi,
                             cell_left, grid_top + label_h, SOLO_CELL_W, img_h)
        else:
            add_textbox(
                slide, "(missing)",
                cell_left, grid_top + label_h + img_h / 2 - 0.15,
                SOLO_CELL_W, 0.3,
                font_pt=14, color=WHITE,
            )
            missing.append(left_label)
        return slide, missing

    if n_chunks == 1:
        cells = [
            (left_label,  left_imgs[0],  CELL_POSITIONS[0][0]),
            (right_label, right_imgs[0], CELL_POSITIONS[1][0]),
        ]
        for label, img_path, cell_left in cells:
            add_textbox(
                slide, label,
                cell_left, grid_top, CELL_W, label_h,
                font_pt=label_font, color=WHITE, bold=True,
            )
            if img_path is not None and _exists_long(img_path):
                add_image_at_ppi(slide, img_path, slide_ppi,
                                 cell_left, grid_top + label_h, CELL_W, img_h)
            else:
                add_textbox(
                    slide, "(missing)",
                    cell_left, grid_top + label_h + img_h / 2 - 0.15,
                    CELL_W, 0.3,
                    font_pt=14, color=WHITE,
                )
                missing.append(label)
        return slide, missing

    # n_chunks > 1: single row with DMSO chunks then LatA chunks side by side.
    col_w, img_h, col_lefts, img_top, banner_lefts, banner_widths = \
        _multichunk_geometry(n_chunks)

    for cond_idx, (cond_label, cond_imgs) in enumerate(
        [(left_label, left_imgs), (right_label, right_imgs)]
    ):
        add_textbox(
            slide, cond_label,
            banner_lefts[cond_idx], GRID_TOP, banner_widths[cond_idx], LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )
        for chunk_idx, img_path in enumerate(cond_imgs):
            col_idx = cond_idx * n_chunks + chunk_idx
            cell_left = col_lefts[col_idx]
            if img_path is not None and _exists_long(img_path):
                add_image_at_ppi(slide, img_path, slide_ppi,
                                 cell_left, img_top, col_w, img_h)
            else:
                add_textbox(
                    slide, "(missing)",
                    cell_left, img_top + img_h / 2 - 0.15,
                    col_w, 0.3,
                    font_pt=14, color=WHITE,
                )
                missing.append(f"{cond_label} (chunk {chunk_idx + 1})")
    return slide, missing


def _exp_date_key(exp):
    """Sort key from the tag. Accepts MM/DD/YYYY and MM/YYYY (day unknown ->
    day 0, so a month-only tag sorts before any dated tag in the same month).
    Stable sort keeps EXPERIMENTS list order for same-key entries (e.g. the two
    concentrations of one experiment)."""
    t = exp["tag"]
    m = re.match(r"(\d{2})/(\d{2})/(\d{4})", t)
    if m:
        return (int(m.group(3)), int(m.group(1)), int(m.group(2)))
    m = re.match(r"(\d{2})/(\d{4})", t)
    if m:
        return (int(m.group(2)), int(m.group(1)), 0)
    return (9999, 99, 99)


def _make_title(base: str, tp_label: str, tag: str) -> str:
    """'<base> (<tp>, <tag>)', dropping the tp part when it's empty (basal
    datasets) -> '<base> (<tag>)'."""
    sub = f"{tp_label}, {tag}" if tp_label else tag
    return f"{base} ({sub})"


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    slide_specs: List[dict] = []
    # Group-major order: emit all of one combo across experiments before the
    # next combo. Experiments are date-sorted within the group (stable).
    for base_combo, base_title, n_chunks, scale_group, fallback, opts in COMBOS:
        for exp in sorted(EXPERIMENTS, key=_exp_date_key):
            root = Path(exp["root"])
            chan_sub = exp["chan_sub"]
            left_folder, left_label = exp["left"]
            right_folder, right_label = exp["right"]
            tag = exp["tag"]
            tp_label = exp["tp_label"]
            # exp_key labels the block in log output (and would key per-experiment
            # PPI if a combo opted into per_exp_scale; none do — pinning is
            # deck-wide per scale_group). Include the LatA condition for a
            # descriptive log line.
            exp_key = f"{tag} {tp_label} {right_folder}"
            combo_folder, base = base_combo, base_title
            primary_dir = montage_dir(root, left_folder, chan_sub, combo_folder)
            primary_ok = bool(list_chunks(primary_dir))
            if not primary_ok:
                if fallback is None:
                    continue
                combo_folder, base = fallback
            title = _make_title(base, tp_label, tag)
            left_dir  = montage_dir(root, left_folder,  chan_sub, combo_folder)
            right_dir = montage_dir(root, right_folder, chan_sub, combo_folder)
            left_imgs  = find_first_chunks(left_dir,  n_chunks)
            right_imgs = find_first_chunks(right_dir, n_chunks)
            sg_key = (f"{scale_group}@{exp_key}"
                      if opts.get("per_exp_scale") else scale_group)
            common = {
                "log_key": f"{exp_key}/{combo_folder}",
                "n_chunks": n_chunks,
                "scale_group": sg_key,
                "exp_key": exp_key,
                "compact_layout": bool(opts.get("compact_layout")),
                "scale_mult": float(opts.get("scale_mult", 1.0)),
            }
            if opts.get("solo_layout"):
                for side_label, side_imgs, side_dir in (
                    (left_label,  left_imgs,  left_dir),
                    (right_label, right_imgs, right_dir),
                ):
                    slide_specs.append({
                        **common,
                        "title": f"{title} — {side_label}",
                        "left_label": side_label, "left_imgs": side_imgs,
                        "right_label": "",        "right_imgs": [None] * n_chunks,
                        "left_dir": side_dir,     "right_dir": side_dir,
                        "is_solo": True,
                    })
            else:
                slide_specs.append({
                    **common,
                    "title": title,
                    "left_label": left_label,   "left_imgs": left_imgs,
                    "right_label": right_label, "right_imgs": right_imgs,
                    "left_dir": left_dir,       "right_dir": right_dir,
                    "is_solo": False,
                })

    def _cell_box(n_chunks: int, compact: bool, is_solo: bool = False):
        if is_solo:
            img_h = COMPACT_IMG_H if compact else IMG_H
            return SOLO_CELL_W, img_h
        if n_chunks == 1:
            img_h = COMPACT_IMG_H if compact else IMG_H
            return CELL_W, img_h
        col_w, img_h, *_ = _multichunk_geometry(n_chunks)
        return col_w, img_h

    group_ppi: dict = {}
    for spec in slide_specs:
        max_w, max_h = _cell_box(spec["n_chunks"], spec["compact_layout"],
                                 spec.get("is_solo", False))
        imgs = [p for p in (*spec["left_imgs"], *spec["right_imgs"])
                if p is not None and _exists_long(p)]
        own_ppi = compute_slide_ppi(imgs, max_w, max_h) if imgs else 0.0
        sg = spec["scale_group"]
        group_ppi[sg] = max(group_ppi.get(sg, 0.0), own_ppi)

    print(
        f"Deck-wide PPI pinning (per scale_group) across {len(slide_specs)} slides. "
        f"Scalebar invariant: {SCALEBAR_UM} μm = {SCALEBAR_PX} px in source "
        f"(PPUM = {PPUM_SOURCE} px/μm — verify with check_scalebar_pixel_widths.py).\n"
    )
    print("Pinned PPI per scale_group:")
    for sg, ppi in sorted(group_ppi.items()):
        if ppi <= 0:
            print(f"  {sg:>40s}  PPI=   0.00  (no montages found yet)")
            continue
        bar = SCALEBAR_PX / ppi
        print(f"  {sg:>40s}  PPI={ppi:>7.2f}  "
              f"scalebar={bar:.3f} in = {bar * 2.54:.3f} cm")
    print(f"\nWriting deck to: {OUTPUT_PATH}\n")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    missing_total = []
    slides_added = 0
    for spec in slide_specs:
        title       = spec["title"]
        left_label  = spec["left_label"]
        left_imgs   = spec["left_imgs"]
        right_label = spec["right_label"]
        right_imgs  = spec["right_imgs"]
        left_dir    = spec["left_dir"]
        right_dir   = spec["right_dir"]
        log_key     = spec["log_key"]
        n_chunks    = spec["n_chunks"]
        scale_group = spec["scale_group"]
        slide_ppi   = group_ppi[scale_group] / spec["scale_mult"]
        compact     = spec["compact_layout"]
        solo        = spec.get("is_solo", False)

        _, missing = build_compare_slide(
            prs, title,
            left_label, left_imgs,
            right_label, right_imgs,
            slide_ppi, compact=compact, solo=solo,
        )
        slides_added += 1

        left_ok  = sum(1 for p in left_imgs  if p is not None)
        right_ok = sum(1 for p in right_imgs if p is not None)
        print(f"[{log_key}]  L:{left_ok}/{n_chunks}  R:{right_ok}/{n_chunks}")

        for cell in missing:
            src = left_dir if cell.startswith(left_label) else right_dir
            missing_total.append(f"{log_key}/{cell}  ({src})")

    prs.save(str(out_path))
    print(f"\nDone. {slides_added} slides written to:\n  {out_path}")

    if slides_added == 0:
        print(
            "\nNo slides written — no LatA montages found yet. Expected until the "
            "MATLAB pipeline produces\n  "
            "<cond>/<chan_sub>/prog_fixed_cells/physical_scale_images/<combo>/montages/ ."
        )
    if missing_total:
        print(f"\nMissing ({len(missing_total)}):")
        for m in missing_total:
            print(f"  - {m}")
    elif slides_added:
        print("\nAll images found - no missing items.")


if __name__ == "__main__":
    main()
