"""
insert_vimkd_by_cent_summary_slides.py

Metric summary deck for the fixed-Jurkat vimentin-knockdown (VimKD) by-centrosome
compilation, siCtrl (control) vs siVim (vimentin KD):
    M:/.../VimentinKD_NucleusData_Fixed/results_compilation/
        VimKD_by_cent_siCtrl_vs_siVim_by_day_violins_20260630
Each `<metric>_by_day.png` is a single-axis siCtrl-vs-siVim violin grouped by day
(all 7 datasets on one x-axis, per-day significance). One slide = one metric
(title + the by-day panel). Companion to the VimKD physical-scale
MONTAGE decks (insert_vimkd_*_phys_scale_slides.py); this is the first VimKD
METRIC deck.

Modeled on insert_bleb_summary_slides.py (title slide, family dividers, curated
FAMILIES, --list dry-run, backup-before-overwrite, blank deck — no template).

Metric selection (per user): the LatA curated set (config_LatA_CD3_combined_*,
the CilioD-derived "no actin / no vim" morphology + invagination + curvature set)
PLUS the noco actin metrics (config_Noco_*). Vimentin metrics are excluded (vim is
the knockdown target) and the IRM synapse-area metric is excluded (no IRM here).

centrosome_center_z_rel_bottom_actin_plane (Centrosome-synapse distance) is the
lead Centrosome metric — present here as `*_by_day.png` (the upstream shared-
loader fix emits it for every day, MT included, so no separate MT variant).

One LatA metric stays omitted: nuc_broadest_slice_area — absent here (only
actin_broadest_slice_area exists).

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_vimkd_by_cent_summary_slides.py
    conda run -n PPT_editing python examples_and_configs/insert_vimkd_by_cent_summary_slides.py --list
"""

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation  # noqa: E402

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(
    "M:/FF/FF_4TB_2_Backup_fullHD/Vimentin_Project_2ndharddrive/"
    "VimentinKD_NucleusData_Fixed/results_compilation/"
    "VimKD_by_cent_siCtrl_vs_siVim_by_day_violins_20260630"
)
# By-day violin panels are flat at the compile root (one PNG per metric), each a
# single-axis siCtrl-vs-siVim violin grouped by day (all 7 datasets on one
# x-axis, with per-day significance) — NOT the faceted grid_panels (empty here).
GRID_DIR = ROOT
GRID_SUFFIX = "_by_day.png"

OUTPUT_PATH = Path(
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/VimentinKD/"
    "VimKD_Jurkats_siCtrl_vs_siVim_by_cent_summary.pptx"
)

DECK_TITLE = "Vimentin knockdown effects on fixed Jurkat nuclei (by centrosome)"
DECK_SUBTITLE = (
    "siCtrl (control) vs siVim (vimentin KD)  ·  fixed Jurkat, αCD3  ·  "
    "by-centrosome compile  ·  compiled 2026-06-30"
)

# ---------------------------------------------------------------------------
# Curated metrics, grouped into families (divider slide per family).
# (panel stem, slide title); stem + GRID_SUFFIX is the PNG under grid_panels/.
# LatA morphology/invagination/curvature set + noco actin metrics.
# ---------------------------------------------------------------------------
FAMILIES = [
    ("Centrosome ↔ nucleus", [
        ("centrosome_center_z_rel_bottom_actin_plane", "Centrosome-synapse distance"),
        ("nuc_cent_closest_dist",         "Nucleus-centrosome closest distance"),
        ("cent_nuc_norm_dist_sphere_rad", "Centrosome-nucleus distance (norm. to nuclear sphere radius)"),
    ]),
    ("Nuclear morphology", [
        ("nuc_aspect_ratio", "Nuclear aspect ratio"),
        ("nuc_solidity",     "Nuclear solidity"),
        ("nuc_volume_mesh",  "Nuclear volume"),
        ("nuc_SA_mesh",      "Nuclear surface area"),
    ]),
    ("Nuclear deformation & invaginations", [
        ("chull_max_D",                       "Max invag depth over full nucleus"),
        ("chull_max_D_by_cent",               "Invag depth near centrosome"),
        ("chull_mean_D_cent_global_ratio",    "Centrosomal Invagination Index"),
        ("centrosome_dist_deepest_real_avg_periphery_ratio", "Centrosome distance to deepest invag vs avg periphery ratio"),
        ("C_min_F_mean_by_cent",              "Curvature near centrosome"),
        ("deepest_invag_volume",              "Deepest invagination volume"),
        ("deepest_invag_fraction_chull_volume", "Deepest invag: frac of convex hull volume"),
        ("deepest_region_periph_ratio_025um", "DNA levels near invag"),
    ]),
    ("Invagination orientation", [
        ("avg_normal_angle_adaptive_region_growth",         "Deepest Invag Orientation"),
        ("avg_normal_angle_adaptive_region_growth_by_cent", "Invag orientation (near centrosome)"),
        ("avg_normal_angle_by_cent",                        "Invag Orientation (by centrosome)"),
    ]),
    ("Actin", [
        ("actin_deform_ratio",        "Cell Aspect Ratio"),
        ("actin_MFI_around_cent_2um", "Actin MFI around centrosome (2 μm)"),
        ("actin_frac_around_cent_2um", "Actin fraction around centrosome (2 μm)"),
        ("actin_bottom_mask_area",    "Synapse Area"),
    ]),
]

# ---------------------------------------------------------------------------
# Colors / layout (matches the bleb/washout summary decks)
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.10

TITLE_LEFT = MARGIN
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.55
TITLE_FONT_PT = 28

IMG_LEFT = MARGIN
IMG_TOP = 0.66
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = SLIDE_H - IMG_TOP - 0.12

# ---------------------------------------------------------------------------


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    """Place an image inside (left, top, w, h), preserving aspect ratio and
    centering on whichever dimension ends up smaller than the box."""
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.6, SLIDE_W - 2 * MARGIN, 1.4,
                font_pt=38, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.1, SLIDE_W - 2 * MARGIN, 1.0,
                font_pt=18, color=BLACK, italic=True)


def build_divider_slide(prs, family_name):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, family_name, MARGIN, 3.1, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=44, color=BLACK, bold=True)


def build_slide(prs, title_text, image_path):
    slide = _new_slide(prs)
    add_textbox(slide, title_text, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title_text), color=BLACK, bold=True)
    missing = not image_path.exists()
    if not missing:
        add_image_in_box(slide, str(image_path), IMG_LEFT, IMG_TOP, IMG_BOX_W, IMG_BOX_H)
    else:
        add_textbox(slide, "(missing)", IMG_LEFT, IMG_TOP + IMG_BOX_H / 2 - 0.2,
                    IMG_BOX_W, 0.4, font_pt=18, color=BLACK)
    return missing


def main():
    list_only = "--list" in sys.argv
    n_metrics = sum(len(items) for _, items in FAMILIES)
    est_slides = 1 + sum(1 + len(items) for _, items in FAMILIES)

    print("Source: {}".format(GRID_DIR))
    print("{} curated metrics across {} families, est. {} slides\n".format(
        n_metrics, len(FAMILIES), est_slides))

    if list_only:
        for fam, items in FAMILIES:
            print("=== {} ({}) ===".format(fam, len(items)))
            for stem, title in items:
                exists = (GRID_DIR / (stem + GRID_SUFFIX)).exists()
                print("  [{}] {:<52s} {}".format(
                    "OK " if exists else "MISS", stem + GRID_SUFFIX, title))
            print("")
        return

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)

    missing = []
    for fam, items in FAMILIES:
        build_divider_slide(prs, fam)
        print("=== {} ===".format(fam))
        for stem, title in items:
            p = GRID_DIR / (stem + GRID_SUFFIX)
            is_missing = build_slide(prs, title, p)
            print("  [{}] {} -> {!r}".format("OK" if not is_missing else "MISSING", stem + GRID_SUFFIX, title))
            if is_missing:
                missing.append(stem + GRID_SUFFIX)
        print("")

    if OUTPUT_PATH.exists():
        backup_dir = OUTPUT_PATH.parent / "backups"
        created = backup_presentation(str(OUTPUT_PATH), backup_base=str(backup_dir))
        if created:
            print("Backed up previous deck to: {}\n".format(backup_dir))

    prs.save(str(OUTPUT_PATH))
    total = len(prs.slides._sldIdLst)
    print("Done. {} metrics, {} slides written to:\n  {}".format(
        n_metrics, total, OUTPUT_PATH))
    if missing:
        print("\nSkipped {} missing panel(s):".format(len(missing)))
        for m in missing:
            print("  - {}".format(m))
    else:
        print("\nAll curated panels found - no missing items.")


if __name__ == "__main__":
    main()
