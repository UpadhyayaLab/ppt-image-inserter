"""
insert_sun1_vs_sun2_comparison_slides.py

SUN1-vs-SUN2 comparison deck from the 20260718 combined compile, which directly
compares the two nuclear-envelope proteins SUN1 and SUN2 (fixed Jurkats,
siControl) in one dataset. Each metric is a single comparison figure with
SUN1 (blue, n=109) and SUN2 (red, n=96) drawn together — so every tile is one
SUN1-vs-SUN2 plot (structurally like the MT DMSO-vs-Noco deck).

Focus: nuclear-envelope metrics + SUN × DNA correlation on the NE, density at
NE (boundary) and perinuc 0.5 μm side by side.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_sun1_vs_sun2_comparison_slides.py
    conda run -n PPT_editing python examples_and_configs/insert_sun1_vs_sun2_comparison_slides.py --list
"""

import math
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

try:
    from PIL import Image
except ImportError:  # pragma: no cover
    Image = None

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation, safe_path, path_exists  # noqa: E402

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(
    "O:/FF_backup/Jurkat_nucleus/from_Ivan_HD/results_compilation/"
    "Jurkats_SUN1_vs_SUN2_siControl_20260718")
SINGLES = ROOT / "geom_density" / "profiles" / "singles"
ENRICH = ROOT / "geom_density" / "enrichment"
GRID = ROOT / "grid_panels" / "SUN1_vs_SUN2"
DEPTH = ROOT / "invag_depth_profiles"

OUTPUT_PATH = Path(
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/SUN/"
    "SUN1 vs SUN2 comparison, siControl (20260718).pptx")

DECK_TITLE = "SUN1 vs SUN2 — nuclear-envelope geometry comparison"
DECK_SUBTITLE = (
    "Fixed Jurkats · siControl · SUN = perinuc 0.5 μm outside NE · "
    "SUN1 (blue, n=109) vs SUN2 (red, n=96), one experiment each · "
    "compiled 2026-07-18 · includes SUN × DNA correlation at the NE")

FOOTER_CMP = "SUN1 (blue) vs SUN2 (red)"

# ---------------------------------------------------------------------------
# Colors / layout (16:9)
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0x55, 0x55, 0x55)
FIELD_COLOR = RGBColor(0x2E, 0x5A, 0x88)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.08
GAP = 0.10

TITLE_LEFT = MARGIN
TITLE_TOP = 0.04
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 27

IMG_TOP = 0.58
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = 6.60

FOOTER_TOP = 7.20
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.28
FOOTER_FONT_PT = 8

CAPTION_HEIGHT = 0.50
MAX_GRID_COLS = 5


# ---------------------------------------------------------------------------
# Slide helpers (shared with the per-channel decks)
# ---------------------------------------------------------------------------
def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER, wrap=True):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    image_path = safe_path(image_path)
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.6, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=34, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.0, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=15, color=GREY, italic=True)


def build_divider_slide(prs, title, subtitle=""):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, title, MARGIN, 3.0, SLIDE_W - 2 * MARGIN, 1.2,
                font_pt=38, color=BLACK, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, MARGIN, 4.25, SLIDE_W - 2 * MARGIN, 0.9,
                    font_pt=18, color=GREY, italic=True)


def add_labeled_caption(slide, lines, left, top, width, height):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    for i, (text, pt, color) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.size = Pt(pt)
        run.font.color.rgb = color
    return box


def _tile_caption(slide, cap, path, left, top, width, height, cap_pt):
    lines = []
    if cap:
        lines.append((cap, min(cap_pt, 10), GREY))
    lines.append((Path(str(path)).stem, 8, FIELD_COLOR))
    add_labeled_caption(slide, lines, left, top, width, height)


_AR_CACHE = {}


def _img_ar(path):
    key = str(path)
    if key in _AR_CACHE:
        return _AR_CACHE[key]
    ar = 1.5
    if Image is not None:
        try:
            with Image.open(safe_path(key)) as im:
                w, h = im.size
                if h:
                    ar = w / float(h)
        except Exception:
            pass
    _AR_CACHE[key] = ar
    return ar


def _uniform_img_size(a, cols, rows, area_w, area_h, cap_h):
    cw = (area_w - (cols - 1) * GAP) / cols
    cell_h = (area_h - (rows - 1) * GAP) / rows
    ih_avail = cell_h - cap_h
    if ih_avail <= 0.2 or cw <= 0.2:
        return None
    iw = min(cw, ih_avail * a)
    return iw, iw / a


def _best_cols(ars, n, area_w, area_h, cap_h, cap_cols):
    uniform = (max(ars) / min(ars)) < 1.20
    a = sorted(ars)[len(ars) // 2]
    best = None
    for cols in range(1, min(cap_cols, n) + 1):
        rows = math.ceil(n / cols)
        if uniform:
            sz = _uniform_img_size(a, cols, rows, area_w, area_h, cap_h)
            if sz is None:
                continue
            score = sz[0] * sz[1]
        else:
            cw = (area_w - (cols - 1) * GAP) / cols
            cell_h = (area_h - (rows - 1) * GAP) / rows
            ih_avail = cell_h - cap_h
            if ih_avail <= 0.2 or cw <= 0.2:
                continue
            score = 0.0
            for ai in ars:
                iw = min(cw, ih_avail * ai)
                score += iw * (iw / ai)
        if best is None or score > best[0] + 1e-6:
            best = (score, cols, rows, uniform)
    if best is None:
        c = min(cap_cols, n)
        return c, math.ceil(n / c), uniform
    return best[1], best[2], best[3]


def build_grid_slide(prs, title, entries, footer_text, cap_pt=11, cap_h=0.30):
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    n = len(entries)
    area_top = IMG_TOP
    area_w = SLIDE_W - 2 * MARGIN
    area_h = FOOTER_TOP - IMG_TOP - 0.02

    ars = [_img_ar(p) for p, _ in entries]
    cap_cols = min(MAX_GRID_COLS, n)
    cols, rows, uniform = _best_cols(ars, n, area_w, area_h, cap_h, cap_cols)

    if uniform:
        a = sorted(ars)[len(ars) // 2]
        iw, ih = _uniform_img_size(a, cols, rows, area_w, area_h, cap_h)
        block_w = cols * iw + (cols - 1) * GAP
        block_h = rows * (ih + cap_h) + (rows - 1) * GAP
        x0 = MARGIN + (area_w - block_w) / 2
        y0 = area_top + max(0.0, (area_h - block_h) / 2)
        for i, (path, cap) in enumerate(entries):
            r, c = divmod(i, cols)
            left = x0 + c * (iw + GAP)
            top = y0 + r * (ih + cap_h + GAP)
            add_image_in_box(slide, str(path), left, top, iw, ih)
            _tile_caption(slide, cap, path, left, top + ih, iw, cap_h, cap_pt)
    else:
        cell_w = (area_w - (cols - 1) * GAP) / cols
        cell_h = (area_h - (rows - 1) * GAP) / rows
        img_h = cell_h - cap_h
        for i, (path, cap) in enumerate(entries):
            r, c = divmod(i, cols)
            left = MARGIN + c * (cell_w + GAP)
            top = area_top + r * (cell_h + GAP)
            add_image_in_box(slide, str(path), left, top, cell_w, img_h)
            _tile_caption(slide, cap, path, left, top + img_h,
                          cell_w, cap_h, cap_pt)

    add_textbox(slide, footer_text, MARGIN, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


# ---------------------------------------------------------------------------
# Content plan
# ---------------------------------------------------------------------------
def _tiles(pairs):
    kept, miss = [], []
    for path, cap in pairs:
        (kept if path_exists(path) else miss).append((path, cap))
    return kept, [p for p, _ in miss]


GEOM_DISP = {"hulldist": "hull distance",
             "mincurv": "min curvature",
             "meancurv": "mean curvature"}


def build_plan():
    plan, missing = [], []

    def add_grid(title, pairs, footer=FOOTER_CMP):
        kept, miss = _tiles(pairs)
        missing.extend(miss)
        if kept:
            plan.append(("grid", title, kept, footer))

    # ===================================================================
    # Sec 1: Density profiles — 1×2 (NE | Perinuclear) combined panel,
    #        one big 1×2 per geometry.
    # ===================================================================
    plan.append(("divider", "SUN density profiles vs NE geometry",
                 "SUN1 vs SUN2 · 1×2: NE | Perinuclear shells"))

    for geom in ["hulldist", "mincurv", "meancurv"]:
        add_grid(
            "SUN density vs {} — NE | Perinuclear".format(GEOM_DISP[geom]),
            [(SINGLES / "SUN_geomdens_{}_shells_OVERLAY_line_all.png".format(
                geom), None)])

    # ===================================================================
    # Sec 2: Enrichment (1×2 NE | Perinuclear violins, one big per slide)
    # ===================================================================
    plan.append(("divider", "SUN enrichment near invaginations & on concave NE",
                 "SUN1 vs SUN2 · 1×2: NE | Perinuclear · "
                 "enrichment = fraction ÷ expected, ref line 1"))

    en_foot = "enrichment = fraction ÷ expected · ref line 1 · " + FOOTER_CMP
    add_grid("SUN enrichment near invaginations (dist > 0.5 μm)",
             [(ENRICH / "SUN_hulldist_gt0.5um_shells.png", None)], en_foot)
    add_grid("SUN enrichment near invaginations (dist > 1 μm)",
             [(ENRICH / "SUN_hulldist_gt1.0um_shells.png", None)], en_foot)
    add_grid("SUN enrichment on concave NE (min curv < 0)",
             [(ENRICH / "SUN_mincurv_lt0_shells.png", None)], en_foot)
    add_grid("SUN enrichment on concave NE (mean curv < 0)",
             [(ENRICH / "SUN_meancurv_lt0_shells.png", None)], en_foot)

    # ===================================================================
    # Sec 3: Correlation (1×2 NE | Perinuclear violins, one big per slide)
    # ===================================================================
    plan.append(("divider", "SUN per-cell correlation vs NE geometry",
                 "SUN1 vs SUN2 · 1×2: NE | Perinuclear · ref line 0"))

    co_foot = "ref line 0 · " + FOOTER_CMP
    add_grid("SUN per-cell correlation vs hull distance",
             [(ENRICH / "SUN_corr_hulldist_shells.png", None)], co_foot)
    add_grid("SUN per-cell correlation vs min curvature",
             [(ENRICH / "SUN_corr_mincurv_shells.png", None)], co_foot)
    add_grid("SUN per-cell correlation vs mean curvature",
             [(ENRICH / "SUN_corr_meancurv_shells.png", None)], co_foot)

    # ===================================================================
    # Sec 3: SUN × DNA (Hoechst) correlation
    # ===================================================================
    plan.append(("divider", "SUN × DNA correlation at the nuclear envelope",
                 "per-cell correlation of SUN with DNA (Hoechst) · SUN1 vs SUN2"))

    add_grid(
        "SUN × DNA (Hoechst) correlation at the nuclear envelope",
        [(GRID / "SUN_NE_Hoechst_corr.png", "SUN vs DNA correlation at the NE")],
        "ref line 0 · " + FOOTER_CMP)

    add_grid(
        "SUN × DNA (Hoechst) correlation in deep invaginations",
        [(GRID / "SUN_deepest_invag_Hoechst_corr_025um.png",
          "deepest invagination · 0.25 μm"),
         (GRID / "SUN_deepest_invag_Hoechst_corr_05um.png",
          "deepest invagination · 0.5 μm")],
        "ref line 0 · " + FOOTER_CMP)

    # ===================================================================
    # Sec 4: Invagination pockets
    # ===================================================================
    plan.append(("divider", "SUN in nuclear invagination pockets",
                 "grooves ÷ reference (MFI ratios) · SUN1 vs SUN2"))

    add_grid(
        "SUN in nuclear invagination pockets",
        [(GRID / "SUN_cyto_in_nuc_hull_vs_near_convex_nuc_MFI_ratio.png",
          "grooves ÷ convex surface"),
         (GRID / "SUN_cyto_in_nuc_hull_vs_all_perinuc_MFI_ratio.png",
          "grooves ÷ whole perinuc shell"),
         (GRID / "SUN_invag_within_chull_vs_convex_within_chull_MFI_ratio.png",
          "invag interior ÷ convex rim")],
        "ref line 1 · " + FOOTER_CMP)

    # ===================================================================
    # Sec 5: Invagination depth (dense figures — one per slide)
    # ===================================================================
    plan.append(("divider", "SUN in invaginations — depth profiles",
                 "signal vs depth within invaginations · SUN1 vs SUN2"))

    add_grid("SUN invagination depth profiles",
             [(DEPTH / "SUN_invag_depth_profiles.png", None)],
             "signal vs depth within nuclear invaginations · " + FOOTER_CMP)
    add_grid("SUN invagination depth — 0.5 μm centrosome-stratified",
             [(DEPTH / "SUN_invag_depth_profiles_0_5um_cent_stratified.png",
               None)],
             "near vs away from centrosome · " + FOOTER_CMP)
    add_grid("SUN invagination depth — 1 μm centrosome-stratified",
             [(DEPTH / "SUN_invag_depth_profiles_1um_cent_stratified.png",
               None)],
             "near vs away from centrosome · " + FOOTER_CMP)
    add_grid("SUN enrichment vs centrosome distance",
             [(DEPTH / "SUN_invag_enrichment_vs_centdist.png", None)],
             "invagination enrichment vs distance from centrosome · "
             + FOOTER_CMP)

    # --- Drop orphan dividers ---
    cleaned = []
    for i, it in enumerate(plan):
        if it[0] == "divider":
            if i + 1 >= len(plan) or plan[i + 1][0] == "divider":
                continue
        cleaned.append(it)
    return cleaned, missing


def main():
    list_only = "--list" in sys.argv
    plan, missing = build_plan()

    print("Output: {}".format(OUTPUT_PATH))
    print("{} content slides (+ title) = {} total\n".format(
        len(plan), 1 + len(plan)))
    for it in plan:
        if it[0] == "divider":
            print("\n=== {} ===".format(it[1]))
        else:
            print("  [{}] {}".format(len(it[2]), it[1]))
    if missing:
        print("\nMISSING ({}):".format(len(missing)))
        for m in missing:
            print("  - {}".format(Path(m).name if hasattr(m, "name") else m))
    if list_only:
        return

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)

    for it in plan:
        if it[0] == "divider":
            build_divider_slide(prs, it[1], it[2])
        else:
            _, title, entries, footer = it
            build_grid_slide(prs, title, entries, footer)

    if OUTPUT_PATH.exists():
        backup_dir = OUTPUT_PATH.parent / "backups"
        created = backup_presentation(str(OUTPUT_PATH),
                                      backup_base=str(backup_dir))
        if created:
            print("\nBacked up previous deck under: {}".format(backup_dir))

    prs.save(str(OUTPUT_PATH))
    total = len(prs.slides._sldIdLst)
    print("\nDone. {} slides written to:\n  {}".format(total, OUTPUT_PATH))
    if missing:
        print("Skipped {} missing panel(s).".format(len(missing)))


if __name__ == "__main__":
    main()
