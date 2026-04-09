"""
insert_IFNg_deepest_invag_montages.py

Creates one slide per condition from the IFNg CTL experiment. Each slide shows:
  - fixed_con montage (top)
  - auto_con montage (bottom)
for the first montage file (cells 0-XX) found per condition.

Usage:
    python examples_and_configs/insert_IFNg_deepest_invag_montages.py
"""

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = "K:/FF/PPT/PPT_autogeneration/IFNg/IFNg_deepest_invag_montages.pptx"

BASE_DIR = "J:/FF/fixed_cell/CTL_nucleus/granuleTestsForFrank/0121026_ctls_ifnG_"

MONTAGE_SUBPATH = "channels/prog_fixed_cells/IFNg/deepest_invag_slice/panels/montages_deepest_invag"

# Slide dimensions (widescreen)
SLIDE_WIDTH_IN  = 13.333
SLIDE_HEIGHT_IN = 7.5

# Title
TITLE_LEFT   = 1.167
TITLE_TOP    = 0.0
TITLE_WIDTH  = 11.0
TITLE_HEIGHT = 0.85
TITLE_FONT_SIZE = Pt(32)

# Images: fixed_con top, auto_con bottom
IMG_LEFT      = 1.167
IMG_MAX_WIDTH = 11.0
FIXED_CON_TOP = 0.9
AUTO_CON_TOP  = 4.3
IMG_MAX_HEIGHT = 3.0   # informational — width-constrained insertion used

# Label (bottom-left, full path)
LABEL_LEFT       = 0.3
LABEL_FONT_SIZE  = Pt(8)
LABEL_FONT_NAME  = "Arial"

# Condition sort order
SUBSTRATE_ORDER  = {"3p0": 0, "8p0": 1, "glass": 2, "pll": 3}
TIMEPOINT_ORDER  = {"2hr": 0, "4hr": 1, "24hr": 2, "48hr": 3}

# ---------------------------------------------------------------------------


def format_condition_name(folder_name: str) -> str:
    """Convert folder name to readable title.

    Examples:
        3p0_4hr  -> '1.5 kPa, 4 hr'
        8p0_24hr -> '12 kPa, 24 hr'
        glass_4hr -> 'Glass, 4 hr'
        pll_2hr  -> 'PLL, 2 hr'
    """
    substrate_map = {
        "3p0":   "1.5 kPa",
        "8p0":   "12 kPa",
        "glass": "Glass",
        "pll":   "PLL",
    }
    timepoint_map = {
        "2hr":  "2 hr",
        "4hr":  "4 hr",
        "24hr": "24 hr",
        "48hr": "48 hr",
    }
    parts = folder_name.strip("_").split("_")
    if len(parts) != 2:
        return folder_name
    substrate, timepoint = parts
    substrate_label = substrate_map.get(substrate.lower(), substrate)
    timepoint_label = timepoint_map.get(timepoint.lower(), timepoint)
    return f"{substrate_label}, {timepoint_label}"


def condition_sort_key(folder_name: str):
    """Sort key: by substrate order, then timepoint order."""
    parts = folder_name.strip("_").split("_")
    if len(parts) != 2:
        return (99, 99, folder_name)
    substrate, timepoint = parts
    return (
        SUBSTRATE_ORDER.get(substrate.lower(), 99),
        TIMEPOINT_ORDER.get(timepoint.lower(), 99),
        folder_name,
    )


def find_first_montage(montage_dir: Path, variant: str) -> Path | None:
    """Return the montage with the lowest starting cell number for the given variant."""
    matches = list(montage_dir.glob(f"montage_cells_*_{variant}.png"))
    if not matches:
        return None
    def start_cell(p: Path) -> int:
        # filename: montage_cells_<start>_<end>_<variant>.png
        try:
            return int(p.stem.split("_")[2])
        except (IndexError, ValueError):
            return 9999
    return sorted(matches, key=start_cell)[0]


def insert_image_width_constrained(slide, image_path: Path, left_in: float, top_in: float, width_in: float) -> None:
    """Insert image at given position; height auto-calculated to preserve aspect ratio."""
    slide.shapes.add_picture(
        str(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(width_in),
        # height omitted -> python-pptx preserves aspect ratio
    )


def add_title(slide, text: str) -> None:
    """Add centered bold title at top of slide."""
    txBox = slide.shapes.add_textbox(
        Inches(TITLE_LEFT), Inches(TITLE_TOP),
        Inches(TITLE_WIDTH), Inches(TITLE_HEIGHT),
    )
    tf = txBox.text_frame
    tf.text = text
    para = tf.paragraphs[0]
    para.font.size = TITLE_FONT_SIZE
    para.font.bold = True
    para.alignment = PP_ALIGN.CENTER


def add_path_label(slide, full_path: str) -> None:
    """Add full image path as small text in bottom-left."""
    label_height = 0.2
    label_width = SLIDE_WIDTH_IN - LABEL_LEFT - 0.3
    top = SLIDE_HEIGHT_IN - 0.3 - label_height
    txBox = slide.shapes.add_textbox(
        Inches(LABEL_LEFT), Inches(top),
        Inches(label_width), Inches(label_height),
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    tf.text = full_path
    para = tf.paragraphs[0]
    para.font.size = LABEL_FONT_SIZE
    para.font.name = LABEL_FONT_NAME


def main() -> None:
    base = Path(BASE_DIR)
    if not base.exists():
        print(f"[ERROR] Base directory not found: {base}")
        sys.exit(1)

    # Collect condition directories (skip non-condition entries like 'results')
    condition_dirs = [
        d for d in base.iterdir()
        if d.is_dir() and d.name not in ("results",)
        and "_" in d.name
    ]
    condition_dirs.sort(key=lambda d: condition_sort_key(d.name))

    if not condition_dirs:
        print(f"[ERROR] No condition directories found in {base}")
        sys.exit(1)

    print(f"Found {len(condition_dirs)} conditions.")

    # Create fresh widescreen presentation
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]  # blank

    slides_added = 0
    missing_report = []

    for cond_dir in condition_dirs:
        cond_name = cond_dir.name
        montage_dir = cond_dir / MONTAGE_SUBPATH

        fixed_img = find_first_montage(montage_dir, "fixed_con") if montage_dir.exists() else None
        auto_img  = find_first_montage(montage_dir, "auto_con")  if montage_dir.exists() else None

        fixed_status = fixed_img.name if fixed_img else "NOT FOUND"
        auto_status  = auto_img.name  if auto_img  else "NOT FOUND"
        print(f"[{cond_name}]")
        print(f"  fixed_con: {fixed_status}")
        print(f"  auto_con:  {auto_status}")

        if not montage_dir.exists():
            missing_report.append(f"{cond_name}: montage folder missing")
        if not fixed_img:
            missing_report.append(f"{cond_name}: fixed_con not found")
        if not auto_img:
            missing_report.append(f"{cond_name}: auto_con not found")

        img_width = 9.0 if cond_name.startswith("glass") else IMG_MAX_WIDTH
        img_left  = (SLIDE_WIDTH_IN - img_width) / 2

        slide = prs.slides.add_slide(blank_layout)
        add_title(slide, f"Nuc+IFN-\u03b3 montages: {format_condition_name(cond_name)}")

        if fixed_img:
            insert_image_width_constrained(slide, fixed_img, img_left, FIXED_CON_TOP, img_width)
        if auto_img:
            insert_image_width_constrained(slide, auto_img,  img_left, AUTO_CON_TOP,  img_width)

        # Label with full path of fixed_con (or auto_con as fallback)
        label_path = str(fixed_img or auto_img or montage_dir)
        add_path_label(slide, label_path)

        slides_added += 1

    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"\nDone. {slides_added} slides saved to:\n  {OUTPUT_PATH}")

    if missing_report:
        print(f"\n[WARNING] Missing images ({len(missing_report)}):")
        for msg in missing_report:
            print(f"  - {msg}")
        sys.exit(1)
    else:
        print("\nAll images found.")


if __name__ == "__main__":
    main()
