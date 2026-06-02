"""
insert_lwi_qc_part1_minEdgeLen_16_p01_slides.py

Part 1 (Overlay / Classified / Edges / Raw) QC deck reading from the
minEdgeLen=16 rerun of the prog_lwi_swi_only pipeline.

For each lmnb1 condition subdirectory, builds one slide with a 2x2 grid
of the four p01 view montages under:

    cropped/channels/prog_lwi_swi_only_minEdgeLen_16/LWI_QC/Montages_filter_pp_load/

Companion to insert_lwi_qc_part1_p01_slides.py (which reads from the
minEdgeLen=10 dir at prog_lwi_swi_only/...). Output filename includes
the minEdgeLen_16 tag so both decks coexist.

Usage:
    python examples_and_configs/insert_lwi_qc_part1_minEdgeLen_16_p01_slides.py
"""

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/"
    "Ctrl, DZNep, bTub, Lamin B, 20260203/"
    "Ctrl_DZNep_upto1hr_LWI_QC_part1_minEdgeLen_16_LaminB1_p01.pptx"
)

PARENT_DIR = (
    "J:/FF/Nucleus_Project_up_to_1hr/"
    "020302026_jurkats_dznep_lmnB1orbTub_561LP55_488LP45_405LP40_"
)

# minEdgeLen=16 rerun output dir.
QC_SUBPATH = "cropped/channels/prog_lwi_swi_only_minEdgeLen_16/LWI_QC/Montages_filter_pp_load"

IMAGE_KINDS = [
    ("Overlay",    "Montage_Overlay_LaminB1_p01.png"),
    ("Classified", "Montage_Classified_LaminB1_p01.png"),
    ("Edges",      "Montage_Edges_LaminB1_p01.png"),
    ("Raw",        "Montage_Raw_LaminB1_p01.png"),
]

MARKER_FILTER = "lmnb1"

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.45
TITLE_FONT_PT = 24

CELL_W = 6.55
CELL_H = 3.40
LABEL_H = 0.25
IMG_H = CELL_H - LABEL_H
LABEL_FONT_PT = 14

GRID_LEFT = 0.10
GRID_TOP = 0.55
COL_GAP = SLIDE_W - 2 * GRID_LEFT - 2 * CELL_W
ROW_GAP = SLIDE_H - GRID_TOP - 2 * CELL_H - 0.1

CELL_POSITIONS = [
    (GRID_LEFT,                       GRID_TOP),                       # TL
    (GRID_LEFT + CELL_W + COL_GAP,    GRID_TOP),                       # TR
    (GRID_LEFT,                       GRID_TOP + CELL_H + ROW_GAP),    # BL
    (GRID_LEFT + CELL_W + COL_GAP,    GRID_TOP + CELL_H + ROW_GAP),    # BR
]

# ---------------------------------------------------------------------------


def format_condition_name(folder_name: str) -> str:
    """G2A3_h2o_8min_lmnb1_ -> 'H2O 8 min: Lamin B1'. Falls back to raw."""
    condition_map = {"h2o": "H2O", "dznep": "DZNep"}
    timepoint_map = {
        "4min": "4 min", "8min": "8 min", "15min": "15 min",
        "30min": "30 min", "1hr": "1 hr",
    }
    marker_map = {"lmnb1": "Lamin B1", "btub": "β-Tub"}

    parts = folder_name.strip("_").split("_")
    if len(parts) != 4:
        return folder_name

    _, condition, timepoint, marker = parts
    return (
        f"{condition_map.get(condition.lower(), condition)} "
        f"{timepoint_map.get(timepoint.lower(), timepoint)}: "
        f"{marker_map.get(marker.lower(), marker)}"
    )


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_image_in_cell(slide, image_path, cell_left, cell_top):
    pic = slide.shapes.add_picture(
        image_path,
        Inches(cell_left),
        Inches(cell_top + LABEL_H),
        width=Inches(CELL_W),
    )
    actual_h_in = pic.height / 914400.0
    if actual_h_in > IMG_H:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path,
            Inches(cell_left),
            Inches(cell_top + LABEL_H),
            height=Inches(IMG_H),
        )
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(cell_left + (CELL_W - actual_w_in) / 2)
    else:
        pic.top = Inches(cell_top + LABEL_H + (IMG_H - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def build_condition_slide(prs, cond_dir: Path):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    title_text = format_condition_name(cond_dir.name)
    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    qc_dir = cond_dir / QC_SUBPATH
    missing = []

    for (kind, filename), (cell_left, cell_top) in zip(IMAGE_KINDS, CELL_POSITIONS):
        add_textbox(
            slide, kind,
            cell_left, cell_top, CELL_W, LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )

        img_path = qc_dir / filename
        if img_path.exists():
            add_image_in_cell(slide, str(img_path), cell_left, cell_top)
        else:
            missing.append(f"{cond_dir.name}: {filename}")
            add_textbox(
                slide, "(missing)",
                cell_left, cell_top + LABEL_H + IMG_H / 2 - 0.15, CELL_W, 0.3,
                font_pt=12, color=WHITE,
            )

    return slide, missing


def main() -> None:
    parent = Path(PARENT_DIR)
    if not parent.exists():
        print(f"ERROR: Parent directory not found: {parent}")
        return

    condition_dirs = sorted([
        d for d in parent.iterdir()
        if d.is_dir() and MARKER_FILTER in d.name.lower()
    ])
    if not condition_dirs:
        print(f"ERROR: No matching subdirectories found in {parent}")
        return

    print(f"Found {len(condition_dirs)} condition directories.")
    print(f"Writing deck to: {OUTPUT_PATH}\n")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    missing_total = []
    for cond_dir in condition_dirs:
        slide, missing = build_condition_slide(prs, cond_dir)
        present = sum(1 for _, fn in IMAGE_KINDS if (cond_dir / QC_SUBPATH / fn).exists())
        print(f"[{cond_dir.name}]  {present}/{len(IMAGE_KINDS)} images present")
        missing_total.extend(missing)

    prs.save(OUTPUT_PATH)
    print(f"\nDone. {len(condition_dirs)} slides written to:\n  {OUTPUT_PATH}")

    if missing_total:
        print(f"\nMissing ({len(missing_total)}):")
        for msg in missing_total:
            print(f"  - {msg}")
    else:
        print("\nAll images found - no missing items.")


if __name__ == "__main__":
    main()
