"""
insert_vim_geomdensity_Noco_CD3_combined_slides.py

Build the Vimentin geometry-density DMSO-vs-Noco summary deck from the
20260717 across-experiments compile (aCD3-activated cells, 3 experiments:
Apr 29 2022, Jan 23 2024, Feb 27 2024).

Structurally a clone of insert_vim_geomdensity_12min_combined_slides.py; the
only difference is the comparison: DMSO (Ctrl) vs Noco 1 μM instead of
PLL vs aCD3, across three experiments instead of two.

Sections (each BYEXP panel already overlays DMSO vs Noco for all
experiments, so no separate OVERLAY slides are needed):
  1. Vim density profiles (by experiment)
  2. DNA density profiles (one slide per geometry)
  3. Vim enrichment & correlation
  4. Invag depth profiles (by experiment)
  5. Curvature profiles (one slide per stratification)
  6. Morphology scatter & centrosome
  7. Scalar violins (grid_panels)
  8. QC

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_vim_geomdensity_Noco_CD3_combined_slides.py
    conda run -n PPT_editing python examples_and_configs/insert_vim_geomdensity_Noco_CD3_combined_slides.py --list
"""

import math
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

try:
    from PIL import Image
except ImportError:  # pragma: no cover
    Image = None

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation, safe_path, path_exists  # noqa: E402

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(
    "J:/FF/fixed_cell/Vimentin/results_across_experiments/"
    "Noco_CD3_20220429_20240123_20240227_20260717"
)
SINGLES = ROOT / "geom_density" / "profiles" / "singles"
ENRICH = ROOT / "geom_density" / "enrichment"
NC = ROOT / "geom_density" / "near_cent"
MORPH = ROOT / "geom_density" / "morphology_scatter"
DEPTH = ROOT / "invag_depth_profiles"
DEPTH_SINGLES = DEPTH / "singles"
CURV = ROOT / "curvature_profiles"
GRID = ROOT / "grid_panels"
HIST = ROOT / "histograms"
SIMP = ROOT / "geom_density" / "simpson"

OUTPUT_PATH = Path(
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/"
    "nuc_mesh_struct_outside_nuc/"
    "Vimentin geom-density vs NE geometry, DMSO vs Noco combined "
    "(20220429 + 20240123 + 20240227).pptx"
)

DECK_TITLE = "Vimentin density vs nuclear-envelope geometry — DMSO vs Noco"
DECK_SUBTITLE = (
    "Fixed Jurkats · aCD3-activated · DMSO (Ctrl) vs Noco 1 μM · "
    "vim = perinuc 0.5 μm · 3 experiments "
    "(Apr 29 2022, Jan 23 2024, Feb 27 2024) · compiled 2026-07-17"
)

COND_A = "CD3_DMSO"
COND_B = "CD3_Noco"
COND_A_DISP = "DMSO"
COND_B_DISP = "Noco 1 μM"

EXP_DATES = [
    ("Apr 29 2022", "Apr_29__2022", "Apr_29,_2022", "Apr_29_2022"),
    ("Jan 23 2024", "Jan_23__2024", "Jan_23,_2024", "Jan_23_2024"),
    ("Feb 27 2024", "Feb_27__2024", "Feb_27,_2024", "Feb_27_2024"),
]

DNA_SHELLS = ["boundary", "inperinuc025", "inperinuc05"]
DNA_SHELL_DISP = {
    "boundary": "boundary",
    "inperinuc025": "perinuc 0.25 μm",
    "inperinuc05": "perinuc 0.5 μm",
}

# ---------------------------------------------------------------------------
# Colors / layout (16:9)
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0x55, 0x55, 0x55)
FIELD_COLOR = RGBColor(0x2E, 0x5A, 0x88)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.08
GAP = 0.10

TITLE_LEFT = MARGIN
TITLE_TOP = 0.04
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 27

IMG_TOP = 0.58
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = 6.60

FOOTER_TOP = 7.20
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.28
FOOTER_FONT_PT = 8

CAPTION_HEIGHT = 0.50

# Upper bound on columns the auto-layout will consider.
MAX_GRID_COLS = 5


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------
def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER, wrap=True):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    image_path = safe_path(image_path)
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.7, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=38, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.1, SLIDE_W - 2 * MARGIN, 1.1,
                font_pt=16, color=GREY, italic=True)


def build_divider_slide(prs, title, subtitle=""):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, title, MARGIN, 3.0, SLIDE_W - 2 * MARGIN, 1.2,
                font_pt=40, color=BLACK, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, MARGIN, 4.25, SLIDE_W - 2 * MARGIN, 0.9,
                    font_pt=18, color=GREY, italic=True)


def build_single_slide(prs, title, image_path, caption, footer=None):
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    cap_h = CAPTION_HEIGHT if caption else 0.0
    img_h = IMG_BOX_H - cap_h
    add_image_in_box(slide, str(image_path), MARGIN, IMG_TOP, IMG_BOX_W, img_h)
    if caption:
        add_textbox(slide, caption, MARGIN, IMG_TOP + img_h, IMG_BOX_W, cap_h,
                    font_pt=12, color=GREY)
    add_textbox(slide, footer or rel_footer(image_path), MARGIN, FOOTER_TOP,
                FOOTER_WIDTH, FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


def add_labeled_caption(slide, lines, left, top, width, height):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    for i, (text, pt, color) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.size = Pt(pt)
        run.font.color.rgb = color
    return box


def _tile_caption(slide, cap, path, left, top, width, height, cap_pt):
    lines = []
    if cap:
        lines.append((cap, min(cap_pt, 10), GREY))
    lines.append((Path(str(path)).stem, 8, FIELD_COLOR))
    add_labeled_caption(slide, lines, left, top, width, height)


_AR_CACHE = {}


def _img_ar(path):
    """Width / height of an image, cached. Falls back to 1.5 if unreadable."""
    key = str(path)
    if key in _AR_CACHE:
        return _AR_CACHE[key]
    ar = 1.5
    if Image is not None:
        try:
            with Image.open(safe_path(key)) as im:
                w, h = im.size
                if h:
                    ar = w / float(h)
        except Exception:
            pass
    _AR_CACHE[key] = ar
    return ar


def _uniform_img_size(a, cols, rows, area_w, area_h, cap_h):
    """Largest image (w, h) of aspect `a` that fits a cols x rows grid."""
    cw = (area_w - (cols - 1) * GAP) / cols
    cell_h = (area_h - (rows - 1) * GAP) / rows
    ih_avail = cell_h - cap_h
    if ih_avail <= 0.2 or cw <= 0.2:
        return None
    iw = min(cw, ih_avail * a)
    return iw, iw / a


def _best_cols(ars, n, area_w, area_h, cap_h, cap_cols):
    """Pick the column count that makes the images largest.

    Uniform-aspect grids (all profile/violin panels) are packed at their
    true aspect ratio so there is no per-cell letterboxing; mixed-aspect
    grids fall back to summed fitted area with stretched cells.
    """
    uniform = (max(ars) / min(ars)) < 1.20
    a = sorted(ars)[len(ars) // 2]
    best = None
    for cols in range(1, min(cap_cols, n) + 1):
        rows = math.ceil(n / cols)
        if uniform:
            sz = _uniform_img_size(a, cols, rows, area_w, area_h, cap_h)
            if sz is None:
                continue
            score = sz[0] * sz[1]
        else:
            cw = (area_w - (cols - 1) * GAP) / cols
            cell_h = (area_h - (rows - 1) * GAP) / rows
            ih_avail = cell_h - cap_h
            if ih_avail <= 0.2 or cw <= 0.2:
                continue
            score = 0.0
            for ai in ars:
                iw = min(cw, ih_avail * ai)
                score += iw * (iw / ai)
        if best is None or score > best[0] + 1e-6:
            best = (score, cols, rows, uniform)
    if best is None:
        return min(cap_cols, n), math.ceil(n / min(cap_cols, n)), uniform
    return best[1], best[2], best[3]


def build_grid_slide(prs, title, entries, footer_text, max_cols,
                     cap_pt=11, cap_h=0.28):
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    n = len(entries)
    area_top = IMG_TOP
    area_w = SLIDE_W - 2 * MARGIN
    area_h = FOOTER_TOP - IMG_TOP - 0.02

    ars = [_img_ar(p) for p, _ in entries]
    # The auto-layout already maximises image size, so let it consider the
    # full column range up to MAX_GRID_COLS rather than the per-slide hint.
    cap_cols = min(MAX_GRID_COLS, n)
    cols, rows, uniform = _best_cols(ars, n, area_w, area_h, cap_h, cap_cols)

    if uniform:
        # Pack images at their true aspect ratio and centre the whole block,
        # so there is no wasted letterbox space inside each cell.
        a = sorted(ars)[len(ars) // 2]
        iw, ih = _uniform_img_size(a, cols, rows, area_w, area_h, cap_h)
        block_w = cols * iw + (cols - 1) * GAP
        block_h = rows * (ih + cap_h) + (rows - 1) * GAP
        x0 = MARGIN + (area_w - block_w) / 2
        y0 = area_top + max(0.0, (area_h - block_h) / 2)
        for i, (path, cap) in enumerate(entries):
            r, c = divmod(i, cols)
            left = x0 + c * (iw + GAP)
            top = y0 + r * (ih + cap_h + GAP)
            add_image_in_box(slide, str(path), left, top, iw, ih)
            _tile_caption(slide, cap, path, left, top + ih, iw, cap_h, cap_pt)
    else:
        cell_w = (area_w - (cols - 1) * GAP) / cols
        cell_h = (area_h - (rows - 1) * GAP) / rows
        img_h = cell_h - cap_h
        for i, (path, cap) in enumerate(entries):
            r, c = divmod(i, cols)
            left = MARGIN + c * (cell_w + GAP)
            top = area_top + r * (cell_h + GAP)
            add_image_in_box(slide, str(path), left, top, cell_w, img_h)
            _tile_caption(slide, cap, path, left, top + img_h,
                          cell_w, cap_h, cap_pt)

    add_textbox(slide, footer_text, MARGIN, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


def build_matrix_slide(prs, title, row_labels, matrix, footer_text,
                       cap_pt=11, cap_h=0.28):
    """Fixed-shape grid: one row per experiment (date labelled on the left),
    one column per metric. Column captions carry the metric name."""
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    rows = len(matrix)
    cols = max(len(r) for r in matrix)
    label_w = 1.05
    area_top = IMG_TOP
    area_w = SLIDE_W - 2 * MARGIN - label_w
    area_h = FOOTER_TOP - IMG_TOP - 0.02

    ars = [_img_ar(p) for r in matrix for p, _ in r]
    a = sorted(ars)[len(ars) // 2]
    sz = _uniform_img_size(a, cols, rows, area_w, area_h, cap_h)
    if sz is None:
        iw = (area_w - (cols - 1) * GAP) / cols
        ih = iw / a
    else:
        iw, ih = sz
    block_w = cols * iw + (cols - 1) * GAP
    block_h = rows * (ih + cap_h) + (rows - 1) * GAP
    x0 = MARGIN + label_w + (area_w - block_w) / 2
    y0 = area_top + max(0.0, (area_h - block_h) / 2)

    for r, row in enumerate(matrix):
        top = y0 + r * (ih + cap_h + GAP)
        add_textbox(slide, row_labels[r], x0 - label_w, top + ih / 2 - 0.18,
                    label_w - 0.08, 0.36, font_pt=13, color=BLACK, bold=True,
                    align=PP_ALIGN.RIGHT, wrap=False)
        for c, (path, cap) in enumerate(row):
            left = x0 + c * (iw + GAP)
            add_image_in_box(slide, str(path), left, top, iw, ih)
            _tile_caption(slide, cap, path, left, top + ih, iw, cap_h, cap_pt)

    add_textbox(slide, footer_text, MARGIN, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


def rel_footer(path):
    try:
        return Path(path).relative_to(ROOT).as_posix()
    except ValueError:
        return Path(path).name


# ---------------------------------------------------------------------------
# Content plan
# ---------------------------------------------------------------------------
def _tiles(pairs):
    kept, miss = [], []
    for path, cap in pairs:
        (kept if path_exists(path) else miss).append((path, cap))
    return kept, [p for p, _ in miss]


def _date_matrix(specs, path_fn):
    """Build a (row_labels, matrix, missing) triple where each row is an
    experiment (date) and each column a metric. `path_fn(stem, date_tuple)`
    returns the image path; `specs` is a list of (stem, column_caption)."""
    labels, matrix, miss = [], [], []
    for dt in EXP_DATES:
        row = []
        for stem, cap in specs:
            p = path_fn(stem, dt)
            if path_exists(p):
                row.append((p, cap))
            else:
                miss.append(p)
        if row:
            labels.append(dt[0])
            matrix.append(row)
    return labels, matrix, miss


GEOM_LABELS_SHORT = {
    "hulldist": "hull dist",
    "mincurv": "min curv",
    "meancurv": "mean curv",
}


def build_plan():
    plan, missing = [], []

    # ===================================================================
    # Sec 1: Vim density profiles — by-experiment (1x2: Apr 29 | Jan 23 | Feb 27)
    #        Each 1x2 panel already overlays DMSO vs Noco for BOTH
    #        experiments, so no separate OVERLAY slides are needed.
    # ===================================================================
    plan.append(("divider",
                 "Vimentin density profiles",
                 "perinuc 0.5 μm shell · DMSO vs Noco, by experiment"))

    byexp = []
    for geom in ["hulldist", "mincurv", "meancurv"]:
        byexp.append((
            SINGLES / "vim_geomdens_{}_perinuc05_BYEXP_line.png".format(geom),
            GEOM_LABELS_SHORT[geom]))
    kept, miss = _tiles(byexp)
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vim density profiles — by experiment",
                     kept,
                     "each 1×2: Ctrl/DMSO (blue) vs Noco 1 μM (yellow), "
                     "Apr 29 | Jan 23 | Feb 27 · perinuc 0.5 μm",
                     3, 10, 0.28))

    # ===================================================================
    # Sec 2: DNA density profiles — one slide per geometry (3 shells)
    # ===================================================================
    plan.append(("divider",
                 "DNA density profiles",
                 "boundary / perinuc 0.25 μm / perinuc 0.5 μm · "
                 "DMSO vs Noco, by experiment"))

    for geom in ["hulldist", "mincurv", "meancurv"]:
        dna = []
        for shell in DNA_SHELLS:
            dna.append((
                SINGLES / "DNA_geomdens_{}_{}_BYEXP_line.png".format(
                    geom, shell),
                DNA_SHELL_DISP[shell]))
        kept, miss = _tiles(dna)
        missing += miss
        if kept:
            plan.append(("grid",
                         "DNA density vs {} — by experiment".format(
                             GEOM_LABELS_SHORT[geom]),
                         kept,
                         "each 1×2: Ctrl/DMSO (blue) vs Noco 1 μM (yellow), "
                         "Apr 29 | Jan 23 | Feb 27 · rows = shell",
                         1, 10, 0.28))

    # ===================================================================
    # Sec 3: Vim enrichment & correlation
    #        Rows = experiment (date on the left), columns = metric.
    # ===================================================================
    plan.append(("divider",
                 "Vimentin enrichment & correlation",
                 "NE + perinuc shells · per experiment · DMSO vs Noco"))

    def _enrich_path(stem, dt):
        return ENRICH / "vim_{}_shells_{}.png".format(stem, dt[2])

    def _add_matrix(specs, path_fn, title, footer):
        labels, matrix, miss = _date_matrix(specs, path_fn)
        missing.extend(miss)
        if matrix:
            plan.append(("matrix", title, labels, matrix, footer, 10, 0.28))

    # Enrichment near invaginations — 2 metrics x 2 experiments
    _add_matrix(
        [("hulldist_gt0.5um", "dist > 0.5 μm"),
         ("hulldist_gt1.0um", "dist > 1.0 μm")],
        _enrich_path,
        "Vim enrichment near invaginations",
        "enrichment = fraction ÷ expected · DMSO vs Noco · rows = experiment")

    # Enrichment on concave NE surfaces — 3 metrics x 2 experiments
    _add_matrix(
        [("mincurv_lt0", "min curv < 0"),
         ("mincurv_ltm0.25", "min curv < −0.25"),
         ("meancurv_lt0", "mean curv < 0")],
        _enrich_path,
        "Vim enrichment on concave NE surfaces",
        "ref line 1 · DMSO vs Noco · rows = experiment")

    # Per-cell correlation — 3 metrics x 2 experiments
    _add_matrix(
        [("corr_hulldist", "corr vs hull dist"),
         ("corr_mincurv", "corr vs min curv"),
         ("corr_meancurv", "corr vs mean curv")],
        _enrich_path,
        "Vim per-cell correlation vs NE geometry",
        "ref line 0 · DMSO vs Noco · rows = experiment")

    # Deep-invagination correlation — 2 metrics x 2 experiments
    _add_matrix(
        [("deepcorr_mincurv", "deep corr vs min curv"),
         ("deepcorr_meancurv", "deep corr vs mean curv")],
        _enrich_path,
        "Vim correlation in deep invaginations",
        "deepest-invagination voxels · ref line 0 · DMSO vs Noco · "
        "rows = experiment")

    # ===================================================================
    # Sec 4: Invag depth profiles — by-experiment (1x2) only
    # ===================================================================
    plan.append(("divider",
                 "Vimentin in invaginations — depth profiles",
                 "centrosome-stratified · DMSO vs Noco, by experiment"))

    byexp_depth = []
    for thresh, tlbl in [("0_5um", "0.5 μm"), ("1um", "1 μm")]:
        for xlim, xlbl in [("", ""), ("_xlim_0_2", " (x-lim 0–2)")]:
            byexp_depth.append((
                DEPTH_SINGLES / "vim_invag_depth_{}_BYEXP_line{}.png".format(
                    thresh, xlim),
                "{} threshold{}".format(tlbl, xlbl)))
    kept, miss = _tiles(byexp_depth)
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vim invagination depth — by experiment",
                     kept,
                     "each 1×2: Ctrl/DMSO (blue) vs Noco 1 μM (yellow), "
                     "Apr 29 | Jan 23 | Feb 27",
                     2, 10, 0.28))

    # ===================================================================
    # Sec 5: Curvature profiles — one slide per stratification
    # ===================================================================
    plan.append(("divider",
                 "Curvature-stratified vimentin profiles",
                 "min & mean curvature · centrosome-stratified"))

    # One curvature type per slide, the two experiments side by side, so the
    # (wide) profile figures render as large as possible.
    for strat, slbl in [("", "unstratified"),
                        ("_0_5um_cent_stratified", "0.5 μm centrosome-strat."),
                        ("_1um_cent_stratified", "1 μm centrosome-strat.")]:
        for curv, clbl in [("min", "min curvature"),
                           ("mean", "mean curvature")]:
            ctiles = []
            for disp, _dbl, _comma, sgl in EXP_DATES:
                ctiles.append((
                    CURV / "vim_{}_curv_profiles{}_{}.png".format(
                        curv, strat, sgl),
                    disp))
            kept, miss = _tiles(ctiles)
            missing += miss
            if kept:
                plan.append(("grid",
                             "Vim {} profiles — {}".format(clbl, slbl),
                             kept,
                             "per experiment · DMSO vs Noco",
                             2, 12, 0.30))

    # ===================================================================
    # Sec 6: Morphology scatter & centrosome level
    # ===================================================================
    plan.append(("divider",
                 "Morphology scatter & centrosome metrics",
                 "per experiment"))

    _add_matrix(
        [("vim_r_hulldist_perinuc05_bycent_vs_chull_max_D_by_cent",
          "r(hulldist) vs max D"),
         ("vim_hull_enrichment_2um_cent_vs_chull_max_D_by_cent",
          "hull enrich 2 μm cent vs max D"),
         ("vim_ratio_invag_away_1um_vs_chull_max_D_by_cent",
          "ratio invag/away 1 μm vs max D")],
        lambda stem, dt: MORPH / "{}_{}.png".format(stem, dt[2]),
        "Vim morphology scatter vs nuclear size",
        "color = centrosome side · DMSO vs Noco · rows = experiment")

    nc_tiles = []
    for disp, _dbl, comma, _sgl in EXP_DATES:
        nc_tiles.append((
            NC / "vim_near_cent_level_{}.png".format(comma),
            disp))
    kept, miss = _tiles(nc_tiles)
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vim intensity level near the centrosome",
                     kept,
                     "per experiment · DMSO vs Noco",
                     2, 12, 0.30))

    # ===================================================================
    # Sec 7: Scalar violins (grid_panels)
    # ===================================================================
    plan.append(("divider",
                 "Vimentin scalars — DMSO vs Noco",
                 "grid panels · DMSO vs Noco 1 μM"))

    def gp(stem, cap):
        return (GRID / (stem + "_grid.png"), cap)

    # Overview + invag pockets — 4 tiles
    kept, miss = _tiles([
        (GRID / "cross_experiment_heatmap.png",
         "condition × experiment heatmap"),
        (GRID / "top_differences_barplot.png",
         "top differences"),
        gp("vim_cyto_in_nuc_hull_vs_near_convex_nuc_MFI_ratio",
           "grooves ÷ convex surface"),
        gp("vim_cyto_in_nuc_hull_vs_all_perinuc_MFI_ratio",
           "grooves ÷ whole perinuc shell"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "Metric overview & invagination pockets",
                     kept,
                     "grid_panels overview + voxel convex-hull "
                     "decomposition · DMSO vs Noco",
                     4, 10, 0.28))

    # Convex-hull supporting
    kept, miss = _tiles([
        gp("vim_cyto_in_nuc_hull_MFI", "grooves MFI (raw)"),
        gp("vim_cyto_in_nuc_hull_sig_fraction",
           "fraction of vim signal in grooves"),
        gp("vim_frac_in_nuc_convex_hull",
           "fraction of vim inside the hull"),
        gp("vim_invag_within_chull_vs_all_chull_MFI_ratio",
           "invag interior ÷ all within-hull"),
        gp("vim_invag_within_chull_vs_convex_within_chull_MFI_ratio",
           "invag interior ÷ convex rim"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vim in the nuclear convex hull — supporting metrics",
                     kept,
                     "DMSO vs Noco",
                     3, 11, 0.30))

    # Deepest invag ratios
    kept, miss = _tiles([
        gp("vim_ratio_by_deepest_invag_all_0_5um",
           "all faces · 0.5 μm"),
        gp("vim_ratio_by_deepest_invag_all_1um",
           "all faces · 1 μm"),
        gp("vim_ratio_by_deepest_invag_away_0_5um",
           "away faces · 0.5 μm"),
        gp("vim_ratio_by_deepest_invag_away_1um",
           "away faces · 1 μm"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vim at the deepest invagination — ratio by side",
                     kept,
                     "vim at deepest invag ÷ reference · DMSO vs Noco · "
                     "ref line 1",
                     4, 11, 0.30))

    # Above/below ratios by centrosome side
    kept, miss = _tiles([
        gp("vim_ratio_above_below_0_5um", "above÷below 0.5 μm"),
        gp("vim_ratio_above_below_1um", "above÷below 1 μm"),
        gp("vim_ratio_above_by_side_0_5um", "above by side 0.5 μm"),
        gp("vim_ratio_above_by_side_1um", "above by side 1 μm"),
        gp("vim_ratio_below_by_side_0_5um", "below by side 0.5 μm"),
        gp("vim_ratio_below_by_side_1um", "below by side 1 μm"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vim above/below ratios by centrosome side",
                     kept,
                     "DMSO vs Noco · ref line 1",
                     3, 10, 0.28))

    # Centrosome fraction & MFI
    kept, miss = _tiles([
        gp("vim_frac_around_cent_1um", "frac around cent 1 μm"),
        gp("vim_frac_around_cent_2um", "frac around cent 2 μm"),
        gp("vim_frac_around_cent_3um", "frac around cent 3 μm"),
        gp("vim_MFI_around_cent_1um", "MFI around cent 1 μm"),
        gp("vim_MFI_around_cent_2um", "MFI around cent 2 μm"),
        gp("vim_MFI_around_cent_3um", "MFI around cent 3 μm"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vim around the centrosome — fraction & intensity",
                     kept,
                     "DMSO vs Noco",
                     3, 10, 0.28))

    # Centrosome polarity & facing
    kept, miss = _tiles([
        gp("vim_cent_nuc_cof_polarity_norm", "COF polarity (norm)"),
        gp("vim_cent_nuc_cof_polarity_um", "COF polarity (μm)"),
        gp("vim_cent_nuc_facing_index_1um", "facing index 1 μm"),
        gp("vim_cent_nuc_facing_index_2um", "facing index 2 μm"),
        gp("vim_cent_nuc_interposition_frac_1um",
           "interposition frac 1 μm"),
        gp("vim_cent_nuc_interposition_frac_2um",
           "interposition frac 2 μm"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vim centrosome polarity & facing index",
                     kept,
                     "DMSO vs Noco",
                     3, 10, 0.28))

    # Signal distribution (3D)
    kept, miss = _tiles([
        gp("vim_3D_CV", "3D CV"),
        gp("vim_3D_norm_mean", "3D norm mean"),
        gp("vim_3D_norm_SD", "3D norm SD"),
        gp("vim_3D_norm_skewness", "3D norm skewness"),
        gp("vim_total_sig", "total signal"),
        gp("vim_all_perinuc_MFI", "all perinuc MFI"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vim signal intensity & distribution",
                     kept,
                     "DMSO vs Noco",
                     3, 10, 0.28))

    # Distance, radial & z metrics
    kept, miss = _tiles([
        gp("vim_avg_dist_to_nuc_3D", "avg dist to nuc 3D"),
        gp("vim_avg_dist_to_nuc_broadest_slice",
           "avg dist broadest slice"),
        gp("vim_r_eff", "r_eff"),
        gp("vim_FDD_3D", "FDD 3D"),
        gp("vim_zCOF", "z center-of-fluorescence"),
        gp("vim_z_FDD", "z FDD"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vim spatial distribution — distance, FDD & z",
                     kept,
                     "DMSO vs Noco",
                     3, 10, 0.28))

    # Vertical (z) fractions & cent/away ratios — 4 tiles
    kept, miss = _tiles([
        gp("vim_frac_above_top_slice_nuc",
           "frac above top nuc slice"),
        gp("vim_frac_below_bottom_slice_nuc",
           "frac below bottom nuc slice"),
        gp("vim_ratio_by_cent_away_0_5um", "cent÷away 0.5 μm"),
        gp("vim_ratio_by_cent_away_1um", "cent÷away 1 μm"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vim vertical fractions & cent/away ratios",
                     kept,
                     "DMSO vs Noco",
                     2, 11, 0.30))

    # Perinuclear, MIP area & entropy — 4 tiles
    kept, miss = _tiles([
        gp("vim_perinuc_sig_fraction", "perinuc signal fraction"),
        gp("vim_MIP_area", "MIP area"),
        gp("vim_entropy_within_half_um_nuc_3D",
           "entropy 0.5 μm near nuc 3D"),
        gp("vim_norm_entropy_at_half_um_nuc_3D",
           "norm entropy 0.5 μm near nuc 3D"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vim perinuclear fraction, MIP area & entropy",
                     kept,
                     "DMSO vs Noco",
                     2, 11, 0.30))

    # Synapse metrics — 4 tiles
    kept, miss = _tiles([
        gp("vim_synapse_MFI", "synapse MFI"),
        gp("vim_synapse_inner_outer_ratio", "inner÷outer ratio"),
        gp("vim_synapse_g_ave", "gradient average"),
        gp("vim_synapse_total_sig", "synapse total signal"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vim at the synapse",
                     kept,
                     "DMSO vs Noco",
                     2, 11, 0.30))

    # Clump metrics — 4 tiles
    kept, miss = _tiles([
        gp("vim_clump_around_cent_is_biggest_1um",
           "clump around cent is biggest 1 μm"),
        gp("vim_clump_closer_to_nuc_1um",
           "clump closer to nuc 1 μm"),
        gp("cent_dist_to_vim_clump_1um",
           "cent dist to vim clump 1 μm"),
        gp("cent_vim_clump_nuc_edge_dist_ratio_1um",
           "cent/vim-clump nuc-edge ratio 1 μm"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vim clump metrics",
                     kept,
                     "DMSO vs Noco",
                     2, 11, 0.30))

    # Enrichment by centrosome hull region
    kept, miss = _tiles([
        gp("vim_enrichment_within_half_um_hull_1_um_cent",
           "0.5 μm hull, 1 μm cent"),
        gp("vim_enrichment_within_half_um_hull_2_um_cent",
           "0.5 μm hull, 2 μm cent"),
        gp("vim_enrichment_within_nuc_hull_1_um_cent",
           "nuc hull, 1 μm cent"),
        gp("vim_enrichment_within_nuc_hull_2_um_cent",
           "nuc hull, 2 μm cent"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "Vim enrichment near centrosome — hull regions",
                     kept,
                     "DMSO vs Noco",
                     4, 11, 0.30))

    # ===================================================================
    # Sec 8: QC — Simpson's paradox check (per experiment)
    # ===================================================================
    plan.append(("divider",
                 "QC",
                 "per experiment"))

    qc_tiles = []
    for disp, _dbl, comma, _sgl in EXP_DATES:
        qc_tiles.append((
            SIMP / "vim_simpson_check_{}.png".format(comma),
            "Simpson check — {}".format(disp)))
    kept, miss = _tiles(qc_tiles)
    missing += miss
    if kept:
        plan.append(("grid",
                     "Simpson's paradox check",
                     kept,
                     "pooled trend vs per-experiment · DMSO vs Noco",
                     3, 10, 0.28))

    # --- Drop orphan dividers ---
    cleaned = []
    for i, it in enumerate(plan):
        if it[0] == "divider":
            if i + 1 >= len(plan) or plan[i + 1][0] == "divider":
                continue
        cleaned.append(it)
    return cleaned, missing


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    list_only = "--list" in sys.argv
    plan, missing = build_plan()

    n_slides = 1 + len(plan)
    print("Output: {}".format(OUTPUT_PATH))
    print("{} content slides (+ title) = {} total\n".format(len(plan), n_slides))
    for it in plan:
        if it[0] == "divider":
            print("\n=== {} ===".format(it[1]))
        elif it[0] == "single":
            print("  [1] {}".format(it[1]))
        elif it[0] == "matrix":
            ntiles = sum(len(r) for r in it[3])
            print("  [{}] {} (matrix {}x{})".format(
                ntiles, it[1], len(it[3]), max(len(r) for r in it[3])))
        else:
            print("  [{}] {}".format(len(it[2]), it[1]))
    if missing:
        print("\nMISSING ({}):".format(len(missing)))
        for m in missing:
            print("  - {}".format(Path(m).name if hasattr(m, "name") else m))
    if list_only:
        return

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)

    for it in plan:
        if it[0] == "divider":
            build_divider_slide(prs, it[1], it[2])
        elif it[0] == "single":
            build_single_slide(prs, it[1], it[2], it[3])
        elif it[0] == "grid":
            _, title, entries, footer, max_cols, cap_pt, cap_h = it
            build_grid_slide(prs, title, entries, footer, max_cols, cap_pt, cap_h)
        elif it[0] == "matrix":
            _, title, row_labels, matrix, footer, cap_pt, cap_h = it
            build_matrix_slide(prs, title, row_labels, matrix, footer,
                               cap_pt, cap_h)

    if OUTPUT_PATH.exists():
        backup_dir = OUTPUT_PATH.parent / "backups"
        created = backup_presentation(str(OUTPUT_PATH),
                                      backup_base=str(backup_dir))
        if created:
            print("\nBacked up previous deck under: {}".format(backup_dir))

    prs.save(str(OUTPUT_PATH))
    total = len(prs.slides._sldIdLst)
    print("\nDone. {} slides written to:\n  {}".format(total, OUTPUT_PATH))
    if missing:
        print("Skipped {} missing panel(s).".format(len(missing)))


if __name__ == "__main__":
    main()
