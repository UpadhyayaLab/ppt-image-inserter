"""
insert_TF_geomdensity_10min_vs_3hr_slides.py

Transcription-factor (NFAT, NFkB) geometry-density decks from the 20260718
CTL 3-siCTL compiles. Each compile has one experiment (donor 1010) with two
timepoint conditions — 10 min vs 3 hr after activation — drawn together in each
OVERLAY figure, so every metric is one 10-min-vs-3-hr comparison tile
(structurally like the SUN1-vs-SUN2 / MT decks).

Builds BOTH decks (NFAT and NFkB) in one run — same structure, different
channel/root. Style follows the vim Noco / SUN comparison decks: the density,
enrichment and correlation panels are the ready-made 1×2 (NE | Perinuclear)
plots, one big per slide.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_TF_geomdensity_10min_vs_3hr_slides.py
    conda run -n PPT_editing python examples_and_configs/insert_TF_geomdensity_10min_vs_3hr_slides.py --list
"""

import math
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

try:
    from PIL import Image
except ImportError:  # pragma: no cover
    Image = None

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation, safe_path, path_exists  # noqa: E402

# ---------------------------------------------------------------------------
# Channels to build (channel prefix -> compile root)
# ---------------------------------------------------------------------------
COMPILE_ROOT = "L:/FF/transcription_factors/compiled_results/CTL_3siCTLs_{}_20260718"
TFS = ["NFAT", "NFkB"]

OUTPUT_DIR = Path("K:/FF/PPT/PPT_autogeneration/Transcription_factors")
FOOTER_CMP = "10 min vs 3 hr"

# ---------------------------------------------------------------------------
# Colors / layout (16:9)
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0x55, 0x55, 0x55)
FIELD_COLOR = RGBColor(0x2E, 0x5A, 0x88)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.08
GAP = 0.10

TITLE_LEFT = MARGIN
TITLE_TOP = 0.04
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 27

IMG_TOP = 0.58
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = 6.60

FOOTER_TOP = 7.20
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.28
FOOTER_FONT_PT = 8

CAPTION_HEIGHT = 0.50
MAX_GRID_COLS = 5


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------
def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER, wrap=True):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    image_path = safe_path(image_path)
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.6, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=34, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.0, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=15, color=GREY, italic=True)


def build_divider_slide(prs, title, subtitle=""):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, title, MARGIN, 3.0, SLIDE_W - 2 * MARGIN, 1.2,
                font_pt=38, color=BLACK, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, MARGIN, 4.25, SLIDE_W - 2 * MARGIN, 0.9,
                    font_pt=18, color=GREY, italic=True)


def add_labeled_caption(slide, lines, left, top, width, height):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    for i, (text, pt, color) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.size = Pt(pt)
        run.font.color.rgb = color
    return box


def _tile_caption(slide, cap, path, left, top, width, height, cap_pt):
    lines = []
    if cap:
        lines.append((cap, min(cap_pt, 10), GREY))
    lines.append((Path(str(path)).stem, 8, FIELD_COLOR))
    add_labeled_caption(slide, lines, left, top, width, height)


_AR_CACHE = {}


def _img_ar(path):
    key = str(path)
    if key in _AR_CACHE:
        return _AR_CACHE[key]
    ar = 1.5
    if Image is not None:
        try:
            with Image.open(safe_path(key)) as im:
                w, h = im.size
                if h:
                    ar = w / float(h)
        except Exception:
            pass
    _AR_CACHE[key] = ar
    return ar


def _uniform_img_size(a, cols, rows, area_w, area_h, cap_h):
    cw = (area_w - (cols - 1) * GAP) / cols
    cell_h = (area_h - (rows - 1) * GAP) / rows
    ih_avail = cell_h - cap_h
    if ih_avail <= 0.2 or cw <= 0.2:
        return None
    iw = min(cw, ih_avail * a)
    return iw, iw / a


def _best_cols(ars, n, area_w, area_h, cap_h, cap_cols):
    uniform = (max(ars) / min(ars)) < 1.20
    a = sorted(ars)[len(ars) // 2]
    best = None
    for cols in range(1, min(cap_cols, n) + 1):
        rows = math.ceil(n / cols)
        if uniform:
            sz = _uniform_img_size(a, cols, rows, area_w, area_h, cap_h)
            if sz is None:
                continue
            score = sz[0] * sz[1]
        else:
            cw = (area_w - (cols - 1) * GAP) / cols
            cell_h = (area_h - (rows - 1) * GAP) / rows
            ih_avail = cell_h - cap_h
            if ih_avail <= 0.2 or cw <= 0.2:
                continue
            score = 0.0
            for ai in ars:
                iw = min(cw, ih_avail * ai)
                score += iw * (iw / ai)
        if best is None or score > best[0] + 1e-6:
            best = (score, cols, rows, uniform)
    if best is None:
        c = min(cap_cols, n)
        return c, math.ceil(n / c), uniform
    return best[1], best[2], best[3]


def build_grid_slide(prs, title, entries, footer_text, cap_pt=11, cap_h=0.30):
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    n = len(entries)
    area_top = IMG_TOP
    area_w = SLIDE_W - 2 * MARGIN
    area_h = FOOTER_TOP - IMG_TOP - 0.02

    ars = [_img_ar(p) for p, _ in entries]
    cap_cols = min(MAX_GRID_COLS, n)
    cols, rows, uniform = _best_cols(ars, n, area_w, area_h, cap_h, cap_cols)

    if uniform:
        a = sorted(ars)[len(ars) // 2]
        iw, ih = _uniform_img_size(a, cols, rows, area_w, area_h, cap_h)
        block_w = cols * iw + (cols - 1) * GAP
        block_h = rows * (ih + cap_h) + (rows - 1) * GAP
        x0 = MARGIN + (area_w - block_w) / 2
        y0 = area_top + max(0.0, (area_h - block_h) / 2)
        for i, (path, cap) in enumerate(entries):
            r, c = divmod(i, cols)
            left = x0 + c * (iw + GAP)
            top = y0 + r * (ih + cap_h + GAP)
            add_image_in_box(slide, str(path), left, top, iw, ih)
            _tile_caption(slide, cap, path, left, top + ih, iw, cap_h, cap_pt)
    else:
        cell_w = (area_w - (cols - 1) * GAP) / cols
        cell_h = (area_h - (rows - 1) * GAP) / rows
        img_h = cell_h - cap_h
        for i, (path, cap) in enumerate(entries):
            r, c = divmod(i, cols)
            left = MARGIN + c * (cell_w + GAP)
            top = area_top + r * (cell_h + GAP)
            add_image_in_box(slide, str(path), left, top, cell_w, img_h)
            _tile_caption(slide, cap, path, left, top + img_h,
                          cell_w, cap_h, cap_pt)

    add_textbox(slide, footer_text, MARGIN, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


# ---------------------------------------------------------------------------
# Content plan (per transcription factor)
# ---------------------------------------------------------------------------
GEOM_DISP = {"hulldist": "hull distance",
             "mincurv": "min curvature",
             "meancurv": "mean curvature"}


def build_plan(tf, root):
    SINGLES = root / "geom_density" / "profiles" / "singles"
    ENRICH = root / "geom_density" / "enrichment"
    GRID = root / "grid_panels" / "across_time"
    DEPTH = root / "invag_depth_profiles"
    SIMP = root / "geom_density" / "simpson"

    plan, missing = [], []

    def add_grid(title, pairs, footer=FOOTER_CMP):
        kept, miss = [], []
        for path, cap in pairs:
            (kept if path_exists(path) else miss).append((path, cap))
        missing.extend(p for p, _ in miss)
        if kept:
            plan.append(("grid", title, kept, footer))

    def gp(stem, cap):
        return (GRID / (stem + ".png"), cap)

    # --- Sec 1: TF density profiles (1×4 shells strip, one big/slide) ---
    shells_note = "shells: Inner 0.5 μm | Inner 0.25 μm | NE | Perinuclear"
    plan.append(("divider", "{} density profiles vs NE geometry".format(tf),
                 "10 min vs 3 hr · " + shells_note))
    for geom in ["hulldist", "mincurv", "meancurv"]:
        add_grid(
            "{} density vs {} (across shells)".format(tf, GEOM_DISP[geom]),
            [(SINGLES / "{}_geomdens_{}_shells_OVERLAY_line_all.png".format(
                tf, geom), None)],
            shells_note + " · " + FOOTER_CMP)

    # --- Sec 2: DNA density profiles (1×4 shells strip) ---
    plan.append(("divider", "DNA density profiles vs NE geometry",
                 "10 min vs 3 hr · " + shells_note))
    for geom in ["hulldist", "mincurv", "meancurv"]:
        add_grid(
            "DNA density vs {} (across shells)".format(GEOM_DISP[geom]),
            [(SINGLES / "DNA_geomdens_{}_shells_OVERLAY_line_all.png".format(
                geom), None)],
            shells_note + " · " + FOOTER_CMP)

    # --- Sec 3: Enrichment (1×2, one big/slide) ---
    plan.append(("divider", "{} enrichment near invaginations & concave NE".format(tf),
                 "10 min vs 3 hr · 1×2: NE | Perinuclear · ref line 1"))
    en_foot = "enrichment = fraction ÷ expected · ref line 1 · " + FOOTER_CMP
    add_grid("{} enrichment near invaginations (dist > 0.5 μm)".format(tf),
             [(ENRICH / "{}_hulldist_gt0.5um_shells.png".format(tf), None)],
             en_foot)
    add_grid("{} enrichment on concave NE (min curv < 0)".format(tf),
             [(ENRICH / "{}_mincurv_lt0_shells.png".format(tf), None)], en_foot)
    add_grid("{} enrichment on concave NE (mean curv < 0)".format(tf),
             [(ENRICH / "{}_meancurv_lt0_shells.png".format(tf), None)], en_foot)

    # --- Sec 4: Correlation (1×2, one big/slide) ---
    plan.append(("divider", "{} per-cell correlation vs NE geometry".format(tf),
                 "10 min vs 3 hr · 1×2: NE | Perinuclear · ref line 0"))
    co_foot = "ref line 0 · " + FOOTER_CMP
    add_grid("{} correlation vs hull distance".format(tf),
             [(ENRICH / "{}_corr_hulldist_shells.png".format(tf), None)], co_foot)
    add_grid("{} correlation vs min curvature".format(tf),
             [(ENRICH / "{}_corr_mincurv_shells.png".format(tf), None)], co_foot)
    add_grid("{} correlation vs mean curvature".format(tf),
             [(ENRICH / "{}_corr_meancurv_shells.png".format(tf), None)], co_foot)

    # --- Sec 5: TF × DNA (Hoechst) correlation ---
    plan.append(("divider", "{} × DNA (Hoechst) correlation".format(tf),
                 "spatial correlation of {} with DNA · 10 min vs 3 hr".format(tf)))
    add_grid("{} × DNA correlation at the nuclear envelope".format(tf),
             [gp("{}_NE_Hoechst_corr".format(tf), "at the NE"),
              gp("{}_Hoechst_corr".format(tf), "whole cell")], co_foot)
    add_grid("{} × DNA correlation in deep invaginations".format(tf),
             [gp("{}_deepest_invag_Hoechst_corr_025um".format(tf),
                 "deepest invag · 0.25 μm"),
              gp("{}_deepest_invag_Hoechst_corr_05um".format(tf),
                 "deepest invag · 0.5 μm")], co_foot)

    # --- Sec 6: Nuclear vs cytoplasmic & invagination-pocket scalars ---
    plan.append(("divider", "{} scalars — distribution & pockets".format(tf),
                 "grid panels · 10 min vs 3 hr"))
    add_grid("{} nuclear vs cytoplasmic distribution".format(tf),
             [gp("{}_nuc_fraction".format(tf), "nuclear fraction"),
              gp("{}_nuc_cyto_ratio".format(tf), "nuc ÷ cyto"),
              gp("{}_nuc_cyto_MFI_ratio".format(tf), "nuc ÷ cyto MFI"),
              gp("{}_MFI".format(tf), "whole-cell MFI"),
              gp("{}_cyto_MFI".format(tf), "cytoplasmic MFI"),
              gp("{}_all_perinuc_MFI".format(tf), "perinuc MFI")])
    add_grid("{} in nuclear invagination pockets".format(tf),
             [gp("{}_frac_in_nuc_convex_hull".format(tf),
                 "fraction inside nuclear hull"),
              gp("{}_cyto_in_nuc_hull_vs_near_convex_nuc_MFI_ratio".format(tf),
                 "grooves ÷ convex surface"),
              gp("{}_cyto_in_nuc_hull_vs_nuc_MFI_ratio".format(tf),
                 "grooves ÷ nucleus MFI"),
              gp("{}_cyto_in_nuc_hull_vs_all_perinuc_MFI_ratio".format(tf),
                 "grooves ÷ whole perinuc"),
              gp("{}_cyto_in_nuc_hull_sig_fraction".format(tf),
                 "fraction of signal in grooves"),
              gp("{}_invag_within_chull_vs_convex_within_chull_MFI_ratio".format(tf),
                 "invag interior ÷ convex rim")],
             "ref line 1 · " + FOOTER_CMP)
    add_grid("{} at the deepest invagination".format(tf),
             [gp("{}_deepest_region_periph_ratio_025um".format(tf),
                 "deepest region ÷ periph (0.25 μm)"),
              gp("{}_deepest_region_periph_ratio_05um".format(tf),
                 "deepest region ÷ periph (0.5 μm)"),
              gp("{}_deepest_point_periph_ratio_025um".format(tf),
                 "deepest point ÷ periph (0.25 μm)"),
              gp("{}_deepest_point_periph_ratio_05um".format(tf),
                 "deepest point ÷ periph (0.5 μm)"),
              gp("{}_deepest_invag_ratio_edge".format(tf), "deepest ÷ edge"),
              gp("{}_deepest_invag_ratio_outer_shell".format(tf),
                 "deepest ÷ outer shell"),
              gp("{}_MFI_deepest_05um".format(tf), "MFI deepest 0.5 μm"),
              gp("{}_deepest_invag_gradient_ratio".format(tf),
                 "deepest gradient ratio")],
             "ref line 1 · " + FOOTER_CMP)

    # --- Sec 7: Invagination depth profiles ---
    plan.append(("divider", "{} in invaginations — depth profiles".format(tf),
                 "signal vs depth within invaginations · 10 min vs 3 hr"))
    add_grid("{} invagination depth profiles".format(tf),
             [(DEPTH / "{}_invag_depth_profiles.png".format(tf), None)],
             "signal vs depth within nuclear invaginations · " + FOOTER_CMP)
    add_grid("{} invagination depth — 0.5 μm centrosome-stratified".format(tf),
             [(DEPTH / "{}_invag_depth_profiles_0_5um_cent_stratified.png".format(
                 tf), None)], "near vs away from centrosome · " + FOOTER_CMP)
    add_grid("{} invagination depth — 1 μm centrosome-stratified".format(tf),
             [(DEPTH / "{}_invag_depth_profiles_1um_cent_stratified.png".format(
                 tf), None)], "near vs away from centrosome · " + FOOTER_CMP)
    add_grid("{} enrichment vs centrosome distance".format(tf),
             [(DEPTH / "{}_invag_enrichment_vs_centdist.png".format(tf), None)],
             "invagination enrichment vs distance from centrosome · " + FOOTER_CMP)

    # --- Sec 8: QC ---
    plan.append(("divider", "QC", "Simpson's paradox check · 10 min vs 3 hr"))
    add_grid("{} — Simpson's paradox check".format(tf),
             [(SIMP / "{}_simpson_check.png".format(tf), None)],
             "pooled trend vs per-condition · " + FOOTER_CMP)

    # --- Drop orphan dividers ---
    cleaned = []
    for i, it in enumerate(plan):
        if it[0] == "divider":
            if i + 1 >= len(plan) or plan[i + 1][0] == "divider":
                continue
        cleaned.append(it)
    return cleaned, missing


def build_deck(tf, list_only):
    root = Path(COMPILE_ROOT.format(tf))
    output = OUTPUT_DIR / (
        "{} vs NE geometry, 10 min vs 3 hr (20260718).pptx".format(tf))
    plan, missing = build_plan(tf, root)

    print("\n########## {} ##########".format(tf))
    print("Output: {}".format(output))
    print("{} content slides (+ title) = {} total".format(
        len(plan), 1 + len(plan)))
    for it in plan:
        if it[0] == "divider":
            print("  === {} ===".format(it[1]))
        else:
            print("    [{}] {}".format(len(it[2]), it[1]))
    if missing:
        print("  MISSING ({}):".format(len(missing)))
        for m in missing:
            print("    - {}".format(Path(m).name if hasattr(m, "name") else m))
    if list_only:
        return

    output.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    title = "{} vs nuclear-envelope geometry — 10 min vs 3 hr".format(tf)
    subtitle = (
        "CTL · 3 siControls · donor 1010 · {} = perinuc 0.5 μm · "
        "10 min (n≈30) vs 3 hr (n≈32) after activation · compiled 2026-07-18"
    ).format(tf)
    build_title_slide(prs, title, subtitle)

    for it in plan:
        if it[0] == "divider":
            build_divider_slide(prs, it[1], it[2])
        else:
            _, ttl, entries, footer = it
            build_grid_slide(prs, ttl, entries, footer)

    if output.exists():
        created = backup_presentation(str(output),
                                      backup_base=str(OUTPUT_DIR / "backups"))
        if created:
            print("  Backed up previous deck.")
    prs.save(str(output))
    print("  Done. {} slides written.".format(len(prs.slides._sldIdLst)))
    if missing:
        print("  Skipped {} missing panel(s).".format(len(missing)))


def main():
    list_only = "--list" in sys.argv
    for tf in TFS:
        build_deck(tf, list_only)


if __name__ == "__main__":
    main()
