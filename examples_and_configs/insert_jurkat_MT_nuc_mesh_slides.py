"""
insert_jurkat_MT_nuc_mesh_slides.py

Per-cell deck for the DMSO control well of the 01/23/2024 nocodazole MT
experiment (prog_fixed_Noco_MT_20240123,
W2_aCD3_E6-1_EGFP-Cen2_DMSO_AF647bTub_535Actin_Hoechst). Two slide types are
interspersed per cell:

  A. Nucleus + MT mesh strips (scene_for_microscopy_nodes/_figures): two
     [nucleus | MT | overlay] 3-panel strips stacked —
       chull_MT_cyan        nucleus mesh colored by convex-hull deviation (red),
                            MT cyan
       curv_min_MT_magenta  nucleus mesh colored by minimum curvature (grey/blue),
                            MT magenta

  B. Raw-with-meshes 4-panel (nucleus/raw_with_meshes_MT/
     panel_raw_region_depth_curv): the raw | region | depth | min-curvature
     breakdown of the nuclear mesh with MT overlaid, two orthogonal views
     stacked (xz over yz).

Deck order interleaves the two per cell: Cell 1 (A), Cell 1 (B), Cell 2 (A),
Cell 2 (B), ... Cells are discovered from the _figures filenames and ordered
numerically. Images are fit-to-box preserving aspect; physical scale is not
pinned (each strip carries its own burned-in scalebar).

Pass --no-rwm to build the nucleus+MT-only version (slide A only) to a separate
output file (OUTPUT_PATH_NO_RWM); the two builds coexist.

Usage:
    python examples_and_configs/insert_jurkat_MT_nuc_mesh_slides.py            # full deck (nuc+MT + rwm 4-panel)
    python examples_and_configs/insert_jurkat_MT_nuc_mesh_slides.py --no-rwm   # nucleus+MT only, separate output
    python examples_and_configs/insert_jurkat_MT_nuc_mesh_slides.py --list     # dry run (add --no-rwm to preview that mode)
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Optional, Sequence, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

CELL_ROOT = (
    "M:/FF/FF_4TB_2_Backup_fullHD/Nucleus Project_2ndharddrive/"
    "prog_fixed_Noco_MT_20240123/"
    "W2_aCD3_E6-1_EGFP-Cen2_DMSO_AF647bTub_535Actin_Hoechst"
)
# Slide A source: [nucleus | MT | overlay] 3-panel strips.
FIG_DIR = f"{CELL_ROOT}/scene_for_microscopy_nodes/_figures"
# Slide B source: raw|region|depth|min-curv 4-panel, MT overlaid.
RWM_DIR = f"{CELL_ROOT}/nucleus/raw_with_meshes_MT/panel_raw_region_depth_curv"

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/"
    "nuc_mesh_struct_outside_nuc/Jurkat_MT_nuc_mesh.pptx"
)
# --no-rwm build writes the nucleus+MT-only version (slide A only) here, so it
# coexists with the full interleaved deck above.
OUTPUT_PATH_NO_RWM = (
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/"
    "nuc_mesh_struct_outside_nuc/Jurkat_MT_nuc_mesh_no_rwm.pptx"
)

DATE_TAG = "01/23/2024"
CONDITION = "DMSO, αCD3"

# Slide A: nucleus + MT 3-panel strip suffixes, rendered top -> bottom.
#   file = <FIG_DIR>/cell_<N>_<suffix>_3panel.png
NUC_MT_VARIANTS = ["chull_MT_cyan", "curv_min_MT_magenta"]

# Slide B: raw-with-meshes 4-panel views, rendered top -> bottom (labeled).
#   file = <RWM_DIR>/Cell_<N>_<view>_white_grey_blue_blackbg.png
RWM_VIEWS = ["xz", "yz"]
RWM_FILE_TAIL = "white_grey_blue_blackbg"

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0xC8, 0xC8, 0xC8)

# Slide layout (inches). 13.333 x 7.5 widescreen.
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.15

TITLE_TOP = 0.06
TITLE_H = 0.54
TITLE_FONT_PT = 24

ROW_LABEL_H = 0.26
ROW_LABEL_FONT_PT = 13
ROW_GAP = 0.14
CONTENT_TOP = TITLE_TOP + TITLE_H
CONTENT_BOTTOM = SLIDE_H - 0.12

# ---------------------------------------------------------------------------


def _winlong(p) -> str:
    """Win32-safe absolute path: prepend \\\\?\\ so paths past MAX_PATH (260)
    work in os.stat / os.listdir / Image.open / add_picture."""
    s = os.path.abspath(str(p))
    if os.name == "nt" and not s.startswith("\\\\?\\"):
        s = "\\\\?\\" + s.replace("/", "\\")
    return s


def _exists_long(p) -> bool:
    return os.path.exists(_winlong(p))


def _png_dims(path) -> Tuple[int, int]:
    with Image.open(_winlong(path)) as im:
        return im.size


def discover_cells(fig_dir: str) -> List[int]:
    """Cell indices present in the _figures dir, sorted numerically. A cell
    qualifies if it has at least one nucleus+MT variant's 3-panel strip."""
    d = _winlong(fig_dir)
    if not os.path.isdir(d):
        return []
    suffixes = "|".join(re.escape(s) for s in NUC_MT_VARIANTS)
    pat = re.compile(rf"^cell_(\d+)_(?:{suffixes})_3panel\.png$")
    cells = {int(m.group(1)) for f in os.listdir(d) if (m := pat.match(f))}
    return sorted(cells)


def nuc_mt_path(cell: int, suffix: str) -> Path:
    return Path(FIG_DIR) / f"cell_{cell}_{suffix}_3panel.png"


def rwm_path(cell: int, view: str) -> Path:
    return Path(RWM_DIR) / f"Cell_{cell}_{view}_{RWM_FILE_TAIL}.png"


def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_image_fit(slide, image_path, area_left, area_top, area_w, area_h):
    """Center image inside (left, top, w, h), fit-to-box preserving aspect."""
    w_px, h_px = _png_dims(image_path)
    scale = min(area_w / w_px, area_h / h_px)
    w_in, h_in = w_px * scale, h_px * scale
    left_in = area_left + (area_w - w_in) / 2
    top_in = area_top + (area_h - h_in) / 2
    return slide.shapes.add_picture(
        _winlong(image_path), Inches(left_in), Inches(top_in), width=Inches(w_in)
    )


def add_rows_slide(prs, blank_layout, title: str, images: Sequence[Path],
                   labels: Optional[Sequence[str]] = None):
    """One slide: white title on black + N fit-to-box image rows stacked.

    ``labels`` (optional, same length as ``images``) draws a small grey label
    above each image row. Missing image paths leave their row blank.
    """
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)
    add_textbox(slide, title, MARGIN, TITLE_TOP, SLIDE_W - 2 * MARGIN, TITLE_H,
                font_pt=TITLE_FONT_PT, color=WHITE, bold=True)

    n = len(images)
    content_h = CONTENT_BOTTOM - CONTENT_TOP
    row_total = (content_h - ROW_GAP * (n - 1)) / n
    label_h = ROW_LABEL_H if labels else 0.0
    img_area_h = row_total - label_h
    img_area_w = SLIDE_W - 2 * MARGIN

    for i, img in enumerate(images):
        row_top = CONTENT_TOP + i * (row_total + ROW_GAP)
        if labels:
            add_textbox(slide, labels[i], MARGIN, row_top, img_area_w, ROW_LABEL_H,
                        font_pt=ROW_LABEL_FONT_PT, color=GREY, bold=True,
                        align=PP_ALIGN.LEFT)
        if img is not None and _exists_long(img):
            add_image_fit(slide, img, MARGIN, row_top + label_h,
                          img_area_w, img_area_h)
    return slide


def main() -> None:
    args = sys.argv[1:]
    dry_run = "--list" in args
    include_rwm = "--no-rwm" not in args
    out_path = Path(OUTPUT_PATH if include_rwm else OUTPUT_PATH_NO_RWM)

    cells = discover_cells(FIG_DIR)
    if not cells:
        print(f"No cell_*_*_3panel.png figures found in:\n  {FIG_DIR}")
        sys.exit(1)

    mode = "nuc+MT + rwm 4-panel, interleaved" if include_rwm else "nuc+MT only"
    print(f"Cells found: {len(cells)}  ({cells[0]}..{cells[-1]})   mode: {mode}")

    if dry_run:
        n_slides = 0
        for n in cells:
            nm = [s for s in NUC_MT_VARIANTS if _exists_long(nuc_mt_path(n, s))]
            line = f"  Cell {n:2d}: nuc+MT {len(nm)}/{len(NUC_MT_VARIANTS)} variant(s)"
            n_slides += 1
            if include_rwm:
                rw = [v for v in RWM_VIEWS if _exists_long(rwm_path(n, v))]
                n_slides += 1 if rw else 0
                line += (f", rwm {len(rw)}/{len(RWM_VIEWS)} view(s)"
                         + ("   [rwm 4-panel slide skipped]" if not rw else ""))
            print(line)
        print(f"\n[dry run] {n_slides} slides ({mode}) -> {out_path}")
        return

    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank_layout = prs.slide_layouts[6]
    print(f"Writing deck to: {out_path}\n")

    total = 0
    for n in cells:
        # A. Nucleus + MT: chull (top) and curv_min (bottom) 3-panel strips.
        add_rows_slide(
            prs, blank_layout,
            f"Cell {n} — nucleus and MT   ({CONDITION} · {DATE_TAG})",
            [nuc_mt_path(n, s) for s in NUC_MT_VARIANTS],
        )
        total += 1

        # B. Raw-with-meshes 4-panel (+MT): xz (top) over yz (bottom). Omitted
        # entirely in --no-rwm mode.
        if include_rwm:
            rwm_imgs = [rwm_path(n, v) for v in RWM_VIEWS]
            if any(_exists_long(p) for p in rwm_imgs):
                add_rows_slide(
                    prs, blank_layout,
                    f"Cell {n} — raw · region · depth · min-curv (+MT)   "
                    f"({CONDITION} · {DATE_TAG})",
                    rwm_imgs, labels=RWM_VIEWS,
                )
                total += 1
            else:
                print(f"  [cell {n}] no raw-with-meshes views — 4-panel slide skipped.")

        print(f"  Cell {n:2d}: done ({total} slides so far)")

    prs.save(str(out_path))
    print(f"\nDone. {total} slides written to:\n  {out_path}")


if __name__ == "__main__":
    main()
