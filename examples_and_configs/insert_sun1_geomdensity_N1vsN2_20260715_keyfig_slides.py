"""
insert_sun1_geomdensity_N1vsN2_20260715_keyfig_slides.py

Build the SUN1 geometry-density N1 vs N2 summary deck from the 20260715
all-in-one compile — KEY FIGURES variant.

Leads with polished key_figures/ plots (pooled N1+N2 comparison panels),
then continues with full per-experiment detail.

Source:
  Jurkats_SUN1_N1vsN2_siControl_geomdensity_20260715

Sections:
  0. Key figures (profiles, enrichment, correlation, MFI ratios)
  1. SUN1 density profiles — per experiment (N1 vs N2 singletons, 2-up)
  2. Per-experiment enrichment & correlation (curated + enrichment/single_condition)
  3. Invagination depth profiles (comparison plots, N1 vs N2 on same axis)
  4. Morphology scatter & centrosome (comparison plots)
  5. Scalar violins (grid_panels/N1_vs_N2 — two violins per plot)
  6. QC raw images (top/low ranked cells per condition)

Self-contained blank 16:9 deck.  Missing panels are skipped; a previous deck
at the output path is backed up first.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_sun1_geomdensity_N1vsN2_20260715_keyfig_slides.py
    conda run -n PPT_editing python examples_and_configs/insert_sun1_geomdensity_N1vsN2_20260715_keyfig_slides.py --list
"""

import math
import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation, safe_path, path_exists  # noqa: E402

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(
    "H:/FF_backup/Jurkat_nucleus/from_Ivan_HD/results_compilation/"
    "Jurkats_SUN1_N1vsN2_siControl_geomdensity_20260715"
)
CURATED = ROOT / "SUN1_loc_wrto_invag"
CURATED_SC = CURATED / "single_condition"
ENRICH = ROOT / "geom_density" / "enrichment"
ENRICH_SC = ENRICH / "single_condition"
SINGLES = ROOT / "geom_density" / "profiles" / "singles"
NC = ROOT / "geom_density" / "near_cent"
MORPH = ROOT / "geom_density" / "morphology_scatter"
DEPTH = ROOT / "invag_depth_profiles"
GRID = ROOT / "grid_panels" / "N1_vs_N2"
KEY_FIG = ROOT / "key_figures"

OUTPUT_PATH = Path(
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/SUN/"
    "SUN1 vs NE geometry, N1 vs N2, key figures (20260715).pptx"
)

DECK_TITLE = "SUN1 density vs nuclear-envelope geometry — N1 vs N2"
DECK_SUBTITLE = (
    "Fixed Jurkats · siControl · N1 (May 18 2022) vs N2 (May 31 2022) · "
    "SUN1 = perinuc 0.5 μm outside NE · "
    "per-experiment where available · compiled 2026-07-15"
)

# Per-experiment prefixes
N1 = "N1_(May_18)"
N2 = "N2_(May_31)"
N1_DISP = "N1 (May 18)"
N2_DISP = "N2 (May 31)"

# ---------------------------------------------------------------------------
# Colors / layout (16:9)
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0x55, 0x55, 0x55)
FIELD_COLOR = RGBColor(0x2E, 0x5A, 0x88)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.10
GAP = 0.16

TITLE_LEFT = MARGIN
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.55
TITLE_FONT_PT = 28

IMG_TOP = 0.66
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = 6.36

FOOTER_TOP = 7.06
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.40
FOOTER_FONT_PT = 9

CAPTION_HEIGHT = 0.55

QC_MAX_COLS = 5


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------
def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    image_path = safe_path(image_path)
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.7, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=38, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.1, SLIDE_W - 2 * MARGIN, 1.1,
                font_pt=16, color=GREY, italic=True)


def build_divider_slide(prs, title, subtitle=""):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, title, MARGIN, 3.0, SLIDE_W - 2 * MARGIN, 1.2,
                font_pt=40, color=BLACK, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, MARGIN, 4.25, SLIDE_W - 2 * MARGIN, 0.9,
                    font_pt=18, color=GREY, italic=True)


def build_single_slide(prs, title, image_path, caption, footer=None):
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    cap_h = CAPTION_HEIGHT if caption else 0.0
    img_h = IMG_BOX_H - cap_h
    add_image_in_box(slide, str(image_path), MARGIN, IMG_TOP, IMG_BOX_W, img_h)
    if caption:
        add_textbox(slide, caption, MARGIN, IMG_TOP + img_h, IMG_BOX_W, cap_h,
                    font_pt=12, color=GREY)
    add_textbox(slide, footer or rel_footer(image_path), MARGIN, FOOTER_TOP,
                FOOTER_WIDTH, FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


def add_labeled_caption(slide, lines, left, top, width, height):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    for i, (text, pt, color) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.size = Pt(pt)
        run.font.color.rgb = color
    return box


def _tile_caption(slide, cap, path, left, top, width, height, cap_pt):
    lines = []
    if cap:
        lines.append((cap, min(cap_pt, 10), GREY))
    lines.append((Path(str(path)).stem, 8, FIELD_COLOR))
    add_labeled_caption(slide, lines, left, top, width, height)


def build_grid_slide(prs, title, entries, footer_text, max_cols,
                     cap_pt=11, cap_h=0.28):
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    n = len(entries)
    cols = min(max_cols, n)
    rows = math.ceil(n / cols)
    area_top = IMG_TOP
    area_h = FOOTER_TOP - IMG_TOP - 0.02
    cell_w = (SLIDE_W - 2 * MARGIN - (cols - 1) * GAP) / cols
    cell_h = (area_h - (rows - 1) * GAP) / rows
    img_h = cell_h - cap_h
    for i, (path, cap) in enumerate(entries):
        r, c = divmod(i, cols)
        left = MARGIN + c * (cell_w + GAP)
        top = area_top + r * (cell_h + GAP)
        add_image_in_box(slide, str(path), left, top, cell_w, img_h)
        _tile_caption(slide, cap, path, left, top + img_h, cell_w, cap_h, cap_pt)
    add_textbox(slide, footer_text, MARGIN, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


def rel_footer(path):
    try:
        return Path(path).relative_to(ROOT).as_posix()
    except ValueError:
        return Path(path).name


# ---------------------------------------------------------------------------
# QC helpers (from insert_sun1_dna_geomdensity_N1vsN2_slides.py)
# ---------------------------------------------------------------------------
def parse_qc(fname):
    """N1_May_18_top1_cell85.png -> (cond, cond_disp, sortkey, caption)."""
    stem = fname[:-4] if fname.lower().endswith(".png") else fname
    m = re.match(r"^(N\d)_(.+?)_(top|low)(\d+)_cell(\d+)$", stem)
    if not m:
        return None
    cond, date, rank, rnum, cell = m.groups()
    cond_disp = "{} ({})".format(cond, date.replace("_", " "))
    sortkey = (0 if rank == "top" else 1, int(rnum))
    caption = "{}{} · cell {}".format(rank, rnum, cell)
    return cond, cond_disp, sortkey, caption


def qc_slides_for(folder):
    """Return list of ("grid", title, entries, footer, cols, cap_pt, cap_h)."""
    qc_dir = folder / "qc_raw"
    if not qc_dir.is_dir():
        return []
    by_cond = {}
    for p in sorted(qc_dir.glob("*.png")):
        parsed = parse_qc(p.name)
        if parsed is None:
            continue
        cond, cond_disp, sortkey, caption = parsed
        by_cond.setdefault((cond, cond_disp), []).append((sortkey, p, caption))
    slides = []
    for (cond, cond_disp), items in sorted(by_cond.items()):
        items.sort(key=lambda t: t[0])
        entries = [(p, cap) for _, p, cap in items]
        title = "SUN1 QC — raw signal at deepest invagination — {}".format(
            cond_disp)
        footer = ("top / low = cells with highest / lowest SUN1 in deep NE "
                  "invaginations (÷ cell mean)  ·  SUN1 orange, nucleus cyan")
        slides.append(("grid", title, entries, footer, QC_MAX_COLS, 10, 0.30))
    return slides


# ---------------------------------------------------------------------------
# Content plan
# ---------------------------------------------------------------------------
def _tiles(pairs):
    """Filter to tiles that exist on disk; return (kept, missing_paths)."""
    kept, miss = [], []
    for path, cap in pairs:
        (kept if path_exists(path) else miss).append((path, cap))
    return kept, [p for p, _ in miss]


def _n1n2_pair_curated(suffix, cap_hint):
    """Return [(n1_path, cap), (n2_path, cap)] from curated single_condition."""
    return [
        (CURATED_SC / "{}_{}.png".format(N1, suffix), "{} — {}".format(cap_hint, N1_DISP)),
        (CURATED_SC / "{}_{}.png".format(N2, suffix), "{} — {}".format(cap_hint, N2_DISP)),
    ]


def _n1n2_pair_enrich(stem, cap_hint):
    """Return [(n1_path, cap), (n2_path, cap)] from enrichment/single_condition."""
    return [
        (ENRICH_SC / "{}_SUN1_{}_shells.png".format(N1, stem), "{} — {}".format(cap_hint, N1_DISP)),
        (ENRICH_SC / "{}_SUN1_{}_shells.png".format(N2, stem), "{} — {}".format(cap_hint, N2_DISP)),
    ]


def _n1n2_pair_profile(geom, shell="perinuc05", suffix=""):
    """Return [(n1_path, cap), (n2_path, cap)] from profiles/singles."""
    return [
        (SINGLES / "SUN1_geomdens_{}_{}_{}_all_line{}.png".format(geom, shell, "N1_siControl", suffix),
         "{}".format(N1_DISP)),
        (SINGLES / "SUN1_geomdens_{}_{}_{}_all_line{}.png".format(geom, shell, "N2_siControl", suffix),
         "{}".format(N2_DISP)),
    ]


def build_plan():
    plan, missing = [], []

    # =======================================================================
    # Sec 0: Key figures (polished comparison plots from key_figures/)
    # =======================================================================
    plan.append(("divider",
                 "Key figures",
                 "polished N1 vs N2 comparison plots"))

    def kf(subdir, fname, cap):
        return (KEY_FIG / subdir / fname, cap)

    # Profiles — N1 | N2 side-by-side, per geometry+shell (from profiles/singles/)
    GEOM_LABELS = {
        "hulldist": "invagination depth (hull-boundary dist)",
        "mincurv": "min curvature (concave)",
        "meancurv": "mean curvature (concave)",
    }
    SHELL_LABELS = {
        "boundary": "NE boundary",
        "perinuc05": "perinuc 0.5 μm",
    }
    for shell in ["boundary", "perinuc05"]:
        for geom in ["hulldist", "mincurv", "meancurv"]:
            suffix = "_concave" if geom != "hulldist" else ""
            n1_path = (SINGLES /
                       "SUN1_geomdens_{}_{}_{}_all_line{}.png".format(
                           geom, shell, "N1_siControl", suffix))
            n2_path = (SINGLES /
                       "SUN1_geomdens_{}_{}_{}_all_line{}.png".format(
                           geom, shell, "N2_siControl", suffix))
            kept, miss = _tiles([
                (n1_path, N1_DISP),
                (n2_path, N2_DISP),
            ])
            missing += miss
            if kept:
                plan.append(("grid",
                             "SUN1 density vs {} — {}".format(
                                 GEOM_LABELS[geom], SHELL_LABELS[shell]),
                             kept,
                             "per-cell traces · N1 vs N2",
                             2, 13, 0.30))

    # SUN1 correlations — N1 (per-experiment singleton from key_figures/)
    kept, miss = _tiles([
        kf("correlations", "{}_04a_correlation_hulldist.png".format(N1),
           "vs hull-boundary dist"),
        kf("correlations", "{}_04b_correlation_minCurvature.png".format(N1),
           "vs min curvature"),
        kf("correlations", "{}_04c_correlation_meanCurvature.png".format(N1),
           "vs mean curvature"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 per-cell correlation vs NE geometry — N1 (May 18)",
                     kept,
                     "key_figures/correlations · ref line 0",
                     3, 10, 0.30))

    # SUN1 correlations — N2 (per-experiment singleton from key_figures/)
    kept, miss = _tiles([
        kf("correlations", "{}_04a_correlation_hulldist.png".format(N2),
           "vs hull-boundary dist"),
        kf("correlations", "{}_04b_correlation_minCurvature.png".format(N2),
           "vs min curvature"),
        kf("correlations", "{}_04c_correlation_meanCurvature.png".format(N2),
           "vs mean curvature"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 per-cell correlation vs NE geometry — N2 (May 31)",
                     kept,
                     "key_figures/correlations · ref line 0",
                     3, 10, 0.30))

    # SUN1 × DNA NE correlation — per experiment (from key_figures/)
    # Expects: key_figures/correlations/N1_(May_18)_SUN1_NE_Hoechst_corr.png etc.
    kept, miss = _tiles([
        kf("correlations", "{}_SUN1_NE_Hoechst_corr.png".format(N1),
           "SUN1 × DNA NE corr"),
        kf("correlations", "{}_SUN1_NE_Hoechst_corr.png".format(N2),
           "SUN1 × DNA NE corr"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 × DNA correlation at the nuclear envelope — key figures",
                     kept,
                     "key_figures/correlations · N1 vs N2",
                     2, 13, 0.30))

    # Enrichment singletons — N1 (from curated single_condition/)
    kept, miss = _tiles([
        (CURATED_SC / "{}_03a_enrichment_hulldist_gt0.5um.png".format(N1),
         "hull dist > 0.5 μm"),
        (CURATED_SC / "{}_03b_enrichment_hulldist_gt1.0um.png".format(N1),
         "hull dist > 1.0 μm"),
        (CURATED_SC / "{}_03c_enrichment_minCurvature_lt0.png".format(N1),
         "min curv < 0"),
        (CURATED_SC / "{}_03d_enrichment_minCurvature_ltm0.25.png".format(N1),
         "min curv < −0.25"),
        (CURATED_SC / "{}_03e_enrichment_meanCurvature.png".format(N1),
         "mean curv < 0"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 enrichment — N1 (May 18)",
                     kept,
                     "SUN1_loc_wrto_invag/single_condition · "
                     "NE + perinuc shell · ref line 1",
                     3, 10, 0.30))

    # Enrichment singletons — N2 (from curated single_condition/)
    kept, miss = _tiles([
        (CURATED_SC / "{}_03a_enrichment_hulldist_gt0.5um.png".format(N2),
         "hull dist > 0.5 μm"),
        (CURATED_SC / "{}_03b_enrichment_hulldist_gt1.0um.png".format(N2),
         "hull dist > 1.0 μm"),
        (CURATED_SC / "{}_03c_enrichment_minCurvature_lt0.png".format(N2),
         "min curv < 0"),
        (CURATED_SC / "{}_03d_enrichment_minCurvature_ltm0.25.png".format(N2),
         "min curv < −0.25"),
        (CURATED_SC / "{}_03e_enrichment_meanCurvature.png".format(N2),
         "mean curv < 0"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 enrichment — N2 (May 31)",
                     kept,
                     "SUN1_loc_wrto_invag/single_condition · "
                     "NE + perinuc shell · ref line 1",
                     3, 10, 0.30))

    # MFI ratios (invagination pockets — N1 vs N2 on same plot)
    kept, miss = _tiles([
        kf("mfi_ratios", "SUN1_cyto_in_nuc_hull_vs_near_convex_nuc_MFI_ratio.png",
           "grooves ÷ convex surface"),
        kf("mfi_ratios", "SUN1_cyto_in_nuc_hull_vs_all_perinuc_MFI_ratio.png",
           "grooves ÷ whole perinuc shell"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 in the nuclear invagination pockets — key figures",
                     kept,
                     "key_figures/mfi_ratios · >1 = enriched in grooves · "
                     "N1 vs N2",
                     2, 13, 0.30))

    # MFI ratio singletons — grooves ÷ convex surface, per experiment
    kept, miss = _tiles([
        kf("mfi_ratios",
           "{}_SUN1_cyto_in_nuc_hull_vs_near_convex_nuc_MFI_ratio.png".format(N1),
           "{}".format(N1_DISP)),
        kf("mfi_ratios",
           "{}_SUN1_cyto_in_nuc_hull_vs_near_convex_nuc_MFI_ratio.png".format(N2),
           "{}".format(N2_DISP)),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 grooves ÷ convex surface — per experiment",
                     kept,
                     "key_figures/mfi_ratios · >1 = enriched in grooves",
                     2, 13, 0.30))

    # =======================================================================
    # Sec 1: SUN1 density profiles — per experiment
    # =======================================================================
    plan.append(("divider",
                 "SUN1 density profiles — per experiment",
                 "per-cell traces · perinuc 0.5 μm shell · N1 vs N2"))

    # Hull-boundary distance
    kept, miss = _tiles(_n1n2_pair_profile("hulldist"))
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 density vs invagination depth — N1 vs N2",
                     kept,
                     "per-cell traces · relative density (÷ cell mean) vs "
                     "hull-boundary distance · perinuc 0.5 μm",
                     2, 10, 0.28))

    # Min curvature (concave half)
    kept, miss = _tiles(_n1n2_pair_profile("mincurv", suffix="_concave"))
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 density vs min curvature — N1 vs N2",
                     kept,
                     "per-cell traces · concave half · perinuc 0.5 μm",
                     2, 10, 0.28))

    # Mean curvature (concave half)
    kept, miss = _tiles(_n1n2_pair_profile("meancurv", suffix="_concave"))
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 density vs mean curvature — N1 vs N2",
                     kept,
                     "per-cell traces · concave half · perinuc 0.5 μm",
                     2, 10, 0.28))

    # =======================================================================
    # Sec 2: SUN1 enrichment — per experiment (curated single_condition)
    # =======================================================================
    plan.append(("divider",
                 "SUN1 enrichment & correlation — per experiment",
                 "each violin = N1 vs N2 within one experiment · "
                 "curated evidence"))

    # Depth enrichment (03a + 03b): N1|N2 × 2 = 4-up
    tiles = (_n1n2_pair_curated("03a_enrichment_hulldist_gt0.5um", "hull dist > 0.5 μm")
             + _n1n2_pair_curated("03b_enrichment_hulldist_gt1.0um", "hull dist > 1.0 μm"))
    kept, miss = _tiles(tiles)
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 levels near invaginations — per experiment",
                     kept,
                     "NE + perinuc shell · ref line 1 · N1 vs N2",
                     4, 10, 0.30))

    # Curvature enrichment (03c + 03d + 03e): N1|N2 × 3 = 6-up
    tiles = (_n1n2_pair_curated("03c_enrichment_minCurvature_lt0", "min curv < 0")
             + _n1n2_pair_curated("03d_enrichment_minCurvature_ltm0.25", "min curv < −0.25")
             + _n1n2_pair_curated("03e_enrichment_meanCurvature", "mean curv < 0"))
    kept, miss = _tiles(tiles)
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 levels on concave NE surfaces — per experiment",
                     kept,
                     "NE + perinuc shell · ref line 1 · N1 vs N2",
                     2, 10, 0.28))

    # Correlation (04a + 04b + 04c): N1|N2 × 3 = 6-up
    tiles = (_n1n2_pair_curated("04a_correlation_hulldist", "vs hull-boundary dist")
             + _n1n2_pair_curated("04b_correlation_minCurvature", "vs min curvature")
             + _n1n2_pair_curated("04c_correlation_meanCurvature", "vs mean curvature"))
    kept, miss = _tiles(tiles)
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 per-cell correlation vs NE geometry — per experiment",
                     kept,
                     "ref line 0 · N1 vs N2",
                     2, 10, 0.28))

    # Deep correlation (enrichment/single_condition): N1|N2 × 2 = 4-up
    tiles = (_n1n2_pair_enrich("deepcorr_mincurv", "vs min curvature (deep)")
             + _n1n2_pair_enrich("deepcorr_meancurv", "vs mean curvature (deep)"))
    kept, miss = _tiles(tiles)
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 correlation vs curvature in deep invaginations "
                     "— per experiment",
                     kept,
                     "ref line 0 · N1 vs N2",
                     2, 10, 0.28))

    # =======================================================================
    # Sec 3: Invagination depth profiles (comparison plots)
    # =======================================================================
    plan.append(("divider",
                 "SUN1 in invaginations — depth profiles",
                 "centrosome-stratified (near vs away) · N1 vs N2"))

    for fname, cap in [
        ("SUN1_invag_depth_profiles_0_5um_cent_stratified.png",
         "depth profile, 0.5 μm perinuc"),
        ("SUN1_invag_depth_profiles_1um_cent_stratified.png",
         "depth profile, 1 μm perinuc"),
        ("SUN1_invag_enrichment_vs_centdist.png",
         "enrichment vs centrosome distance"),
    ]:
        p = DEPTH / fname
        if path_exists(p):
            plan.append(("single",
                         "SUN1 in invaginations — {}".format(cap),
                         p,
                         "near vs away from centrosome · N1 vs N2"))
        else:
            missing.append(p)

    # =======================================================================
    # Sec 4: Morphology scatter & centrosome
    # =======================================================================
    plan.append(("divider",
                 "Morphology scatter & centrosome metrics",
                 "enrichment / correlation vs deepest invagination depth, "
                 "by centrosome side · N1 vs N2"))

    # Morphology scatter (4 SUN1 plots, 2×2)
    scat_tiles = []
    for fname, cap in [
        ("SUN1_r_hulldist_boundary_bycent_vs_chull_max_D_by_cent.png",
         "r(hulldist) boundary by cent vs max D"),
        ("SUN1_r_hulldist_perinuc05_bycent_vs_chull_max_D_by_cent.png",
         "r(hulldist) perinuc05 by cent vs max D"),
        ("SUN1_hull_enrichment_2um_cent_vs_chull_max_D_by_cent.png",
         "hull enrich 2μm cent vs max D"),
        ("SUN1_ratio_invag_away_1um_vs_chull_max_D_by_cent.png",
         "ratio invag/away 1μm vs max D"),
    ]:
        scat_tiles.append((MORPH / fname, cap))
    kept, miss = _tiles(scat_tiles)
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 vs invagination depth — morphology scatter",
                     kept,
                     "color = centrosome side · N1 vs N2",
                     2, 9, 0.30))

    # Near centrosome level + SUN1 × DNA correlation at centrosome
    nc_tiles = []
    for fname, cap in [
        ("SUN1_near_cent_level.png", "SUN1 near centrosome NE"),
        ("SUN1_x_DNA_near_cent_boundary_corr.png",
         "SUN1 × DNA boundary corr near cent"),
    ]:
        nc_tiles.append((NC / fname, cap))
    kept, miss = _tiles(nc_tiles)
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 near the centrosome — NE-facing metrics",
                     kept,
                     "SUN1 ÷ cell mean near centrosome NE · N1 vs N2",
                     2, 12, 0.30))

    # =======================================================================
    # Sec 5: Scalar violins (grid_panels/N1_vs_N2)
    # =======================================================================
    plan.append(("divider",
                 "SUN1 scalars — N1 vs N2",
                 "two violins per panel: N1 (May 18) vs N2 (May 31)"))

    def gp(stem, cap):
        return (GRID / (stem + ".png"), cap)

    # Featured invagination pockets
    kept, miss = _tiles([
        gp("SUN1_cyto_in_nuc_hull_vs_near_convex_nuc_MFI_ratio",
           "grooves ÷ convex surface"),
        gp("SUN1_cyto_in_nuc_hull_vs_all_perinuc_MFI_ratio",
           "grooves ÷ whole perinuc shell"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 in the nuclear invagination pockets",
                     kept,
                     "voxel convex-hull decomposition · "
                     ">1 = enriched in grooves · N1 vs N2",
                     2, 13, 0.30))

    # SUN1 × DNA NE correlation
    kept, miss = _tiles([
        gp("SUN1_NE_Hoechst_corr", "SUN1 × DNA NE correlation"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 × DNA correlation at the nuclear envelope",
                     kept,
                     "NE boundary shell · N1 vs N2",
                     2, 13, 0.30))

    # Convex-hull supporting
    kept, miss = _tiles([
        gp("SUN1_cyto_in_nuc_hull_MFI", "grooves MFI (raw)"),
        gp("SUN1_cyto_in_nuc_hull_sig_fraction",
           "fraction of SUN1 signal in grooves"),
        gp("SUN1_frac_in_nuc_convex_hull",
           "fraction of SUN1 inside the hull"),
        gp("SUN1_invag_within_chull_vs_all_chull_MFI_ratio",
           "invag interior ÷ all within-hull"),
        gp("SUN1_invag_within_chull_vs_convex_within_chull_MFI_ratio",
           "invag interior ÷ convex rim"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 in the nuclear convex hull — supporting metrics",
                     kept,
                     "N1 vs N2",
                     3, 11, 0.30))

    # Deepest invagination
    kept, miss = _tiles([
        gp("SUN1_deepest_invag_ratio_edge", "deepest invag ÷ edge"),
        gp("SUN1_deepest_invag_ratio_outer_shell",
           "deepest invag ÷ outer shell"),
        gp("SUN1_deepest_invag_Hoechst_corr_025um",
           "SUN1×DNA corr 0.25 μm"),
        gp("SUN1_deepest_invag_Hoechst_corr_05um",
           "SUN1×DNA corr 0.5 μm"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 at the deepest invagination",
                     kept,
                     "N1 vs N2",
                     4, 11, 0.30))

    # Invag-vs-other ratios
    kept, miss = _tiles([
        gp("SUN1_invag_other_ratio_edge_025um",
           "invag÷other edge 0.25 μm"),
        gp("SUN1_invag_other_ratio_edge_05um",
           "invag÷other edge 0.5 μm"),
        gp("SUN1_invag_other_ratio_edge_1um",
           "invag÷other edge 1 μm"),
        gp("SUN1_invag_other_ratio_outer_shell_025um",
           "invag÷other outer shell 0.25 μm"),
        gp("SUN1_invag_other_ratio_outer_shell_05um",
           "invag÷other outer shell 0.5 μm"),
        gp("SUN1_invag_other_ratio_outer_shell_1um",
           "invag÷other outer shell 1 μm"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 invagination vs other NE ratios",
                     kept,
                     "N1 vs N2 · ref line 1",
                     3, 10, 0.28))

    # Ratio above/below by centrosome side
    kept, miss = _tiles([
        gp("SUN1_ratio_above_below_0_5um",
           "above÷below 0.5 μm"),
        gp("SUN1_ratio_above_below_1um",
           "above÷below 1 μm"),
        gp("SUN1_ratio_above_by_side_0_5um",
           "above by side 0.5 μm"),
        gp("SUN1_ratio_above_by_side_1um",
           "above by side 1 μm"),
        gp("SUN1_ratio_below_by_side_0_5um",
           "below by side 0.5 μm"),
        gp("SUN1_ratio_below_by_side_1um",
           "below by side 1 μm"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 above/below ratios by centrosome side",
                     kept,
                     "N1 vs N2 · ref line 1",
                     3, 10, 0.28))

    # By deepest invagination
    kept, miss = _tiles([
        gp("SUN1_ratio_by_deepest_invag_all_0_5um",
           "all faces · 0.5 μm"),
        gp("SUN1_ratio_by_deepest_invag_all_1um",
           "all faces · 1 μm"),
        gp("SUN1_ratio_by_deepest_invag_away_0_5um",
           "away faces · 0.5 μm"),
        gp("SUN1_ratio_by_deepest_invag_away_1um",
           "away faces · 1 μm"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 at the deepest invagination — ratio by side",
                     kept,
                     "SUN1 at deepest invag ÷ reference · N1 vs N2 · "
                     "ref line 1",
                     4, 11, 0.30))

    # Signal distribution / peak / boundary / outer shell metrics
    kept, miss = _tiles([
        gp("SUN1_MFI_at_nuc_boundary", "MFI at boundary"),
        gp("SUN1_MFI_outer_shell", "MFI outer shell"),
        gp("SUN1_peak_sig", "peak signal"),
        gp("SUN1_all_perinuc_MFI", "all perinuc MFI"),
        gp("SUN1_perinuc_sig_fraction", "perinuc signal fraction"),
        gp("SUN1_total_sig", "total signal"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 signal intensity & distribution",
                     kept,
                     "N1 vs N2",
                     3, 10, 0.28))

    # Heterogeneity metrics (CV, SD, skewness, entropy)
    kept, miss = _tiles([
        gp("SUN1_CV_at_nuc_boundary", "CV at boundary"),
        gp("SUN1_CV_outer_shell", "CV outer shell"),
        gp("SUN1_SD_at_nuc_boundary", "SD at boundary"),
        gp("SUN1_SD_outer_shell", "SD outer shell"),
        gp("SUN1_norm_entropy_at_nuc_boundary",
           "norm entropy at boundary"),
        gp("SUN1_norm_entropy_outer_shell",
           "norm entropy outer shell"),
        gp("SUN1_skewness_at_nuc_boundary",
           "skewness at boundary"),
        gp("SUN1_skewness_outer_shell",
           "skewness outer shell"),
    ])
    missing += miss
    if kept:
        plan.append(("grid",
                     "SUN1 heterogeneity at the nuclear envelope",
                     kept,
                     "N1 vs N2",
                     4, 10, 0.28))

    # =======================================================================
    # Sec 6: QC raw images
    # =======================================================================
    plan.append(("divider",
                 "QC — raw signal images",
                 "SUN1 signal at deepest invagination · top / low ranked cells"))

    for qc_it in qc_slides_for(CURATED):
        plan.append(qc_it)

    # --- Drop orphan dividers ---
    cleaned = []
    for i, it in enumerate(plan):
        if it[0] == "divider":
            if i + 1 >= len(plan) or plan[i + 1][0] == "divider":
                continue
        cleaned.append(it)
    return cleaned, missing


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    list_only = "--list" in sys.argv
    plan, missing = build_plan()

    n_slides = 1 + len(plan)
    print("Output: {}".format(OUTPUT_PATH))
    print("{} content slides (+ title) = {} total\n".format(len(plan), n_slides))
    for it in plan:
        if it[0] == "divider":
            print("\n=== {} ===".format(it[1]))
        elif it[0] == "single":
            print("  [1] {}".format(it[1]))
        else:
            print("  [{}] {}".format(len(it[2]), it[1]))
    if missing:
        print("\nMISSING ({}):".format(len(missing)))
        for m in missing:
            print("  - {}".format(Path(m).name if hasattr(m, "name") else m))
    if list_only:
        return

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)

    for it in plan:
        if it[0] == "divider":
            build_divider_slide(prs, it[1], it[2])
        elif it[0] == "single":
            build_single_slide(prs, it[1], it[2], it[3])
        elif it[0] == "grid":
            _, title, entries, footer, max_cols, cap_pt, cap_h = it
            build_grid_slide(prs, title, entries, footer, max_cols, cap_pt, cap_h)

    if OUTPUT_PATH.exists():
        backup_dir = OUTPUT_PATH.parent / "backups"
        created = backup_presentation(str(OUTPUT_PATH),
                                      backup_base=str(backup_dir))
        if created:
            print("\nBacked up previous deck under: {}".format(backup_dir))

    prs.save(str(OUTPUT_PATH))
    total = len(prs.slides._sldIdLst)
    print("\nDone. {} slides written to:\n  {}".format(total, OUTPUT_PATH))
    if missing:
        print("Skipped {} missing panel(s).".format(len(missing)))


if __name__ == "__main__":
    main()
