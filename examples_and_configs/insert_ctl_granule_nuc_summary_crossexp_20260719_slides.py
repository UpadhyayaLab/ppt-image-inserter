"""
insert_ctl_granule_nuc_summary_crossexp_20260719_slides.py

Cross-experiment condition-comparison summary deck for the fixed activated-CTL
granule/nucleus experiments, pooling the two datasets:
  - June 17 2026 (20260617_Fixed_CTLs_glass_...3min_12min)
  - July 16 2026 (adds a 5 min timepoint)
compiled into
compiled_results/CTL_nuc_MT_granules_20260617_20260716_crossexp_20260719.
Conditions are TIMEPOINTS after activation (3 min / 12 min, plus 5 min for the
July 16 dataset) on αCD3/ICAM1/3SI glass. Channels: LAMP1 (lytic granules),
MT (β-tubulin; also the centrosome-context stain — no dedicated centrosome
marker), actin, and Hoechst/DNA. Each grid panel is a cross-group comparison.

This is the cross-experiment variant of insert_ctl_granule_nuc_summary_20260617_
slides.py (which builds the single-experiment June-17 deck). It reuses that deck's
curated FAMILIES verbatim (all grid-panel stems resolve in the crossexp compile)
and appends a granule geometry-density ("geom_density") section at the end,
modeled on the MT/vim/SUN geom-density decks.

The pairwise families and the appended geom_density families point at the crossexp
compile's own pairwise_plots/ and geom_density/ subfolders. Those are not present
yet in this compile, so those slides render "no data yet" and auto-fill once the
crossexp compile is regenerated with those outputs. The geom_density panel names
are taken from the single-experiment reference compile
(.../CTL_Glass_nuc_MT_granules_20260617_20260718/geom_density/) — the same MATLAB
pipeline, so the crossexp run emits the same Lamp1_* / pooled stems.

A FAMILIES entry is (stem, title) for a single panel, or ([(stem, sublabel), ...],
title) to place related panels side by side on one slide with a sublabel above
each. A plain metric stem maps to grid_panels/<stem>_grid.png; a stem ending in
".png" is taken relative to the compilation ROOT (used for the pairwise and
geom_density panels).

Self-contained: builds a blank deck (no template .pptx). Missing panels render
"no data yet" rather than failing. A previous deck is backed up before overwrite.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_ctl_granule_nuc_summary_crossexp_20260719_slides.py
    # dry run (print planned families/titles, build nothing):
    conda run -n PPT_editing python examples_and_configs/insert_ctl_granule_nuc_summary_crossexp_20260719_slides.py --list
"""

import os
import sys
import uuid
from pathlib import Path
from xml.sax.saxutils import escape

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation, safe_path, path_exists  # noqa: E402

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
COMPILE_ROOT = "L:/FF/Nucleus_granules/CTL_fixed/compiled_results"
OUT_DIR = "K:/FF/PPT/PPT_autogeneration/CTL_Glass_Nucleus_Centrosome/CTL_fixed_LAMP1"

# Variant registry — two compiles that share the same structure and this same deck
# layout, written to separate .pptx files:
#   base    = the standard cross-experiment compile (default)
#   ncdn05  = the nucleus-centrosome-distance filter variant (--ncdn05)
# Select with the --ncdn05 flag; otherwise the base variant is built.
_VARIANTS = {
    "base": {
        "root": "CTL_nuc_MT_granules_20260617_20260716_crossexp_20260720",
        "out": "CTL_fixed_granule_nuc_summary_crossexp_20260617_20260716.pptx",
        "note": "",
    },
    "ncdn05": {
        "root": "CTL_nuc_MT_granules_20260617_20260716_crossexp_ncdn05_20260724",
        "out": "CTL_fixed_granule_nuc_summary_crossexp_ncdn05_20260617_20260716.pptx",
        "note": "  ·  ncdn05 nucleus-centrosome-distance filter variant",
    },
}
VARIANT = "ncdn05" if "--ncdn05" in sys.argv else "base"
_V = _VARIANTS[VARIANT]

ROOT = Path(COMPILE_ROOT) / _V["root"]
GRID_DIR = ROOT / "grid_panels"
CELL_COUNTS_PNG = ROOT / "cell_counts_barplot.png"   # context slide (optional)
OUTPUT_PATH = Path(OUT_DIR) / _V["out"]

# Per-experiment stacked-rows violins (June 17 row: 3/12 min; July 16 row: 3/5/12 min)
# — the desired layout, matching the LatA per-experiment-rows pattern. The single-axis
# "_grid.png" variant (all groups on one axis) also exists but is not used here.
GRID_SUFFIX = "_per_experiment_grid.png"

# Compile date parsed from the dated ROOT folder (…_YYYYMMDD), shown in the footer.
_d = ROOT.name.rsplit("_", 1)[-1]          # e.g. "20260720"
COMPILE_DATE = "{}-{}-{}".format(_d[:4], _d[4:6], _d[6:8])

DECK_TITLE = "Granule polarization and nuclear morphology in activated CTLs"
# Cross-experiment: two datasets pooled; the cell-counts slide carries the per-group
# n's (June 17: 3/12 min; July 16: 3/5/12 min), so they are not hardcoded here.
DECK_SUBTITLE = (
    "Cross-experiment (Jun 17 + Jul 16 2026), pooled  ·  3 / 12 min "
    "(Jul 16 also 5 min)  ·  αCD3/ICAM1/3SI, glass  ·  LAMP1 / MT / actin / DNA  ·  "
    "compiled " + COMPILE_DATE + _V["note"]
)

# ---------------------------------------------------------------------------
# Curated metrics, grouped into families (divider slide per family). Each entry
# is (grid-panel stem, slide title) for one panel, or ([(stem, sublabel), ...],
# slide title) to show related panels side by side on one slide. The stem +
# GRID_SUFFIX is the PNG under grid_panels/. Titles are for navigation; each plot
# carries its own authoritative y-axis label and per-group x-labels.
# ---------------------------------------------------------------------------
FAMILIES = [
    ("Cell and nuclear spreading", [
        ([("nuc_aspect_ratio", "nucleus"),
          ("actin_deform_ratio", "cell")], "Nuclear and cell aspect ratio"),
        ([("actin_bottom_mask_area", "synapse area"),
          ("nuc_broadest_slice_area", "nuclear broadest slice")],
         "Synapse and nuclear broadest-slice area"),
    ]),
    ("Granules — polarization and synapse delivery", [
        ([("centrosome_center_z_rel_bottom_actin_plane", "centrosome"),
          ("Lamp1_zCOF_cell_bottom_distance", "granules")],
         "Centrosome and granule distance to synapse"),
        ([("Lamp1_z50_cell_bottom_distance", "z₅₀"),
          ("Lamp1_z75_cell_bottom_distance", "z₇₅")], "Granule z₅₀ / z₇₅ distance to synapse"),
        ("Lamp1_synapse_g_ave",              "Granule clustering at synapse (g)"),
        ("Lamp1_synapse_inner_outer_ratio",  "Granule synapse inner/outer ratio"),
    ]),
    ("Granules — dispersion", [
        ([("Lamp1_FDD_3D", "3D"),
          ("Lamp1_z_FDD", "axial (z)")], "Granule dispersion (FDD): 3D vs axial"),
        ([("Lamp1_FDD_3D_rel_cent", "3D"),
          ("Lamp1_z_FDD_rel_cent", "axial (z)")], "Granule dispersion rel. centrosome: 3D vs axial"),
    ]),
    ("Granules — signal and centrosome localization", [
        ("Lamp1_total_sig",   "Granule total signal (whole cell)"),
        ("Lamp1_peak_sig",    "Granule peak signal"),
        ("Lamp1_synapse_MFI", "Granule MFI at synapse"),
        ([("Lamp1_synapse_total_sig", "single slice"),
          ("Lamp1_synapse_total_sig_3mip", "3-slice MIP")],
         "Total granule signal at synapse"),
        ([("Lamp1_frac_around_cent_1um", "1 μm"),
          ("Lamp1_frac_around_cent_2um", "2 μm"),
          ("Lamp1_frac_around_cent_3um", "3 μm")], "Granule fraction around centrosome"),
        ([("Lamp1_MFI_around_cent_1um", "1 μm"),
          ("Lamp1_MFI_around_cent_2um", "2 μm"),
          ("Lamp1_MFI_around_cent_3um", "3 μm")], "Granule MFI around centrosome"),
    ]),
    ("Granules — perinuclear and centrosome enrichment", [
        ([("Lamp1_all_perinuc_MFI", "MFI"),
          ("Lamp1_perinuc_sig_fraction", "signal fraction")],
         "Granule perinuclear MFI and signal fraction"),
        ([("Lamp1_frac_perinuc_within_1_um_cent", "1 μm cent"),
          ("Lamp1_frac_perinuc_within_2_um_cent", "2 μm cent")],
         "Granule perinuclear fraction near centrosome"),
        ([("Lamp1_cyto_in_nuc_hull_MFI", "MFI"),
          ("Lamp1_cyto_in_nuc_hull_sig_fraction", "signal fraction")],
         "Granule MFI and fraction in cytoplasm within nuclear hull"),
        ("Lamp1_frac_in_nuc_convex_hull", "Granule fraction in nuclear convex hull"),
        ([("Lamp1_enrichment_within_half_um_nuc_2_um_cent", "granule"),
          ("MT_enrichment_within_half_um_nuc_2_um_cent", "MT")],
         "Enrichment near centrosome (0.5 μm of nucleus, 2 μm of cent)"),
    ]),
    # Granule polarization relative to the MTOC (2 μm sphere). Each is one figure
    # per slide with a descriptive caption (item = (stem, title, caption)).
    ("Granules around the MTOC (1 & 2 μm)", [
        ([("Lamp1_cent_facing_mass_ratio_1um", "1 μm"),
          ("Lamp1_cent_facing_mass_ratio_2um", "2 μm")],
         "Granules are not biased to the nucleus-facing side of the centrosome",
         "Ratio of granule intensity per available cytoplasm, nucleus-facing vs. far "
         "hemisphere of the MTOC (nucleus excluded), for 1 μm and 2 μm spheres. Dashed "
         "line at 1 = no side preference. June 17 and July 16, 2026; 3/5/12 min. The ratio "
         "sits at ~1 in every timepoint, both experiments and both radii (all n.s.) — no "
         "facing preference, and the null reproduces across replicates."),
        ([("Lamp1_cent_nuc_dist_ratio_1um", "1 μm"),
          ("Lamp1_cent_nuc_dist_ratio_2um", "2 μm")],
         "Granules sit slightly closer to the nucleus in July, but not June",
         "Intensity-weighted mean granule distance to the nucleus, over the chance "
         "distance across available cytoplasm, for 1 μm and 2 μm spheres. Dashed line at "
         "1 = chance; <1 = closer than chance. July sits below chance (~0.92 at 2 μm, more "
         "modest at 1 μm); June sits at ~1 (n.s.). The effect does not reproduce between "
         "experiments and is not explained by centrosome–nucleus distance (near-identical "
         "between the two), so it is a between-experiment difference, not a claimable result."),
    ]),
    ("Centrosome ↔ nucleus", [
        ("nuc_cent_closest_dist",            "Nucleus-centrosome closest distance"),
        ("cent_nuc_norm_dist_sphere_rad",    "Centrosome-to-nuclear-centroid distance (norm. to equiv sphere radius)"),
        ("centrosome_dist_deepest_real_avg_periphery_ratio", "Centrosome distance to deepest invag vs avg periphery ratio"),
        ([("centrosome_r_norm_bottom_plane_from_MT", "synapse plane"),
          ("centrosome_r_norm_MIP_from_MT", "MIP")],
         "Centrosome radial position (0 = center, 1 = periphery)"),
    ]),
    ("Nuclear deformation and invaginations", [
        ("chull_max_D",                       "Max invag depth over full nucleus"),
        ("chull_max_D_by_cent",               "Invagination depth near centrosome"),
        ("chull_mean_D_cent_global_ratio",    "Centrosomal Invagination Index (global)"),
        ([("C_min_F_mean_by_cent", "min principal"),
          ("C_mean_F_mean_by_cent", "mean")],
         "Nuclear surface curvature near centrosome"),
        ("deepest_invag_fraction_chull_volume", "Deepest invag: frac of convex hull volume"),
        ("deepest_region_periph_ratio_025um", "DNA levels near invag"),
        ([("invag_by_cent_centroid_z_syn_from_MT", "centroid"),
          ("invag_by_cent_tip_z_syn_from_MT", "tip")],
         "Invagination region (near centrosome): height above synapse"),
    ]),
    ("Invagination orientation", [
        ("avg_normal_angle_adaptive_region_growth",         "Deepest invag orientation"),
        ("avg_normal_angle_adaptive_region_growth_by_cent", "Invag orientation (adaptive) near centrosome"),
        ("avg_normal_angle_by_cent",                        "Invag orientation near centrosome"),
    ]),
    ("Nuclear morphology", [
        ("nuc_solidity",    "Nuclear solidity"),
        ("nuc_volume_mesh", "Nuclear volume"),
        ("nuc_SA_mesh",     "Nuclear surface area"),
    ]),
    ("Microtubules (β-tubulin)", [
        ("MT_frac_in_nuc_convex_hull",                  "MT fraction in nuclear convex hull"),
        ("MT_frac_around_cent_2um",  "MT fraction around centrosome (2 μm)"),
        ("MT_MFI_around_cent_2um",   "MT MFI around centrosome (2 μm)"),
    ]),
    ("Actin — levels and localization", [
        ("actin_total_sig", "Total actin signal (whole cell)"),
        ([("actin_bottom_MFI", "single slice"),
          ("actin_bottom_MFI_3mip", "3-slice MIP")], "Actin MFI at synapse"),
        ([("actin_bottom_total_sig", "single slice"),
          ("actin_bottom_total_sig_3mip", "3-slice MIP")], "Total actin signal at synapse"),
        ("actin_bottom_inner_outer_ratio", "Actin synapse inner/outer ratio"),
        ([("actin_MFI_around_cent_1um", "1 μm"),
          ("actin_MFI_around_cent_2um", "2 μm")], "Actin MFI around centrosome"),
        ([("actin_frac_around_cent_1um", "1 μm"),
          ("actin_frac_around_cent_2um", "2 μm")], "Actin fraction around centrosome"),
    ]),
    # --- Appendix: additional context metrics ---------------------------------
    ("Cell and nuclear flattening (context)", [
        ("actin_height",        "Cell height"),
        ("nuc_height",          "Nuclear height"),
        ("nuc_centroid_z",      "Nuclear centroid height above synapse"),
        ("nuc_mesh_sphericity", "Nuclear sphericity"),
        ("actin_MIP_circularity", "Cell footprint circularity"),
    ]),
    ("Chromatin / DNA distribution", [
        ("nuc_all_CV",          "DNA intensity CV (heterogeneity)"),
        ("nuc_all_prop_gr_2med", "DNA fraction > 2× median (bright foci)"),
        ("nuc_all_skewness",    "DNA intensity skewness"),
        ("nuc_all_norm_entropy", "DNA distribution normalized entropy"),
    ]),
]

# Pairwise (scatter) section drawn from the compilation's pairwise_plots/ output.
# In the CROSSEXP compile these are nested per experiment:
# pairwise_plots/<Experiment>/<Suite>/<X>_VS_<Y>.png (there is no pooled cross-
# experiment scatter), so each family produces one slide per experiment.
_PW = "pairwise_plots"
_PW_EXPTS = [("June_17,_2026", "Jun 17"), ("July_16,_2026", "Jul 16")]

# Suite GranuleDelivery_vs_InvagZ: the centrosome-region invagination-height metric
# (X) vs the granule distance-to-synapse metrics (Y), the Y's side by side. In the
# crossexp compile the X-metric is the adaptive-region centroid height (renamed
# from the single-exp invag_by_cent_* stem).
_INVAGZ_X = "adaptive_region_by_cent_centroid_z_syn_from_MT"
_GRANULE_DIST_YS = [
    ("Lamp1_zCOF_cell_bottom_distance", "zCOF"),
    ("Lamp1_z50_cell_bottom_distance",  "z₅₀"),
    ("Lamp1_z75_cell_bottom_distance",  "z₇₅"),
]
FAMILIES.append((
    "Pairwise — invagination-region height vs granule delivery",
    [([("{}/{}/GranuleDelivery_vs_InvagZ/{}_VS_{}.png".format(
          _PW, _exdir, _INVAGZ_X, _ystem), _ylab)
       for _ystem, _ylab in _GRANULE_DIST_YS],
      "Granule distance to synapse vs invag-region centroid height — {}".format(_exlab))
     for _exdir, _exlab in _PW_EXPTS],
))

# Suite GranuleCent_vs_InvagDepth: granule clustering at the centrosome (Y) vs
# invagination depth near the centrosome (X = chull_max_D_by_cent_from_MT), one
# slide per (experiment × granule-metric group).
def _pw_centdepth(exdir, ys):
    """Pairwise panels: chull_max_D_by_cent (X) vs each granule y-metric (Y)."""
    return [("{}/{}/GranuleCent_vs_InvagDepth/chull_max_D_by_cent_from_MT_VS_{}.png".format(
                _PW, exdir, _y), _lab)
            for _y, _lab in ys]


FAMILIES.append((
    "Pairwise — granule clustering at centrosome vs invagination depth",
    [item
     for _exdir, _exlab in _PW_EXPTS
     for item in [
        (_pw_centdepth(_exdir, [("Lamp1_FDD_3D_rel_cent", "avg dist to cent"),
                                ("Lamp1_z_FDD_rel_cent",  "axial dist to cent")]),
         "Granule distance to centrosome vs invag depth — {}".format(_exlab)),
        (_pw_centdepth(_exdir, [("Lamp1_MFI_around_cent_1um", "1 μm"),
                                ("Lamp1_MFI_around_cent_2um", "2 μm")]),
         "Granule MFI around centrosome vs invag depth — {}".format(_exlab)),
        (_pw_centdepth(_exdir, [("Lamp1_frac_around_cent_1um", "1 μm"),
                                ("Lamp1_frac_around_cent_2um", "2 μm")]),
         "Granule fraction around centrosome vs invag depth — {}".format(_exlab)),
     ]],
))

# Suite CentRadialPos_vs_InvagOrient: centrosome radial position (X) vs centrosome-
# facing invagination orientation (Y), one slide per experiment.
FAMILIES.append((
    "Pairwise — centrosome radial position vs invagination orientation",
    [("{}/{}/CentRadialPos_vs_InvagOrient/"
      "centrosome_r_norm_bottom_plane_from_MT_VS_avg_normal_angle_by_cent_from_MT.png".format(
          _PW, _exdir),
      "Invag orientation vs centrosome radial position — {}".format(_exlab))
     for _exdir, _exlab in _PW_EXPTS],
))

# ---------------------------------------------------------------------------
# Granule geometry-density ("geom_density") section — appended at the end.
# Modeled on the MT/vim/SUN geom-density decks, curated to the granule (Lamp1_*)
# panels. Density / enrichment / correlation of granules as a function of nuclear-
# envelope geometry (invagination depth = hull-boundary distance, min curvature,
# mean curvature), sampled in the perinuc 0.5 μm shell (0.5 μm outside the NE).
# Uses the "pooled" (cross-condition) variants so a single plot carries all groups.
# Panels live under the crossexp compile's geom_density/ (not present yet → these
# render "no data yet" until the crossexp compile is regenerated with geom_density).
# ---------------------------------------------------------------------------
_GD_PROF_S = "geom_density/profiles/singles"
_GD_ENR = "geom_density/enrichment"

# Granule geometry-density section, modeled on the vim/MT geom-density comparison
# decks. One slide per metric, with the two experiments side by side (1×2): the
# first experiment (June 17, 3/12 min) next to the second (July 16, 3/5/12 min), so
# both experiments are visible together on every slide. Density profiles use the
# full-range OVERLAY line (perinuc 0.5 μm); enrichment/correlation use the per-cell
# shell violins. Panels come from geom_density/ (the curated Lamp1_loc_wrto_invag/
# set is pooled-only, so it can't provide the per-experiment split).
#   (label, enrichment token [comma], OVERLAY-singles token [double underscore])
_GD_EXP = [
    ("June 17", "June_17,_2026", "June_17__2026"),
    ("July 16", "July_16,_2026", "July_16__2026"),
]


def _gd_profile_pair(geom, channel="Lamp1"):
    """June|July OVERLAY line profile panels for one geometry (hulldist/mincurv/meancurv).
    channel='Lamp1' (granules) or 'MT' (microtubules)."""
    return [("{}/{}_geomdens_{}_perinuc05_OVERLAY_line_{}.png".format(_GD_PROF_S, channel, geom, ov), lab)
            for lab, _enr, ov in _GD_EXP]


def _gd_shell_pair(stem):
    """June|July shell-violin panels for one enrichment/correlation metric stem."""
    return [("{}/{}_{}.png".format(_GD_ENR, stem, enr), lab)
            for lab, enr, _ov in _GD_EXP]


# Intersperse granule (LAMP1) and MT slides for each metric — LAMP1 slide
# immediately followed by MT slide with the equivalent metric.
def _lamp_mt_profile_slides(geom, title):
    return [
        (_gd_profile_pair(geom, "Lamp1"), "Granule "  + title),
        (_gd_profile_pair(geom, "MT"),    "MT "        + title),
    ]


def _lamp_mt_shell_slides(stem_no_ch, title_lamp, title_mt):
    """stem_no_ch is the metric stem after the 'Lamp1_'/'MT_' prefix, e.g.
    'hulldist_gt0.5um_shells'. Returns a LAMP1 slide immediately followed by
    the MT-analog slide."""
    return [
        (_gd_shell_pair("Lamp1_" + stem_no_ch), title_lamp),
        (_gd_shell_pair("MT_"    + stem_no_ch), title_mt),
    ]


def _dp(title):  # "density" title suffix
    return "density " + title


FAMILIES.append((
    "Granule + MT density vs NE geometry (June 17 vs July 16)",
    _lamp_mt_profile_slides("hulldist",  "density vs invagination depth (perinuc 0.5 μm)") +
    _lamp_mt_profile_slides("mincurv",   "density vs min curvature (perinuc 0.5 μm)") +
    _lamp_mt_profile_slides("meancurv",  "density vs mean curvature (perinuc 0.5 μm)"),
))

FAMILIES.append((
    "Granule + MT enrichment vs NE geometry (June 17 vs July 16)",
    _lamp_mt_shell_slides("hulldist_gt0.5um_shells",
        "Granule levels in deep invaginations (hull dist > 0.5 μm)",
        "MT levels in deep invaginations (hull dist > 0.5 μm)") +
    _lamp_mt_shell_slides("hulldist_gt1.0um_shells",
        "Granule levels in deep invaginations (hull dist > 1.0 μm)",
        "MT levels in deep invaginations (hull dist > 1.0 μm)") +
    _lamp_mt_shell_slides("mincurv_lt0_shells",
        "Granule levels on concave surface (min curv < 0)",
        "MT levels on concave surface (min curv < 0)") +
    _lamp_mt_shell_slides("mincurv_ltm0.25_shells",
        "Granule levels on strongly concave surface (min curv < −0.25)",
        "MT levels on strongly concave surface (min curv < −0.25)") +
    _lamp_mt_shell_slides("meancurv_lt0_shells",
        "Granule levels on concave surface (mean curv < 0)",
        "MT levels on concave surface (mean curv < 0)"),
))

FAMILIES.append((
    "Granule + MT correlation vs NE geometry (June 17 vs July 16)",
    _lamp_mt_shell_slides("corr_hulldist_shells",
        "Granule per-cell correlation vs hull-boundary distance",
        "MT per-cell correlation vs hull-boundary distance") +
    _lamp_mt_shell_slides("corr_mincurv_shells",
        "Granule per-cell correlation vs min curvature",
        "MT per-cell correlation vs min curvature") +
    _lamp_mt_shell_slides("corr_meancurv_shells",
        "Granule per-cell correlation vs mean curvature",
        "MT per-cell correlation vs mean curvature") +
    _lamp_mt_shell_slides("deepcorr_mincurv_shells",
        "Granule correlation vs min curvature within deep invaginations",
        "MT correlation vs min curvature within deep invaginations") +
    _lamp_mt_shell_slides("deepcorr_meancurv_shells",
        "Granule correlation vs mean curvature within deep invaginations",
        "MT correlation vs mean curvature within deep invaginations"),
))

# ---------------------------------------------------------------------------
# Colors / layout (matches the bleb/noco/vimkd summary decks)
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.10

TITLE_LEFT = MARGIN
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.55
TITLE_FONT_PT = 28

IMG_LEFT = MARGIN
IMG_TOP = 0.66
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = 6.36

# Side-by-side (multi-panel) slides: a sublabel band above the figures, and a
# gap between columns.
COL_GAP = 0.15
SUBLABEL_H = 0.30
SUBLABEL_GAP = 0.04
SUBLABEL_FONT_PT = 16

FOOTER_LEFT = MARGIN
FOOTER_TOP = 7.06
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.40
FOOTER_FONT_PT = 9

# Optional caption block (between image and footer) for slides that carry a
# descriptive paragraph (e.g. the MTOC granule figures).
CAPTION_BLOCK_H = 1.25
CAPTION_FONT_PT = 12

# ---------------------------------------------------------------------------


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    """Place an image inside (left, top, w, h), preserving aspect ratio and
    centering on whichever dimension ends up smaller than the box. The grid
    panels are near-square, so they fit to height and center horizontally."""
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def rel_footer(path):
    try:
        rel = path.relative_to(ROOT).as_posix()
    except ValueError:
        rel = path.as_posix()
    return "{} / {}".format(COMPILE_DATE, rel)


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.7, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=40, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.1, SLIDE_W - 2 * MARGIN, 1.0,
                font_pt=18, color=BLACK, italic=True)


def build_divider_slide(prs, family_name):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, family_name, MARGIN, 3.1, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=44, color=BLACK, bold=True)



def _place_missing(slide, left, top, width):
    add_textbox(slide, "no data yet", left, top + IMG_BOX_H / 2 - 0.2,
                width, 0.4, font_pt=18, color=BLACK)


def build_slide(prs, title_text, panels, footer_text, caption=None):
    """Title + one or more panels (aspect preserved) + optional caption + footer.
    `panels` is a list of (path, sublabel); a single entry fills the image box,
    multiple entries are laid out in side-by-side columns with a sublabel above
    each. An optional multi-line `caption` sits between the image and the footer
    (shrinking the image box). Returns the list of missing panel paths."""
    slide = _new_slide(prs)
    add_textbox(slide, title_text, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title_text), color=BLACK, bold=True)

    cap_h = CAPTION_BLOCK_H if caption else 0.0
    box_h = IMG_BOX_H - cap_h

    missing = []
    if len(panels) == 1:
        path, _ = panels[0]
        if path_exists(path):
            add_image_in_box(slide, safe_path(path), IMG_LEFT, IMG_TOP, IMG_BOX_W, box_h)
        else:
            _place_missing(slide, IMG_LEFT, IMG_TOP, IMG_BOX_W)
            missing.append(path)
    else:
        n = len(panels)
        col_w = (IMG_BOX_W - (n - 1) * COL_GAP) / n
        band_top = IMG_TOP + SUBLABEL_H + SUBLABEL_GAP
        band_h = box_h - SUBLABEL_H - SUBLABEL_GAP
        for i, (path, sublabel) in enumerate(panels):
            left = IMG_LEFT + i * (col_w + COL_GAP)
            sub_box = add_textbox(slide, sublabel or "", left, IMG_TOP, col_w, SUBLABEL_H,
                                  font_pt=SUBLABEL_FONT_PT, color=BLACK, bold=True)
            # Bottom-align so the label sits just above its plot, not floating high.
            sub_box.text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
            if path_exists(path):
                add_image_in_box(slide, safe_path(path), left, band_top, col_w, band_h)
            else:
                _place_missing(slide, left, band_top, col_w)
                missing.append(path)

    if caption:
        add_textbox(slide, caption, MARGIN, IMG_TOP + box_h + 0.04, IMG_BOX_W,
                    cap_h - 0.04, font_pt=CAPTION_FONT_PT, color=BLACK, align=PP_ALIGN.LEFT)

    add_textbox(slide, footer_text, FOOTER_LEFT, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=BLACK)
    return missing


def _panel_path(stem):
    """Resolve a panel reference to a file path. A plain metric stem maps to
    GRID_DIR/<stem>_grid.png; a reference ending in '.png' is taken relative to
    the compilation ROOT (used for pairwise_plots/ and geom_density/ panels)."""
    if stem.endswith(".png"):
        return ROOT / stem
    return GRID_DIR / (stem + GRID_SUFFIX)


def entry_panels(entry_stem):
    """Normalize a FAMILIES entry's stem field into a list of (path, sublabel)."""
    if isinstance(entry_stem, str):
        return [(_panel_path(entry_stem), None)]
    return [(_panel_path(stem), sublabel) for stem, sublabel in entry_stem]


def add_sections(prs, section_spec):
    """Add native PowerPoint sections (the collapsible, named groups shown as
    tabs in the slide navigator / slide sorter). python-pptx has no API for
    this, so inject the p14:sectionLst extension into the presentation part.

    `section_spec` is an ordered list of (section_name, n_slides); the counts
    must sum to the total slide count in build order, and the first section
    must contain the first slide (a PowerPoint requirement)."""
    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    slide_ids = [sldId.get("id") for sldId in prs.slides._sldIdLst]

    parts = ['<p:extLst xmlns:p="{}">'.format(P),
             '<p:ext uri="{{521415D9-36F7-43E2-AB2F-B90AF26B5E84}}">',
             '<p14:sectionLst xmlns:p14="{}">'.format(P14)]
    i = 0
    for name, count in section_spec:
        parts.append('<p14:section name="{}" id="{{{}}}">'.format(
            escape(name), str(uuid.uuid4()).upper()))
        parts.append('<p14:sldIdLst>')
        for sid in slide_ids[i:i + count]:
            parts.append('<p14:sldId id="{}"/>'.format(sid))
        parts.append('</p14:sldIdLst></p14:section>')
        i += count
    parts.append('</p14:sectionLst></p:ext></p:extLst>')

    # The sldIdLst's parent is the <p:presentation> element; extLst goes last.
    prs.slides._sldIdLst.getparent().append(parse_xml("".join(parts)))


def main():
    list_only = "--list" in sys.argv

    n_metrics = sum(len(items) for _, items in FAMILIES)
    # title + (cell-counts if present) + per-family (divider + metric slides)
    est_slides = 1 + (1 if CELL_COUNTS_PNG.exists() else 0) + \
        sum(1 + len(items) for _, items in FAMILIES)

    print("Source: {}".format(GRID_DIR))
    print("{} metric slides across {} families, est. {} slides\n".format(
        n_metrics, len(FAMILIES), est_slides))

    if list_only:
        for fam, items in FAMILIES:
            print("=== {} ({}) ===".format(fam, len(items)))
            for entry in items:
                entry_stem, title = entry[0], entry[1]
                panels = entry_panels(entry_stem)
                flags = " ".join(
                    "{}:{}".format(sub or "-", "OK" if path_exists(p) else "MISS")
                    for p, sub in panels)
                print("  [{}] {:<50s} {}".format(
                    "OK " if all(path_exists(p) for p, _ in panels) else "MISS",
                    title, flags))
            print("")
        return

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)

    if CELL_COUNTS_PNG.exists():
        build_slide(prs, "Cell counts (per experiment and timepoint)",
                    [(CELL_COUNTS_PNG, None)], rel_footer(CELL_COUNTS_PNG))
    else:
        print("Note: {} not found - skipping cell-counts slide.\n".format(
            CELL_COUNTS_PNG.name))

    # Native PowerPoint sections: an intro section (title + cell counts) then
    # one section per family (its divider slide + metric slides).
    section_spec = [("Overview", len(prs.slides._sldIdLst))]

    missing = []
    for fam, items in FAMILIES:
        build_divider_slide(prs, fam)
        print("=== {} ===".format(fam))
        for entry in items:
            entry_stem, title = entry[0], entry[1]
            caption = entry[2] if len(entry) > 2 else None
            panels = entry_panels(entry_stem)
            footer = " | ".join(rel_footer(p) for p, _ in panels)
            miss = build_slide(prs, title, panels, footer, caption)
            status = "OK" if not miss else "MISSING"
            print("  [{}] {} -> {!r}".format(
                status, ", ".join(p.name for p, _ in panels), title))
            missing.extend(p.name for p in miss)
        section_spec.append((fam, 1 + len(items)))
        print("")

    add_sections(prs, section_spec)

    if OUTPUT_PATH.exists():
        backup_dir = OUTPUT_PATH.parent / "backups"
        created = backup_presentation(str(OUTPUT_PATH), backup_base=str(backup_dir))
        if created:
            print("Backed up previous deck to: {}\n".format(backup_dir))

    prs.save(str(OUTPUT_PATH))
    total = len(prs.slides._sldIdLst)
    print("Done. {} metric slides, {} slides written to:\n  {}".format(
        n_metrics, total, OUTPUT_PATH))
    if missing:
        print("\nSkipped {} missing panel(s) (not on disk):".format(len(missing)))
        for m in missing:
            print("  - {}".format(m))
    else:
        print("\nAll curated panels found - no missing items.")


if __name__ == "__main__":
    main()
