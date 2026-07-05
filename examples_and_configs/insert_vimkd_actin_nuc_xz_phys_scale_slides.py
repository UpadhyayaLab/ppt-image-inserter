"""
insert_vimkd_actin_nuc_xz_phys_scale_slides.py

siCtrl (left) vs siVim (right) ACTIN + NUCLEUS x-z MIP physical-scale montage deck
spanning ALL the fixed-Jurkat VimentinKD experiments — one slide per experiment.

Why actin+nuc (not vimentin): vimentin is the knockdown *target* and is imaged in only
2 of the 10 VimentinKD experiments (the KD-validation runs). Actin, by contrast, is
imaged as the cell-bottom / spreading marker in the Vim, MT, AcTub and pMLC runs, so its
x-z MIP exists for those — `physical_scale_images/actin_nuc_xz_nolines/`. This deck uses
that one combo to give a uniform actin+nucleus axial view, with the siCtrl-vs-siVim
columns carrying the knockdown comparison. (The vimentin x-z MIP, where it exists, is a
separate 2-slide deck: insert_vimkd_vim_nuc_xz_slides.py.)

PERICENTRIN RUNS INCLUDED — their actin_nuc_xz is genuine actin. The 3 pericentrin runs
(20220406/20220520/20220628) are C1=pericentrin, C2=actin. At first glance their
actin_nuc_xz looks like a centrosome punctum, but that was a misread: the channel role is
assigned BY NAME (config/load_config.m via config/channel_registry.m — Actin→'actin',
Pericentrin→'centrosome'), so actin_channel_prefix is the Actin def's prefix (C2-cell), not
a swapped centrosome. The single-channel montages confirm it: actin_xz (C2) is faint diffuse
actin + a bright basal actin focus (likely synaptic actin in these αCD3 cells), while cent_xz
(C1) is sparse scattered puncta — different channels. So the panel IS actin; it just reads as
a "dot" because the actin is weaker/more focal than the crisp cortical sheet in the
Vim/AcTub/pMLC runs and the bright cyan nucleus swamps the diffuse part in the merge.

The experiments span 4 staining assays and 3 drives — AcTub (L:), MT / pMLC (J:),
Vim / Pericentrin (M:) — so chan_sub and condition-folder names differ per experiment. Each
slide title is DATE-ONLY (the staining assay is deliberately not named, so the panel reads
purely as actin+nucleus). Two runs share 04/06/2022 (a Vim and a pericentrin acquisition),
so those two slides carry the identical title and simply sit adjacent — accepted by design.

Coloring: actin=RED, nucleus=cyan (combo `actin_nuc_xz_planes_nolines`); the pericentrin
runs also carry a magenta centrosome overlay from that renderer. (No blue exists in the
XZ palette — cyan is the nucleus color in every red-actin render.)

Same machinery as the sibling phys-scale decks (noco/bleb/LatA, vimkd MT/pericentrin):
long-path-safe enumeration, group-major ordering, deck-wide per-scale-group PPI pinning
(104 px = 5 μm scalebar invariant), one representative montage chunk per condition,
experiments whose montages aren't present are skipped / shown "(missing)".

Usage:
    python examples_and_configs/insert_vimkd_actin_nuc_xz_phys_scale_slides.py
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/VimentinKD/"
    "VimKD_Jurkats_siCtrl_vs_siVim_actin_nuc_xz_phys_scale_montages.pptx"
)

# Actin (RED) + nucleus (cyan) x-z MIP, physical-scale, no dashed plane lines.
# The `_planes` renderer (save_nuc_actin_planes_xz_mip_physical.m) colors actin RED
# (vs the plain actin_nuc_xz* combo, which uses orange); `_nolines` drops the dashed
# actin/nucleus plane markers. On the pericentrin runs this renderer also overlays the
# centrosome in MAGENTA (it adds I_cent when a centrosome channel exists); the
# Vim/pMLC/AcTub runs have no centrosome, so they're red actin + cyan nucleus only.
# (There is no blue in the XZ color palette — cyan is the nucleus color for all the
# red-actin renders.)
COMBO = "actin_nuc_xz_planes_nolines"

# Drive/experiment roots.
ROOT_ACTUB = "L:/FF/Nucleus_MT/Jurkat_fixed/vimentinKD_tubulin-acetylation_fixed"
ROOT_VIM_M = ("M:/FF/FF_4TB_2_Backup_fullHD/Vimentin_Project_2ndharddrive/"
              "VimentinKD_NucleusData_Fixed")

# One entry per experiment (siCtrl left / siVim right). Fields:
#   tag      — MM/DD/YYYY, used for chronological sort + woven into the title.
#   marker   — staining assay, names the experiment in the slide title (the actin+nuc
#              panel is the same kind of view; this just gives provenance).
#   root     — dir that directly contains the condition folders.
#   chan_sub — intermediate path before prog_fixed_cells (differs per drive/pipeline).
#   tp_label — activation ("αCD3" or "" for the basal runs).
#   left/right — (condition_folder, display_label) for the siCtrl / siVim columns.
EXPERIMENTS = [
    {
        "tag": "04/06/2022", "marker": "Vim", "tp_label": "",
        "root": ROOT_VIM_M + "/20220406 - Vimentin siRNA Experiments/transfection2_48h",
        "chan_sub": "cells/individual-channels",
        "left": ("Control", "siCtrl"), "right": ("siVIM", "siVim"),
    },
    {
        "tag": "05/04/2022", "marker": "Vim", "tp_label": "αCD3",
        "root": (ROOT_VIM_M + "/20220504 - Vimentin siRNA Experiments/"
                 "Vimentin knockdown validation - 48h"),
        "chan_sub": "cells/individual-channels",
        "left": ("siCtrl - aCD3", "siCtrl"), "right": ("siVim - aCD3", "siVim"),
    },
    {
        "tag": "05/23/2022", "marker": "MT", "tp_label": "αCD3",
        "root": "J:/FF/vim_data/from_AP_04122024/20220523 - Vim KD bTub AcTub",
        "chan_sub": "cells/individual-channels",
        "left": ("siCtrl", "siCtrl"), "right": ("siVim", "siVim"),
    },
    {
        "tag": "08/04/2023", "marker": "pMLC", "tp_label": "αCD3",
        "root": "J:/FF/vim_data/PMLC/20230804_KN_pMLC",
        "chan_sub": "cells/individual_channels",
        "left": ("W1_CD3_siCtrl_7min_R-Phalloidin_488-Vim_Hoechst", "siCtrl"),
        "right": ("W2_CD3_siVim_7min_R-Phalloidin_488-Vim_Hoechst", "siVim"),
    },
    {
        "tag": "10/24/2023", "marker": "AcTub", "tp_label": "αCD3",
        "root": ROOT_ACTUB + "/20231024_MG_AcTub",
        "chan_sub": "channels",
        "left": ("siCtrl_aCD3_640BetaTub_535Actin_488AcTub_405Hoechst", "siCtrl"),
        "right": ("siVim_aCD3_640BetaTub_535Actin_488AcTub_405Hoechst", "siVim"),
    },
    {
        "tag": "01/17/2024", "marker": "AcTub", "tp_label": "αCD3",
        "root": ROOT_ACTUB + "/20240117_MG_AcTub",
        "chan_sub": "channels",
        "left": ("siCtrl_aCD3_E6-1_647BTub_535Actin_488AcTub_Hoechst", "siCtrl"),
        "right": ("siVim_aCD3_E6-1_647BTub_535Actin_488AcTub_Hoechst", "siVim"),
    },
    {
        "tag": "01/29/2024", "marker": "AcTub", "tp_label": "αCD3",
        "root": ROOT_ACTUB + "/20240129_MG_AcTub",
        "chan_sub": "channels",
        "left": ("siCtrl_aCD3_E6-1_647BTub_535Actin_488AcTub_Hoechst", "siCtrl"),
        "right": ("siVim_aCD3_E6-1_647BTub_535Actin_488AcTub_Hoechst", "siVim"),
    },
    # Pericentrin runs (C1=pericentrin, C2=actin). actin_nuc_xz IS genuine actin here —
    # just weaker/more focal than the crisp cortical sheet in the Vim/AcTub/pMLC runs
    # (verified: the C2 actin channel is loaded by name, not a swapped centrosome). Sorted
    # by date, so the 04/06 pericentrin slide sits right after the 04/06 Vim slide — both
    # carry the identical date-only title "Actin + Nuc XZ MIP (04/06/2022)", by design.
    {
        "tag": "04/06/2022", "marker": "Pericentrin", "tp_label": "",
        "root": ROOT_VIM_M + "/20220406 - Vimentin siRNA Experiments/Pericentrin",
        "chan_sub": "cells/individual-channels",
        "left": ("siCtrl", "siCtrl"), "right": ("siVim", "siVim"),
    },
    {
        "tag": "05/20/2022", "marker": "Pericentrin", "tp_label": "αCD3",
        "root": (ROOT_VIM_M + "/20220519 - Vimentin siRNA Experiments/"
                 "20220520 - Vim KD Pericentrin"),
        "chan_sub": "cells/individual-channels",
        "left": ("E6-1_siCtrl_48h_aCD3_Pericentrin", "siCtrl"),
        "right": ("E6-1_siVim_48h_aCD3_Pericentrin", "siVim"),
    },
    {
        "tag": "06/28/2022", "marker": "Pericentrin", "tp_label": "αCD3",
        "root": (ROOT_VIM_M + "/20220628 - Vimentin siRNA Experiments/"
                 "20220628 - Vim KD Pericentrin"),
        "chan_sub": "cells/individual-channels",
        "left": ("E6-1_siCtrl_tf2_48h_aCD3_Pericentrin", "siCtrl"),
        "right": ("E6-1_siVim_tf2_48h_aCD3_Pericentrin", "siVim"),
    },
]

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Slide layout (inches). 13.333 x 7.5 widescreen.
SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

GRID_LEFT = 0.10
GRID_TOP = 0.60
CELL_W = 6.50
CELL_H = SLIDE_H - GRID_TOP - 0.10   # 6.80"
LABEL_H = 0.30
IMG_H = CELL_H - LABEL_H             # 6.50"
LABEL_FONT_PT = 16
COL_GAP = SLIDE_W - 2 * GRID_LEFT - 2 * CELL_W

CELL_POSITIONS = [
    (GRID_LEFT,                    GRID_TOP),
    (GRID_LEFT + CELL_W + COL_GAP, GRID_TOP),
]

# Scalebar invariant for the Jurkat nucleus fixed-cell physical-scale pipeline:
# 104 px = 5 μm (20.8 px/μm) in the rendered PNG, same as the noco/LatA/vimkd decks.
# Deck-wide PPI pinning per scale_group → one physical scale across the deck. The
# constant is diagnostic (printed scalebar-cm); layout uses measured PNG dims + PPI.
SCALEBAR_PX = 104
SCALEBAR_UM = 5
PPUM_SOURCE = SCALEBAR_PX / SCALEBAR_UM

# ---------------------------------------------------------------------------


def _winlong(p) -> str:
    """Win32-safe absolute path; prepends \\\\?\\ so paths near MAX_PATH (260) work
    in os.listdir / Image.open / add_picture. Load-bearing for the longer M:/ paths."""
    s = os.path.abspath(str(p))
    if os.name == "nt" and not s.startswith("\\\\?\\"):
        s = "\\\\?\\" + s.replace("/", "\\")
    return s


def _exists_long(p) -> bool:
    return os.path.exists(_winlong(p))


def montage_dir(root, cond_folder: str, chan_sub: str, combo: str) -> Path:
    """physical_scale_images montages dir for one (experiment, condition, combo)."""
    return (Path(root) / cond_folder / chan_sub /
            "prog_fixed_cells" / "physical_scale_images" / combo / "montages")


def mask_montage_dir(root, cond_folder: str, chan_sub: str) -> Path:
    """Actin synapse-plane segmentation-mask montages (NON-physical-scale):
    prog_fixed_cells/actin/bottom_slice_seg/montages — grayscale actin at the actin
    bottom (synapse) plane with the segmented cell outline drawn in red. Appended as a
    second section at the end of the deck, in its own PPI group (not physical-scale)."""
    return (Path(root) / cond_folder / chan_sub /
            "prog_fixed_cells" / "actin" / "bottom_slice_seg" / "montages")


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _png_dims(path: Path) -> Tuple[int, int]:
    with Image.open(_winlong(path)) as im:
        return im.size


def compute_slide_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_image_at_ppi(slide, image_path: Path, ppi: float,
                     area_left: float, area_top: float,
                     area_w: float, area_h: float):
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    left_in = area_left + (area_w - w_in) / 2
    top_in  = area_top  + (area_h - h_in) / 2
    return slide.shapes.add_picture(
        _winlong(image_path), Inches(left_in), Inches(top_in), width=Inches(w_in),
    )


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _chunk_start_index(p: Path) -> int:
    m = re.match(r"montage_cells_(\d+)", p.name)
    return int(m.group(1)) if m else 0


def list_chunks(montages_dir) -> List[Path]:
    """Montage chunk PNGs sorted by chunk-start index, long-path-safe (os.listdir on
    the \\\\?\\-prefixed dir). [] if absent → experiment simply skipped."""
    d = _winlong(montages_dir)
    if not os.path.isdir(d):
        return []
    names = [f for f in os.listdir(d)
             if f.startswith("montage_cells_") and f.endswith(".png")]
    return sorted((Path(montages_dir) / n for n in names), key=_chunk_start_index)


def find_first_chunks(montages_dir: Path, n: int) -> List[Optional[Path]]:
    chunks = list_chunks(montages_dir)
    return (chunks + [None] * n)[:n]


def build_compare_slide(prs, title_text, left_label, left_img, right_label, right_img,
                        slide_ppi):
    """siCtrl | siVim comparison slide, one montage chunk per side, shared PPI."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BLACK)
    add_textbox(slide, title_text, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=TITLE_FONT_PT, color=WHITE, bold=True)

    missing = []
    for label, img_path, cell_left in (
        (left_label,  left_img,  CELL_POSITIONS[0][0]),
        (right_label, right_img, CELL_POSITIONS[1][0]),
    ):
        add_textbox(slide, label, cell_left, GRID_TOP, CELL_W, LABEL_H,
                    font_pt=LABEL_FONT_PT, color=WHITE, bold=True)
        if img_path is not None and _exists_long(img_path):
            add_image_at_ppi(slide, img_path, slide_ppi,
                             cell_left, GRID_TOP + LABEL_H, CELL_W, IMG_H)
        else:
            add_textbox(slide, "(missing)",
                        cell_left, GRID_TOP + LABEL_H + IMG_H / 2 - 0.15, CELL_W, 0.3,
                        font_pt=14, color=WHITE)
            missing.append(label)
    return slide, missing


def _exp_date_key(exp):
    m = re.match(r"(\d{2})/(\d{2})/(\d{4})", exp["tag"])
    return (int(m.group(3)), int(m.group(1)), int(m.group(2))) if m else (9999, 99, 99)


def _make_title(prefix: str, marker: str, tp_label: str, tag: str) -> str:
    # "<prefix> (<tp>, <date>) (<marker> staining)". The staining tag disambiguates the
    # two acquisitions that share a date (Vim-stain vs pericentrin-stain, both 04/06/2022)
    # and names which experiment each slide came from.
    sub = f"{tp_label}, {tag}" if tp_label else tag
    return f"{prefix} ({sub}) ({marker} staining)"


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    def build_section(dir_fn, prefix, scale_group, skip_msg):
        """One siCtrl|siVim slide per experiment (date-sorted) for a montage source.
        dir_fn(root, cond_folder, chan_sub) -> montages dir. Skips an experiment when
        neither condition has montages."""
        specs = []
        for exp in sorted(EXPERIMENTS, key=_exp_date_key):
            root = Path(exp["root"])
            chan_sub = exp["chan_sub"]
            left_folder, left_label = exp["left"]
            right_folder, right_label = exp["right"]
            log_key = f"{exp['marker']} {exp['tag']}"
            left_dir  = dir_fn(root, left_folder,  chan_sub)
            right_dir = dir_fn(root, right_folder, chan_sub)
            if not list_chunks(left_dir) and not list_chunks(right_dir):
                print(f"[{scale_group}/{log_key}]  {skip_msg} — skipped")
                continue
            specs.append({
                "log_key": log_key, "scale_group": scale_group,
                "title": _make_title(prefix, exp["marker"], exp["tp_label"], exp["tag"]),
                "left_label": left_label,   "left_img": find_first_chunks(left_dir, 1)[0],
                "right_label": right_label, "right_img": find_first_chunks(right_dir, 1)[0],
                "left_dir": left_dir, "right_dir": right_dir,
            })
        return specs

    # Section 1: actin + nucleus XZ MIP (physical-scale). Section 2 (appended at the
    # end): actin synapse-plane segmentation masks (non-physical-scale). Each section is
    # its own PPI group, so the masks size uniformly among themselves and independently
    # of the XZ MIPs.
    slide_specs: List[dict] = []
    slide_specs += build_section(
        lambda r, c, cs: montage_dir(r, c, cs, COMBO),
        "Actin + Nuc XZ MIP", "xz", "no actin_nuc_xz montages")
    slide_specs += build_section(
        mask_montage_dir,
        "Actin synapse mask", "mask", "no actin synapse-mask montages")

    # PPI pinning per scale_group.
    group_ppi: dict = {}
    for spec in slide_specs:
        imgs = [p for p in (spec["left_img"], spec["right_img"])
                if p is not None and _exists_long(p)]
        if imgs:
            sg = spec["scale_group"]
            group_ppi[sg] = max(group_ppi.get(sg, 0.0),
                                compute_slide_ppi(imgs, CELL_W, IMG_H))

    print(f"\nPPI pinning per scale_group across {len(slide_specs)} slides:")
    for sg, ppi in sorted(group_ppi.items()):
        print(f"  {sg:>5s}  PPI={ppi:.2f}")
    print(f"\nWriting deck to: {OUTPUT_PATH}\n")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    missing_total = []
    slides_added = 0
    for spec in slide_specs:
        _, missing = build_compare_slide(
            prs, spec["title"],
            spec["left_label"], spec["left_img"],
            spec["right_label"], spec["right_img"],
            group_ppi[spec["scale_group"]],
        )
        slides_added += 1
        l_ok = 1 if spec["left_img"] is not None else 0
        r_ok = 1 if spec["right_img"] is not None else 0
        print(f"[{spec['scale_group']}/{spec['log_key']}]  L:{l_ok}/1  R:{r_ok}/1")
        for cell in missing:
            src = spec["left_dir"] if cell == spec["left_label"] else spec["right_dir"]
            missing_total.append(f"{spec['log_key']}/{cell}  ({src})")

    prs.save(str(out_path))
    print(f"\nDone. {slides_added} slides written to:\n  {out_path}")

    if missing_total:
        print(f"\nMissing panels ({len(missing_total)}):")
        for m in missing_total:
            print(f"  - {m}")
    elif slides_added:
        print("\nAll panels found - no missing items.")


if __name__ == "__main__":
    main()
