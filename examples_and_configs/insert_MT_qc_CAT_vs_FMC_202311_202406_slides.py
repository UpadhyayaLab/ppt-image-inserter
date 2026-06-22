"""
insert_MT_qc_CAT_vs_FMC_202311_202406_slides.py

CAT (left) vs FMC (right) side-by-side MT (microtubule / βTub) QC deck
for the J: drive CART experiments. Stage AG pipeline output under
`prog_fixed_cells_foci/MT/`.

Same DATASETS + helpers as
insert_actin_qc_CAT_vs_FMC_202311_202406_slides.py; the only deltas
are PROG_SUBPATH, KIND_BLOCKS, and OUTPUT_PATH.

Kinds (3 per deck):
    Block 1 (interleaved):
        MT at Synapse  -> synapse/1slice/raw/montages
        MT Centrosome  -> centrosome/montages (MT-specific QC)
    Block 2:
        MT XZ MIP      -> xz_mip/montages

Slide math: (1 + 2 + 2) timepoints × 3 kinds = 15 slides.

Usage:
    python examples_and_configs/insert_MT_qc_CAT_vs_FMC_202311_202406_slides.py
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/CART/MT_CatB/"
    "CART_MT_QC_CAT_vs_FMC_202311_202406.pptx"
)

CTSB_ROOT = "J:/FF/fixed_cell/CAR_TCell"

# (date_tag, dataset_folder, cell_templates, timepoints)
# Each cell_templates entry: CAR_label -> (CAR_subdir, condition_format)
DATASETS = [
    (
        "20231127",
        "20231127_Fixed_CAR-Tcells_CTSB-mCherry_bTub",
        {
            "CAT": ("CAT",   "W3_NA_bCD19_CAT_{tp}_CTSB-mCh_647bTub_488Actin_Hoechst"),
            "FMC": ("FMC63", "W1_NA_bCD19_FMC63_{tp}_CTSB-mCh_647bTub_488Actin_Hoechst"),
        },
        ["15min"],
    ),
    (
        "20240620",
        "20240620_day3_Fixed_CAR-Tcells_CTSB-mCherry_bTub",
        {
            "CAT": ("CAT",   "CAT{tp}"),
            "FMC": ("FMC63", "FMC{tp}"),
        },
        ["5min", "15min"],
    ),
    (
        "20240624",
        "20240624_day5_Fixed_CAR-Tcells_CTSB-mCherry_bTub",
        {
            "CAT": ("CAT",   "CAT_{tp}"),
            "FMC": ("FMC63", "FMC63_{tp}"),
        },
        ["5min", "15min"],
    ),
]

CONDITION_SUBPATH = "cells/channels"
PROG_SUBPATH = "prog_fixed_cells_foci/MT"

KIND_BLOCKS = [
    [
        ("MT at Synapse", "synapse/1slice/raw/montages"),
        ("MT",            "centrosome/montages"),
    ],
    [
        ("MT XZ MIP", "xz_mip/montages"),
    ],
]

CHUNK_GLOB = "montage_cells_*.png"

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Slide layout (inches). 13.333 x 7.5 widescreen.
SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

GRID_LEFT = 0.10
GRID_TOP = 0.60
CELL_W = 6.50
CELL_H = SLIDE_H - GRID_TOP - 0.10
LABEL_H = 0.30
IMG_H = CELL_H - LABEL_H
LABEL_FONT_PT = 16
COL_GAP = SLIDE_W - 2 * GRID_LEFT - 2 * CELL_W

CELL_POSITIONS = [
    (GRID_LEFT,                    GRID_TOP),
    (GRID_LEFT + CELL_W + COL_GAP, GRID_TOP),
]

# Scalebar invariant (Stage AF/AG): per-tile 5 μm scalebar = 150 px nominal.
PPUM_SOURCE = 30
SCALEBAR_UM = 5
SCALEBAR_PX = PPUM_SOURCE * SCALEBAR_UM

# ---------------------------------------------------------------------------


def format_timepoint(tp: str) -> str:
    tp_map = {"5min": "5 min", "15min": "15 min"}
    return tp_map.get(tp.lower(), tp)


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _parse_chunk_range(p: Path) -> Tuple[int, int]:
    """Return (start, end) cell-id range for a montage_cells_*.png filename.
    Handles both the 4-int FOV-padded pattern (e.g. montage_cells_001_0_012_25.png)
    and the 2-int J:-drive pattern (e.g. montage_cells_1_26.png). Returns
    (0, 0) if no pattern matches."""
    m4 = re.match(r"montage_cells_(\d+)_(\d+)_(\d+)_(\d+)\.png$", p.name)
    if m4:
        f_a, c_a, f_b, c_b = (int(x) for x in m4.groups())
        return (f_a * 1000 + c_a, f_b * 1000 + c_b)
    m2 = re.match(r"montage_cells_(\d+)_(\d+)\.png$", p.name)
    if m2:
        return (int(m2.group(1)), int(m2.group(2)))
    return (0, 0)


def find_first_chunk(montages_dir: Path) -> Optional[Path]:
    """Pick the lowest-start chunk, BUT first drop any chunk whose [start, end]
    range is strictly contained in another chunk's range. That catches leftover
    smoke chunks (e.g. `montage_cells_1_12.png` next to `montage_cells_1_26.png`
    — the 1-12 range is shadowed by 1-26 and gets dropped). Same logic also
    handles the CTL/Kiet 4-int smoke pattern."""
    if not montages_dir.is_dir():
        return None
    chunks = list(montages_dir.glob(CHUNK_GLOB))
    if not chunks:
        return None
    parsed = [(p, *_parse_chunk_range(p)) for p in chunks]
    keep = []
    for (p, s, e) in parsed:
        shadowed = any(
            s2 <= s and e <= e2 and (s2 < s or e < e2)
            for (p2, s2, e2) in parsed
            if p2 is not p
        )
        if not shadowed:
            keep.append((p, s, e))
    if not keep:
        return None
    keep.sort(key=lambda x: x[1])  # by start
    return keep[0][0]


def _png_dims(path: Path) -> Tuple[int, int]:
    with Image.open(str(path)) as im:
        return im.size


def compute_deck_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_image_in_cell_at_ppi(slide, image_path: Path, ppi: float,
                             cell_left: float, cell_top: float):
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    img_area_top = cell_top + LABEL_H
    left_in = cell_left + (CELL_W - w_in) / 2
    top_in  = img_area_top + (IMG_H - h_in) / 2
    return slide.shapes.add_picture(
        str(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(w_in),
    )


def build_compare_slide(prs, title_text, cat_img, fmc_img, deck_ppi):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    cells = [
        ("CAT", cat_img, CELL_POSITIONS[0]),
        ("FMC", fmc_img, CELL_POSITIONS[1]),
    ]
    missing = []
    for label, img_path, (cell_left, cell_top) in cells:
        add_textbox(
            slide, label,
            cell_left, cell_top, CELL_W, LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )
        if img_path is not None and img_path.exists():
            add_image_in_cell_at_ppi(slide, img_path, deck_ppi, cell_left, cell_top)
        else:
            add_textbox(
                slide, "(missing)",
                cell_left, cell_top + LABEL_H + IMG_H / 2 - 0.15, CELL_W, 0.3,
                font_pt=14, color=WHITE,
            )
            missing.append(label)
    return slide, missing


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    root = Path(CTSB_ROOT)

    # slide_specs entries now carry kind_label so we can pin a PER-KIND PPI.
    # Within a kind, every slide shares the same px/inch (scalebar consistent
    # across all slides of that kind). Across kinds the PPI can differ — that
    # lets wide-aspect kinds like centrosome and XZ MIP fill more of the slide
    # without dragging the more-square synapse kind down to their size.
    slide_specs: List[Tuple[str, Optional[Path], Optional[Path], Path, Path, str, str]] = []
    for block in KIND_BLOCKS:
        for (date_tag, dataset_folder, cell_templates, timepoints) in DATASETS:
            for tp in timepoints:
                tp_pretty = format_timepoint(tp)
                for kind_label, kind_subpath in block:
                    def resolve(cell):
                        car_sub, cond_template = cell_templates[cell]
                        cond_folder = cond_template.format(tp=tp)
                        return (
                            root / dataset_folder / car_sub / cond_folder
                            / CONDITION_SUBPATH / PROG_SUBPATH / kind_subpath
                        )
                    cat_dir = resolve("CAT")
                    fmc_dir = resolve("FMC")
                    cat_img = find_first_chunk(cat_dir)
                    fmc_img = find_first_chunk(fmc_dir)
                    title = f"{kind_label}: {tp_pretty} ({date_tag})"
                    log_key = f"{kind_label}/{date_tag}/{tp}"
                    slide_specs.append((title, cat_img, fmc_img, cat_dir, fmc_dir, log_key, kind_label))

    # Group present images by kind and compute one PPI per kind.
    present_by_kind: dict = {}
    for (_, cat_img, fmc_img, _, _, _, kind_label) in slide_specs:
        bucket = present_by_kind.setdefault(kind_label, [])
        for p in (cat_img, fmc_img):
            if p is not None and p.exists():
                bucket.append(p)

    ppi_by_kind = {}
    for kind_label, paths in present_by_kind.items():
        if not paths:
            ppi_by_kind[kind_label] = 100.0
        else:
            ppi_by_kind[kind_label] = compute_deck_ppi(paths, CELL_W, IMG_H)

    print("Per-kind PPI (each kind internally consistent, cross-kind may differ):")
    for kind_label in [k for block in KIND_BLOCKS for (k, _) in block]:
        ppi = ppi_by_kind.get(kind_label)
        if ppi is None:
            continue
        bar_in = SCALEBAR_PX / ppi
        n = len(present_by_kind.get(kind_label, []))
        print(
            f"  {kind_label:<22} -> PPI {ppi:7.2f}  "
            f"(5 μm = {bar_in:.3f} in = {bar_in * 2.54:.3f} cm)  [{n} cells]"
        )
    print(f"\nSource PPUM = {PPUM_SOURCE} px/μm (locked).\n")
    print(f"Writing deck to: {OUTPUT_PATH}\n")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    missing_total = []
    slides_added = 0
    for (title, cat_img, fmc_img, cat_dir, fmc_dir, log_key, kind_label) in slide_specs:
        kind_ppi = ppi_by_kind[kind_label]
        _, missing = build_compare_slide(prs, title, cat_img, fmc_img, kind_ppi)
        slides_added += 1

        status_parts = [
            "CAT:OK" if cat_img else "CAT:MISSING",
            "FMC:OK" if fmc_img else "FMC:MISSING",
        ]
        print(f"[{log_key}]  " + "  ".join(status_parts))

        for cell in missing:
            src = cat_dir if cell == "CAT" else fmc_dir
            missing_total.append(f"{log_key}/{cell}  ({src})")

    prs.save(str(out_path))
    print(f"\nDone. {slides_added} slides written to:\n  {out_path}")

    if missing_total:
        print(f"\nMissing ({len(missing_total)}):")
        for m in missing_total:
            print(f"  - {m}")
    else:
        print("\nAll images found - no missing items.")


if __name__ == "__main__":
    main()
