"""
insert_xzmip_nuc_lamin_slides.py

Generate the XZ-MIP nucleus + lamin merge montage deck from a styled
template. For each lmnb1 condition subdirectory of the experiment, the
template slide is duplicated, the placeholder image is swapped for the
condition's XZ-MIP montage, and the title textbox text is replaced in
place (preserving the template's font, size, alignment). Title color is
forced to white and the slide background to black so the dark-deck look
is consistent regardless of template state.

Companion to insert_xymip_nuc_lamin_slides.py — only the source image
directory, file-name pattern, and template path differ.

Usage:
    python examples_and_configs/insert_xzmip_nuc_lamin_slides.py
"""

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

# Add repo root to path so we can import ppt_image_inserter when running
# this script directly (mirrors batch_insert_images.py).
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter.slide_utils import (  # noqa: E402
    duplicate_slide,
    remove_pictures_from_slide,
)

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

TEMPLATE_PATH = (
    "K:/FF/PPT/PPT_autogeneration/"
    "Ctrl, DZNep, bTub, Lamin B, 20260203/"
    "Ctrl_DZNep_upto1hr_XZMIP_montages_nuc_lamin_template.pptx"
)

# Output: same folder, "_template" stripped
OUTPUT_PATH = TEMPLATE_PATH.replace("_template.pptx", ".pptx")

PARENT_DIR = (
    "J:/FF/Nucleus_Project_up_to_1hr/"
    "020302026_jurkats_dznep_lmnB1orbTub_561LP55_488LP45_405LP40_"
)

MONTAGE_SUBPATH = "cropped/xz_mip_images/montages_nuc_lamin_merge"
IMAGE_GLOB = "xz_mip_montage_nuc_lamin_merge_*.png"

# Only process lmnb1 wells (skip bTub directories)
MARKER_FILTER = "lmnb1"

# Template structure: slide 0 is the cover, slide 1 is the per-condition
# template (with placeholder image + pre-formatted title textbox).
TEMPLATE_SLIDE_INDEX = 1

# Override title text color so titles read on a dark/black background.
TITLE_COLOR_RGB = RGBColor(0xFF, 0xFF, 0xFF)  # white

# Generated-slide background fill.
SLIDE_BG_RGB = RGBColor(0x00, 0x00, 0x00)  # black

# 914400 EMUs = 1 inch (python-pptx internal unit)
EMUS_PER_INCH = 914400.0

# ---------------------------------------------------------------------------


def format_condition_name(folder_name: str) -> str:
    """Convert a raw folder name to a readable slide title.

    Expected format: [replicate]_[condition]_[timepoint]_[marker]_
    Example: G2A3_h2o_8min_lmnb1_ -> 'H2O 8 min: Lamin B1'

    Falls back to the raw folder name if the pattern isn't recognised.
    """
    condition_map = {"h2o": "H2O", "dznep": "DZNep"}
    timepoint_map = {
        "4min": "4 min", "8min": "8 min", "15min": "15 min",
        "30min": "30 min", "1hr": "1 hr",
    }
    marker_map = {"lmnb1": "Lamin B1", "btub": "β-Tub"}

    parts = folder_name.strip("_").split("_")
    if len(parts) != 4:
        return folder_name

    _, condition, timepoint, marker = parts
    return (
        f"{condition_map.get(condition.lower(), condition)} "
        f"{timepoint_map.get(timepoint.lower(), timepoint)}: "
        f"{marker_map.get(marker.lower(), marker)}"
    )


def find_image(folder: Path, pattern: str) -> Path | None:
    """Return first glob match in folder, or None."""
    matches = sorted(folder.glob(pattern))
    return matches[0] if matches else None


def get_first_picture_position(slide) -> dict | None:
    """Return position of the first picture on slide as
    {'left','top','width','height'} in inches, or None if no picture."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return {
                "left": shape.left / EMUS_PER_INCH,
                "top": shape.top / EMUS_PER_INCH,
                "width": shape.width / EMUS_PER_INCH,
                "height": shape.height / EMUS_PER_INCH,
            }
    return None


def set_first_textbox_text(slide, text: str) -> bool:
    """Replace the text of the first TEXT_BOX shape on the slide,
    preserving the first run's formatting (font, size, color, bold).

    Targets MSO_SHAPE_TYPE.TEXT_BOX specifically so placeholders
    (slide numbers, layout titles) are skipped.

    Returns True if a textbox was updated, False otherwise.
    """
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
            continue
        if not getattr(shape, "has_text_frame", False):
            continue

        tf = shape.text_frame
        if not tf.paragraphs:
            continue
        para = tf.paragraphs[0]
        if not para.runs:
            # Empty paragraph: assign via .text (no formatting to preserve)
            para.text = text
        else:
            para.runs[0].text = text
            # Clear any subsequent runs in the first paragraph so only
            # the first run's text/formatting survives.
            for extra_run in para.runs[1:]:
                extra_run.text = ""

        # Override color so the title is legible on a dark/black background.
        if para.runs:
            para.runs[0].font.color.rgb = TITLE_COLOR_RGB

        # Clear any additional paragraphs (e.g. blank lines in the template).
        for extra_para in tf.paragraphs[1:]:
            if extra_para.runs:
                extra_para.runs[0].text = ""
                for r in extra_para.runs[1:]:
                    r.text = ""
            else:
                extra_para.text = ""
        return True

    return False


def add_image_at(slide, image_path: str, position: dict) -> None:
    """Add an image to slide at the given position dict (inches), preserving
    aspect ratio by specifying width only."""
    slide.shapes.add_picture(
        image_path,
        Inches(position["left"]),
        Inches(position["top"]),
        width=Inches(position["width"]),
    )


def set_slide_background(slide, rgb: RGBColor) -> None:
    """Set the slide background to a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def main() -> None:
    parent = Path(PARENT_DIR)
    if not parent.exists():
        print(f"ERROR: Parent directory not found: {parent}")
        return

    if not Path(TEMPLATE_PATH).exists():
        print(f"ERROR: Template not found: {TEMPLATE_PATH}")
        return

    # Filter to lmnb1 condition dirs only, sorted
    condition_dirs = sorted([
        d for d in parent.iterdir()
        if d.is_dir() and MARKER_FILTER in d.name.lower()
    ])
    if not condition_dirs:
        print(f"ERROR: No matching subdirectories found in {parent}")
        return

    print(f"Found {len(condition_dirs)} condition directories.")
    print(f"Opening template: {TEMPLATE_PATH}\n")

    prs = Presentation(TEMPLATE_PATH)
    template_slide = prs.slides[TEMPLATE_SLIDE_INDEX]

    # Capture the placeholder image position once from the template;
    # reuse for every duplicated slide.
    image_position = get_first_picture_position(template_slide)
    if image_position is None:
        print(
            f"ERROR: Template slide {TEMPLATE_SLIDE_INDEX} has no picture "
            "to use as a position reference."
        )
        return
    print(f"Template image position (in): {image_position}\n")

    slides_added = 0
    missing_report = []

    for cond_dir in condition_dirs:
        cond_name = cond_dir.name
        montage_dir = cond_dir / MONTAGE_SUBPATH
        img = find_image(montage_dir, IMAGE_GLOB) if montage_dir.exists() else None

        title = format_condition_name(cond_name)
        img_status = img.name if img else "NOT FOUND"
        print(f"[{cond_name}]  {img_status}  ->  {title!r}")

        if not montage_dir.exists():
            missing_report.append(f"{cond_name}: montage folder missing")
        elif not img:
            missing_report.append(f"{cond_name}: image not found")

        new_slide = duplicate_slide(prs, TEMPLATE_SLIDE_INDEX)

        if not set_first_textbox_text(new_slide, title):
            print(f"  WARNING: no textbox found on duplicated slide for {cond_name}")

        # Drop the template's placeholder picture from the duplicate,
        # then add the condition's montage at the template's position.
        remove_pictures_from_slide(new_slide)
        if img:
            add_image_at(new_slide, str(img), image_position)

        set_slide_background(new_slide, SLIDE_BG_RGB)

        slides_added += 1

    # Remove the original template slide so the output is:
    # [cover, generated_1, generated_2, ..., generated_N]
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[TEMPLATE_SLIDE_INDEX])

    prs.save(OUTPUT_PATH)
    print(f"\nDone. {slides_added} slides added. Saved to:\n  {OUTPUT_PATH}")

    if missing_report:
        print(f"\nMissing ({len(missing_report)}):")
        for msg in missing_report:
            print(f"  - {msg}")
    else:
        print("\nAll images found - no missing items.")


if __name__ == "__main__":
    main()
