"""
insert_cart_catb_nuc_summary_20260706_slides.py

Construct-comparison summary deck for the fixed CAR-T dataset compiled into
compiled_results/CART_MT_CatB_nuc_across_dates_20260706. The two conditions are
the CAR constructs — CAT vs FMC63 — pooled across 3 experiments / 5 date×
timepoint groups (Nov 27 2023 15 min; Jun 20 2024 d3 and Jun 24 2024 d5, each at
5 min and 15 min; CAT n = 432, FMC63 n = 514 from cell_counts.csv). Channels:
CatB (Cathepsin B, lytic granules), MT (β-tubulin; also the centrosome-context
stain — no dedicated centrosome marker), actin, and Hoechst/DNA. Each grid panel
is a CAT vs FMC63 comparison plot.

Direct port of the activated-CTL granule/nucleus deck
(insert_ctl_granule_nuc_summary_20260617_slides.py): same per-slide layout,
title slide, family dividers, native PowerPoint sections, and --list dry-run.
Every curated CTL metric maps 1:1 here after substituting the granule marker
Lamp1 -> CatB; slide titles name the marker "CatB". The reference deck's pairwise
scatter section is dropped — this compile's pairwise plots are nested per-date
and per-construct (date -> suite -> CAT/FMC63), which does not fit the pooled
single-scatter summary layout.

A FAMILIES entry is (stem, title) for a single panel, or ([(stem, sublabel), ...],
title) to place related panels (e.g. 1 μm / 2 μm thresholds) side by side on one
slide with a sublabel above each. The compile has 300+ grid panels total;
FAMILIES is this curated subset — add a row to grow the deck.

Self-contained: builds a blank deck (no template .pptx). Missing panels render
"(missing)" rather than failing. A previous deck is backed up before overwrite.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_cart_catb_nuc_summary_20260706_slides.py
    # dry run (print planned families/titles, build nothing):
    conda run -n PPT_editing python examples_and_configs/insert_cart_catb_nuc_summary_20260706_slides.py --list
"""

import os
import sys
import uuid
from pathlib import Path
from xml.sax.saxutils import escape

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation  # noqa: E402

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(
    "J:/FF/fixed_cell/CAR_TCell/compiled_results/"
    "CART_MT_CatB_nuc_across_dates_20260706"
)
BYDAY_DIR = ROOT / "by_day_panels"
CELL_COUNTS_PNG = ROOT / "cell_counts_barplot.png"   # context slide (optional)

OUTPUT_PATH = Path(
    "K:/FF/PPT/PPT_autogeneration/CART/nucleus/"
    "CART_CATB_MT_nuc_summary_across_dates_20260706.pptx"
)

# Each metric is one single-axis "by day" panel: CAT and FMC63 as an adjacent
# violin pair per experiment-date/timepoint, with a per-pair significance bracket
# and a CAT/FMC63 legend (produced by the compile's plot_by_day_condition_violin).
PANEL_SUFFIX = "_by_day.png"

# Compile date parsed from the dated ROOT folder (…_YYYYMMDD), shown in the footer.
_d = ROOT.name.rsplit("_", 1)[-1]          # e.g. "20260706"
COMPILE_DATE = "{}-{}-{}".format(_d[:4], _d[4:6], _d[6:8])

DECK_TITLE = "Granule polarization and nuclear morphology in CAR T cells (CAT vs FMC63)"
DECK_SUBTITLE = (
    "CAT vs FMC63 by date/timepoint  ·  5 experiments (Nov 2023, Jun 2024 d3/d5; 5 & 15 min)  ·  "
    "CatB / MT / actin / DNA  ·  compiled 2026-07-06"
)

# ---------------------------------------------------------------------------
# Curated metrics, grouped into families (divider slide per family). Each entry
# is (metric stem, slide title) for ONE by_day panel per slide. The stem +
# PANEL_SUFFIX is the PNG under by_day_panels/. Titles are for navigation; each
# panel carries its own authoritative y-axis label, per-date CAT/FMC63 x-labels,
# significance brackets, and legend.
# ---------------------------------------------------------------------------
FAMILIES = [
    ("Cell and nuclear spreading", [
        ("nuc_aspect_ratio",        "Nuclear aspect ratio"),
        ("actin_deform_ratio",      "Cell aspect ratio"),
        ("actin_bottom_mask_area",  "Synapse area"),
        ("nuc_broadest_slice_area", "Nuclear broadest-slice area"),
    ]),
    ("CatB granules — polarization and synapse delivery", [
        ("centrosome_center_z_rel_bottom_actin_plane", "Centrosome distance to synapse"),
        ("CatB_zCOF_cell_bottom_distance",  "CatB distance to synapse (zCOF)"),
        ("CatB_z50_cell_bottom_distance",   "CatB z₅₀ distance to synapse"),
        ("CatB_z75_cell_bottom_distance",   "CatB z₇₅ distance to synapse"),
        ("CatB_synapse_g_ave",              "CatB clustering at synapse (g)"),
        ("CatB_synapse_inner_outer_ratio",  "CatB synapse inner/outer ratio"),
    ]),
    ("CatB granules — dispersion", [
        ("CatB_FDD_3D",          "CatB dispersion (FDD, 3D)"),
        ("CatB_z_FDD",           "CatB dispersion (FDD, axial z)"),
        ("CatB_FDD_3D_rel_cent", "CatB dispersion rel. centrosome (3D)"),
        ("CatB_z_FDD_rel_cent",  "CatB dispersion rel. centrosome (axial z)"),
    ]),
    ("CatB granules — signal and centrosome localization", [
        ("CatB_total_sig",   "CatB total signal (whole cell)"),
        ("CatB_peak_sig",    "CatB peak signal"),
        ("CatB_synapse_MFI", "CatB MFI at synapse"),
        ("CatB_synapse_total_sig",      "Total CatB signal at synapse (single slice)"),
        ("CatB_synapse_total_sig_3mip", "Total CatB signal at synapse (3-slice MIP)"),
        ("CatB_frac_around_cent_1um", "CatB fraction around centrosome (1 μm)"),
        ("CatB_frac_around_cent_2um", "CatB fraction around centrosome (2 μm)"),
        ("CatB_frac_around_cent_3um", "CatB fraction around centrosome (3 μm)"),
        ("CatB_MFI_around_cent_1um",  "CatB MFI around centrosome (1 μm)"),
        ("CatB_MFI_around_cent_2um",  "CatB MFI around centrosome (2 μm)"),
        ("CatB_MFI_around_cent_3um",  "CatB MFI around centrosome (3 μm)"),
    ]),
    ("CatB granules — perinuclear and centrosome enrichment", [
        ("CatB_all_perinuc_MFI",             "CatB perinuclear MFI"),
        ("CatB_perinuc_sig_fraction",        "CatB perinuclear signal fraction"),
        ("CatB_frac_perinuc_within_1_um_cent", "CatB perinuclear fraction near centrosome (1 μm)"),
        ("CatB_frac_perinuc_within_2_um_cent", "CatB perinuclear fraction near centrosome (2 μm)"),
        ("CatB_cyto_in_nuc_hull_MFI",          "CatB MFI in cytoplasm within nuclear hull"),
        ("CatB_cyto_in_nuc_hull_sig_fraction", "CatB signal fraction in cytoplasm within nuclear hull"),
        ("CatB_frac_in_nuc_convex_hull",       "CatB fraction in nuclear convex hull"),
        ("CatB_enrichment_within_half_um_nuc_2_um_cent", "CatB enrichment near centrosome (0.5 μm nuc, 2 μm cent)"),
        ("MT_enrichment_within_half_um_nuc_2_um_cent",   "MT enrichment near centrosome (0.5 μm nuc, 2 μm cent)"),
    ]),
    ("Centrosome ↔ nucleus", [
        ("nuc_cent_closest_dist",            "Nucleus-centrosome closest distance"),
        ("cent_nuc_norm_dist_sphere_rad",    "Centrosome-to-nuclear-centroid distance (norm. to equiv sphere radius)"),
        ("centrosome_dist_deepest_real_avg_periphery_ratio", "Centrosome distance to deepest invag vs avg periphery ratio"),
        # Centrosome radial position in the cell footprint (0 = center, 1 = edge),
        # computed on the MT-derived centrosome (process_MT_channel.m).
        ("centrosome_r_norm_bottom_plane_from_MT", "Centrosome radial position (synapse plane)"),
        ("centrosome_r_norm_MIP_from_MT",          "Centrosome radial position (MIP)"),
    ]),
    ("Nuclear deformation and invaginations", [
        ("chull_max_D",                       "Max invag depth over full nucleus"),
        ("chull_max_D_by_cent",               "Invagination depth near centrosome"),
        ("chull_mean_D_cent_global_ratio",    "Centrosomal Invagination Index (global)"),
        ("C_min_F_mean_by_cent",  "Nuclear surface curvature near centrosome (min principal)"),
        ("C_mean_F_mean_by_cent", "Nuclear surface curvature near centrosome (mean)"),
        ("deepest_invag_fraction_chull_volume", "Deepest invag: frac of convex hull volume"),
        ("deepest_region_periph_ratio_025um", "DNA levels near invag"),
        ("invag_by_cent_centroid_z_syn_from_MT", "Invag region (near cent): centroid height above synapse"),
        ("invag_by_cent_tip_z_syn_from_MT",      "Invag region (near cent): tip height above synapse"),
    ]),
    ("Invagination orientation", [
        ("avg_normal_angle_adaptive_region_growth",         "Deepest invag orientation"),
        ("avg_normal_angle_adaptive_region_growth_by_cent", "Invag orientation (adaptive) near centrosome"),
        ("avg_normal_angle_by_cent",                        "Invag orientation near centrosome"),
    ]),
    ("Nuclear morphology", [
        ("nuc_solidity",    "Nuclear solidity"),
        ("nuc_volume_mesh", "Nuclear volume"),
        ("nuc_SA_mesh",     "Nuclear surface area"),
    ]),
    ("Microtubules (β-tubulin)", [
        ("MT_frac_in_nuc_convex_hull", "MT fraction in nuclear convex hull"),
        ("MT_frac_around_cent_2um",    "MT fraction around centrosome (2 μm)"),
        ("MT_MFI_around_cent_2um",     "MT MFI around centrosome (2 μm)"),
    ]),
    ("Actin — levels and localization", [
        ("actin_total_sig",              "Total actin signal (whole cell)"),
        ("actin_bottom_MFI",             "Actin MFI at synapse (single slice)"),
        ("actin_bottom_MFI_3mip",        "Actin MFI at synapse (3-slice MIP)"),
        ("actin_bottom_total_sig",       "Total actin signal at synapse (single slice)"),
        ("actin_bottom_total_sig_3mip",  "Total actin signal at synapse (3-slice MIP)"),
        ("actin_bottom_inner_outer_ratio", "Actin synapse inner/outer ratio"),
        ("actin_MFI_around_cent_1um",    "Actin MFI around centrosome (1 μm)"),
        ("actin_MFI_around_cent_2um",    "Actin MFI around centrosome (2 μm)"),
        ("actin_frac_around_cent_1um",   "Actin fraction around centrosome (1 μm)"),
        ("actin_frac_around_cent_2um",   "Actin fraction around centrosome (2 μm)"),
    ]),
    # --- Appendix: additional context metrics -----------------------------------
    ("Cell and nuclear flattening (context)", [
        ("actin_height",        "Cell height"),
        ("nuc_height",          "Nuclear height"),
        ("nuc_centroid_z",      "Nuclear centroid height above synapse"),
        ("nuc_mesh_sphericity", "Nuclear sphericity"),
        ("actin_MIP_circularity", "Cell footprint circularity"),
    ]),
    ("Chromatin / DNA distribution", [
        ("nuc_all_CV",          "DNA intensity CV (heterogeneity)"),
        ("nuc_all_prop_gr_2med", "DNA fraction > 2× median (bright foci)"),
        ("nuc_all_skewness",    "DNA intensity skewness"),
        ("nuc_all_norm_entropy", "DNA distribution normalized entropy"),
    ]),
]

# ---------------------------------------------------------------------------
# Colors / layout (matches the bleb/noco/vimkd/CTL summary decks)
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.10

TITLE_LEFT = MARGIN
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.55
TITLE_FONT_PT = 28

IMG_LEFT = MARGIN
IMG_TOP = 0.66
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = 6.36

# Side-by-side (multi-panel) slides: a sublabel band above the figures, and a
# gap between columns.
COL_GAP = 0.15
SUBLABEL_H = 0.30
SUBLABEL_GAP = 0.04
SUBLABEL_FONT_PT = 16

FOOTER_LEFT = MARGIN
FOOTER_TOP = 7.06
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.40
FOOTER_FONT_PT = 9

# ---------------------------------------------------------------------------


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    """Place an image inside (left, top, w, h), preserving aspect ratio and
    centering on whichever dimension ends up smaller than the box. The grid
    panels are near-square, so they fit to height and center horizontally."""
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def rel_footer(path):
    try:
        rel = path.relative_to(ROOT).as_posix()
    except ValueError:
        rel = path.as_posix()
    return "{} / {}".format(COMPILE_DATE, rel)


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.7, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=40, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.1, SLIDE_W - 2 * MARGIN, 1.0,
                font_pt=18, color=BLACK, italic=True)


def build_divider_slide(prs, family_name):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, family_name, MARGIN, 3.1, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=44, color=BLACK, bold=True)


def _ext(path):
    """Windows extended-length path (\\\\?\\...) so files whose absolute path
    exceeds the 260-char MAX_PATH are still found and opened. Long metric names
    under the deep compiled_results tree bust MAX_PATH; without this prefix
    Path.exists()/PIL.open silently fail. No-op off Windows."""
    p = os.path.abspath(str(path))
    if os.name == "nt" and not p.startswith("\\\\?\\"):
        p = "\\\\?\\" + p.replace("/", "\\")
    return p


def _exists(path):
    return os.path.exists(_ext(path))


def _place_missing(slide, left, top, width):
    add_textbox(slide, "no data yet", left, top + IMG_BOX_H / 2 - 0.2,
                width, 0.4, font_pt=18, color=BLACK)


def build_slide(prs, title_text, panels, footer_text):
    """Title + one by_day panel (aspect preserved, full width) + source-path
    footer. `panels` is a single-element list of (path, _). Returns the list of
    missing panel paths (not on disk)."""
    slide = _new_slide(prs)
    add_textbox(slide, title_text, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title_text), color=BLACK, bold=True)

    missing = []
    path, _ = panels[0]
    if _exists(path):
        add_image_in_box(slide, _ext(path), IMG_LEFT, IMG_TOP, IMG_BOX_W, IMG_BOX_H)
    else:
        _place_missing(slide, IMG_LEFT, IMG_TOP, IMG_BOX_W)
        missing.append(path)

    add_textbox(slide, footer_text, FOOTER_LEFT, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=BLACK)
    return missing


def _panel_path(stem):
    """Resolve a metric stem to its by_day panel: BYDAY_DIR/<stem>_by_day.png."""
    return BYDAY_DIR / (stem + PANEL_SUFFIX)


def entry_panels(entry_stem):
    """Normalize a FAMILIES entry's stem field into a list of (path, sublabel).
    Every entry is one by_day panel per slide, so this is a single-element list."""
    return [(_panel_path(entry_stem), None)]


def add_sections(prs, section_spec):
    """Add native PowerPoint sections (the collapsible, named groups shown as
    tabs in the slide navigator / slide sorter). python-pptx has no API for
    this, so inject the p14:sectionLst extension into the presentation part.

    `section_spec` is an ordered list of (section_name, n_slides); the counts
    must sum to the total slide count in build order, and the first section
    must contain the first slide (a PowerPoint requirement)."""
    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    slide_ids = [sldId.get("id") for sldId in prs.slides._sldIdLst]

    parts = ['<p:extLst xmlns:p="{}">'.format(P),
             '<p:ext uri="{{521415D9-36F7-43E2-AB2F-B90AF26B5E84}}">',
             '<p14:sectionLst xmlns:p14="{}">'.format(P14)]
    i = 0
    for name, count in section_spec:
        parts.append('<p14:section name="{}" id="{{{}}}">'.format(
            escape(name), str(uuid.uuid4()).upper()))
        parts.append('<p14:sldIdLst>')
        for sid in slide_ids[i:i + count]:
            parts.append('<p14:sldId id="{}"/>'.format(sid))
        parts.append('</p14:sldIdLst></p14:section>')
        i += count
    parts.append('</p14:sectionLst></p:ext></p:extLst>')

    # The sldIdLst's parent is the <p:presentation> element; extLst goes last.
    prs.slides._sldIdLst.getparent().append(parse_xml("".join(parts)))


def main():
    list_only = "--list" in sys.argv

    n_metrics = sum(len(items) for _, items in FAMILIES)
    # title + (cell-counts if present) + per-family (divider + metric slides)
    est_slides = 1 + (1 if CELL_COUNTS_PNG.exists() else 0) + \
        sum(1 + len(items) for _, items in FAMILIES)

    print("Source: {}".format(BYDAY_DIR))
    print("{} metric slides across {} families, est. {} slides\n".format(
        n_metrics, len(FAMILIES), est_slides))

    if list_only:
        for fam, items in FAMILIES:
            print("=== {} ({}) ===".format(fam, len(items)))
            for entry_stem, title in items:
                panels = entry_panels(entry_stem)
                flags = " ".join(
                    "{}:{}".format(sub or "-", "OK" if _exists(p) else "MISS")
                    for p, sub in panels)
                print("  [{}] {:<50s} {}".format(
                    "OK " if all(_exists(p) for p, _ in panels) else "MISS",
                    title, flags))
            print("")
        return

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)

    if CELL_COUNTS_PNG.exists():
        build_slide(prs, "Cell counts (CAT vs FMC63)",
                    [(CELL_COUNTS_PNG, None)], rel_footer(CELL_COUNTS_PNG))
    else:
        print("Note: {} not found - skipping cell-counts slide.\n".format(
            CELL_COUNTS_PNG.name))

    # Native PowerPoint sections: an intro section (title + cell counts) then
    # one section per family (its divider slide + metric slides).
    section_spec = [("Overview", len(prs.slides._sldIdLst))]

    missing = []
    for fam, items in FAMILIES:
        build_divider_slide(prs, fam)
        print("=== {} ===".format(fam))
        for entry_stem, title in items:
            panels = entry_panels(entry_stem)
            footer = " | ".join(rel_footer(p) for p, _ in panels)
            miss = build_slide(prs, title, panels, footer)
            status = "OK" if not miss else "MISSING"
            print("  [{}] {} -> {!r}".format(
                status, ", ".join(p.name for p, _ in panels), title))
            missing.extend(p.name for p in miss)
        section_spec.append((fam, 1 + len(items)))
        print("")

    add_sections(prs, section_spec)

    if OUTPUT_PATH.exists():
        backup_dir = OUTPUT_PATH.parent / "backups"
        created = backup_presentation(str(OUTPUT_PATH), backup_base=str(backup_dir))
        if created:
            print("Backed up previous deck to: {}\n".format(backup_dir))

    prs.save(str(OUTPUT_PATH))
    total = len(prs.slides._sldIdLst)
    print("Done. {} metric slides, {} slides written to:\n  {}".format(
        n_metrics, total, OUTPUT_PATH))
    if missing:
        print("\nSkipped {} missing panel(s) (not on disk):".format(len(missing)))
        for m in missing:
            print("  - {}".format(m))
    else:
        print("\nAll curated panels found - no missing items.")


if __name__ == "__main__":
    main()
