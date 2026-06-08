"""
insert_actin_qc_naive_CTL_synapse_mask_xz_mip_slides.py

Naive-CTL actin QC deck — one image per slide, two kinds (synapse mask
+ XZ MIP) across 12 substrate / timepoint conditions from the
06/17/2025 GzmB + βTub imaging session.

Pipeline subdir is `prog_CTL_analysis/actin/<kind>/montages/` (the CTL
analog of the CART pipeline's `prog_fixed_cells_actin_only/actin/`).
Kind structure and montage filename pattern
(`montage_cells_<FOV>_<cell>_<FOV>_<cell>.png`) match the CART
pipeline.

Result: 2 kinds x 12 conditions = 24 slides.
    Slides  1-12: Actin at Synapse, conditions in display order
    Slides 13-24: Actin XZ MIP, same conditions

Modeled on insert_actin_qc_synapse_mask_xz_mip_slides.py.

The pick logic skips leftover smoke-test montages
(`montage_cells_<a>_<b>_<a>_<c>.png` where c - b <= 5; e.g.
`montage_cells_01_2_01_4.png` = 3 cells in FOV 01) and uses the first
real chunk instead.

Usage:
    python examples_and_configs/insert_actin_qc_naive_CTL_synapse_mask_xz_mip_slides.py
"""

import os
import re
import sys
from pathlib import Path
from typing import Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/Naive_CTL/"
    "Naive_CTL_actin_QC_synapse_mask_xz_mip_06172025.pptx"
)

DATASET_ROOT = (
    "L:/FF/Naive_CTLs/GzB_bTub/06172025_naiveCTLs_grzmB_bTub_"
)

# Date tag shown in slide titles.
DATE_TAG = "06/17/2025"

# (display_label, raw_folder_name) — user-specified ordering.
CONDITIONS = [
    ("1.5 kPa 24 hr", "auto-gel3p0-24hr-0617"),
    ("1.5 kPa 48 hr", "f-gel3p0-48hr-0617"),
    ("1.5 kPa 72 hr", "auto-gel3p0-72hr-0617"),
    ("12 kPa 24 hr",  "auto-gel8p0-24hr-0617"),
    ("12 kPa 48 hr",  "auto-gel8p0-48hr-0617-weirdchannel"),
    ("12 kPa 72 hr",  "auto-gel8p0-72hr-0617"),
    ("glass 2 hr",    "f-actGlass-101010-2hr-0617"),
    ("glass 6 hr",    "f-actGlass-101010-6hr-0617"),
    ("glass 24 hr",   "f-actglass-101010-24hr-0617"),
    ("glass 48 hr",   "f-actglass-101010-48hr-0617"),
    ("glass 72 hr",   "f-actGlass-101010-72hr-0617"),
    ("PLL 2 hr",      "f-nonactGlassPLL-2hr-0617"),
]

PROG_SUBPATH = "prog_CTL_analysis/actin"

# (kind label, subpath under PROG_SUBPATH/)
KINDS = [
    ("Actin at Synapse", "synapse/mask/montages"),
    ("Actin XZ MIP",     "xz_mip/montages"),
]

CHUNK_GLOB = "montage_cells_*.png"

# Smoke-test threshold: a chunk is considered a smoke artifact if it
# fits in a single FOV and covers <= this many cells of index range.
SMOKE_RANGE_THRESHOLD = 5

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Slide layout (inches). 13.333 x 7.5 widescreen.
SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

IMG_LEFT = 0.10
IMG_TOP = 0.60
IMG_BOX_W = SLIDE_W - 2 * 0.10           # 13.13"
IMG_BOX_H = SLIDE_H - IMG_TOP - 0.10     # 6.80"

# ---------------------------------------------------------------------------


def _chunk_start_index(p: Path) -> int:
    """Extract the first integer after 'montage_cells_' for natural sort."""
    m = re.match(r"montage_cells_(\d+)", p.name)
    return int(m.group(1)) if m else 0


def _is_smoke_chunk(p: Path) -> bool:
    """Detect leftover smoke-test montages.

    Filename pattern: `montage_cells_<FOV_a>_<cell_a>_<FOV_b>_<cell_b>.png`.
    Treated as smoke iff FOV_a == FOV_b AND (cell_b - cell_a) <= threshold.
    Filenames not matching the 4-integer pattern are NOT treated as smoke.
    """
    m = re.match(r"montage_cells_(\d+)_(\d+)_(\d+)_(\d+)\.png$", p.name)
    if not m:
        return False
    fov_a, cell_a, fov_b, cell_b = (int(x) for x in m.groups())
    return fov_a == fov_b and (cell_b - cell_a) <= SMOKE_RANGE_THRESHOLD


def find_first_real_chunk(montages_dir: Path) -> Tuple[Optional[Path], bool]:
    """Return (chunk_path, fallback_was_used) for the first non-smoke chunk.

    If every chunk in the dir looks like a smoke, fall back to the
    lowest-numbered one but set fallback_was_used=True so it gets flagged.
    """
    if not montages_dir.is_dir():
        return None, False
    chunks = sorted(montages_dir.glob(CHUNK_GLOB), key=_chunk_start_index)
    if not chunks:
        return None, False
    real = [c for c in chunks if not _is_smoke_chunk(c)]
    if real:
        return real[0], False
    return chunks[0], True


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    """Add a centered textbox with given text and styling."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    """Place an image inside the given bounding box, preserving aspect ratio."""
    pic = slide.shapes.add_picture(
        image_path,
        Inches(box_left),
        Inches(box_top),
        width=Inches(box_w),
    )
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path,
            Inches(box_left),
            Inches(box_top),
            height=Inches(box_h),
        )
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def build_slide(prs, title_text: str, image_path: Optional[Path]):
    """Build one full-image slide with title. Returns (slide, missing_flag)."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    if image_path is not None and image_path.exists():
        add_image_in_box(slide, str(image_path), IMG_LEFT, IMG_TOP, IMG_BOX_W, IMG_BOX_H)
        return slide, False

    add_textbox(
        slide, "(missing)",
        IMG_LEFT, IMG_TOP + IMG_BOX_H / 2 - 0.2, IMG_BOX_W, 0.4,
        font_pt=18, color=WHITE,
    )
    return slide, True


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    root = Path(DATASET_ROOT)
    missing = []
    smoke_fallback_notes = []
    slides_added = 0

    print(f"Writing deck to: {OUTPUT_PATH}\n")

    # Loop order: kind -> condition. Slides 1-12 Synapse, 13-24 XZ MIP.
    for kind_label, kind_subpath in KINDS:
        for display_label, raw_folder in CONDITIONS:
            montages_dir = root / raw_folder / PROG_SUBPATH / kind_subpath
            chunk, fell_back = find_first_real_chunk(montages_dir)

            title = f"{kind_label}: {display_label}  ({DATE_TAG})"
            _, is_missing = build_slide(prs, title, chunk)
            slides_added += 1

            if chunk is None:
                status = "MISSING"
                missing.append(f"{kind_label}/{display_label}  ({montages_dir})")
            elif fell_back:
                status = f"OK (smoke-fallback: {chunk.name})"
                smoke_fallback_notes.append(
                    f"{kind_label}/{display_label}: only smoke-sized chunks "
                    f"present, used {chunk.name}"
                )
            else:
                status = f"OK ({chunk.name})"
            print(f"[{kind_label}/{display_label}]  {status}")

    prs.save(str(out_path))
    print(f"\nDone. {slides_added} slides written to:\n  {out_path}")

    if missing:
        print(f"\nMissing ({len(missing)}):")
        for m in missing:
            print(f"  - {m}")
    if smoke_fallback_notes:
        print(f"\nSmoke fallback used ({len(smoke_fallback_notes)}):")
        for n in smoke_fallback_notes:
            print(f"  - {n}")
    if not missing and not smoke_fallback_notes:
        print("\nAll real chunks found — no smoke fallbacks or missing items.")


if __name__ == "__main__":
    main()
