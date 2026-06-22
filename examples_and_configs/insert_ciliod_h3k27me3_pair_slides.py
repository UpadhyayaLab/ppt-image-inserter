"""
insert_ciliod_h3k27me3_pair_slides.py

Build a deck where each slide shows ONE metric with the 06/12/2026 grid panel
on the left and the 06/13/2026 panel on the right, separated by a small gap.
The metric list and titles are pulled from the two existing single-date
YAMLs so the three CilioD H3K27me3 decks stay in lockstep — change one YAML
and re-run all the rebuilds.

Modeled after build_pair_slide() in insert_histone_mods_summary_slides.py.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_ciliod_h3k27me3_pair_slides.py
"""

import io
import os
import sys
from pathlib import Path

import yaml
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation  # noqa: E402


def safe_exists(path):
    """Bypass Windows MAX_PATH (260 chars) for the .exists() check."""
    s = str(path)
    if sys.platform == "win32" and len(s) > 240:
        return os.path.exists("\\\\?\\" + s.replace("/", "\\"))
    return os.path.exists(s)


def safe_open_rb(path):
    """Read a file with Windows long-path support — returns BytesIO so we can
    hand it to python-pptx's add_picture without it doing its own open()."""
    s = str(path)
    if sys.platform == "win32" and len(s) > 240:
        s = "\\\\?\\" + s.replace("/", "\\")
    with open(s, "rb") as f:
        return io.BytesIO(f.read())

# ---------------------------------------------------------------------------

REPO = Path(__file__).resolve().parents[1]
CONFIG_A = REPO / "examples_and_configs/configs/fixed/Jurkats/config_CilioD_H3K27me3_06122026_fixed_cells.yaml"
CONFIG_B = REPO / "examples_and_configs/configs/fixed/Jurkats/config_CilioD_H3K27me3_06132026_fixed_cells.yaml"

OUTPUT_PATH = Path(
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/CilioD/"
    "Effects of CilioD, H3K27me3, 06122026 vs 06132026 side-by-side.pptx"
)

LABEL_LEFT = "06/12/2026"
LABEL_RIGHT = "06/13/2026"

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Slide layout (inches). 13.333 x 7.5 widescreen.
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.10

TITLE_LEFT = MARGIN
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

PAIR_GAP = 0.20
PAIR_LABEL_TOP = 0.55
PAIR_LABEL_HEIGHT = 0.32
PAIR_LABEL_FONT_PT = 16
PAIR_IMG_TOP = 0.92
PAIR_IMG_H_MAX = 6.08  # upper bound; actual height shared by L/R, picked per-slide

FOOTER_LEFT = MARGIN
FOOTER_TOP = 7.05
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.40
FOOTER_FONT_PT = 8

# ---------------------------------------------------------------------------


def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_image_at_height(slide, image_path, left_in, top_in, height_in):
    """Place image at exact height (inches); width follows aspect ratio.
    Reads via safe_open_rb to bypass Windows MAX_PATH."""
    image_bytes = safe_open_rb(image_path)
    return slide.shapes.add_picture(
        image_bytes, Inches(left_in), Inches(top_in), height=Inches(height_in),
    )


def get_image_aspect(image_path):
    """Width-over-height ratio. Uses safe_open_rb for long Windows paths."""
    image_bytes = safe_open_rb(image_path)
    w, h = PILImage.open(image_bytes).size
    return w / h


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def build_pair_slide(prs, title, left_path, right_path):
    """Build one slide: title, two images sharing the same height
    (06/12 left, 06/13 right — wider because of more conditions),
    and a two-line footer."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, WHITE)

    add_textbox(
        slide, title,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=BLACK, bold=True,
    )

    labels = (LABEL_LEFT, LABEL_RIGHT)
    paths = (left_path, right_path)
    missing = [not safe_exists(p) for p in paths]

    # Both images share the same rendered height H. Width per image follows
    # its aspect ratio. H is capped both by available width (so the pair fits
    # horizontally) and PAIR_IMG_H_MAX (so it doesn't bleed into the footer).
    usable_w = SLIDE_W - 2 * MARGIN - PAIR_GAP
    aspects = [get_image_aspect(p) if not m else 1.0
               for p, m in zip(paths, missing)]
    height_by_width = usable_w / sum(aspects)
    h_in = min(height_by_width, PAIR_IMG_H_MAX)
    widths = [a * h_in for a in aspects]

    total_w = sum(widths) + PAIR_GAP
    start_x = (SLIDE_W - total_w) / 2
    lefts = (start_x, start_x + widths[0] + PAIR_GAP)

    for i in (0, 1):
        add_textbox(
            slide, labels[i],
            lefts[i], PAIR_LABEL_TOP, widths[i], PAIR_LABEL_HEIGHT,
            font_pt=PAIR_LABEL_FONT_PT, color=BLACK, bold=True,
        )
        if not missing[i]:
            add_image_at_height(slide, paths[i], lefts[i], PAIR_IMG_TOP, h_in)
        else:
            add_textbox(
                slide, "(missing)",
                lefts[i], PAIR_IMG_TOP + h_in / 2 - 0.2,
                widths[i], 0.4, font_pt=14, color=BLACK,
            )

    return slide, missing


def main():
    with open(CONFIG_A, "r", encoding="utf-8") as f:
        cfg_a = yaml.safe_load(f)
    with open(CONFIG_B, "r", encoding="utf-8") as f:
        cfg_b = yaml.safe_load(f)

    base_a = Path(cfg_a["base_dir"])
    base_b = Path(cfg_b["base_dir"])
    images_a = cfg_a["images"]
    images_b = cfg_b["images"]
    assert len(images_a) == len(images_b), \
        f"YAMLs differ in length: {len(images_a)} vs {len(images_b)}"

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    print(f"Writing pair deck to: {OUTPUT_PATH}\n")

    missing_files = []
    for ea, eb in zip(images_a, images_b):
        assert ea["path"] == eb["path"] and ea["title"] == eb["title"], \
            f"entry mismatch: {ea} vs {eb}"
        title = ea["title"]
        left_path = base_a / ea["path"]
        right_path = base_b / eb["path"]
        _, miss = build_pair_slide(prs, title, left_path, right_path)
        status_l = "OK" if not miss[0] else "MISSING"
        status_r = "OK" if not miss[1] else "MISSING"
        print(f"[{ea['path']}]  L:{status_l} R:{status_r}  -> {title!r}")
        if miss[0]:
            missing_files.append(f"L: {left_path}")
        if miss[1]:
            missing_files.append(f"R: {right_path}")

    if OUTPUT_PATH.exists():
        backup_dir = OUTPUT_PATH.parent / "backups"
        created = backup_presentation(str(OUTPUT_PATH), backup_base=str(backup_dir))
        if created:
            print(f"\nBacked up previous deck to: {backup_dir}")

    prs.save(str(OUTPUT_PATH))
    print(f"\nDone. {len(images_a)} pair slides written.")
    if missing_files:
        print(f"\nMissing ({len(missing_files)}):")
        for m in missing_files:
            print(f"  - {m}")
        sys.exit(1)


if __name__ == "__main__":
    main()
