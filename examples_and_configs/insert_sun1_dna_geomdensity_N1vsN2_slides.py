"""
insert_sun1_dna_geomdensity_N1vsN2_slides.py

Assemble the SUN1 / DNA "localization with respect to nuclear-envelope geometry"
summary deck for the siControl Jurkats N1-vs-N2 geometry-density compilation
(Jurkats_SUN1_N1vsN2_siControl_geomdensity_20260712).

Three sections, built from the curated result folders:

  1. SUN1 vs NE geometry (curated SUN1_loc_wrto_invag/ + NE-only singles)
       - averaged Mean+/-SEM density vs geometry (hulldist / min / mean curv)
       - enrichment in deep invaginations (hull dist > 0.5 & > 1.0 um), and an
         NE-only companion (boundary shell) from geom_density/enrichment/singles/
       - enrichment on concave surface (min curv < 0, < -0.25; mean curv < 0)
       - per-cell correlation vs geometry (hulldist / min / mean curv)
       then raw QC panels (qc_raw/, tiled).
       Story: SUN1 accumulates in invaginated / concave NE, reproduced N1 <-> N2.
  2. DNA vs NE geometry (DNA_loc_wrto_invag/)  -- same layout for DNA.
  3. grid_panels/N1_vs_N2/  -- SUN1 x DNA correlation at the NE plus the N1-vs-N2
       comparison violins for every SUN1 geometry-correlation scalar (convex-hull
       interior, deepest invagination, invag-vs-other per shell).

Self-contained: builds a blank 16:9 deck (no template .pptx). Missing panels are
skipped (no placeholder, no failure); a previous deck at the output path is
backed up first.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_sun1_dna_geomdensity_N1vsN2_slides.py
    # dry run (print the planned sections/slides, build nothing):
    conda run -n PPT_editing python examples_and_configs/insert_sun1_dna_geomdensity_N1vsN2_slides.py --list
"""

import math
import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import backup_presentation  # noqa: E402

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(
    "H:/FF_backup/Jurkat_nucleus/from_Ivan_HD/results_compilation/"
    "Jurkats_SUN1_N1vsN2_siControl_geomdensity_20260712"
)
SUN1_DIR = ROOT / "SUN1_loc_wrto_invag"
DNA_DIR = ROOT / "DNA_loc_wrto_invag"
SINGLES_DIR = ROOT / "geom_density" / "enrichment" / "singles"   # NE-only (boundary) panels
GRID_DIR = ROOT / "grid_panels" / "N1_vs_N2"

OUTPUT_PATH = Path(
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/SUN/"
    "SUN1 and DNA vs NE geometry, N1 vs N2.pptx"
)

DECK_TITLE = "SUN1 & DNA localization vs nuclear-envelope geometry"
DECK_SUBTITLE = (
    "Jurkats, siControl  ·  per topic: 05/18 (N1) then N1-vs-N2  ·  NE + perinuc 0.5 μm"
    "  ·  + centrosome metrics (N2)  ·  geometry–density analysis, compiled 2026-07-12"
)

# Second compilation: the single-condition 05/18/2022 (N1) results, added ahead
# of the N1-vs-N2 comparison. Same folder layout as ROOT.
SINGLE_ROOT = Path(
    "H:/FF_backup/Jurkat_nucleus/from_Ivan_HD/05182022 - KDs Validation N1/"
    "results_compilation/Jurkats_SUN1_20220518_siControl_geomdensity_20260712"
)
COMPARE_ROOT = ROOT


def curated_dir(root, channel):
    return root / (channel + "_loc_wrto_invag")


def singles_dir(root):
    return root / "geom_density" / "enrichment" / "singles"


# Compile C (N2 alone, 05/31/2022) -- source of the centrosome-proximal SUN1×DNA
# metric only (no per-topic N2 slides).
COMPILE_C_ROOT = Path(
    "H:/FF_backup/Jurkat_nucleus/from_Ivan_HD/05312022 - KDs Validation N2/"
    "results_compilation/Jurkats_SUN1_N2_20220531_siControl_geomdensity_20260712"
)

# Per channel the 05/18/2022 single condition (N1) and the N1-vs-N2 comparison
# are interleaved per topic (a 0518 slide, then the comparison slide); comparison
# QC follows. Grid-panel scalars (+ the N2 centrosome metric) are appended last.
# (channel, divider_title, divider_subtitle)
# 0518 QC is skipped: those raw panels are byte-identical to the comparison's
# N1 (May 18) QC already in the deck.
CHANNEL_SECTIONS = [
    ("SUN1", "SUN1 vs NE geometry", "per topic: 05/18/2022 (N1), then N1 vs N2 · siControl"),
    ("DNA", "DNA vs NE geometry", "per topic: 05/18/2022 (N1), then N1 vs N2 · siControl"),
]

# ---------------------------------------------------------------------------
# Section 3 : grid_panels/N1_vs_N2 geometry-correlation scalars.
# The single leading NE correlation is built separately; the groups below are
# every SUN1 scalar that scores signal against nuclear geometry
# (invaginations / convex hull). Pure intensity, texture, synapse, z-position
# and above/below & by-side ratios are deliberately excluded.
# ---------------------------------------------------------------------------
GRID_NE_CORR = "SUN1_NE_Hoechst_corr"       # SUN1 x DNA correlation at the NE (single)
GRID_GROUPS = [
    # Supporting variants of the voxel convex-hull "SUN1 in invaginations"
    # family; the two clean ">1 = grooves" ratios are featured separately above.
    ("SUN1 in the nuclear convex hull — supporting metrics", [
        "SUN1_frac_in_nuc_convex_hull",
        "SUN1_max_distance_to_chull",
        "SUN1_cyto_in_nuc_hull_MFI",
        "SUN1_cyto_in_nuc_hull_sig_fraction",
        "SUN1_invag_within_chull_vs_all_chull_MFI_ratio",
        "SUN1_invag_within_chull_vs_convex_within_chull_MFI_ratio",
    ]),
    ("SUN1 at the deepest invagination", [
        "SUN1_deepest_invag_ratio_edge",
        "SUN1_deepest_invag_ratio_outer_shell",
        "SUN1_ratio_by_deepest_invag_all_0_5um",
        "SUN1_ratio_by_deepest_invag_all_1um",
        "SUN1_ratio_by_deepest_invag_away_0_5um",
        "SUN1_ratio_by_deepest_invag_away_1um",
    ]),
    ("SUN1 invagination-vs-other ratios (per shell)", [
        "SUN1_invag_other_ratio_edge_025um",
        "SUN1_invag_other_ratio_edge_05um",
        "SUN1_invag_other_ratio_edge_1um",
        "SUN1_invag_other_ratio_outer_shell_025um",
        "SUN1_invag_other_ratio_outer_shell_05um",
        "SUN1_invag_other_ratio_outer_shell_1um",
    ]),
]
GRID_MAX_PER_SLIDE = 6      # violins tiled 3 x 2
QC_MAX_COLS = 5

# ---------------------------------------------------------------------------
# Colors / layout (16:9)
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0x55, 0x55, 0x55)
DIVIDER_BG = RGBColor(0xF0, 0xF0, 0xF0)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.10
GAP = 0.16

TITLE_LEFT = MARGIN
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * MARGIN
TITLE_HEIGHT = 0.55
TITLE_FONT_PT = 28

IMG_TOP = 0.66
IMG_BOX_W = SLIDE_W - 2 * MARGIN
IMG_BOX_H = 6.36

FOOTER_TOP = 7.06
FOOTER_WIDTH = SLIDE_W - 2 * MARGIN
FOOTER_HEIGHT = 0.40
FOOTER_FONT_PT = 9

CAPTION_HEIGHT = 0.55      # under a single big image
TILE_CAP_HEIGHT = 0.30     # default under each grid tile


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------
def add_textbox(slide, text, left, top, width, height, font_pt, color,
                bold=False, italic=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def title_font_for(text):
    n = len(text)
    if n <= 52:
        return TITLE_FONT_PT
    if n <= 70:
        return 24
    if n <= 90:
        return 20
    return 18


def add_image_in_box(slide, image_path, box_left, box_top, box_w, box_h):
    """Add an image scaled to fit within (box_w, box_h), centered in the box."""
    pic = slide.shapes.add_picture(
        image_path, Inches(box_left), Inches(box_top), width=Inches(box_w))
    actual_h_in = pic.height / 914400.0
    if actual_h_in > box_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path, Inches(box_left), Inches(box_top), height=Inches(box_h))
        actual_w_in = pic.width / 914400.0
        pic.left = Inches(box_left + (box_w - actual_w_in) / 2)
    else:
        pic.top = Inches(box_top + (box_h - actual_h_in) / 2)
    return pic


def set_slide_background(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, bg)
    return slide


def build_title_slide(prs, title, subtitle):
    slide = _new_slide(prs)
    add_textbox(slide, title, MARGIN, 2.7, SLIDE_W - 2 * MARGIN, 1.3,
                font_pt=38, color=BLACK, bold=True)
    add_textbox(slide, subtitle, MARGIN, 4.1, SLIDE_W - 2 * MARGIN, 1.0,
                font_pt=17, color=GREY, italic=True)


def build_divider_slide(prs, title, subtitle=""):
    slide = _new_slide(prs, bg=DIVIDER_BG)
    add_textbox(slide, title, MARGIN, 3.0, SLIDE_W - 2 * MARGIN, 1.2,
                font_pt=40, color=BLACK, bold=True)
    if subtitle:
        add_textbox(slide, subtitle, MARGIN, 4.25, SLIDE_W - 2 * MARGIN, 0.9,
                    font_pt=18, color=GREY, italic=True)


def build_single_slide(prs, title, image_path, caption, footer=None):
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    cap_h = CAPTION_HEIGHT if caption else 0.0
    img_h = IMG_BOX_H - cap_h
    add_image_in_box(slide, str(image_path), MARGIN, IMG_TOP, IMG_BOX_W, img_h)
    if caption:
        add_textbox(slide, caption, MARGIN, IMG_TOP + img_h, IMG_BOX_W, cap_h,
                    font_pt=12, color=GREY)
    add_textbox(slide, footer or rel_footer(image_path), MARGIN, FOOTER_TOP,
                FOOTER_WIDTH, FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


def build_hero_slide(prs, title, entries, footer_text, cap_h=0.0,
                     cap_pt=11, top_frac=0.58):
    """One "hero" image on top spanning the full content width, the remaining
    images in a single row beneath it. Maximises the lead plot; keeps the
    companion plots side by side and larger than a 3-row stack."""
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    area_top = IMG_TOP
    area_h = FOOTER_TOP - IMG_TOP - 0.02
    top_h = (area_h - GAP) * top_frac
    bot_h = (area_h - GAP) * (1 - top_frac)

    def place(path, cap, left, top, w, h):
        add_image_in_box(slide, str(path), left, top, w, h - cap_h)
        if cap and cap_h:
            add_textbox(slide, cap, left, top + h - cap_h, w, cap_h,
                        font_pt=cap_pt, color=GREY)

    place(entries[0][0], entries[0][1], MARGIN, area_top, IMG_BOX_W, top_h)
    rest = entries[1:]
    if rest:
        nb = len(rest)
        cw = (IMG_BOX_W - (nb - 1) * GAP) / nb
        bot_top = area_top + top_h + GAP
        for i, (path, cap) in enumerate(rest):
            place(path, cap, MARGIN + i * (cw + GAP), bot_top, cw, bot_h)
    add_textbox(slide, footer_text, MARGIN, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


def build_grid_slide(prs, title, entries, footer_text, max_cols,
                     cap_pt=11, cap_h=TILE_CAP_HEIGHT):
    """Tile (image_path, caption) entries in a grid (max_cols=1 => vertical
    stack), each fit + centered in its cell with a caption underneath."""
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    n = len(entries)
    cols = min(max_cols, n)
    rows = math.ceil(n / cols)
    area_top = IMG_TOP
    area_h = FOOTER_TOP - IMG_TOP - 0.02
    cell_w = (SLIDE_W - 2 * MARGIN - (cols - 1) * GAP) / cols
    cell_h = (area_h - (rows - 1) * GAP) / rows
    img_h = cell_h - cap_h
    for i, (path, cap) in enumerate(entries):
        r, c = divmod(i, cols)
        left = MARGIN + c * (cell_w + GAP)
        top = area_top + r * (cell_h + GAP)
        add_image_in_box(slide, str(path), left, top, cell_w, img_h)
        if cap:
            add_textbox(slide, cap, left, top + img_h, cell_w, cap_h,
                        font_pt=cap_pt, color=GREY)
    add_textbox(slide, footer_text, MARGIN, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


def build_row_slide(prs, title, entries, footer_text, cap_h=0.44, cap_pt=11):
    """Place images in a single row scaled to a COMMON height (so panels with
    different aspect ratios line up), filling the content width, with a caption
    under each. Best for related plots that share a pixel height."""
    slide = _new_slide(prs)
    add_textbox(slide, title, TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
                font_pt=title_font_for(title), color=BLACK, bold=True)
    n = len(entries)
    area_top = IMG_TOP
    area_h = FOOTER_TOP - IMG_TOP - 0.02
    box_h = area_h - cap_h
    # Native aspect ratios (add, measure, remove).
    ars = []
    for path, _ in entries:
        pic = slide.shapes.add_picture(str(path), Inches(0), Inches(0))
        ars.append(pic.width / pic.height)
        pic._element.getparent().remove(pic._element)
    avail_w = IMG_BOX_W - (n - 1) * GAP
    h = min(avail_w / sum(ars), box_h)
    widths = [h * ar for ar in ars]
    total_w = sum(widths) + (n - 1) * GAP
    x = MARGIN + (IMG_BOX_W - total_w) / 2.0
    top = area_top + (area_h - (h + cap_h)) / 2.0
    for (path, cap), w in zip(entries, widths):
        slide.shapes.add_picture(str(path), Inches(x), Inches(top), height=Inches(h))
        if cap:
            add_textbox(slide, cap, x, top + h + 0.03, w, cap_h,
                        font_pt=cap_pt, color=GREY)
        x += w + GAP
    add_textbox(slide, footer_text, MARGIN, FOOTER_TOP, FOOTER_WIDTH,
                FOOTER_HEIGHT, font_pt=FOOTER_FONT_PT, color=GREY)
    return slide


# ---------------------------------------------------------------------------
# Content assembly
# ---------------------------------------------------------------------------
def rel_to(path, root):
    try:
        return Path(path).relative_to(root).as_posix()
    except ValueError:
        return Path(path).name


def rel_footer(path):
    for base in (ROOT, SINGLE_ROOT):
        try:
            return Path(path).relative_to(base).as_posix()
        except ValueError:
            pass
    return Path(path).as_posix()


def pretty_scalar(stem):
    """Human caption for a grid violin stem (the plot's own y-axis label is the
    authoritative label; this just tags each tile)."""
    s = stem
    if s.startswith("SUN1_"):
        s = s[len("SUN1_"):]
    s = (s.replace("_025um", " 0.25μm").replace("_05um", " 0.5μm")
          .replace("_1um", " 1μm").replace("_0_5um", " 0.5μm"))
    s = s.replace("Hoechst_corr", "DNA corr").replace("_corr", " corr")
    s = s.replace("_", " ")
    return "SUN1 " + s


def parse_qc(fname):
    """N1_May_18_top1_cell85.png -> (cond, cond_disp, sortkey, caption)."""
    stem = fname[:-4] if fname.lower().endswith(".png") else fname
    m = re.match(r"^(N\d)_(.+?)_(top|low)(\d+)_cell(\d+)$", stem)
    if not m:
        return None
    cond, date, rank, rnum, cell = m.groups()
    cond_disp = "{} ({})".format(cond, date.replace("_", " "))
    sortkey = (0 if rank == "top" else 1, int(rnum))
    caption = "{}{} · cell {}".format(rank, rnum, cell)
    return cond, cond_disp, sortkey, caption


def qc_slides_for(channel, folder):
    """Return list of ("grid", title, entries, footer, cols, cap_pt, cap_h)."""
    qc_dir = folder / "qc_raw"
    if not qc_dir.is_dir():
        return []
    by_cond = {}
    for p in sorted(qc_dir.glob("*.png")):
        parsed = parse_qc(p.name)
        if parsed is None:
            continue
        cond, cond_disp, sortkey, caption = parsed
        by_cond.setdefault((cond, cond_disp), []).append((sortkey, p, caption))
    slides = []
    for (cond, cond_disp), items in sorted(by_cond.items()):
        items.sort(key=lambda t: t[0])
        entries = [(p, cap) for _, p, cap in items]
        title = "{} QC — raw signal at deepest invagination — {}".format(
            channel, cond_disp)
        footer = ("top / low = cells with highest / lowest {ch} in deep NE "
                  "invaginations (÷ cell mean)  ·  {ch} orange, nucleus cyan").format(ch=channel)
        slides.append(("grid", title, entries, footer, QC_MAX_COLS, 10,
                       TILE_CAP_HEIGHT))
    return slides


def channel_items(channel, root):
    """The five combined curated slides for a channel of one compilation. Every
    item is a plan tuple ("grid"/"hero", ...); missing files are dropped from
    tiles (recorded by the caller). max_cols=1 => vertical stack."""
    F = curated_dir(root, channel)
    S = singles_dir(root)
    cf = rel_to(F, root)
    sf = rel_to(S, root)

    def g(name, folder=F):
        return folder / name

    items = [
        # 1 (lead) -- deep invaginations, NE only (boundary singles)
        ("grid", "{} levels and invagination depth — NE only".format(channel), [
            (g("{}_boundary_hulldist_gt0.5um.png".format(channel), S), "hull dist > 0.5 μm"),
            (g("{}_boundary_hulldist_gt1.0um.png".format(channel), S), "hull dist > 1.0 μm"),
        ], sf, 2, 14, 0.30),
        # 2 -- deep invaginations, NE + perinuc 0.5 µm (03a + 03b, stacked)
        ("grid", "{} levels and invagination depth — NE + perinuc 0.5 μm".format(channel), [
            (g("03a_{}_enrichment_hulldist_gt0.5um.png".format(channel)), "hull dist > 0.5 μm"),
            (g("03b_{}_enrichment_hulldist_gt1.0um.png".format(channel)), "hull dist > 1.0 μm"),
        ], cf, 1, 13, 0.26),
        # 3 -- averaged Mean+/-SEM density vs geometry (hull on top, curvatures below)
        ("hero", "{} density vs NE geometry — averaged Mean±SEM".format(channel), [
            (g("02a_{}_avg_hulldist.png".format(channel)), "vs hull-boundary distance"),
            (g("02b_{}_avg_minCurvature_concave.png".format(channel)), "vs min curvature"),
            (g("02c_{}_avg_meanCurvature_concave.png".format(channel)), "vs mean curvature"),
        ], cf, 0.0),
        # 4 -- levels vs curvature (min curv < 0 on top, variants below)
        ("hero", "{} levels and nuclear curvature".format(channel), [
            (g("03c_{}_enrichment_minCurvature_lt0.png".format(channel)), "min curv < 0"),
            (g("03d_{}_enrichment_minCurvature_ltm0.25.png".format(channel)), "min curv < −0.25"),
            (g("03e_{}_enrichment_meanCurvature.png".format(channel)), "mean curv < 0"),
        ], cf, 0.0),
        # 5 -- per-cell correlation vs geometry (hull on top, curvatures below)
        ("hero", "{} per-cell correlation vs NE geometry".format(channel), [
            (g("04a_{}_correlation_hulldist.png".format(channel)), "vs hull-boundary distance"),
            (g("04b_{}_correlation_minCurvature.png".format(channel)), "vs min curvature"),
            (g("04c_{}_correlation_meanCurvature.png".format(channel)), "vs mean curvature"),
        ], cf, 0.0),
    ]
    return items


def filter_tiles(item, missing):
    """Drop tiles whose file is absent; return the item or None if all gone."""
    kind, title, tiles = item[0], item[1], item[2]
    kept = [(p, c) for p, c in tiles if Path(p).exists()]
    for p, _ in tiles:
        if not Path(p).exists():
            missing.append("{}: {}".format(title, Path(p).name))
    if not kept:
        return None
    return (kind, title, kept) + tuple(item[3:])


def grid_entries(stems):
    entries, missing = [], []
    for stem in stems:
        path = GRID_DIR / (stem + ".png")
        if path.exists():
            entries.append((path, pretty_scalar(stem)))
        else:
            missing.append(stem)
    return entries, missing


def chunk(seq, size):
    return [seq[i:i + size] for i in range(0, len(seq), size)]


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    list_only = "--list" in sys.argv

    plan = []
    missing = []

    # Per channel: interleave the 05/18 (N1) and N1-vs-N2 comparison versions of
    # each topic (0518 slide, then comparison slide), then the comparison QC.
    for channel, div_title, div_sub in CHANNEL_SECTIONS:
        plan.append(("divider", div_title, div_sub))
        single_items = channel_items(channel, SINGLE_ROOT)
        compare_items = channel_items(channel, COMPARE_ROOT)
        for s_item, c_item in zip(single_items, compare_items):
            for item, tag in ((s_item, "05/18 (N1)"), (c_item, "N1 vs N2")):
                tagged = (item[0], "{} — {}".format(item[1], tag)) + tuple(item[2:])
                kept = filter_tiles(tagged, missing)
                if kept is not None:
                    plan.append(kept)
        plan.extend(qc_slides_for(channel, curated_dir(COMPARE_ROOT, channel)))

    # Section 3 -- SUN1 invagination-localization scalars + SUN1 × DNA correlation.
    plan.append(("divider", "SUN1 invagination scalars & SUN1 × DNA correlation",
                 "grid_panels / N1_vs_N2  -  N1-vs-N2 comparison violins"))

    # Featured (headline): SUN1 in the nuclear invagination pockets -- the two
    # clean ">1 = SUN1 in grooves" voxel convex-hull ratios (an independent
    # confirmation of the mesh-face density story). "Pocket" = inside the nuclear
    # convex hull but outside the nucleus. Supporting variants sit in the group below.
    pockets = [
        ("SUN1_cyto_in_nuc_hull_vs_near_convex_nuc_MFI_ratio", "grooves ÷ convex surface (0.5 μm shell)"),
        ("SUN1_cyto_in_nuc_hull_vs_all_perinuc_MFI_ratio", "grooves ÷ whole perinuc 0.5 μm shell"),
    ]
    p_entries = []
    for stem, cap in pockets:
        pth = GRID_DIR / (stem + ".png")
        if pth.exists():
            p_entries.append((pth, cap))
        else:
            missing.append("grid: " + stem)
    if p_entries:
        plan.append(("grid", "SUN1 in the nuclear invagination pockets — N1 vs N2", p_entries,
                     "voxel convex-hull decomposition · >1 = perinuclear SUN1 enriched in the grooves "
                     "between the folded NE and its hull · one dot per cell", 2, 13, 0.32))

    ne = GRID_DIR / (GRID_NE_CORR + ".png")
    if ne.exists():
        plan.append(("single", "SUN1 × DNA correlation at the NE — N1 vs N2", ne,
                     "Per-cell area-weighted Pearson r · SUN1 vs DNA at the nuclear envelope.",
                     "grid_panels/N1_vs_N2/" + GRID_NE_CORR + ".png"))
    else:
        missing.append("grid: " + GRID_NE_CORR)
    for group_title, stems in GRID_GROUPS:
        entries, miss = grid_entries(stems)
        missing.extend("grid: " + m for m in miss)
        if not entries:
            continue
        parts = chunk(entries, GRID_MAX_PER_SLIDE)
        for j, part in enumerate(parts):
            suffix = "" if len(parts) == 1 else "  ({}/{})".format(j + 1, len(parts))
            title = "{}  —  N1 vs N2{}".format(group_title, suffix)
            plan.append(("grid", title, part, "grid_panels/N1_vs_N2  ·  one dot per cell",
                         min(3, len(part)), 11, TILE_CAP_HEIGHT))

    # Centrosome metrics (Compile C, N2 alone), grouped by biology. All N2-only,
    # appended near the end. (A) SUN1 on the NE FACING the centrosome: the
    # geom_density/near_cent standard-shell metrics (SUN1 level at NE + perinuc
    # 0.5 µm outside; DNA level at NE + inner shells; SUN1×DNA corr at NE
    # boundary only) + the half-vs-half cent-side/away-side NE ratio.
    nc = COMPILE_C_ROOT / "geom_density" / "near_cent"
    vp = COMPILE_C_ROOT / "violin_plots"
    ne_item = (
        "row", "SUN1 & DNA on the NE facing the centrosome — N2 (May 31)", [
            (nc / "SUN1_near_cent_level.png", "SUN1 level ÷ cell mean · NE, perinuc 0.5 μm"),
            (nc / "DNA_near_cent_level.png", "DNA level ÷ cell mean · NE, inner 0.25 / 0.5 μm"),
            (nc / "SUN1_x_DNA_near_cent_boundary_corr.png", "SUN1×DNA correlation · NE only"),
            (vp / "SUN1_ratio_by_cent_away_0_5um.png", "SUN1 cent-side ÷ away-side · 0.5 μm shell"),
        ], "geom_density/near_cent + violin_plots · N2 alone (05/31/2022) · ref 1 (level/ratio) / 0 (corr)",
        0.44, 10)
    kept = filter_tiles(ne_item, missing)
    if kept is not None:
        plan.append(kept)

    # (B) A DIFFERENT pool -- SUN1 clustered in the cytoplasm near the MTOC (a 3D
    # ball around the centrosome, NOT NE-restricted).
    mtoc_item = (
        "grid", "SUN1 clustered near the centrosome (cytoplasmic MTOC) — N2 (May 31)", [
            (vp / "SUN1_frac_around_cent_2um.png", "fraction of SUN1 within 2 μm of centrosome"),
            (vp / "SUN1_FDD_3D_rel_cent.png", "SUN1 3D spread relative to centrosome (FDD)"),
        ], "violin_plots · N2 alone (05/31/2022) · cytoplasmic 3D ball around the MTOC, not NE-restricted",
        2, 12, 0.30)
    kept = filter_tiles(mtoc_item, missing)
    if kept is not None:
        plan.append(kept)

    # ---- report / dry run ----
    n_slides = 1 + len(plan)  # + title slide
    print("Output: {}".format(OUTPUT_PATH))
    print("{} sections, {} content slides (+ title) = {} total\n".format(
        sum(1 for it in plan if it[0] == "divider"), len(plan), n_slides))
    for it in plan:
        if it[0] == "divider":
            print("\n=== {} ===".format(it[1]))
        elif it[0] == "single":
            print("  [1] {}".format(it[1]))
        elif it[0] in ("grid", "hero", "row"):
            tag = {"hero": "▲", "row": "▭"}.get(it[0], "")
            print("  [{}{}] {}".format(len(it[2]), tag, it[1]))
    if missing:
        print("\nMISSING ({}):".format(len(missing)))
        for m in missing:
            print("  - {}".format(m))
    if list_only:
        return

    # ---- build ----
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_title_slide(prs, DECK_TITLE, DECK_SUBTITLE)
    for it in plan:
        if it[0] == "divider":
            build_divider_slide(prs, it[1], it[2])
        elif it[0] == "single":
            _, title, path, cap = it[:4]
            footer = it[4] if len(it) > 4 else None
            build_single_slide(prs, title, path, cap, footer)
        elif it[0] == "hero":
            _, title, entries, footer, cap_h = it
            build_hero_slide(prs, title, entries, footer, cap_h)
        elif it[0] == "row":
            _, title, entries, footer, cap_h, cap_pt = it
            build_row_slide(prs, title, entries, footer, cap_h, cap_pt)
        elif it[0] == "grid":
            _, title, entries, footer, max_cols, cap_pt, cap_h = it
            build_grid_slide(prs, title, entries, footer, max_cols, cap_pt, cap_h)

    if OUTPUT_PATH.exists():
        backup_dir = OUTPUT_PATH.parent / "backups"
        created = backup_presentation(str(OUTPUT_PATH), backup_base=str(backup_dir))
        if created:
            print("\nBacked up previous deck under: {}".format(backup_dir))

    prs.save(str(OUTPUT_PATH))
    total = len(prs.slides._sldIdLst)
    print("\nDone. {} slides written to:\n  {}".format(total, OUTPUT_PATH))
    if missing:
        print("Skipped {} missing panel(s).".format(len(missing)))


if __name__ == "__main__":
    main()
