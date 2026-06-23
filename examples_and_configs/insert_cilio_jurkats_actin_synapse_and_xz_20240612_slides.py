"""
insert_cilio_jurkats_actin_cent_xz_20240612_slides.py

Cilio-D Jurkats XZ MIP + synapse mask deck for the 06/12/2024 E6-1
Pericentrin/Actin Cilio-D test
(L:/FF/Centrosome_Jurkats/20240612_E6-1_Cilio-D_Test).

Three blocks × two timepoints = 6 slides. Each slide is DMSO (left) vs CilioD
(right). Timepoints:
  30 min αCD3 — W1 DMSO | W2 CilioD
   1 hr αCD3 — W4 DMSO | W3 CilioD

Blocks:
  1. Actin synapse mask (bottom slice)  — actin/bottom_slice_seg
       Scale group 'synapse' — layout-pin only (no embedded scalebar).
  2. Actin XZ MIP, cell top/bottom marked — physical_scale_images/actin_xz
       Scale group 'xz_phys' — 5 µm = 104 px embedded scalebar.
  3. Actin + Cent XZ MIP (no lines) — physical_scale_images/actin_cent_xz_nolines
       Scale group 'xz_phys' — pinned with the actin XZ above so both render
       at the same cm-on-page.

No nucleus channel exists in this experiment.

Usage:
    conda run -n PPT_editing python examples_and_configs/insert_cilio_jurkats_actin_cent_xz_20240612_slides.py
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

EXPERIMENT_ROOT = "L:/FF/Centrosome_Jurkats/20240612_E6-1_Cilio-D_Test"

OUTPUT_PATH = (
    "K:/FF/PPT/PPT_autogeneration/Fixed Jurkats, Miscellaneous/CilioD/"
    "Cilio_Jurkats_actin_synapse_and_xz_20240612.pptx"
)

# Side-by-side condition pairs per timepoint (DMSO left, CilioD right).
TIMEPOINTS = [
    (
        "30 min αCD3",
        "W1_aCD3_E6-1_0p5pcDMSO_30min_AF488Pericentrin_535Actin",  "0.5% DMSO, 30 min αCD3",
        "W2_aCD3_E6-1_50uMCilioD_30min_AF488Pericentrin_535Actin", "50 µM CilioD, 30 min αCD3",
    ),
    (
        "1 hr αCD3",
        "W4_aCD3_E6-1_0p5pcDMSO_1hr_AF488Pericentrin_535Actin",    "0.5% DMSO, 1 hr αCD3",
        "W3_aCD3_E6-1_50uMCilioD_1hr_AF488Pericentrin_535Actin",   "50 µM CilioD, 1 hr αCD3",
    ),
]

# Blocks: each contributes one slide per TIMEPOINT. Title template gets
# .format(tp=...) at build time. subpath uses {cond}.
# prog_fixed_cells is at the EXP root (not under each cond):
#   <root>/prog_fixed_cells/<cond>/<subpath>
BLOCKS = [
    (
        "prog_fixed_cells/{cond}/actin/bottom_slice_raw/montages",
        "Actin at synapse (bottom slice, raw) — DMSO vs CilioD ({tp}, 06/12/2024)",
        "synapse",
    ),
    (
        "prog_fixed_cells/{cond}/actin/bottom_slice_seg/montages",
        "Actin synapse mask (bottom slice, segmented) — DMSO vs CilioD ({tp}, 06/12/2024)",
        "synapse",
    ),
    (
        "prog_fixed_cells/{cond}/physical_scale_images/actin_xz/montages",
        "Actin XZ MIP, cell top/bottom marked — DMSO vs CilioD ({tp}, 06/12/2024)",
        "xz_phys",
    ),
    (
        "prog_fixed_cells/{cond}/physical_scale_images/actin_cent_xz_nolines/montages",
        "Actin + Cent XZ MIP — DMSO vs CilioD ({tp}, 06/12/2024)",
        "xz_phys",
    ),
]

CHUNK_GLOB = "montage_cells_*.png"

# Scalebar invariant — measure with check_scalebar_pixel_widths.py against an
# emitted montage to verify. The CilioD-pipeline convention has been 104 px =
# 5 µm; if this dataset's render differs, adjust SCALEBAR_PX here.
SCALEBAR_PX = 104
SCALEBAR_UM = 5
PPUM_SOURCE = SCALEBAR_PX / SCALEBAR_UM

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = 13.333
SLIDE_H = 7.5

TITLE_LEFT = 0.10
TITLE_TOP = 0.05
TITLE_WIDTH = SLIDE_W - 2 * 0.10
TITLE_HEIGHT = 0.50
TITLE_FONT_PT = 28

GRID_LEFT = 0.10
GRID_TOP = 0.60
CELL_W = 6.50
CELL_H = SLIDE_H - GRID_TOP - 0.10
LABEL_H = 0.30
IMG_H = CELL_H - LABEL_H
LABEL_FONT_PT = 14
COL_GAP = SLIDE_W - 2 * GRID_LEFT - 2 * CELL_W

CELL_POSITIONS = [
    (GRID_LEFT,                    GRID_TOP),
    (GRID_LEFT + CELL_W + COL_GAP, GRID_TOP),
]

# ---------------------------------------------------------------------------


def _winlong(p) -> str:
    s = os.path.abspath(str(p))
    if os.name == "nt" and not s.startswith("\\\\?\\"):
        s = "\\\\?\\" + s.replace("/", "\\")
    return s


def _exists_long(p) -> bool:
    return os.path.exists(_winlong(p))


def _png_dims(path: Path) -> Tuple[int, int]:
    with Image.open(_winlong(path)) as im:
        return im.size


def add_textbox(slide, text, left, top, width, height, font_pt, color, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0]
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def set_slide_background(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def compute_deck_ppi(image_paths: List[Path], max_w_in: float, max_h_in: float) -> float:
    ppi = 0.0
    for p in image_paths:
        w_px, h_px = _png_dims(p)
        ppi = max(ppi, w_px / max_w_in, h_px / max_h_in)
    return ppi


def add_image_at_ppi(slide, image_path: Path, ppi: float,
                     area_left: float, area_top: float,
                     area_w: float, area_h: float):
    w_px, h_px = _png_dims(image_path)
    w_in = w_px / ppi
    h_in = h_px / ppi
    left_in = area_left + (area_w - w_in) / 2
    top_in  = area_top  + (area_h - h_in) / 2
    return slide.shapes.add_picture(
        _winlong(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(w_in),
    )


def _chunk_start_index(p) -> int:
    m = re.match(r"montage_cells_(\d+)", p.name if hasattr(p, "name") else os.path.basename(p))
    return int(m.group(1)) if m else 0


def find_first_chunk(montages_dir: Path) -> Optional[Path]:
    # Long-path-safe: use os.listdir over the \\?\ prefix and join back.
    import fnmatch
    long_dir = _winlong(montages_dir)
    if not os.path.isdir(long_dir):
        return None
    try:
        names = [n for n in os.listdir(long_dir) if fnmatch.fnmatch(n, CHUNK_GLOB)]
    except OSError:
        return None
    if not names:
        return None
    names.sort(key=lambda n: int(re.match(r"montage_cells_(\d+)", n).group(1))
               if re.match(r"montage_cells_(\d+)", n) else 0)
    return montages_dir / names[0]


def build_compare_slide(prs, title_text, left_label, left_img,
                        right_label, right_img, deck_ppi):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, BLACK)

    add_textbox(
        slide, title_text,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT,
        font_pt=TITLE_FONT_PT, color=WHITE, bold=True,
    )

    cells = [
        (left_label,  left_img,  CELL_POSITIONS[0]),
        (right_label, right_img, CELL_POSITIONS[1]),
    ]
    missing = []
    for label, img_path, (cell_left, cell_top) in cells:
        add_textbox(
            slide, label,
            cell_left, cell_top, CELL_W, LABEL_H,
            font_pt=LABEL_FONT_PT, color=WHITE, bold=True,
        )
        if img_path is not None and _exists_long(img_path):
            add_image_at_ppi(slide, img_path, deck_ppi,
                             cell_left, cell_top + LABEL_H, CELL_W, IMG_H)
        else:
            add_textbox(
                slide, "(missing)",
                cell_left, cell_top + LABEL_H + IMG_H / 2 - 0.15, CELL_W, 0.3,
                font_pt=14, color=WHITE,
            )
            missing.append(label)
    return slide, missing


def main() -> None:
    out_path = Path(OUTPUT_PATH)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    root = Path(EXPERIMENT_ROOT)

    # Pre-pass: walk every (block, timepoint) -> collect slide specs.
    # Slides are grouped by block first, then timepoint, so the deck reads:
    # [synapse 30min, synapse 1hr, actin_xz 30min, actin_xz 1hr, ...].
    slide_specs: List[dict] = []
    for subpath_tmpl, title_tmpl, scale_group in BLOCKS:
        for tp_label, l_cond, l_label, r_cond, r_label in TIMEPOINTS:
            l_dir = root / Path(subpath_tmpl.format(cond=l_cond))
            r_dir = root / Path(subpath_tmpl.format(cond=r_cond))
            slide_specs.append({
                "title": title_tmpl.format(tp=tp_label),
                "left_label": l_label,  "left_img": find_first_chunk(l_dir),  "left_dir": l_dir,
                "right_label": r_label, "right_img": find_first_chunk(r_dir), "right_dir": r_dir,
                "scale_group": scale_group,
            })

    # One PPI per scale_group, pinned to the largest source image in that group.
    group_ppi: dict = {}
    for spec in slide_specs:
        imgs = [p for p in (spec["left_img"], spec["right_img"])
                if p is not None and _exists_long(p)]
        own = compute_deck_ppi(imgs, CELL_W, IMG_H) if imgs else 0.0
        sg = spec["scale_group"]
        group_ppi[sg] = max(group_ppi.get(sg, 0.0), own)

    # Fallback for any group that had zero present images.
    for sg in {s["scale_group"] for s in slide_specs}:
        if group_ppi.get(sg, 0.0) <= 0:
            group_ppi[sg] = 400.0

    PHYS_GROUPS = {"xz_phys"}
    print(f"Pinned PPI per scale_group ({len(slide_specs)} slides total):")
    for sg, ppi in sorted(group_ppi.items()):
        if sg in PHYS_GROUPS:
            bar = SCALEBAR_PX / ppi
            note = f"  scalebar = {bar:.3f} in = {bar * 2.54:.3f} cm"
        else:
            note = "  (layout-pin only — no embedded scalebar)"
        print(f"  {sg:>9s}  PPI={ppi:>7.2f}{note}")
    print(f"\nWriting deck to: {OUTPUT_PATH}\n")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    missing_total = []
    for spec in slide_specs:
        slide_ppi = group_ppi[spec["scale_group"]]
        _, missing = build_compare_slide(
            prs, spec["title"],
            spec["left_label"],  spec["left_img"],
            spec["right_label"], spec["right_img"],
            slide_ppi,
        )
        l = "OK" if spec["left_img"] else "MISSING"
        r = "OK" if spec["right_img"] else "MISSING"
        print(f"  [{spec['scale_group']:>7s}]  DMSO:{l}  CilioD:{r}  {spec['title']}")
        for cell in missing:
            src = spec["left_dir"] if cell == spec["left_label"] else spec["right_dir"]
            missing_total.append(f"{spec['title']}/{cell}  ({src})")

    prs.save(str(out_path))
    print(f"\nDone. {len(slide_specs)} slides written to:\n  {out_path}")
    if missing_total:
        print(f"\nMissing ({len(missing_total)}):")
        for m in missing_total:
            print(f"  - {m}")
    else:
        print("\nAll panels found.")


if __name__ == "__main__":
    main()
