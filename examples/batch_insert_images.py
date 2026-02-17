#!/usr/bin/env python
"""
Example: Batch Insert Images into PowerPoint

This script demonstrates how to use a YAML config file to batch-insert
images into a PowerPoint presentation.

Usage:
    python batch_insert_images.py config.yaml

Requirements:
    - PowerPoint file with title slide (slide 1) and template slide (slide 2)
    - Template slide should have one image at the desired position
    - All images listed in the config file should exist
"""

import sys
import os
import yaml
from pptx import Presentation

# Add parent directory to path to import ppt_image_inserter
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from ppt_image_inserter import delete_slide, copy_slide_replace_image, copy_slide_replace_images, backup_presentation, get_all_image_positions


def validate_config(config, config_path):
    """
    Validate configuration file has all required fields with valid values.

    Args:
        config (dict): Parsed configuration dictionary
        config_path (str): Path to config file (for error messages)

    Raises:
        SystemExit: If validation fails, exits with code 1
    """
    required_fields = ['presentation', 'template_slide', 'images']

    # Check required fields exist
    for field in required_fields:
        if field not in config:
            print(f"[ERROR] Required field '{field}' missing in {config_path}")
            print(f"Config must include: {', '.join(required_fields)}")
            sys.exit(1)

    # Validate template_slide is integer
    if not isinstance(config['template_slide'], int):
        print(f"[ERROR] template_slide must be an integer (0-based index)")
        print(f"Found: {config['template_slide']} ({type(config['template_slide']).__name__})")
        sys.exit(1)

    # Validate template_slide is non-negative
    if config['template_slide'] < 0:
        print(f"[ERROR] template_slide must be non-negative (found: {config['template_slide']})")
        sys.exit(1)

    # Validate images is a list
    if not isinstance(config['images'], list):
        print(f"[ERROR] 'images' must be a list")
        sys.exit(1)

    # Validate images list is not empty
    if len(config['images']) == 0:
        print(f"[ERROR] 'images' list is empty - nothing to process")
        sys.exit(1)

    # Validate presentation file exists
    if not os.path.exists(config['presentation']):
        print(f"[ERROR] Presentation file not found: {config['presentation']}")
        sys.exit(1)

    # Validate presentation is .pptx format
    if not config['presentation'].endswith('.pptx'):
        print(f"[ERROR] Presentation must be .pptx format (not .ppt)")
        print(f"Found: {config['presentation']}")
        sys.exit(1)

    # Validate base_dir exists if specified
    if 'base_dir' in config and config['base_dir']:
        if not os.path.isdir(config['base_dir']):
            print(f"[ERROR] base_dir is not a valid directory: {config['base_dir']}")
            sys.exit(1)

    # Validate output_path if specified
    if 'output_path' in config and config['output_path']:
        # Check output directory exists
        output_dir = os.path.dirname(config['output_path'])
        if output_dir and not os.path.exists(output_dir):
            print(f"[ERROR] Output directory does not exist: {output_dir}")
            sys.exit(1)

        # Validate output_path is .pptx format
        if not config['output_path'].endswith('.pptx'):
            print(f"[ERROR] output_path must be .pptx format")
            print(f"Found: {config['output_path']}")
            sys.exit(1)


def main(config_path):
    """Process images according to config file."""

    # Load configuration
    with open(config_path, 'r') as f:
        config = yaml.safe_load(f)

    # Validate configuration
    validate_config(config, config_path)

    ppt_file = config['presentation']
    base_dir = config.get('base_dir', '')
    images = config['images']
    template_slide_index = config.get('template_slide', 1)

    # Get preserve_slides, default to [0, template_slide] if not specified
    preserve_slides = config.get('preserve_slides', [0, template_slide_index])

    # Get backup directory, default to 'backups' subfolder in presentation's directory
    if 'backup_dir' in config:
        backup_dir = config['backup_dir']
    else:
        # Create backups folder in same directory as presentation
        ppt_dir = os.path.dirname(os.path.abspath(ppt_file))
        backup_dir = os.path.join(ppt_dir, 'backups')

    # Handle output_path if specified (template preservation mode)
    output_path = config.get('output_path', None)
    if output_path:
        # Validate output_path
        if os.path.normpath(output_path) == os.path.normpath(ppt_file):
            print(f"[ERROR] output_path cannot be the same as presentation path")
            sys.exit(1)

        # Check output directory exists
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            print(f"[ERROR] Output directory does not exist: {output_dir}")
            sys.exit(1)

        # Backup existing output file before overwriting (if it exists)
        if os.path.exists(output_path):
            print(f"[WARNING] Output file exists and will be overwritten")
            print(f"Creating backup of existing output file in {backup_dir}...")
            try:
                backups_created = backup_presentation(output_path, backup_base=backup_dir)
                print(f"  Backed up to: {', '.join(backups_created.keys())}")
            except Exception as e:
                print(f"[WARNING] Backup failed: {e}")

        # Copy template to output location
        try:
            import shutil
            shutil.copy2(ppt_file, output_path)
        except Exception as e:
            print(f"[ERROR] Failed to copy template: {e}")
            sys.exit(1)

        # Update ppt_file to point to output for all subsequent operations
        ppt_file = output_path

    # Validate template slide exists and has images
    prs = Presentation(ppt_file)
    if template_slide_index >= len(prs.slides):
        print(f"[ERROR] template_slide index {template_slide_index} out of range")
        print(f"Presentation has {len(prs.slides)} slide(s) (indices 0-{len(prs.slides)-1})")
        print(f"Note: Slide indices are 0-based (Slide 1 in UI = index 0)")
        sys.exit(1)

    # Validate template slide has images for auto-position
    try:
        positions = get_all_image_positions(ppt_file, template_slide_index)
        if not positions:
            print(f"[ERROR] Template slide {template_slide_index} has no images")
            print("Add placeholder images to the template slide for auto-positioning")
            sys.exit(1)
        print(f"Template slide validated: {len(positions)} placeholder image(s) found")
    except Exception as e:
        print(f"[ERROR] Could not read template slide: {e}")
        sys.exit(1)

    # Ensure template_slide is in preserve_slides if not explicitly excluded
    if template_slide_index not in preserve_slides:
        print(f"[WARNING] template_slide {template_slide_index} not in preserve_slides")

    # Step 0: Pre-validate ALL images exist before any destructive operations
    print("Validating all images exist...")
    validation_errors = []

    for i, image_spec in enumerate(images):
        if isinstance(image_spec, list):
            # Multi-image case
            for img_filename in image_spec:
                if os.path.isabs(img_filename):
                    img_path = img_filename
                else:
                    img_path = os.path.join(base_dir, img_filename)

                if not os.path.exists(img_path):
                    validation_errors.append(f"Image {i+1} (multi): {img_path}")

        elif isinstance(image_spec, dict):
            # Legacy dict format
            img_path = image_spec['path']
            if not os.path.exists(img_path):
                validation_errors.append(f"Image {i+1} (dict): {img_path}")

        else:
            # Single string
            img_path = os.path.join(base_dir, image_spec) if not os.path.isabs(image_spec) else image_spec
            if not os.path.exists(img_path):
                validation_errors.append(f"Image {i+1}: {img_path}")

    if validation_errors:
        print(f"\n[ERROR] {len(validation_errors)} image(s) not found:")
        for error in validation_errors:
            print(f"  - {error}")
        print("\nNo slides were modified (validation failed before processing)")
        sys.exit(1)

    print(f"All {len(images)} image spec(s) validated successfully\n")

    # Step 1: Delete old slides (except preserved ones)

    prs = Presentation(ppt_file)
    total_slides = len(prs.slides)

    # Collect slides to delete (all except preserved)
    slides_to_delete = [idx for idx in range(total_slides) if idx not in preserve_slides]

    if slides_to_delete:
        # Delete in reverse order to maintain indices
        for slide_idx in reversed(slides_to_delete):
            delete_slide(ppt_file, slide_idx, backup_base=backup_dir)

    # Step 2: Create new slides from images
    # Skip first image (should already be in template slide)
    remaining_images = images[1:]

    print(f"Inserting {len(remaining_images)} images...")

    success_count = 0
    error_count = 0

    for i, image_spec in enumerate(remaining_images):
        # Handle three cases: list (multi-image), dict (legacy), or string (single image)

        if isinstance(image_spec, list):
            # Multiple images per slide (NEW FEATURE)
            if not image_spec:
                # Empty list - skip with warning
                print(f"[WARNING] Empty image list at index {i+1}, skipping")
                continue

            # Build full paths for all images in the list
            image_paths = []
            for img_filename in image_spec:
                if os.path.isabs(img_filename):
                    img_path = img_filename
                else:
                    img_path = os.path.join(base_dir, img_filename)
                image_paths.append(img_path)

            # Validate all images exist
            all_exist = True
            for img_path in image_paths:
                if not os.path.exists(img_path):
                    print(f"[ERROR] Image not found: {img_path}")
                    error_count += 1
                    all_exist = False

            if not all_exist:
                continue  # Skip this slide if any image is missing

            try:
                # Copy template slide and insert multiple images
                new_idx = copy_slide_replace_images(
                    ppt_file,
                    template_slide_index,
                    image_paths,
                    positions=None,  # Auto-detect from template
                    store_metadata=True,
                    add_label=False  # Don't add labels for multi-image slides
                )
                success_count += 1
                print(f"  Created slide with {len(image_paths)} images: {[os.path.basename(p) for p in image_paths]}")
            except Exception as e:
                print(f"[ERROR] Failed on multi-image slide: {e}")
                error_count += 1

        elif isinstance(image_spec, dict):
            # Legacy dict format (single image with metadata)
            image_path = image_spec['path']

            # Check if image exists
            if not os.path.exists(image_path):
                print(f"[ERROR] Image not found: {image_path}")
                error_count += 1
                continue

            try:
                # Copy template slide and insert image
                new_idx = copy_slide_replace_image(
                    ppt_file,
                    template_slide_index,
                    image_path,
                    position=None,  # Auto-detect from template
                    store_metadata=True,
                    add_label=True
                )
                success_count += 1
            except Exception as e:
                print(f"[ERROR] Failed on {os.path.basename(image_path)}: {e}")
                error_count += 1

        else:
            # Single string (backwards compatible - most common case)
            image_path = os.path.join(base_dir, image_spec)

            # Check if image exists
            if not os.path.exists(image_path):
                print(f"[ERROR] Image not found: {image_path}")
                error_count += 1
                continue

            try:
                # Copy template slide and insert image
                new_idx = copy_slide_replace_image(
                    ppt_file,
                    template_slide_index,
                    image_path,
                    position=None,  # Auto-detect from template
                    store_metadata=True,
                    add_label=True
                )
                success_count += 1
            except Exception as e:
                print(f"[ERROR] Failed on {os.path.basename(image_path)}: {e}")
                error_count += 1

    # Summary
    print(f"\nComplete: {success_count}/{len(remaining_images)} slides created")
    if output_path:
        print(f"Output: {ppt_file}")
    if error_count > 0:
        print(f"[WARNING] {error_count} error(s) occurred")

    # Exit with error code if any failures
    if error_count > 0:
        sys.exit(1)


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python batch_insert_images.py <config.yaml>")
        print("\nExample:")
        print("  python batch_insert_images.py example_config.yaml")
        sys.exit(1)

    config_path = sys.argv[1]

    if not os.path.exists(config_path):
        print(f"Error: Config file not found: {config_path}")
        sys.exit(1)

    main(config_path)
